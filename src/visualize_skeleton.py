"""
3D Neuron Skeleton Visualization Module
========================================

This module provides the `VisualizeSkeleton` class for interactive 3D visualization
of neuron skeletons, synapses, and brain region meshes across multiple connectome datasets.

Supported Datasets
------------------
- **NeuPrint datasets**: hemibrain:v1.2.1, optic-lobe:v1.1, manc:v1.0, male-cns:v0.9
- **FlyWire/FAFB datasets**: flywire_FAFB_v783, flywire_BANC_v626

Key Features
------------
- **Multi-dataset support**: Seamlessly work with NeuPrint and FlyWire datasets
- **Skeleton visualization**: Render neurons as meshes or lines with automatic simplification
- **Synapse plotting**: Visualize synapses as scatter points, spheres, cones, or tetrahedrons
- **ROI meshes**: Display brain region meshes with automatic bilateral expansion
- **Soma-aware simplification**: Preserve cell body detail while simplifying branches (FAFB)
- **Extrusion auto-fix**: Detect and replace distorted skeletons automatically (FAFB)
- **Caching system**: Efficient caching for skeletons, meshes, and analysis results
- **Export options**: Save as HTML, PNG images, or video animations

Quick Start
-----------
```python
from coana import VisualizeSkeleton

# Basic neuron visualization
vs = VisualizeSkeleton(
    dataset='hemibrain:v1.2.1',
    neuron_layers=['DA1_lPN'],
)
vs.plot_skeleton()

# FlyWire with automatic extrusion fixes
vs = VisualizeSkeleton(
    dataset='flywire_FAFB_v783',
    neuron_layers=['MTe50'],
    auto_fix_extrusions=True,  # Auto-detect and fix distorted skeletons
)
vs.plot_skeleton()
```

Main Class
----------
VisualizeSkeleton : dataclass
    3D visualization of neuron skeletons with synapses and brain ROI meshes.
    
    Key attributes:
    - dataset: Dataset identifier (e.g., 'hemibrain:v1.2.1', 'flywire_FAFB_v783')
    - neuron_layers: List of neuron types or body IDs to visualize
    - skeleton_mode: 'tube' (mesh) or 'line' rendering
    - auto_fix_extrusions: Automatically detect and replace distorted skeletons (FAFB)
    - cache_neurons: Enable persistent caching for faster subsequent loads
    
    Key methods:
    - plot_skeleton(): Generate 3D visualization
    - list_available_rois(): Show available brain regions for current dataset

FAFB-Specific Features
----------------------
For FlyWire/FAFB datasets, the system includes:
- **Soma-aware mesh simplification**: Preserves cell body detail
  (skeleton_mesh_simplification=0.95, soma_mesh_simplification=0.8)
- **Automatic extrusion detection**: Identifies and replaces distorted skeletons
  (auto_fix_extrusions=True by default)
- **Parquet-based caching**: Efficient storage of extrusion check results
- **API fallback**: Fetches fresh skeletons from CAVE API when needed

Performance Notes
-----------------
- First run downloads and caches data (may take several minutes)
- Subsequent runs use cached data (typically <10 seconds)
- For FAFB: Auto-extrusion detection adds ~30-60s on first run, cached thereafter
- Use skeleton_mesh_simplification=0.95-0.99 for faster rendering of large scenes

See Also
--------
- FAFB_INTEGRATION.md: Detailed guide for FlyWire/FAFB usage
- INSTALLATION.md: Setup and configuration instructions
- OUTPUT_FILES.md: Documentation of output formats and file structure
"""

import os
import sys
import shutil
import time
import signal
import gc
from dataclasses import dataclass, field
from datetime import datetime
from typing import List
import logging
import warnings
from contextlib import contextmanager
from tqdm import tqdm

# Suppress FutureWarning from neuprint about Series.__getitem__
warnings.filterwarnings("ignore", category=FutureWarning, module="neuprint")

# Suppress navis warnings about missing/invalid radii
# These are expected when skeletons have no radius info and are handled automatically
warnings.filterwarnings("ignore", message=".*radii are missing or <= 0.*")
warnings.filterwarnings("ignore", message=".*Mesh will look funny.*")

import numpy as np
import pandas as pd
import cv2
import navis
# Configure navis logger to suppress warnings about missing radii
logging.getLogger("navis").setLevel(logging.ERROR)

import navis.interfaces.neuprint as neu
from neuprint import Client, fetch_synapse_connections, SynapseCriteria, fetch_meta
import plotly.graph_objects as go
import bokeh.palettes

# Import color utilities
from utils.color_utils import (
    standardize_color,
    color_to_hex,
    extract_rgba_tuple,
)


def _configure_roi_mesh_traces(mesh_traces, roi_name):
    """Give one resolved ROI its own Plotly legend entry and trace group."""
    legend_group = f'roi_mesh:{roi_name}'
    display_name = f'brain region [{roi_name}]'

    for trace_index, trace in enumerate(mesh_traces):
        trace.legendgroup = legend_group
        trace.showlegend = trace_index == 0
        trace.name = display_name
        trace.hovertemplate = '<b>%{fullData.name}</b><extra></extra>'
        trace.hoverinfo = 'name'

    return mesh_traces


# Timeout exception for PNG export

def _detect_content_bounds(img, background_color=(255, 255, 255)):
        """
        Detect the bounding box of content in an image.
        
        Parameters
        ----------
        img : PIL.Image
            Input image
        background_color : tuple
            RGB tuple for background color
            
        Returns
        -------
        tuple or None
            (row_min, row_max, col_min, col_max) or None if no content found
        """
        from PIL import Image
        import numpy as np
        
        # Convert to RGB if needed
        if img.mode == 'RGBA':
            bg = Image.new('RGB', img.size, background_color)
            bg.paste(img, mask=img.split()[3])
            img_rgb = bg
        else:
            img_rgb = img.convert('RGB')
        
        arr = np.array(img_rgb)
        
        tolerance = 10
        bg_r, bg_g, bg_b = background_color
        
        non_bg_mask = (
            (np.abs(arr[:, :, 0].astype(int) - bg_r) > tolerance) |
            (np.abs(arr[:, :, 1].astype(int) - bg_g) > tolerance) |
            (np.abs(arr[:, :, 2].astype(int) - bg_b) > tolerance)
        )
        
        rows = np.any(non_bg_mask, axis=1)
        cols = np.any(non_bg_mask, axis=0)
        
        if not rows.any() or not cols.any():
            return None
        
        row_min, row_max = np.where(rows)[0][[0, -1]]
        col_min, col_max = np.where(cols)[0][[0, -1]]
        
        return (row_min, row_max, col_min, col_max)

def _compute_unified_crop_bounds(image_paths, background_color=(255, 255, 255), sample_count=None):
        """
        Compute unified crop bounds across multiple images.
        
        This method finds the maximum extent (union) of content across all images
        to ensure consistent cropping during video rotation.
        
        Parameters
        ----------
        image_paths : list
            List of image file paths to analyze
        background_color : tuple
            RGB tuple for background color
        sample_count : int, optional
            Number of evenly-spaced frames to sample (None = all frames)
            
        Returns
        -------
        tuple or None
            (row_min, row_max, col_min, col_max) unified bounds or None if no content
        """
        from PIL import Image
        
        if not image_paths:
            return None
        
        # Determine which frames to sample
        if sample_count is not None and sample_count < len(image_paths):
            # Sample evenly across all frames
            indices = [int(i * len(image_paths) / sample_count) for i in range(sample_count)]
            sampled_paths = [image_paths[i] for i in indices]
        else:
            sampled_paths = image_paths
        
        # Initialize with None (will expand with each frame)
        unified_row_min = None
        unified_row_max = None
        unified_col_min = None
        unified_col_max = None
        
        for path in sampled_paths:
            try:
                img = Image.open(path)
                bounds = _detect_content_bounds(img, background_color)
                img.close()
                
                if bounds is None:
                    continue
                    
                row_min, row_max, col_min, col_max = bounds
                
                # Expand unified bounds
                if unified_row_min is None:
                    unified_row_min = row_min
                    unified_row_max = row_max
                    unified_col_min = col_min
                    unified_col_max = col_max
                else:
                    unified_row_min = min(unified_row_min, row_min)
                    unified_row_max = max(unified_row_max, row_max)
                    unified_col_min = min(unified_col_min, col_min)
                    unified_col_max = max(unified_col_max, col_max)
            except Exception:
                # Skip frames that can't be read
                continue
        
        if unified_row_min is None:
            return None
            
        return (unified_row_min, unified_row_max, unified_col_min, unified_col_max)

def _apply_consistent_crop(pic_folder, margin=20, background_color=(255, 255, 255)):
        """
        Apply consistent cropping to all images in a folder based on unified bounds.
        
        This method:
        1. Detects content bounds across all frames
        2. Computes the union (maximum extent) of all bounds
        3. Crops all frames to the same unified bounds + margin
        
        Parameters
        ----------
        pic_folder : str
            Folder containing the frame images
        margin : int
            Margin to preserve around content
        background_color : tuple
            RGB tuple for background color
            
        Returns
        -------
        tuple
            (width, height) of the cropped images, or None if failed
        """
        from PIL import Image
        import glob
        
        # Find all JPEG images in folder
        image_paths = sorted(glob.glob(os.path.join(pic_folder, 'deg_*.jpeg')))
        
        if not image_paths:
            return None
        
        # Compute unified bounds across all frames
        # Sample 36 frames (every 10 degrees) for efficiency on long videos
        unified_bounds = _compute_unified_crop_bounds(
            image_paths, background_color, sample_count=min(36, len(image_paths))
        )
        
        if unified_bounds is None:
            return None
        
        row_min, row_max, col_min, col_max = unified_bounds
        
        # Get image dimensions from first image
        with Image.open(image_paths[0]) as img:
            img_height = img.height
            img_width = img.width
        
        # Add margin to unified bounds
        row_min = max(0, row_min - margin)
        row_max = min(img_height - 1, row_max + margin)
        col_min = max(0, col_min - margin)
        col_max = min(img_width - 1, col_max + margin)
        
        final_width = col_max - col_min + 1
        final_height = row_max - row_min + 1
        
        # Apply consistent crop to all images
        for path in image_paths:
            try:
                with Image.open(path) as img:
                    # Crop to unified bounds (no additional margin - already added)
                    cropped = img.crop((col_min, row_min, col_max + 1, row_max + 1))
                    # Convert to RGB for JPEG
                    if cropped.mode == 'RGBA':
                        rgb_img = Image.new('RGB', cropped.size, background_color)
                        rgb_img.paste(cropped, mask=cropped.split()[3])
                        cropped = rgb_img
                    elif cropped.mode != 'RGB':
                        cropped = cropped.convert('RGB')
                    cropped.save(path, 'JPEG', quality=95)
            except Exception:
                # Log but continue
                pass
        
        return (final_width, final_height)

class PNGExportTimeout(Exception):
    """Raised when PNG export times out."""
    pass

def _timeout_handler(signum, frame):
    """Signal handler for PNG export timeout."""
    raise PNGExportTimeout("PNG export timed out")


class WebDriverExportSession:
    """
    Context manager for WebDriver-based exports that keeps the browser open
    for efficient multi-view or video frame exports.
    
    Opens the browser once and allows:
    - Loading HTML once and exporting multiple views by rotating camera
    - Exporting video frames without reopening browser
    
    File Loading
    ------------
    Files are served via a local HTTP server for maximum reliability.
    This approach works with files of any size (tested up to 200MB+).
    
    Export Method
    -------------
    Uses canvas.toDataURL() for fast, high-quality screenshots.
    This method is memory-efficient and produces excellent quality output.
    
    Usage:
    ------
    ```python
    with WebDriverExportSession(width=1200, height=900, scale=2) as session:
        session.load_html('/path/to/figure.html')
        for view_name, camera in cameras.items():
            session.set_camera(camera['eye'], camera['up'], camera.get('center'))
            session.screenshot(f'output_{view_name}.png')
    ```
    """
    
    def __init__(self, width=1200, height=900, scale=2, timeout=300, render_wait=None):
        """
        Initialize WebDriver session parameters.
        
        Parameters
        ----------
        width : int
            Browser viewport width
        height : int
            Browser viewport height  
        scale : int
            Scale factor for high DPI rendering (2 = 144 DPI, 3 = 216 DPI)
        timeout : int
            Timeout in seconds for page load (default 300s for large HTML files)
        render_wait : float, optional
            Fixed render wait time in seconds. If None, auto-calibrated.
        """
        self.width = width
        self.height = height
        self.scale = scale
        self.timeout = timeout
        self.driver = None
        self._loaded_url = None
        self._render_wait = render_wait if render_wait is not None else 0.3
        self._render_wait_fixed = render_wait is not None  # Skip calibration if fixed
        self._initial_camera = None  # Store initial camera from HTML
        self._screenshot_count = 0  # Track screenshot count for periodic memory cleanup
        self._html_size_mb = 0.0  # Track loaded HTML file size for memory management
        self._cleanup_interval = 15  # Default cleanup interval (frames)
        self._http_server = None  # Local HTTP server for large files
        self._http_server_port = None
        
    def __enter__(self):
        """Initialize the WebDriver and return self."""
        try:
            from selenium import webdriver
            from selenium.webdriver.chrome.service import Service as ChromeService
            from selenium.webdriver.chrome.options import Options as ChromeOptions
        except ImportError:
            raise ImportError("selenium is required. Install with: pip install selenium webdriver-manager")
        
        # Scale the window size directly for high-resolution output
        # This is more reliable than deviceScaleFactor for WebGL content
        scaled_width = self.width * self.scale
        scaled_height = self.height * self.scale
        
        chrome_options = ChromeOptions()
        chrome_options.add_argument('--no-sandbox')
        chrome_options.add_argument('--disable-dev-shm-usage')
        chrome_options.add_argument(f'--window-size={scaled_width},{scaled_height}')
        
        # Use new headless mode (Chrome 109+) which supports WebGL better than old headless
        # This replaces the offscreen/position hack which can be flaky on macOS
        chrome_options.add_argument('--headless=new')
        
        # Enable high-quality rendering for WebGL
        chrome_options.add_argument('--use-gl=angle')  # Use ANGLE for better WebGL
        chrome_options.add_argument('--enable-webgl')
        chrome_options.add_argument('--enable-webgl2')
        chrome_options.add_argument('--disable-software-rasterizer')
        
        # Disable animations and GPU compositing to reduce flashing
        chrome_options.add_argument('--disable-gpu-compositing')
        chrome_options.add_argument('--disable-smooth-scrolling')
        # Memory and stability settings for large HTML files (>100MB)
        # --expose-gc enables window.gc() for explicit garbage collection
        chrome_options.add_argument('--js-flags=--max-old-space-size=8192 --expose-gc')
        chrome_options.add_argument('--disable-features=RendererCodeIntegrity')
        chrome_options.add_argument('--disable-backgrounding-occluded-windows')
        chrome_options.add_argument('--disable-renderer-backgrounding')
        chrome_options.add_argument('--memory-pressure-off')
        # Increase shared memory size (default /dev/shm is often too small)
        chrome_options.add_argument('--shm-size=2gb')
        # GPU memory for large WebGL scenes
        chrome_options.add_argument('--enable-gpu-rasterization')
        chrome_options.add_argument('--ignore-gpu-blocklist')
        # Limit GPU memory to prevent runaway allocation
        chrome_options.add_argument('--gpu-memory-buffer-count=4')
        
        # Initialize ChromeDriver using webdriver-manager (cross-platform) or system Chrome
        # webdriver-manager automatically downloads and caches the correct ChromeDriver version
        # Cache location: ~/.wdm/drivers/chromedriver/ (managed automatically)
        try:
            from webdriver_manager.chrome import ChromeDriverManager
            service = ChromeService(ChromeDriverManager().install())
            self.driver = webdriver.Chrome(service=service, options=chrome_options)
        except ImportError:
            # webdriver-manager not installed, try system Chrome directly
            try:
                self.driver = webdriver.Chrome(options=chrome_options)
            except Exception as e:
                raise RuntimeError(
                    f"Could not initialize Chrome WebDriver: {e}\n\n"
                    f"WebDriver export requires:\n"
                    f"  1. Google Chrome browser (version 109+)\n"
                    f"  2. selenium package: pip install selenium\n"
                    f"  3. webdriver-manager package: pip install webdriver-manager\n\n"
                    f"Install with: pip install selenium webdriver-manager\n"
                    f"Or use export_method='kaleido' as fallback."
                )
        except Exception as e:
            # webdriver-manager failed (version mismatch, network error, etc.)
            # Try system Chrome directly
            try:
                self.driver = webdriver.Chrome(options=chrome_options)
            except Exception as e2:
                raise RuntimeError(
                    f"Could not initialize Chrome WebDriver: {e2}\n\n"
                    f"Common causes:\n"
                    f"  - Chrome not installed or not found in PATH\n"
                    f"  - Chrome version mismatch with ChromeDriver\n"
                    f"  - Network error downloading ChromeDriver\n\n"
                    f"Solutions:\n"
                    f"  1. Install/update Google Chrome (version 109+ required for --headless=new)\n"
                    f"  2. Ensure webdriver-manager is installed: pip install webdriver-manager\n"
                    f"  3. Check internet connection for ChromeDriver download\n"
                    f"  4. Use export_method='kaleido' as fallback\n\n"
                    f"Original error: {e}"
                )
        
        # Set up viewport size at initialization using CDP
        # Use scaled dimensions directly for accurate high-resolution WebGL export
        try:
            self.driver.execute_cdp_cmd('Emulation.setDeviceMetricsOverride', {
                'width': self.width * self.scale,
                'height': self.height * self.scale,
                'deviceScaleFactor': 1,  # Use 1 since we already scaled the viewport
                'mobile': False
            })
            self._viewport_configured = True
        except Exception:
            self._viewport_configured = False
        
        return self
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        """Clean up WebDriver and HTTP server."""
        # Stop HTTP server if running
        self._stop_http_server()
        
        if self.driver:
            try:
                self.driver.quit()
            except:
                pass
        return False
    
    def load_html(self, html_path, wait_for_render=True, render_wait=3, background_color='white'):
        """
        Load an HTML file into the browser via local HTTP server.
        
        Uses a local HTTP server to serve the file, which is more reliable than
        document.write() injection and works with files of any size.
        
        Parameters
        ----------
        html_path : str
            Path to the HTML file
        wait_for_render : bool
            If True, wait for Plotly to render
        render_wait : float
            Additional seconds to wait after Plotly detected
        background_color : str
            Background color for the page body (to match plot background)
        """
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC
        from selenium.webdriver.common.by import By
        import time
        
        # Set page load and script timeouts for large HTML files
        self.driver.set_page_load_timeout(self.timeout)
        self.driver.set_script_timeout(self.timeout)
        
        # Track HTML file size for adaptive memory management
        # Larger files require more aggressive cleanup intervals
        self._html_size_mb = os.path.getsize(html_path) / (1024 * 1024)
        if self._html_size_mb > 30:
            # Very large files (>30MB): cleanup every 5 frames
            self._cleanup_interval = 5
        elif self._html_size_mb > 10:
            # Large files (>10MB): cleanup every 10 frames
            self._cleanup_interval = 10
        else:
            # Normal files: cleanup every 15 frames
            self._cleanup_interval = 15
        
        # Always use HTTP server for reliability with any file size
        self._start_http_server_for_file(html_path)
        url = f'http://127.0.0.1:{self._http_server_port}/{os.path.basename(html_path)}'
        self.driver.get(url)
        
        # Ensure no body margin/padding and set background color
        # This prevents white margins around the plot
        js_bg = f"document.body.style.margin = '0'; document.body.style.padding = '0'; document.body.style.backgroundColor = '{background_color}';"
        self.driver.execute_script(js_bg)
        
        self._loaded_url = os.path.abspath(html_path)
        
        if wait_for_render:
            wait = WebDriverWait(self.driver, self.timeout)
            # Wait for js-plotly-plot class (not just 'plotly') as this is added after render
            wait.until(EC.presence_of_element_located((By.CLASS_NAME, "js-plotly-plot")))
            time.sleep(render_wait)
            
            # Re-apply background color styling just in case page load reset it
            self.driver.execute_script(js_bg)
            
            # Set up the viewport and figure size ONCE to avoid flashing on every screenshot
            # This must be done after Plotly is loaded
            self._setup_viewport_size(background_color)
            
            # Force an initial draw to ensure WebGL canvas has content
            # This prevents blank images on first screenshot
            self._force_initial_draw()
            
            # Read initial camera from HTML - this is the "front view" reference
            self._initial_camera = self.get_current_camera()
            
            # Calibrate render wait time based on machine performance (unless fixed)
            if not self._render_wait_fixed:
                self._calibrate_render_wait()
    
    def _start_http_server_for_file(self, html_path):
        """
        Start a local HTTP server to serve a large HTML file.
        
        This is used for files ≥50MB where document.write() injection
        would exceed WebDriver protocol content-length limits.
        
        Parameters
        ----------
        html_path : str
            Path to the HTML file to serve
        """
        import http.server
        import socketserver
        import threading
        
        # Stop any existing server
        self._stop_http_server()
        
        # Get the directory containing the HTML file
        html_dir = os.path.dirname(os.path.abspath(html_path))
        
        # Find an available port
        for port in range(8765, 8865):
            try:
                # Create a handler that serves from the HTML file's directory
                handler = lambda *args, directory=html_dir, **kwargs: \
                    http.server.SimpleHTTPRequestHandler(*args, directory=directory, **kwargs)
                
                # Disable logging to avoid cluttering output
                class QuietHandler(http.server.SimpleHTTPRequestHandler):
                    def __init__(self, *args, **kwargs):
                        super().__init__(*args, directory=html_dir, **kwargs)
                    def log_message(self, format, *args):
                        pass  # Suppress logging
                
                self._http_server = socketserver.TCPServer(("127.0.0.1", port), QuietHandler)
                self._http_server_port = port
                
                # Start server in a daemon thread
                server_thread = threading.Thread(target=self._http_server.serve_forever)
                server_thread.daemon = True
                server_thread.start()
                
                return
            except OSError:
                continue  # Port in use, try next
        
        raise RuntimeError("Could not find an available port for HTTP server")
    
    def _stop_http_server(self):
        """Stop the local HTTP server if running."""
        if self._http_server is not None:
            try:
                self._http_server.shutdown()
                self._http_server.server_close()
            except Exception:
                pass
            self._http_server = None
            self._http_server_port = None

    def _setup_viewport_size(self, background_color='white'):
        """
        Resize the Plotly figure to the full scaled resolution for high-quality export.
        This ensures the WebGL canvas renders at the target resolution.
        """
        import time
        
        # Resize the Plotly figure to fill the viewport at scaled resolution
        scaled_width = self.width * self.scale
        scaled_height = self.height * self.scale
        
        try:
            self.driver.execute_script("""
                var gd = document.querySelector('.js-plotly-plot');
                if (gd && Plotly) {
                    // Resize the figure to fill the viewport
                    Plotly.relayout(gd, {
                        width: arguments[0],
                        height: arguments[1],
                        'paper_bgcolor': arguments[2],
                        'plot_bgcolor': arguments[2]
                    });
                }
            """, scaled_width, scaled_height, background_color)
            
            # Wait for resize to complete
            time.sleep(0.5)
        except Exception:
            pass  # Non-critical, continue with original size
    
    def _force_initial_draw(self):
        """
        Force an initial WebGL draw to ensure canvas has content.
        This prevents blank images on first screenshot due to WebGL buffer clearing.
        """
        try:
            self.driver.execute_script("""
                var gd = document.querySelector('.js-plotly-plot');
                if (gd && gd._fullLayout && gd._fullLayout.scene) {
                    var scene = gd._fullLayout.scene._scene;
                    if (scene && scene.glplot) {
                        // Force WebGL to draw and keep the buffer
                        scene.glplot.draw();
                    }
                }
            """)
        except Exception:
            pass  # Not critical - screenshot will fall back to CDP
    
    def _calibrate_render_wait(self):
        """Calibrate render wait time based on machine performance."""
        import time
        
        if self._initial_camera is None:
            self._render_wait = 0.4
            return
        
        # Measure time for a small camera change
        eye = self._initial_camera.get('eye', {'x': 0, 'y': 0, 'z': -2.5})
        
        t0 = time.time()
        # Make a tiny camera adjustment
        self.driver.execute_script("""
            var gd = document.querySelector('.js-plotly-plot');
            if (gd && Plotly) {
                Plotly.relayout(gd, {'scene.camera.eye.x': arguments[0]});
            }
        """, eye.get('x', 0) + 0.001)
        
        # Wait for render to stabilize
        time.sleep(0.1)
        
        # Restore original
        self.driver.execute_script("""
            var gd = document.querySelector('.js-plotly-plot');
            if (gd && Plotly) {
                Plotly.relayout(gd, {'scene.camera.eye.x': arguments[0]});
            }
        """, eye.get('x', 0))
        
        elapsed = time.time() - t0
        
        # Set render wait based on observed performance
        # Minimum 0.2s, scale up for slower machines
        self._render_wait = max(0.2, min(0.6, elapsed * 1.5))
    
    def get_current_camera(self):
        """
        Get the current camera settings from the Plotly figure.
        
        Returns
        -------
        dict
            Camera settings with 'eye', 'up', 'center' keys
        """
        try:
            camera = self.driver.execute_script("""
                var gd = document.querySelector('.js-plotly-plot');
                if (gd && gd.layout && gd.layout.scene && gd.layout.scene.camera) {
                    return gd.layout.scene.camera;
                }
                return null;
            """)
            return camera if camera else {}
        except Exception:
            return {}
    
    def set_camera(self, eye, up=None, center=None):
        """
        Set the camera position using JavaScript.
        
        Parameters
        ----------
        eye : dict
            Camera eye position, e.g., {'x': 0, 'y': 2.5, 'z': 0}
        up : dict, optional
            Camera up vector, e.g., {'x': 0, 'y': 0, 'z': -1}
        center : dict, optional
            Camera center point
        """
        import time
        import json
        
        # Build camera object
        camera_obj = {}
        if eye:
            camera_obj['eye'] = eye
        if up:
            camera_obj['up'] = up
        if center:
            camera_obj['center'] = center
        
        camera_json = json.dumps(camera_obj)
        
        # Use Plotly.relayout which is the official API for camera updates
        # The key insight: flashing comes from deviceScaleFactor being applied
        # during relayout. We use relayout but ensure it's instant with no animation.
        js_code = f"""
        (function() {{
            var gd = document.querySelector('.js-plotly-plot');
            if (!gd || !Plotly) return;
            
            var camera = {camera_json};
            
            // Use relayout - the official Plotly API for updating camera
            // This is synchronous and doesn't cause resize if we only update camera
            Plotly.relayout(gd, {{'scene.camera': camera}});
        }})();
        """
        self.driver.execute_script(js_code)
        
        # Wait for rendering using calibrated time
        time.sleep(self._render_wait)
    
    def set_camera_for_rotation(self, eye_x, eye_y, eye_z, up_x=0, up_y=-1, up_z=0):
        """
        Set camera for rotation animation (convenience method).
        
        Parameters
        ----------
        eye_x, eye_y, eye_z : float
            Camera eye position
        up_x, up_y, up_z : float
            Camera up vector
        """
        self.set_camera(
            eye={'x': eye_x, 'y': eye_y, 'z': eye_z},
            up={'x': up_x, 'y': up_y, 'z': up_z}
        )
    
    def screenshot(self, output_path, convert_to_jpeg=False, jpeg_quality=100,
                   auto_crop=False, margin=20, background_color=(255, 255, 255),
                   use_webgl_export=True, fast_mode=True, auto_fast_mode=True):
        """
        Take a screenshot and save to file using canvas.toDataURL().
        
        Export method:
        1. WebGL canvas.toDataURL() - Fast, high-quality direct canvas capture
        2. CDP screenshot - Chrome DevTools Protocol screenshot as fallback
        
        Parameters
        ----------
        output_path : str
            Output file path
        convert_to_jpeg : bool
            If True, convert PNG to JPEG (smaller files)
        jpeg_quality : int
            JPEG quality (1-100)
        auto_crop : bool
            If True, automatically crop whitespace/background and add margin
        margin : int
            Margin (in pixels) to preserve around content when auto_crop=True
        background_color : tuple
            RGB tuple for background color detection (default white)
        use_webgl_export : bool, default True
            If True, try WebGL export method for higher quality.
            Falls back to CDP screenshot if WebGL export fails.
        fast_mode : bool, default True
            Legacy parameter, always True. Uses canvas.toDataURL() directly.
        auto_fast_mode : bool, default True
            Legacy parameter, ignored. Canvas mode is always used.
        """
        from PIL import Image
        import base64
        import time
        
        # Track screenshot count for periodic memory cleanup
        self._screenshot_count += 1
        
        # Always use fast_mode (canvas.toDataURL) - it's fast and high quality
        fast_mode = True
        
        img = None
        temp_png = output_path.rsplit('.', 1)[0] + '_temp.png'
        scaled_width = self.width * self.scale
        scaled_height = self.height * self.scale
        
        # Method 1: Try Plotly.toImage() - the official high-quality export API
        # Skip this in fast_mode for better performance
        if use_webgl_export and not fast_mode:
            try:
                # Use Plotly's toImage function which produces high-quality exports
                # This function properly handles WebGL rendering and antialiasing
                # Note: We store in window._lastImageData to allow explicit cleanup
                img_data = self.driver.execute_async_script("""
                    var callback = arguments[arguments.length - 1];
                    var gd = document.querySelector('.js-plotly-plot');
                    if (!gd || !Plotly) {
                        callback(null);
                        return;
                    }
                    
                    // Force a render before capture
                    if (gd._fullLayout && gd._fullLayout.scene && gd._fullLayout.scene._scene) {
                        var scene = gd._fullLayout.scene._scene;
                        if (scene.render) scene.render();
                    }
                    
                    // Use Plotly.toImage for high-quality export
                    Plotly.toImage(gd, {
                        format: 'png',
                        width: arguments[0],
                        height: arguments[1],
                        scale: 1  // We already scaled the dimensions
                    }).then(function(dataUrl) {
                        // Store temporarily for potential cleanup
                        window._lastImageData = dataUrl;
                        callback(dataUrl);
                    }).catch(function(err) {
                        callback(null);
                    });
                """, scaled_width, scaled_height)
                
                if img_data and img_data.startswith('data:image/png;base64,'):
                    base64_data = img_data.split(',')[1]
                    img_bytes = base64.b64decode(base64_data)
                    
                    # Immediately clear the data URL from browser memory
                    try:
                        self.driver.execute_script("window._lastImageData = null;")
                    except:
                        pass
                    
                    # Clear Python reference to large string
                    img_data = None
                    base64_data = None
                    
                    # Verify non-empty image
                    if len(set(img_bytes[:1000])) > 1:
                        with open(temp_png, 'wb') as f:
                            f.write(img_bytes)
                        img = Image.open(temp_png)
                        
                        # Clear the large bytes object
                        img_bytes = None
                        
                        # Verify image has content
                        import numpy as np
                        arr = np.array(img)
                        if arr.max() == 0:  # All black
                            img = None
            except Exception:
                pass  # Will try next method
        
        # Method 2: Try WebGL canvas.toDataURL() with forced render
        if img is None and use_webgl_export:
            try:
                # Force render and immediately capture the WebGL canvas
                canvas_data = self.driver.execute_script("""
                    var gd = document.querySelector('.js-plotly-plot');
                    if (!gd) return null;
                    
                    // Force the scene to render
                    if (gd._fullLayout && gd._fullLayout.scene && gd._fullLayout.scene._scene) {
                        var scene = gd._fullLayout.scene._scene;
                        if (scene.render) scene.render();
                        
                        // Get the canvas
                        if (scene.glplot && scene.glplot.canvas) {
                            return scene.glplot.canvas.toDataURL('image/png');
                        }
                    }
                    
                    // Fallback: try direct canvas access from gl-container
                    var glCanvas = gd.querySelector('.gl-container canvas');
                    if (glCanvas) {
                        return glCanvas.toDataURL('image/png');
                    }
                    
                    return null;
                """)
                
                if canvas_data and canvas_data.startswith('data:image/png;base64,'):
                    base64_data = canvas_data.split(',')[1]
                    img_bytes = base64.b64decode(base64_data)
                    
                    if len(set(img_bytes[:1000])) > 1:
                        with open(temp_png, 'wb') as f:
                            f.write(img_bytes)
                        img = Image.open(temp_png)
                        
                        import numpy as np
                        arr = np.array(img)
                        if arr.max() == 0:
                            img = None
            except Exception:
                pass  # Will fall back to CDP
        
        # Method 3: Fall back to CDP screenshot (captures entire viewport)
        if img is None:
            try:
                # Wait a moment for any pending renders
                time.sleep(0.2)
                
                result = self.driver.execute_cdp_cmd('Page.captureScreenshot', {
                    'format': 'png',
                    'captureBeyondViewport': False,
                    'fromSurface': True,  # Capture from surface for better WebGL support
                    'clip': {
                        'x': 0,
                        'y': 0,
                        'width': scaled_width,
                        'height': scaled_height,
                        'scale': 1
                    }
                })
                screenshot_data = base64.b64decode(result['data'])
                with open(temp_png, 'wb') as f:
                    f.write(screenshot_data)
                img = Image.open(temp_png)
            except Exception:
                # Final fallback to regular screenshot
                self.driver.save_screenshot(temp_png)
                img = Image.open(temp_png)
        
        if auto_crop:
            img = self._auto_crop_image(img, margin, background_color)
        
        if convert_to_jpeg:
            img = img.convert('RGB')
            img.save(output_path, 'JPEG', quality=jpeg_quality)
        else:
            # For PNG output with auto_crop, save directly
            if auto_crop:
                img.save(output_path, 'PNG')
            else:
                # No processing needed, just rename
                os.rename(temp_png, output_path)
                return
        
        # Clean up temp file if it still exists
        if os.path.exists(temp_png):
            os.remove(temp_png)
        
        # Close PIL image to free memory
        if img is not None:
            try:
                img.close()
            except:
                pass
        img = None
        
        # Periodic memory cleanup to prevent Chrome crashes from accumulated WebGL resources
        # Chrome WebGL accumulates texture/buffer memory that isn't freed automatically
        # Interval is adaptive: more frequent for large HTML files (30MB+ = every 5 frames)
        # Note: For very large files (>30MB), using fast_mode=True is recommended
        if self._screenshot_count % self._cleanup_interval == 0:
            self._cleanup_chrome_memory()
            gc.collect()
    
    def _cleanup_chrome_memory(self):
        """
        Clean up Chrome memory to prevent crashes from accumulated WebGL resources.
        
        Chrome's WebGL implementation can accumulate texture and buffer memory
        that isn't automatically garbage collected. This forces cleanup.
        Also clears JavaScript variables that hold large data URLs.
        """
        try:
            # Clear WebGL context caches and force garbage collection in browser
            self.driver.execute_script("""
                // Clear any global variables that might hold image data
                if (window._lastImageData) {
                    window._lastImageData = null;
                }
                
                // Force WebGL context cleanup
                var gd = document.querySelector('.js-plotly-plot');
                if (gd && gd._fullLayout && gd._fullLayout.scene && gd._fullLayout.scene._scene) {
                    var scene = gd._fullLayout.scene._scene;
                    if (scene.glplot) {
                        var gl = scene.glplot.gl;
                        if (gl) {
                            // Flush pending GL commands
                            gl.flush();
                            gl.finish();
                            
                            // Clear WebGL caches by triggering texture/buffer cleanup
                            // This helps release GPU memory
                            var numTextureUnits = gl.getParameter(gl.MAX_COMBINED_TEXTURE_IMAGE_UNITS);
                            for (var i = 0; i < numTextureUnits; i++) {
                                gl.activeTexture(gl.TEXTURE0 + i);
                                gl.bindTexture(gl.TEXTURE_2D, null);
                                gl.bindTexture(gl.TEXTURE_CUBE_MAP, null);
                            }
                            gl.bindBuffer(gl.ARRAY_BUFFER, null);
                            gl.bindBuffer(gl.ELEMENT_ARRAY_BUFFER, null);
                            gl.bindRenderbuffer(gl.RENDERBUFFER, null);
                            gl.bindFramebuffer(gl.FRAMEBUFFER, null);
                        }
                    }
                }
                
                // Clear Plotly's internal image cache if it exists
                if (window.Plotly && window.Plotly.Plots && window.Plotly.Plots.resize) {
                    // Trigger a resize to flush internal caches
                    try {
                        window.Plotly.Plots.resize(gd);
                    } catch(e) {}
                }
                
                // Request browser garbage collection if available (Chrome with --js-flags=--expose-gc)
                if (window.gc) {
                    window.gc();
                }
            """)
        except Exception:
            pass  # Non-critical, continue if cleanup fails
    
    def _auto_crop_image(self, img, margin=20, background_color=(255, 255, 255)):
        """
        Auto-crop image by detecting content boundaries against background.
        
        Parameters
        ----------
        img : PIL.Image
            Input image
        margin : int
            Margin to preserve around content
        background_color : tuple
            RGB tuple for background color
            
        Returns
        -------
        PIL.Image
            Cropped image with margin
        """
        bounds = self._detect_content_bounds(img, background_color)
        if bounds is None:
            return img
        
        row_min, row_max, col_min, col_max = bounds
        
        # Add margin
        row_min = max(0, row_min - margin)
        row_max = min(img.height - 1, row_max + margin)
        col_min = max(0, col_min - margin)
        col_max = min(img.width - 1, col_max + margin)
        
        # Crop the original image (preserve RGBA if present)
        cropped = img.crop((col_min, row_min, col_max + 1, row_max + 1))
        
        return cropped
    
    def _detect_content_bounds(self, img, background_color=(255, 255, 255)):
        return _detect_content_bounds(img, background_color)
    def _compute_unified_crop_bounds(self, image_paths, background_color=(255, 255, 255), 
                                      sample_count=None):
        return _compute_unified_crop_bounds(image_paths, background_color, sample_count)
    def _apply_consistent_crop(self, pic_folder, margin=20, background_color=(255, 255, 255)):
        return _apply_consistent_crop(pic_folder, margin, background_color)
    def set_trace_visibility(self, visible_indices: list, total_traces: int):
        """
        Set visibility of traces by index using JavaScript.
        
        Parameters
        ----------
        visible_indices : list
            List of trace indices to make visible (0-indexed)
        total_traces : int
            Total number of traces in the figure
        """
        import json
        import time
        
        # Build visibility array
        visibility = [False] * total_traces
        for idx in visible_indices:
            if 0 <= idx < total_traces:
                visibility[idx] = True
        
        visibility_json = json.dumps(visibility)
        
        js_code = f"""
        var gd = document.querySelector('.js-plotly-plot');
        if (gd && Plotly) {{
            Plotly.restyle(gd, {{'visible': {visibility_json}}});
        }}
        """
        self.driver.execute_script(js_code)
        time.sleep(0.15)  # Brief wait for restyle
    
    def update_layout(self, layout_update: dict):
        """
        Update layout properties using JavaScript.
        
        Parameters
        ----------
        layout_update : dict
            Dictionary of layout properties to update
        """
        import json
        import time
        
        layout_json = json.dumps(layout_update)
        
        js_code = f"""
        var gd = document.querySelector('.js-plotly-plot');
        if (gd && Plotly) {{
            Plotly.relayout(gd, {layout_json});
        }}
        """
        self.driver.execute_script(js_code)
        time.sleep(0.1)


def export_individuals_webdriver(
    html_path: str,
    output_dir: str,
    legend_entries: dict,
    background_indices: list,
    total_traces: int,
    views: list,
    view_cameras: dict,
    scale: int = 2,
    width: int = 900,
    height: int = 900,
    timeout: int = 60,
    verbose: bool = True,
    auto_crop: bool = True,
    crop_margin: int = 30,
    background_color: str = 'white'
) -> dict:
    """
    Export individual neuron plots efficiently using WebDriver.
    
    Opens the HTML once and toggles trace visibility via JavaScript
    to export each individual neuron/type with all requested views.
    This avoids reopening the browser for each export, significantly
    improving performance for large numbers of individuals.
    
    Parameters
    ----------
    html_path : str
        Path to the main HTML figure file
    output_dir : str
        Directory to save exported PNG files
    legend_entries : dict
        Dictionary mapping legend names to lists of trace indices
        e.g., {'MTe04': [0, 1, 2], 'MTe50': [3, 4, 5]}
    background_indices : list
        List of trace indices to always keep visible (meshes, ROIs)
    total_traces : int
        Total number of traces in the figure
    views : list
        List of view names to export, e.g., ['front', 'top']
    view_cameras : dict
        Dictionary mapping view names to camera settings
        e.g., {'front': {'eye': {...}, 'up': {...}}}
    scale : int, default 2
        Scale factor for export resolution
    width : int, default 900
        Base viewport width
    height : int, default 900
        Base viewport height
    timeout : int, default 60
        Timeout for page load
    verbose : bool, default True
        Whether to print progress messages
    auto_crop : bool, default True
        If True, automatically crop whitespace and preserve margin
    crop_margin : int, default 30
        Margin (in pixels) to preserve around content when auto_crop=True
    
    Returns
    -------
    dict
        Dictionary with results:
        {
            'success': bool,
            'files': {safe_name: [(png_path, view_name), ...]},
            'failed': [legend_name, ...],
            'error': str or None
        }
    """
    from tqdm import tqdm
    
    result = {
        'success': True,
        'files': {},
        'failed': [],
        'error': None
    }
    
    def _sanitize_filename(name):
        """Sanitize legend name for use as filename."""
        safe_name = "".join(c if c.isalnum() or c in '.+_- ' else '_' for c in str(name))
        safe_name = safe_name.strip().replace(' ', '_')
        while '__' in safe_name:
            safe_name = safe_name.replace('__', '_')
        return safe_name.rstrip('_')
    
    # Retry logic for Chrome crashes
    max_retries = 3
    last_error = None
    
    for retry_attempt in range(max_retries):
        if retry_attempt > 0:
            if verbose:
                print(f'   🔄 Retry attempt {retry_attempt + 1}/{max_retries}...')
            import time
            time.sleep(2)  # Brief pause before retry
            # Reset result for retry
            result['files'] = {}
            result['failed'] = []
            result['success'] = True
            result['error'] = None
        
        try:
            with WebDriverExportSession(
                width=width, 
                height=height, 
                scale=scale, 
                timeout=timeout
            ) as session:
                # Load HTML once with background color support
                session.load_html(html_path, wait_for_render=True, render_wait=2, background_color=background_color)
                
                # Hide legend and clean up layout for export
                # Remove interactive UI elements (view selector, controls hint)
                session.update_layout({
                    'showlegend': False,
                    'title.text': '',
                    'margin': {'l': 0, 'r': 0, 'b': 0, 't': 0},
                    'scene.domain': {'x': [0, 1], 'y': [0, 1]},  # Full viewport, no margin
                    'updatemenus': [],  # Remove view selection dropdown
                    'annotations': []   # Remove controls hint text
                })
                
                legend_names = list(legend_entries.keys())
                progress_iter = tqdm(legend_names, desc='Exporting individuals') if verbose else legend_names
                
                for legend_name in progress_iter:
                    trace_indices = legend_entries[legend_name]
                    safe_name = _sanitize_filename(legend_name)
                    
                    # Set visibility: show this legend's traces + background
                    visible_indices = list(trace_indices) + list(background_indices)
                    session.set_trace_visibility(visible_indices, total_traces)
                    
                    result['files'][safe_name] = []
                    
                    # Export each view
                    for view_name in views:
                        camera = view_cameras.get(view_name, {})
                        
                        # Set camera position
                        session.set_camera(
                            eye=camera.get('eye'),
                            up=camera.get('up'),
                            center=camera.get('center')
                        )
                        
                        # Export PNG with auto-crop
                        png_filename = f'{view_name}_{safe_name}.png'
                        png_path = os.path.join(output_dir, png_filename)
                        
                        # Convert background string to RGB tuple for screenshot
                        bg_rgb = (255, 255, 255)
                        if background_color.lower().strip() == 'black':
                            bg_rgb = (0, 0, 0)
                        
                        session.screenshot(png_path, auto_crop=auto_crop, margin=crop_margin, background_color=bg_rgb)
                        
                        # Verify export
                        if os.path.exists(png_path) and os.path.getsize(png_path) > 1024:
                            result['files'][safe_name].append((png_path, view_name))
                        else:
                            result['failed'].append(legend_name)
                            if verbose:
                                print(f'   ⚠️  PNG export failed for {legend_name} ({view_name})')
                            break
            
            # If we got here without exception, export succeeded
            break  # Exit retry loop on success
            
        except Exception as e:
            last_error = e
            error_msg = str(e)
            # Check if this is a Chrome crash
            is_chrome_crash = (
                'Message: \n' in str(e) or 
                error_msg == '' or
                'chrome not reachable' in error_msg.lower() or
                'session deleted' in error_msg.lower() or
                'no such window' in error_msg.lower()
            )
            
            if is_chrome_crash and retry_attempt < max_retries - 1:
                if verbose:
                    print(f'   ⚠️  Chrome crashed unexpectedly. Will retry...')
                continue
            else:
                # Not a crash or out of retries
                result['success'] = False
                result['error'] = str(e)
                if verbose:
                    print(f'   ⚠️  WebDriver export failed: {e}')
                break
    
    return result


# Local imports
try:
    import statvis as sv
    import FAFB_file_converter
    import BANC_file_converter
except ImportError:
    # Fallback for when running from different context
    from . import statvis as sv
    from . import FAFB_file_converter
    from . import BANC_file_converter

@dataclass
class VisualizeSkeleton:
    """
    Interactive 3D visualization of neuron skeletons, synapses, and brain regions.
    
    This class provides comprehensive tools for visualizing connectome data from multiple
    datasets including NeuPrint (hemibrain, male-cns, manc, optic-lobe) and FlyWire/FAFB.
    It handles data fetching, coordinate transformations, mesh simplification, caching,
    and export to various formats.
    
    Key Capabilities
    ----------------
    - **Multi-dataset support**: Seamlessly visualize neurons from different connectome datasets
    - **Flexible neuron selection**: Query by type, instance name, bodyId, or custom layers
    - **Mesh rendering**: High-quality tube meshes with automatic simplification
    - **Synapse visualization**: Multiple modes (scatter, sphere, cone, tetrahedron)
    - **ROI meshes**: Display brain region boundaries with bilateral auto-expansion
    - **Smart caching**: Efficient storage of skeletons, meshes, and analysis results
    - **Export options**: HTML (interactive), PNG (multiple views), video animations
    
    FAFB/FlyWire Features
    ---------------------
    For FlyWire datasets, specialized features handle high-resolution mesh data:
    
    - **Soma-aware simplification**: Preserves cell body detail while simplifying branches
    - **Automatic extrusion detection**: Identifies and replaces distorted skeletons
    - **Parquet caching**: Efficient storage of large-scale analysis results
    - **API fallback**: Fetches fresh meshes from CAVE when ZIP data has artifacts
    
    Attributes
    ----------
    dataset : str, default='hemibrain:v1.2.1'
        Dataset identifier. Supported values:
        - NeuPrint: 'hemibrain:v1.2.1', 'optic-lobe:v1.1', 'manc:v1.0', 'male-cns:v0.9'
        - FlyWire: 'flywire_FAFB_v783', 'flywire_BANC_v626'
    
    neuron_layers : str | list
        Neuron layers to visualize. Can be:
        - List of layers: ['L1', 'L2', 'L3']
        - String with '->' separator: 'L1->L2->L3'
        - Each layer can contain types, instances (regex), bodyIds, or lists thereof
    
    backend : str, default='plotly'
        Visualization backend. Options: 'plotly' (interactive HTML), 'k3d' (WebGL)
    
    skeleton_mode : str, default='tube'
        Rendering mode: 'tube' (mesh with radius), 'line' (simple lines)
    
    auto_fix_extrusions : bool, default=True
        (FAFB only) Automatically detect and replace distorted skeletons.
        Uses edge length analysis to identify extrusion artifacts from aggressive
        simplification. Results cached in parquet format for fast subsequent runs.
    
    skeleton_mesh_simplification : float, optional
        Mesh simplification factor (0.0-1.0). Higher = more simplification.
        Default: 0.95 for FAFB (high-detail), 0.9 for NeuPrint datasets.
        Example: 0.95 removes 95% of faces, keeping 5%.
    
    soma_mesh_simplification : float, default=0.9
        (FAFB only) Gentler simplification for soma region to prevent artifacts.
        Applied within soma_region_radius (default 15µm) of detected soma position.
    
    cache_neurons : bool, default=False
        Enable persistent caching of fetched skeletons to disk.
        Cache location: cache/{dataset}/skeletons/{bodyId}.pkl
    
    mesh_roi : list, default=[]
        List of brain region names to display. ROI names without (L)/(R) suffix
        are auto-expanded to bilateral variants. Use list_available_rois() to see options.
    
    brain_mesh : str, default='none'
        Brain/VNC envelope mesh. Options:
        - 'none': No envelope (only ROIs in mesh_roi)
        - 'template': Dataset's native template (fast, no H5 transforms needed)
        - 'whole': Standard brain mesh (may require H5 transform download)
    
    legend_mode : str, default='layer'
        Controls how neurons appear in the legend. Options:
        - 'single': Each neuron gets its own legend entry ({bodyId}_{layer_name})
        - 'type': Group by neuron type within each layer. If a layer has multiple types,
                  each type gets a separate legend entry.
        - 'layer': Merge all neurons in a layer into one legend entry.
                   Auto-named as {type1}_{type2}_etc if 3+ types present.
    
    expand_colors : str, default='interpolation'
        Method for generating extra colors when more layers than colors available.
        - 'interpolation': Create a smooth colormap and sample extra colors (recommended)
        - 'darken': Recycle colors with progressive darkening (100% to 70% brightness)
        - 'cycle': Simply cycle through colors (color1, color2, ..., color1, color2, ...)
    
    min_synapse_num : int, default=10
        Minimum synapse count threshold for fetching/plotting connections.
    
    synapse_mode : str, default='scatter'
        Synapse rendering mode: 'scatter', 'sphere', 'cone', 'tetrahedron'
    
    skip_synapse : bool, default=False
        Skip all synapse operations for faster initialization and smaller files.
    
    export_views : bool | list, default=True
        PNG export configuration:
        - True: Export all 6 views (front, back, top, bottom, left, right)
        - False: No PNG export
        - List[str]: Specific views to export, e.g., ['front', 'top']
    
    export_scale : int, default=3
        PNG resolution scale factor (1-4). Higher = larger/better quality but slower.
        scale=3 produces ~3600x2700 pixel images.
    
    verbose : bool | str, default='full'
        Verbosity level: 'full' (all messages), 'simple' (essential only), False (silent)
    
    output_dir : str, optional
        Directory for output files. If None, uses data_folder.
    
    include_timestamp : bool, default=True
        Include timestamp in output folder names for unique versioning.
    
    Methods
    -------
    plot_skeleton()
        Generate and display/save the 3D visualization.
    
    list_available_rois(refresh=False, fetch_online=True)
        List all brain regions available for the current dataset.
    
    Examples
    --------
    Basic visualization with hemibrain data:
    
    >>> vs = VisualizeSkeleton(
    ...     dataset='hemibrain:v1.2.1',
    ...     neuron_layers=['DA1_lPN', 'DA1_adPN'],
    ...     mesh_roi=['AL', 'LH', 'MB'],
    ... )
    >>> vs.plot_skeleton()
    
    FlyWire with automatic artifact correction:
    
    >>> vs = VisualizeSkeleton(
    ...     dataset='flywire_FAFB_v783',
    ...     neuron_layers=['MTe50', 'MTe51'],
    ...     auto_fix_extrusions=True,  # Auto-detect distorted neurons
    ...     skeleton_mesh_simplification=0.95,
    ...     soma_mesh_simplification=0.8,  # Gentler on cell bodies
    ... )
    >>> vs.plot_skeleton()
    
    Custom layer mapping from CSV:
    
    >>> vs = VisualizeSkeleton(
    ...     dataset='male-cns:v0.9',
    ...     layer_map_csv='my_layers.csv',  # CSV with 'layer' and 'id_type_instance' columns
    ...     cache_neurons=True,
    ... )
    >>> vs.plot_skeleton()
    
    Video export with multiple views:
    
    >>> vs = VisualizeSkeleton(
    ...     dataset='hemibrain:v1.2.1',
    ...     neuron_layers=['MBON14'],
    ...     export_views=['front', 'top', 'lateral'],
    ...     export_scale=2,
    ... )
    >>> vs.plot_skeleton()
    
    Notes
    -----
    - First run downloads and caches data (may take several minutes)
    - Subsequent runs use cached data (typically <10 seconds)
    - For FAFB: First run with auto_fix_extrusions takes ~30-60s for analysis,
      then results are cached in parquet format
    - Use higher simplification (0.95-0.99) for scenes with many neurons
    - Set legend_mode='single' to see individual neurons in legend
    
    See Also
    --------
    FAFB_INTEGRATION.md : Detailed guide for FlyWire/FAFB usage
    INSTALLATION.md : Setup and configuration instructions
    OUTPUT_FILES.md : Documentation of output formats
    """
    
    backend: str = 'plotly'
    '''
    visualization backend: 'plotly' (default) or 'k3d'
    'plotly': interactive HTML with plotly (good for small/medium scenes)
    'k3d': WebGL-based, faster for large scenes, supports binary export
    '''

    dataset: str = 'hemibrain:v1.2.1'
    '''dataset to use, default is hemibrain:v1.2.1'''

    client_type: str = 'neuprint'
    '''client type: 'neuprint' (default) or 'flywire' '''

    client_flywire: object = None
    '''flywire client instance'''

    client: object = None
    '''neuprint client instance (optional, to reuse existing client)'''

    server: str = 'https://neuprint.janelia.org'
    '''the neuprint server to visit'''
    
    token: str = None
    '''neuprint auth token'''

    version: int | None = None
    '''Materialization version for FlyWire (e.g. 783). If None, uses default/latest.'''

    source_path: str = os.path.dirname(os.path.abspath(__file__))
    '''absolute path to the src/ directory where coana.py is located'''
    
    script_path: str = os.path.dirname(source_path)
    '''absolute path to the project root directory (parent of src/)'''

    data_folder: str = os.path.join(os.path.expanduser('~'), 'connectome_analysis')
    '''
    folder to save all data (subfolders auto-generated based on neuron_layers)
    Default: ~/connectome_analysis/
    '''

    output_dir: str = None
    '''Directory to save output. If None, uses data_folder.'''

    verbose: bool | str = 'full'
    '''
    Verbosity level:
    'full' or True: Print all messages
    'simple': Print only essential progress
    False: Silent
    '''
    
    save_folder: str = ''
    '''
    folder to save the current data (auto-generated from neuron_layers)
    # initialized in __post_init__
    # You can set the "saveas" parameter to customize the folder name'''

    neuron_layers: str | list = ''
    '''
    layers of neurons to plot, can be:
        list of neuron layers: e.g. ['L1', 'L2', 'L3']; or \n
        str of neuron layers separated by '->': \n
        e.g. 'L1->L2->L3'. All type, instance (in regular expression), and bodyId are compatible.\n
    when use list, each layer can be neuron bodyIds, types, instances in regular expressions, or a list of them\n
    e.g. [['L1_0','L1_1'], ['L2_0','L2_1'], ['L3_0','L3_1']]
    '''

    search_columns: str = 'auto'
    '''
    Which columns to search when resolving neuron names (same backend as
    pathfinding, statvis.getNeurons):
        'auto' (default): search all columns with priority
        bodyId -> type -> instance -> other string columns
        (e.g. flywireType, hemibrainType, mancType); 'type', 'instance'
        and 'bodyId' restrict the search to that single column.
    '''

    hemisphere: str = 'both'
    '''
    Restrict the plotted neurons to one hemisphere:
        'both' (default): all neurons.
        'left': only left-hemisphere neurons.
        'right': only right-hemisphere neurons.
    Hemisphere assignment comes from the 'Soma side' / 'hemisphere' columns
    or the instance suffix (_L/_R).  Neurons WITHOUT an explicit hemisphere
    are treated as unclassified ('U') and are ALWAYS included - in 'left',
    'right' AND 'both' - so data that lacks hemisphere notation is never
    silently dropped.
    '''

    custom_layer_names: list = field(default_factory=list)
    '''
    Optional custom names for layers. If empty, auto-generates smart names from neuron types.
    
    Supports partial specification:
    - If list is shorter than neuron_layers, remaining layers use auto-generated names
    - Use empty string '' to skip a layer and use auto-naming for that position
    
    Examples:
        custom_layer_names=['DN1', 'LN', 'Output']  # Full specification for 3 layers
        custom_layer_names=['MyLayer1']             # Only name first layer, auto-name rest
        custom_layer_names=['', '', 'Output']       # Skip first two, name only third layer
    '''

    layer_map_csv: str = None
    '''
    Path to CSV file that defines neuron layers mapping.
    CSV format: columns 'layer' and 'id_type_instance'
    - 'layer': custom layer name (neurons with same layer value are grouped together)
    - 'id_type_instance': neuron identifier (bodyId, type, or instance name)
    
    When provided, this overrides `neuron_layers` and `custom_layer_names`.
    The CSV is parsed to construct layers automatically.
    
    Example CSV:
        layer,id_type_instance
        DN1p,DN1pA
        DN1p,DN1pB
        DN2,DN2
        l-LNv,l-LNv
    
    This creates 3 layers: DN1p (with DN1pA, DN1pB), DN2 (with DN2), l-LNv (with l-LNv)
    '''

    soma_radius_cap: float = None
    '''
    Maximum radius for soma node (in nm) to prevent extrusion artifacts.
    When set, skeleton nodes near the soma with radius > soma_radius_cap will be capped.
    Useful for FAFB skeletons where soma detection may create exaggerated radii.
    Example: soma_radius_cap=2000 caps soma radius to 2 microns
    None (default): No capping, use original skeleton radii
    '''

    smooth_skeleton: bool = False
    '''
    Whether to apply iterative smoothing to skeleton radii.
    When True, applies aggressive smoothing to prevent extrusion artifacts from
    chains of large-radius nodes. Requires soma_radius_cap to be set.
    False (default): Only apply hard cap without smoothing.
    '''

    min_synapse_num: int = 10
    '''minimum number of synapses to fetch and plot'''

    saveas: str = None
    '''filename to save the plot, if an absolute path is given, ignore data_folder'''

    export_views: bool | list = True
    '''
    Which views to export as PNG.
    True: Export all 6 views (front, back, top, bottom, left, right)
    False: Do not export PNGs
    List[str]: List of views to export, e.g. ['front', 'top']
    Default: True (all 6 views)
    '''

    export_scale: int = 3
    '''
    Scale factor for PNG export resolution (1-4 recommended).
    Higher values produce larger, higher-quality images but take longer to export.
    - scale=1: 1200x900 pixels (fast, ~135KB)
    - scale=2: 2400x1800 pixels (~400KB)
    - scale=3: 3600x2700 pixels (default, ~1MB+)
    - scale=4: 4800x3600 pixels (high quality, ~2MB+)
    
    Note: PNG export uses a 60-second timeout per image. If export times out:
    1. Automatic retry with scale=1
    2. If still fails, remaining exports are skipped
    3. Recommendation: increase skeleton_mesh_simplification (e.g., 0.95)
    
    For large figures (>50MB HTML), consider using scale=2 to avoid timeouts.
    For export_video, scale is capped at 3 by default. Use scale=4 explicitly to override.
    '''

    export_method: str = 'webdriver'
    '''
    PNG export method. Options:
    - 'webdriver': Use Selenium WebDriver with Chrome (default, fast and high quality)
    - 'kaleido': Use kaleido engine (slower but reliable fallback if WebDriver fails)
    
    Note: 'webdriver' requires selenium and webdriver-manager packages and Chrome 109+.
    If WebDriver fails, the export automatically falls back to kaleido.
    
    The 'webdriver' method uses canvas.toDataURL() which is fast (~0.3s/frame at scale 1)
    and produces high-quality output. Files are served via local HTTP server for reliability.
    
    Performance Guide:
    ┌───────────────────────────────────────────────────────────────────────┐
    │ Scale   │ webdriver       │ kaleido   │ Recommended                   │
    ├─────────┼─────────────────┼───────────┼───────────────────────────────┤
    │ 1-3     │ 0.3-0.9s ✓      │ 1.6-2.2s  │ webdriver (2-7x faster)       │
    │ 4-10    │ 1.5-2.5s ✓      │ 2.4-3.8s  │ webdriver (faster, stable)    │
    └───────────────────────────────────────────────────────────────────────┘
    
    Recommendations:
    - General use: 'webdriver' (default, fast and reliable)
    - Video export: 'webdriver' with scale≤3 (best throughput)
    - Maximum reliability: 'kaleido' (1.5-3x slower but works without Chrome)
    
    Legacy note: 'webdriver-fast' is still accepted but maps to 'webdriver'.
    '''

    webdriver_render_wait: float = 0
    '''
    Render wait time (seconds) between camera updates in WebDriver export.
    Default: 0 (fastest, works well with modern Chrome).
    Set to None for auto-calibration, or a larger value (e.g., 0.3, 0.5) if
    export aborts mid-way through (not at the beginning).
    '''

    export_timeout: int = 60
    '''
    Timeout in seconds for kaleido-based PNG/video export per frame.
    Default: 60 seconds for first frame, subsequent frames use 2x this value
    if frame 1 exports successfully (verified working).
    
    For large/complex figures, increase this value or use export_method='webdriver'.
    '''

    html_size_cap: int = None
    '''
    HTML file size threshold (in MB) for automatic figure simplification before export.
    When HTML file size exceeds this threshold, the figure is simplified before PNG/video export.
    
    Options:
    - None (default): Auto-determine based on export_method:
        • kaleido: 100 MB (kaleido struggles with large figures)
        • webdriver: 200 MB (Chrome can handle larger files but may still crash)
    - int: Explicit threshold in MB (e.g., 50, 100, 200)
    
    When threshold is exceeded:
    1. First export attempt uses original figure
    2. If export fails (timeout/crash), automatically simplifies and retries
    3. Simplified HTML is saved as {name}_simplified.html for reference
    
    To disable automatic simplification, set a very high value (e.g., 9999).
    '''

    export_simplified_png: bool | int = False
    '''
    Whether to create a simplified figure for PNG export (reduces mesh complexity).
    
    Options:
    - False: Never simplify. Export the full figure directly (default).
            If export times out, will retry with simplified figure.
    - True: Always simplify figures > 50MB HTML before export.
    - int (e.g., 50): Target size threshold in MB. Simplify figures larger than this.
    
    Simplification reduces mesh vertex/face count by ~90% for faster export.
    The original HTML is always preserved at full quality.
    
    Note: Simplification only affects PNG export, not the interactive HTML.
    '''

    include_timestamp: bool = True
    '''Whether to include timestamp in the output folder name. Default True for unique folders.'''

    interactive_html: bool = True
    '''
    Enable interactive controls in exported HTML files.
    
    When True (default):
    - Displays the Plotly mode bar (toolbar) with pan/zoom/rotate buttons
    - Adds a "View" dropdown menu for preset camera angles (Front, Back, Top, Bottom, Left, Right)
    - Enables scroll zoom
    - Shows download button for saving PNG
    - Displays camera angle indicator (azimuth/elevation)
    
    When False:
    - Hides the mode bar for cleaner appearance
    - No view selection dropdown
    - Simpler HTML file (slightly smaller)
    
    Note: Mouse controls (drag to rotate, scroll to zoom) always work regardless of this setting.
    '''

    background_color: str = 'white'
    '''
    Background color for HTML and exported PNG/video views.
    
    Options:
    - 'white' (default): White background, good for publications and printing
    - 'black': Black background, good for presentations and dark mode
    - Any valid CSS color: '#f0f0f0', 'rgb(240,240,240)', 'lightgray', etc.
    
    The background color is applied to:
    - Interactive HTML files
    - Exported PNG images
    - Video frames
    '''

    neuron_colors: tuple | list | str = None
    '''
    Colors for neuron layers. Supports multiple input formats that are automatically
    standardized to rgba format internally.
    
    None (default): auto-pick a categorical palette from the background -
    bokeh Category10 on a white background, bokeh Set3 on a black background.
    
    Supported Formats
    -----------------
    - **Named colors**: 'red', 'blue', 'lightgray', 'darkslategray', etc.
    - **Hex colors**: '#ff0000', '#f00', '#FF0000FF' (with alpha)
    - **RGB tuples**: (255, 0, 0) or (1.0, 0.0, 0.0) (normalized 0-1)
    - **RGBA tuples**: (255, 0, 0, 0.5) or (1.0, 0.0, 0.0, 0.5)
    - **CSS rgb/rgba strings**: 'rgb(255, 0, 0)', 'rgba(255, 0, 0, 0.5)'
    - **Bokeh palettes**: bokeh.palettes.Category10[10], Category20[20], etc.
    - **Matplotlib colormap names**: 'viridis', 'plasma', etc. (requires matplotlib)
    
    Usage
    -----
    - List/tuple of colors: Each color corresponds to a neuron layer
    - Single color: Applied to all layers
    - If fewer colors than layers, colors are expanded via `expand_colors` method
    
    Examples
    --------
    >>> neuron_colors = ['red', 'blue', 'green']  # Named colors
    >>> neuron_colors = ['#ff0000', '#00ff00', '#0000ff']  # Hex colors
    >>> neuron_colors = [(255, 0, 0), (0, 255, 0)]  # RGB tuples
    >>> neuron_colors = [(255, 0, 0, 0.5), (0, 255, 0, 0.8)]  # RGBA tuples
    >>> neuron_colors = bokeh.palettes.Category20[20]  # Bokeh palette
    >>> neuron_colors = 'rgba(255, 0, 0, 0.5)'  # Single RGBA color for all
    '''

    color_mode: str = 'per_layer'
    '''
    Controls how `neuron_colors` are assigned to neurons.

    Options:
    - 'per_layer' (default): One color per layer. All neurons in the same layer share a color.
    - 'per_neuron': One color per neuron across the full selection. Colors are assigned
      in layer order, then neuron order within each layer.

    This is independent of `legend_mode`, which still controls how legend entries are grouped.
    '''

    neuron_alpha: float = 0.2
    '''Alpha (transparency) for neurons. 0.0 = transparent, 1.0 = opaque.
    
    This is a single value applied uniformly to all neuron layers.
    
    **Override Behavior**: If `neuron_colors` contains colors with explicit
    alpha channels (e.g., 'rgba(255,0,0,0.5)', '#ff000080', or (255,0,0,0.5)),
    those alpha values will be used instead and this setting is ignored.
    A warning will be shown when this happens.
    
    **Per-Layer Alpha**: Since `neuron_alpha` only supports a single value,
    use `neuron_colors` with embedded alpha to set different transparencies
    per layer:
        neuron_colors = ['rgba(255,0,0,0.3)', 'rgba(0,255,0,0.7)', 'rgba(0,0,255,0.5)']
    
    Note: Only applies when skeleton_mode='tube'. For skeleton_mode='line',
    alpha is always taken from the color values.
    '''

    synapse_colors: tuple | list | str = bokeh.palettes.Category10[10]
    '''
    Colors for synapse connections between layers. Same format options as neuron_colors.
    
    Supported Formats
    -----------------
    - **Named colors**: 'red', 'blue', 'lightgray', etc.
    - **Hex colors**: '#ff0000', '#f00', '#FF0000FF' (with alpha)
    - **RGB tuples**: (255, 0, 0) or (1.0, 0.0, 0.0) (normalized 0-1)
    - **RGBA tuples**: (255, 0, 0, 0.5) or (1.0, 0.0, 0.0, 0.5)
    - **CSS rgb/rgba strings**: 'rgb(255, 0, 0)', 'rgba(255, 0, 0, 0.5)'
    - **Bokeh palettes**: bokeh.palettes.Category10[10], etc.
    
    Note: Number of synapse colors needed = number of neuron layers - 1
    (for connections between adjacent layers)
    '''

    synapse_size: int | str = 1
    '''
    size of synapse\n
    when synapse_mode='scatter': size in pixels (1–3 recommended)\n
    when synapse_mode='sphere'/'cone'/'tetrahedron': multiplier of the real distance between pre- and post-synaptic sites.\n
    e.g., 1 or 'real' = exact distance size. 2 = 2x distance size.\n
    '''

    synapse_criteria: SynapseCriteria = None
    '''criteria to filter synapses'''

    synapse_mode: str = 'scatter'
    '''
    mode to plot synapses, 'scatter', 'sphere', 'cone', or 'tetrahedron'\n
    'scatter': plot synapses as scatter points, relative size to the view\n
    'sphere': plot synapses as spheres, absolute size in the figure \n
    'cone': plot synapses as cones pointing from pre to post\n
    'tetrahedron': plot synapses as tetrahedrons pointing from pre to post\n
    '''
    
    synapse_alpha: float = 0.6
    '''Alpha (transparency) for synapses. 0.0 = transparent, 1.0 = opaque.
    
    This is a single value applied uniformly to all synapse layers.
    Only works when synapse_mode='sphere'.
    
    **Override Behavior**: If `synapse_colors` contains colors with explicit
    alpha channels (e.g., 'rgba(255,0,0,0.5)', '#ff000080', or (255,0,0,0.5)),
    those alpha values will be used instead and this setting is ignored.
    A warning will be shown when this happens.
    
    **Per-Layer Alpha**: Since `synapse_alpha` only supports a single value,
    use `synapse_colors` with embedded alpha to set different transparencies
    per synapse layer:
        synapse_colors = ['rgba(255,0,0,0.3)', 'rgba(0,255,0,0.7)']
    '''

    mesh_roi: list | str = field(default_factory=list)
    '''
    Meshes of brain ROIs to plot. Accepts a list of ROI names, special keywords,
    regex patterns, or nested lists for color grouping. Also accepts a single string.
    
    Input Formats
    -------------
    - **Single ROI**: 'EB' or ['EB']
    - **Multiple ROIs**: ['LH', 'AL', 'EB']
    - **Special Keywords**:
        - 'primary': All primary brain regions (~50-100 major ROIs)
        - 'all': Every available ROI for the current dataset
    - **Regex patterns**: 'ME.*' matches all ROIs starting with 'ME'
    - **Nested lists for color grouping**: ['AME', ['aL', 'bL', 'gL'], 'EB']
        ROIs in nested lists share the same color but keep separate legend entries.
    
    Auto-Expansion
    --------------
    ROI names without (L)/(R) suffix are automatically expanded to include
    both bilateral variants if available:
    - 'LH' → ['LH(L)', 'LH(R)']
    - 'AL' → ['AL(L)', 'AL(R)']
    - 'EB' → ['EB'] (unpaired, no expansion)
    
    Examples
    --------
    >>> mesh_roi = 'EB'                    # Single ROI
    >>> mesh_roi = ['LH', 'AL', 'EB']      # Multiple ROIs
    >>> mesh_roi = ['primary']             # All primary brain regions
    >>> mesh_roi = ['all']                 # All available ROIs
    >>> mesh_roi = ['ME.*', 'LO.*']        # Regex patterns
    >>> mesh_roi = ['AME', ['aL', 'bL', 'gL'], 'EB']  # Nested for shared color
    >>> mesh_roi = ['.*\\(R\\)']            # All right-hemisphere ROIs
    
    FAFB/FlyWire Note
    -----------------
    FAFB/FlyWire datasets do not have native ROI meshes. When visualizing FAFB data,
    ROI meshes from male-cns are automatically transformed to FAFB coordinates.
    This allows ROI context visualization but may have minor alignment differences.
    
    Finding Available ROIs
    ----------------------
    - Use `vs.list_available_rois()` to query available ROIs programmatically
    - Check the cached list at: `cache/{dataset}/available_rois.json`
    - For FAFB: Uses `cache/male-cns_v0_9/available_rois.json` (transformed)
    
    Common ROIs
    -----------
    - Central Complex (unpaired): EB, FB, PB, NO, AB
    - Mushroom Body: CA, PED, aL, bL, gL, a'L, b'L (bilateral)
    - Optic Lobe: ME, LO, LOP, AME (bilateral)
    - Antennal Lobe: AL (bilateral)
    - Lateral Horn: LH (bilateral)
    
    Set mesh_roi=None or [] to hide all ROI meshes.
    Use brain_mesh parameter to show whole brain/hemibrain envelope.
    '''

    mesh_color: tuple | list | str = (100, 100, 100)
    '''
    Colors for brain ROI meshes. Supports multiple input formats.
    
    Supported Formats
    -----------------
    - **Named colors**: 'red', 'blue', 'lightgray', etc.
    - **Hex colors**: '#ff0000', '#f00', '#FF0000FF' (with alpha)
    - **RGB tuples**: (255, 0, 0) or (1.0, 0.0, 0.0) (normalized 0-1)
    - **RGBA tuples**: (255, 0, 0, 0.5) or (1.0, 0.0, 0.0, 0.5)
    - **CSS rgb/rgba strings**: 'rgb(255, 0, 0)', 'rgba(255, 0, 0, 0.5)'
    - **Bokeh palettes**: bokeh.palettes.Category10[10], etc.
    - **List of colors**: Each color for corresponding ROI in mesh_roi
    
    **Override Behavior**: If `mesh_color` contains colors with explicit
    alpha channels (e.g., 'rgba(100,100,100,0.1)', '#64646419', or (100,100,100,0.1)),
    those alpha values will be used instead and `mesh_alpha` is ignored for those colors.
    A warning will be shown when this happens.
    
    Note: ROI mesh colors are separate from brain_mesh_color and vnc_mesh_color.
    Use brain_mesh_color and vnc_mesh_color to customize the brain/VNC outline meshes.
    
    Examples
    --------
    >>> mesh_color = (100, 100, 100)  # Gray (uses mesh_alpha for transparency)
    >>> mesh_color = 'gray'  # Named color
    >>> mesh_color = 'rgba(100, 0, 100, 0.05)'  # Purple with explicit 5% opacity
    >>> mesh_color = [(255, 0, 0, 0.2), (0, 0, 255, 0.2)]  # Red and blue for 2 ROIs
    '''

    mesh_alpha: float = 0.1
    '''Alpha (transparency) for ROI meshes. 0.0 = transparent, 1.0 = opaque.
    
    This is a single value applied uniformly to all ROI meshes.
    
    **Override Behavior**: If `mesh_color` contains colors with explicit
    alpha channels (e.g., 'rgba(100,100,100,0.1)', '#64646419', or (100,100,100,0.1)),
    those alpha values will be used instead and this setting is ignored.
    A warning will be shown when this happens.
    
    **Per-ROI Alpha**: Since `mesh_alpha` only supports a single value,
    use `mesh_color` with embedded alpha to set different transparencies
    per ROI:
        mesh_color = ['rgba(255,0,0,0.1)', 'rgba(0,255,0,0.2)', 'rgba(0,0,255,0.05)']
    
    Note: This setting only affects ROI meshes (mesh_roi). Brain mesh and VNC mesh
    transparency are controlled by brain_mesh_color and vnc_mesh_color respectively.
    
    Recommended values: 0.05-0.2 for subtle background meshes.
    '''

    legend_mode: str = 'layer'
    '''
    Controls how neurons appear in the legend. Options:
    - 'single': Each neuron gets its own legend entry ({bodyId}_{layer_name}).
                Full detail for identifying individual neurons.
    - 'type': Group by neuron type within each layer. If a layer has multiple
              neuron types, each type gets a separate legend entry.
    - 'layer': Merge all neurons in a layer into one legend entry.
               Auto-named as {type1}_{type2}_etc if 3+ types present.
    '''

    expand_colors: str = 'interpolation'
    '''
    Method for generating extra colors when more layers than colors available.
    - 'interpolation': Create a smooth colormap and sample extra colors (recommended).
                       Ensures visually distinct colors even with many layers.
    - 'darken': Recycle colors with progressive darkening (100% to 70% brightness).
                Same base colors but darker on each cycle.
    - 'cycle': Simple color cycling without modification. Same colors repeat.
    '''

    mirror_on_contralateral: bool = False
    '''
    Whether to mirror neurons and ROIs to the contralateral hemisphere.
    True: Mirror neurons and ROIs (e.g. 'ME(R)' -> 'ME(L)') to the other side.
          Useful for visualizing the full brain structure from hemibrain data.
    False: Only show the original data (default).
    '''

    skeleton_mesh_simplification: float = None
    '''
    Mesh simplification factor for neuron skeletons (0.0 to 1.0).
    Only applies when skeleton_mode='tube'.
    0.0: No simplification (keep all faces).
    0.8: Remove 80% of faces (keep 20%).
    Higher values reduce file size but may lose detail.
    
    Default:
    - FAFB/FlyWire: 0.95 (remove 95% of faces - high detail meshes)
    - NeuPrint datasets: 0.9 (remove 90% of faces)
    
    Recommended: 0.5 - 0.95 for large populations.
    
    NeuPrint note: cached skeletons are already stored at the fixed 0.9
    (simp90) level, so at exactly 0.9 no additional mesh decimation is
    applied. Values ABOVE 0.9 apply the remaining relative reduction
    (e.g., 0.95 keeps 50% of the cached-level faces, 0.99 keeps 10%),
    which is the effective way to shrink output files further.
    Values below 0.9 re-fetch the raw skeletons transiently.
    '''
    
    soma_mesh_simplification: float = 0.9
    '''
    Mesh simplification factor specifically for the soma (cell body) region (0.0 to 1.0).
    Only applies when skeleton_mode='tube' and for FAFB/FlyWire datasets.
    
    The soma region often has high vertex density that can cause extrusion artifacts
    when using high simplification levels. This parameter allows applying gentler 
    simplification to the soma while keeping higher simplification for the skeleton.
    
    None (default): Use skeleton_mesh_simplification for the entire neuron.
    0.8: Recommended for FAFB - reduces soma artifacts while preserving shape.
    
    Example usage:
        skeleton_mesh_simplification=0.95,  # Aggressive on skeleton branches
        soma_mesh_simplification=0.8,       # Gentler on cell body
    
    Note: The soma region is defined as vertices within soma_region_radius (default 15µm)
    of the detected soma position.
    '''
    
    soma_region_radius: float = 15000
    '''
    Radius (in nm) around the soma position to define the soma region for 
    differential simplification. Default is 15000nm (15 microns).
    
    Vertices within this distance from the soma center will use soma_mesh_simplification.
    Vertices outside will use skeleton_mesh_simplification.
    '''

    roi_mesh_simplification: float = 0.9
    '''
    Mesh simplification factor for ROI meshes (0.0 to 1.0).
    0.0: No simplification (keep all faces).
    0.9: Remove 90% of faces (keep 10%).
    Higher values reduce file size but may lose detail.
    Recommended: 0.9 - 0.99 for large ROI meshes.
    '''

    FAFB_template_correction: bool = True
    '''
    Whether to apply tilt correction for FAFB/FlyWire datasets when using 'template' brain mesh.
    
    The FAFB/FlyWire template mesh has a slight tilt relative to the standard view axes.
    When True (default), a rotation correction is applied to align the brain:
    - Z-axis rotation: -4 degrees (corrects left-right tilt in front view)
    - Y-axis rotation: -3 degrees (corrects tilt in top view)
    - Rotation is applied around the brain center to preserve position.
    
    Set to False to use original coordinates (no rotation), which may be preferred if
    aligning with other raw FAFB data or if the rotation causes issues.
    '''

    show_soma: bool = True
    '''whether to show soma'''

    show_fig: bool = True
    '''whether to show the figure'''

    skeleton_mode: str = 'tube'
    '''
    whether to plot the radius of skeleton or only skeleton lines\n
    'tube': plot the radius of skeleton\n
    'line': only plot skeleton lines\n
    when 'line', the file size will be significantly smaller and the rendering will be faster
    '''

    show_connectors: bool = False
    '''whether to fetch and plot the connectors, all pre- and post-synaptic sites of the neurons, for single layer of neurons'''

    skip_synapse: bool = False
    '''
    whether to skip synapse fetching and plotting between layers
    True: skip all synapse operations (faster initialization, smaller file size)
    False: fetch and plot synapses between layers (default behavior)
    Note: This only affects inter-layer synapses, not show_connectors (neuron connectors)
    '''


    
    transforms_dir: str = '~/flybrain-data'
    '''
    Directory for brain transform files (used by flybrains package)\n
    Default: ~/flybrain-data (flybrains default location)\n
    To use a custom location:\n
    1. Set this attribute to your preferred path\n
    2. Ensure the flybrains package uses this path\n
    Note: Changing this requires setting the FLYBRAINS_DATA environment variable\n
    before importing flybrains, or manually moving existing transform files.\n
    '''
    
    cache_neurons: bool = False
    '''
    Whether to cache fetched neuron skeletons to disk\n
    True: Save fetched skeletons as individual {bodyId}.pkl files to cache/{dataset}/skeletons/\n
    False: Fetch from NeuPrint every time (default)\n
    Cache location: cache/{dataset}/skeletons/{bodyId}.pkl\n
    Individual files allow better reuse across different neuron layer selections.\n
    '''
    
    cache_synapses: bool = False
    '''
    Whether to cache fetched synapse data to disk\n
    True: Use synapse table from datasets/{dataset}/*_synapse_table.parquet if available\n
    False: Fetch from NeuPrint every time (default)\n
    For FlyWire/FAFB: Always uses datasets/{dataset}/flywire_FAFB_v783_synapse_table.parquet\n
    '''
    
    force_API_fetching: bool = False
    '''
    Force fetching skeletons from CAVE API instead of downloaded ZIP files (FAFB only).\n
    \n
    True: Fetch meshes from CloudVolume API and skeletonize using navis.\n
          Avoids extrusion artifacts from pre-generated ZIP skeletons.\n
          Results are cached in cache/{dataset}/API_cache/skeletons/\n
    False: Use downloaded ZIP skeletons if available (default).\n
    \n
    Note: API fetching is slower (~5-10s per neuron) but produces cleaner skeletons.\n
    Only applies to FlyWire/FAFB datasets. Has no effect on NeuPrint datasets.\n
    '''
    
    auto_fix_extrusions: bool = True
    '''
    Automatically detect and replace skeletons with extrusion artifacts (FAFB only).\n
    \n
    True: When loading skeletons from ZIP, automatically check for extrusion artifacts\n
          (spiky protrusions from aggressive mesh simplification). If detected, fetch\n
          fresh skeletons from CAVE API to replace them. Results are cached for future use.\n
          This is the default to ensure visual quality.\n
    False: Skip extrusion detection. Faster but may show distorted neurons.\n
    \n
    Extrusion check results are cached in cache/{dataset}/extrusion_check_results.parquet\n
    to avoid repeated analysis. Previously checked neurons are skipped.\n
    \n
    Note: First run with this enabled may take longer due to mesh analysis.\n
    Subsequent runs use cached results and only check new neurons.\n
    Only applies to FlyWire/FAFB datasets. Has no effect on NeuPrint datasets.\n
    '''
    
    brain_mesh: str = 'none'
    ''' 
    Brain/VNC mesh visualization options (dataset-specific):\n
    - 'none': Only plot meshes specified in mesh_roi parameter\n
    - 'template': Plot the dataset's native template mesh (EM resolution)\n
      • hemibrain → JRCFIB2018F (affine transform only, fast)\n
      • optic-lobe → JRCFIB2022M (affine transform only, fast)\n
      • manc → MANC (male adult nerve cord VNC, affine transform only, fast)\n
      • male-cns → JRCFIB2022M (full male CNS: brain + VNC, affine transform only, fast)\n
      • flywire/FAFB → FLYWIRE (native FAFB coordinates, NO transform needed)\n
    - 'whole': Plot standard whole-brain/VNC envelope mesh\n
      • hemibrain → JRC2018F (REQUIRES H5 transforms ~13GB download)\n
      • optic-lobe → JRCFIB2022M (affine transform only, fast)\n
      • manc → MANC VNC envelope (affine transform only, fast)\n
      • male-cns → JRCFIB2022M CNS envelope (affine transform only, fast)\n
      • flywire/FAFB → JRC2018F (REQUIRES H5 transforms, standard female brain)\n
    \n
    ⚡ Performance Tip: Use 'template' for fast visualization with native coordinates.\n
    Only hemibrain and FAFB with brain_mesh='whole' require H5 transform downloads.\n
    \n
    See https://github.com/navis-org/navis-flybrains
    '''
    
    brain_mesh_color: str | tuple | list = 'auto'
    '''
    Color of the brain/VNC mesh. Supports multiple input formats.
    
    Supported Formats
    -----------------
    - **'auto'** (default): Automatically selects optimal color based on background_color:
        • White background: 'rgba(200, 230, 240, 0.1)' (light blue, 10% opacity)
        • Black background: 'rgba(60, 60, 70, 0.1)' (dark gray, 10% opacity)
    - **Named colors**: 'lightblue', 'gray', etc. (alpha can be set separately)
    - **Hex colors**: '#c8e6f0' (use with alpha in tuple form for transparency)
    - **RGB/RGBA tuples**: (200, 230, 240, 0.1) - R,G,B 0-255, alpha 0-1
    - **CSS rgba strings**: 'rgba(200, 230, 240, 0.1)'
    
    Recommendations
    ---------------
    - White background: Light blue/gray with 5-15% opacity
    - Black background: Dark gray with 5-15% opacity to avoid mesh fragment highlights
    
    Examples
    --------
    >>> brain_mesh_color = 'auto'  # Adaptive (default)
    >>> brain_mesh_color = 'rgba(200, 230, 240, 0.1)'  # Light blue, 10% opacity
    >>> brain_mesh_color = (60, 60, 70, 0.1)  # Dark gray tuple
    >>> brain_mesh_color = 'rgba(40, 40, 50, 0.05)'  # Very subtle for dark backgrounds
    '''
    
    vnc_mesh: bool = False
    '''
    Whether to show the VNC (Ventral Nerve Cord) mesh.\n
    Available for datasets with VNC data (requires flybrains >= 0.6.3):\n
    - male-cns → JRCFIB2022M.mesh_vnc (VNC portion of male CNS)\n
    - manc → MANC template (native VNC mesh)\n
    For other datasets (hemibrain, optic-lobe, flywire), this option is ignored.\n
    Note: For MANC with brain_mesh='template', the VNC is already shown\n
    (MANC template IS the VNC, so it ignores vnc_mesh value). 
    Use brain_mesh='none' to hide brain and VNC mesh.\n
    Default: False\n
    '''
    
    vnc_mesh_color: str | tuple | list = 'auto'
    '''
    Color of the VNC mesh. Supports multiple input formats (same as brain_mesh_color).
    
    Supported Formats
    -----------------
    - **'auto'** (default): Automatically selects optimal color based on background_color:
        • White background: 'rgba(200, 230, 240, 0.1)' (light green, 10% opacity)
        • Black background: 'rgba(60, 60, 70, 0.1)' (dark green-gray, 10% opacity)
    - **Named colors**: 'lightgreen', 'gray', etc.
    - **Hex colors**: '#c8f0e6'
    - **RGB/RGBA tuples**: (200, 240, 230, 0.1)
    - **CSS rgba strings**: 'rgba(200, 240, 230, 0.1)'
    
    Note: Default 'auto' uses slightly different hue from brain_mesh_color to distinguish.
    '''

    def _get_effective_mesh_color(self, mesh_type='brain'):
        """
        Get the effective mesh color, resolving 'auto' based on background.
        
        Parameters
        ----------
        mesh_type : str
            'brain' or 'vnc' to select which mesh color to resolve
            
        Returns
        -------
        str
            RGBA color string
        """
        if mesh_type == 'brain':
            color = self.brain_mesh_color
            if color == 'auto':
                if self._is_dark_background():
                    return 'rgba(60, 60, 70, 0.1)'  # Subtle dark gray for dark backgrounds
                else:
                    return 'rgba(200, 230, 240, 0.1)'  # Light blue for light backgrounds
            return color
        else:  # vnc
            color = self.vnc_mesh_color
            if color == 'auto':
                if self._is_dark_background():
                    return 'rgba(60, 60, 70, 0.1)'  # Subtle dark green-gray for dark backgrounds
                else:
                    return 'rgba(200, 230, 240, 0.1)'  # Light green for light backgrounds
            return color

    def list_available_rois(self, refresh=False, fetch_online=True):
        """List all available ROIs for the current dataset.
        
        Parameters
        ----------
        refresh : bool
            If True, force refresh from NeuPrint API. If False, use cached data if available.
        fetch_online : bool
            If True, attempt to fetch from NeuPrint online database. If False, only use local cache.
        
        Returns
        -------
        list
            Sorted list of available ROI names.
        
        Examples
        --------
        >>> vs = VisualizeSkeleton(dataset='hemibrain:v1.2.1', neuron_layers=['EB'])
        >>> available_rois = vs.list_available_rois()
        >>> print(f"Found {len(available_rois)} available ROIs")
        >>> print(available_rois[:10])  # Show first 10 ROIs
        
        >>> # Force refresh from online database
        >>> fresh_rois = vs.list_available_rois(refresh=True, fetch_online=True)
        """
        self._vprint(f'\\n' + '='*70)
        self._vprint(f'Available ROIs for {self.dataset}')
        self._vprint('='*70)
        
        rois = self._get_available_rois(use_cache=not refresh, fetch_online=fetch_online)
        
        if rois:
            self._vprint(f'\\n📊 Total: {len(rois)} ROIs')
            self._vprint(f'\\n🔹 First 30 ROIs:')
            for i in range(0, min(30, len(rois)), 5):
                self._vprint('  ' + ', '.join(rois[i:i+5]))
            if len(rois) > 30:
                self._vprint(f'  ... and {len(rois) - 30} more')
            self._vprint(f'\\n💡 Use these ROI names in the mesh_roi parameter')
            self._vprint('='*70)
        else:
            self._vprint('⚠️  No ROIs found')
            self._vprint('='*70)
        
        return rois
    
    def _vprint(self, msg, level='simple', use_tqdm=False, **kwargs):
        """
        Print message based on verbosity level.
        level: 'simple' (default) or 'full'
        use_tqdm: if True, use tqdm.write() to avoid progress bar conflicts
        """
        if not self.verbose:
            return
        
        # If verbose is 'simple', only print 'simple' messages
        if self.verbose == 'simple' and level == 'full':
            return
            
        # If verbose is 'full', print everything
        if use_tqdm:
            from tqdm import tqdm
            # tqdm.write doesn't support 'end' kwarg, handle it separately
            end = kwargs.pop('end', '\n')
            if end != '\n':
                # For partial lines, just print normally (will be on same line)
                print(msg, end=end, **kwargs)
            else:
                tqdm.write(msg, **kwargs)
        else:
            print(msg, **kwargs)

    def _add_view_selection_menu(self):
        """
        Add interactive view selection dropdown and camera angle display to the figure.
        
        Adds a dropdown menu with preset camera positions (Front, Back, Top, Bottom, Left, Right)
        and annotations showing current camera angles. The camera presets are adjusted based on
        the dataset coordinate system.
        """
        # Define camera positions based on dataset
        # Default: Standard fly brain - X: Left-Right, Y: Dorsal-Ventral, Z: Anterior-Posterior
        view_cameras = {
            'Front': dict(eye=dict(x=0, y=0, z=-2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
            'Back': dict(eye=dict(x=0, y=0, z=2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
            'Top': dict(eye=dict(x=0, y=-2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=1)),
            'Bottom': dict(eye=dict(x=0, y=2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
            'Left': dict(eye=dict(x=-2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
            'Right': dict(eye=dict(x=2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
        }
        
        # Adjust for MANC (Male Adult Nerve Cord)
        if 'manc' in self.dataset.lower():
            view_cameras = {
                'Front': dict(eye=dict(x=0, y=0, z=2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                'Back': dict(eye=dict(x=0, y=0, z=-2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                'Top': dict(eye=dict(x=0, y=-2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=1)),
                'Bottom': dict(eye=dict(x=0, y=2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=1)),
                'Left': dict(eye=dict(x=-2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                'Right': dict(eye=dict(x=2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
            }
        
        # Adjust for hemibrain template (JRCFIB2018F)
        if 'hemibrain' in self.dataset.lower() and self.brain_mesh == 'template':
            view_cameras = {
                'Front': dict(eye=dict(x=0, y=2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
                'Back': dict(eye=dict(x=0, y=-2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
                'Top': dict(eye=dict(x=0, y=0, z=-2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=1, z=0)),
                'Bottom': dict(eye=dict(x=0, y=0, z=2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=1, z=0)),
                'Left': dict(eye=dict(x=-2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
                'Right': dict(eye=dict(x=2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
            }
        
        # Create dropdown buttons for view selection
        view_buttons = []
        for view_name, camera in view_cameras.items():
            view_buttons.append(
                dict(
                    args=[{'scene.camera': camera}],
                    label=view_name,
                    method='relayout'
                )
            )
        
        # Determine text color based on background
        is_dark_bg = self._is_dark_background()
        text_color = 'white' if is_dark_bg else 'black'
        hint_color = 'lightgray' if is_dark_bg else 'gray'
        dropdown_bg = 'rgba(50,50,50,0.9)' if is_dark_bg else 'rgba(255,255,255,0.9)'
        hint_bg = 'rgba(50,50,50,0.7)' if is_dark_bg else 'rgba(255,255,255,0.7)'
        
        # Add dropdown menu for view selection
        self.fig_3d.update_layout(
            updatemenus=[
                dict(
                    type='dropdown',
                    showactive=True,
                    active=0,  # Default to Front view
                    buttons=view_buttons,
                    x=0.07,  # Offset to make room for title
                    y=0.98,
                    xanchor='left',
                    yanchor='top',
                    bgcolor=dropdown_bg,
                    bordercolor='rgba(128,128,128,0.5)',
                    borderwidth=1,
                    font=dict(size=12, color=text_color),
                    pad=dict(l=5, r=5, t=5, b=5),
                )
            ],
            # Add annotations: title for dropdown and controls hint
            annotations=[
                # Title for view selection dropdown
                dict(
                    text="<b>View:</b>",
                    x=0.01, y=0.975,
                    xref='paper', yref='paper',
                    xanchor='left', yanchor='top',
                    showarrow=False,
                    font=dict(size=12, color=text_color),
                ),
                # Controls hint at bottom
                dict(
                    text="🖱️ Drag: Rotate | Scroll: Zoom | Ctrl+Drag: Pan",
                    x=0.5, y=0.01,
                    xref='paper', yref='paper',
                    xanchor='center', yanchor='bottom',
                    showarrow=False,
                    font=dict(size=10, color=hint_color),
                    bgcolor=hint_bg,
                    borderpad=3,
                )
            ],
        )

    def _is_dark_background(self, color=None):
        """
        Check if the background color is dark.
        
        Parameters
        ----------
        color : str, optional
            Color to check. If None, uses self.background_color.
        
        Returns True if the background color is considered dark (for text contrast).
        """
        bg = (color if color is not None else self.background_color).lower().strip()
        
        # Common dark color names
        dark_colors = {'black', 'darkgray', 'darkgrey', 'dimgray', 'dimgrey', 
                       'gray', 'grey', 'darkblue', 'navy', 'darkgreen', 'darkred',
                       'maroon', 'purple', 'indigo', 'midnightblue', 'darkslategray',
                       'darkslategrey', 'darkolivegreen', 'darkmagenta', 'darkcyan'}
        
        if bg in dark_colors:
            return True
        
        # Check hex colors
        if bg.startswith('#'):
            try:
                hex_color = bg.lstrip('#')
                if len(hex_color) == 3:
                    hex_color = ''.join([c*2 for c in hex_color])
                r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
                # Calculate luminance
                luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
                return luminance < 0.5
            except:
                pass
        
        # Check rgb/rgba colors
        if bg.startswith('rgb'):
            try:
                import re
                numbers = re.findall(r'[\d.]+', bg)
                r, g, b = float(numbers[0]), float(numbers[1]), float(numbers[2])
                if r <= 1 and g <= 1 and b <= 1:  # Normalized
                    r, g, b = r * 255, g * 255, b * 255
                luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
                return luminance < 0.5
            except:
                pass
        
        # Default: assume light background
        return False

    def _get_html_size_cap(self):
        """
        Get the effective HTML size cap in MB for figure simplification.
        
        Returns the configured html_size_cap, or auto-determines based on export_method:
        - kaleido: 100 MB (kaleido struggles with large figures)
        - webdriver: 200 MB (Chrome can handle larger files)
        
        Returns
        -------
        int
            HTML size threshold in MB
        """
        if self.html_size_cap is not None:
            return self.html_size_cap
        
        # Auto-determine based on export method
        if self.export_method in ('webdriver', 'webdriver-fast'):
            return 200
        else:  # kaleido
            return 100

    def _simplify_mesh_open3d(self, trimesh_obj, target_faces):
        """
        Simplify a trimesh mesh using open3d's quadric decimation.
        
        This uses open3d directly instead of trimesh's wrapper because:
        - trimesh 4.x uses fast_simplification which achieves ~60% max reduction on tube meshes
        - open3d achieves exact target face count (true 90%+ reduction)
        
        Parameters
        ----------
        trimesh_obj : trimesh.Trimesh
            The mesh to simplify.
        target_faces : int
            Target number of faces after simplification.
            
        Returns
        -------
        trimesh.Trimesh
            Simplified mesh, or original if simplification fails.
        """
        try:
            import open3d as o3d
            import trimesh
            import numpy as np
            
            # Convert trimesh to open3d
            o3d_mesh = o3d.geometry.TriangleMesh()
            o3d_mesh.vertices = o3d.utility.Vector3dVector(trimesh_obj.vertices)
            o3d_mesh.triangles = o3d.utility.Vector3iVector(trimesh_obj.faces)
            
            # Simplify with open3d (achieves exact target)
            simplified = o3d_mesh.simplify_quadric_decimation(int(target_faces))
            
            # Convert back to trimesh
            return trimesh.Trimesh(
                vertices=np.asarray(simplified.vertices),
                faces=np.asarray(simplified.triangles)
            )
        except ImportError:
            # Fallback to trimesh if open3d not available
            try:
                return trimesh_obj.simplify_quadric_decimation(face_count=target_faces)
            except Exception:
                return trimesh_obj
        except Exception:
            # Return original on any error
            return trimesh_obj

    def _simplify_mesh_with_soma_awareness(self, trimesh_obj, skeleton_simp, soma_simp, soma_pos, soma_radius=15000):
        """
        Simplify a mesh with different simplification levels for soma vs skeleton regions.
        
        This method splits the mesh into soma region and skeleton region, applies different
        simplification levels to each, then recombines them. This helps prevent extrusion
        artifacts around the cell body that can occur with aggressive simplification.
        
        Parameters
        ----------
        trimesh_obj : trimesh.Trimesh
            The mesh to simplify.
        skeleton_simp : float
            Simplification factor for skeleton region (0.0-1.0, e.g., 0.95 = remove 95% faces).
        soma_simp : float
            Simplification factor for soma region (0.0-1.0, e.g., 0.8 = remove 80% faces).
        soma_pos : array-like
            [x, y, z] position of the soma center in the same units as the mesh.
        soma_radius : float
            Radius around soma_pos defining the soma region (default 15000nm = 15µm).
            
        Returns
        -------
        trimesh.Trimesh
            Combined simplified mesh with region-specific simplification.
        """
        import trimesh
        import numpy as np
        
        if soma_pos is None or len(soma_pos) == 0:
            # No soma info, fall back to uniform simplification
            n_faces = len(trimesh_obj.faces)
            target_faces = max(100, int(n_faces * (1 - skeleton_simp)))
            return self._simplify_mesh_open3d(trimesh_obj, target_faces)
        
        # Ensure soma_pos is a 1D array
        soma_pos = np.array(soma_pos).flatten()[:3]
        
        # Calculate face centroids
        vertices = trimesh_obj.vertices
        faces = trimesh_obj.faces
        face_centroids = vertices[faces].mean(axis=1)
        
        # Calculate distance from each face centroid to soma
        distances = np.linalg.norm(face_centroids - soma_pos, axis=1)
        
        # Split faces into soma region and skeleton region
        soma_mask = distances <= soma_radius
        skeleton_mask = ~soma_mask
        
        soma_faces = faces[soma_mask]
        skeleton_faces = faces[skeleton_mask]
        
        # If one region is empty, simplify the whole mesh with appropriate level
        if len(soma_faces) == 0:
            n_faces = len(faces)
            target_faces = max(100, int(n_faces * (1 - skeleton_simp)))
            return self._simplify_mesh_open3d(trimesh_obj, target_faces)
        
        if len(skeleton_faces) == 0:
            n_faces = len(faces)
            target_faces = max(100, int(n_faces * (1 - soma_simp)))
            return self._simplify_mesh_open3d(trimesh_obj, target_faces)
        
        # Create sub-meshes for each region
        # We need to keep all vertices and just select faces
        try:
            # Create soma sub-mesh
            soma_mesh = trimesh.Trimesh(vertices=vertices, faces=soma_faces, process=False)
            # Remove unreferenced vertices to clean up
            soma_mesh.remove_unreferenced_vertices()
            
            # Create skeleton sub-mesh  
            skeleton_mesh = trimesh.Trimesh(vertices=vertices, faces=skeleton_faces, process=False)
            skeleton_mesh.remove_unreferenced_vertices()
            
            # Simplify each region with its target level
            soma_target = max(50, int(len(soma_mesh.faces) * (1 - soma_simp)))
            skeleton_target = max(50, int(len(skeleton_mesh.faces) * (1 - skeleton_simp)))
            
            simplified_soma = self._simplify_mesh_open3d(soma_mesh, soma_target)
            simplified_skeleton = self._simplify_mesh_open3d(skeleton_mesh, skeleton_target)
            
            # Combine the simplified meshes
            combined = trimesh.util.concatenate([simplified_soma, simplified_skeleton])
            
            return combined
            
        except Exception as e:
            # Fall back to uniform simplification on any error
            self._vprint(f'      ⚠️ Soma-aware simplification failed ({e}), using uniform', level='full')
            n_faces = len(faces)
            target_faces = max(100, int(n_faces * (1 - skeleton_simp)))
            return self._simplify_mesh_open3d(trimesh_obj, target_faces)

    def _export_png_with_timeout(self, fig, output_path, width=1200, height=900, scale=3, timeout=120,
                                   auto_crop=False, crop_margin=30):
        """
        Export PNG using kaleido with timeout protection and optional auto-crop.
        
        Parameters
        ----------
        fig : plotly.graph_objects.Figure
            The figure to export.
        output_path : str
            Path to save the PNG file.
        width : int
            Image width in pixels.
        height : int
            Image height in pixels.
        scale : int
            Scale factor for resolution.
        timeout : int
            Timeout in seconds (default 120).
        auto_crop : bool, default False
            If True, automatically crop whitespace/background from the exported image.
        crop_margin : int, default 30
            Margin (in pixels) to preserve around content when auto_crop=True.
            
        Returns
        -------
        tuple
            (success: bool, message: str, final_scale: int)
            If success, message contains file size info.
            If failed, message contains error/recommendation.
        """
        import signal
        
        # Get export_timeout from attribute if set, otherwise use parameter
        actual_timeout = getattr(self, 'export_timeout', timeout)
        
        # Set up timeout handler (Unix only)
        old_handler = None
        has_alarm_support = False
        try:
            old_handler = signal.signal(signal.SIGALRM, _timeout_handler)
            has_alarm_support = True
        except (AttributeError, ValueError):
            has_alarm_support = False
        
        try:
            if has_alarm_support:
                signal.alarm(actual_timeout)

            fig.write_image(output_path, width=width, height=height, scale=scale, validate=False)

            if has_alarm_support:
                signal.alarm(0)

            if not os.path.exists(output_path):
                return (False, "Export failed - no file created", scale)

            size_kb = os.path.getsize(output_path) / 1024
            if size_kb < 10:
                return (False, f"Export produced blank image ({size_kb:.1f}KB)", scale)

            if auto_crop:
                try:
                    from PIL import Image

                    img = Image.open(output_path)
                    bounds = self._detect_content_bounds(img, (255, 255, 255))
                    if bounds:
                        row_min, row_max, col_min, col_max = bounds
                        row_min = max(0, row_min - crop_margin)
                        row_max = min(img.height - 1, row_max + crop_margin)
                        col_min = max(0, col_min - crop_margin)
                        col_max = min(img.width - 1, col_max + crop_margin)
                        cropped = img.crop((col_min, row_min, col_max + 1, row_max + 1))
                        cropped.save(output_path, 'PNG')
                        size_kb = os.path.getsize(output_path) / 1024
                except Exception:
                    pass

            return (True, f"{size_kb:.1f}KB", scale)

        except PNGExportTimeout:
            return (False, f"Export timed out after {actual_timeout}s", scale)
        except Exception as e:
            return (False, str(e), scale)
        finally:
            if has_alarm_support:
                try:
                    signal.alarm(0)
                except Exception:
                    pass
            try:
                signal.signal(signal.SIGALRM, old_handler)
            except Exception:
                pass

    def _export_views_with_webdriver_session(self, export_fig, views_to_export, view_cameras, output_folder):
        """
        Export multiple views efficiently using a single WebDriver session.
        
        Opens the browser once, loads the HTML, then rotates camera via JavaScript
        to export all views without reopening the browser.
        
        Parameters
        ----------
        export_fig : plotly.graph_objects.Figure
            The figure to export
        views_to_export : list
            List of view names to export (e.g., ['front', 'back', 'top'])
        view_cameras : dict
            Dict mapping view names to camera dicts with 'eye', 'up', 'center' keys
        output_folder : str
            Folder to save the PNG files
            
        Returns
        -------
        list
            List of successfully exported view names
        """
        exported_views = []
        
        # Save figure to temporary HTML
        temp_html = os.path.join(output_folder, "_temp_export.html")
        export_fig.write_html(
            temp_html, 
            auto_open=False, 
            include_plotlyjs='cdn',
            config={'displayModeBar': False}
        )
        
        try:
            self._vprint(f'     Using WebDriver session for {len(views_to_export)} views...')
            
            # Get render_wait from attribute (None = auto-calibrate)
            render_wait = getattr(self, 'webdriver_render_wait', None)
            
            # Retry logic for Chrome crashes
            max_retries = 3
            last_error = None
            
            for retry_attempt in range(max_retries):
                if retry_attempt > 0:
                    self._vprint(f'      🔄 Retry attempt {retry_attempt + 1}/{max_retries}...')
                    import time
                    time.sleep(2)  # Brief pause before retry
                
                try:
                    with WebDriverExportSession(
                        width=1200, height=900, 
                        scale=self.export_scale, 
                        timeout=300,
                        render_wait=render_wait
                    ) as session:
                        # Load HTML once
                        session.load_html(temp_html, wait_for_render=True, render_wait=3, background_color=self.background_color)
                        self._vprint(f'      ✓ HTML loaded in browser (render_wait={session._render_wait:.2f}s)')
                        
                        for view_name in views_to_export:
                            if view_name not in view_cameras:
                                self._vprint(f'      ⚠️  Skipping invalid view: {view_name}')
                                continue
                            
                            camera = view_cameras[view_name]
                            view_path = os.path.join(output_folder, f"{self.saveas}_{view_name}.png")
                            
                            try:
                                # Rotate camera via JavaScript
                                session.set_camera(
                                    eye=camera.get('eye'),
                                    up=camera.get('up'),
                                    center=camera.get('center')
                                )
                                
                                # Take screenshot
                                session.screenshot(view_path)
                                
                                # Verify file
                                if os.path.exists(view_path):
                                    size_kb = os.path.getsize(view_path) / 1024
                                    if size_kb > 10:
                                        exported_views.append(view_name)
                                        self._vprint(f'      {view_name}: {size_kb:.1f}KB', level='full')
                                    else:
                                        self._vprint(f'      ⚠️  {view_name}: blank image ({size_kb:.1f}KB)')
                                else:
                                    self._vprint(f'      ⚠️  {view_name}: no file created')
                                    
                            except Exception as e:
                                self._vprint(f'      ⚠️  {view_name} export failed: {e}')
                        
                        # If we got any views, consider it a success
                        if exported_views:
                            break  # Exit retry loop
                            
                except Exception as e:
                    last_error = e
                    error_msg = str(e)
                    # Check if this is a Chrome crash
                    is_chrome_crash = (
                        'Message: \n' in str(e) or 
                        error_msg == '' or
                        'chrome not reachable' in error_msg.lower() or
                        'session deleted' in error_msg.lower() or
                        'no such window' in error_msg.lower()
                    )
                    
                    if is_chrome_crash and retry_attempt < max_retries - 1:
                        self._vprint(f'      ⚠️  Chrome crashed unexpectedly. Will retry...')
                        continue
                    else:
                        raise  # Re-raise for outer exception handler
                
        except Exception as e:
            self._vprint(f'      ⚠️  WebDriver session failed: {e}')
            
            # Check HTML size and try with simplified figure
            html_size_mb = os.path.getsize(temp_html) / 1024 / 1024 if os.path.exists(temp_html) else 0
            html_size_cap = self._get_html_size_cap()
            
            if html_size_mb > html_size_cap:
                # First check if we already have a simplified figure
                existing_simplified = getattr(self, '_simplified_export_fig', None)
                existing_html = getattr(self, '_simplified_html_path', None)
                
                if existing_simplified is not None and existing_html and os.path.exists(existing_html):
                    self._vprint(f'      💡 Trying with existing simplified figure...')
                    try:
                        with WebDriverExportSession(
                            width=1200, height=900, 
                            scale=self.export_scale, 
                            timeout=300,
                            render_wait=render_wait
                        ) as session:
                            session.load_html(existing_html, wait_for_render=True, render_wait=3, background_color=self.background_color)
                            self._vprint(f'      ✓ Existing simplified HTML loaded in browser')
                            
                            for view_name in views_to_export:
                                if view_name not in view_cameras:
                                    continue
                                
                                camera = view_cameras[view_name]
                                view_path = os.path.join(output_folder, f"{self.saveas}_{view_name}.png")
                                
                                try:
                                    session.set_camera(
                                        eye=camera.get('eye'),
                                        up=camera.get('up'),
                                        center=camera.get('center')
                                    )
                                    session.screenshot(view_path)
                                    
                                    if os.path.exists(view_path):
                                        size_kb = os.path.getsize(view_path) / 1024
                                        if size_kb > 10:
                                            exported_views.append(view_name)
                                            self._vprint(f'      {view_name}: {size_kb:.1f}KB', level='full')
                                except Exception as inner_e:
                                    self._vprint(f'      ⚠️  {view_name} failed: {inner_e}')
                    except Exception as retry_e:
                        self._vprint(f'      ⚠️  Existing simplified failed: {retry_e}')
                
                # If still no success, try creating new simplified figure with 50% target, then 25%
                if not exported_views:
                    html_overhead_mb = 25  # Estimated fixed overhead
                    for reduction_target in [0.5, 0.25]:  # 50% then 25% of original
                        if exported_views:
                            break
                        
                        target_size_mb = html_size_mb * reduction_target
                        data_size = html_size_mb - html_overhead_mb
                        target_data_size = target_size_mb - html_overhead_mb
                        simplification_factor = max(0.1, target_data_size / data_size) if data_size > 0 else reduction_target
                        
                        self._vprint(f'      💡 Retrying with simplified figure (target ~{target_size_mb:.0f}MB, factor={simplification_factor:.2f})...')
                        
                        try:
                            simplified_fig = self._simplify_figure_for_kaleido(export_fig, simplification_factor)
                            
                            # Save simplified HTML to permanent file for reuse
                            simplified_html_path = os.path.join(output_folder, f"{self.saveas}_simplified.html")
                            simplified_fig.write_html(simplified_html_path, auto_open=False, 
                                                     include_plotlyjs='cdn',
                                                     config={'displayModeBar': False})
                            
                            new_size_mb = os.path.getsize(simplified_html_path) / 1024 / 1024
                            self._vprint(f'      ✓ Simplified HTML saved: {os.path.basename(simplified_html_path)} ({new_size_mb:.0f}MB)')
                            
                            # Store simplified figure and path for reuse in subsequent exports
                            self._simplified_export_fig = simplified_fig
                            self._simplified_html_path = simplified_html_path
                            
                            with WebDriverExportSession(
                                width=1200, height=900, 
                                scale=self.export_scale, 
                                timeout=300,
                                render_wait=render_wait
                            ) as session:
                                session.load_html(simplified_html_path, wait_for_render=True, render_wait=3, background_color=self.background_color)
                                self._vprint(f'      ✓ Simplified HTML loaded in browser')
                                
                                for view_name in views_to_export:
                                    if view_name not in view_cameras:
                                        continue
                                    
                                    camera = view_cameras[view_name]
                                    view_path = os.path.join(output_folder, f"{self.saveas}_{view_name}.png")
                                    
                                    try:
                                        session.set_camera(
                                            eye=camera.get('eye'),
                                            up=camera.get('up'),
                                            center=camera.get('center')
                                        )
                                        session.screenshot(view_path)
                                        
                                        if os.path.exists(view_path):
                                            size_kb = os.path.getsize(view_path) / 1024
                                            if size_kb > 10:
                                                exported_views.append(view_name)
                                                self._vprint(f'      {view_name}: {size_kb:.1f}KB', level='full')
                                    except Exception as inner_e:
                                        self._vprint(f'      ⚠️  {view_name} failed: {inner_e}')
                            
                            if exported_views:
                                break  # Success with this simplification level
                                
                        except Exception as retry_e:
                            self._vprint(f'      ⚠️  Retry failed: {retry_e}')
                            continue
            
            if not exported_views:
                self._vprint(f'      ⚠️  All WebDriver export attempts failed')
                self._vprint(f'      💡 Try: export_method="kaleido" or increase skeleton_mesh_simplification')
        
        finally:
            # Clean up temp HTML (but not the permanent simplified HTML)
            if temp_html and not temp_html.endswith('_simplified.html'):
                try:
                    os.remove(temp_html)
                except:
                    pass
        
        return exported_views

    def _create_simplified_export_figure(self):
        """
        Create a simplified copy of the figure for PNG export.
        
        For large/complex figures, this dramatically reduces export time by:
        1. Reducing mesh complexity (decimating vertices/faces)
        2. Simplifying line traces (reducing point count)
        3. Keeping visual appearance similar but with less data
        
        Returns
        -------
        plotly.graph_objects.Figure
            Simplified figure optimized for static image export.
        """
        import trimesh
        
        self._vprint('      Creating simplified figure for PNG export...')
        
        # Create a new figure with same layout
        simplified_fig = go.Figure()
        simplified_fig.update_layout(self.fig_3d.layout)
        
        # Statistics
        total_original_vertices = 0
        total_simplified_vertices = 0
        mesh_count = 0
        
        for trace in self.fig_3d.data:
            trace_dict = trace.to_plotly_json()
            trace_type = trace_dict.get('type', '')
            
            if trace_type == 'mesh3d':
                # Simplify mesh traces - this is the main bottleneck
                x = trace_dict.get('x', [])
                y = trace_dict.get('y', [])
                z = trace_dict.get('z', [])
                i = trace_dict.get('i', [])
                j = trace_dict.get('j', [])
                k = trace_dict.get('k', [])
                
                if x is not None and len(x) > 0 and i is not None and len(i) > 0:
                    original_verts = len(x)
                    original_faces = len(i)
                    total_original_vertices += original_verts
                    mesh_count += 1
                    
                    # Only simplify large meshes
                    if original_faces > 5000:
                        try:
                            import numpy as np
                            vertices = np.array([x, y, z]).T
                            faces = np.array([i, j, k]).T
                            
                            mesh = trimesh.Trimesh(vertices=vertices, faces=faces, process=False)
                            
                            # Target: reduce to ~10% of faces (aggressive for export)
                            target_faces = max(500, int(original_faces * 0.1))
                            simplified_mesh = self._simplify_mesh_open3d(mesh, target_faces)
                            
                            # Create new trace with simplified data
                            new_trace_dict = {k: v for k, v in trace_dict.items() 
                                             if k not in ['x', 'y', 'z', 'i', 'j', 'k', 'type'] and v is not None}
                            new_trace_dict['x'] = simplified_mesh.vertices[:, 0].tolist()
                            new_trace_dict['y'] = simplified_mesh.vertices[:, 1].tolist()
                            new_trace_dict['z'] = simplified_mesh.vertices[:, 2].tolist()
                            new_trace_dict['i'] = simplified_mesh.faces[:, 0].tolist()
                            new_trace_dict['j'] = simplified_mesh.faces[:, 1].tolist()
                            new_trace_dict['k'] = simplified_mesh.faces[:, 2].tolist()
                            
                            simplified_fig.add_trace(go.Mesh3d(**new_trace_dict))
                            total_simplified_vertices += len(simplified_mesh.vertices)
                            continue
                            
                        except Exception as e:
                            # Fall back to original if simplification fails
                            self._vprint(f'        ⚠️ Mesh simplification failed: {e}', level='full')
                    
                    total_simplified_vertices += original_verts
                
                # Keep original if not simplified
                simplified_fig.add_trace(trace)
                
            elif trace_type == 'scatter3d':
                # For scatter/line traces, reduce point count
                x = trace_dict.get('x', [])
                y = trace_dict.get('y', [])
                z = trace_dict.get('z', [])
                
                if x is not None and len(x) > 1000:
                    # Subsample large traces
                    step = max(1, len(x) // 500)
                    trace_dict['x'] = list(x[::step])
                    trace_dict['y'] = list(y[::step])
                    trace_dict['z'] = list(z[::step])
                    
                    new_trace = go.Scatter3d(**{k: v for k, v in trace_dict.items() 
                                               if k != 'type' and v is not None})
                    simplified_fig.add_trace(new_trace)
                else:
                    simplified_fig.add_trace(trace)
                
            else:
                # Keep other trace types as-is
                simplified_fig.add_trace(trace)
        
        if total_original_vertices > 0:
            reduction = (1 - total_simplified_vertices / total_original_vertices) * 100
            self._vprint(f'      Simplified {mesh_count} meshes: {total_original_vertices:,} → {total_simplified_vertices:,} vertices ({reduction:.0f}% reduction)')
        
        return simplified_fig

    def _simplify_figure_for_kaleido(self, fig, simplification_factor):
        """
        Simplify a figure for kaleido export based on HTML size ratio.
        
        For figures where HTML size > 100MB, this applies proportional simplification
        to reduce memory/time requirements for kaleido rendering.
        
        Parameters
        ----------
        fig : plotly.graph_objects.Figure
            The figure to simplify
        simplification_factor : float
            Target ratio for data reduction (e.g., 0.25 = reduce to 25% of original)
            
        Returns
        -------
        plotly.graph_objects.Figure
            Simplified figure
        """
        import trimesh
        import numpy as np
        
        # Create new figure with same layout
        simplified_fig = go.Figure()
        simplified_fig.update_layout(fig.layout)
        
        mesh_count = 0
        total_original_faces = 0
        total_simplified_faces = 0
        
        for trace in fig.data:
            trace_dict = trace.to_plotly_json()
            trace_type = trace_dict.get('type', '')
            
            if trace_type == 'mesh3d':
                x = trace_dict.get('x', [])
                y = trace_dict.get('y', [])
                z = trace_dict.get('z', [])
                i = trace_dict.get('i', [])
                j = trace_dict.get('j', [])
                k = trace_dict.get('k', [])
                
                if x is not None and len(x) > 0 and i is not None and len(i) > 0:
                    original_faces = len(i)
                    total_original_faces += original_faces
                    mesh_count += 1
                    
                    # Target faces based on simplification factor
                    target_faces = max(100, int(original_faces * simplification_factor))
                    
                    # Only simplify if it's worth it
                    if target_faces < original_faces * 0.8:
                        try:
                            vertices = np.array([x, y, z]).T
                            faces = np.array([i, j, k]).T
                            
                            mesh = trimesh.Trimesh(vertices=vertices, faces=faces, process=False)
                            simplified_mesh = self._simplify_mesh_open3d(mesh, target_faces)
                            
                            # Create new trace with simplified data
                            new_trace_dict = {k: v for k, v in trace_dict.items() 
                                             if k not in ['x', 'y', 'z', 'i', 'j', 'k', 'type'] and v is not None}
                            new_trace_dict['x'] = simplified_mesh.vertices[:, 0].tolist()
                            new_trace_dict['y'] = simplified_mesh.vertices[:, 1].tolist()
                            new_trace_dict['z'] = simplified_mesh.vertices[:, 2].tolist()
                            new_trace_dict['i'] = simplified_mesh.faces[:, 0].tolist()
                            new_trace_dict['j'] = simplified_mesh.faces[:, 1].tolist()
                            new_trace_dict['k'] = simplified_mesh.faces[:, 2].tolist()
                            
                            simplified_fig.add_trace(go.Mesh3d(**new_trace_dict))
                            total_simplified_faces += len(simplified_mesh.faces)
                            continue
                            
                        except Exception as e:
                            # Fall back to original if simplification fails
                            pass
                    
                    total_simplified_faces += original_faces
                
                simplified_fig.add_trace(trace)
                
            elif trace_type == 'scatter3d':
                # Reduce point count for line traces
                x = trace_dict.get('x', [])
                
                if x is not None and len(x) > 500:
                    # Subsample based on factor
                    target_points = max(100, int(len(x) * simplification_factor))
                    step = max(1, len(x) // target_points)
                    
                    for key in ['x', 'y', 'z']:
                        if key in trace_dict and trace_dict[key] is not None:
                            trace_dict[key] = list(trace_dict[key][::step])
                    
                    # Also handle marker colors if present
                    if 'marker' in trace_dict and isinstance(trace_dict['marker'], dict):
                        if 'color' in trace_dict['marker'] and isinstance(trace_dict['marker']['color'], (list, tuple)):
                            trace_dict['marker']['color'] = list(trace_dict['marker']['color'][::step])
                    
                    new_trace = go.Scatter3d(**{k: v for k, v in trace_dict.items() 
                                               if k != 'type' and v is not None})
                    simplified_fig.add_trace(new_trace)
                else:
                    simplified_fig.add_trace(trace)
            else:
                simplified_fig.add_trace(trace)
        
        if total_original_faces > 0:
            reduction = (1 - total_simplified_faces / total_original_faces) * 100
            self._vprint(f'   Simplified {mesh_count} meshes: {total_original_faces:,} → {total_simplified_faces:,} faces ({reduction:.0f}% reduction)')
        
        return simplified_fig

    def _detect_content_bounds(self, img, background_color=(255, 255, 255)):
        return _detect_content_bounds(img, background_color)
    def _compute_unified_crop_bounds(self, image_paths, background_color=(255, 255, 255), 
                                      sample_count=None):
        return _compute_unified_crop_bounds(image_paths, background_color, sample_count)
    def _apply_consistent_crop(self, pic_folder, margin=20, background_color=(255, 255, 255)):
        return _apply_consistent_crop(pic_folder, margin, background_color)
    def _filter_neuron_df_by_hemisphere(self, ndf, rdf=None):
        """Keep only neurons of the selected hemisphere (self.hemisphere).

        Hemisphere assignment mirrors the pathfinding backend: the
        'Soma side' / 'hemisphere' column wins, then the instance suffix
        (_L/_R).  Neurons WITHOUT an explicit hemisphere are treated as
        unclassified ('U') and are ALWAYS kept - in 'left', 'right' AND
        'both' - so data that lacks hemisphere notation is never silently
        dropped.  Returns (filtered_ndf, filtered_rdf).
        """
        df = ndf.copy()
        code = pd.Series('U', index=df.index, dtype=object)

        # Locate the side column regardless of naming variant ('Soma side',
        # 'somaSide', 'rootSide', 'hemisphere') - male-cns CSVs use camelCase.
        side_col = None
        lowered = {str(c).strip().lower(): c for c in df.columns}
        for candidate in ('hemisphere', 'soma side', 'somaside', 'rootside'):
            if candidate in lowered:
                side_col = lowered[candidate]
                break
        if side_col is not None:
            side = df[side_col].fillna('').astype(str).str.strip().str.lower()
            code[side.isin(['l', 'left', 'lhs', 'left hemisphere'])] = 'L'
            code[side.isin(['r', 'right', 'rhs', 'right hemisphere'])] = 'R'

        if 'instance' in df.columns:
            inst = df['instance'].fillna('').astype(str)
            unassigned = code == 'U'
            code[unassigned & inst.str.endswith('_R')] = 'R'
            code[unassigned & inst.str.endswith('_L')] = 'L'

        if self.hemisphere == 'left':
            df = df[code.isin(['L', 'U'])].copy()
        elif self.hemisphere == 'right':
            df = df[code.isin(['R', 'U'])].copy()

        if rdf is not None and not rdf.empty and 'bodyId' in df.columns:
            rdf = rdf[rdf.index.isin(df['bodyId'].astype(str))].copy() if rdf.index.name == 'bodyId' else rdf
        return df, rdf

    def _parse_layer_map_csv(self):
        """
        Parse layer_map_csv file to construct neuron_layers and custom_layer_names.
        
        The CSV must have columns 'layer' and 'id_type_instance'.
        Optionally can have a 'color' column for per-neuron color overrides.
        Rows with the same 'layer' value are grouped together into a single layer.
        
        This method overrides self.neuron_layers and self.custom_layer_names.
        If 'color' column exists, populates self._neuron_color_overrides.
        """
        import pandas as pd
        # standardize_color already imported at module level from utils.color_utils
        
        csv_path = self.layer_map_csv
        if not os.path.exists(csv_path):
            raise FileNotFoundError(f"layer_map_csv not found: {csv_path}")
        
        self._vprint(f"Loading layer map from: {csv_path}", level='full')
        
        df = pd.read_csv(csv_path)
        
        # Validate columns
        required_cols = ['layer', 'id_type_instance']
        for col in required_cols:
            if col not in df.columns:
                raise ValueError(f"layer_map_csv must have column '{col}'. Found: {list(df.columns)}")
        
        # Check for optional color column
        has_color_col = 'color' in df.columns
        if has_color_col:
            self._vprint(f"  Found 'color' column - will apply per-neuron colors", level='simple')
        
        # Build per-neuron color overrides if color column exists
        self._neuron_color_overrides = {}  # bodyId/type -> rgba_string
        if has_color_col:
            for _, row in df.iterrows():
                id_val = row['id_type_instance']
                color_val = row['color']
                if pd.notna(color_val) and str(color_val).strip():
                    # Normalize id
                    id_str = str(id_val).strip()
                    id_key = int(id_str) if id_str.isdigit() else id_str
                    # Standardize color
                    try:
                        rgba_str = standardize_color(str(color_val).strip(), default_alpha=self.neuron_alpha)
                        self._neuron_color_overrides[id_key] = rgba_str
                        # Also store string version for lookup flexibility
                        self._neuron_color_overrides[str(id_key)] = rgba_str
                    except Exception as e:
                        self._vprint(f"  ⚠️ Failed to parse color '{color_val}' for {id_val}: {e}", level='full')
            
            if self._neuron_color_overrides:
                self._vprint(f"  Loaded {len(self._neuron_color_overrides) // 2} per-neuron color overrides", level='simple')
        
        # Group by layer name to create neuron_layers
        layer_groups = df.groupby('layer', sort=False)['id_type_instance'].apply(list).to_dict()
        
        # Construct neuron_layers and custom_layer_names
        self.neuron_layers = []
        self.custom_layer_names = []
        
        for layer_name, identifiers in layer_groups.items():
            # Convert identifiers: if it looks like a bodyId (all digits), convert to int
            processed_ids = []
            for id_val in identifiers:
                id_str = str(id_val).strip()
                if id_str.isdigit():
                    processed_ids.append(int(id_str))
                else:
                    processed_ids.append(id_str)
            
            # If single item, use it directly; if multiple, keep as list
            if len(processed_ids) == 1:
                self.neuron_layers.append(processed_ids[0])
            else:
                self.neuron_layers.append(processed_ids)
            
            self.custom_layer_names.append(str(layer_name))
        
        self._vprint(f"  Loaded {len(self.neuron_layers)} layers from CSV:")
        for i, (name, neurons) in enumerate(zip(self.custom_layer_names, self.neuron_layers)):
            n_count = len(neurons) if isinstance(neurons, list) else 1
            self._vprint(f"    Layer {i}: {name} ({n_count} neurons)")

    def _apply_soma_radius_cap(self, neuron_vols):
        """
        Apply radius capping and optional smoothing to skeleton radii.
        
        When smooth_skeleton=False (default): Only applies hard cap to radii.
        When smooth_skeleton=True: Also applies iterative smoothing to prevent
        extrusion artifacts from chains of large-radius nodes.
        
        Parameters
        ----------
        neuron_vols : navis.NeuronList
            List of neurons to process (modified in place)
        """
        cap = self.soma_radius_cap
        total_capped = 0
        total_smoothed = 0
        
        for n in neuron_vols:
            if not hasattr(n, 'nodes') or not isinstance(n.nodes, pd.DataFrame):
                continue
            if 'radius' not in n.nodes.columns:
                continue
            
            nodes = n.nodes
            radii = nodes['radius'].values.copy().astype(float)
            original_radii = radii.copy()
            
            # Step 1: Hard cap all radii above threshold
            over_cap = radii > cap
            if over_cap.any():
                radii[over_cap] = cap
                total_capped += over_cap.sum()
            
            # Step 2: Optional iterative smoothing (only if smooth_skeleton=True)
            if self.smooth_skeleton and 'parent_id' in nodes.columns and 'node_id' in nodes.columns:
                node_ids = nodes['node_id'].values
                parent_ids = nodes['parent_id'].values
                id_to_idx = {nid: idx for idx, nid in enumerate(node_ids)}
                
                # Build child map
                children = {idx: [] for idx in range(len(radii))}
                for idx, pid in enumerate(parent_ids):
                    if pid in id_to_idx:
                        children[id_to_idx[pid]].append(idx)
                
                # Aggressive smoothing: 20 passes with very strong neighbor influence
                for pass_num in range(20):
                    new_radii = radii.copy()
                    for idx in range(len(radii)):
                        if radii[idx] <= 0:
                            continue
                        
                        # Collect neighbor radii (parent + children)
                        neighbors = []
                        pid = parent_ids[idx]
                        if pid in id_to_idx:
                            neighbors.append(radii[id_to_idx[pid]])
                        for child_idx in children[idx]:
                            if radii[child_idx] > 0:
                                neighbors.append(radii[child_idx])
                        
                        if neighbors:
                            # Very strong neighbor influence: 10% self, 90% neighbors
                            neighbor_avg = np.mean(neighbors)
                            new_radii[idx] = 0.1 * radii[idx] + 0.9 * neighbor_avg
                    
                    radii = new_radii
                
                # Count how many were significantly changed
                total_smoothed += np.sum(np.abs(radii - original_radii) > 1)
                
                # Final cap check after smoothing
                radii = np.minimum(radii, cap)
            
            n.nodes['radius'] = radii
        
        if total_capped > 0:
            if self.smooth_skeleton:
                self._vprint(f"  ✓ Radius capping: capped {total_capped}, smoothed {total_smoothed} nodes (cap={cap:.0f}nm)", level='full')
            else:
                self._vprint(f"  ✓ Radius capping: capped {total_capped} nodes (cap={cap:.0f}nm)", level='full')

    @contextmanager
    def _suppress_output(self):
        """Suppress stdout and stderr if verbose is not full."""
        if self.verbose == 'full':
            yield
        else:
            with open(os.devnull, "w") as devnull:
                old_stdout = sys.stdout
                old_stderr = sys.stderr
                sys.stdout = devnull
                sys.stderr = devnull
                try:  
                    yield
                finally:
                    sys.stdout = old_stdout
                    sys.stderr = old_stderr

    def _xform_neurons_safe(self, neuron_vols, source, target, layer_label='layer'):
        """
        Transform neurons one at a time with explicit progress and per-neuron fallback.

        ``navis.xform_brain(neuron_vols, ...)`` is a single black-box call: when
        one pathological neuron (or a stall inside the transform backend) freezes
        the run, the only clue is a tqdm bar stuck at some percentage. This helper
        instead:

          * resolves the bridging transform sequence once and prints it,
          * shows the neuron currently being transformed in the bar postfix, so a
            stall immediately identifies the culprit neuron,
          * prints an explicit heartbeat line every HEARTBEAT neurons (flushed via
            tqdm.write so it always reaches the UI log, even in non-TTY pipes),
          * keeps the untransformed original on a per-neuron failure instead of
            aborting the whole layer,
          * releases each original neuron as soon as its copy is transformed to
            keep peak memory flat on very large layers.
        """
        from navis.transforms import registry as _registry

        # flybrains registers the template brains + transforms with navis;
        # make sure it is imported before resolving the bridging sequence.
        try:
            import flybrains  # noqa: F401
        except ImportError:
            pass

        neurons = list(neuron_vols) if isinstance(neuron_vols, navis.NeuronList) else [neuron_vols]
        total = len(neurons)

        n_nodes_total = 0
        for n in neurons:
            if hasattr(n, 'n_nodes'):
                n_nodes_total += int(n.n_nodes)
            elif hasattr(n, 'vertices') and getattr(n, 'vertices', None) is not None:
                n_nodes_total += len(n.vertices)
        tqdm.write(f'  🔀 {layer_label}: transforming {total} neurons, '
                   f'{n_nodes_total:,} points  [{source} -> {target}]')

        # Resolve the transform sequence once (raises if no path exists, so the
        # caller can trigger the transform download + retry flow).
        path, seq = _registry.shortest_bridging_seq(source, target)
        tqdm.write(f'     Resolved transform path: {" -> ".join(str(p) for p in path)} '
                   f'({len(seq)} step(s))')

        HEARTBEAT = 250
        start = time.time()
        xf_neurons = []
        failed = []
        pbar = tqdm(total=total, desc='  Transforming', unit='neuron', mininterval=0.5)
        try:
            for i, n in enumerate(neurons):
                label = getattr(n, 'name', None) or str(getattr(n, 'id', i))
                pbar.set_postfix_str(f'last={label}')
                if not isinstance(n, navis.BaseNeuron):
                    # Volumes or other non-neuron objects pass through untouched
                    xf_neurons.append(n)
                    neurons[i] = None
                    pbar.update(1)
                    continue
                try:
                    xf_neurons.append(navis.xform(n, transform=seq))
                except Exception as e:
                    failed.append((label, str(e)))
                    tqdm.write(f'     ⚠️  Transform failed for {label}: {e} '
                               f'(keeping original coordinates)')
                    xf_neurons.append(n)
                # Free the original as soon as its copy is done to keep peak
                # memory flat (large layers double memory inside navis.xform).
                neurons[i] = None
                pbar.update(1)
                if (i + 1) % HEARTBEAT == 0 or (i + 1) == total:
                    elapsed = time.time() - start
                    rate = (i + 1) / max(elapsed, 1e-6)
                    tqdm.write(f'     ... {i + 1}/{total} neurons transformed '
                               f'({elapsed:.1f}s elapsed, {rate:.0f} neurons/s)')
        finally:
            pbar.close()

        if failed:
            tqdm.write(f'  ⚠️  {layer_label}: {len(failed)}/{total} neurons could not '
                       f'be transformed and were kept in source coordinates')
        tqdm.write(f'  ✓ {layer_label}: transform finished in {time.time() - start:.1f}s')
        return navis.NeuronList(xf_neurons)

    def _validate_inputs(self):
        """
        Validate all input parameters before processing.
        Raises ValueError with descriptive messages for invalid inputs.
        """
        errors = []
        
        # === String validations ===
        if not isinstance(self.backend, str):
            errors.append(f"backend must be a string, got {type(self.backend).__name__}")
        elif self.backend not in ('plotly', 'k3d'):
            errors.append(f"backend must be 'plotly' or 'k3d', got '{self.backend}'")
            
        if not isinstance(self.dataset, str):
            errors.append(f"dataset must be a string, got {type(self.dataset).__name__}")
        elif not self.dataset:
            errors.append("dataset cannot be empty")
            
        if not isinstance(self.client_type, str):
            errors.append(f"client_type must be a string, got {type(self.client_type).__name__}")
        elif self.client_type not in ('neuprint', 'flywire'):
            errors.append(f"client_type must be 'neuprint' or 'flywire', got '{self.client_type}'")

        if not isinstance(self.search_columns, str):
            errors.append(f"search_columns must be a string, got {type(self.search_columns).__name__}")
        elif self.search_columns not in ('auto', 'type', 'instance', 'bodyId'):
            errors.append(
                f"search_columns must be 'auto', 'type', 'instance' or 'bodyId', "
                f"got '{self.search_columns}'"
            )

        if not isinstance(self.hemisphere, str):
            errors.append(f"hemisphere must be a string, got {type(self.hemisphere).__name__}")
        elif self.hemisphere not in ('both', 'left', 'right'):
            errors.append(
                f"hemisphere must be 'both', 'left' or 'right', got '{self.hemisphere}'"
            )
            
        if not isinstance(self.server, str):
            errors.append(f"server must be a string, got {type(self.server).__name__}")
            
        if self.token is not None and not isinstance(self.token, str):
            errors.append(f"token must be a string or None, got {type(self.token).__name__}")
            
        if self.version is not None and not isinstance(self.version, int):
            errors.append(f"version must be an integer or None, got {type(self.version).__name__}")
        elif self.version is not None and self.version < 0:
            errors.append(f"version must be positive, got {self.version}")
            
        # === Path validations ===
        if not isinstance(self.data_folder, str):
            errors.append(f"data_folder must be a string, got {type(self.data_folder).__name__}")
            
        if self.output_dir is not None and not isinstance(self.output_dir, str):
            errors.append(f"output_dir must be a string or None, got {type(self.output_dir).__name__}")
            
        if self.saveas is not None and not isinstance(self.saveas, str):
            errors.append(f"saveas must be a string or None, got {type(self.saveas).__name__}")
            
        if self.layer_map_csv is not None:
            if not isinstance(self.layer_map_csv, str):
                errors.append(f"layer_map_csv must be a string or None, got {type(self.layer_map_csv).__name__}")
            elif not os.path.exists(self.layer_map_csv):
                errors.append(f"layer_map_csv file not found: {self.layer_map_csv}")
        
        # === Neuron layers validation ===
        if not isinstance(self.neuron_layers, (str, list)):
            errors.append(f"neuron_layers must be a string or list, got {type(self.neuron_layers).__name__}")
        else:
            # Check if neuron_layers is empty (empty string or empty list)
            is_empty = (isinstance(self.neuron_layers, str) and not self.neuron_layers) or \
                       (isinstance(self.neuron_layers, list) and len(self.neuron_layers) == 0)
            
            if is_empty and self.layer_map_csv is None:
                # Allow empty neuron_layers if mesh_roi or brain_mesh is specified (mesh-only mode)
                has_mesh = self.mesh_roi or (self.brain_mesh and self.brain_mesh.lower() != 'none') or self.vnc_mesh
                if not has_mesh:
                    errors.append("neuron_layers cannot be empty (or provide layer_map_csv, mesh_roi, or brain_mesh)")
            
        if not isinstance(self.custom_layer_names, list):
            errors.append(f"custom_layer_names must be a list, got {type(self.custom_layer_names).__name__}")
            
        # === Boolean validations ===
        bool_params = [
            ('cache_neurons', self.cache_neurons),
            ('cache_synapses', self.cache_synapses),
            ('skip_synapse', self.skip_synapse),
            ('include_timestamp', self.include_timestamp),
            ('auto_fix_extrusions', self.auto_fix_extrusions),
            ('smooth_skeleton', self.smooth_skeleton),
        ]
        for name, value in bool_params:
            if not isinstance(value, bool):
                errors.append(f"{name} must be a boolean, got {type(value).__name__}")
        
        # === Numeric validations ===
        if not isinstance(self.min_synapse_num, int):
            errors.append(f"min_synapse_num must be an integer, got {type(self.min_synapse_num).__name__}")
        elif self.min_synapse_num < 0:
            errors.append(f"min_synapse_num must be non-negative, got {self.min_synapse_num}")
            
        if not isinstance(self.export_scale, int):
            errors.append(f"export_scale must be an integer, got {type(self.export_scale).__name__}")
        elif self.export_scale < 1 or self.export_scale > 10:
            errors.append(f"export_scale must be between 1 and 10, got {self.export_scale}")
            
        if not isinstance(self.export_timeout, int):
            errors.append(f"export_timeout must be an integer, got {type(self.export_timeout).__name__}")
        elif self.export_timeout < 1:
            errors.append(f"export_timeout must be at least 1 second, got {self.export_timeout}")
            
        # === Float validations ===
        if not isinstance(self.neuron_alpha, (int, float)):
            errors.append(f"neuron_alpha must be a number, got {type(self.neuron_alpha).__name__}")
        elif not 0 <= self.neuron_alpha <= 1:
            errors.append(f"neuron_alpha must be between 0 and 1, got {self.neuron_alpha}")
            
        if not isinstance(self.synapse_alpha, (int, float)):
            errors.append(f"synapse_alpha must be a number, got {type(self.synapse_alpha).__name__}")
        elif not 0 <= self.synapse_alpha <= 1:
            errors.append(f"synapse_alpha must be between 0 and 1, got {self.synapse_alpha}")
            
        # Optional floats (can be None)
        if self.skeleton_mesh_simplification is not None:
            if not isinstance(self.skeleton_mesh_simplification, (int, float)):
                errors.append(f"skeleton_mesh_simplification must be a number or None, got {type(self.skeleton_mesh_simplification).__name__}")
            elif not 0 <= self.skeleton_mesh_simplification <= 1:
                errors.append(f"skeleton_mesh_simplification must be between 0 and 1, got {self.skeleton_mesh_simplification}")
                
        if self.soma_mesh_simplification is not None:
            if not isinstance(self.soma_mesh_simplification, (int, float)):
                errors.append(f"soma_mesh_simplification must be a number or None, got {type(self.soma_mesh_simplification).__name__}")
            elif not 0 <= self.soma_mesh_simplification <= 1:
                errors.append(f"soma_mesh_simplification must be between 0 and 1, got {self.soma_mesh_simplification}")
                
        if self.soma_radius_cap is not None:
            if not isinstance(self.soma_radius_cap, (int, float)):
                errors.append(f"soma_radius_cap must be a number or None, got {type(self.soma_radius_cap).__name__}")
            elif self.soma_radius_cap <= 0:
                errors.append(f"soma_radius_cap must be positive, got {self.soma_radius_cap}")
                
        if self.webdriver_render_wait is not None:
            if not isinstance(self.webdriver_render_wait, (int, float)):
                errors.append(f"webdriver_render_wait must be a number or None, got {type(self.webdriver_render_wait).__name__}")
            elif self.webdriver_render_wait < 0:
                errors.append(f"webdriver_render_wait must be non-negative, got {self.webdriver_render_wait}")
        
        # === Synapse size validation ===
        if not isinstance(self.synapse_size, (int, float, str)):
            errors.append(f"synapse_size must be a number or 'real', got {type(self.synapse_size).__name__}")
        elif isinstance(self.synapse_size, str) and self.synapse_size != 'real':
            errors.append(f"synapse_size string must be 'real', got '{self.synapse_size}'")
        elif isinstance(self.synapse_size, (int, float)) and self.synapse_size < 0:
            errors.append(f"synapse_size must be non-negative, got {self.synapse_size}")
            
        # === Mode validations (done in detail later, basic type check here) ===
        if not isinstance(self.skeleton_mode, str):
            errors.append(f"skeleton_mode must be a string, got {type(self.skeleton_mode).__name__}")
            
        if not isinstance(self.synapse_mode, str):
            errors.append(f"synapse_mode must be a string, got {type(self.synapse_mode).__name__}")
            
        if not isinstance(self.brain_mesh, str):
            errors.append(f"brain_mesh must be a string, got {type(self.brain_mesh).__name__}")
            
        if not isinstance(self.legend_mode, str):
            errors.append(f"legend_mode must be a string, got {type(self.legend_mode).__name__}")

        if not isinstance(self.color_mode, str):
            errors.append(f"color_mode must be a string, got {type(self.color_mode).__name__}")
            
        if not isinstance(self.expand_colors, str):
            errors.append(f"expand_colors must be a string, got {type(self.expand_colors).__name__}")
            
        if not isinstance(self.export_method, str):
            errors.append(f"export_method must be a string, got {type(self.export_method).__name__}")
        elif self.export_method not in ('webdriver', 'webdriver-fast', 'kaleido'):
            errors.append(f"export_method must be 'webdriver' or 'kaleido', got '{self.export_method}'")
            
        # === export_views validation ===
        valid_views = {'front', 'back', 'top', 'bottom', 'left', 'right', 'lateral', 'all'}
        if isinstance(self.export_views, bool):
            pass  # Valid
        elif isinstance(self.export_views, list):
            for view in self.export_views:
                if not isinstance(view, str):
                    errors.append(f"export_views list items must be strings, got {type(view).__name__}")
                elif view.lower() not in valid_views:
                    errors.append(f"Invalid view in export_views: '{view}'. Valid options: {sorted(valid_views)}")
        else:
            errors.append(f"export_views must be a boolean or list, got {type(self.export_views).__name__}")
            
        # === verbose validation ===
        if self.verbose not in (True, False, 'full', 'simple'):
            errors.append(f"verbose must be True, False, 'full', or 'simple', got {repr(self.verbose)}")
            
        # === Color validations (now more flexible - accepts multiple formats) ===
        # Empty/invalid values are coerced to defaults in _normalize_color_inputs();
        # anything still malformed here warns and falls back instead of failing.
        # neuron_colors: can be tuple, list, or single color string
        if isinstance(self.neuron_colors, str):
            # Single color string - will be wrapped in a list later
            pass  # Valid
        elif isinstance(self.neuron_colors, (tuple, list)):
            if len(self.neuron_colors) == 0:
                self._vprint("\033[33m⚠️  Warning: neuron_colors is empty; falling back to default palette.\033[0m", level='simple')
                self.neuron_colors = None
        else:
            self._vprint(f"\033[33m⚠️  Warning: neuron_colors must be a tuple, list, or color string, got {type(self.neuron_colors).__name__}; falling back to default palette.\033[0m", level='simple')
            self.neuron_colors = None
            
        # synapse_colors: same flexibility
        if isinstance(self.synapse_colors, str):
            pass  # Valid - single color string
        elif isinstance(self.synapse_colors, (tuple, list)):
            if len(self.synapse_colors) == 0:
                self._vprint("\033[33m⚠️  Warning: synapse_colors is empty; falling back to default palette.\033[0m", level='simple')
                self.synapse_colors = list(bokeh.palettes.Category10[10])
        else:
            self._vprint(f"\033[33m⚠️  Warning: synapse_colors must be a tuple, list, or color string, got {type(self.synapse_colors).__name__}; falling back to default palette.\033[0m", level='simple')
            self.synapse_colors = list(bokeh.palettes.Category10[10])
            
        # === mesh_roi validation ===
        if self.mesh_roi is not None and not isinstance(self.mesh_roi, (list, str)):
            errors.append(f"mesh_roi must be a list, string, or None, got {type(self.mesh_roi).__name__}")
            
        # === html_size_cap validation ===
        if self.html_size_cap is not None:
            if not isinstance(self.html_size_cap, int):
                errors.append(f"html_size_cap must be an integer or None, got {type(self.html_size_cap).__name__}")
            elif self.html_size_cap < 1:
                errors.append(f"html_size_cap must be at least 1 MB, got {self.html_size_cap}")
                
        # === export_simplified_png validation ===
        if not isinstance(self.export_simplified_png, (bool, int)):
            errors.append(f"export_simplified_png must be a boolean or integer, got {type(self.export_simplified_png).__name__}")
        elif isinstance(self.export_simplified_png, int) and self.export_simplified_png < 0:
            errors.append(f"export_simplified_png threshold must be non-negative, got {self.export_simplified_png}")
        
        # === Additional boolean validations ===
        additional_bools = [
            ('force_API_fetching', self.force_API_fetching),
            ('vnc_mesh', self.vnc_mesh),
            ('mirror_on_contralateral', self.mirror_on_contralateral),
            ('show_soma', self.show_soma),
            ('show_fig', self.show_fig),
            ('show_connectors', self.show_connectors),
            ('FAFB_template_correction', self.FAFB_template_correction),
        ]
        for name, value in additional_bools:
            if not isinstance(value, bool):
                errors.append(f"{name} must be a boolean, got {type(value).__name__}")
        
        # === Additional float validations ===
        if self.soma_region_radius is not None:
            if not isinstance(self.soma_region_radius, (int, float)):
                errors.append(f"soma_region_radius must be a number, got {type(self.soma_region_radius).__name__}")
            elif self.soma_region_radius <= 0:
                errors.append(f"soma_region_radius must be positive, got {self.soma_region_radius}")
                
        if self.roi_mesh_simplification is not None:
            if not isinstance(self.roi_mesh_simplification, (int, float)):
                errors.append(f"roi_mesh_simplification must be a number, got {type(self.roi_mesh_simplification).__name__}")
            elif not 0 <= self.roi_mesh_simplification <= 1:
                errors.append(f"roi_mesh_simplification must be between 0 and 1, got {self.roi_mesh_simplification}")
        
        # === String path validations ===
        if not isinstance(self.transforms_dir, str):
            errors.append(f"transforms_dir must be a string, got {type(self.transforms_dir).__name__}")
            
        # === Color validations for brain_mesh_color and vnc_mesh_color (flexible formats) ===
        # These now accept: 'auto', named colors, hex, rgba strings, tuples.
        # Malformed values warn and fall back to 'auto' instead of failing.
        def _validate_mesh_color(color, name):
            """Helper to validate flexible color format."""
            if isinstance(color, str):
                return  # Strings are validated later when standardizing
            elif isinstance(color, (tuple, list)):
                if len(color) >= 3:
                    return  # Valid tuple format
                self._vprint(f"\033[33m⚠️  Warning: {name} tuple must have at least 3 values (RGB), got {len(color)}; falling back to 'auto'.\033[0m", level='simple')
                setattr(self, name, 'auto')
            else:
                self._vprint(f"\033[33m⚠️  Warning: {name} must be a string, tuple, or list, got {type(color).__name__}; falling back to 'auto'.\033[0m", level='simple')
                setattr(self, name, 'auto')
        
        _validate_mesh_color(self.brain_mesh_color, 'brain_mesh_color')
        _validate_mesh_color(self.vnc_mesh_color, 'vnc_mesh_color')
            
        # === mesh_color validation (flexible formats now) ===
        if self.mesh_color is None:
            # None means "not specified": use the default gray tuple
            self.mesh_color = (100, 100, 100)
        if isinstance(self.mesh_color, str):
            pass  # Will be standardized later
        elif isinstance(self.mesh_color, (tuple, list)):
            # Check if it's a single color or list of colors
            if len(self.mesh_color) >= 3:
                # Could be (R, G, B, alpha) or [(color1), (color2), ...]
                if all(isinstance(x, (int, float)) for x in self.mesh_color):
                    # Single color tuple
                    if len(self.mesh_color) >= 4:
                        r, g, b, a = self.mesh_color[0], self.mesh_color[1], self.mesh_color[2], self.mesh_color[3]
                        if not 0 <= a <= 1:
                            self._vprint(f"\033[33m⚠️  Warning: mesh_color alpha must be 0-1, got {a}; using mesh_alpha instead.\033[0m", level='simple')
                            self.mesh_color = (r, g, b)
                # else: list of colors, validated later when standardizing
        else:
            self._vprint(f"\033[33m⚠️  Warning: mesh_color must be a tuple, list, or string, got {type(self.mesh_color).__name__}; falling back to default gray.\033[0m", level='simple')
            self.mesh_color = (100, 100, 100)
        
        # === Raise all errors together ===
        if errors:
            error_msg = "VisualizeSkeleton input validation failed:\n  - " + "\n  - ".join(errors)
            raise ValueError(error_msg)

    def _is_valid_color(self, color) -> bool:
        """Return True if `color` parses to a usable color (any supported format)."""
        try:
            standardize_color(color)
            return True
        except Exception:
            return False

    def _normalize_color_inputs(self):
        """
        Coerce empty or invalid color inputs to their defaults before validation.

        Empty lists (e.g., the UI sends synapse_colors=[] for a single layer)
        and unparseable values never hard-fail: each falls back to the
        parameter's default palette with a warning, regardless of skip_synapse.
        """
        def warn(name, detail, fallback_label):
            self._vprint(
                f"\033[33m⚠️  Warning: {name} {detail}; falling back to default ({fallback_label}).\033[0m",
                level='simple'
            )

        def normalize_sequence(value, name, fallback, fallback_label, allow_none=False, silent_empty=False):
            """Normalize a parameter that accepts a single color or a list of colors."""
            if value is None:
                return None if allow_none else fallback
            # Single color string
            if isinstance(value, str):
                if self._is_valid_color(value):
                    return value
                warn(name, f"contains an invalid color '{value}'", fallback_label)
                return fallback
            if isinstance(value, (tuple, list)):
                # Single RGB(A) color tuple: 3-4 numeric values
                if len(value) in (3, 4) and all(isinstance(x, (int, float)) for x in value):
                    if self._is_valid_color(value):
                        return value
                    warn(name, f"contains an invalid color {value!r}", fallback_label)
                    return fallback
                if len(value) == 0:
                    if not silent_empty:
                        warn(name, "is empty", fallback_label)
                    return fallback
                # List of colors: drop unparseable entries, keep the rest
                valid = [c for c in value if self._is_valid_color(c)]
                dropped = len(value) - len(valid)
                if dropped and valid:
                    self._vprint(
                        f"\033[33m⚠️  Warning: {name} contains {dropped} invalid color(s) that were dropped.\033[0m",
                        level='simple'
                    )
                elif dropped:
                    warn(name, f"contains {dropped} invalid color(s)", fallback_label)
                return valid if valid else fallback
            warn(name, f"has unsupported type {type(value).__name__}", fallback_label)
            return fallback

        def normalize_single(value, name, fallback='auto'):
            """Normalize a parameter that accepts exactly one color (or 'auto')."""
            if isinstance(value, str):
                if value.strip().lower() == 'auto' or self._is_valid_color(value):
                    return value
                warn(name, f"contains an invalid color '{value}'", f"'{fallback}'")
                return fallback
            if isinstance(value, (tuple, list)):
                if len(value) in (3, 4) and all(isinstance(x, (int, float)) for x in value) \
                        and self._is_valid_color(value):
                    return value
                warn(name, f"contains an invalid color {value!r}", f"'{fallback}'")
                return fallback
            warn(name, f"has unsupported type {type(value).__name__}", f"'{fallback}'")
            return fallback

        # neuron_colors: None triggers the background-dependent default palette below
        self.neuron_colors = normalize_sequence(
            self.neuron_colors, 'neuron_colors', None, 'default palette', allow_none=True
        )
        # synapse_colors: fall back to the field default palette. An empty
        # list is the expected input when there is only one neuron layer
        # (no inter-layer connections), so fall back silently in that case.
        n_layers = len(self.neuron_layers) if isinstance(self.neuron_layers, list) else 1
        self.synapse_colors = normalize_sequence(
            self.synapse_colors, 'synapse_colors',
            list(bokeh.palettes.Category10[10]), 'Category10 palette',
            silent_empty=(n_layers <= 1)
        )
        # mesh_color: fall back to the default gray (None is resolved by validation)
        self.mesh_color = normalize_sequence(
            self.mesh_color, 'mesh_color', (100, 100, 100), '(100, 100, 100)', allow_none=True
        )
        # brain/VNC mesh colors: 'auto' or a single valid color
        self.brain_mesh_color = normalize_single(self.brain_mesh_color, 'brain_mesh_color')
        self.vnc_mesh_color = normalize_single(self.vnc_mesh_color, 'vnc_mesh_color')
        # background_color: any parseable color string, else white
        if not isinstance(self.background_color, str) or not self._is_valid_color(self.background_color):
            warn('background_color', f"value {self.background_color!r} is not a valid color", "'white'")
            self.background_color = 'white'

    def __post_init__(self):
        # Empty or invalid color inputs fall back to their default palettes
        # (with warnings) instead of failing validation.
        self._normalize_color_inputs()

        # Default neuron palette follows the background color: bokeh
        # Category10 on white, bokeh Set3 on black (resolved before the
        # color validation below runs).
        if self.neuron_colors is None:
            if str(self.background_color).strip().lower() == 'black':
                self.neuron_colors = bokeh.palettes.Set3[12]
            else:
                self.neuron_colors = bokeh.palettes.Category10[10]

        # Apply dataset-specific defaults BEFORE validation
        if self.dataset and 'manc' in self.dataset.lower():
            # For MANC, enable VNC mesh by default if not explicitly disabled
            # We assume if the user passed False explicitly, they know what they are doing.
            # But dataclasses don't track "default vs explicit" easily.
            # So we'll just set it to True if it's currently False, assuming default was False.
            # The default in field definition is False.
            # EXCEPTION: If brain_mesh='none', we assume user wants no context meshes at all,
            # so we do NOT enable VNC mesh by default.
            if self.vnc_mesh is False and self.brain_mesh != 'none':
                # We can't distinguish "User set False" vs "Default False".
                # To follow user request "set vnc_mesh=True by default":
                self.vnc_mesh = True
        
        # Check for elastix dependency if MANC 'whole' mode is requested
        if 'manc' in self.dataset.lower() and self.brain_mesh == 'whole':
            import shutil
            if shutil.which('elastix') is None:
                print('⚠️  Elastix not found: Cannot transform MANC to Male-CNS space.')
                print('   automatically changing brain_mesh from "whole" to "template".')
                print('   Tip: Install elastix or use male-cns dataset for full CNS context.')
                self.brain_mesh = 'template'
        
        # === INPUT VALIDATION ===
        # Validate all input parameters before processing
        self._validate_inputs()
        
        # Normalize verbose parameter
        if self.verbose is True:
            self.verbose = 'full'
        elif self.verbose is False:
            self.verbose = False
        
        # Silence navis INFO messages (like "Use the `.show()` method to plot the figure.")
        # These are not useful for automated visualization and clutter output
        try:
            # Set to ERROR to suppress "radii are missing" warnings which are common/harmless
            # when generating meshes from skeletons without radius info
            navis.set_loggers('ERROR') 
        except Exception:
            pass  # Ignore if function not available in older versions
            
        # Ensure cleanup of logging handlers if needed
        logging.getLogger('navis').setLevel(logging.ERROR)

        # Silence navis and other libraries' debug output if verbose is not full
        if self.verbose != 'full':
            logging.getLogger('navis').setLevel(logging.ERROR)
            logging.getLogger('trimesh').setLevel(logging.ERROR)
        
        # Initialize output_dir if not set
        if self.output_dir is None:
            self.output_dir = self.data_folder

        # Initialize list to store meshes for export
        self.exportable_meshes = []
        
        # Auto-detect client_type from dataset if not explicitly set to flywire
        if self.client_type == 'neuprint' and ('flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()):
            self.client_type = 'flywire'
            self._vprint(f"Auto-detected client_type='flywire' from dataset '{self.dataset}'", level='full')

        # For FlyWire/FAFB: Enable mesh caching (transformed+meshed), disable raw skeleton caching
        if self.client_type == 'flywire' or 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower():
            # Raw skeleton pkl caching is disabled (files too large and need transformation anyway)
            if self.cache_neurons:
                self._vprint("  ℹ️  FlyWire/FAFB: Using mesh cache (simplified) instead of raw skeletons", level='full')
            if self.cache_synapses:
                self._vprint("  ℹ️  Disabling synapse caching for FlyWire/FAFB (files too large)", level='full')
                self.cache_synapses = False
        
        # Set default skeleton_mesh_simplification based on dataset if not specified
        if self.skeleton_mesh_simplification is None:
            if 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower():
                # FAFB meshes are very high detail, use 0.95 (keep 5% of faces)
                self.skeleton_mesh_simplification = 0.95
                self._vprint(f"  ℹ️  Using default skeleton_mesh_simplification=0.95 for FAFB (high-detail meshes)", level='full')
            else:
                # NeuPrint datasets (hemibrain, male-cns, manc, optic-lobe) use 0.9
                self.skeleton_mesh_simplification = 0.9
                self._vprint(f"  ℹ️  Using default skeleton_mesh_simplification=0.9 for {self.dataset}", level='full')

        # Auto-detect version from dataset if not provided
        if self.client_type == 'flywire' and self.version is None:
            import re
            # Look for v783 or version 783
            match = re.search(r'v(\d+)', self.dataset)
            if match:
                self.version = int(match.group(1))
                self._vprint(f"Auto-detected version={self.version} from dataset '{self.dataset}'", level='full')

        # Initialize client if needed
        if self.client_type == 'neuprint':
            import neuprint
            
            # Use provided client if available
            if self.client is not None:
                self._vprint(f'Using provided client for {self.dataset}', level='full')
            else:
                # Check if global client exists AND matches our dataset
                # This is critical: if a default client exists for a DIFFERENT dataset,
                # we MUST create a new client for the correct dataset, otherwise
                # fetch_skeletons will fail with "No neurons matching the given criteria found!"
                use_existing_client = False
                try:
                    existing_client = neuprint.default_client()
                    if existing_client is not None:
                        # Check if datasets match (normalize for comparison)
                        existing_ds = existing_client.dataset.lower().replace('_', ':').replace('.', '.')
                        target_ds = self.dataset.lower().replace('_', ':').replace('.', '.')
                        if existing_ds == target_ds:
                            use_existing_client = True
                            self._vprint(f'Using existing default client for {self.dataset}', level='full')
                        else:
                            self._vprint(f'Default client is for {existing_client.dataset}, need new client for {self.dataset}', level='full')
                except RuntimeError:
                    pass

                if not use_existing_client:
                    # Use TokenManager to get token
                    try:
                        from utils.token_manager import token_manager
                        self.token = token_manager.get_token('NEUPRINT_TOKEN', self.token)
                    except ImportError:
                        pass

                    if self.token:
                        self.client = Client(self.server, dataset=self.dataset, token=self.token)
                        self.client.fetch_version()
                        # Set as default to avoid "multiple clients" error
                        neuprint.set_default_client(self.client)
                        self._vprint(f'Client initialized for {self.dataset} (set as default)', level='full')
                    elif os.environ.get('NEUPRINT_APPLICATION_CREDENTIALS'):
                        # Auto-detect from env
                        self.client = Client(self.server, dataset=self.dataset)
                        self.client.fetch_version()
                        # Set as default to avoid "multiple clients" error
                        neuprint.set_default_client(self.client)
                        self._vprint(f'Client initialized from env for {self.dataset} (set as default)', level='full')
                    else:
                        # Only warn if we are not using local cache/files exclusively
                        # But we don't know that yet.
                        pass
        
        # Initialize FlyWire client if needed
        if self.client_type == 'flywire' and self.client_flywire is None:
            # FlyWire API fetching removed
            pass

        # Check FlyWire visualization files
        if 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower():
            # Ensure data is prepared using the converter
            dataset_dir = os.path.join(self.script_path, 'datasets', self.dataset)
            
            # Use the converter module to ensure data is ready
            if 'BANC' in self.dataset:
                success = BANC_file_converter.ensure_banc_data(self.dataset, dataset_dir)
            else:
                success = FAFB_file_converter.ensure_flywire_data(self.dataset, dataset_dir)

            if not success:
                print("\\n\033[31mCRITICAL ERROR: FlyWire/BANC data preparation failed.\033[0m")
                print("Please follow the instructions above to download the required files.")
                sys.exit(1)
            
            try:
                import fafb_utils
                # Check for skeleton zip
                if not os.path.exists(dataset_dir):
                    dataset_dir = os.path.join(self.script_path, 'datasets', 'flywire_FAFB_v783')
                
                if os.path.exists(dataset_dir):
                    sk_zip = fafb_utils.get_fafb_skeleton_zip(dataset_dir)
                    if not sk_zip:
                        print(f'\033[31mWarning: FlyWire skeleton zip not found in {dataset_dir}\033[0m')
                        if 'BANC' in self.dataset:
                            print('Skeleton visualization not available for BANC, because the skeleton data not available in flywire codex')
                        else:
                            print(f'Please download sk_lod1_783_healed.zip from https://codex.flywire.ai/api/download?dataset=fafb')
                        print(f'Visualization might fail or be incomplete.')
                        sys.exit(0)
            except ImportError:
                pass

        if self.synapse_mode not in ['scatter', 'sphere', 'cone', 'tetrahedron']:
            raise ValueError('synapse_mode can only be "scatter", "sphere", "cone", or "tetrahedron"')
        
        # Validate legend_mode
        if self.legend_mode not in ['single', 'type', 'layer']:
            raise ValueError('legend_mode must be "single", "type", or "layer"')

        if self.color_mode not in ['per_layer', 'per_neuron']:
            raise ValueError('color_mode must be "per_layer" or "per_neuron"')
        
        # Validate expand_colors
        if self.expand_colors not in ['interpolation', 'darken', 'cycle']:
            raise ValueError('expand_colors must be "interpolation", "darken", or "cycle"')
        
        if self.skeleton_mode not in ['line','tube']:
            raise ValueError('skeleton_mode can only be "line" or "tube"')
        if self.brain_mesh not in ['none', 'whole', 'template']:
            raise ValueError('brain_mesh must be "none", "template", or "whole"')
        if self.backend not in ['plotly', 'k3d']:
            raise ValueError('backend must be "plotly" or "k3d"')
        
        # Early transform check - advise user on transformation-free options
        self._check_transform_requirements_early()
        
        # Parse layer_map_csv if provided (overrides neuron_layers and custom_layer_names)
        if self.layer_map_csv is not None:
            self._parse_layer_map_csv()
        
        # convert neuron_layers str to list, if is str
        if type(self.neuron_layers) is str:
            if self.neuron_layers.strip():
                # Non-empty string: parse by '->' separator
                self.neuron_layers = self.neuron_layers.replace(' ','').split('->')
                for i,layer in enumerate(self.neuron_layers): # convert bodyId str to int
                    if layer.isnumeric():
                        self.neuron_layers[i] = int(layer)
            else:
                # Empty string: convert to empty list for mesh-only mode
                self.neuron_layers = []
        
        if self.synapse_mode == 'scatter' and self.synapse_size == 0:
            self.synapse_size = 2
        # elif self.synapse_mode in ['sphere', 'cone', 'tetrahedron']:
        #     # Only check size limit if synapse_size is a number (not 'real')
        #     # For these modes, synapse_size is a multiplier, so small values (e.g. 1.0) are valid
        #     if isinstance(self.synapse_size, (int, float)) and self.synapse_size < 20 and self.brain_mesh != 'whole':
        #         self.synapse_size = 20
        #         self._vprint('\033[33mSynapse size is too small (< 20) for sphere, cone, or tetrahedron mode, automatically reset to 20\033[0m', level='full')
            
        if self.mesh_roi == None:
            self.mesh_roi = []
        
        # Convert single string to list for mesh_roi
        if isinstance(self.mesh_roi, str):
            self.mesh_roi = [self.mesh_roi]
        
        # Initialize nested ROI groups (will be populated if nested lists are used)
        self._nested_roi_groups = {}
        
        # Handle nested lists in mesh_roi for color grouping
        # e.g., ['AME', ['aL', 'bL', 'gL'], 'EB'] -> flattened list with same colors for grouped ROIs
        if self.mesh_roi:
            self.mesh_roi, self.mesh_color, self._nested_roi_groups = self._flatten_nested_roi_groups(
                self.mesh_roi, self.mesh_color
            )
        
        # Expand special keywords and regex patterns in mesh_roi
        # 'primary' -> all primary ROIs, 'all' -> all available ROIs
        # 'ME.*' -> all ROIs matching regex pattern
        if self.mesh_roi:
            self.mesh_roi = self._expand_mesh_roi_patterns(self.mesh_roi)
        
        # Expand ROI names to include bilateral (L/R) variants
        # e.g., 'LH' -> ['LH(L)', 'LH(R)']
        # Also expand colors to match the expanded ROIs
        if self.mesh_roi:
            original_rois = list(self.mesh_roi)
            original_colors = self.mesh_color if isinstance(self.mesh_color, list) else self.mesh_color
            
            # Use the new function that expands both ROIs and colors together
            expanded_rois, expanded_colors = self._expand_roi_names_with_colors(
                self.mesh_roi, 
                self.mesh_color
            )
            
            self.mesh_roi = expanded_rois
            # Only update mesh_color if it was a list (preserving single color behavior)
            if isinstance(original_colors, list):
                self.mesh_color = expanded_colors
            elif len(set(expanded_colors)) == 1:
                # If all expanded colors are the same, keep as single color
                self.mesh_color = expanded_colors[0]
            else:
                self.mesh_color = expanded_colors
            
            if self.mesh_roi != original_rois:
                self._vprint(f"   🔄 ROI expansion: {original_rois} → {self.mesh_roi}", level='simple')
                if isinstance(original_colors, list) and len(original_colors) != len(expanded_colors):
                    self._vprint(f"   🔄 Color expansion: {len(original_colors)} colors → {len(expanded_colors)} colors", level='full')
        
        # === Standardize color inputs ===
        # Standardize neuron_colors to a list of rgba strings
        # Check if neuron_colors have explicit alpha values that override neuron_alpha
        self._neuron_colors_have_explicit_alpha = self._colors_have_explicit_alpha(self.neuron_colors)
        if self._neuron_colors_have_explicit_alpha:
            self._vprint(
                f"\033[33m⚠️  Warning: neuron_colors contains explicit alpha values. "
                f"These will override neuron_alpha={self.neuron_alpha}. "
                f"To use uniform alpha, remove alpha from colors.\033[0m",
                level='simple'
            )
            # Use alpha=1.0 as placeholder, the explicit alpha from colors will be preserved
            self.neuron_colors = self._standardize_color_input(self.neuron_colors, 'neuron_colors', default_alpha=1.0)
        else:
            self.neuron_colors = self._standardize_color_input(self.neuron_colors, 'neuron_colors', self.neuron_alpha)
        self._base_neuron_colors = tuple(self.neuron_colors)
        
        # Standardize synapse_colors
        # Check if synapse_colors have explicit alpha values that override synapse_alpha
        self._synapse_colors_have_explicit_alpha = self._colors_have_explicit_alpha(self.synapse_colors)
        if self._synapse_colors_have_explicit_alpha:
            self._vprint(
                f"\033[33m⚠️  Warning: synapse_colors contains explicit alpha values. "
                f"These will override synapse_alpha={self.synapse_alpha}. "
                f"To use uniform alpha, remove alpha from colors.\033[0m",
                level='simple'
            )
            self.synapse_colors = self._standardize_color_input(self.synapse_colors, 'synapse_colors', default_alpha=1.0)
        else:
            self.synapse_colors = self._standardize_color_input(self.synapse_colors, 'synapse_colors', self.synapse_alpha)
        
        # Standardize mesh_color (for ROI meshes)
        # Store original for detecting custom colors
        self._original_mesh_color = self.mesh_color
        
        # Check if mesh_color has explicit alpha values that override mesh_alpha
        self._mesh_colors_have_explicit_alpha = self._colors_have_explicit_alpha(self.mesh_color)
        if self._mesh_colors_have_explicit_alpha:
            self._vprint(
                f"\033[33m⚠️  Warning: mesh_color contains explicit alpha values. "
                f"These will override mesh_alpha={self.mesh_alpha}. "
                f"To use uniform alpha, remove alpha from colors.\033[0m",
                level='simple'
            )
            # Use alpha=1.0 as placeholder, the explicit alpha from colors will be preserved
            self.mesh_color = self._standardize_mesh_color_input(self.mesh_color, default_alpha=1.0)
        else:
            self.mesh_color = self._standardize_mesh_color_input(self.mesh_color, default_alpha=self.mesh_alpha)
        
        # Warn if custom mesh colors are provided that brain/VNC mesh colors are separate
        if self._is_custom_mesh_color_specified() and self.mesh_roi:
            self._vprint(
                f"\033[34mℹ️  Note: mesh_color applies to ROI meshes only. "
                f"Brain mesh and VNC mesh colors are controlled by brain_mesh_color and vnc_mesh_color.\033[0m",
                level='simple'
            )
        
        # Standardize brain_mesh_color and vnc_mesh_color if not 'auto'
        if isinstance(self.brain_mesh_color, (tuple, list)):
            self.brain_mesh_color = standardize_color(self.brain_mesh_color, default_alpha=0.1)
        if isinstance(self.vnc_mesh_color, (tuple, list)):
            self.vnc_mesh_color = standardize_color(self.vnc_mesh_color, default_alpha=0.1)
        
        # Ensure enough colors for all layers by expanding if needed
        n_layers = len(self.neuron_layers)
        self.neuron_colors = self._expand_color_sequence(
            self._base_neuron_colors,
            n_layers,
            target_label='layers',
            tip_parameter='neuron_colors',
            warn=self.color_mode == 'per_layer',
        )
        
        # Same for synapse colors (one fewer than neuron layers for connections between layers)
        n_synapse_needed = max(0, n_layers - 1)
        self.synapse_colors = self._expand_color_sequence(
            tuple(self.synapse_colors),
            n_synapse_needed,
            target_label='synapse layers',
        )
        
        if self.skeleton_mode == 'line':
            self.show_skeleton_radius = False
            # neuron_alpha is now supported for line mode via opacity
        elif self.skeleton_mode == 'tube':
            self.show_skeleton_radius = True
        
        # fetch neurons and automatically generate layer names
        self.neuron_dfs = []
        self.roi_dfs = []
        self.layer_criteria = []
        self.layer_names = []
        
        n_layers = len(self.neuron_layers)
        self._vprint(f'\n📊 Fetching neuron info for {n_layers} layer(s)...')
        
        # Use tqdm for progress bar
        from tqdm import tqdm
        layer_iter = tqdm(range(n_layers), desc="Loading layers", disable=self.verbose != 'full')
        
        total_neurons = 0
        for i in layer_iter:
            layer_input = self.neuron_layers[i]
            if not isinstance(layer_input, list):
                layer_input = [layer_input]
            
            # Update progress bar description
            layer_desc = str(layer_input[0])[:20] if layer_input else f"layer_{i}"
            layer_iter.set_description(f"Layer {i}: {layer_desc}")
            
            ndf, rdf, auto_name, cri = sv.getNeurons(
                layer_input,
                dataset=self.dataset,
                client=self.client,
                verbose=False,
                search_columns=self.search_columns,
            )
            if self.hemisphere != 'both' and ndf is not None and not ndf.empty:
                ndf, rdf = self._filter_neuron_df_by_hemisphere(ndf, rdf)
            self.neuron_dfs.append(ndf)
            self.roi_dfs.append(rdf)
            self.layer_criteria.append(cri)
            self.layer_names.append(auto_name)
            
            n_neurons = len(ndf) if ndf is not None else 0
            total_neurons += n_neurons
            
            # Update postfix with neuron count
            layer_iter.set_postfix(neurons=n_neurons, total=total_neurons)
        
        # Print summary
        self._vprint(f'✓ Loaded {total_neurons:,} neurons across {n_layers} layers')
        
        # Show detailed breakdown if full verbose
        if self.verbose == 'full':
            self._vprint('\n  Layer summary:')
            for i, (ndf, name) in enumerate(zip(self.neuron_dfs, self.layer_names)):
                n = len(ndf) if ndf is not None else 0
                if n > 0 and 'type' in ndf.columns:
                    types = ndf['type'].dropna().unique()
                    n_types = len(types)
                    type_preview = ', '.join(str(t) for t in types[:3])
                    if n_types > 3:
                        type_preview += f' (+{n_types-3} more)'
                    self._vprint(f'    [{i}] {name}: {n} neurons, {n_types} types ({type_preview})')
                else:
                    self._vprint(f'    [{i}] {name}: {n} neurons')

        # Generate smart layer names based on types (if not using custom names)
        if not self.custom_layer_names:
            self.layer_names = self._generate_smart_layer_names()
        else:
            # Support partial custom_layer_names - merge with auto-generated names
            auto_names = self._generate_smart_layer_names()
            n_layers = len(self.neuron_layers)
            n_custom = len(self.custom_layer_names)
            
            merged_names = []
            for i in range(n_layers):
                if i < n_custom and self.custom_layer_names[i]:
                    # Use custom name if provided and non-empty
                    merged_names.append(self.custom_layer_names[i])
                elif i < len(auto_names):
                    # Use auto-generated name
                    merged_names.append(auto_names[i])
                else:
                    # Fallback
                    merged_names.append(f"layer_{i}")
            
            self.layer_names = merged_names

        self._per_neuron_colors = {}
        if self.color_mode == 'per_neuron':
            self._per_neuron_colors = self._build_per_neuron_color_map()
            
        if self.saveas is None:
            # Generate saveas from layer names
            # Use exact names for ≤3 layers, first 2 + etc{remaining} for >3 layers
            n_layers = len(self.layer_names) if self.layer_names else 0
            if n_layers == 0:
                # Mesh-only mode without neurons
                if self.mesh_roi:
                    n_rois = len(self.mesh_roi)
                    if n_rois == 0:
                        self.saveas = "brain_mesh"
                    elif n_rois <= 3:
                        # Use all ROI names for ≤3 ROIs
                        roi_names = []
                        for roi in self.mesh_roi[:3]:
                            # Clean up ROI name for filename (remove parentheses)
                            roi_clean = str(roi).replace('(', '').replace(')', '')
                            roi_names.append(roi_clean)
                        self.saveas = '_'.join(roi_names) + '_mesh'
                    else:
                        # Use first 2 ROI names + count for >3 ROIs
                        roi_names = []
                        for roi in self.mesh_roi[:2]:
                            roi_clean = str(roi).replace('(', '').replace(')', '')
                            roi_names.append(roi_clean)
                        self.saveas = '_'.join(roi_names) + f'_etc{n_rois - 2}_mesh'
                else:
                    self.saveas = "brain_mesh"
            elif n_layers <= 3:
                # Use all layer names for ≤3 layers
                self.saveas = '_'.join(self.layer_names)
            else:
                # Use first 2 names + count of remaining layers for >3 layers
                first_two = '_'.join(self.layer_names[:2])
                remaining = n_layers - 2
                self.saveas = f"{first_two}_etc{remaining}"
        
        # Ensure saveas doesn't exceed reasonable length (max 80 chars)
        if len(self.saveas) > 80:
            # Truncate and add hash for uniqueness
            import hashlib
            hash_suffix = hashlib.md5('_'.join(self.layer_names).encode()).hexdigest()[:6]
            self.saveas = self.saveas[:70] + f"_{hash_suffix}"
        
        # Get dataset abbreviation for folder naming
        dataset_abbrev = self._get_dataset_abbreviation()
        
        # Create output subfolder (with or without timestamp based on include_timestamp)
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        base_folder_name = f'plot3d_{dataset_abbrev}_' + self.saveas.split('.')[0]
        if self.include_timestamp:
            self.save_folder = os.path.join(self.output_dir, base_folder_name + '_' + timestamp)
        else:
            self.save_folder = os.path.join(self.output_dir, base_folder_name)
        if not os.path.exists(self.save_folder): os.makedirs(self.save_folder)
        
        # Save parameters to text file with comprehensive formatting
        param_file = os.path.join(self.save_folder, 'parameters.txt')
        with open(param_file, 'w') as f:
            f.write("=" * 70 + "\n")
            f.write("VisualizeSkeleton Parameters - Complete Configuration\n")
            f.write("=" * 70 + "\n\n")
            
            # Basic Info
            f.write("[Basic Info]\n")
            f.write(f"  Dataset:          {self.dataset}\n")
            f.write(f"  Timestamp:        {timestamp}\n")
            if self.version:
                f.write(f"  Version:          {self.version}\n")
            f.write(f"  Client Type:      {self.client_type}\n")
            f.write(f"  Server:           {self.server}\n")
            f.write(f"  Output Directory: {self.output_dir}\n")
            f.write(f"  Save Filename:    {self.saveas}\n")
            f.write(f"  Include Timestamp:{self.include_timestamp}\n")
            f.write(f"  Verbose:          {self.verbose}\n")
            f.write("\n")
            
            # Layer Info
            f.write("[Layers]\n")
            for i, (layer, name) in enumerate(zip(self.neuron_layers, self.layer_names)):
                n_neurons = len(self.neuron_dfs[i]) if i < len(self.neuron_dfs) and self.neuron_dfs[i] is not None else 0
                f.write(f"  Layer {i}: {name} ({n_neurons} neurons)\n")
                # Show first few neuron IDs if available
                if n_neurons > 0 and 'bodyId' in self.neuron_dfs[i].columns:
                    body_ids = self.neuron_dfs[i]['bodyId'].tolist()[:5]
                    ids_str = ', '.join(str(bid) for bid in body_ids)
                    if n_neurons > 5:
                        ids_str += f", ... (+{n_neurons - 5} more)"
                    f.write(f"           IDs: {ids_str}\n")
            f.write(f"  Total Layers:     {len(self.neuron_layers)}\n")
            f.write("\n")
            
            # Color Settings
            f.write("[Color Settings]\n")
            if self.neuron_colors:
                if isinstance(self.neuron_colors, list):
                    colors_preview = self.neuron_colors[:5]
                    f.write(f"  Neuron Colors:    {colors_preview}")
                    if len(self.neuron_colors) > 5:
                        f.write(f" ... (+{len(self.neuron_colors) - 5} more)")
                    f.write("\n")
                else:
                    f.write(f"  Neuron Colors:    {self.neuron_colors}\n")
            else:
                f.write(f"  Neuron Colors:    (auto-assigned)\n")
            f.write(f"  Neuron Alpha:     {self.neuron_alpha}\n")
            f.write(f"  Synapse Alpha:    {self.synapse_alpha}\n")
            f.write(f"  Brain Mesh Color: {self.brain_mesh_color}\n")
            f.write(f"  VNC Mesh Color:   {self.vnc_mesh_color}\n")
            if self.mesh_color:
                f.write(f"  Mesh ROI Color:   {self.mesh_color}\n")
            f.write("\n")
            
            # Visualization Settings
            f.write("[Visualization]\n")
            f.write(f"  Skeleton Mode:    {self.skeleton_mode}\n")
            f.write(f"  Backend:          {self.backend}\n")
            f.write(f"  Brain Mesh:       {self.brain_mesh}\n")
            f.write(f"  VNC Mesh:         {self.vnc_mesh}\n")
            if self.mesh_roi:
                f.write(f"  Mesh ROI:         {self.mesh_roi}\n")
            f.write(f"  Legend Mode:      {self.legend_mode}\n")
            f.write(f"  Expand Colors:    {self.expand_colors}\n")
            f.write(f"  Show Soma:        {self.show_soma}\n")
            f.write("\n")
            
            # Synapse Settings
            f.write("[Synapse Settings]\n")
            f.write(f"  Synapse Mode:     {self.synapse_mode}\n")
            f.write(f"  Synapse Size:     {self.synapse_size}\n")
            f.write(f"  Min Synapse Num:  {self.min_synapse_num}\n")
            f.write(f"  Synapse Alpha:    {self.synapse_alpha}\n")
            f.write("\n")
            
            f.write("=" * 70 + "\n")
        if self.backend == 'plotly':
            self.fig_3d = go.Figure()
        elif self.backend == 'k3d':
            try:
                import k3d
                self.fig_3d = k3d.plot()
            except ImportError:
                self._vprint("⚠️  k3d not installed. Please install it with `pip install k3d`")
                self._vprint("   Falling back to plotly backend")
                self.backend = 'plotly'
                self.fig_3d = go.Figure()
        
        # save neuron dataframes to excel file
        file_path = os.path.join(self.save_folder, self.saveas+'_neuron_info.xlsx')
        for i in range(len(self.neuron_layers)):
            if i == 0:
                mode = 'w'
            else:
                mode = 'a'
            with pd.ExcelWriter(file_path,mode=mode,engine='openpyxl') as writer:
                self.neuron_dfs[i].to_excel(writer, sheet_name=f'neuron_df{i}')
                self.roi_dfs[i].to_excel(writer, sheet_name=f'roi_count_df{i}')
    
    def _get_cache_path(self, cache_type):
        """Get the cache directory for skeletons or synapses
        
        Uses project cache/ folder for organized storage:
        cache/{dataset}/skeletons/ - for individual skeleton .pkl files
        cache/{dataset}/synapses/ - for synapse cache files
        
        For datasets folder resources:
        datasets/{dataset}/*_synapse_table.parquet - synapse table
        
        Example:
        - cache/hemibrain_v1_2_1/skeletons/{bodyId}.pkl
        - datasets/flywire_FAFB_v783/flywire_FAFB_v783_synapse_table.parquet
        """
        dataset_normalized = self.dataset.replace(':', '_').replace('.', '_')
        cache_dir = os.path.join(self.script_path, 'cache', dataset_normalized, cache_type)
        os.makedirs(cache_dir, exist_ok=True)
        return cache_dir
    
    def _get_dataset_abbreviation(self) -> str:
        """Get a short abbreviation for the dataset name for use in folder naming.
        
        Maps common dataset names to short abbreviations:
        - hemibrain -> HEMI
        - male-cns -> MCNS
        - flywire_FAFB -> FAFB
        - flywire_BANC -> BANC
        - optic-lobe -> OL
        - manc -> MANC
        
        Returns
        -------
        str
            Short abbreviation for the dataset, or first 4 chars if not recognized
        """
        # Dataset abbreviation mapping
        abbrev_map = {
            'hemibrain': 'HEMI',
            'male-cns': 'MCNS',
            'male_cns': 'MCNS',
            'flywire_fafb': 'FAFB',
            'fafb': 'FAFB',
            'flywire_banc': 'BANC',
            'banc': 'BANC',
            'optic-lobe': 'OL',
            'optic_lobe': 'OL',
            'manc': 'MANC',
        }
        
        if self.dataset is None:
            return 'UNKN'
        
        # Try exact match first (lowercase)
        dataset_lower = self.dataset.lower()
        for key, abbrev in abbrev_map.items():
            if key in dataset_lower:
                return abbrev
        
        # If no match, use first 4 chars uppercase
        return self.dataset[:4].upper().replace('-', '').replace('_', '')
    
    def _generate_smart_layer_names(self) -> List[str]:
        """Generate smart layer names based on neuron types.
        
        For each layer, generates a name in format:
        - {type} if all neurons in layer are the same type (even if multiple neurons)
        - {type}_etc if multiple neurons with different types (uses most common type)
        - {bodyId} if single untyped neuron
        - {bodyId}_etc if multiple untyped neurons with different IDs
        
        This method looks at the neuron_dfs to determine types and counts.
        
        Returns:
            List[str]: Smart layer names for each layer
        """
        smart_names = []
        
        for i, ndf in enumerate(self.neuron_dfs):
            if ndf is None or len(ndf) == 0:
                # Fallback to original auto-generated name
                smart_names.append(self.layer_names[i] if i < len(self.layer_names) else f"layer_{i}")
                continue
            
            n_neurons = len(ndf)
            
            # Get type column (different datasets may use different column names)
            type_col = None
            for col in ['type', 'cell_type', 'neuronType']:
                if col in ndf.columns:
                    type_col = col
                    break
            
            # Get types from the dataframe
            if type_col and type_col in ndf.columns:
                types = ndf[type_col].dropna().unique().tolist()
                # Filter out empty strings and None
                types = [t for t in types if t and str(t).strip()]
            else:
                types = []
            
            if types:
                # Count unique types
                n_unique_types = len(types)
                
                # Use the most common type as the representative
                if type_col in ndf.columns:
                    type_counts = ndf[type_col].value_counts()
                    primary_type = type_counts.index[0] if len(type_counts) > 0 else types[0]
                else:
                    primary_type = types[0]
                
                # Only add _etc if there are multiple different types
                if n_unique_types > 1:
                    smart_names.append(f"{primary_type}_etc")
                else:
                    # All neurons are the same type
                    smart_names.append(str(primary_type))
            else:
                # No type info - use bodyId
                body_ids = ndf['bodyId'].tolist() if 'bodyId' in ndf.columns else []
                if body_ids:
                    first_id = body_ids[0]
                    # Multiple untyped neurons with different IDs -> _etc
                    if n_neurons > 1:
                        smart_names.append(f"{first_id}_etc")
                    else:
                        smart_names.append(str(first_id))
                else:
                    # Ultimate fallback
                    smart_names.append(self.layer_names[i] if i < len(self.layer_names) else f"layer_{i}")
        
        return smart_names

    def _get_synapse_table_path(self):
        """Get path to synapse table in datasets folder.
        
        Returns the path to the synapse table parquet file.
        For FlyWire/FAFB: datasets/flywire_FAFB_v783/flywire_FAFB_v783_synapse_table.parquet
        For NeuPrint: datasets/{dataset}/{dataset}_synapse_table.parquet
        
        Returns:
            str: Path to synapse table, or None if not found
        """
        dataset_normalized = self.dataset.replace(':', '_').replace('.', '_')
        datasets_dir = os.path.join(self.script_path, 'datasets', dataset_normalized)
        
        # Look for synapse table file
        parquet_file = os.path.join(datasets_dir, f"{dataset_normalized}_synapse_table.parquet")
        
        if os.path.exists(parquet_file):
            return parquet_file
        
        # Fallback: try FAFB naming if dataset includes flywire
        if 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower():
            fafb_dir = os.path.join(self.script_path, 'datasets', 'flywire_FAFB_v783')
            fafb_file = os.path.join(fafb_dir, 'flywire_FAFB_v783_synapse_table.parquet')
            if os.path.exists(fafb_file):
                return fafb_file
        
        return None
    
    def _preload_fafb_skeletons(self, body_ids_filter=None):
        """Pre-load FAFB skeletons from ZIP file in a single batch.
        
        This is much faster than opening the ZIP file for each layer.
        
        Parameters
        ----------
        body_ids_filter : set, optional
            If provided, only load these bodyIds from the ZIP.
            If None, load all bodyIds from self.neuron_dfs.
        
        Returns:
            dict: bodyId -> TreeNeuron mapping
        """
        from tqdm import tqdm
        import sys
        
        # Collect all body IDs needed
        if body_ids_filter is not None:
            all_body_ids = set(body_ids_filter)
        else:
            # Collect from all layers
            all_body_ids = set()
            for df in self.neuron_dfs:
                if df is not None:
                    all_body_ids.update(df['bodyId'].tolist())
        
        if not all_body_ids:
            return {}
        
        skeleton_cache = {}
        
        try:
            import fafb_utils
            project_root = os.path.dirname(os.path.dirname(__file__))
            
            # Try to find dataset directory by name
            data_dir = os.path.join(project_root, "datasets", self.dataset)
            if not os.path.exists(data_dir):
                data_dir = os.path.join(project_root, "datasets", "flywire_FAFB_v783")
            
            zip_path = fafb_utils.get_fafb_skeleton_zip(data_dir)
            
            if zip_path:
                import zipfile
                import io
                
                self._vprint(f'  📦 Loading {len(all_body_ids)} skeletons from ZIP...')
                
                with zipfile.ZipFile(zip_path, 'r') as z:
                    zip_files = set(z.namelist())
                    
                    # Progress bar for skeleton loading
                    pbar = tqdm(all_body_ids, desc="  Loading skeletons", 
                               disable=self.verbose != 'full', leave=False, file=sys.stdout)
                    
                    for bid in pbar:
                        filename = f"{bid}.swc"
                        try:
                            if filename in zip_files:
                                with z.open(filename) as f:
                                    content = f.read().decode('utf-8')
                                    n = navis.read_swc(io.StringIO(content))
                                    n.units = 'nm'
                                    n.id = bid
                                    n.name = str(bid)
                                    skeleton_cache[bid] = n
                        except Exception:
                            pass  # Skip errors silently
                
                self._vprint(f'  ✓ Loaded {len(skeleton_cache)}/{len(all_body_ids)} skeletons from ZIP')
        except ImportError:
            pass
        except Exception as e:
            self._vprint(f'  ⚠️  Error pre-loading FAFB skeletons: {e}')
        
        return skeleton_cache

    def _load_extrusion_check_cache(self):
        """Load cached extrusion check results from parquet.
        
        Returns
        -------
        dict
            Dictionary of bodyId -> bool (True if has extrusion)
        """
        import pandas as pd
        
        project_root = os.path.dirname(os.path.dirname(__file__))
        dataset_safe = self.dataset.replace(':', '_').replace('.', '_')
        cache_file = os.path.join(project_root, 'cache', dataset_safe, 'extrusion_check_results.parquet')
        
        if os.path.exists(cache_file):
            try:
                df = pd.read_parquet(cache_file)
                # Convert DataFrame to dict with bodyId as key
                return dict(zip(df['bodyId'].astype(str), df['has_extrusion']))
            except Exception:
                return {}
        return {}
    
    def _save_extrusion_check_cache(self, results_dict):
        """Save extrusion check results to parquet cache.
        
        Parameters
        ----------
        results_dict : dict
            Dictionary of bodyId -> bool (True if has extrusion)
        """
        import pandas as pd
        
        project_root = os.path.dirname(os.path.dirname(__file__))
        dataset_safe = self.dataset.replace(':', '_').replace('.', '_')
        cache_dir = os.path.join(project_root, 'cache', dataset_safe)
        os.makedirs(cache_dir, exist_ok=True)
        cache_file = os.path.join(cache_dir, 'extrusion_check_results.parquet')
        
        # Load existing cache and merge
        existing = self._load_extrusion_check_cache()
        existing.update({str(k): v for k, v in results_dict.items()})
        
        try:
            # Convert to DataFrame with bodyId index for efficient lookup
            df = pd.DataFrame([
                {'bodyId': str(k), 'has_extrusion': v}
                for k, v in existing.items()
            ])
            df.to_parquet(cache_file, index=False)
        except Exception as e:
            self._vprint(f'  ⚠️ Failed to save extrusion cache: {e}', level='full')

    def _detect_extrusions_in_skeletons(self, skeletons_dict, simplification=0.95):
        """Detect extrusion artifacts in a batch of skeletons.
        
        Converts each skeleton to a simplified mesh and checks for extrusion
        artifacts using edge length analysis. Returns list of body IDs that
        have moderate or severe extrusions and should be replaced via API.
        
        Results are cached in cache/{dataset}/extrusion_check_results.parquet to
        avoid repeated analysis. Previously checked neurons are skipped.
        
        Parameters
        ----------
        skeletons_dict : dict
            Dictionary of bodyId -> TreeNeuron
        simplification : float
            Simplification level to apply for testing (default 0.95)
            
        Returns
        -------
        list
            List of body IDs with extrusion issues that need API replacement
        """
        from tqdm import tqdm
        import sys
        import numpy as np
        
        # Load cached results to skip already-checked neurons
        cached_results = self._load_extrusion_check_cache()
        
        # Filter to only check neurons not in cache
        to_check = {bid: skel for bid, skel in skeletons_dict.items() 
                    if bid not in cached_results and str(bid) not in cached_results}
        
        # Return cached extrusion IDs for neurons we already know about
        extrusion_ids = [bid for bid in skeletons_dict.keys() 
                        if cached_results.get(bid, False) or cached_results.get(str(bid), False)]
        
        if not to_check:
            if extrusion_ids:
                self._vprint(f'  ℹ️  Using cached extrusion results: {len(extrusion_ids)} known issues', level='simple')
            return extrusion_ids
        
        self._vprint(f'  🔍 Checking {len(to_check)} skeletons for extrusions (skipping {len(skeletons_dict) - len(to_check)} cached)...', level='simple')
        
        new_results = {}  # bodyId -> bool (True if has extrusion)
        
        # Progress bar for extrusion checking
        pbar = tqdm(to_check.items(), desc="  Checking extrusions", 
                   disable=not self.verbose, leave=False, file=sys.stdout)
        
        for body_id, skeleton in pbar:
            has_extrusion = False
            try:
                # Get soma position if available
                soma_pos = None
                if hasattr(skeleton, 'soma') and skeleton.soma is not None:
                    soma_idx = skeleton.soma
                    if isinstance(soma_idx, (list, np.ndarray)) and len(soma_idx) > 0:
                        soma_idx = soma_idx[0]
                    if soma_idx is not None and hasattr(skeleton, 'nodes'):
                        soma_node = skeleton.nodes[skeleton.nodes['node_id'] == soma_idx]
                        if len(soma_node) > 0:
                            soma_pos = soma_node[['x', 'y', 'z']].values[0]
                
                # Convert to mesh with simplification
                if hasattr(skeleton, 'nodes') and 'radius' in skeleton.nodes.columns:
                    invalid_mask = (skeleton.nodes['radius'] <= 0) | (skeleton.nodes['radius'].isna())
                    if invalid_mask.any():
                        skeleton.nodes.loc[invalid_mask, 'radius'] = 1
                elif hasattr(skeleton, 'nodes'):
                    skeleton.nodes['radius'] = 1
                
                mesh_neuron = navis.conversion.tree2meshneuron(
                    skeleton,
                    tube_points=8,
                    use_normals=True
                )
                
                if not hasattr(mesh_neuron, 'trimesh'):
                    new_results[body_id] = False
                    continue
                    
                # Apply simplification
                original_faces = len(mesh_neuron.trimesh.faces)
                target_faces = max(100, int(original_faces * (1 - simplification)))
                
                try:
                    import open3d as o3d
                    o3d_mesh = o3d.geometry.TriangleMesh()
                    o3d_mesh.vertices = o3d.utility.Vector3dVector(mesh_neuron.trimesh.vertices)
                    o3d_mesh.triangles = o3d.utility.Vector3iVector(mesh_neuron.trimesh.faces)
                    simplified = o3d_mesh.simplify_quadric_decimation(target_number_of_triangles=target_faces)
                    
                    import trimesh
                    simplified_trimesh = trimesh.Trimesh(
                        vertices=np.asarray(simplified.vertices),
                        faces=np.asarray(simplified.triangles)
                    )
                except ImportError:
                    simplified_trimesh = mesh_neuron.trimesh.simplify_quadric_decimation(target_faces)
                
                # Analyze for extrusions using edge length analysis
                vertices = simplified_trimesh.vertices
                faces = simplified_trimesh.faces
                
                edges = set()
                for face in faces:
                    for j in range(3):
                        v1, v2 = face[j], face[(j + 1) % 3]
                        edges.add((min(v1, v2), max(v1, v2)))
                
                edge_lengths = [np.linalg.norm(vertices[e[0]] - vertices[e[1]]) for e in edges]
                edge_lengths = np.array(edge_lengths)
                
                if len(edge_lengths) == 0:
                    new_results[body_id] = False
                    continue
                
                median_edge = np.median(edge_lengths)
                max_edge = np.max(edge_lengths)
                edge_ratio = max_edge / median_edge if median_edge > 0 else 0
                
                # Check if edge ratio indicates extrusions (threshold=10)
                # Also check if max edge is unreasonably long (>50000nm indicates extrusion spike)
                if edge_ratio > 10 or max_edge > 50000:
                    has_extrusion = True
                    extrusion_ids.append(body_id)
                    
            except Exception:
                pass  # Skip errors silently
            
            new_results[body_id] = has_extrusion
        
        # Save new results to cache
        if new_results:
            self._save_extrusion_check_cache(new_results)
            n_new_extrusions = sum(1 for v in new_results.values() if v)
            self._vprint(f'  💾 Cached extrusion results for {len(new_results)} neurons ({n_new_extrusions} with issues)', level='full')
        
        return extrusion_ids
    
    def _skeleton_cache_is_simplified(self) -> bool:
        """True when the skeleton cache holds the fixed simplified level.

        Reads the ``skeletons/.level`` marker written by the cache pipeline
        (morphology.fetch_skeleton_on_demand and ``_save_cached_neurons``);
        a missing marker means the cache predates the marker (raw).
        """
        marker = os.path.join(self._get_cache_path('skeletons'), '.level')
        try:
            with open(marker) as f:
                return f.read().strip() == 'simp90'
        except Exception:
            return False

    def _effective_render_simplification(self, is_fafb: bool) -> float:
        """Effective render-time mesh decimation fraction for tube mode.

        For NeuPrint datasets whose skeleton cache already holds the fixed
        simp90 level, the cache-level reduction is baked into the tube mesh.
        At exactly that level no further decimation is applied (decimating
        again would double-reduce: 10% nodes -> ~1% faces). Levels ABOVE
        the cache level apply the *remaining* relative reduction (same
        semantics as the FAFB path), so e.g. 0.95 still halves the face
        count instead of silently doing nothing.

        Returns the fraction of faces to remove (0.0 = keep all).
        """
        target = self.skeleton_mesh_simplification
        if (not is_fafb and self._skeleton_cache_is_simplified()
                and target >= self.NEUPRINT_SKELETON_CACHE_LEVEL):
            remaining_after_cache = 1 - self.NEUPRINT_SKELETON_CACHE_LEVEL
            remaining_target = 1 - target
            keep_factor = remaining_target / remaining_after_cache
            return max(0.0, 1.0 - keep_factor)
        return target

    def _load_cached_neurons(self, neuron_df, transformed_target=None,
                             ignore_cache=False):
        """Load cached neuron skeletons if available.

        Loads individual {bodyId}.pkl files from cache/{dataset}/skeletons/
        (NeuPrint caches hold the fixed 90%-simplified skeletons). With
        ``ignore_cache=True`` (less-simplified render, NeuPrint only) every
        neuron is reported missing so the caller re-fetches RAW skeletons
        transiently; the simplified cache is still written from that fetch.

        Returns:
            tuple: (navis.NeuronList or None, list of missing bodyIds)
        """
        if not self.cache_neurons:
            return None, neuron_df['bodyId'].tolist()
        
        cache_dir = self._get_cache_path('skeletons')
        body_ids = neuron_df['bodyId'].tolist()
        
        import pickle
        neurons = []
        loaded_ids = []
        missing_ids = []
        
        for bid in body_ids:
            if ignore_cache:
                missing_ids.append(bid)
                continue
            cache_file = os.path.join(cache_dir, f'{bid}.pkl')
            if os.path.exists(cache_file):
                try:
                    with open(cache_file, 'rb') as f:
                        neuron = pickle.load(f)
                    neurons.append(neuron)
                    loaded_ids.append(bid)
                except Exception as e:
                    self._vprint(f'  ⚠ Failed to load cached skeleton {bid}: {e}')
                    missing_ids.append(bid)
            else:
                missing_ids.append(bid)
        
        if neurons:
            self._vprint(f'  ✓ Loaded {len(neurons)} neurons from cache', level='full')
            if missing_ids:
                self._vprint(f'  ℹ  {len(missing_ids)} neurons not in cache, will fetch', level='full')
            # Return loaded neurons plus info about missing ones
            return navis.NeuronList(neurons), missing_ids
        
        return None, body_ids  # All missing
    
    def _save_cached_neurons(self, neuron_df, neuron_vols):
        """Save neuron skeletons to cache as individual {bodyId}.pkl files.

        NeuPrint skeletons are cached ONLY at the fixed 90%-simplified level
        (``navis.downsample_neuron`` factor 10): the raw skeleton is never
        persisted, and the folder's ``.level`` marker is written so the
        morphology pipeline's level guards stay consistent.

        Saves each neuron as a separate file for better reusability.
        """
        if not self.cache_neurons:
            return
        
        is_neuprint = not (self.client_type == 'flywire'
                           or 'flywire' in self.dataset.lower()
                           or 'fafb' in self.dataset.lower())
        cache_dir = self._get_cache_path('skeletons')
        
        import pickle
        saved_count = 0
        
        # Handle both NeuronList and list of neurons
        if hasattr(neuron_vols, '__iter__'):
            for neuron in neuron_vols:
                try:
                    # Get bodyId from neuron
                    bid = getattr(neuron, 'id', None) or getattr(neuron, 'bodyId', None)
                    if bid is None:
                        continue
                    
                    cache_file = os.path.join(cache_dir, f'{bid}.pkl')
                    
                    # Skip if already cached
                    if os.path.exists(cache_file):
                        continue
                    
                    # NeuPrint: persist ONLY the simplified skeleton (raw is
                    # never cached); the morphology pipeline vectorizes the
                    # raw skeleton on its next fetch instead.
                    to_store = neuron
                    if is_neuprint and hasattr(neuron, 'nodes'):
                        try:
                            # A multi-node "soma" (navis' radius>=1 detection
                            # on nm radii) would freeze the skeleton at full
                            # resolution during downsampling.
                            soma = getattr(neuron, 'soma', None)
                            if soma is not None and hasattr(soma, '__len__') and len(soma) > 1:
                                neuron = neuron.copy()
                                neuron.soma = None
                            to_store = navis.downsample_neuron(
                                neuron, downsampling_factor=self.NEUPRINT_SKELETON_DOWNSAMPLE
                            )
                        except Exception:
                            to_store = neuron
                    
                    with open(cache_file, 'wb') as f:
                        pickle.dump(to_store, f)
                    saved_count += 1
                except Exception as e:
                    self._vprint(f'  ⚠ Failed to save skeleton {bid}: {e}')
        
        if saved_count > 0:
            if is_neuprint:
                # Mark the folder's simplification level (idempotent);
                # only the NeuPrint path writes simplified skeletons.
                try:
                    marker = os.path.join(cache_dir, '.level')
                    if not os.path.exists(marker):
                        os.makedirs(cache_dir, exist_ok=True)
                        with open(marker, 'w') as f:
                            f.write('simp90\n')
                except Exception:
                    pass
            self._vprint(f'  💾 Saved {saved_count} new neurons to cache', level='full')
    
    # Cache stores meshes simplified with soma-aware parameters
    # Skeleton: 0.95 simplification (keep 5% of faces) 
    # Soma: 0.8 simplification (keep 20% of faces) to prevent extrusion artifacts
    FAFB_MESH_CACHE_SIMPLIFICATION = 0.95  # Skeleton simplification level
    FAFB_MESH_CACHE_SOMA_SIMPLIFICATION = 0.8  # Gentler simplification for soma region

    # NeuPrint skeletons are cached ONLY at the fixed 90%-simplified level
    # (``navis.downsample_neuron`` factor 10, keeps ~10% of nodes): raw
    # skeletons are never persisted. Rendering at >= this level skips the
    # render-time decimation (the tube mesh is already at the cache level);
    # rendering below it transiently re-fetches RAW skeletons instead.
    NEUPRINT_SKELETON_CACHE_LEVEL = 0.9
    NEUPRINT_SKELETON_DOWNSAMPLE = 10
    FAFB_MESH_CACHE_SOMA_RADIUS = 20000  # 20µm radius around soma for gentler simplification
    
    def _get_fafb_mesh_cache_key(self):
        """Generate a cache key based on coordinate space and simplification settings.
        
        Returns a subfolder name like 'FLYWIRE_simp95_soma80_r20' for caching purposes.
        Cache stores meshes with soma-aware simplification (gentler on cell body).
        Note: FAFB cache stores UN-TRANSFORMED, UN-ROTATED meshes (native FLYWIRE).
        """
        # Always use 'FLYWIRE' as base since we cache raw simplified meshes
        target = 'FLYWIRE'
        
        # Include simplification levels and soma radius in cache key
        skel_simp = int(self.FAFB_MESH_CACHE_SIMPLIFICATION * 100)
        soma_simp = int(self.FAFB_MESH_CACHE_SOMA_SIMPLIFICATION * 100)
        soma_r = int(self.FAFB_MESH_CACHE_SOMA_RADIUS / 1000)  # In µm for readability
        return f"{target}_simp{skel_simp}_soma{soma_simp}_r{soma_r}"
    
    def _load_cached_fafb_meshes(self, body_ids):
        """Load simplified and meshed FAFB neurons from cache.
        
        Cache contains meshes with soma-aware simplification:
        - Skeleton: 0.95 simplification (keep 5% of faces)
        - Soma region (within 20µm): 0.8 simplification (keep 20% of faces)
        
        No coordinate transformation is needed since FAFB uses native FLYWIRE coordinates.
        Only used when skeleton_mesh_simplification >= 0.95.
        If simplification > 0.95, additional simplification is applied after loading.
        
        Parameters
        ----------
        body_ids : list
            List of bodyIds to load
            
        Returns
        -------
        tuple: (dict of bodyId -> MeshNeuron, list of missing bodyIds)
        """
        if not self.cache_neurons:
            return {}, body_ids
        
        # Only use cache when simplification >= cache level (0.95)
        if self.skeleton_mesh_simplification < self.FAFB_MESH_CACHE_SIMPLIFICATION:
            return {}, body_ids
        
        # Check for flywire/fafb dataset
        if not ('flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()):
            return {}, body_ids
        
        import pickle
        
        cache_key = self._get_fafb_mesh_cache_key()
        # Store in skeletons folder as individual simplified meshes
        cache_dir = os.path.join(self._get_cache_path('skeletons'), cache_key)
        os.makedirs(cache_dir, exist_ok=True)
        
        loaded = {}
        missing = []
        
        for bid in body_ids:
            cache_file = os.path.join(cache_dir, f'{bid}.pkl')
            if os.path.exists(cache_file):
                try:
                    with open(cache_file, 'rb') as f:
                        mesh_neuron = pickle.load(f)
                    loaded[bid] = mesh_neuron
                except Exception as e:
                    self._vprint(f'  ⚠ Failed to load cached mesh {bid}: {e}', level='full')
                    missing.append(bid)
            else:
                missing.append(bid)
        
        if loaded:
            self._vprint(f'  ✓ Loaded {len(loaded)} neurons from mesh cache (skel={self.FAFB_MESH_CACHE_SIMPLIFICATION}, soma={self.FAFB_MESH_CACHE_SOMA_SIMPLIFICATION})', level='full')
        
        return loaded, missing
    
    def _save_cached_fafb_meshes(self, mesh_neurons_dict):
        """Save transformed and meshed FAFB neurons to cache.
        
        Parameters
        ----------
        mesh_neurons_dict : dict
            Dictionary of bodyId -> MeshNeuron to save
        """
        if not self.cache_neurons:
            return
        
        # Check for flywire/fafb dataset
        if not ('flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()):
            return
        
        import pickle
        
        cache_key = self._get_fafb_mesh_cache_key()
        # Store in skeletons folder as individual simplified meshes
        cache_dir = os.path.join(self._get_cache_path('skeletons'), cache_key)
        os.makedirs(cache_dir, exist_ok=True)
        
        saved_count = 0
        for bid, mesh_neuron in mesh_neurons_dict.items():
            cache_file = os.path.join(cache_dir, f'{bid}.pkl')
            if os.path.exists(cache_file):
                continue  # Skip if already cached
            
            try:
                with open(cache_file, 'wb') as f:
                    pickle.dump(mesh_neuron, f)
                saved_count += 1
            except Exception as e:
                self._vprint(f'  ⚠ Failed to save mesh {bid}: {e}', level='full')
        
        if saved_count > 0:
            self._vprint(f'  💾 Saved {saved_count} new meshes to cache', level='full')

    def _load_api_cached_skeletons(self, body_ids: list) -> tuple:
        """Load skeletons from API cache (cache/{dataset}/API_cache/skeletons/).
        
        These are skeletons previously fetched via CAVE API and cached locally.
        API cache takes priority over ZIP files as it contains more up-to-date data.
        
        Parameters
        ----------
        body_ids : list
            List of body IDs to look for in cache
            
        Returns
        -------
        tuple
            (dict of loaded skeletons {bodyId: TreeNeuron}, list of missing bodyIds)
        """
        import pickle
        
        project_root = os.path.dirname(os.path.dirname(__file__))
        dataset_safe = self.dataset.replace(':', '_').replace('.', '_')
        cache_dir = os.path.join(project_root, 'cache', dataset_safe, 'API_cache', 'skeletons')
        
        if not os.path.exists(cache_dir):
            return {}, list(body_ids)
        
        cached = {}
        missing = []
        
        for bid in body_ids:
            # Check for both int and str versions of bodyId
            cache_file = os.path.join(cache_dir, f'{bid}.pkl')
            if os.path.exists(cache_file):
                try:
                    with open(cache_file, 'rb') as f:
                        neuron = pickle.load(f)
                        cached[bid] = neuron
                        cached[str(bid)] = neuron  # Also store with str key for compatibility
                except Exception:
                    missing.append(bid)
            else:
                missing.append(bid)
        
        return cached, missing

    def _fetch_fafb_skeletons_via_api(self, body_ids: list) -> dict:
        """Fetch FAFB skeletons via CAVE API (CloudVolume mesh + navis skeletonization).
        
        Uses CAVEDataFetcher to fetch meshes and generate skeletons.
        Results are cached in cache/{dataset}/API_cache/skeletons/.
        
        Parameters
        ----------
        body_ids : list
            List of body IDs to fetch (can be strings or ints)
            
        Returns
        -------
        dict
            Dictionary of bodyId -> TreeNeuron (keys are both int and str for compatibility)
        """
        try:
            from cave_data_fetcher import CAVEDataFetcher
        except ImportError:
            try:
                from .cave_data_fetcher import CAVEDataFetcher
            except ImportError:
                self._vprint("  ⚠️  cave_data_fetcher module not found, falling back to ZIP", use_tqdm=True)
                return {}
        
        self._vprint(f'  🌐 Fetching {len(body_ids)} skeletons via CAVE API...', use_tqdm=True)
        
        # Convert body_ids to integers for API call
        int_body_ids = [int(bid) for bid in body_ids]
        
        fetcher = CAVEDataFetcher(
            dataset=self.dataset,
            verbose=self.verbose == 'full'
        )
        
        # Fetch skeletons with soma-aware simplification
        # Use the class cache parameters for consistency
        skeleton_simp = self.FAFB_MESH_CACHE_SIMPLIFICATION
        soma_simp = self.FAFB_MESH_CACHE_SOMA_SIMPLIFICATION
        soma_radius = self.FAFB_MESH_CACHE_SOMA_RADIUS
        
        skeleton_cache = {}
        # Fetch with skeleton simplification level (soma-aware applied after)
        neurons = fetcher.fetch_skeletons(int_body_ids, use_cache=True, simplify_mesh=skeleton_simp)
        
        for n in neurons:
            if hasattr(n, 'id'):
                # Apply soma-aware simplification if soma position is available
                if hasattr(n, 'trimesh') and hasattr(n, 'soma_pos') and n.soma_pos is not None:
                    try:
                        simplified_trimesh = self._simplify_mesh_with_soma_awareness(
                            n.trimesh,
                            skeleton_simp=skeleton_simp,
                            soma_simp=soma_simp,
                            soma_pos=n.soma_pos,
                            soma_radius=soma_radius
                        )
                        n._trimesh = simplified_trimesh
                    except Exception:
                        pass  # Keep original if soma-aware simplification fails
                
                # Store with both int and str keys for compatibility
                skeleton_cache[n.id] = n
                skeleton_cache[str(n.id)] = n
        
        self._vprint(f'  ✓ Fetched {len(neurons)}/{len(body_ids)} skeletons via API (soma-aware simplification)', use_tqdm=True)
        return skeleton_cache

    @staticmethod
    def detect_mesh_extrusions(mesh, soma_pos=None, soma_radius=20000, 
                               distance_threshold_factor=3.0, verbose=False) -> dict:
        """Detect extrusion artifacts in a mesh (spikes/protrusions from simplification).
        
        Extrusions typically manifest as:
        1. Vertices with very long edge connections (spiky protrusions)
        2. Isolated vertex clusters far from the main mesh body
        3. Non-manifold geometry issues
        
        Parameters
        ----------
        mesh : trimesh.Trimesh or navis.MeshNeuron
            The mesh to analyze for extrusions.
        soma_pos : array-like, optional
            [x, y, z] position of the soma center. If provided, analysis focuses
            on the soma region where extrusions are most common.
        soma_radius : float
            Radius (in nm) around soma_pos to focus analysis (default 20000nm = 20µm).
        distance_threshold_factor : float
            Vertices with average edge length > (median_edge_length * this_factor) 
            are flagged as potential extrusions (default 3.0).
        verbose : bool
            Whether to print detailed analysis info.
            
        Returns
        -------
        dict
            Analysis results with keys:
            - 'has_extrusions': bool - Whether extrusions were detected
            - 'severity': str - 'none', 'mild', 'moderate', 'severe'
            - 'extrusion_count': int - Number of extrusion vertices detected
            - 'extrusion_vertices': array - Indices of extrusion vertices
            - 'max_edge_length': float - Maximum edge length in mesh
            - 'median_edge_length': float - Median edge length
            - 'edge_length_ratio': float - Ratio of max to median edge length
            - 'soma_region_issues': bool - Whether issues are near soma
            - 'recommendation': str - Suggested action
            
        Example
        -------
        >>> from coana import VisualizeSkeleton
        >>> import trimesh
        >>> # Analyze a mesh
        >>> result = VisualizeSkeleton.detect_mesh_extrusions(mesh, soma_pos=[100, 200, 300])
        >>> if result['has_extrusions']:
        ...     print(f"Extrusions detected! Severity: {result['severity']}")
        ...     print(result['recommendation'])
        """
        import numpy as np
        
        # Extract trimesh object if MeshNeuron
        if hasattr(mesh, 'trimesh'):
            trimesh_obj = mesh.trimesh
        elif hasattr(mesh, 'vertices') and hasattr(mesh, 'faces'):
            trimesh_obj = mesh
        else:
            return {
                'has_extrusions': False,
                'severity': 'unknown',
                'extrusion_count': 0,
                'extrusion_vertices': np.array([]),
                'max_edge_length': 0,
                'median_edge_length': 0,
                'edge_length_ratio': 0,
                'soma_region_issues': False,
                'recommendation': 'Unable to analyze - invalid mesh format'
            }
        
        vertices = trimesh_obj.vertices
        faces = trimesh_obj.faces
        
        # Calculate edge lengths
        edges = []
        for face in faces:
            for i in range(3):
                v1, v2 = face[i], face[(i + 1) % 3]
                edges.append((min(v1, v2), max(v1, v2)))
        
        edges = list(set(edges))  # Unique edges
        edge_lengths = []
        edge_to_length = {}
        for e in edges:
            length = np.linalg.norm(vertices[e[0]] - vertices[e[1]])
            edge_lengths.append(length)
            edge_to_length[e] = length
        
        edge_lengths = np.array(edge_lengths)
        median_edge = np.median(edge_lengths)
        max_edge = np.max(edge_lengths)
        edge_ratio = max_edge / median_edge if median_edge > 0 else 0
        
        # Calculate per-vertex average edge length
        vertex_edge_lengths = {i: [] for i in range(len(vertices))}
        for e, length in edge_to_length.items():
            vertex_edge_lengths[e[0]].append(length)
            vertex_edge_lengths[e[1]].append(length)
        
        vertex_avg_edge = np.array([
            np.mean(vertex_edge_lengths[i]) if vertex_edge_lengths[i] else 0 
            for i in range(len(vertices))
        ])
        
        # Detect extrusion vertices (abnormally long edges)
        threshold = median_edge * distance_threshold_factor
        extrusion_mask = vertex_avg_edge > threshold
        extrusion_vertices = np.where(extrusion_mask)[0]
        
        # Check if extrusions are near soma region
        soma_region_issues = False
        if soma_pos is not None and len(extrusion_vertices) > 0:
            soma_pos = np.array(soma_pos).flatten()[:3]
            extrusion_positions = vertices[extrusion_vertices]
            distances_to_soma = np.linalg.norm(extrusion_positions - soma_pos, axis=1)
            soma_region_issues = np.any(distances_to_soma <= soma_radius)
        
        # Determine severity
        n_extrusions = len(extrusion_vertices)
        total_vertices = len(vertices)
        extrusion_ratio = n_extrusions / total_vertices if total_vertices > 0 else 0
        
        if n_extrusions == 0:
            severity = 'none'
        elif extrusion_ratio < 0.001 or n_extrusions < 5:
            severity = 'mild'
        elif extrusion_ratio < 0.01 or n_extrusions < 50:
            severity = 'moderate'
        else:
            severity = 'severe'
        
        has_extrusions = n_extrusions > 0 and (edge_ratio > distance_threshold_factor)
        
        # Generate recommendation
        if not has_extrusions:
            recommendation = 'No extrusions detected. Mesh appears clean.'
        elif severity == 'mild':
            recommendation = (
                'Mild extrusions detected. Consider using soma-aware simplification '
                '(soma_mesh_simplification=0.8) or fetch fresh skeleton via CAVE API.'
            )
        elif severity == 'moderate':
            if soma_region_issues:
                recommendation = (
                    'Moderate extrusions near soma region. Strongly recommend using '
                    'VisualizeSkeleton.fix_fafb_extrusions([bodyId]) to fetch fresh '
                    'skeleton from CAVE API.'
                )
            else:
                recommendation = (
                    'Moderate extrusions detected. Use soma-aware simplification or '
                    'VisualizeSkeleton.fix_fafb_extrusions([bodyId]) to fix.'
                )
        else:  # severe
            recommendation = (
                'Severe extrusions detected! This skeleton should be replaced. Use '
                'VisualizeSkeleton.fix_fafb_extrusions([bodyId]) to fetch fresh '
                'skeleton from CAVE API.'
            )
        
        if verbose:
            print(f"Mesh Analysis Results:")
            print(f"  Total vertices: {total_vertices:,}")
            print(f"  Total faces: {len(faces):,}")
            print(f"  Median edge length: {median_edge:.1f}nm")
            print(f"  Max edge length: {max_edge:.1f}nm")
            print(f"  Edge ratio (max/median): {edge_ratio:.1f}x")
            print(f"  Extrusion vertices: {n_extrusions} ({extrusion_ratio*100:.2f}%)")
            print(f"  Severity: {severity}")
            print(f"  Near soma: {soma_region_issues}")
            print(f"  Recommendation: {recommendation}")
        
        return {
            'has_extrusions': has_extrusions,
            'severity': severity,
            'extrusion_count': n_extrusions,
            'extrusion_vertices': extrusion_vertices,
            'max_edge_length': max_edge,
            'median_edge_length': median_edge,
            'edge_length_ratio': edge_ratio,
            'soma_region_issues': soma_region_issues,
            'recommendation': recommendation
        }

    @staticmethod
    def fix_fafb_extrusions(body_ids: list, dataset: str = 'flywire_FAFB_v783', 
                            verbose: bool = True) -> dict:
        """Fix FAFB skeleton extrusion issues by fetching fresh data from CAVE API.
        
        The downloaded FAFB skeleton ZIP (sk_lod1_783_healed.zip) may contain neurons
        with extrusion artifacts (mesh errors appearing as spikes/protrusions). This
        method fetches fresh skeletons from CAVE API and caches them locally.
        
        Once cached, VisualizeSkeleton will automatically use the API-cached versions
        instead of the problematic ZIP versions.
        
        Parameters
        ----------
        body_ids : list
            List of body IDs (as integers or strings) with extrusion issues
        dataset : str
            Dataset name, default 'flywire_FAFB_v783'
        verbose : bool
            Whether to print progress messages
            
        Returns
        -------
        dict
            Dictionary of fixed skeletons {bodyId: TreeNeuron}
            
        Example
        -------
        >>> # Fix specific neurons with extrusion issues
        >>> from coana import VisualizeSkeleton
        >>> fixed = VisualizeSkeleton.fix_fafb_extrusions([720575940596125868, 720575940597856265])
        >>> print(f"Fixed {len(fixed)} neurons")
        
        >>> # Now use them in visualization (automatically uses fixed versions)
        >>> vs = VisualizeSkeleton(
        ...     dataset='flywire_FAFB_v783',
        ...     neuron_layers=['l-LNv'],
        ...     force_API_fetching=False,  # API cache still takes priority
        ... )
        >>> vs.plot_neurons()
        """
        try:
            from cave_data_fetcher import CAVEDataFetcher
        except ImportError:
            try:
                from .cave_data_fetcher import CAVEDataFetcher
            except ImportError:
                raise ImportError(
                    "CAVE data fetcher not available. Install with: "
                    "pip install caveclient cloud-volume"
                )
        
        if verbose:
            print(f"🔧 Fixing FAFB skeleton extrusions for {len(body_ids)} neurons...")
        
        # Convert to integers
        int_body_ids = [int(bid) for bid in body_ids]
        
        fetcher = CAVEDataFetcher(
            dataset=dataset,
            verbose=verbose
        )
        
        # Fetch and cache fresh skeletons (0.95 simplification)
        neurons = fetcher.fetch_skeletons(int_body_ids, use_cache=False, simplify_mesh=0.95)
        
        result = {}
        for n in neurons:
            if hasattr(n, 'id'):
                result[n.id] = n
                result[str(n.id)] = n
        
        if verbose:
            print(f"✓ Fixed and cached {len(neurons)}/{len(body_ids)} skeletons")
            print(f"  Cache location: cache/{dataset.replace(':', '_').replace('.', '_')}/API_cache/skeletons/")
            print(f"  These will automatically be used instead of ZIP versions")
        
        return result

    @staticmethod
    def check_fafb_skeleton_for_extrusions(body_id, dataset: str = 'flywire_FAFB_v783',
                                           simplification: float = 0.95,
                                           verbose: bool = True,
                                           auto_fix: bool = False) -> dict:
        """Load a FAFB skeleton from ZIP and check for extrusion artifacts.
        
        This is a convenience method that combines loading a skeleton from the local
        ZIP file, applying simplification, and running extrusion detection.
        
        Parameters
        ----------
        body_id : int or str
            The body ID to check.
        dataset : str
            Dataset name, default 'flywire_FAFB_v783'.
        simplification : float
            Simplification level to apply before checking (default 0.95).
        verbose : bool
            Whether to print detailed analysis results.
        auto_fix : bool
            If True and extrusions are detected, automatically fetch fresh
            skeleton from CAVE API (default False).
            
        Returns
        -------
        dict
            Analysis results with additional keys:
            - All keys from detect_mesh_extrusions()
            - 'body_id': The body ID that was checked
            - 'auto_fixed': bool - Whether auto-fix was applied
            - 'skeleton': The loaded skeleton (TreeNeuron or MeshNeuron)
            
        Example
        -------
        >>> from coana import VisualizeSkeleton
        >>> # Check a specific neuron
        >>> result = VisualizeSkeleton.check_fafb_skeleton_for_extrusions(
        ...     720575940624086675, 
        ...     simplification=0.95,
        ...     verbose=True
        ... )
        >>> if result['has_extrusions']:
        ...     # Auto-fix it
        ...     VisualizeSkeleton.fix_fafb_extrusions([720575940624086675])
        
        >>> # Or use auto_fix=True
        >>> result = VisualizeSkeleton.check_fafb_skeleton_for_extrusions(
        ...     720575940624086675,
        ...     auto_fix=True  # Automatically fetch from API if extrusions found
        ... )
        """
        import navis
        from pathlib import Path
        import zipfile
        import re
        import numpy as np
        
        body_id_str = str(body_id)
        body_id_int = int(body_id)
        
        if verbose:
            print(f"🔍 Checking FAFB skeleton {body_id} for extrusions...")
        
        # Find the skeleton ZIP file
        dataset_clean = dataset.replace(':', '_').replace('.', '_')
        zip_path = Path('datasets') / dataset_clean / 'sk_lod1_783_healed.zip'
        
        if not zip_path.exists():
            return {
                'has_extrusions': False,
                'severity': 'unknown',
                'extrusion_count': 0,
                'extrusion_vertices': np.array([]),
                'max_edge_length': 0,
                'median_edge_length': 0,
                'edge_length_ratio': 0,
                'soma_region_issues': False,
                'recommendation': f'ZIP file not found: {zip_path}',
                'body_id': body_id_int,
                'auto_fixed': False,
                'skeleton': None
            }
        
        # Find matching file in ZIP
        skeleton = None
        soma_pos = None
        
        try:
            with zipfile.ZipFile(zip_path, 'r') as zf:
                # Find file matching this body ID
                pattern = re.compile(rf'^{body_id_str}\..*\.swc$', re.IGNORECASE)
                matching_files = [f for f in zf.namelist() if pattern.match(f)]
                
                if not matching_files:
                    # Try without extension pattern
                    matching_files = [f for f in zf.namelist() if f.startswith(f'{body_id_str}.')]
                
                if not matching_files:
                    return {
                        'has_extrusions': False,
                        'severity': 'unknown',
                        'extrusion_count': 0,
                        'extrusion_vertices': np.array([]),
                        'max_edge_length': 0,
                        'median_edge_length': 0,
                        'edge_length_ratio': 0,
                        'soma_region_issues': False,
                        'recommendation': f'Body ID {body_id} not found in ZIP',
                        'body_id': body_id_int,
                        'auto_fixed': False,
                        'skeleton': None
                    }
                
                swc_filename = matching_files[0]
                
                # Read SWC content
                with zf.open(swc_filename) as swc_file:
                    swc_content = swc_file.read().decode('utf-8')
                
                # Parse with navis
                from io import StringIO
                skeleton = navis.read_swc(StringIO(swc_content))
                
                if isinstance(skeleton, navis.NeuronList):
                    skeleton = skeleton[0]
                
                skeleton.id = body_id_int
                
                # Get soma position if available
                if hasattr(skeleton, 'soma') and skeleton.soma is not None:
                    soma_idx = skeleton.soma
                    if isinstance(soma_idx, (list, np.ndarray)) and len(soma_idx) > 0:
                        soma_idx = soma_idx[0]
                    if soma_idx is not None:
                        soma_node = skeleton.nodes[skeleton.nodes['node_id'] == soma_idx]
                        if len(soma_node) > 0:
                            soma_pos = soma_node[['x', 'y', 'z']].values[0]
                
                if verbose:
                    print(f"  ✓ Loaded skeleton with {len(skeleton.nodes):,} nodes")
                    if soma_pos is not None:
                        print(f"  ✓ Soma position: [{soma_pos[0]:.0f}, {soma_pos[1]:.0f}, {soma_pos[2]:.0f}]")
        
        except Exception as e:
            return {
                'has_extrusions': False,
                'severity': 'unknown',
                'extrusion_count': 0,
                'extrusion_vertices': np.array([]),
                'max_edge_length': 0,
                'median_edge_length': 0,
                'edge_length_ratio': 0,
                'soma_region_issues': False,
                'recommendation': f'Error loading skeleton: {e}',
                'body_id': body_id_int,
                'auto_fixed': False,
                'skeleton': None
            }
        
        # Convert to mesh with simplification
        if verbose:
            print(f"  ⚙️  Converting to mesh (simplification={simplification})...")
        
        try:
            mesh_neuron = navis.conversion.tree2meshneuron(
                skeleton,
                tube_points=8,
                use_normals=True
            )
            
            if hasattr(mesh_neuron, 'trimesh') and simplification > 0:
                import trimesh
                original_faces = len(mesh_neuron.trimesh.faces)
                target_faces = max(100, int(original_faces * (1 - simplification)))
                
                # Simplify using quadric decimation
                try:
                    import open3d as o3d
                    o3d_mesh = o3d.geometry.TriangleMesh()
                    o3d_mesh.vertices = o3d.utility.Vector3dVector(mesh_neuron.trimesh.vertices)
                    o3d_mesh.triangles = o3d.utility.Vector3iVector(mesh_neuron.trimesh.faces)
                    simplified = o3d_mesh.simplify_quadric_decimation(target_number_of_triangles=target_faces)
                    simplified_trimesh = trimesh.Trimesh(
                        vertices=np.asarray(simplified.vertices),
                        faces=np.asarray(simplified.triangles)
                    )
                    mesh_neuron._trimesh = simplified_trimesh
                except ImportError:
                    # Fallback to trimesh simplification
                    mesh_neuron._trimesh = mesh_neuron.trimesh.simplify_quadric_decimation(target_faces)
                
                if verbose:
                    print(f"  ✓ Simplified: {original_faces:,} → {len(mesh_neuron.trimesh.faces):,} faces")
        
        except Exception as e:
            return {
                'has_extrusions': False,
                'severity': 'unknown',
                'extrusion_count': 0,
                'extrusion_vertices': np.array([]),
                'max_edge_length': 0,
                'median_edge_length': 0,
                'edge_length_ratio': 0,
                'soma_region_issues': False,
                'recommendation': f'Error converting to mesh: {e}',
                'body_id': body_id_int,
                'auto_fixed': False,
                'skeleton': skeleton
            }
        
        # Run extrusion detection
        result = VisualizeSkeleton.detect_mesh_extrusions(
            mesh_neuron,
            soma_pos=soma_pos,
            soma_radius=20000,
            verbose=verbose
        )
        
        result['body_id'] = body_id_int
        result['auto_fixed'] = False
        result['skeleton'] = mesh_neuron
        
        # Auto-fix if requested
        if auto_fix and result['has_extrusions']:
            if verbose:
                print(f"\n  🔧 Auto-fixing: Fetching fresh skeleton from CAVE API...")
            try:
                fixed = VisualizeSkeleton.fix_fafb_extrusions([body_id_int], dataset=dataset, verbose=verbose)
                if body_id_int in fixed or body_id_str in fixed:
                    result['auto_fixed'] = True
                    result['skeleton'] = fixed.get(body_id_int) or fixed.get(body_id_str)
                    result['recommendation'] = 'Extrusions fixed! Fresh skeleton fetched and cached from CAVE API.'
            except Exception as e:
                if verbose:
                    print(f"  ⚠️  Auto-fix failed: {e}")
                result['recommendation'] += f' Auto-fix failed: {e}'
        
        return result

    def plot_skeleton(self):
        from tqdm import tqdm
        import sys
        
        n_layers = len(self.neuron_layers)
        total_skeletons = sum(len(df) if df is not None else 0 for df in self.neuron_dfs)
        self._vprint(f'\n🔬 Fetching skeletons for {n_layers} layers ({total_skeletons:,} neurons total)...')
        
        # For FAFB: Check mesh cache first (transformed + meshed neurons)
        # Cache stores pre-simplified meshes at FAFB_MESH_CACHE_SIMPLIFICATION (0.9 = keep 10% faces)
        # 
        # Cache usage decision:
        # - If user wants simplification >= 0.9 (keep ≤10% faces): use cache, apply additional simplification if needed
        # - If user wants simplification < 0.9 (keep >10% faces): bypass cache, load from ZIP and apply user's simplification
        #
        # Example scenarios:
        # - simplification=0.95 (keep 5%): load from cache (10%), simplify to 5% → additional_keep = 0.05/0.1 = 50%
        # - simplification=0.9 (keep 10%): load from cache (10%), no additional simplification needed
        # - simplification=0.5 (keep 50%): cannot use cache (only has 10%), load from ZIP and apply 0.5 simplification
        is_fafb = 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()
        fafb_mesh_cache = {}  # bodyId -> MeshNeuron (from cache)
        fafb_mesh_missing = []  # bodyIds that need processing
        use_fafb_cache = is_fafb and self.cache_neurons and self.skeleton_mesh_simplification >= self.FAFB_MESH_CACHE_SIMPLIFICATION
        
        # Check for force_API_fetching - bypasses ZIP loading for FAFB
        use_api_fetching = is_fafb and self.force_API_fetching
        
        if is_fafb:
            if use_api_fetching:
                self._vprint(f'  ℹ️  force_API_fetching=True: Using CAVE API instead of ZIP', level='simple')
            elif use_fafb_cache:
                self._vprint(f'  ℹ️  FAFB mesh cache enabled (simplification={self.skeleton_mesh_simplification} >= cache level {self.FAFB_MESH_CACHE_SIMPLIFICATION})', level='full')
            else:
                self._vprint(f'  ℹ️  FAFB mesh cache bypassed (simplification={self.skeleton_mesh_simplification} < cache level {self.FAFB_MESH_CACHE_SIMPLIFICATION})', level='full')
        
        # Load from mesh cache when simplification >= cache level
        # This applies regardless of force_API_fetching, since mesh cache stores processed meshes
        if use_fafb_cache:
            # Collect all body IDs across layers
            all_fafb_body_ids = []
            for df in self.neuron_dfs:
                if df is not None and 'bodyId' in df.columns:
                    all_fafb_body_ids.extend(df['bodyId'].tolist())
            all_fafb_body_ids = list(set(all_fafb_body_ids))
            
            # Load from mesh cache
            fafb_mesh_cache, fafb_mesh_missing = self._load_cached_fafb_meshes(all_fafb_body_ids)
        
        # Pre-load all FAFB skeletons
        fafb_skeleton_cache = {}  # bodyId -> TreeNeuron
        if is_fafb:
            # Collect all body IDs first
            all_fafb_body_ids = []
            for df in self.neuron_dfs:
                if df is not None and 'bodyId' in df.columns:
                    all_fafb_body_ids.extend(df['bodyId'].tolist())
            all_fafb_body_ids = list(set(all_fafb_body_ids))
            
            # ALWAYS check API cache first for VisualizeSkeleton
            # This allows fixing individual neurons with extrusion issues by fetching via API
            # API-cached skeletons take priority over ZIP data for better morphology
            api_cached_skeletons, api_cache_missing = self._load_api_cached_skeletons(all_fafb_body_ids)
            
            if api_cached_skeletons:
                fafb_skeleton_cache.update(api_cached_skeletons)
                self._vprint(f'  ✓ Loaded {len(api_cached_skeletons)} skeletons from API cache (extrusion-fixed)', level='simple')
            
            if use_api_fetching:
                # force_API_fetching=True: Fetch remaining via CAVE API
                if api_cache_missing:
                    self._vprint(f'  🌐 Fetching {len(api_cache_missing)} skeletons via CAVE API...')
                    api_fetched = self._fetch_fafb_skeletons_via_api(api_cache_missing)
                    fafb_skeleton_cache.update(api_fetched)
            else:
                # force_API_fetching=False: Load remaining from local ZIP
                remaining_ids = set(api_cache_missing)
                
                if remaining_ids:
                    if use_fafb_cache and fafb_mesh_missing:
                        # Load only those that are missing from both API cache and mesh cache
                        zip_needed = remaining_ids & set(fafb_mesh_missing)
                        if zip_needed:
                            self._vprint(f'  ℹ️  {len(zip_needed)} neurons loading from ZIP')
                            zip_skeletons = self._preload_fafb_skeletons(body_ids_filter=zip_needed)
                            fafb_skeleton_cache.update(zip_skeletons)
                            
                            # Auto-detect extrusions and fetch via API if enabled
                            if self.auto_fix_extrusions:
                                extrusion_ids = self._detect_extrusions_in_skeletons(zip_skeletons)
                                if extrusion_ids:
                                    self._vprint(f'  🔍 Detected {len(extrusion_ids)} skeletons with extrusion artifacts, fetching fresh from API...', level='simple')
                                    api_fixed = self._fetch_fafb_skeletons_via_api(extrusion_ids)
                                    if api_fixed:
                                        fafb_skeleton_cache.update(api_fixed)
                                        self._vprint(f'  ✓ Replaced {len(api_fixed)} extrusion-affected skeletons', level='simple')
                    elif not use_fafb_cache:
                        # Cache not used - load remaining from ZIP
                        self._vprint(f'  ℹ️  Loading {len(remaining_ids)} neurons from ZIP')
                        zip_skeletons = self._preload_fafb_skeletons(body_ids_filter=remaining_ids)
                        fafb_skeleton_cache.update(zip_skeletons)
                        
                        # Auto-detect extrusions and fetch via API if enabled
                        if self.auto_fix_extrusions:
                            extrusion_ids = self._detect_extrusions_in_skeletons(zip_skeletons)
                            if extrusion_ids:
                                self._vprint(f'  🔍 Detected {len(extrusion_ids)} skeletons with extrusion artifacts, fetching fresh from API...', level='simple')
                                api_fixed = self._fetch_fafb_skeletons_via_api(extrusion_ids)
                                if api_fixed:
                                    fafb_skeleton_cache.update(api_fixed)
                                    self._vprint(f'  ✓ Replaced {len(api_fixed)} extrusion-affected skeletons', level='simple')
                
                # Auto-fallback to API for any still missing (graceful degradation)
                # Skip IDs that are already in mesh cache (they don't need skeleton processing)
                # Build a type-robust lookup set for fafb_mesh_cache
                mesh_cache_lookup = set()
                for mid in fafb_mesh_cache.keys():
                    mesh_cache_lookup.add(mid)
                    if isinstance(mid, (int, np.integer)):
                        mesh_cache_lookup.add(str(mid))
                    elif isinstance(mid, str) and mid.isdigit():
                        mesh_cache_lookup.add(int(mid))
                
                still_missing = [bid for bid in all_fafb_body_ids 
                                if bid not in fafb_skeleton_cache 
                                and str(bid) not in fafb_skeleton_cache
                                and bid not in mesh_cache_lookup]  # Don't fetch if mesh cache has them
                if still_missing:
                    self._vprint(f'  ⚠️  {len(still_missing)} neurons not in local cache/ZIP, auto-fetching via CAVE API...')
                    api_fetched = self._fetch_fafb_skeletons_via_api(still_missing)
                    fafb_skeleton_cache.update(api_fetched)
        
        
        # Note: For legend_mode='type', neurons keep their layer colors but get separate legend entries
        # Per-neuron colors from CSV are stored in self._neuron_color_overrides (bodyId -> color)
        self._type_color_map = {}  # Not used anymore - types keep layer colors
        
        # Main progress bar for layers - always show when verbose is enabled
        layer_pbar = tqdm(range(n_layers), desc="Processing layers", 
                          disable=not self.verbose, leave=True, file=sys.stdout)
        
        for i in layer_pbar:
            layer_name = self.layer_names[i] if i < len(self.layer_names) else f"layer_{i}"
            n_in_layer = len(self.neuron_dfs[i]) if self.neuron_dfs[i] is not None else 0
            layer_pbar.set_postfix_str(f"{layer_name} ({n_in_layer} neurons)")
            
            # Determine if we need transformation
            # Use _needs_skeleton_transform() which checks for skip_transform flag (FAFB uses native coords)
            needs_transform = self._needs_skeleton_transform()
            template_info = None
            if self.brain_mesh in ['whole', 'template']:
                template_info = self._get_template_info()
            
            neuron_vols = None
            
            # For FAFB with caching: check which neurons already have cached meshes
            layer_body_ids = self.neuron_dfs[i]['bodyId'].tolist() if self.neuron_dfs[i] is not None else []
            cached_mesh_neurons = []  # MeshNeurons loaded from cache
            mesh_missing_ids = layer_body_ids  # IDs that need processing
            
            if use_fafb_cache and fafb_mesh_cache:
                # Separate cached vs missing with type-robust matching
                cached_mesh_neurons = []
                mesh_missing_ids = []
                for bid in layer_body_ids:
                    # Check both int and str versions of the ID
                    if bid in fafb_mesh_cache:
                        cached_mesh_neurons.append(fafb_mesh_cache[bid])
                    elif str(bid) in fafb_mesh_cache:
                        cached_mesh_neurons.append(fafb_mesh_cache[str(bid)])
                    elif isinstance(bid, str) and bid.isdigit() and int(bid) in fafb_mesh_cache:
                        cached_mesh_neurons.append(fafb_mesh_cache[int(bid)])
                    else:
                        mesh_missing_ids.append(bid)
                
                if cached_mesh_neurons:
                    self._vprint(f'    ✓ {len(cached_mesh_neurons)}/{len(layer_body_ids)} from mesh cache', level='full', use_tqdm=True)
            
            # Load from raw cache (for non-FAFB datasets). NeuPrint caches
            # hold the fixed 90%-simplified skeletons; a less-simplified
            # render (< 0.9) is too coarse for them, so the cache is ignored
            # and the RAW skeletons are re-fetched transiently (the
            # simplified cache is still written from the same fetch).
            ignore_skeleton_cache = (not is_fafb and self.skeleton_mode == 'tube'
                                     and self.skeleton_mesh_simplification
                                     < self.NEUPRINT_SKELETON_CACHE_LEVEL)
            cache_result = self._load_cached_neurons(self.neuron_dfs[i],
                                                     ignore_cache=ignore_skeleton_cache)
            cached_neurons, missing_ids = cache_result
            
            # Check and fix MANC scaling for cached neurons (handling legacy cache)
            if cached_neurons is not None and 'manc' in self.dataset.lower():
                try:
                    # Check first neuron
                    bbox = cached_neurons[0].bbox if hasattr(cached_neurons[0], 'bbox') else None
                    if bbox is not None and np.max(bbox) < 150000:
                        cached_neurons = cached_neurons * 8
                        self._vprint(f'  ℹ️  Applied 8x scaling to cached MANC skeletons', level='full')
                except Exception as e:
                    self._vprint(f'  ⚠️  Failed to check/apply scaling to cache: {e}', level='full')
            
            raw_neuron_vols = None
            
            # Fetch missing neurons (only those not in mesh cache for FAFB when cache is used)
            fetch_ids = mesh_missing_ids if is_fafb else missing_ids
            if fetch_ids:
                # Special handling for FAFB local data - use pre-loaded cache
                if fafb_skeleton_cache:
                    neurons = []
                    for bid in fetch_ids:
                        # Handle both int and string types for body ID lookup
                        if bid in fafb_skeleton_cache:
                            neurons.append(fafb_skeleton_cache[bid])
                        elif str(bid) in fafb_skeleton_cache:
                            neurons.append(fafb_skeleton_cache[str(bid)])
                        elif isinstance(bid, str) and bid.isdigit() and int(bid) in fafb_skeleton_cache:
                            neurons.append(fafb_skeleton_cache[int(bid)])
                    if neurons:
                        raw_neuron_vols = navis.NeuronList(neurons)

                # Fetch from API if not loaded locally
                if raw_neuron_vols is None and fetch_ids:
                    if self.client_type == 'flywire' and self.client_flywire:
                        missing_df = self.neuron_dfs[i][self.neuron_dfs[i]['bodyId'].isin(fetch_ids)]
                        # Retry logic for network errors
                        max_retries = 5
                        for attempt in range(max_retries):
                            try:
                                raw_neuron_vols = self.client_flywire.fetch_skeletons(self.layer_criteria[i], with_synapses=self.show_connectors)
                                break  # Success
                            except Exception as e:
                                error_msg = str(e)
                                is_network_error = any(x in error_msg.lower() for x in 
                                    ['timeout', 'connection', 'network', 'refused', 'reset', 'temporary'])
                                
                                if is_network_error and attempt < max_retries - 1:
                                    import time
                                    wait_time = (attempt + 1) * 2
                                    tqdm.write(f'  ⚠️  Network error, retrying in {wait_time}s (attempt {attempt + 1}/{max_retries}): {e}')
                                    time.sleep(wait_time)
                                else:
                                    tqdm.write(f'  ⚠️  FlyWire fetch failed for layer {layer_name}: {e}')
                                    raw_neuron_vols = None
                                    break
                    else:
                        missing_df = self.neuron_dfs[i][self.neuron_dfs[i]['bodyId'].isin(fetch_ids)].copy()
                        if not missing_df.empty:
                            # Ensure bodyId is int64 for neuprint compatibility
                            # NeuPrint/navis expects bodyId as int, not string
                            if missing_df['bodyId'].dtype == object or str(missing_df['bodyId'].dtype) == 'string':
                                try:
                                    missing_df['bodyId'] = missing_df['bodyId'].astype('int64')
                                except (ValueError, TypeError):
                                    pass  # Keep original type if conversion fails
                            kwargs = {
                                'with_synapses': self.show_connectors,
                                'missing_swc': 'warn',  # Skip missing skeletons instead of raising
                            }
                            if self.client:
                                kwargs['client'] = self.client
                            
                            # Retry logic for network errors
                            max_retries = 5
                            for attempt in range(max_retries):
                                try:
                                    raw_neuron_vols = neu.fetch_skeletons(missing_df, **kwargs)
                                    break  # Success
                                except Exception as e:
                                    error_msg = str(e)
                                    # Check if it's a network/connection error that might be retried
                                    is_network_error = any(x in error_msg.lower() for x in 
                                        ['timeout', 'connection', 'network', 'refused', 'reset', 'temporary'])
                                    
                                    if is_network_error and attempt < max_retries - 1:
                                        import time
                                        wait_time = (attempt + 1) * 2  # 2, 4, 6 seconds
                                        tqdm.write(f'  ⚠️  Network error, retrying in {wait_time}s (attempt {attempt + 1}/{max_retries}): {e}')
                                        time.sleep(wait_time)
                                    else:
                                        # Handle "No neurons matching the given criteria found!" and other errors
                                        # This can happen if neurons exist in NeuronBridge but not in NeuPrint (different versions)
                                        tqdm.write(f'  ⚠️  NeuPrint fetch failed for layer {layer_name}: {e}')
                                        raw_neuron_vols = None
                                        break
                
                # Apply scaling for MANC datasets (Raw -> NM)
                # MANC skeletons from NeuPrint are in 8nm voxels, but meshes are in nm
                if raw_neuron_vols is not None and 'manc' in self.dataset.lower():
                    # Check if scaling is needed (max coord < 150000 indicates raw units)
                    # Typical MANC nm extent is ~300k-500k
                    try:
                        bbox = raw_neuron_vols[0].bbox if hasattr(raw_neuron_vols[0], 'bbox') else None
                        if bbox is not None and np.max(bbox) < 150000:
                            raw_neuron_vols = raw_neuron_vols * 8
                            self._vprint(f'  ℹ️  Applied 8x scaling to fetched MANC skeletons', level='full')
                    except Exception as e:
                        self._vprint(f'  ⚠️  Failed to check/apply scaling: {e}', level='full')

                # Save to raw cache (for non-FAFB datasets)
                if raw_neuron_vols is not None and not is_fafb:
                    self._save_cached_neurons(self.neuron_dfs[i], raw_neuron_vols)
            
            # Combine cached and newly fetched neurons
            if cached_neurons is not None and raw_neuron_vols is not None:
                all_neurons = list(cached_neurons) + list(raw_neuron_vols)
                neuron_vols = navis.NeuronList(all_neurons)
            elif cached_neurons is not None:
                neuron_vols = cached_neurons
            elif raw_neuron_vols is not None:
                neuron_vols = raw_neuron_vols
            else:
                neuron_vols = None

            # Normalize to NeuronList so downstream len()/iteration works for single TreeNeuron
            if neuron_vols is not None and not isinstance(neuron_vols, (list, navis.NeuronList)):
                neuron_vols = navis.NeuronList([neuron_vols])
            
            # For FAFB with all meshes cached, we can skip skeleton processing
            if is_fafb and cached_mesh_neurons and (neuron_vols is None or len(neuron_vols) == 0):
                # All neurons loaded from mesh cache - neuron_vols stays None/empty
                # The combine block below will handle adding cached_mesh_neurons with simplification
                pass
            elif neuron_vols is None or len(neuron_vols) == 0:
                if cached_mesh_neurons:
                    # Partial cache hit - neuron_vols stays None/empty
                    # The combine block below will handle cached_mesh_neurons
                    pass
                else:
                    tqdm.write(f'  ⚠️  Failed to fetch skeletons for layer {i}: {layer_name}')
                    continue

            # Apply soma radius capping to prevent extrusion artifacts
            if self.soma_radius_cap is not None and self.skeleton_mode == 'tube':
                self._apply_soma_radius_cap(neuron_vols)

            # Transformation logic moved to after caching to ensure cache stores raw FlyWire coordinates

            
            # For FAFB: convert to mesh, apply soma-aware simplification, and cache (only when cache is used)
            # Cache stores meshes with soma-aware simplification for reuse
            if use_fafb_cache and mesh_missing_ids and neuron_vols is not None and self.skeleton_mode == 'tube':
                try:
                    import trimesh
                    meshes_to_cache = {}
                    mesh_neurons_list = []
                    cache_skel_simp = self.FAFB_MESH_CACHE_SIMPLIFICATION
                    cache_soma_simp = self.FAFB_MESH_CACHE_SOMA_SIMPLIFICATION
                    cache_soma_radius = self.FAFB_MESH_CACHE_SOMA_RADIUS
                    
                    # Convert mesh_missing_ids to a set with both int and str versions for robust matching
                    mesh_missing_ids_set = set()
                    for mid in mesh_missing_ids:
                        mesh_missing_ids_set.add(mid)
                        if isinstance(mid, (int, np.integer)):
                            mesh_missing_ids_set.add(str(mid))
                        elif isinstance(mid, str) and mid.isdigit():
                            mesh_missing_ids_set.add(int(mid))
                    
                    # Setup progress bar if verbose
                    iterator = neuron_vols
                    if self.verbose == 'full' or self.verbose is True:
                        iterator = tqdm(neuron_vols, desc="Simplifying meshes (soma-aware)", leave=False)

                    for n in iterator:
                        if hasattr(n, 'id') and n.id in mesh_missing_ids_set:
                            # Get soma position before conversion (TreeNeuron has this info)
                            soma_pos = None
                            if hasattr(n, 'soma_pos') and n.soma_pos is not None:
                                soma_pos = n.soma_pos
                            
                            # Convert TreeNeuron to MeshNeuron if needed
                            if isinstance(n, navis.TreeNeuron):
                                # Fix radii if needed
                                if hasattr(n, 'nodes') and 'radius' in n.nodes.columns:
                                    invalid_mask = (n.nodes['radius'] <= 0) | (n.nodes['radius'].isna())
                                    if invalid_mask.any():
                                        n.nodes.loc[invalid_mask, 'radius'] = 1
                                elif hasattr(n, 'nodes'):
                                    n.nodes['radius'] = 1
                                # Convert
                                if hasattr(navis, 'conversion') and hasattr(navis.conversion, 'tree2meshneuron'):
                                    mesh_n = navis.conversion.tree2meshneuron(n)
                                else:
                                    mesh_neurons_list.append(n)
                                    continue
                            elif isinstance(n, navis.MeshNeuron):
                                mesh_n = n
                            else:
                                mesh_neurons_list.append(n)
                                continue
                            
                            # Apply soma-aware simplification for caching
                            if mesh_n and hasattr(mesh_n, 'trimesh'):
                                try:
                                    # Use soma-aware simplification with cache parameters
                                    simplified_trimesh = self._simplify_mesh_with_soma_awareness(
                                        mesh_n.trimesh,
                                        skeleton_simp=cache_skel_simp,
                                        soma_simp=cache_soma_simp,
                                        soma_pos=soma_pos,
                                        soma_radius=cache_soma_radius
                                    )
                                    # Create new MeshNeuron with simplified mesh to ensure proper storage
                                    mesh_n = navis.MeshNeuron(simplified_trimesh)
                                    mesh_n.id = n.id  # Preserve original ID
                                    if hasattr(n, 'name'):
                                        mesh_n.name = n.name
                                except Exception as e:
                                    self._vprint(f'      ⚠️ Simplification failed for {n.id}: {e}', level='full', use_tqdm=True)
                            
                            meshes_to_cache[n.id] = mesh_n
                            mesh_neurons_list.append(mesh_n)
                        else:
                            mesh_neurons_list.append(n)
                    
                    # Save soma-aware simplified meshes to cache
                    if meshes_to_cache:
                        self._save_cached_fafb_meshes(meshes_to_cache)
                        self._vprint(f'    ✓ Cached {len(meshes_to_cache)} meshes (skel={cache_skel_simp}, soma={cache_soma_simp})', level='full', use_tqdm=True)
                    
                    # Apply additional simplification to newly cached neurons if target > cache level
                    target_simp = self.skeleton_mesh_simplification
                    if target_simp > cache_skel_simp and mesh_neurons_list:
                        remaining_after_cache = 1 - cache_skel_simp
                        remaining_target = 1 - target_simp
                        additional_keep_factor = remaining_target / remaining_after_cache
                        self._vprint(f'    ⚡ Applying additional simplification to new meshes: {target_simp} (keep {additional_keep_factor:.1%})', level='full', use_tqdm=True)
                        
                        further_simplified = []
                        for mesh_n in mesh_neurons_list:
                            if hasattr(mesh_n, 'trimesh'):
                                n_faces = len(mesh_n.trimesh.faces)
                                target_faces = max(100, int(n_faces * additional_keep_factor))  # Keep at least 100 faces
                                if target_faces < n_faces:
                                    try:
                                        simplified_trimesh = self._simplify_mesh_open3d(mesh_n.trimesh, target_faces)
                                        new_mesh = navis.MeshNeuron(simplified_trimesh)
                                        new_mesh.id = mesh_n.id if hasattr(mesh_n, 'id') else None
                                        if hasattr(mesh_n, 'name'):
                                            new_mesh.name = mesh_n.name
                                        further_simplified.append(new_mesh)
                                        continue
                                    except Exception:
                                        pass
                            further_simplified.append(mesh_n)
                        mesh_neurons_list = further_simplified
                    
                    # Update neuron_vols with mesh versions
                    if mesh_neurons_list:
                        neuron_vols = navis.NeuronList(mesh_neurons_list)
                except Exception as e:
                    self._vprint(f'    ⚠️ FAFB mesh caching failed: {e}', level='full')
            
            # For FAFB: combine cached + newly processed neurons, then merge by layer if needed
            # This block handles:
            # 1. When cache is used (simplification >= 0.9): combine cached + new, apply additional simp if > 0.9
            # 2. When cache is not used (simplification < 0.9): just process neuron_vols for merging
            # Set flag to skip generic simplification block below (FAFB is already simplified here)
            fafb_already_simplified = False
            if is_fafb and self.skeleton_mode == 'tube':
                import trimesh
                
                all_mesh_neurons = []
                
                # Add cached neurons if available
                if cached_mesh_neurons:
                    # Apply additional simplification if user wants > 0.9
                    target_simp = self.skeleton_mesh_simplification
                    cache_simp = self.FAFB_MESH_CACHE_SIMPLIFICATION
                    
                    if target_simp > cache_simp:
                        # Calculate additional simplification factor
                        remaining_after_cache = 1 - cache_simp  # e.g., 0.1 for 90%
                        remaining_target = 1 - target_simp  # e.g., 0.05 for 95%
                        additional_keep_factor = remaining_target / remaining_after_cache
                        
                        self._vprint(f'    ⚡ Applying additional simplification: {target_simp} (keep {additional_keep_factor:.1%} of cached)', level='full', use_tqdm=True)
                        
                        simplified_cached = []
                        for mesh_n in cached_mesh_neurons:
                            if hasattr(mesh_n, 'trimesh'):
                                n_faces = len(mesh_n.trimesh.faces)
                                target_faces = max(100, int(n_faces * additional_keep_factor))  # Keep at least 100 faces
                                if target_faces < n_faces:
                                    try:
                                        simplified_trimesh = self._simplify_mesh_open3d(mesh_n.trimesh, target_faces)
                                        new_mesh = navis.MeshNeuron(simplified_trimesh)
                                        new_mesh.id = mesh_n.id if hasattr(mesh_n, 'id') else None
                                        if hasattr(mesh_n, 'name'):
                                            new_mesh.name = mesh_n.name
                                        simplified_cached.append(new_mesh)
                                        continue
                                    except Exception:
                                        pass
                            simplified_cached.append(mesh_n)
                        cached_mesh_neurons = simplified_cached
                    
                    all_mesh_neurons.extend(cached_mesh_neurons)
                
                # Add newly processed neurons
                if neuron_vols is not None and len(neuron_vols) > 0:
                    neurons_list = list(neuron_vols) if isinstance(neuron_vols, navis.NeuronList) else [neuron_vols]
                    
                    # When cache not used (simplification < 0.9), need to convert and simplify here
                    # This path processes neurons loaded directly from ZIP with user's actual simplification setting
                    if not use_fafb_cache:
                        processed_neurons = []
                        target_simp = self.skeleton_mesh_simplification
                        soma_simp = self.soma_mesh_simplification
                        use_soma_aware = soma_simp is not None and soma_simp != target_simp
                        
                        if use_soma_aware:
                            self._vprint(f'    ⚡ Processing {len(neurons_list)} neurons from ZIP (skeleton simp={target_simp}, soma simp={soma_simp})', level='full', use_tqdm=True)
                        else:
                            self._vprint(f'    ⚡ Processing {len(neurons_list)} neurons from ZIP (target simplification={target_simp})', level='full', use_tqdm=True)
                        
                        for n in neurons_list:
                            # Get soma position before conversion (TreeNeuron has this info)
                            soma_pos = None
                            if use_soma_aware and hasattr(n, 'soma_pos') and n.soma_pos is not None:
                                soma_pos = n.soma_pos
                            
                            # Convert TreeNeuron to MeshNeuron if needed
                            if isinstance(n, navis.TreeNeuron):
                                if hasattr(n, 'nodes') and 'radius' in n.nodes.columns:
                                    invalid_mask = (n.nodes['radius'] <= 0) | (n.nodes['radius'].isna())
                                    if invalid_mask.any():
                                        n.nodes.loc[invalid_mask, 'radius'] = 1
                                elif hasattr(n, 'nodes'):
                                    n.nodes['radius'] = 1
                                if hasattr(navis, 'conversion') and hasattr(navis.conversion, 'tree2meshneuron'):
                                    mesh_n = navis.conversion.tree2meshneuron(n)
                                else:
                                    processed_neurons.append(n)
                                    continue
                            elif isinstance(n, navis.MeshNeuron):
                                mesh_n = n
                            else:
                                processed_neurons.append(n)
                                continue
                            
                            # Apply simplification (soma-aware if configured)
                            if target_simp > 0 and mesh_n and hasattr(mesh_n, 'trimesh'):
                                n_faces = len(mesh_n.trimesh.faces)
                                try:
                                    if use_soma_aware and soma_pos is not None:
                                        # Use soma-aware simplification with different levels
                                        simplified_trimesh = self._simplify_mesh_with_soma_awareness(
                                            mesh_n.trimesh, 
                                            skeleton_simp=target_simp,
                                            soma_simp=soma_simp,
                                            soma_pos=soma_pos,
                                            soma_radius=self.soma_region_radius
                                        )
                                    else:
                                        # Standard uniform simplification
                                        target_faces = max(100, int(n_faces * (1 - target_simp)))
                                        simplified_trimesh = self._simplify_mesh_open3d(mesh_n.trimesh, target_faces)
                                    
                                    mesh_n = navis.MeshNeuron(simplified_trimesh)
                                    mesh_n.id = n.id if hasattr(n, 'id') else None
                                    if hasattr(n, 'name'):
                                        mesh_n.name = n.name
                                except Exception:
                                    pass
                            
                            processed_neurons.append(mesh_n)
                        
                        all_mesh_neurons.extend(processed_neurons)
                    else:
                        # Cache path: convert TreeNeurons to MeshNeurons if needed
                        converted_list = []
                        for n in neurons_list:
                            if isinstance(n, navis.MeshNeuron):
                                converted_list.append(n)
                            elif isinstance(n, navis.TreeNeuron):
                                # Convert TreeNeuron to MeshNeuron
                                try:
                                    if hasattr(n, 'nodes') and 'radius' in n.nodes.columns:
                                        invalid_mask = (n.nodes['radius'] <= 0) | (n.nodes['radius'].isna())
                                        if invalid_mask.any():
                                            n.nodes.loc[invalid_mask, 'radius'] = 1
                                    elif hasattr(n, 'nodes'):
                                        n.nodes['radius'] = 1
                                    if hasattr(navis, 'conversion') and hasattr(navis.conversion, 'tree2meshneuron'):
                                        mesh_n = navis.conversion.tree2meshneuron(n)
                                        mesh_n.id = n.id if hasattr(n, 'id') else None
                                        if hasattr(n, 'name'):
                                            mesh_n.name = n.name
                                        converted_list.append(mesh_n)
                                    else:
                                        converted_list.append(n)  # Keep original if conversion unavailable
                                except Exception:
                                    converted_list.append(n)  # Keep original if conversion fails
                            else:
                                converted_list.append(n)
                        all_mesh_neurons.extend(converted_list)
                
                # Set neuron_vols from all_mesh_neurons (no mesh merging - legend grouping handles display)
                if all_mesh_neurons:
                    neuron_vols = navis.NeuronList(all_mesh_neurons)
                
                # Mark FAFB as already simplified to skip generic simplification below
                fafb_already_simplified = True

            # Transform if needed (Now applies to ALL neurons: cached + new)
            needs_actual_transform = needs_transform
            if needs_actual_transform and neuron_vols is not None:
                layer_pbar.set_postfix_str(f"{layer_name} (transforming {len(neuron_vols)}...)")
                try:
                    # Ensure float64 coordinates to avoid dtype warnings in navis
                    if isinstance(neuron_vols, (list, navis.NeuronList)):
                        for n in neuron_vols:
                            if hasattr(n, 'nodes') and isinstance(n.nodes, pd.DataFrame):
                                for col in ['x', 'y', 'z']:
                                    if col in n.nodes.columns:
                                        n.nodes[col] = n.nodes[col].astype('float64')
                            # Also handle MeshNeuron vertices
                            if hasattr(n, 'vertices') and n.vertices is not None:
                                n.vertices = n.vertices.astype('float64')
                    elif hasattr(neuron_vols, 'nodes') and isinstance(neuron_vols.nodes, pd.DataFrame):
                         for col in ['x', 'y', 'z']:
                            if col in neuron_vols.nodes.columns:
                                neuron_vols.nodes[col] = neuron_vols.nodes[col].astype('float64')

                    # Per-neuron transform with explicit progress: a single
                    # navis.xform_brain() call gives no visibility into WHICH
                    # neuron it is on, so a stall looked like a silent hang.
                    neuron_vols = self._xform_neurons_safe(
                        neuron_vols,
                        source=template_info['source'],
                        target=template_info['target'],
                        layer_label=f"Layer {i} ({layer_name})",
                    )
                except Exception as e:
                    tqdm.write(f'  ⚠️  Layer {i} transform failed: {e}')
                    if self._dataset_needs_transform() and not self._check_and_download_transforms():
                        self.brain_mesh = 'none'
                    else:
                        # Retry transformation after download
                        try:
                            neuron_vols = self._xform_neurons_safe(
                                neuron_vols,
                                source=template_info['source'],
                                target=template_info['target'],
                                layer_label=f"Layer {i} ({layer_name}, retry)",
                            )
                        except Exception as retry_e:
                            tqdm.write(f'  ⚠️  Transformation still failed, setting brain_mesh to "none"')
                            self.brain_mesh = 'none'
            
            # Apply FAFB tilt correction if using template mode
            # This corrects the left-right tilt in the FLYWIRE template mesh
            if is_fafb and self.brain_mesh == 'template' and neuron_vols is not None:
                neuron_vols = self._apply_fafb_tilt_correction(neuron_vols)
            
            # Ensure iterable after potential transforms (navis may return TreeNeuron)
            if neuron_vols is not None and not isinstance(neuron_vols, (list, navis.NeuronList)):
                neuron_vols = navis.NeuronList([neuron_vols])

            # Mirror neurons if requested
            if self.mirror_on_contralateral:
                try:
                    template = None
                    if self.brain_mesh == 'whole':
                        template_info = self._get_template_info()
                        template = template_info['target']
                    elif self.brain_mesh == 'template':
                         if 'hemibrain' in self.dataset or 'optic-lobe' in self.dataset:
                             template = 'JRCFIB2018F'
                         elif 'male-cns' in self.dataset:
                             template = 'JRCFIB2022M'
                    
                    if template:
                        mirrored = navis.mirror_brain(neuron_vols, template, mirror_axis='x')
                        if isinstance(neuron_vols, navis.NeuronList):
                            neuron_vols = neuron_vols + mirrored
                        else:
                            neuron_vols = navis.NeuronList([neuron_vols, mirrored])
                except Exception as e:
                    tqdm.write(f'  ⚠️ Mirror failed for layer {i}: {e}')

            # Simplify individual neurons if requested (and not merging)
            # Skip for FAFB - already handled in the FAFB-specific block above
            # NeuPrint: on-disk skeletons are ALREADY at the fixed
            # 90%-simplified cache level; at exactly that level the
            # render-time decimation is skipped (the tube mesh is already
            # at the cache level and decimating again would double-reduce:
            # 10% nodes -> ~1% faces). Levels above the cache level apply
            # the remaining relative reduction; below it the render ran on
            # transient RAW fetches and decimates at the user's level.
            render_simplification = self._effective_render_simplification(is_fafb)
            if render_simplification > 0 and self.skeleton_mode == 'tube' and not fafb_already_simplified:
                try:
                    import trimesh
                    simplified_neurons = []
                    total_original_faces = 0
                    total_simplified_faces = 0
                    # Ensure iterable
                    neurons_to_simplify = neuron_vols if isinstance(neuron_vols, navis.NeuronList) else [neuron_vols]
                    
                    for n in neurons_to_simplify:
                        try:
                            # Convert to mesh if needed (TreeNeuron -> MeshNeuron)
                            mesh_n = None
                            if isinstance(n, navis.TreeNeuron):
                                # Fix radii if needed
                                if hasattr(n, 'nodes') and 'radius' in n.nodes.columns:
                                    invalid_mask = (n.nodes['radius'] <= 0) | (n.nodes['radius'].isna())
                                    if invalid_mask.any():
                                        n.nodes.loc[invalid_mask, 'radius'] = 1
                                elif hasattr(n, 'nodes'):
                                    n.nodes['radius'] = 1
                                
                                # Convert
                                if hasattr(navis, 'conversion') and hasattr(navis.conversion, 'tree2meshneuron'):
                                    mesh_n = navis.conversion.tree2meshneuron(n)
                            elif isinstance(n, navis.MeshNeuron):
                                mesh_n = n
                                
                            # Simplify if we have a mesh neuron
                            # MeshNeuron.trimesh is read-only, so we must create a new MeshNeuron
                            if mesh_n and hasattr(mesh_n, 'trimesh'):
                                n_faces = len(mesh_n.trimesh.faces)
                                total_original_faces += n_faces
                                target_faces = max(100, int(n_faces * (1 - render_simplification)))  # Keep at least 100 faces
                                if target_faces < n_faces:
                                    # Use open3d for accurate simplification (trimesh 4.x fast_simplification only achieves ~60%)
                                    # Create NEW MeshNeuron from simplified trimesh (can't just assign to .trimesh)
                                    try:
                                        simplified_tm = self._simplify_mesh_open3d(mesh_n.trimesh, target_faces)
                                    except Exception as simp_err:
                                        # Fallback to original if simplification fails
                                        simplified_tm = mesh_n.trimesh
                                    
                                    new_mesh_n = navis.MeshNeuron(simplified_tm)
                                    # Copy over attributes
                                    new_mesh_n.id = mesh_n.id if hasattr(mesh_n, 'id') else n.id
                                    if hasattr(mesh_n, 'name'):
                                        new_mesh_n.name = mesh_n.name
                                    total_simplified_faces += len(new_mesh_n.trimesh.faces)
                                    simplified_neurons.append(new_mesh_n)
                                else:
                                    total_simplified_faces += n_faces
                                    simplified_neurons.append(mesh_n)
                            else:
                                # Keep original if conversion failed or not applicable
                                simplified_neurons.append(n)
                        except Exception as e:
                            # Keep original if failed
                            simplified_neurons.append(n if mesh_n is None else mesh_n)
                    
                    neuron_vols = navis.NeuronList(simplified_neurons)
                    
                    # Log simplification results
                    if total_original_faces > 0:
                        reduction = (1 - total_simplified_faces / total_original_faces) * 100
                        self._vprint(f'    ✓ Simplified: {total_original_faces:,} → {total_simplified_faces:,} faces ({reduction:.1f}% reduction)', level='full', use_tqdm=True)
                except Exception as e:
                    self._vprint(f'    ⚠️ Simplification failed: {e}', level='full', use_tqdm=True)
                    pass  # Keep original neurons if simplification fails

            # Update status and plot
            layer_pbar.set_postfix_str(f"{layer_name} (plotting...)")
            
            # Determine soma rendering
            show_soma_here = self.show_soma if not isinstance(neuron_vols, navis.Volume) else False
            
            # Determine alpha: use color's alpha if explicit, otherwise use neuron_alpha
            if self._neuron_colors_have_explicit_alpha:
                layer_neuron_alpha = self._extract_alpha_from_color(self.neuron_colors[i])
            else:
                layer_neuron_alpha = self.neuron_alpha

            if self.skeleton_mode == 'tube' and neuron_vols is not None:
                neurons_for_export = neuron_vols if isinstance(neuron_vols, (list, navis.NeuronList)) else [neuron_vols]
                for unit_index, neuron in enumerate(neurons_for_export):
                    neuron_id = str(getattr(neuron, 'id', f'neuron_{unit_index}'))
                    neuron_color = self._resolve_neuron_color(neuron_id, i)
                    neuron_alpha = self._extract_alpha_from_color(neuron_color)
                    self._append_exportable_mesh(
                        neuron,
                        color=neuron_color,
                        alpha=neuron_alpha,
                        name=getattr(neuron, 'name', neuron_id),
                        role='neuron',
                    )
            
            if self.backend == 'plotly':
                trace_entries = []
                if self.color_mode == 'per_neuron' and not isinstance(neuron_vols, navis.Volume):
                    for unit_index, neuron in enumerate(neuron_vols):
                        neuron_id = str(getattr(neuron, 'id', f'neuron_{unit_index}'))
                        neuron_color = self._resolve_neuron_color(neuron_id, i)
                        neuron_alpha = self._extract_alpha_from_color(neuron_color)
                        neuron_color_hex = self._rgba_to_hex(neuron_color)

                        with self._suppress_output():
                            fig_layer = navis.plot3d(
                                neuron,
                                backend='plotly',
                                color=neuron_color_hex,
                                alpha=neuron_alpha,
                                soma=show_soma_here,
                                radius=self.show_skeleton_radius,
                                connectors=self.show_connectors,
                            )

                        for trace in fig_layer.data:
                            trace_entries.append((trace, neuron_id, unit_index, neuron_color))
                else:
                    # Convert rgba color to hex for navis compatibility
                    layer_color_hex = self._rgba_to_hex(self.neuron_colors[i])

                    with self._suppress_output():
                        fig_layer = navis.plot3d(
                            neuron_vols,
                            backend='plotly',
                            color=layer_color_hex,
                            alpha=layer_neuron_alpha,
                            soma=show_soma_here,
                            # fig=self.fig_3d,
                            radius=self.show_skeleton_radius,
                            connectors=self.show_connectors if not isinstance(neuron_vols, navis.Volume) else False,
                        )
                    fig_traces = fig_layer.data
                    for j, trace in enumerate(fig_traces):
                        existing_name = getattr(trace, 'name', None)
                        if existing_name:
                            neuron_id = str(existing_name)
                        elif j < len(neuron_vols):
                            neuron_id = str(neuron_vols[j].id)
                        else:
                            neuron_id = f"neuron_{j}"

                        neuron_color = self._resolve_neuron_color(neuron_id, i)
                        trace_entries.append((trace, neuron_id, j, neuron_color))

                # Build a mapping of neuron ID to type for 'type' legend mode
                neuron_type_map = {}
                if self.legend_mode == 'type' and self.neuron_dfs[i] is not None:
                    ndf = self.neuron_dfs[i]
                    type_col = None
                    for col in ['type', 'cell_type', 'neuronType']:
                        if col in ndf.columns:
                            type_col = col
                            break
                    if type_col and 'bodyId' in ndf.columns:
                        for _, row in ndf.iterrows():
                            body_id = str(row['bodyId'])
                            neuron_type = str(row[type_col]) if pd.notna(row[type_col]) else None
                            if neuron_type:
                                neuron_type_map[body_id] = neuron_type
                                # Also map neuron name/instance to type if available (helpful when trace name != bodyId)
                                for name_col in ['name', 'instance', 'roi']:
                                    if name_col in row and pd.notna(row[name_col]):
                                        neuron_type_map[str(row[name_col])] = neuron_type

                # Track which legend groups we've shown (for type/layer modes)
                shown_legend_groups = set()
                # Track legend info for fixing opacity later
                legend_color_map = {}  # legend_group -> (color, should_show)

                for trace, neuron_id, source_index, neuron_color in trace_entries:
                    # Enforce opacity for lines if not already set or if we want to override
                    if self.skeleton_mode == 'line':
                        trace.opacity = self._extract_alpha_from_color(neuron_color)

                    neuron_color = self._resolve_neuron_color(neuron_id, i)
                    if neuron_color != self.neuron_colors[i]:
                        self._apply_plotly_trace_color(trace, neuron_color)

                    if self.legend_mode == 'layer':
                        # Group all neurons in layer under one legend entry
                        legend_group = self.layer_names[i]
                        trace.name = legend_group
                        trace.legendgroup = legend_group
                        should_show = legend_group not in shown_legend_groups
                        trace.showlegend = should_show
                        if should_show:
                            legend_color_map[legend_group] = (self.neuron_colors[i], True)
                        shown_legend_groups.add(legend_group)
                        trace.hovertemplate = '<b>%{fullData.name}</b><extra></extra>'
                        trace.hoverinfo = 'name'
                        self.fig_3d.add_trace(trace)

                    elif self.legend_mode == 'type':
                        # Group by neuron type - each type gets separate legend but keeps layer color
                        neuron_type = neuron_type_map.get(neuron_id, None)
                        
                        # Fallback: try different ID strategies if type not found
                        if not neuron_type and source_index < len(neuron_vols):
                            # Try using the ID from the source neuron object
                            try:
                                vid = str(neuron_vols[source_index].id)
                                neuron_type = neuron_type_map.get(vid, None)
                            except:
                                pass
                            
                            # If still not found, try using the name from source neuron
                            if not neuron_type and hasattr(neuron_vols[source_index], 'name'):
                                try:
                                    vname = str(neuron_vols[source_index].name)
                                    neuron_type = neuron_type_map.get(vname, None)
                                except:
                                    pass

                        if neuron_type:
                            legend_group = f"{neuron_type}"
                        else:
                            # Fallback to layer name if type unknown
                            legend_group = self.layer_names[i]
                        
                        trace.name = legend_group
                        trace.legendgroup = legend_group  # Same type shares legend group
                        should_show = legend_group not in shown_legend_groups
                        trace.showlegend = should_show
                        if should_show:
                            legend_color_map[legend_group] = (neuron_color, True)
                        shown_legend_groups.add(legend_group)
                        trace.hovertemplate = '<b>%{fullData.name}</b><extra></extra>'
                        trace.hoverinfo = 'name'
                        self.fig_3d.add_trace(trace)

                    elif self.legend_mode == 'single':
                        # Each neuron gets its own legend entry with layer color
                        new_trace_name = f"{neuron_id}_{self.layer_names[i]}"
                        
                        trace.name = new_trace_name
                        trace.legendgroup = new_trace_name
                        trace.showlegend = True
                        legend_color_map[new_trace_name] = (neuron_color, True)
                        trace.hoverinfo = 'name'
                        trace.hovertemplate = '<b>%{fullData.name}</b><extra></extra>'
                        self.fig_3d.add_trace(trace)
                    else:
                        raise ValueError(f'legend_mode {self.legend_mode} not supported')
                
                # Fix legend opacity: Add invisible marker traces with full opacity for legend display
                # Needed when alpha < 1.0 (either from neuron_alpha or explicit alpha in colors)
                # to ensure legend color patches are clearly visible
                needs_legend_fix = (layer_neuron_alpha < 1.0 or self._neuron_colors_have_explicit_alpha) and legend_color_map
                if needs_legend_fix:
                    import plotly.graph_objects as go
                    for legend_group, (color, _) in legend_color_map.items():
                        # Get opaque version of color for legend display
                        opaque_color = self._get_opaque_color(color)
                        # Add an invisible scatter point (no coords = not rendered) for legend only
                        legend_trace = go.Scatter3d(
                            x=[None], y=[None], z=[None],  # No coordinates = invisible in plot
                            mode='markers',
                            marker=dict(size=10, color=opaque_color, opacity=1.0),  # Full opacity for legend
                            name=legend_group,
                            legendgroup=legend_group,
                            showlegend=True,
                            hoverinfo='skip',
                        )
                        # Hide original traces from legend (but keep them visible in plot)
                        for existing_trace in self.fig_3d.data:
                            if getattr(existing_trace, 'legendgroup', None) == legend_group:
                                existing_trace.showlegend = False
                        # Add the legend-only trace
                        self.fig_3d.add_trace(legend_trace)
            
            elif self.backend == 'k3d':
                try:
                    if self.color_mode == 'per_neuron' and not isinstance(neuron_vols, navis.Volume):
                        for unit_index, neuron in enumerate(neuron_vols):
                            neuron_id = str(getattr(neuron, 'id', f'neuron_{unit_index}'))
                            neuron_color = self._resolve_neuron_color(neuron_id, i)
                            neuron_alpha = self._extract_alpha_from_color(neuron_color)
                            neuron_color_hex = self._rgba_to_hex(neuron_color)

                            with self._suppress_output():
                                temp_plot = navis.plot3d(
                                    neuron,
                                    backend='k3d',
                                    color=neuron_color_hex,
                                    alpha=neuron_alpha,
                                    soma=show_soma_here,
                                    radius=self.show_skeleton_radius,
                                    connectors=self.show_connectors,
                                    inline=False
                                )

                            for obj in temp_plot.objects:
                                self._apply_k3d_object_color(obj, neuron_color)
                                if hasattr(obj, 'name'):
                                    obj.name = self.layer_names[i]
                                self.fig_3d += obj
                    else:
                        # navis.plot3d with k3d backend returns a k3d.Plot object
                        # Convert rgba color to hex for navis compatibility
                        layer_color_hex = self._rgba_to_hex(self.neuron_colors[i])
                        
                        with self._suppress_output():
                            temp_plot = navis.plot3d(
                                neuron_vols,
                                backend='k3d',
                                color=layer_color_hex,
                                alpha=layer_neuron_alpha,  # Use layer-specific alpha (from color or neuron_alpha)
                                soma=show_soma_here,
                                radius=self.show_skeleton_radius,
                                connectors=self.show_connectors if not isinstance(neuron_vols, navis.Volume) else False,
                                inline=False
                            )
                        
                        for j, obj in enumerate(temp_plot.objects):
                            neuron_id = getattr(obj, 'name', None)
                            if not neuron_id and j < len(neuron_vols):
                                neuron_id = str(neuron_vols[j].id)
                            neuron_color = self._resolve_neuron_color(neuron_id, i)
                            self._apply_k3d_object_color(obj, neuron_color)
                            if hasattr(obj, 'name'):
                                obj.name = self.layer_names[i]
                            self.fig_3d += obj
                except Exception as e:
                    self._vprint(f'⚠️  k3d plotting failed: {e}', level='full')

        return 0
    
    def _get_synapse_cache_path(self, pre_id, post_id):
        """Get cache file path for synapses between a specific pre/post neuron pair.
        
        Cache structure: cache/{dataset}/synapses/{pre_id}_{post_id}.parquet
        
        This caches by neuron pair rather than by layer, because:
        1. The same synapse data is reusable across different queries
        2. Layer indices are arbitrary and session-specific
        3. Avoids duplicate storage of the same synaptic connections
        """
        cache_dir = self._get_cache_path('synapses')  # Note: 'synapses' not 'synapse'
        return os.path.join(cache_dir, f'{pre_id}_{post_id}.parquet')
    
    def _load_cached_synapses(self, source_ids, target_ids):
        """Load cached synapse connections for given source/target neuron pairs.
        
        For FlyWire/FAFB datasets, loads from the master synapse table at:
            datasets/{dataset}/{dataset}_synapse_table.parquet
        and filters by source_ids and target_ids.
        
        For other datasets, loads individual cache files per neuron pair from:
            cache/{dataset}/synapses/{pre_id}_{post_id}.parquet
        
        Args:
            source_ids: Set/list of source (presynaptic) body IDs
            target_ids: Set/list of target (postsynaptic) body IDs
            
        Returns:
            Tuple of (cached_df, missing_pairs) where:
            - cached_df: DataFrame of cached synapses (may be None if nothing cached)
            - missing_pairs: List of (pre_id, post_id) tuples not found in cache
        """
        if not self.cache_synapses:
            # Return all pairs as missing
            all_pairs = [(s, t) for s in source_ids for t in target_ids]
            return None, all_pairs
        
        source_ids = set(str(s) for s in source_ids)
        target_ids = set(str(t) for t in target_ids)
        
        # For FlyWire/FAFB, use the master synapse table from datasets folder
        if 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower():
            synapse_table_path = self._get_synapse_table_path()
            if os.path.exists(synapse_table_path):
                try:
                    # Load master synapse table
                    synapse_df = pd.read_parquet(synapse_table_path)
                    self._vprint(f'  ✓ Loaded synapse table from {synapse_table_path}', level='full')
                    
                    # Determine column names (may vary by dataset)
                    pre_col = 'pre_pt_root_id' if 'pre_pt_root_id' in synapse_df.columns else 'bodyId_pre'
                    post_col = 'post_pt_root_id' if 'post_pt_root_id' in synapse_df.columns else 'bodyId_post'
                    
                    # Convert to string for matching
                    synapse_df[pre_col] = synapse_df[pre_col].astype(str)
                    synapse_df[post_col] = synapse_df[post_col].astype(str)
                    
                    filtered_df = synapse_df[
                        (synapse_df[pre_col].isin(source_ids)) & 
                        (synapse_df[post_col].isin(target_ids))
                    ]
                    self._vprint(f'  ✓ Filtered to {len(filtered_df)} synapses between {len(source_ids)} sources and {len(target_ids)} targets', level='full')
                    # For FlyWire, master table has all data - no missing pairs
                    return filtered_df, []
                except Exception as e:
                    self._vprint(f'  ⚠ Failed to load synapse table: {e}', level='full')
                    all_pairs = [(s, t) for s in source_ids for t in target_ids]
                    return None, all_pairs
            else:
                self._vprint(f'  ⚠ Synapse table not found at {synapse_table_path}', level='full')
                all_pairs = [(s, t) for s in source_ids for t in target_ids]
                return None, all_pairs
        
        # For other datasets, load from individual cache files per neuron pair
        cached_dfs = []
        missing_pairs = []
        
        for pre_id in source_ids:
            for post_id in target_ids:
                cache_file = self._get_synapse_cache_path(pre_id, post_id)
                if os.path.exists(cache_file):
                    try:
                        df = pd.read_parquet(cache_file)
                        if not df.empty:
                            cached_dfs.append(df)
                    except Exception as e:
                        self._vprint(f'  ⚠ Cache load failed for {pre_id}→{post_id}: {e}', level='full')
                        missing_pairs.append((pre_id, post_id))
                else:
                    missing_pairs.append((pre_id, post_id))
        
        if cached_dfs:
            cached_df = pd.concat(cached_dfs, ignore_index=True)
            self._vprint(f'  ✓ Loaded {len(cached_df)} synapses from cache ({len(cached_dfs)} pairs cached, {len(missing_pairs)} pairs missing)', level='full')
        else:
            cached_df = None
            
        return cached_df, missing_pairs
    
    def _save_cached_synapses(self, conn_df, attempted_pairs=None):
        """Save synapse connections to cache, organized by pre/post neuron pairs.
        
        Each unique (pre_id, post_id) pair gets its own cache file at:
            cache/{dataset}/synapses/{pre_id}_{post_id}.parquet
            
        This approach ensures:
        1. Synapses are cached by their actual content (neuron pairs + positions)
        2. Same synapse data is reusable across different queries/layers
        3. Incremental caching - only fetch what's not already cached
        
        Args:
            conn_df: DataFrame containing synapse data
            attempted_pairs: Optional list of (pre_id, post_id) tuples that were queried.
                             Used to cache empty results for pairs with no synapses.
        """
        if not self.cache_synapses:
            return
            
        # Do not cache for FlyWire/FAFB - they use the master synapse table
        if 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower():
            return
        
        # Track which pairs we have saved
        saved_pairs = set()
        
        if conn_df is not None and not conn_df.empty:
            # Determine column names for pre/post body IDs
            pre_col = 'bodyId_pre' if 'bodyId_pre' in conn_df.columns else 'pre_pt_root_id'
            post_col = 'bodyId_post' if 'bodyId_post' in conn_df.columns else 'post_pt_root_id'
            
            if pre_col not in conn_df.columns or post_col not in conn_df.columns:
                self._vprint(f'  ⚠ Cannot cache synapses: missing {pre_col} or {post_col} columns', level='full')
                return
            
            # Group by pre/post pairs and save each group
            saved_count = 0
            for (pre_id, post_id), group_df in conn_df.groupby([pre_col, post_col]):
                pre_id_str = str(pre_id)
                post_id_str = str(post_id)
                cache_file = self._get_synapse_cache_path(pre_id_str, post_id_str)
                
                try:
                    group_df.to_parquet(cache_file, index=False)
                    saved_count += 1
                    saved_pairs.add((pre_id_str, post_id_str))
                except Exception as e:
                    self._vprint(f'  ⚠ Cache save failed for {pre_id}→{post_id}: {e}', level='full')
            
            self._vprint(f'  💾 Saved synapses to cache ({saved_count} neuron pairs)', level='full')

        # Handle empty results for attempted pairs
        if attempted_pairs:
            empty_saved_count = 0
            for pre_id, post_id in attempted_pairs:
                pre_id_str = str(pre_id)
                post_id_str = str(post_id)
                if (pre_id_str, post_id_str) not in saved_pairs:
                    # Save empty dataframe
                    cache_file = self._get_synapse_cache_path(pre_id_str, post_id_str)
                    try:
                        # Create empty DF
                        pd.DataFrame().to_parquet(cache_file)
                        empty_saved_count += 1
                    except Exception as e:
                        self._vprint(f'  ⚠ Cache save failed for empty {pre_id}→{post_id}: {e}', level='full')
            
            if empty_saved_count > 0:
                self._vprint(f'  💾 Cached {empty_saved_count} empty synapse pairs', level='full')
    
    def plot_synapses(self):
        if self.skip_synapse:
            self._vprint('Skipping synapse plotting as requested.', level='full')
            return

        for i in range(len(self.neuron_layers) - 1):
            source_criteria = self.layer_criteria[i]
            target_criteria = self.layer_criteria[i + 1]
            # Use a single file for all synapse layers, consistent with neuron_info.xlsx
            file_path = os.path.join(self.save_folder, self.saveas + '_synapses.xlsx')
            conn_df = None

            # --- Begin FlyWire/NeuPrint synapse loading logic ---
            if self.client_type == 'flywire':
                # Try loading from local file first
                # Find dataset folder and synapse file dynamically
                dataset_normalized = self.dataset.replace(':', '_').replace('.', '_')
                dataset_dir = os.path.join(self.script_path, 'datasets', dataset_normalized)
                
                # Explicitly look for the file generated by FAFB_file_converter
                parquet_file = os.path.join(dataset_dir, f"{dataset_normalized}_synapse_table.parquet")
                
                # Use Parquet if available
                if os.path.exists(parquet_file):
                    try:
                        self._vprint(f'  Reading synapses from {parquet_file} (Parquet)...', level='full')
                        source_ids = set(self.neuron_dfs[i]['bodyId'].astype(str))
                        target_ids = set(self.neuron_dfs[i+1]['bodyId'].astype(str))
                        
                        # Read parquet with filters (requires pyarrow)
                        import pyarrow.parquet as pq
                        schema = pq.read_schema(parquet_file)
                        pre_col = next((c for c in schema.names if c.startswith('pre_root_id')), None)
                        post_col = next((c for c in schema.names if c.startswith('post_root_id')), None)
                        
                        if pre_col and post_col:
                            # Check for coordinate columns
                            coord_cols = ['pre_x', 'pre_y', 'pre_z', 'post_x', 'post_y', 'post_z']
                            available_cols = schema.names
                            missing_coords = [c for c in coord_cols if c not in available_cols]
                            
                            if missing_coords:
                                self._vprint(f"  ⚠ Missing coordinate columns in Parquet: {missing_coords}", level='full')
                                # Try to find alternatives (e.g. x_pre vs pre_x)
                                alt_map = {
                                    'pre_x': ['x_pre', 'pre_pt_x'],
                                    'pre_y': ['y_pre', 'pre_pt_y'],
                                    'pre_z': ['z_pre', 'pre_pt_z'],
                                    'post_x': ['x_post', 'post_pt_x'],
                                    'post_y': ['y_post', 'post_pt_y'],
                                    'post_z': ['z_post', 'post_pt_z']
                                }
                                found_map = {}
                                for target, alts in alt_map.items():
                                    if target in available_cols:
                                        found_map[target] = target
                                    else:
                                        for alt in alts:
                                            if alt in available_cols:
                                                found_map[target] = alt
                                                break
                                
                                if len(found_map) == 6:
                                    self._vprint("  ✓ Found alternative coordinate columns", level='full')
                                    columns = list(found_map.values()) + [pre_col, post_col]
                                    df = pd.read_parquet(parquet_file, columns=columns)
                                    # Rename to standard
                                    inv_map = {v: k for k, v in found_map.items()}
                                    df = df.rename(columns=inv_map)
                                else:
                                    self._vprint("  ❌ Could not resolve all coordinate columns. Skipping.", level='full')
                                    conn_df = None
                            else:
                                columns = coord_cols + [pre_col, post_col]
                                df = pd.read_parquet(parquet_file, columns=columns)

                            if conn_df is None and 'df' in locals():
                                df[pre_col] = df[pre_col].astype(str)
                                df[post_col] = df[post_col].astype(str)
                                
                                mask = (df[pre_col].isin(source_ids)) & (df[post_col].isin(target_ids))
                                conn_df = df[mask].copy()
                                
                                if not conn_df.empty:
                                    rename_map = {
                                        'pre_x': 'x_pre', 'pre_y': 'y_pre', 'pre_z': 'z_pre',
                                        'post_x': 'x_post', 'post_y': 'y_post', 'post_z': 'z_post',
                                        pre_col: 'bodyId_pre',
                                        post_col: 'bodyId_post'
                                    }
                                    conn_df = conn_df.rename(columns=rename_map)
                                    
                                    # Check coordinate scale
                                    # If Z > 10000, assume nm and DO NOT scale
                                    if conn_df['z_pre'].max() > 10000:
                                        self._vprint('  ✓ Detected coordinates in nanometers (no scaling applied)', level='full')
                                    else:
                                        self._vprint('  ✓ Detected coordinates in voxels (scaling 4x4x40)', level='full')
                                        conn_df['x_pre'] = conn_df['x_pre'] * 4
                                        conn_df['y_pre'] = conn_df['y_pre'] * 4
                                        conn_df['z_pre'] = conn_df['z_pre'] * 40
                                        conn_df['x_post'] = conn_df['x_post'] * 4
                                        conn_df['y_post'] = conn_df['y_post'] * 4
                                        conn_df['z_post'] = conn_df['z_post'] * 40

                                    self._vprint(f'  ✓ Found {len(conn_df)} synapses in Parquet file', level='full')
                                else:
                                    self._vprint('  No matching synapses found in Parquet file', level='full')
                                    conn_df = None
                        else:
                            self._vprint("  ⚠️ Could not find root_id columns in Parquet schema", level='full')
                            conn_df = None
                    except Exception as e:
                        self._vprint(f'  ⚠️ Failed to read Parquet file: {e}', level='full')
                        conn_df = None
                else:
                    # Fallback or warning
                    self._vprint(f"  ℹ️  Synapse table not found: {parquet_file}", level='full')
                    self._vprint("     If you have the raw CSV, please ensure FAFB_file_converter has run successfully.", level='full')
                    conn_df = None

                
                # Fallback to client if local failed or returned nothing
                if conn_df is None and self.client_flywire:
                    self._vprint(f"\\n  ⚠️  Local synapse file not found for dataset '{self.dataset}'.", level='full')
                    if 'fafb' in self.dataset.lower():
                        self._vprint("  Please download the synapse table from: https://codex.flywire.ai/api/download?dataset=fafb", level='full')
                    self._vprint(f"  Save the file to: {dataset_dir}", level='full')
                    self._vprint("  Skipping synapse plotting for this layer.", level='full')
                    continue
            else:
                # Fetch from NeuPrint - use new caching strategy
                source_ids = set(self.neuron_dfs[i]['bodyId'].astype(str))
                target_ids = set(self.neuron_dfs[i+1]['bodyId'].astype(str))
                
                # Try to load from cache first
                cached_df, missing_pairs = self._load_cached_synapses(source_ids, target_ids)
                
                if not missing_pairs:
                    # All data cached
                    conn_df = cached_df
                elif cached_df is not None and len(missing_pairs) < len(source_ids) * len(target_ids):
                    # Partial cache - fetch missing and combine
                    self._vprint(f'  Fetching {len(missing_pairs)} missing neuron pairs from NeuPrint...')
                    fetched_df = fetch_synapse_connections(
                        source_criteria=source_criteria,
                        target_criteria=target_criteria,
                        min_total_weight=self.min_synapse_num,
                        synapse_criteria=self.synapse_criteria,
                        client=self.client,
                    )
                    if fetched_df is not None and not fetched_df.empty:
                        conn_df = pd.concat([cached_df, fetched_df], ignore_index=True)
                        # Save newly fetched data to cache
                        self._save_cached_synapses(fetched_df, attempted_pairs=missing_pairs)
                    else:
                        conn_df = cached_df
                        # Also save empty results for missing pairs
                        self._save_cached_synapses(None, attempted_pairs=missing_pairs)
                else:
                    # No cache - fetch all
                    conn_df = fetch_synapse_connections(
                        source_criteria=source_criteria,
                        target_criteria=target_criteria,
                        min_total_weight=self.min_synapse_num,
                        synapse_criteria=self.synapse_criteria,
                        client=self.client,
                    )
                    # Save to cache
                    all_pairs = [(s, t) for s in source_ids for t in target_ids]
                    self._save_cached_synapses(conn_df, attempted_pairs=all_pairs)
        
            if conn_df is None or conn_df.empty:
                self._vprint('  No synapses found.', level='full')
                continue

            # Check if file exists to determine mode (handle skipped layers)
            if os.path.exists(file_path):
                mode = 'a'
            else:
                mode = 'w'
                
            with pd.ExcelWriter(file_path, mode=mode, engine='openpyxl') as writer:
                conn_df.to_excel(writer, sheet_name=f'conn_df{i}_{i+1}')
            
            self._vprint('plotting...', end='', level='full')
            
            if self.synapse_mode == 'scatter' or self.backend == 'k3d':
                X = (conn_df['x_pre']+conn_df['x_post'])/2
                Y = (conn_df['y_pre']+conn_df['y_post'])/2
                Z = (conn_df['z_pre']+conn_df['z_post'])/2
                xyz_df = pd.DataFrame({'x':X, 'y':Y, 'z':Z})
                
                # Ensure coordinates are float to avoid dtype warnings during transform
                xyz_df = xyz_df.astype(float)

                # Attach colors to dataframe to preserve order during transform
                c_val = self.synapse_colors[i]
                is_color_array = False
                if isinstance(c_val, (list, np.ndarray)) and len(c_val) == len(xyz_df):
                     # Check if it's not just a single RGB tuple
                     if len(xyz_df) != 3 or (isinstance(c_val[0], (str, list, tuple, np.ndarray))):
                         xyz_df['__color'] = c_val
                         is_color_array = True
                
                # Transform synapses only if needed (skip for FAFB native coords)
                if self._needs_skeleton_transform():
                    template_info = self._get_template_info()
                    self._vprint(f'Transforming synapses of layer {i} -> {i+1}...', end='', level='full')
                    with self._suppress_output():
                        xyz_df = navis.xform_brain(xyz_df, source=template_info['source'], target=template_info['target'])
                
                # Apply FAFB tilt correction if using template mode
                # This corrects the left-right tilt in the FLYWIRE template mesh
                is_fafb = 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()
                if is_fafb and self.brain_mesh == 'template':
                    xyz_df = self._apply_fafb_tilt_correction(xyz_df)
                
                # Retrieve colors
                if is_color_array and '__color' in xyz_df.columns:
                    plot_colors = xyz_df['__color'].tolist()
                else:
                    plot_colors = self.synapse_colors[i]
                
                if self.backend == 'plotly':

                    # Create 3 layers for gradient effect (Outer -> Inner)
                    # Center: synapse_alpha, Surround: synapse_alpha/10
                    base_alpha = self.synapse_alpha
                    outer_alpha = base_alpha / 10.0
                    layers = 3
                    
                    for l in range(layers):
                        # Calculate size and alpha for this layer
                        # l=0 (Outer): Size=100%, Alpha=Low
                        # l=2 (Inner): Size=33%, Alpha=High
                        
                        # Size factor: 1.0 -> 0.33
                        size_factor = (layers - l) / layers 
                        current_size = self.synapse_size * size_factor
                        
                        # Alpha interpolation: outer_alpha -> base_alpha
                        if layers > 1:
                            t = l / (layers - 1)
                            current_alpha = outer_alpha + t * (base_alpha - outer_alpha)
                        else:
                            current_alpha = base_alpha
                            
                        # Only show legend for the inner-most layer (most representative color)
                        show_legend = (l == layers - 1)
                        
                        sp = go.Scatter3d(
                            x = xyz_df['x'],
                            y = xyz_df['y'],
                            z = xyz_df['z'],
                            mode = 'markers',
                            name = f'synapses {i} -> {i+1} ({len(conn_df)})',
                            hoverinfo = 'name',
                            hovertemplate = 'x: %{x}<br>y: %{y}<br>z: %{z}<br>name: %{fullData.name}<extra></extra>',
                            legendgroup = f'synapses {i} -> {i+1} ({len(conn_df)})',
                            showlegend = show_legend,
                            marker = dict(
                                size = current_size,
                                color = plot_colors,
                                symbol = 'circle',
                                opacity = current_alpha
                            ),
                        )
                        self.fig_3d.add_trace(sp)
                elif self.backend == 'k3d':
                    try:
                        import k3d
                        # import numpy as np # Removed to avoid UnboundLocalError
                        import matplotlib.colors as mcolors
                        
                        # Color conversion helper
                        def to_int_color(c):
                            color_int = 0xff0000 # Default red
                            try:
                                if isinstance(c, str):
                                    if not c.startswith('#'):
                                        c = mcolors.to_hex(c)
                                    color_int = int(c.replace('#', ''), 16)
                                elif isinstance(c, (tuple, list, np.ndarray)):
                                    if len(c) >= 3:
                                        if isinstance(c[0], float) and c[0] <= 1.0:
                                            r, g, b = int(c[0]*255), int(c[1]*255), int(c[2]*255)
                                        else:
                                            r, g, b = int(c[0]), int(c[1]), int(c[2])
                                        color_int = (r << 16) + (g << 8) + b
                                    elif len(c) == 1: # Handle single element array
                                        return to_int_color(c[0])
                            except Exception:
                                pass
                            return color_int

                        # Determine if we have per-point colors or single color
                        c_val = plot_colors
                        colors_to_pass = None
                        
                        # Check if c_val is a list/array of colors matching the number of points
                        # Note: A single RGB tuple (r,g,b) has len 3, but we shouldn't treat it as 3 points if len(xyz_df) != 3
                        is_array_of_colors = False
                        if isinstance(c_val, (list, np.ndarray)):
                            if len(c_val) == len(xyz_df) and len(xyz_df) > 0:
                                # It matches length, but is it a list of colors or a single RGB tuple?
                                # If len(xyz_df) == 3, it's ambiguous. Assume RGB tuple if elements are numbers.
                                first_elem = c_val[0]
                                if isinstance(first_elem, (str, list, tuple, np.ndarray)):
                                    is_array_of_colors = True
                                elif len(xyz_df) != 3: # If not 3 points, it must be array of colors
                                    is_array_of_colors = True
                                # If len is 3 and elements are numbers, assume single RGB color (default behavior)

                        if is_array_of_colors:
                            # Convert each color to int
                            colors_to_pass = [to_int_color(c) for c in c_val]
                            # k3d expects uint32 array for per-point colors
                            colors_to_pass = np.array(colors_to_pass, dtype=np.uint32)
                        else:
                            # Single color
                            colors_to_pass = to_int_color(c_val)

                        # Determine opacity: use color's alpha if explicit, otherwise use synapse_alpha
                        if self._synapse_colors_have_explicit_alpha:
                            k3d_synapse_opacity = self._extract_alpha_from_color(self.synapse_colors[i])
                        else:
                            k3d_synapse_opacity = self.synapse_alpha

                        pts = k3d.points(
                            positions=xyz_df[['x', 'y', 'z']].values.astype(np.float32),
                            point_size=float(self.synapse_size) if self.synapse_mode == 'scatter' else float(self.synapse_size)/10.0,
                            color=colors_to_pass,
                            opacity=k3d_synapse_opacity,
                            name=f'synapses {i} -> {i+1} ({len(conn_df)})'
                        )
                        self.fig_3d += pts
                    except Exception as e:
                        self._vprint(f'⚠️  k3d synapse plotting failed: {e}', level='full')
            
            elif self.synapse_mode in ['sphere', 'cone', 'tetrahedron'] and self.backend == 'plotly':
                pre_coords = conn_df[['x_pre', 'y_pre', 'z_pre']].rename(columns={'x_pre':'x', 'y_pre':'y', 'z_pre':'z'})
                post_coords = conn_df[['x_post', 'y_post', 'z_post']].rename(columns={'x_post':'x', 'y_post':'y', 'z_post':'z'})
                
                # Transform synapses only if needed (skip for FAFB native coords)
                if self._needs_skeleton_transform():
                    template_info = self._get_template_info()
                    self._vprint(f'Transforming synapses of layer {i} -> {i+1}...', end='', level='full')
                    pre_coords = navis.xform_brain(pre_coords, source=template_info['source'], target=template_info['target'])
                    post_coords = navis.xform_brain(post_coords, source=template_info['source'], target=template_info['target'])
                
                # Apply FAFB tilt correction if using template mode
                # This corrects the left-right tilt in the FLYWIRE template mesh
                is_fafb = 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()
                if is_fafb and self.brain_mesh == 'template':
                    pre_coords = self._apply_fafb_tilt_correction(pre_coords)
                    post_coords = self._apply_fafb_tilt_correction(post_coords)
                
                # Calculate sizes
                # Calculate Euclidean distance
                diff = pre_coords[['x', 'y', 'z']].values - post_coords[['x', 'y', 'z']].values
                dists = np.linalg.norm(diff, axis=1)
                
                if self.synapse_size == 'real':
                    multiplier = 1.0
                else:
                    try:
                        multiplier = float(self.synapse_size)
                    except ValueError:
                        multiplier = 1.0 # Default fallback
                
                # Adjust for template geometry in statvis.py to ensure size=1.0 matches real distance
                # Cone: template height 2 (-1 to 1) -> divide by 2
                # Sphere: template diameter 2 (radius 1) -> divide by 2
                # Tetrahedron: template height 3 (-1.5 to 1.5) -> divide by 3
                if self.synapse_mode == 'tetrahedron':
                    current_size = (dists * multiplier) / 3.0
                else:
                    # Cone and Sphere (diameter)
                    current_size = (dists * multiplier) / 2.0
                
                # Determine opacity: use color's alpha if explicit, otherwise use synapse_alpha
                if self._synapse_colors_have_explicit_alpha:
                    synapse_opacity = self._extract_alpha_from_color(self.synapse_colors[i])
                else:
                    synapse_opacity = self.synapse_alpha
                
                mesh = sv.build_synapse_mesh(
                    pre_coords, 
                    post_coords, 
                    mode=self.synapse_mode, 
                    size=current_size, 
                    color=self.synapse_colors[i], 
                    opacity=synapse_opacity,
                    name=f'synapses {i} -> {i+1} ({len(conn_df)})'
                )
                mesh.hoverinfo = 'name'
                mesh.legendgroup = f'synapses {i} -> {i+1} ({len(conn_df)})'
                mesh.hovertemplate = '<b>%{fullData.name}</b><extra></extra>'
                mesh.showlegend = False
                self.fig_3d.add_trace(mesh)
                self._append_exportable_mesh(
                    mesh,
                    color=self.synapse_colors[i],
                    alpha=synapse_opacity,
                    name=f'synapses {i} -> {i+1} ({len(conn_df)})',
                    role='synapse',
                )

                # Add dummy scatter trace for legend with opaque color for visibility
                opaque_synapse_color = self._get_opaque_color(self.synapse_colors[i])
                dummy_legend = go.Scatter3d(
                    x=[None], y=[None], z=[None],
                    mode='markers',
                    name=f'synapses {i} -> {i+1} ({len(conn_df)})',
                    legendgroup=f'synapses {i} -> {i+1} ({len(conn_df)})',
                    showlegend=True,
                    marker=dict(
                        size=10,
                        color=opaque_synapse_color,  # Use opaque color for legend visibility
                        symbol='circle'
                    )
                )
                self.fig_3d.add_trace(dummy_legend)
            self._vprint('Done', level='full')
        return 0
    
    def _get_dataset_mesh_dir(self):
        """Get dataset-specific mesh directory path.
        
        Uses cache/ folder for ROI meshes:
        - hemibrain:v1.2.1 -> cache/hemibrain_v1_2_1/meshes/
        - optic-lobe:v1.1 -> cache/optic-lobe_v1_1/meshes/
        
        References:
        - navis mesh handling: https://navis.readthedocs.io/en/latest/source/api.html#navis.Volume
        - mesh compression: use navis.Volume.to_json() with compression for storage optimization
        """
        dataset_normalized = self.dataset.replace(':', '_').replace('.', '_')
        cache_mesh_dir = os.path.join(self.script_path, 'cache', dataset_normalized, 'meshes')
        os.makedirs(cache_mesh_dir, exist_ok=True)
        return cache_mesh_dir
    
    def _roi_to_filename(self, roi_name: str) -> str:
        """Convert ROI name to a case-safe filename for storage.
        
        macOS and Windows filesystems are case-insensitive by default, so 'aL(L)' and 
        'AL(L)' would be treated as the same file. This function creates unique filenames
        by encoding lowercase letters with an underscore prefix.
        
        Parameters
        ----------
        roi_name : str
            The ROI name (e.g., 'AL(L)', 'aL(L)', 'LH(R)')
            
        Returns
        -------
        str
            Filesystem-safe filename with .json extension (e.g., 'AL(L).json', '_a_L(L).json')
            
        Notes
        -----
        Encoding scheme:
        - Lowercase letters are prefixed with '_' (underscore)
        - Example: 'aL(L)' -> '_aL(L).json' (the 'a' gets underscore prefix)
        - Example: 'AL(L)' -> 'AL(L).json' (no lowercase, no encoding)
        - This preserves readability while ensuring unique filenames on case-insensitive systems
        """
        # Always encode lowercase letters to ensure unique filenames
        if any(c.islower() for c in roi_name):
            encoded_name = ''
            for char in roi_name:
                if char.islower():
                    encoded_name += f'_{char}'
                else:
                    encoded_name += char
            return encoded_name + '.json'
        else:
            return roi_name + '.json'
    
    def _get_mesh_file_path(self, mesh_dir: str, roi_name: str) -> str:
        """Get the mesh file path for an ROI, handling case-sensitivity.
        
        Tries to find the mesh file using case-safe encoding first, then falls back
        to the legacy direct naming for backward compatibility (only for all-uppercase ROI names).
        
        Parameters
        ----------
        mesh_dir : str
            Directory containing mesh files
        roi_name : str
            ROI name to look up
            
        Returns
        -------
        str
            Path to the mesh file (may or may not exist)
            
        Notes
        -----
        On case-insensitive filesystems (macOS, Windows), 'aL(L).json' and 'AL(L).json' 
        would be treated as the same file. To avoid this:
        - ROI names with lowercase letters ONLY use encoded filenames (e.g., '_aL(L).json')
        - ROI names that are all uppercase can use legacy fallback for backward compatibility
        """
        has_lowercase = any(c.islower() for c in roi_name)
        
        # Get encoded filename path
        encoded_file = os.path.join(mesh_dir, self._roi_to_filename(roi_name))
        if os.path.exists(encoded_file):
            return encoded_file
        
        # Only use legacy fallback for all-uppercase ROI names
        # This prevents case-insensitive filesystem issues where 'aL(L).json' would match 'AL(L).json'
        if not has_lowercase:
            legacy_file = os.path.join(mesh_dir, roi_name + '.json')
            if os.path.exists(legacy_file):
                return legacy_file
        
        # If neither exists, return the encoded path for new files
        return encoded_file

    def _colors_have_explicit_alpha(self, colors) -> bool:
        """
        Check if any color in the input has an explicit alpha channel.
        
        Used to determine if user-provided colors should override the separate
        neuron_alpha/synapse_alpha settings.
        
        Parameters
        ----------
        colors : str, tuple, list
            Color input in any format
            
        Returns
        -------
        bool
            True if any color has explicit alpha, False otherwise
            
        Notes
        -----
        Explicit alpha is detected in these formats:
        - RGBA tuple: (255, 0, 0, 0.5) - 4th value
        - Hex with alpha: '#ff000080' or '#f008' - 8 or 4 chars
        - CSS rgba string: 'rgba(255, 0, 0, 0.5)'
        - Named colors with alpha modifier (not common, but supported)
        
        Colors without explicit alpha (returns False):
        - RGB tuple: (255, 0, 0) - only 3 values  
        - Standard hex: '#ff0000' or '#f00' - 6 or 3 chars
        - CSS rgb string: 'rgb(255, 0, 0)'
        - Named colors: 'red', 'blue'
        - Bokeh palettes: typically hex without alpha
        """
        def _single_color_has_alpha(color) -> bool:
            """Check if a single color value has explicit alpha."""
            if isinstance(color, str):
                color_stripped = color.strip().lower()
                
                # Check rgba() string
                if color_stripped.startswith('rgba('):
                    return True
                
                # Check hex with alpha (8 chars for #RRGGBBAA or 4 chars for #RGBA)
                if color_stripped.startswith('#'):
                    hex_part = color_stripped[1:]
                    if len(hex_part) == 8 or len(hex_part) == 4:
                        return True
                elif len(color_stripped) in [8, 4] and all(c in '0123456789abcdef' for c in color_stripped):
                    # Hex without # prefix
                    return True
                    
                return False
            
            elif isinstance(color, (tuple, list)):
                # RGBA has 4 values
                if len(color) == 4 and all(isinstance(x, (int, float)) for x in color):
                    return True
                return False
            
            return False
        
        # Handle single color
        if isinstance(colors, str):
            return _single_color_has_alpha(colors)
        
        # Handle single RGBA tuple
        if isinstance(colors, tuple):
            if len(colors) == 4 and all(isinstance(x, (int, float)) for x in colors):
                return True
            # Tuple of colors - check each
            for c in colors:
                if _single_color_has_alpha(c):
                    return True
            return False
        
        # Handle list of colors
        if isinstance(colors, list):
            # Check if it's a single RGB(A) color disguised as list
            if len(colors) in [3, 4] and all(isinstance(x, (int, float)) for x in colors):
                return len(colors) == 4
            # List of colors
            for c in colors:
                if _single_color_has_alpha(c):
                    return True
            return False
        
        return False
    
    def _extract_alpha_from_color(self, color_str: str) -> float:
        """
        Extract the alpha value from a standardized rgba color string.
        
        Parameters
        ----------
        color_str : str
            Standardized color string in 'rgba(r, g, b, a)' format
            
        Returns
        -------
        float
            Alpha value between 0.0 and 1.0
            
        Examples
        --------
        >>> self._extract_alpha_from_color('rgba(255, 0, 0, 0.5)')
        0.5
        >>> self._extract_alpha_from_color('rgba(128, 128, 128, 1.0)')
        1.0
        """
        import re
        match = re.match(r'rgba\s*\(\s*\d+\s*,\s*\d+\s*,\s*\d+\s*,\s*([\d.]+)\s*\)', color_str)
        if match:
            return float(match.group(1))
        return 1.0  # Default to fully opaque if parsing fails
    
    def _rgba_to_hex(self, color_str: str) -> str:
        """
        Convert a standardized rgba color string to hex format for navis compatibility.
        
        Parameters
        ----------
        color_str : str
            Standardized color string in 'rgba(r, g, b, a)' format
            
        Returns
        -------
        str
            Hex color string like '#ff0000'
            
        Examples
        --------
        >>> self._rgba_to_hex('rgba(255, 0, 0, 0.5)')
        '#ff0000'
        >>> self._rgba_to_hex('rgba(128, 128, 128, 1.0)')
        '#808080'
        """
        import re
        match = re.match(r'rgba\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*[\d.]+\s*\)', color_str)
        if match:
            r, g, b = int(match.group(1)), int(match.group(2)), int(match.group(3))
            return f'#{r:02x}{g:02x}{b:02x}'
        # Fallback: try to use color_to_hex from color_utils
        try:
            return color_to_hex(color_str)
        except:
            return '#808080'  # Default gray
    
    def _get_opaque_color(self, color_str: str) -> str:
        """
        Get opaque version of a color (alpha=1.0) for legend display.
        
        Used to ensure legend color patches are fully visible regardless
        of the actual alpha used for plotting.
        
        Parameters
        ----------
        color_str : str
            Standardized color string in 'rgba(r, g, b, a)' format
            
        Returns
        -------
        str
            RGBA color string with alpha=1.0
            
        Examples
        --------
        >>> self._get_opaque_color('rgba(255, 0, 0, 0.2)')
        'rgba(255, 0, 0, 1.0)'
        """
        import re
        match = re.match(r'rgba\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*[\d.]+\s*\)', color_str)
        if match:
            r, g, b = int(match.group(1)), int(match.group(2)), int(match.group(3))
            return f'rgba({r}, {g}, {b}, 1.0)'
        # Fallback: return as-is
        return color_str
    
    def _standardize_color_input(self, colors, name='colors', default_alpha=1.0):
        """
        Standardize color input to a list of rgba strings.
        
        Accepts multiple input formats and converts them to a consistent
        list of 'rgba(r, g, b, a)' strings for internal use.
        
        Parameters
        ----------
        colors : str, tuple, list
            Color input in any supported format:
            - Single color string: 'red', '#ff0000', 'rgba(255, 0, 0, 0.5)'
            - Tuple/list of colors: ['red', 'blue'], [(255, 0, 0), (0, 0, 255)]
            - Bokeh palette: bokeh.palettes.Category10[10]
        name : str
            Name of the parameter (for error messages)
        default_alpha : float
            Default alpha value if not specified in colors
            
        Returns
        -------
        list
            List of standardized rgba color strings
        """
        # Handle single color string - wrap in list
        if isinstance(colors, str):
            try:
                std_color = standardize_color(colors, default_alpha=default_alpha)
                return [std_color]
            except ValueError as e:
                self._vprint(f"⚠️  Warning: Could not parse {name} '{colors}': {e}", level='simple')
                return ['rgba(128, 128, 128, 1.0)']  # Fallback gray
        
        # Handle tuple that might be a single RGB(A) color
        if isinstance(colors, tuple):
            # Check if it's a single color (3-4 numeric values)
            if len(colors) in [3, 4] and all(isinstance(x, (int, float)) for x in colors):
                try:
                    std_color = standardize_color(colors, default_alpha=default_alpha)
                    return [std_color]
                except ValueError:
                    pass
            # Otherwise treat as list of colors
            colors = list(colors)
        
        # Handle list of colors
        if isinstance(colors, list):
            result = []
            for i, c in enumerate(colors):
                try:
                    std_color = standardize_color(c, default_alpha=default_alpha)
                    result.append(std_color)
                except ValueError as e:
                    self._vprint(f"⚠️  Warning: Could not parse {name}[{i}] '{c}': {e}", level='simple')
                    result.append('rgba(128, 128, 128, 1.0)')  # Fallback gray
            return result if result else ['rgba(128, 128, 128, 1.0)']
        
        # Fallback
        return ['rgba(128, 128, 128, 1.0)']
    
    def _standardize_mesh_color_input(self, color, default_alpha=0.1):
        """
        Standardize mesh_color input (can be single color or list of colors).
        
        Parameters
        ----------
        color : str, tuple, list
            Single color or list of colors
        default_alpha : float
            Default alpha value to use if not specified in color
            
        Returns
        -------
        str or list
            Standardized color(s) - single rgba string or list of rgba strings
        """
        # Handle string input
        if isinstance(color, str):
            try:
                return standardize_color(color, default_alpha=default_alpha)
            except ValueError:
                return f'rgba(100, 100, 100, {default_alpha})'  # Default gray
        
        # Handle tuple input
        if isinstance(color, tuple):
            # Check if it's a single RGBA color
            if len(color) >= 3 and all(isinstance(x, (int, float)) for x in color):
                try:
                    # Use explicit alpha from tuple if present, otherwise use default_alpha
                    alpha = color[3] if len(color) > 3 else default_alpha
                    return standardize_color(color[:min(len(color), 4)], default_alpha=alpha)
                except ValueError:
                    return f'rgba(100, 100, 100, {default_alpha})'
            # Tuple of colors
            return [self._standardize_mesh_color_input(c, default_alpha=default_alpha) for c in color]
        
        # Handle list input (list of colors for multiple ROIs)
        if isinstance(color, list):
            # Check if it might be a single RGB color (list of 3-4 numbers)
            if len(color) in [3, 4] and all(isinstance(x, (int, float)) for x in color):
                try:
                    alpha = color[3] if len(color) > 3 else default_alpha
                    return standardize_color(color[:min(len(color), 4)], default_alpha=alpha)
                except ValueError:
                    return f'rgba(100, 100, 100, {default_alpha})'
            # List of colors
            return [self._standardize_mesh_color_input(c, default_alpha=default_alpha) for c in color]
        
        # Fallback
        return f'rgba(100, 100, 100, {default_alpha})'

    def _is_custom_mesh_color_specified(self) -> bool:
        """
        Check if the user specified a custom mesh_color (not the default).
        
        Returns
        -------
        bool
            True if mesh_color was customized from default, False otherwise
        """
        default_color = (100, 100, 100)  # Default RGB without alpha
        
        # Get the original mesh_color before standardization
        # We check if it differs from the default gray
        if hasattr(self, '_original_mesh_color'):
            original = self._original_mesh_color
        else:
            original = self.mesh_color
        
        # If it's the default tuple, it's not custom
        if isinstance(original, tuple) and len(original) in [3, 4]:
            if original[:3] == default_color:
                return False
        
        # If it's a standardized rgba string matching default, it's not custom
        if isinstance(original, str):
            if 'rgba(100, 100, 100,' in original or 'rgb(100, 100, 100)' in original:
                return False
        
        # Otherwise, it's custom
        return True

    def _darken_color(self, color, brightness):
        """Darken a color by the given brightness factor.
        
        Parameters
        ----------
        color : str or tuple
            Color in any supported format (now uses standardize_color internally):
            - Hex string: '#RRGGBB', '#RGB'
            - RGB tuple: (r, g, b)
            - RGBA tuple: (r, g, b, a)
            - CSS string: 'rgb(r, g, b)', 'rgba(r, g, b, a)'
            - Named color: 'red', 'blue', etc.
        brightness : float
            Brightness factor from 0.0 (black) to 1.0 (original color)
        
        Returns
        -------
        str
            Darkened color in rgba format
        """
        # Use the standardize_color utility to parse any format
        try:
            rgba = extract_rgba_tuple(color)
            r, g, b, a = rgba
        except ValueError:
            # Fallback to manual parsing for legacy support
            if isinstance(color, str):
                if color.startswith('#'):
                    hex_str = color.lstrip('#')
                    if len(hex_str) == 6:
                        r = int(hex_str[0:2], 16)
                        g = int(hex_str[2:4], 16)
                        b = int(hex_str[4:6], 16)
                        a = 1.0
                    else:
                        return 'rgba(128, 128, 128, 1.0)'
                else:
                    return 'rgba(128, 128, 128, 1.0)'
            elif isinstance(color, (tuple, list)) and len(color) >= 3:
                r, g, b = int(color[0]), int(color[1]), int(color[2])
                a = color[3] if len(color) > 3 else 1.0
            else:
                return 'rgba(128, 128, 128, 1.0)'
        
        # Apply brightness (darken)
        r = int(r * brightness)
        g = int(g * brightness)
        b = int(b * brightness)
        
        # Clamp values
        r = max(0, min(255, r))
        g = max(0, min(255, g))
        b = max(0, min(255, b))
        
        return f'rgba({r}, {g}, {b}, {a})'
    
    def _interpolate_colors(self, base_colors, n_needed):
        """Generate n_needed colors by keeping original colors first, then appending interpolated ones.
        
        Uses a round-based approach to generate new colors:
        - Round 1: Generate midpoints between consecutive pairs (0-1, 1-2, ..., 8-9, 9-0)
          e.g., 10 colors → 10 new midpoints → 20 total candidates
        - Round 2: If more needed, generate midpoints between all consecutive pairs in expanded set
        - Continue until enough candidates are generated
        
        Final result: [original colors] + [interpolated colors as needed]
        
        Parameters
        ----------
        base_colors : tuple or list
            Base color palette to preserve. Accepts any format supported by
            standardize_color: hex strings, RGB tuples, RGBA tuples, named colors, rgba strings
        n_needed : int
            Number of colors needed
        
        Returns
        -------
        list
            List of n_needed rgba color strings, with original colors first, then interpolated
        """
        import numpy as np
        
        # Convert base colors to normalized RGB tuples using our standardize_color utility
        rgb_colors = []
        alphas = []
        for c in base_colors:
            try:
                rgba = extract_rgba_tuple(c)
                # Normalize to 0-1 for interpolation
                rgb_colors.append(np.array([rgba[0]/255, rgba[1]/255, rgba[2]/255]))
                alphas.append(rgba[3])
            except ValueError:
                rgb_colors.append(np.array([0.5, 0.5, 0.5]))  # fallback gray
                alphas.append(1.0)
        
        # Use average alpha from base colors
        avg_alpha = sum(alphas) / len(alphas) if alphas else 1.0
        
        n_base = len(rgb_colors)
        
        # Handle single color case - just repeat the color n_needed times
        if n_base == 1:
            r, g, b = int(rgb_colors[0][0]*255), int(rgb_colors[0][1]*255), int(rgb_colors[0][2]*255)
            return [f'rgba({r}, {g}, {b}, {avg_alpha})' for _ in range(n_needed)]
        
        # If we need fewer or equal colors than base, just use the base colors
        if n_needed <= n_base:
            result = []
            for i in range(n_needed):
                r, g, b = int(rgb_colors[i][0]*255), int(rgb_colors[i][1]*255), int(rgb_colors[i][2]*255)
                result.append(f'rgba({r}, {g}, {b}, {avg_alpha})')
            return result
        
        # Start with original colors, then generate interpolated colors to append
        # Generate interpolated colors using round-based midpoint approach
        interpolated_pool = []  # Pool of interpolated colors to draw from
        current_set = rgb_colors.copy()  # Colors to generate midpoints from
        
        while len(rgb_colors) + len(interpolated_pool) < n_needed:
            # Generate midpoints between consecutive pairs (with wrap-around)
            new_midpoints = []
            n_current = len(current_set)
            for i in range(n_current):
                next_idx = (i + 1) % n_current
                midpoint = (current_set[i] + current_set[next_idx]) / 2
                new_midpoints.append(midpoint)
            
            # Add new midpoints to the pool
            interpolated_pool.extend(new_midpoints)
            
            # For next round, use interleaved set (original + midpoints) as the new current_set
            # This creates finer gradations in the next round
            interleaved = []
            for i in range(n_current):
                interleaved.append(current_set[i])
                interleaved.append(new_midpoints[i])
            current_set = interleaved
        
        # Build final result: original colors first, then interpolated as needed
        result = []
        
        # Add all original colors first
        for c in rgb_colors:
            r, g, b = int(c[0]*255), int(c[1]*255), int(c[2]*255)
            result.append(f'rgba({r}, {g}, {b}, {avg_alpha})')
        
        # Add interpolated colors until we reach n_needed
        n_extra_needed = n_needed - n_base
        for i in range(n_extra_needed):
            c = interpolated_pool[i]
            r, g, b = int(c[0]*255), int(c[1]*255), int(c[2]*255)
            result.append(f'rgba({r}, {g}, {b}, {avg_alpha})')
        
        return result

    def _expand_color_sequence(self, colors, n_needed, target_label='items', tip_parameter=None, warn=True):
        """Expand a color sequence to the requested size using the configured policy."""
        colors = tuple(colors)
        n_colors = len(colors)

        if n_needed <= 0:
            return tuple()
        if n_needed <= n_colors:
            return tuple(colors[:n_needed])

        if self.expand_colors == 'interpolation':
            expanded = self._interpolate_colors(colors, n_needed)
            warning = f'\033[33m⚠️  Warning: {n_needed} {target_label} but only {n_colors} colors. Generated {n_needed} colors via interpolation.\033[0m'
        elif self.expand_colors == 'darken':
            n_cycles = (n_needed - 1) // n_colors
            expanded = []
            for i in range(n_needed):
                cycle_num = i // n_colors
                color_idx = i % n_colors
                base_color = colors[color_idx]
                brightness = 1.0 - (0.3 * cycle_num / n_cycles) if n_cycles > 0 else 1.0
                expanded.append(self._darken_color(base_color, brightness))
            warning = f'\033[33m⚠️  Warning: {n_needed} {target_label} but only {n_colors} colors. Recycling with darkening (100%→70%).\033[0m'
        else:
            expanded = [colors[i % n_colors] for i in range(n_needed)]
            warning = f'\033[33m⚠️  Warning: {n_needed} {target_label} but only {n_colors} colors. Cycling colors (repeating pattern).\033[0m'

        if warn:
            self._vprint(warning)
            if tip_parameter:
                self._vprint(f'\033[33m   💡 Tip: Use {tip_parameter} parameter with custom palette to specify more colors.\033[0m')

        return tuple(expanded)

    def _normalize_neuron_lookup_keys(self, value):
        """Normalize neuron identifiers for robust lookup by bodyId or name-like fields."""
        if value is None:
            return []

        if isinstance(value, (int, np.integer)):
            int_value = int(value)
            return [int_value, str(int_value)]

        try:
            if pd.isna(value):
                return []
        except TypeError:
            pass

        value_str = str(value).strip()
        if not value_str or value_str.lower() == 'nan':
            return []

        keys = [value_str]
        if value_str.isdigit():
            keys.insert(0, int(value_str))
        return keys

    def _build_per_neuron_color_map(self):
        """Build a lookup of neuron identifiers to per-neuron colors."""
        total_neurons = sum(len(ndf) for ndf in self.neuron_dfs if ndf is not None)
        if total_neurons == 0:
            return {}

        neuron_palette = self._expand_color_sequence(
            self._base_neuron_colors,
            total_neurons,
            target_label='neurons',
            tip_parameter='neuron_colors',
        )

        neuron_color_map = {}
        color_index = 0
        for ndf in self.neuron_dfs:
            if ndf is None or ndf.empty:
                continue
            for _, row in ndf.iterrows():
                neuron_color = neuron_palette[color_index]
                for column in ['bodyId', 'name', 'instance', 'roi']:
                    if column in row.index:
                        for key in self._normalize_neuron_lookup_keys(row[column]):
                            neuron_color_map[key] = neuron_color
                color_index += 1

        return neuron_color_map

    def _resolve_neuron_color(self, neuron_id, layer_index):
        """Resolve the display color for a neuron trace, including overrides."""
        neuron_color = self.neuron_colors[layer_index]
        lookup_keys = self._normalize_neuron_lookup_keys(neuron_id)

        if self.color_mode == 'per_neuron' and hasattr(self, '_per_neuron_colors'):
            for key in lookup_keys:
                if key in self._per_neuron_colors:
                    neuron_color = self._per_neuron_colors[key]
                    break

        if hasattr(self, '_neuron_color_overrides') and self._neuron_color_overrides:
            for key in lookup_keys:
                if key in self._neuron_color_overrides:
                    return self._neuron_color_overrides[key]

        return neuron_color

    def _apply_plotly_trace_color(self, trace, neuron_color):
        """Apply a resolved neuron color to a Plotly trace."""
        color_hex = self._rgba_to_hex(neuron_color)
        color_alpha = self._extract_alpha_from_color(neuron_color)

        if hasattr(trace, 'color'):
            try:
                trace.color = color_hex
            except Exception:
                pass
        if hasattr(trace, 'line') and trace.line is not None:
            trace.line.color = color_hex
        if hasattr(trace, 'marker') and trace.marker is not None:
            trace.marker.color = color_hex
        trace.opacity = color_alpha

    def _apply_k3d_object_color(self, obj, neuron_color):
        """Apply a resolved neuron color to a k3d object when supported."""
        try:
            r, g, b, a = extract_rgba_tuple(neuron_color)
            if hasattr(obj, 'color'):
                obj.color = (int(r) << 16) + (int(g) << 8) + int(b)
            if hasattr(obj, 'opacity'):
                obj.opacity = a
        except Exception:
            pass

    def _flatten_nested_roi_groups(self, roi_list, color_input):
        """Flatten nested ROI lists while assigning same color to grouped ROIs.
        
        Supports nested lists for grouping ROIs that should share the same color:
        e.g., ['AME', ['aL', 'bL', 'gL'], 'EB'] -> 
              flat list with same color for all ROIs in the nested list.
        
        Parameters
        ----------
        roi_list : list
            List of ROI names, potentially with nested lists for grouping.
            e.g., ['AME', ['aL', 'bL', 'gL', "a'L", "b'L"], 'EB', 'PB']
        color_input : str, list, or tuple
            Colors for each ROI group. Can be a single color (applies to all),
            or a list matching the number of top-level items in roi_list.
            
        Returns
        -------
        tuple
            (flattened_rois, expanded_colors, nested_groups)
            - flattened_rois: Flat list of all ROI names
            - expanded_colors: List of colors matching flattened_rois
            - nested_groups: Dict recording the source color group for each ROI
            
        Examples
        --------
        >>> _flatten_nested_roi_groups(['AME', ['aL', 'bL'], 'EB'], ['red', 'green', 'blue'])
        (['AME', 'aL', 'bL', 'EB'], ['red', 'green', 'green', 'blue'], 
         {'AME': 'AME', 'aL': 'MB_lobes', 'bL': 'MB_lobes', 'EB': 'EB'})
        """
        if not roi_list:
            return roi_list, color_input, {}
        
        # Normalize color_input to a list matching roi_list length
        if isinstance(color_input, str):
            colors = [color_input] * len(roi_list)
        elif isinstance(color_input, (tuple, list)):
            # Check if it's a single color tuple (RGB/RGBA with numeric values) or a list/tuple of colors
            # RGB/RGBA tuple: (255, 0, 0) or (1.0, 0.5, 0.0, 0.5)
            # List of colors: ['red', 'blue'] or ('#ff0000', '#0000ff') or [(255,0,0), (0,0,255)]
            is_single_rgb_color = (
                len(color_input) in [3, 4] and 
                all(isinstance(x, (int, float)) for x in color_input)
            )
            
            if is_single_rgb_color:
                # Single RGB/RGBA color tuple - apply to all ROIs
                colors = [color_input] * len(roi_list)
            else:
                # List/tuple of colors (strings, hex codes, or color tuples)
                colors = list(color_input)
                # Extend if needed
                while len(colors) < len(roi_list):
                    colors.append(colors[-1] if colors else 'gray')
        else:
            colors = ['gray'] * len(roi_list)
        
        flattened_rois = []
        expanded_colors = []
        nested_groups = {}
        
        for i, item in enumerate(roi_list):
            color = colors[i] if i < len(colors) else colors[-1]
            
            if isinstance(item, list):
                # Nested list - all items share the same color. Keep metadata
                # about the source group without coupling their legend entries.
                group_label = '+'.join(str(r) for r in item[:3])  # Use first 3 names
                if len(item) > 3:
                    group_label += f'+{len(item)-3}more'
                
                for roi in item:
                    flattened_rois.append(roi)
                    expanded_colors.append(color)
                    nested_groups[roi] = group_label
                    
                self._vprint(f"   🔗 Grouped [{', '.join(str(r) for r in item)}] → same color", level='simple')
            else:
                # Single ROI
                flattened_rois.append(item)
                expanded_colors.append(color)
                nested_groups[item] = str(item)
        
        return flattened_rois, expanded_colors, nested_groups
    
    def _expand_roi_names_with_colors(self, roi_list, color_list, available_rois=None):
        """Expand ROI names AND their corresponding colors to include bilateral (L/R) variants.
        
        When a user specifies 'LH' with color 'red', this function will automatically expand
        to ['LH(L)', 'LH(R)'] with colors ['red', 'red']. This ensures colors match expanded ROIs.
        
        Parameters
        ----------
        roi_list : list
            List of ROI names to expand
        color_list : list or single color
            Colors corresponding to each ROI. If a single color, applies to all.
            If a list shorter than roi_list, extra ROIs get the last color.
        available_rois : list, optional
            List of available ROI names. If None, will be fetched from cache/API.
            
        Returns
        -------
        tuple
            (expanded_rois, expanded_colors) - Both lists with matching lengths
            
        Examples
        --------
        >>> _expand_roi_names_with_colors(['LH', 'EB'], ['red', 'blue'])
        (['LH(L)', 'LH(R)', 'EB'], ['red', 'red', 'blue'])
        
        >>> _expand_roi_names_with_colors(['LH'], 'green')  # Single color
        (['LH(L)', 'LH(R)'], ['green', 'green'])
        """
        if not roi_list:
            return roi_list, color_list
        
        # Normalize color_list to a list
        if not isinstance(color_list, list):
            color_list = [color_list] * len(roi_list)
        
        # Ensure color_list matches roi_list length
        if len(color_list) < len(roi_list):
            # Extend with the last color
            last_color = color_list[-1] if color_list else None
            color_list = list(color_list) + [last_color] * (len(roi_list) - len(color_list))
        
        # Get available ROIs if not provided
        if available_rois is None:
            is_fafb = 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()
            if is_fafb:
                malecns_cache = os.path.join(self.script_path, 'cache', 'male-cns_v0_9', 'available_rois.json')
                if os.path.exists(malecns_cache):
                    import json
                    with open(malecns_cache, 'r') as f:
                        available_rois = json.load(f)
                else:
                    available_rois = self._get_available_rois(use_cache=True, fetch_online=False)
            else:
                available_rois = self._get_available_rois(use_cache=True, fetch_online=False)
        
        available_set = set(available_rois) if available_rois else set()
        
        expanded_rois = []
        expanded_colors = []
        seen = set()
        
        for i, roi in enumerate(roi_list):
            color = color_list[i] if i < len(color_list) else color_list[-1]
            
            # Check if ROI already has (L) or (R) suffix
            if roi.endswith('(L)') or roi.endswith('(R)'):
                if roi not in seen:
                    expanded_rois.append(roi)
                    expanded_colors.append(color)
                    seen.add(roi)
                continue
            
            # Check if the ROI exists as-is (like 'EB' which is unpaired)
            if roi in available_set:
                if roi not in seen:
                    expanded_rois.append(roi)
                    expanded_colors.append(color)
                    seen.add(roi)
                continue
            
            # Try to expand to bilateral variants
            left_variant = f'{roi}(L)'
            right_variant = f'{roi}(R)'
            
            found_left = left_variant in available_set
            found_right = right_variant in available_set
            
            if found_left and found_right:
                # Both sides exist, expand to both with same color
                if left_variant not in seen:
                    expanded_rois.append(left_variant)
                    expanded_colors.append(color)
                    seen.add(left_variant)
                if right_variant not in seen:
                    expanded_rois.append(right_variant)
                    expanded_colors.append(color)
                    seen.add(right_variant)
            elif found_left:
                if left_variant not in seen:
                    expanded_rois.append(left_variant)
                    expanded_colors.append(color)
                    seen.add(left_variant)
            elif found_right:
                if right_variant not in seen:
                    expanded_rois.append(right_variant)
                    expanded_colors.append(color)
                    seen.add(right_variant)
            else:
                # No bilateral variants found, keep original
                if roi not in seen:
                    expanded_rois.append(roi)
                    expanded_colors.append(color)
                    seen.add(roi)
        
        return expanded_rois, expanded_colors

    def _expand_mesh_roi_patterns(self, roi_list):
        """Expand special keywords and regex patterns in mesh_roi list.
        
        Supports:
        - 'primary': All primary ROIs (major brain regions)
        - 'all': All available ROIs for the current dataset
        - Regex patterns: 'ME.*' matches all ROIs starting with 'ME'
        
        Parameters
        ----------
        roi_list : list
            List of ROI names, keywords, or regex patterns
            
        Returns
        -------
        list
            Expanded list of ROI names
            
        Examples
        --------
        >>> _expand_mesh_roi_patterns(['primary'])
        ['AL(L)', 'AL(R)', 'EB', 'FB', ...]  # All primary ROIs
        
        >>> _expand_mesh_roi_patterns(['ME.*'])
        ['ME(L)', 'ME(R)', 'ME_glomerulus(L)', ...]  # All ROIs matching ME.*
        
        >>> _expand_mesh_roi_patterns(['LH', 'ME.*', 'EB'])
        ['LH', 'ME(L)', 'ME(R)', ..., 'EB']  # Mixed: literal + regex
        """
        import re
        
        if not roi_list:
            return roi_list
        
        # Convert single string to list
        if isinstance(roi_list, str):
            roi_list = [roi_list]
        
        # Get available ROIs for pattern matching
        available_rois = self._get_available_rois(use_cache=True, fetch_online=True)
        available_set = set(available_rois) if available_rois else set()
        
        # Define primary ROIs - major brain regions that are commonly used
        # These are regions that typically exist across datasets
        primary_roi_patterns = [
            # Central Complex
            'EB', 'FB', 'PB', 'NO', 'AB',
            # Mushroom Body  
            'MB.*', 'CA.*', 'PED.*',
            # Antennal Lobe
            'AL\\(.*\\)',
            # Lateral Horn
            'LH\\(.*\\)',
            # Optic Lobe
            'ME\\(.*\\)', 'LO\\(.*\\)', 'LOP\\(.*\\)', 'AME\\(.*\\)',
            # Subesophageal Zone
            'SEZ.*', 'GNG.*',
            # Other major regions
            'CRE.*', 'SCL.*', 'ICL.*', 'IB.*', 'ATL.*', 'AVLP.*', 'PVLP.*',
            'PLP.*', 'WED.*', 'SLP.*', 'SIP.*', 'SMP.*', 'CAN.*',
            'FLA.*', 'EPA.*', 'GOR.*', 'SPS.*', 'IPS.*',
        ]
        
        expanded = []
        seen = set()
        
        for item in roi_list:
            item_str = str(item).strip()
            
            # Handle special keywords
            if item_str.lower() == 'all':
                # Add all available ROIs
                for roi in sorted(available_rois):
                    if roi not in seen:
                        expanded.append(roi)
                        seen.add(roi)
                self._vprint(f"   🌐 'all': Added {len(available_rois)} ROIs", level='simple')
                continue
                
            elif item_str.lower() == 'primary':
                # Add ROIs matching primary patterns
                primary_count = 0
                for pattern in primary_roi_patterns:
                    try:
                        regex = re.compile(f'^{pattern}$')
                        for roi in available_rois:
                            if regex.match(roi) and roi not in seen:
                                expanded.append(roi)
                                seen.add(roi)
                                primary_count += 1
                    except re.error:
                        # If pattern is invalid regex, try exact match
                        if pattern in available_set and pattern not in seen:
                            expanded.append(pattern)
                            seen.add(pattern)
                            primary_count += 1
                self._vprint(f"   🏛️ 'primary': Added {primary_count} primary ROIs", level='simple')
                continue
            
            # Check if it's a regex pattern (contains regex special chars)
            # Common regex patterns: .* .+ [abc] ^ $ etc.
            is_regex = any(c in item_str for c in ['*', '+', '?', '[', ']', '^', '$', '|', '\\'])
            
            if is_regex:
                # Treat as regex pattern
                try:
                    # Anchor the pattern if not already anchored
                    pattern = item_str
                    if not pattern.startswith('^'):
                        pattern = '^' + pattern
                    if not pattern.endswith('$'):
                        pattern = pattern + '$'
                    
                    regex = re.compile(pattern)
                    matched = []
                    for roi in available_rois:
                        if regex.match(roi) and roi not in seen:
                            expanded.append(roi)
                            seen.add(roi)
                            matched.append(roi)
                    
                    if matched:
                        self._vprint(f"   🔍 '{item_str}': Matched {len(matched)} ROIs", level='simple')
                    else:
                        self._vprint(f"   ⚠️ '{item_str}': No matching ROIs found", level='simple')
                        
                except re.error as e:
                    self._vprint(f"   ⚠️ Invalid regex '{item_str}': {e}, treating as literal", level='simple')
                    if item_str not in seen:
                        expanded.append(item_str)
                        seen.add(item_str)
            else:
                # Literal ROI name
                if item_str not in seen:
                    expanded.append(item_str)
                    seen.add(item_str)
        
        return expanded

    def _get_available_rois(self, use_cache=True, fetch_online=True):
        """Query NeuPrint database for available ROIs in the current dataset.
        
        Caches results locally to avoid repeated API calls. Returns a list of ROI names
        that are available in the NeuPrint database for the current dataset.
        
        Parameters
        ----------
        use_cache : bool
            If True, use cached ROI list if available. If False, force refresh from API.
        fetch_online : bool
            If True, attempt to fetch from NeuPrint online. If False, only use local cache/meshes.
        
        Returns
        -------
        list
            List of available ROI names for the current dataset.
        
        References:
        - NeuPrint ROI documentation: https://neuprint.janelia.org/
        - navis neuprint interface: https://navis-org.github.io/navis/reference/navis/interfaces/neuprint/
        - neuprint-python API: https://github.com/connectome-neuprint/neuprint-python
        """
        # Cache file path in organized cache/ structure
        dataset_normalized = self.dataset.replace(':', '_').replace('.', '_')
        cache_dir = os.path.join(self.script_path, 'cache', dataset_normalized)
        cache_file = os.path.join(cache_dir, 'available_rois.json')
        
        # Try to load from cache first
        if use_cache and os.path.exists(cache_file):
            try:
                import json
                with open(cache_file, 'r') as f:
                    cached_data = json.load(f)
                    self._vprint(f'✓ Loaded {len(cached_data)} available ROIs from cache', level='full')
                    return cached_data
            except Exception as e:
                self._vprint(f'⚠️ Failed to load ROI cache: {e}, fetching from API...', level='full')
        
        # Fetch from NeuPrint API
        if fetch_online:
            # Special handling for FlyWire/FAFB: Do not use API, use local primary_rois or hemibrain cache
            if 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower():
                self._vprint('ℹ️  FlyWire/FAFB dataset detected: Skipping online API fetch for ROIs.', level='full')
                self._vprint('   Scanning local ROI meshes...', level='full')
                
                found_rois = set()
                
                # Scan primary_rois
                primary_dir = os.path.join(self.script_path, 'navis_roi_meshes_json', 'primary_rois')
                if os.path.exists(primary_dir):
                    for f in os.listdir(primary_dir):
                        if f.endswith('.json'):
                            found_rois.add(f[:-5])
                            
                # Scan hemibrain cache
                hb_cache = os.path.join(self.script_path, 'cache', 'hemibrain_v1_2_1', 'meshes')
                if os.path.exists(hb_cache):
                    for f in os.listdir(hb_cache):
                        if f.endswith('.json'):
                            found_rois.add(f[:-5])
                            
                roi_list = sorted(list(found_rois))
                self._vprint(f'✓ Found {len(roi_list)} available ROIs from local storage', level='full')
                
                # Cache the results
                if roi_list:
                    try:
                        import json
                        os.makedirs(cache_dir, exist_ok=True)
                        with open(cache_file, 'w') as f:
                            json.dump(roi_list, f, indent=2)
                    except Exception as e:
                        self._vprint(f'⚠️ Failed to cache ROI list: {e}', level='full')
                        
                return roi_list

            try:
                self._vprint('📥 Fetching available ROIs from NeuPrint online database...', level='full')
                
                # Initialize neuprint client using environment variable or global client
                from neuprint import Client, fetch_meta
                
                client = self.client
                
                if client is None:
                    # Try to get token from environment variable first
                    token = os.environ.get('NEUPRINT_APPLICATION_CREDENTIALS')
                    
                    if token:
                        # Determine server URL based on dataset
                        if 'optic' in self.dataset.lower():
                            server = 'https://neuprint-optic-lobe.janelia.org'
                            dataset_name = self.dataset.split(':')[0]  # 'optic-lobe'
                        else:
                            server = 'https://neuprint.janelia.org'
                            dataset_name = 'hemibrain:v1.2.1'  # default
                        
                        try:
                            client = Client(server, dataset=dataset_name, token=token)
                        except Exception as e:
                            self._vprint(f'   Warning: Failed to create client with token: {e}', level='full')
                            self._vprint(f'   Attempting to use default/global client...', level='full')
                            client = None
                
                # Fetch metadata (will use client if provided, otherwise global)
                meta = fetch_meta(client=client)
                
                roi_list = []
                # Extract ROI list from meta info
                if 'roiInfo' in meta:
                    roi_list = list(meta['roiInfo'].keys())
                    self._vprint(f'   Found {len(roi_list)} ROIs from roiInfo', level='full')
                elif 'primaryRois' in meta:
                    roi_list = list(meta['primaryRois'])
                    self._vprint(f'   Found {len(roi_list)} primary ROIs', level='full')
                else:
                    self._vprint(f'   Warning: No roiInfo/primaryRois in metadata, falling back to local cache', level='full')
                
                roi_list = sorted(roi_list)
                
                # Cache the results (create directory only when needed)
                if roi_list:
                    try:
                        import json
                        os.makedirs(cache_dir, exist_ok=True)
                        with open(cache_file, 'w') as f:
                            json.dump(roi_list, f, indent=2)
                        self._vprint(f'✓ Cached {len(roi_list)} available ROIs to {cache_file}', level='full')
                    except Exception as e:
                        self._vprint(f'⚠️ Failed to cache ROI list: {e}', level='full')
                
                return roi_list
                
            except Exception as e:
                self._vprint(f'⚠️ Failed to fetch available ROIs from NeuPrint: {e}', level='full')
                self._vprint(f'   Tip: Set NEUPRINT_APPLICATION_CREDENTIALS environment variable', level='full')
                self._vprint(f'   Using ROIs from local mesh directory instead.', level='full')
        
        # Fallback: list available meshes from local directory
        mesh_dir = self._get_dataset_mesh_dir()
        if os.path.exists(mesh_dir):
            roi_list = [f.replace('.json', '') for f in os.listdir(mesh_dir) if f.endswith('.json')]
            roi_list = sorted(roi_list)
            self._vprint(f'✓ Found {len(roi_list)} ROIs in local cache: {mesh_dir}', level='full')
            
            # Cache the results from local scan
            if roi_list:
                try:
                    import json
                    os.makedirs(cache_dir, exist_ok=True)
                    with open(cache_file, 'w') as f:
                        json.dump(roi_list, f, indent=2)
                    self._vprint(f'✓ Cached {len(roi_list)} available ROIs to {cache_file}', level='full')
                except Exception as e:
                    self._vprint(f'⚠️ Failed to cache ROI list: {e}', level='full')
            
            return roi_list
        else:
            self._vprint(f'⚠️ No ROI data available (online fetch failed and no local cache)', level='full')
            return []
    
    def _check_transform_requirements_early(self):
        """Check transform requirements at startup and advise user on options.
        
        This method checks if the current dataset and brain_mesh settings require
        coordinate transforms, and if so, verifies transform availability and
        advises the user on transformation-free alternatives.
        
        For hemibrain or FAFB with brain_mesh='whole', H5 transforms are required. 
        The user will be prompted to either:
        1. Download transforms (~13GB, enables JRC2018F whole brain view)
        2. Use brain_mesh='template' for transformation-free native view
        3. Use brain_mesh='none' for no brain mesh
        
        For all other datasets, native templates are used without transforms.
        """
        dataset_lower = self.dataset.lower()
        
        # Check if transforms are needed
        needs_transform = self._dataset_needs_transform()
        
        if not needs_transform:
            # Dataset uses native template - no transforms needed
            if self.brain_mesh in ['template', 'whole']:
                template_info = self._get_template_info()
                if template_info.get('skip_transform', False):
                    self._vprint(f'✓ Using native {template_info["mesh_name"]} - no coordinate transforms needed', level='full')
            return
        
        # Only hemibrain and FAFB with brain_mesh='whole' require H5 transforms
        is_hemibrain = 'hemibrain' in dataset_lower
        is_fafb = 'flywire' in dataset_lower or 'fafb' in dataset_lower
        
        if not (is_hemibrain or is_fafb):
            return
            
        # brain_mesh='whole' needs transform check
        if self.brain_mesh == 'none':
            return
            
        # Check if transforms are available
        import flybrains
        import navis
        
        template_info = self._get_template_info()
        source = template_info['source']
        target = template_info['target']
        
        # Try to find the transform path
        transforms_available = False
        try:
            path = navis.transforms.registry.find_bridging_path(source, target)
            transforms_available = True
        except (ValueError, KeyError):
            pass
        
        if transforms_available:
            if self.brain_mesh == 'whole':
                self._vprint(f'✓ Transforms available for {source} → {target}', level='full')
            return
            
        # Transforms not available - prompt user
        YELLOW = '\\033[93m'
        CYAN = '\\033[96m'
        GREEN = '\\033[92m'
        RED = '\\033[91m'
        RESET = '\\033[0m'
        
        dataset_name = 'FlyWire/FAFB' if is_fafb else 'Hemibrain'
        native_template = 'FLYWIRE' if is_fafb else 'JRCFIB2018F'
        
        print(f'\\n{YELLOW}{"="*70}')
        print(f'⚠️  Coordinate Transform Required for {dataset_name}')
        print(f'{"="*70}{RESET}')
        print()
        print(f'Your settings: dataset={self.dataset}, brain_mesh={self.brain_mesh}')
        print()
        print(f'The {dataset_name} dataset requires coordinate transforms for brain_mesh="whole":')
        print(f'  • Transform path: {source} → {target}')
        print(f'  • Requires downloading ~{CYAN}13 GB{RESET} of transform files')
        print(f'  • Transformation adds processing time to visualization')
        print()
        print(f'{GREEN}💡 Transformation-Free Alternatives:{RESET}')
        print()
        print(f'  Option 1: Use brain_mesh="template" (recommended)')
        print(f'            → Uses native {native_template} template mesh')
        print(f'            → No transforms needed, fast visualization')
        print()
        print(f'  Option 2: Use brain_mesh="none"')
        print(f'            → No brain mesh, only neurons and synapses')
        print()
        print(f'  Option 3: Use a different dataset')
        if is_fafb:
            print(f'            → male-cns: Native JRCFIB2022M template (male CNS)')
            print(f'            → hemibrain: Native JRCFIB2018F template')
        else:
            print(f'            → FlyWire/FAFB: Native FLYWIRE template (female brain)')
            print(f'            → male-cns: Native JRCFIB2022M template (male CNS)')
        print(f'            → These datasets have no transform requirements with brain_mesh="template"')
        print()
        print(f'{"="*70}')
        
        # Prompt user for choice
        print(f'\\nHow would you like to proceed?')
        print(f'  [1] Download transforms (~13GB) and continue with brain_mesh="whole"')
        print(f'  [2] Use brain_mesh="template" instead (no download, fast)')
        print(f'  [3] Use brain_mesh="none" (no brain mesh)')
        print(f'  [q] Quit')
        print()
        
        try:
            choice = input('Enter choice [1/2/3/q] (default: 2): ').strip().lower()
        except (EOFError, KeyboardInterrupt):
            choice = '2'
        
        if choice == '' or choice == '2':
            print(f'\\n{GREEN}✓ Using brain_mesh="template" (transformation-free){RESET}')
            self.brain_mesh = 'template'
        elif choice == '3':
            print(f'\\n{GREEN}✓ Using brain_mesh="none"{RESET}')
            self.brain_mesh = 'none'
        elif choice == '1':
            # Try to download transforms
            if self._check_and_download_transforms():
                print(f'\\n{GREEN}✓ Transforms downloaded successfully{RESET}')
            else:
                print(f'\\n{YELLOW}⚠️ Transform download failed or cancelled, using brain_mesh="template"{RESET}')
                self.brain_mesh = 'template'
        elif choice == 'q':
            print(f'\\n{RED}Exiting...{RESET}')
            sys.exit(0)
        else:
            print(f'\\n{YELLOW}⚠️ Invalid choice, using brain_mesh="template"{RESET}')
            self.brain_mesh = 'template'
    
    def _dataset_needs_transform(self):
        """Check if current dataset needs H5 transforms that require file downloads.
        
        Returns
        -------
        bool
            True if H5 transforms requiring file downloads are needed,
            False if only built-in affine transforms are needed
            
        Notes
        -----
        H5 transforms required (need ~13GB download):
        - hemibrain with brain_mesh='whole': JRCFIB2018Fraw → JRC2018F path includes H5transform
        - FlyWire/FAFB with brain_mesh='whole': FAFB → JRC2018F path includes H5transform
        
        Only affine transforms (built-in, no download):
        - hemibrain with brain_mesh='template': JRCFIB2018Fraw → JRCFIB2018F
        - male-cns: JRCFIB2022Mraw → JRCFIB2022M
        - manc: MANCraw → MANC
        - optic-lobe: JRCFIB2022Mraw → JRCFIB2022M
        
        No transforms at all:
        - FlyWire/FAFB with brain_mesh='template': Native FLYWIRE template (identity transform)
        """
        dataset_lower = self.dataset.lower()
        
        # Hemibrain with brain_mesh='whole' requires H5 transforms
        # because it needs to go from JRCFIB2018Fraw to JRC2018F (involves H5transform)
        if 'hemibrain' in dataset_lower and self.brain_mesh == 'whole':
            return True
        
        # FlyWire/FAFB with brain_mesh='whole' requires H5 transforms
        # because it needs to go from FAFB to JRC2018F (involves H5transform)
        if ('flywire' in dataset_lower or 'fafb' in dataset_lower) and self.brain_mesh == 'whole':
            return True
        
        # All other cases use only affine transforms or no transforms:
        # - hemibrain with template: Affine only
        # - FlyWire/FAFB with template: No transform (native FLYWIRE)
        # - male-cns, manc, optic-lobe: Affine only
        return False
    
    def _get_fafb_tilt_correction_matrix(self):
        """Get the affine transformation matrix to correct FAFB/FLYWIRE left-right tilt.
        
        The FLYWIRE template mesh has an intrinsic left-right tilt (~8-10 degrees rotation 
        around the Z-axis when viewed from the front). This method returns a rotation matrix
        that corrects this tilt.
        
        Returns
        -------
        numpy.ndarray
            4x4 affine transformation matrix for tilt correction.
            Returns identity matrix if not FAFB dataset or not using template mode.
        
        Notes
        -----
        The rotation is applied around the center of the FLYWIRE brain mesh to avoid
        shifting objects. The rotation combines:
        - Z-axis rotation (-4 degrees) to correct left-right tilt in front view
        - Y-axis rotation (-3 degrees) to correct tilt in top view
        """
        import numpy as np
        
        # Check if correction is needed
        is_fafb = 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()
        if not is_fafb or self.brain_mesh != 'template':
            return np.eye(4)  # Return identity - no correction needed
        
        import math
        
        # Rotation angles
        angle_z = math.radians(-3)  # Front view tilt correction (Z-axis)
        angle_y = math.radians(-3)  # Top view tilt correction (Y-axis)
        
        # FLYWIRE brain mesh approximate center (computed from flybrains.FLYWIRE.mesh)
        # This ensures rotation happens around the brain center, not the origin
        center_x = 527652.0
        center_y = 240039.0
        center_z = 148110.0
        
        # Build transformation: translate to origin -> rotate -> translate back
        # T_back @ R @ T_to_origin
        
        # Translation to origin
        T_to_origin = np.array([
            [1, 0, 0, -center_x],
            [0, 1, 0, -center_y],
            [0, 0, 1, -center_z],
            [0, 0, 0, 1]
        ])
        
        # Translation back
        T_back = np.array([
            [1, 0, 0, center_x],
            [0, 1, 0, center_y],
            [0, 0, 1, center_z],
            [0, 0, 0, 1]
        ])
        
        # Rotation matrix around Z-axis (front view correction)
        cos_z = math.cos(angle_z)
        sin_z = math.sin(angle_z)
        Rz = np.array([
            [cos_z, -sin_z, 0, 0],
            [sin_z,  cos_z, 0, 0],
            [0,      0,     1, 0],
            [0,      0,     0, 1]
        ])
        
        # Rotation matrix around Y-axis (top view correction)
        cos_y = math.cos(angle_y)
        sin_y = math.sin(angle_y)
        Ry = np.array([
            [cos_y,  0, sin_y, 0],
            [0,      1, 0,     0],
            [-sin_y, 0, cos_y, 0],
            [0,      0, 0,     1]
        ])
        
        # Combined rotation: first Z, then Y
        R = Ry @ Rz
        
        # Full transformation: translate to origin, rotate, translate back
        rotation_matrix = T_back @ R @ T_to_origin
        
        return rotation_matrix
    
    def _apply_fafb_tilt_correction(self, obj):
        """Apply FAFB tilt correction to navis objects or DataFrames with xyz coordinates.
        
        Parameters
        ----------
        obj : TreeNeuron, MeshNeuron, Volume, NeuronList, or pd.DataFrame
            Object to transform. DataFrames must have 'x', 'y', 'z' columns.
            
        Returns
        -------
        same type as input
            Transformed object (in-place for most objects)
        """
        import numpy as np
        
        # Check if correction is needed
        is_fafb = 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()
        if not is_fafb or self.brain_mesh != 'template' or not self.FAFB_template_correction:
            return obj  # No correction needed
        
        rotation_matrix = self._get_fafb_tilt_correction_matrix()
        
        # Check if it's identity (no transform needed)
        if np.allclose(rotation_matrix, np.eye(4)):
            return obj
        
        # Handle DataFrame separately (for synapse coordinates)
        if isinstance(obj, pd.DataFrame):
            try:
                # Apply full 4x4 affine transform to x, y, z columns
                if 'x' in obj.columns and 'y' in obj.columns and 'z' in obj.columns:
                    # Convert to homogeneous coordinates and apply full transform
                    coords = obj[['x', 'y', 'z']].values
                    n_points = coords.shape[0]
                    
                    # Add homogeneous coordinate (1s column)
                    homogeneous = np.hstack([coords, np.ones((n_points, 1))])
                    
                    # Apply transform: result = coords @ transform.T
                    transformed = homogeneous @ rotation_matrix.T
                    
                    obj = obj.copy()
                    obj['x'] = transformed[:, 0]
                    obj['y'] = transformed[:, 1]
                    obj['z'] = transformed[:, 2]
                return obj
            except Exception as e:
                self._vprint(f'⚠️  Failed to apply FAFB tilt correction to DataFrame: {e}', level='full')
                return obj
        
        # Apply transform using navis.xform for neurons/volumes
        try:
            # navis.xform requires an AffineTransform object, not a raw numpy array
            from navis.transforms import AffineTransform
            affine_transform = AffineTransform(rotation_matrix)
            transformed = navis.xform(obj, affine_transform)
            return transformed
        except Exception as e:
            self._vprint(f'⚠️  Failed to apply FAFB tilt correction: {e}', level='full')
            return obj
    
    def _get_template_info(self):
        """Get template brain/VNC information for current dataset.
        
        Handles transform paths for all NeuPrint datasets:
        - Brain datasets: hemibrain, optic-lobe
        - VNC datasets: manc (various versions)
        - Brain+VNC datasets: male-cns
        
        Returns
        -------
        dict
            Dictionary with 'source', 'target', 'template_obj', and 'mesh_name' keys
            
        Notes
        -----
        Transform paths by dataset:
        - hemibrain: JRCFIB2018Fraw → JRCFIB2018F → JRCFIB2018Fum → JRC2018F
        - optic-lobe: JRCFIB2018Fraw → JRCFIB2018F → JRCFIB2018Fum → JRC2018F (same as hemibrain)
        - manc: MANCraw → MANC (VNC only, no brain transform)
        - male-cns: JRCFIB2022Mraw → JRCFIB2022M (brain + VNC)
        
        Note: optic-lobe uses the same coordinate system as hemibrain because it's
        a focused reconstruction of the optic lobe region within the hemibrain volume.
        """
        dataset_lower = self.dataset.lower()
        import flybrains
        
        # Brain datasets
        if 'hemibrain' in dataset_lower:
            return {
                'source': 'JRCFIB2018Fraw',
                'target': 'JRC2018F' if self.brain_mesh == 'whole' else 'JRCFIB2018F',
                'template_obj': flybrains.JRC2018F if self.brain_mesh == 'whole' else flybrains.JRCFIB2018F,
                'mesh_name': 'JRC2018F (whole brain)' if self.brain_mesh == 'whole' else 'JRCFIB2018F (hemibrain)'
            }
        elif 'optic' in dataset_lower:
            # Optic-lobe dataset is part of the Male CNS (JRCFIB2022M) volume
            # It is NOT part of the hemibrain (JRCFIB2018F) volume
            # Stored in JRCFIB2022Mraw coordinates
            return {
                'source': 'JRCFIB2022Mraw',
                'target': 'JRCFIB2022M',  # Male CNS template
                'template_obj': flybrains.JRCFIB2022M,
                'mesh_name': 'JRCFIB2022M (Male CNS)'
            }
        
        # VNC datasets
        elif 'manc' in dataset_lower:
            # MANC (Male Adult Nerve Cord)
            if self.brain_mesh == 'whole':
                # Transform to Male CNS (brain + VNC) space
                return {
                    'source': 'MANC',
                    'target': 'JRCFIB2022M',
                    'template_obj': flybrains.JRCFIB2022M,
                    'mesh_name': 'JRCFIB2022M (male CNS: brain + VNC)'
                }
            else:
                # Use native MANC template (VNC only)
                return {
                    'source': 'MANC',
                    'target': 'MANC',
                    'template_obj': flybrains.MANC,
                    'mesh_name': 'MANC (VNC envelope)'
                }
        
        # Brain + VNC datasets
        elif 'male-cns' in dataset_lower or 'malecns' in dataset_lower:
            # Male CNS (JRCFIB2022M) - Brain + VNC
            # 'whole' shows full CNS envelope (brain + VNC)
            return {
                'source': 'JRCFIB2022Mraw',
                'target': 'JRCFIB2022M',
                'template_obj': flybrains.JRCFIB2022M,
                'mesh_name': 'JRCFIB2022M (male CNS: brain + VNC)'
            }
        
        # FlyWire / FAFB datasets
        # For 'template': No transform needed - use native FLYWIRE coordinates
        # For 'whole': Transform to JRC2018F (standard female brain template)
        elif 'flywire' in dataset_lower or 'fafb' in dataset_lower:
            if self.brain_mesh == 'whole':
                # Transform to JRC2018F standard whole-brain template
                # This requires H5 transforms (~580MB download)
                return {
                    'source': 'FAFB',  # FAFB native coordinates
                    'target': 'JRC2018F',  # Standard female brain template
                    'template_obj': flybrains.JRC2018F,
                    'mesh_name': 'JRC2018F (standard whole brain)',
                    'skip_transform': False  # Need to transform skeletons/synapses
                }
            else:
                # 'template' mode: Use native FLYWIRE coordinates (no transform)
                # FLYWIRE and FAFB14 share the same bounding box and are effectively the same space
                # for visualization purposes. Using native coordinates avoids:
                # 1. Downloading ~580MB transform file (JRC2018F_FAFB.h5)
                # 2. Slow transformation of all skeleton vertices and synapse coordinates
                # 3. Potential coordinate precision loss from warping
                return {
                    'source': 'FLYWIRE',  # Native space - no transform
                    'target': 'FLYWIRE',  # Same as source - identity transform
                    'template_obj': flybrains.FLYWIRE,
                    'mesh_name': 'FLYWIRE (native FAFB coordinates)',
                    'skip_transform': True  # Flag to skip skeleton/synapse transforms
                }
        
        # Fallback to hemibrain for unknown datasets
        else:
            self._vprint(f'⚠️  Unknown dataset "{self.dataset}", defaulting to hemibrain template')
            return {
                'source': 'JRCFIB2018Fraw',
                'target': 'JRCFIB2018F',
                'template_obj': flybrains.JRCFIB2018F,
                'mesh_name': 'JRCFIB2018F (hemibrain)'
            }
    
    def _get_vnc_template_info(self):
        """Get VNC template information for current dataset.
        
        Available for datasets with VNC data:
        - male-cns: VNC portion of JRCFIB2022M (extracted geometrically, Y >= 210000)
        - manc: MANC template (native VNC mesh)
        
        Returns
        -------
        dict or None
            Dictionary with 'mesh' (trimesh object) and 'mesh_name' keys,
            or None if VNC mesh is not available for the current dataset.
        """
        dataset_lower = self.dataset.lower()
        import flybrains
        import trimesh
        
        # Male CNS and MANC datasets - extract VNC portion from JRCFIB2022M mesh
        if 'male-cns' in dataset_lower or 'malecns' in dataset_lower or 'manc' in dataset_lower:
            # Check if flybrains provides mesh_vnc (flybrains >= 0.6.3)
            if hasattr(flybrains.JRCFIB2022M, 'mesh_vnc'):
                return {
                    'mesh': flybrains.JRCFIB2022M.mesh_vnc,
                    'mesh_name': 'JRCFIB2022M (VNC)'
                }
            else:
                # Extract VNC portion geometrically
                # VNC is in the posterior portion (higher Z values in JRCFIB2022M coordinates)
                # Z cutoff ~340000 separates brain from VNC (Neck region)
                try:
                    full_mesh = flybrains.JRCFIB2022M.mesh
                    vnc_z_cutoff = 340000  # Z coordinate separating brain from VNC
                    
                    # Get vertices in the VNC region (Z >= cutoff)
                    vnc_mask = full_mesh.vertices[:, 2] >= vnc_z_cutoff
                    
                    # Extract submesh by filtering faces that have all vertices in VNC region
                    vnc_faces = []
                    for face in full_mesh.faces:
                        if all(vnc_mask[v] for v in face):
                            vnc_faces.append(face)
                    
                    if vnc_faces:
                        # Create new mesh from VNC vertices only
                        vnc_mesh_trimesh = trimesh.Trimesh(
                            vertices=full_mesh.vertices,
                            faces=vnc_faces
                        )
                        # Clean up unused vertices
                        vnc_mesh_trimesh.remove_unreferenced_vertices()
                        vnc_mesh = navis.Volume(vnc_mesh_trimesh, name='JRCFIB2022M_vnc')
                        return {
                            'mesh': vnc_mesh,
                            'mesh_name': 'JRCFIB2022M (VNC)'
                        }
                except Exception as e:
                    self._vprint(f'⚠️  VNC mesh extraction failed: {e}', level='simple')
                    return None
        
        # MANC dataset - VNC only (has proper VNC mesh)
        elif 'manc' in dataset_lower:
            return {
                'mesh': flybrains.MANC.mesh,
                'mesh_name': 'MANC (VNC)'
            }
        
        # VNC mesh not available for other datasets
        return None
    
    def _needs_skeleton_transform(self):
        """Check if skeleton/synapse coordinate transforms are needed.
        
        This method determines whether coordinate transforms should be applied
        to skeleton and synapse data during visualization.
        
        Returns
        -------
        bool
            True if transforms should be applied, False if data is already in target space
            
        Notes
        -----
        Returns False (skip transform) for:
        - FlyWire/FAFB: Data and template mesh are both in FLYWIRE space (identity transform)
        - brain_mesh='none': No template mesh, no transform needed
        
        Returns True (apply transform) for:
        - hemibrain, male-cns, manc, optic-lobe: Data needs affine transform to template space
          (These are fast, built-in affine transforms, no file download needed)
        
        Note: This is different from _dataset_needs_transform() which checks if heavy
        H5 transforms requiring file downloads are needed.
        """
        if self.brain_mesh == 'none':
            return False
            
        template_info = self._get_template_info()
        
        # Check for skip_transform flag (set for FAFB/FlyWire - identity transform)
        if template_info.get('skip_transform', False):
            return False
        
        # No transform needed if source == target (identity transform)
        if template_info['source'] == template_info['target']:
            return False
            
        return True
    
    def _check_and_download_transforms(self):
        """Check if flybrains transforms exist locally, prompt user before downloading.
        
        Brain transforms are large files (multiple files, ~10GB total uncompressed). 
        This method checks if the required transforms exist locally before attempting 
        to download them, and prompts the user for confirmation.
        
        Transforms are stored in the default flybrains data directory:
        ~/flybrain-data/
        
        Returns
        -------
        bool
            True if transforms are available (already exist or successfully downloaded),
            False otherwise.
        
        References:
        - flybrains package: https://github.com/navis-org/navis-flybrains
        - JRC2018F brain template: https://www.janelia.org/open-science/jrc-2018-brain-templates
        """
        if not self.verbose:
            return False

        try:
            import flybrains
            
            # Get the transform directory from attribute or use default
            transforms_dir = os.path.expanduser(self.transforms_dir)
            
            # Set environment variable if custom path is specified
            if self.transforms_dir != '~/flybrain-data':
                os.environ['FLYBRAINS_DATA'] = transforms_dir
                self._vprint(f'Using custom transform directory: {transforms_dir}', level='full')
            
            # Get dataset-specific template info
            template_info = self._get_template_info()
            source = template_info['source']
            target = template_info['target']
            
            # ANSI color codes
            YELLOW = '\033[93m'
            RESET = '\033[0m'
            
            # Check if the transformation path exists by attempting to find bridging path
            try:
                path = navis.transforms.registry.find_bridging_path(source, target)
                self._vprint(f'✓ Brain transforms already available', level='full')
                self._vprint(f'  Location: {YELLOW}{transforms_dir}{RESET}', level='full')
                self._vprint(f'  Transform path: {" -> ".join([str(p) for p in path])}', level='full')
                return True
            except (ValueError, KeyError):
                # Transform path not found, need to download
                pass
            
            # ANSI color codes
            YELLOW = '\033[93m'
            RESET = '\033[0m'
            
            # Prompt user for download confirmation
            self._vprint('\\n' + '='*70)
            self._vprint('⚠️  Brain Transformation Required')
            self._vprint('='*70)
            self._vprint(f'To use brain_mesh="whole" for {self.dataset}, you need brain transforms.')
            self._vprint(f'Transform path needed: {source} → JRCFIB2018F → JRCFIB2018Fum → {target}')
            self._vprint('')
            self._vprint('⚠️  IMPORTANT: flybrains downloads ALL JRC transforms as a bundle:')
            self._vprint('   • JRC2018F_JRCFIB2018F.h5   (~1.29 GB)  ← YOU NEED THIS for hemibrain/optic-lobe')
            self._vprint('   • JRC2018F_FAFB.h5          (~580 MB)   (enables FAFB dataset support)')
            self._vprint('   • JRC2018F_JFRC2013.h5      (~1.39 GB)  (enables JFRC2013 template)')
            self._vprint('   • JRC2018F_FCWB.h5          (~1.29 GB)  (enables FCWB template)')
            self._vprint('   • JRC2018U_JRC2018F.h5      (~717 MB)   (enables unisex template)')
            self._vprint('   • JRC2018U_JRC2018M.h5      (~1.10 GB)  (enables male template)')
            self._vprint('   • JRC2018F_JFRC2010.h5      (~1.65 GB)  (enables legacy template)')
            self._vprint('   • JRCFIB2022M_JRC2018M.h5   (~2.12 GB)  (enables male CNS registration)')
            self._vprint('')
            self._vprint('   Total download: ~10 GB (but only ~1.3 GB used for your dataset)')
            self._vprint('   Download time: ~1-2 hours (cannot download individual files)')
            self._vprint('   Why all files? The flybrains package bundles all transforms together.')
            self._vprint('')
            self._vprint('The transforms will be cached in:')
            self._vprint(f'  {YELLOW}{transforms_dir}/{RESET}')
            
            # Save transform path info to file
            info_file = os.path.join(self.output_dir, 'brain_transforms_info.txt')
            os.makedirs(self.output_dir, exist_ok=True)
            with open(info_file, 'w', encoding='utf-8') as f:
                f.write('Brain Transforms Information\\n')
                f.write('='*70 + '\\n\\n')
                f.write(f'Dataset: {self.dataset}\\n')
                f.write(f'Transform path: {source} → JRCFIB2018F → JRCFIB2018Fum → {target}\\n\\n')
                f.write('Storage Location:\\n')
                f.write(f'  {transforms_dir}/\\n\\n')
                f.write('Transform Files (8 files, ~10 GB total):\\n')
                f.write('  • JRC2018F_JRCFIB2018F.h5   (~1.29 GB)\\n')
                f.write('  • JRC2018F_FAFB.h5          (~580 MB)\\n')
                f.write('  • JRC2018F_JFRC2013.h5      (~1.39 GB)\\n')
                f.write('  • JRC2018F_FCWB.h5          (~1.29 GB)\\n')
                f.write('  • JRC2018U_JRC2018F.h5      (~717 MB)\\n')
                f.write('  • JRC2018U_JRC2018M.h5      (~1.10 GB)\\n')
                f.write('  • JRC2018F_JFRC2010.h5      (~1.65 GB)\\n')
                f.write('  • JRCFIB2022M_JRC2018M.h5   (~2.12 GB)\\n\\n')
                f.write('To change the storage location:\\n')
                f.write('  1. Set transforms_dir attribute when creating VisualizeSkeleton\\n')
                f.write('  2. Set FLYBRAINS_DATA environment variable before importing flybrains\\n')
                f.write('  3. Or manually move files to the new location\\n\\n')
                f.write('More information:\\n')
                f.write('  https://github.com/navis-org/navis-flybrains\\n')
            self._vprint(f'\\n📄 Transform info saved to: {info_file}')
            self._vprint('')
            self._vprint('💡 Note: The flybrains.download_jrc_transforms() function downloads')
            self._vprint('   ALL 8 files as a bundle with no selective download option.')
            self._vprint('   This is by design in the flybrains library to provide complete')
            self._vprint('   cross-dataset registration capabilities.')
            self._vprint('')
            self._vprint('For more information, see:')
            self._vprint('  https://github.com/navis-org/navis-flybrains')
            self._vprint('='*70)
            
            response = input('Download all transforms now? [y/N]: ').strip().lower()
            
            if response in ['y', 'yes']:
                self._vprint('\\n📥 Downloading brain transforms...')
                self._vprint('This may take several minutes depending on your connection.')
                flybrains.download_jrc_transforms()
                
                # Re-register transforms after download
                self._vprint('📝 Registering downloaded transforms...')
                flybrains.register_transforms()
                
                # Verify the transform path is now available
                try:
                    path = navis.transforms.registry.find_bridging_path(source, target)
                    self._vprint(f'✓ Transforms downloaded and registered successfully!')
                    self._vprint(f'  Location: {YELLOW}{transforms_dir}{RESET}')
                    self._vprint(f'  Transform path: {" -> ".join([str(p) for p in path])}')
                    
                    # Update the saved info file with success status
                    info_file = os.path.join(self.output_dir, 'brain_transforms_info.txt')
                    with open(info_file, 'a', encoding='utf-8') as f:
                        f.write(f'\\nDownload Status: SUCCESS\\n')
                        f.write(f'Downloaded at: {pd.Timestamp.now()}\\n')
                    return True
                except (ValueError, KeyError) as e:
                    self._vprint(f'⚠️  Transforms downloaded but bridging path not found: {e}')
                    self._vprint(f'   This may indicate the transforms do not include {source} → {target}')
                    return False
            else:
                self._vprint('\\n⚠️  Download cancelled. Setting brain_mesh to "none".')
                return False
                
        except ImportError:
            self._vprint('\\n⚠️  flybrains package not installed.')
            self._vprint('   Install it with: pip install navis[flybrains]')
            self._vprint('   Setting brain_mesh to "none".')
            return False
        except Exception as e:
            self._vprint(f'\\n⚠️  Error checking brain transforms: {e}')
            self._vprint('   Setting brain_mesh to "none".')
            return False
    
    def plot_mesh(self):
        """Plot ROI meshes and brain meshes.
        
        Loads ROI meshes from dataset-specific cache directories, with fallback to
        primary_rois/ for backward compatibility. Supports brain mesh visualization
        with automatic transform handling.
        
        Dataset-specific mesh caching:
        - hemibrain:v1.2.1 -> navis_roi_meshes_json/hemibrain_v1_2_1/
        - optic-lobe:v1.1 -> navis_roi_meshes_json/optic-lobe_v1_1/
        - Fallback: navis_roi_meshes_json/primary_rois/
        
        Brain mesh options (dataset-aware):
        - 'none': Only plot ROI meshes specified in mesh_roi parameter
        - 'template': Plot native EM template mesh (JRCFIB2018F, MANC, or JRCFIB2022M)
        - 'whole': Plot standard template mesh (may require transforms for some datasets)
        
        Behavior with mesh_roi=[]:
        - When mesh_roi is an empty list [], no ROI meshes are plotted
        - But brain_mesh='whole' or 'template' will still plot the brain mesh
        - This allows showing neurons with only the whole brain outline
        
        References:
        - navis Volume API: https://navis.readthedocs.io/en/latest/source/api.html#navis.Volume
        - flybrains templates: https://github.com/navis-org/navis-flybrains
        - mesh optimization: use Volume.simplify() to reduce mesh complexity for faster rendering
        """
        # Skip if mesh_roi is None (explicitly disabled)
        # Note: Empty list [] means "no ROI meshes but maybe brain mesh"
        if self.mesh_roi is None:
            return
        
        # Check if we have any work to do (ROI meshes, brain mesh, or VNC mesh)
        has_roi_meshes = len(self.mesh_roi) > 0
        has_brain_mesh = self.brain_mesh in ['template', 'whole']
        has_vnc_mesh = self.vnc_mesh
        
        if not has_roi_meshes and not has_brain_mesh and not has_vnc_mesh:
            return
        
        # For FAFB with brain_mesh='whole' and mesh_roi specified
        # ROI transforms from male-cns (JRCFIB2022Mraw) to JRC2018F require elastix which is problematic
        is_flywire = 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()
        if is_flywire and self.brain_mesh == 'whole' and has_roi_meshes:
            # Always skip ROI meshes in whole mode for FAFB/FlyWire to avoid elastix dependency
            self._vprint('')
            self._vprint('⚠️  WARNING: ROI mesh transformation is not supported in "whole" mode for FAFB.', level='simple')
            self._vprint('   ROI meshes (mesh_roi) will be skipped.', level='simple')
            self._vprint('', level='simple')
            self._vprint('   📌 Recommendation:', level='simple')
            self._vprint('   Use brain_mesh="template" instead. This mode supports ROI meshes natively.', level='simple')
            self._vprint('')
            # Clear mesh_roi to skip ROI plotting but continue with brain mesh
            self.mesh_roi = []
            has_roi_meshes = False
        
        # Ensure available_rois.json exists (generate if missing)
        # This checks cache first, and if missing, fetches from API or scans local meshes
        self._get_available_rois(use_cache=True, fetch_online=True)
        
        # Get dataset-specific mesh directory
        mesh_dir = self._get_dataset_mesh_dir()
        self._vprint(f'Using mesh directory: {mesh_dir}', level='full')
        
        roiunits = []
        roi_names = []
        roi_colors = []
        
        # Use mesh_roi list directly (no auto-mirroring suffix expansion)
        final_mesh_roi = self.mesh_roi
        
        # Handle colors - mesh_color is already standardized to rgba strings
        final_mesh_colors = []
        for i, roi in enumerate(final_mesh_roi):
            if isinstance(self.mesh_color, list):
                if i < len(self.mesh_color):
                    color = self.mesh_color[i]
                else:
                    # Fallback to default gray with current mesh_alpha
                    color = f'rgba(100, 100, 100, {self.mesh_alpha})'
            else:
                color = self.mesh_color
            
            # Ensure color is a string (standardize if needed)
            if not isinstance(color, str):
                color = standardize_color(color, default_alpha=self.mesh_alpha)
            
            final_mesh_colors.append(color)
        
        for i, roi in enumerate(final_mesh_roi):
            color = final_mesh_colors[i]
            source_info = "Dataset Cache"
            roi_source_space = None # Track the coordinate space of the ROI
            roi_needs_transform = False  # Track if ROI needs transform after loading
            
            # Determine if this is FlyWire/FAFB
            is_flywire = 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()

            # Try dataset-specific directory first (with case-safe filename)
            mesh_file = self._get_mesh_file_path(mesh_dir, roi)
            
            # For FAFB: Check for pre-transformed ROI mesh cache first
            # Transformed meshes are stored in cache/{dataset}/meshes_transformed/{TARGET}/
            # where TARGET is FLYWIRE (for template mode) or JRC2018F (for whole mode)
            transformed_mesh_file = None
            if is_flywire:
                # Determine target space based on brain_mesh mode
                target_space = 'JRC2018F' if self.brain_mesh == 'whole' else 'FLYWIRE'
                transformed_cache_dir = os.path.join(self._get_cache_path('meshes_transformed'), target_space)
                transformed_mesh_file = self._get_mesh_file_path(transformed_cache_dir, roi)
                
                if os.path.exists(transformed_mesh_file):
                    # Load pre-transformed mesh - no further transform needed
                    mesh_file = transformed_mesh_file
                    source_info = f"FAFB Transformed Cache ({target_space})"
                    roi_needs_transform = False
                    self._vprint(f'  ✓ Loading "{roi}" from transformed cache ({target_space})', level='full')
            
            # Special handling for FlyWire/FAFB - fetch from male-cns if not found
            if is_flywire and not os.path.exists(mesh_file):
                    self._vprint(f'📥 ROI mesh "{roi}" not found locally, attempting to download...', level='full')
                    mesh_found = False
                    
                    # 1. Try male-cns:v0.9 (NeuPrint)
                    try:
                        import navis.interfaces.neuprint as neu
                        from neuprint import Client
                        
                        token = os.environ.get('NEUPRINT_APPLICATION_CREDENTIALS') or self.token
                        if token:
                            try:
                                self._vprint(f'   Checking male-cns:v0.9...', level='full')
                                mc_client = Client('https://neuprint.janelia.org', dataset='male-cns:v0.9', token=token)
                                mesh = neu.fetch_roi(roi, client=mc_client)
                                if mesh:
                                    os.makedirs(mesh_dir, exist_ok=True)
                                    mesh.to_json(mesh_file)
                                    self._vprint(f'   ✓ Found in male-cns:v0.9', level='full')
                                    source_info = "male-cns:v0.9 (Downloaded)"
                                    roi_source_space = 'JRCFIB2022Mraw' # Use raw coordinates for male-cns ROIs
                                    mesh_found = True
                            except Exception as e:
                                # print(f'   (male-cns check failed: {e})')
                                pass
                    except ImportError:
                        pass
                    
                    # 2. Try generic navis fetch (fallback) - REMOVED as it causes errors if navis doesn't have fetch_roi
                    # if not mesh_found:
                    #     try:
                    #         print(f'   Attempting generic navis fetch...')
                    #         # This tries to use whatever client is default or configured in navis
                    #         # Usually fetches from Hemibrain if no dataset specified, or checks available clients
                    #         mesh = navis.fetch_roi(roi)
                    #         if mesh:
                    #             os.makedirs(mesh_dir, exist_ok=True)
                    #             mesh.to_json(mesh_file)
                    #             print(f'   ✓ Found via navis.fetch_roi')
                    #             source_info = "navis.fetch_roi"
                    #             roi_source_space = 'JRCFIB2018F' # Default for Hemibrain ROIs
                    #             mesh_found = True
                    #     except Exception as e:
                    #         print(f'   Warning: Failed to fetch "{roi}" via navis: {e}')

            # Standard logic for non-FlyWire or if file exists
            # Fallback to primary_rois if not found (only for non-FlyWire or if we want to support it)
            if not os.path.exists(mesh_file) and not is_flywire:
                # Try legacy fallback path (doesn't use case-safe encoding)
                mesh_file_fallback = os.path.join(self.script_path, 'navis_roi_meshes_json', 'primary_rois', roi + '.json')
                if os.path.exists(mesh_file_fallback):
                    mesh_file = mesh_file_fallback
                    source_info = "Primary ROIs (Local)"
                    roi_source_space = 'JRCFIB2018F'
                else:
                    # Try to download from NeuPrint (Hemibrain/Optic Lobe/Male CNS)
                    self._vprint(f'📥 ROI mesh "{roi}" not found locally, attempting to download from NeuPrint...', level='full')
                    source_info = "NeuPrint (Downloaded)"
                    try:
                        import navis.interfaces.neuprint as neu
                        from neuprint import Client
                        
                        client = self.client
                        
                        if client is None:
                            token = os.environ.get('NEUPRINT_APPLICATION_CREDENTIALS') or self.token
                            
                            if token:
                                if 'optic' in self.dataset.lower():
                                    server = 'https://neuprint-optic-lobe.janelia.org'
                                    dataset_name = self.dataset.split(':')[0]
                                    roi_source_space = 'JRCFIB2022Mraw' # Optic lobe
                                elif 'male-cns' in self.dataset.lower() or 'malecns' in self.dataset.lower():
                                    server = 'https://neuprint.janelia.org'
                                    dataset_name = 'male-cns:v0.9' # Default for male-cns
                                    roi_source_space = 'JRCFIB2022Mraw' # Male CNS raw
                                else:
                                    server = 'https://neuprint.janelia.org'
                                    dataset_name = 'hemibrain:v1.2.1'
                                    roi_source_space = 'JRCFIB2018F'
                                
                                try:
                                    client = Client(server, dataset=dataset_name, token=token)
                                except Exception as e:
                                    self._vprint(f'   Warning: Failed to create client: {e}', level='full')
                        
                        mesh = neu.fetch_roi(roi, client=client)
                        os.makedirs(mesh_dir, exist_ok=True)
                        mesh.to_json(mesh_file)
                        self._vprint(f'✓ Downloaded and cached "{roi}" mesh to {mesh_file}', level='full')
                        
                        # Transform if needed (Hemibrain specific)
                        if self.brain_mesh in ['whole', 'template']:
                            template_info = self._get_template_info()
                            self._vprint(f'Transforming brain region {roi}...', end='', level='full')
                            with self._suppress_output():
                                mesh = navis.xform_brain(mesh, source=template_info['source'], target=template_info['target'])
                            # Note: We don't save the transformed mesh back to cache here to keep cache pure?
                            # Actually previous code didn't save transformed.
                    except Exception as e:
                        self._vprint(f'⚠️  Failed to download "{roi}" mesh: {e}', level='full')
            
            # Load and plot
            if os.path.exists(mesh_file):
                try:
                    mesh = navis.Volume.from_json(mesh_file)
                    self._vprint(f'✓ Loaded "{roi}" from {source_info}', level='full')
                    
                    # Transform if needed (skip if loaded from transformed cache)
                    if self.brain_mesh in ['whole', 'template'] and "Transformed Cache" not in source_info:
                        template_info = self._get_template_info()
                        target = template_info['target']
                        
                        # For FAFB with ROIs from male-cns:
                        # - 'template' mode: transform to FLYWIRE (native FAFB space)
                        # - 'whole' mode: transform to JRC2018F (standard whole brain template)
                        if is_flywire:
                            # ROIs from male-cns need to be transformed
                            if roi_source_space:
                                source = roi_source_space
                            else:
                                # If loading from raw cache (roi_source_space is None), assume it's from male-cns
                                source = 'JRCFIB2022Mraw'
                            
                            # Target depends on brain_mesh mode
                            if self.brain_mesh == 'whole':
                                target = 'JRC2018F'  # Match skeleton/template space
                            else:
                                target = 'FLYWIRE'  # Native FAFB space for template mode
                            
                            self._vprint(f'Transforming ROI {roi} ({source} -> {target})...', end='', level='full')
                            try:
                                with self._suppress_output():
                                    mesh = navis.xform_brain(mesh, source=source, target=target)
                                self._vprint(' Done', level='full')
                                
                                # Cache the transformed mesh for future use
                                # Use target space name for cache directory (FLYWIRE or JRC2018F)
                                transformed_cache_dir = os.path.join(self._get_cache_path('meshes_transformed'), target)
                                os.makedirs(transformed_cache_dir, exist_ok=True)
                                # Use case-safe filename for transformed mesh
                                transformed_mesh_file = os.path.join(transformed_cache_dir, self._roi_to_filename(roi))
                                try:
                                    mesh.to_json(transformed_mesh_file)
                                    self._vprint(f'  💾 Cached transformed ROI to {transformed_mesh_file}', level='full')
                                except Exception as cache_e:
                                    self._vprint(f'  ⚠️ Failed to cache transformed ROI: {cache_e}', level='full')
                            except Exception as e:
                                self._vprint(f' Failed: {e}', level='full')
                        else:
                            # Non-FAFB datasets: use standard transform
                            source = template_info['source']
                            
                            # Skip transform if source == target (identity)
                            if source != target:
                                self._vprint(f'Transforming brain region {roi} ({source} -> {target})...', end='', level='full')
                                try:
                                    with self._suppress_output():
                                        mesh = navis.xform_brain(mesh, source=source, target=target)
                                    self._vprint(' Done', level='full')
                                except Exception as e:
                                    self._vprint(f' Failed: {e}', level='full')
                    
                    # Simplify mesh if requested
                    if self.roi_mesh_simplification > 0:
                        try:
                            import trimesh
                            # Access underlying trimesh object
                            tm = None
                            if hasattr(mesh, 'trimesh'):
                                tm = mesh.trimesh
                            elif hasattr(mesh, 'mesh'):
                                tm = mesh.mesh
                            elif hasattr(mesh, 'vertices') and hasattr(mesh, 'faces'):
                                # Fallback: create trimesh from vertices/faces
                                # Note: navis.Volume properties might be numpy arrays
                                tm = trimesh.Trimesh(vertices=mesh.vertices, faces=mesh.faces)
                            
                            if tm:
                                n_faces = len(tm.faces)
                                target_faces = int(n_faces * (1 - self.roi_mesh_simplification))
                                if target_faces < n_faces:
                                    # Use open3d for accurate simplification (trimesh 4.x fast_simplification only achieves ~60%)
                                    new_tm = self._simplify_mesh_open3d(tm, target_faces)
                                    
                                    # Re-instantiate Volume to ensure it's clean and updated
                                    # Preserving attributes
                                    old_name = getattr(mesh, 'name', roi)
                                    old_id = getattr(mesh, 'id', None)
                                    
                                    mesh = navis.Volume(new_tm, name=old_name, id=old_id)
                                    
                                    self._vprint(f' (simplified {self.roi_mesh_simplification*100:.0f}%: {n_faces}->{len(new_tm.faces)} faces)', end='', level='full')
                                else:
                                    self._vprint(f' (simplification skipped: target {target_faces} >= {n_faces} faces)', end='', level='full')
                            else:
                                # Debug: print available attributes to help diagnose
                                attrs = [a for a in dir(mesh) if not a.startswith('_')]
                                self._vprint(f' (simplification skipped: could not extract mesh from {type(mesh)}. Available attrs: {attrs[:10]}...)', end='', level='full')
                        except Exception as e:
                            self._vprint(f' (simplification failed: {e})', end='', level='full')

                    # Collect for export
                    try:
                        self._append_exportable_mesh(mesh, color=color, name=getattr(mesh, 'name', roi), role='roi')
                    except Exception as e:
                        self._vprint(f' (export collection failed: {e})', end='', level='full')

                    # Apply FAFB tilt correction if using template mode
                    # This corrects the left-right tilt in the FLYWIRE template mesh
                    if is_flywire and self.brain_mesh == 'template':
                        mesh = self._apply_fafb_tilt_correction(mesh)

                    roiunits.append(mesh)
                    roi_names.append(roi)
                    roi_colors.append(color)

                    # Mirror logic: SKIP if FlyWire
                    if not is_flywire:
                        contralateral_roi = roi.replace('(R)', '(L)')
                        should_mirror = (
                            self.mirror_on_contralateral and 
                            roi.endswith('(R)') and 
                            contralateral_roi not in final_mesh_roi
                        )
                        
                        if should_mirror:
                            try:
                                template = None
                                if self.brain_mesh == 'whole':
                                    template_info = self._get_template_info()
                                    template = template_info['target']
                                elif self.brain_mesh == 'template':
                                    if 'hemibrain' in self.dataset or 'optic-lobe' in self.dataset:
                                        template = 'JRCFIB2018F'
                                    elif 'male-cns' in self.dataset:
                                        template = 'JRCFIB2022M'
                                
                                if template:
                                    mirrored_mesh = navis.mirror_brain(mesh, template, mirror_axis='x')
                                    roiunits.append(mirrored_mesh)
                                    roi_names.append(contralateral_roi)
                                    roi_colors.append(color)
                            except Exception as e:
                                self._vprint(f' (mirror failed: {e})', end='', level='full')

                except Exception as e:
                    self._vprint(f'⚠️  Failed to load mesh {roi}: {e}', level='full')
            else:
                if not is_flywire: # Only warn if we expected to find it (FlyWire might just fail silently if not found)
                     self._vprint(f'⚠️  ROI mesh "{roi}" not found.', level='full')
        
        # Plot ROI meshes if any were loaded
        if roiunits:
            self._vprint('plotting mesh of brain regions...', level='full')
            for roi_i in range(len(roiunits)):
                # Colors are now standardized to rgba strings, convert to hex + alpha for navis
                color_str = roi_colors[roi_i]
                
                # Safeguard: If color is not a string (could be tuple/list from unexpected path),
                # standardize it now
                if not isinstance(color_str, str):
                    color_str = standardize_color(color_str, default_alpha=self.mesh_alpha)
                
                # Extract hex color and alpha from standardized rgba string
                color_hex = self._rgba_to_hex(color_str)
                alpha = self._extract_alpha_from_color(color_str)
                
                if self.backend == 'plotly':
                    with self._suppress_output():
                        fig_mesh = navis.plot3d(roiunits[roi_i], backend='plotly', color=color_hex, alpha=alpha)
                    mesh_traces = fig_mesh.data
                    
                    roi_name = roi_names[roi_i]
                    _configure_roi_mesh_traces(mesh_traces, roi_name)
                    self.fig_3d.add_traces(mesh_traces)
                elif self.backend == 'k3d':
                    try:
                        with self._suppress_output():
                            temp_plot = navis.plot3d(roiunits[roi_i], backend='k3d', inline=False, color=color_hex, alpha=alpha)
                        for obj in temp_plot.objects:
                            obj.name = f'brain region [{roi_names[roi_i]}]'
                            self.fig_3d += obj
                    except Exception as e:
                        self._vprint(f'⚠️  k3d mesh plotting failed: {e}', level='full')
        elif has_roi_meshes:
            # Only warn if user specified ROI meshes but none loaded
            self._vprint('⚠️  No valid ROI meshes loaded', level='full')

        # Plot brain mesh (whole brain or template) regardless of ROI mesh status
        if self.brain_mesh in ['template', 'whole']:
            template_info = self._get_template_info()
            mesh_display_name = template_info['mesh_name']
            
            # For male-cns with brain_mesh='template', use brain-only mesh if vnc_mesh is also True
            # This allows independent show/hide of brain and VNC in the interactive HTML
            # If vnc_mesh=False, show the full CNS mesh
            dataset_lower = self.dataset.lower()
            is_male_cns = 'male-cns' in dataset_lower or 'malecns' in dataset_lower
            # Note: 'manc' is NOT treated as male-cns here because its native template is VNC-only,
            # so we don't need to extract a "brain" portion from it.
            
            use_brain_only = is_male_cns and self.brain_mesh == 'template' and self.vnc_mesh
            
            if use_brain_only:
                mesh_display_name = 'JRCFIB2022M (brain)'
            
            self._vprint(f'Plotting {mesh_display_name} mesh...', level='full')
            try:
                import flybrains
                import trimesh
                
                # Select appropriate mesh
                if use_brain_only:
                    # For male-cns with vnc_mesh=True, extract brain-only portion
                    # flybrains doesn't provide mesh_brain, so we clip the mesh geometrically
                    # Brain is in the anterior portion (lower Z values in JRCFIB2022M coordinates)
                    # Z cutoff ~340000 separates brain from VNC based on geometry
                    full_mesh = flybrains.JRCFIB2022M.mesh
                    brain_z_cutoff = 340000  # Z coordinate separating brain from VNC
                    
                    # Get vertices in the brain region (Z < cutoff)
                    brain_mask = full_mesh.vertices[:, 2] < brain_z_cutoff
                    
                    # Extract submesh by filtering faces that have all vertices in brain region
                    brain_faces = []
                    for face in full_mesh.faces:
                        if all(brain_mask[v] for v in face):
                            brain_faces.append(face)
                    
                    if brain_faces:
                        # Create new mesh from brain vertices only
                        brain_mesh_trimesh = trimesh.Trimesh(
                            vertices=full_mesh.vertices,
                            faces=brain_faces
                        )
                        # Clean up unused vertices
                        brain_mesh_trimesh.remove_unreferenced_vertices()
                        brain_mesh = navis.Volume(brain_mesh_trimesh, name='JRCFIB2022M_brain')
                        self._vprint(f'   Extracted brain mesh: {len(brain_mesh_trimesh.vertices)} vertices', level='full')
                    else:
                        # Fallback to full mesh if extraction fails
                        brain_mesh = full_mesh
                        self._vprint('   ⚠️  Brain mesh extraction failed, using full CNS mesh', level='full')
                elif is_male_cns and hasattr(flybrains.JRCFIB2022M, 'mesh_brain'):
                    brain_mesh = flybrains.JRCFIB2022M.mesh_brain
                else:
                    brain_mesh = template_info['template_obj'].mesh if hasattr(template_info['template_obj'], 'mesh') else template_info['template_obj']
                
                # Apply FAFB tilt correction if using template mode
                # This corrects the left-right tilt in the FLYWIRE template mesh
                is_fafb = 'flywire' in self.dataset.lower() or 'fafb' in self.dataset.lower()
                if is_fafb and self.brain_mesh == 'template':
                    brain_mesh = self._apply_fafb_tilt_correction(brain_mesh)
                
                if self.backend == 'plotly':
                    with self._suppress_output():
                        fig_brain = navis.plot3d(brain_mesh, backend='plotly')
                    brain_traces = fig_brain.data
                    for trace in brain_traces:
                        trace.showlegend = True
                        trace.name = mesh_display_name
                        trace.hoverinfo = 'none'
                        trace.color = self._get_effective_mesh_color('brain')
                    self.fig_3d.add_traces(brain_traces)
                elif self.backend == 'k3d':
                    with self._suppress_output():
                        temp_plot = navis.plot3d(brain_mesh, backend='k3d', inline=False)
                    for obj in temp_plot.objects:
                        obj.name = mesh_display_name
                        self.fig_3d += obj

                self._append_exportable_mesh(
                    brain_mesh,
                    color=self._get_effective_mesh_color('brain'),
                    name=mesh_display_name,
                    role='brain',
                )
                        
                self._vprint(f'✓ {mesh_display_name} mesh loaded successfully', level='full')
            except Exception as e:
                self._vprint(f'⚠️  Failed to load {mesh_display_name} mesh: {e}', level='full')
                if self._dataset_needs_transform() and not self._check_and_download_transforms():
                    self._vprint('   Skipping brain/VNC mesh visualization', level='full')
                else:
                    # Retry after download - use template object mesh
                    try:
                        retry_mesh = template_info['template_obj'].mesh if hasattr(template_info['template_obj'], 'mesh') else template_info['template_obj']
                        if self.backend == 'plotly':
                            with self._suppress_output():
                                fig_brain = navis.plot3d(retry_mesh, backend='plotly')
                            brain_traces = fig_brain.data
                            for trace in brain_traces:
                                trace.showlegend = True
                                trace.name = mesh_display_name
                                trace.hoverinfo = 'none'
                                trace.color = self._get_effective_mesh_color('brain')
                            self.fig_3d.add_traces(brain_traces)
                        elif self.backend == 'k3d':
                            with self._suppress_output():
                                temp_plot = navis.plot3d(retry_mesh, backend='k3d', inline=False)
                            for obj in temp_plot.objects:
                                obj.name = mesh_display_name
                                self.fig_3d += obj
                        self._append_exportable_mesh(
                            retry_mesh,
                            color=self._get_effective_mesh_color('brain'),
                            name=mesh_display_name,
                            role='brain',
                        )
                        self._vprint(f'✓ {mesh_display_name} mesh loaded successfully after download', level='full')
                    except Exception as retry_e:
                        self._vprint(f'⚠️  Still failed to load {mesh_display_name} mesh: {retry_e}', level='full')
                        self._vprint('   Skipping brain/VNC mesh visualization', level='full')
        
        # Plot VNC mesh if requested (only for manc and male-cns datasets)
        if self.vnc_mesh:
            dataset_lower = self.dataset.lower()
            
            # Check if VNC is already shown by brain_mesh (for male-cns/manc datasets using JRCFIB2022M)
            is_male_cns = 'male-cns' in dataset_lower or 'malecns' in dataset_lower
            is_manc = 'manc' in dataset_lower
            
            # use_brain_only logic from above (replicated here for clarity)
            use_brain_only = is_male_cns and self.brain_mesh == 'template'
            
            # MANC with template mode: The "brain mesh" (template) IS the VNC mesh.
            if is_manc and self.brain_mesh == 'template':
                self._vprint('ℹ️  VNC mesh already shown as template', level='full')
            
            # If using JRCFIB2022M and NOT splitting brain (e.g. brain_mesh='whole'), VNC is already included
            elif (is_male_cns or is_manc) and self.brain_mesh in ['template', 'whole'] and not use_brain_only:
                self._vprint('ℹ️  VNC mesh already shown via brain_mesh', level='full')
            else:
                vnc_info = self._get_vnc_template_info()
                if vnc_info:
                    vnc_display_name = vnc_info['mesh_name']
                    self._vprint(f'Plotting {vnc_display_name} mesh...', level='full')
                    try:
                        vnc_mesh = vnc_info['mesh']
                        
                        if self.backend == 'plotly':
                            with self._suppress_output():
                                fig_vnc = navis.plot3d(vnc_mesh, backend='plotly')
                            vnc_traces = fig_vnc.data
                            for trace in vnc_traces:
                                trace.showlegend = True
                                trace.name = vnc_display_name
                                trace.hoverinfo = 'none'
                                trace.color = self._get_effective_mesh_color('vnc')
                            self.fig_3d.add_traces(vnc_traces)
                        elif self.backend == 'k3d':
                            with self._suppress_output():
                                temp_plot = navis.plot3d(vnc_mesh, backend='k3d', inline=False)
                            for obj in temp_plot.objects:
                                obj.name = vnc_display_name
                                self.fig_3d += obj

                        self._append_exportable_mesh(
                            vnc_mesh,
                            color=self._get_effective_mesh_color('vnc'),
                            name=vnc_display_name,
                            role='vnc',
                        )
                        
                        self._vprint(f'✓ {vnc_display_name} mesh loaded successfully', level='full')
                    except Exception as e:
                        self._vprint(f'⚠️  Failed to load VNC mesh: {e}', level='full')
                else:
                    self._vprint('⚠️  VNC mesh is only available for manc and male-cns datasets', level='full')
        
        self._vprint('Done', level='full')
        return 0
    
    def save_figure(self):
        if self.backend == 'plotly':
            # No sliders currently used
            sliders = []
            
            # set layout
            # Always use frontal view camera regardless of brain_mesh setting
            # This ensures consistent viewing angle for all visualizations
            # Standard fly brain orientation: X: Left-Right, Y: Dorsal-Ventral, Z: Anterior-Posterior
            # Frontal view: Look from Anterior (negative Z direction)
            
            scene_camera_parameters = dict(
                up=dict(x=0, y=-1, z=0),  # Y is up (inverted in some templates)
                eye=dict(x=0, y=0, z=-2.5),  # Look from front
                # center=dict(x=0, y=0, z=0), # Let Plotly auto-center
            )
            
            # Fix for MANC dataset (VNC)
            # Default Front (-Z) shows Tail. We want to look from Neck (+Z).
            if 'manc' in self.dataset.lower():
                 scene_camera_parameters = dict(
                    up=dict(x=0, y=-1, z=0),  # Dorsal (-Y) is up
                    eye=dict(x=0, y=0, z=2.5),  # Look from Anterior (+Z)
                )

            # Fix for hemibrain template mode (JRCFIB2018F)
            # JRCFIB2018F has Y-axis pointing posterior→anterior, Z is Dorsal-Ventral.
            # To show Frontal view, we need to look from Anterior (positive Y axis).
            if 'hemibrain' in self.dataset.lower() and self.brain_mesh == 'template':
                 scene_camera_parameters = dict(
                    up=dict(x=0, y=0, z=-1),  # Z is up (Dorsal)
                    eye=dict(x=0, y=2.0, z=0),  # Look from Front (Anterior is +Y)
                )
            
            self.fig_3d.update_layout(
                colorway = self.synapse_colors,
                sliders=sliders,
                paper_bgcolor=self.background_color,
                plot_bgcolor=self.background_color,
                scene=dict(
                    dragmode='orbit',
                    xaxis={'visible':False}, 
                    yaxis={'visible':False},
                    zaxis={'visible':False},
                    # Use 'data' aspectmode to ensure equal axis scaling
                    # This prevents distortion when no meshes are plotted
                    aspectmode='data',
                    bgcolor=self.background_color,
                ),
                scene_camera=scene_camera_parameters,
                # Legend settings: use constant sizing so alpha doesn't affect legend swatches
                legend=dict(
                    itemsizing='constant',  # Fixed legend swatch size regardless of trace properties
                    font=dict(color='white' if self._is_dark_background() else 'black'),
                    bgcolor='rgba(0,0,0,0)',  # Transparent legend background
                ),
            )

            # save figure
            self.fig_path = os.path.join(self.save_folder,self.saveas)
            
            # Ensure save folder exists
            if not os.path.exists(self.save_folder):
                os.makedirs(self.save_folder, exist_ok=True)
            
            self._vprint(f'saving figure to \033[34m{self.fig_path}.html\033[0m...', end='')
            
            # Configure interactive controls based on interactive_html setting
            if self.interactive_html:
                # Add view selection dropdown menu to the figure
                self._add_view_selection_menu()
                
                # Interactive config: show toolbar, enable scroll zoom, customize buttons
                html_config = {
                    'displayModeBar': True,
                    'displaylogo': False,  # Hide Plotly logo
                    'scrollZoom': True,
                    'modeBarButtonsToRemove': ['sendDataToCloud', 'lasso2d', 'select2d'],
                    'toImageButtonOptions': {
                        'format': 'png',
                        'filename': self.saveas,
                        'height': 1800,
                        'width': 2400,
                        'scale': 2
                    }
                }
            else:
                # Minimal config: hide toolbar for cleaner appearance
                html_config = {'displayModeBar': False}
            
            # Optimization: use 'cdn' for smaller file size (loads plotly.js from CDN)
            # This reduces HTML file size significantly compared to 'directory' or including full plotly.js
            # Fix: Set auto_open=False to prevent hanging, handle opening manually
            # Reverted 'cdn' to default (embed) as user reported issues with subsequent PNG export
            self.fig_3d.write_html(
                self.fig_path+'.html',
                auto_open=False, 
                # include_plotlyjs='cdn',  # Reverted to default to avoid potential issues
                config=html_config
            )
            
            if self.show_fig:
                try:
                    import webbrowser
                    webbrowser.open('file://' + os.path.abspath(self.fig_path+'.html'))
                except Exception as e:
                    self._vprint(f'\\n⚠️  Failed to open browser: {e}')
            
            self._vprint('Done (HTML saved)')
            
            # Export multiple view angles as PNG
            # Skip if export_views is False
            if self.export_views is not False:
                try:
                    self._vprint('   Exporting static PNGs (multiple views)...')
                    
                    # Create exported_views subfolder
                    views_folder = os.path.join(self.save_folder, 'exported_views')
                    os.makedirs(views_folder, exist_ok=True)
                    
                    # Define camera angles for different views
                    # Based on the default front view: eye=(0, 0, -2), up=(0, -1, 0)
                    # X: Left-Right, Y: Dorsal-Ventral (up), Z: Anterior-Posterior (front-back)
                    view_cameras = {
                        'front': dict(eye=dict(x=0, y=0, z=-2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                        'back': dict(eye=dict(x=0, y=0, z=2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                        'top': dict(eye=dict(x=0, y=-2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=1)),
                        'bottom': dict(eye=dict(x=0, y=2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
                        'left': dict(eye=dict(x=-2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                        'right': dict(eye=dict(x=2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                    }

                    # Adjust for MANC (Male Adult Nerve Cord)
                    # User reports default Front view (-Z) shows Tail (Posterior).
                    # This implies Tail is at -Z (closest to camera), so Anterior is at +Z.
                    # Fix: Reverse Z axis for Front/Back views.
                    if 'manc' in self.dataset.lower():
                         view_cameras = {
                            'front': dict(eye=dict(x=0, y=0, z=2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                            'back': dict(eye=dict(x=0, y=0, z=-2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                            'top': dict(eye=dict(x=0, y=-2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=1)),
                            'bottom': dict(eye=dict(x=0, y=2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=1)),
                            'left': dict(eye=dict(x=-2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                            'right': dict(eye=dict(x=2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                        }

                    # Adjust for hemibrain template (JRCFIB2018F)
                    # JRCFIB2018F coordinate system:
                    #   X-axis: Left-Right (as normal fly brain)
                    #   Y-axis: Posterior→Anterior (front is +Y direction)
                    #   Z-axis: Ventral→Dorsal (up is +Z direction, but we use -Z as "up" for standard brain viewing)
                    #
                    # For front view: eye at +Y, looking at origin, with -Z as up (dorsal up)
                    # Key insight: When looking from +Y, the up vector should be -Z to match standard viewing
                    if 'hemibrain' in self.dataset.lower() and self.brain_mesh == 'template':
                         view_cameras = {
                            'front': dict(eye=dict(x=0, y=2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
                            'back': dict(eye=dict(x=0, y=-2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
                            'top': dict(eye=dict(x=0, y=0, z=-2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=1, z=0)),
                            'bottom': dict(eye=dict(x=0, y=0, z=2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=1, z=0)),
                            'left': dict(eye=dict(x=-2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
                            'right': dict(eye=dict(x=2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
                        }
                    
                    # Update layout for static export to remove UI elements
                    self.fig_3d.update_layout(
                        margin=dict(l=0, r=0, b=0, t=0),
                        sliders=[],      # Remove sliders
                        updatemenus=[],  # Remove view selection dropdown
                        annotations=[],  # Remove controls hint and view label
                        paper_bgcolor=self.background_color,
                        plot_bgcolor=self.background_color,
                        scene=dict(bgcolor=self.background_color),
                        legend=dict(
                            bgcolor=self.background_color,
                            font=dict(color='white' if self._is_dark_background() else 'black'),
                        ),
                    )
                    
                    import shutil
                    front_view_path = None
                    
                    # Determine which views to export
                    views_to_export = []
                    if self.export_views is True:
                        views_to_export = list(view_cameras.keys())
                    elif isinstance(self.export_views, list):
                        views_to_export = self.export_views
                    elif isinstance(self.export_views, str):
                        views_to_export = [self.export_views]
                    
                    if not views_to_export:
                        self._vprint('   Skipping PNG export (export_views=False)')
                    else:
                        # Check HTML size
                        html_path = self.fig_path + '.html'
                        html_size_mb = 0
                        if os.path.exists(html_path):
                            html_size_mb = os.path.getsize(html_path) / 1024 / 1024
                        
                        # Determine if we should pre-simplify based on export_simplified_png setting
                        # Default (False): Only simplify on timeout
                        # True: Simplify if > 50MB
                        # int: Simplify if > that many MB
                        simplify_threshold_mb = None
                        if self.export_simplified_png is True:
                            simplify_threshold_mb = 50
                        elif isinstance(self.export_simplified_png, (int, float)) and self.export_simplified_png > 0:
                            simplify_threshold_mb = self.export_simplified_png
                        
                        # Create export figure (may be simplified)
                        export_fig = self.fig_3d
                        using_simplified = False
                        html_size_cap = self._get_html_size_cap()
                        
                        # Auto-simplify for kaleido if HTML > size cap (prevents timeout)
                        if self.export_method not in ('webdriver', 'webdriver-fast') and html_size_mb > html_size_cap:
                            # Target 50% file size reduction
                            # HTML overhead (Plotly.js, layout, etc.) is ~20-30MB
                            # To get 50% total reduction, we need more aggressive mesh simplification
                            # Factor = (target_size - overhead) / (original_size - overhead)
                            html_overhead_mb = 25  # Estimated fixed overhead
                            target_size_mb = html_size_mb * 0.5  # 50% target
                            data_size = html_size_mb - html_overhead_mb
                            target_data_size = target_size_mb - html_overhead_mb
                            simplification_factor = max(0.1, target_data_size / data_size) if data_size > 0 else 0.5
                            
                            self._vprint(f'      ⚠️  HTML size ({html_size_mb:.0f}MB) exceeds {html_size_cap}MB - auto-simplifying for kaleido')
                            self._vprint(f'         Target: {target_size_mb:.0f}MB (50% reduction), simplification factor: {simplification_factor:.2f}')
                            export_fig = self._simplify_figure_for_kaleido(self.fig_3d, simplification_factor)
                            using_simplified = True
                            
                            # Save simplified HTML for reuse by export_video, plot_individuals, etc.
                            simplified_html_path = os.path.join(self.save_folder, f"{self.saveas}_simplified.html")
                            export_fig.write_html(simplified_html_path, auto_open=False,
                                                 include_plotlyjs='cdn',
                                                 config={'displayModeBar': False})
                            simplified_size_mb = os.path.getsize(simplified_html_path) / 1024 / 1024
                            self._vprint(f'      ✓ Saved simplified HTML: {os.path.basename(simplified_html_path)} ({simplified_size_mb:.0f}MB)')
                            
                            # Store for reuse
                            self._simplified_export_fig = export_fig
                            self._simplified_html_path = simplified_html_path
                            
                            self._vprint(f'      💡 For better quality, use export_method="webdriver"')
                        elif simplify_threshold_mb and html_size_mb > simplify_threshold_mb:
                            self._vprint(f'      ⚠️  HTML size: {html_size_mb:.1f}MB > {simplify_threshold_mb}MB - creating simplified copy for export')
                            export_fig = self._create_simplified_export_figure()
                            using_simplified = True
                        elif html_size_mb > 50:
                            self._vprint(f'      ℹ️  HTML size: {html_size_mb:.1f}MB - export may take 1-2 minutes')
                        
                        # Track if any export fails to skip remaining views
                        export_failed = False
                        exported_views = []
                        use_kaleido_fallback = False
                        
                        # Choose export method
                        if self.export_method in ('webdriver', 'webdriver-fast'):
                            # WebDriver method - OPTIMIZED: open browser once for all views
                            exported_views = self._export_views_with_webdriver_session(
                                export_fig, views_to_export, view_cameras, views_folder
                            )
                            export_failed = len(exported_views) == 0
                            if exported_views and 'front' in exported_views:
                                front_view_path = os.path.join(views_folder, f"{self.saveas}_front.png")
                            
                            # Fallback to kaleido if webdriver failed
                            if export_failed:
                                self._vprint(f'      ⚠️  WebDriver export failed. Falling back to kaleido...')
                                self._vprint(f'      💡 Tip: Set export_method="kaleido" for future use if WebDriver continues to fail.')
                                use_kaleido_fallback = True
                                export_failed = False  # Reset to try kaleido
                        
                        if self.export_method == 'kaleido' or use_kaleido_fallback:
                            # Default: kaleido export with timeout (per-view)
                            for view_name in views_to_export:
                                if export_failed:
                                    self._vprint(f'      ⚠️  Skipping {view_name} (previous export failed)')
                                    continue
                                    
                                if view_name not in view_cameras:
                                    self._vprint(f'      ⚠️  Skipping invalid view: {view_name}')
                                    continue
                                    
                                camera = view_cameras[view_name]
                                view_path = os.path.join(views_folder, f"{self.saveas}_{view_name}.png")
                                export_fig.update_layout(scene_camera=camera)
                                
                                success, msg, final_scale = self._export_png_with_timeout(
                                    export_fig, view_path, 
                                    width=1200, height=900, 
                                    scale=self.export_scale, 
                                    timeout=300
                                )
                                
                                if success:
                                    self._vprint(f'      {view_name}: {msg}', level='full')
                                    exported_views.append(view_name)
                                    
                                    # Save front view path for copying to root
                                    if view_name == 'front':
                                        front_view_path = view_path
                                else:
                                    self._vprint(f'      ⚠️  {view_name} export failed: {msg}')
                                    export_failed = True
                                    self._vprint(f'      💡 Skipping remaining views. Try reducing export_scale or increasing skeleton_mesh_simplification.')
                        
                        # Copy front view to root folder without '_front' suffix
                        if front_view_path and os.path.exists(front_view_path):
                            root_png_path = os.path.join(self.save_folder, f"{self.saveas}.png")
                            shutil.copy2(front_view_path, root_png_path)
                            self._vprint(f'   ✓ Copied front view to root: {self.saveas}.png')
                        
                        if exported_views:
                            self._vprint(f'   ✓ Exported {len(exported_views)} view PNGs to exported_views/ ({", ".join(exported_views)})')
                        
                except Exception as e:
                    self._vprint(f'\\n   ⚠️  PNG/SVG export failed: {e}. Continuing without static images...')
            
        elif self.backend == 'k3d':
            self.fig_path = os.path.join(self.save_folder,self.saveas)
            self._vprint(f'saving figure to \033[34m{self.fig_path}.html\033[0m...', end='')
            
            try:
                from ipywidgets.embed import embed_minimal_html
                embed_minimal_html(
                    self.fig_path+'.html', 
                    views=[self.fig_3d], 
                    title=self.saveas
                )
                self._vprint('Done')
                
                if self.show_fig:
                    self._vprint('Note: k3d plots cannot be automatically opened from script. Please open the HTML file manually.')
                    
            except ImportError:
                self._vprint('\\n⚠️  ipywidgets not installed. Cannot save k3d plot to HTML.')
                self._vprint('   Please install it with `pip install ipywidgets`')
            except Exception as e:
                self._vprint(f'\\n⚠️  Failed to save k3d plot: {e}')
    
    def plot_neurons(self):
        import time
        start_time = time.time()
        
        # Reset simplified figure and HTML path from previous export
        self._simplified_export_fig = None
        self._simplified_html_path = None
        
        self._vprint('\n' + '='*60)
        self._vprint(f'🧠 VisualizeSkeleton: Plotting {self.dataset}')
        self._vprint('='*60)
        
        self.plot_skeleton()
        self.plot_synapses()
        self.plot_mesh()
        self.save_figure()
        
        elapsed = time.time() - start_time
        self._vprint(f'\n✅ Complete! Total time: {elapsed:.1f}s')
        self._vprint(f'📁 Output: {self.save_folder}')
        self._vprint('='*60 + '\n')

    def plot_individuals(
        self,
        output_format: str | list = 'png',
        views: str | list = 'front',
        scale: int = None,
        pdf_images_per_page: tuple = (3, 2),
        pdf_title: str = None,
        neuron_alpha: float = None,
        summary_format: str | list = 'pdf',
        auto_crop: bool = True,
        crop_margin: int = 30,
        export_method: str = None,
    ):
        """
        Plot individual neurons/types independently based on the main figure's legend entries.
        
        This method should be called AFTER plot_neurons() to ensure all necessary data is available.
        It iterates through the legend entries in the main figure and generates separate plots
        for each individual legend item by hiding other neuron traces (efficient, no duplication).
        
        Behavior varies by legend_mode:
        - 'single': plots individual neurons (each neuron separate)
        - 'type': plots by neuron type (grouped by type within layers)
        - 'layer': plots by layer (all neurons in a layer grouped)
        
        Parameters
        ----------
        output_format : str or list, default 'png'
            Output format(s) for individual plots.
            Options: 'png', 'html', or list like ['png', 'html']
        views : str or list, default 'front'
            View angle(s) for PNG exports.
            Options: 'front', 'back', 'top', 'bottom', 'left', 'right'
            Can be a single string or list like ['front', 'top']
        scale : int, default None (uses self.export_scale)
            Scale factor for PNG export resolution. Overrides self.export_scale if provided.
            Higher values produce larger, higher-quality images.
        pdf_images_per_page : tuple, default (3, 2)
            (columns, rows) - number of images per page when generating PDF/PPTX.
        pdf_title : str, optional
            Custom title for PDF/PPTX pages. If None, uses the layer/neuron name.
        neuron_alpha : float, optional
            Opacity for neuron traces in individual plots (0.0-1.0).
            If None, defaults to 0.8 for better visibility in individual views.
        summary_format : str or list, default 'pdf'
            Format(s) for summary file generation.
            Options: 'pdf', 'pptx', or list like ['pdf', 'pptx']
        auto_crop : bool, default True
            If True, automatically crop whitespace/background from PNG exports
            to minimize empty margins around content.
        crop_margin : int, default 30
            Margin (in pixels) to preserve around content when auto_crop=True.
        export_method : str, optional
            Export method for PNG generation. Options: 'webdriver', 'kaleido'.
            If None, uses self.export_method (default from class initialization).
            - 'webdriver': Uses Selenium Chrome for rendering (better for large figures)
            - 'kaleido': Uses kaleido for rendering (faster for small figures)
            
        Returns
        -------
        str or None
            Path to the output folder containing individual plots,
            or None if no plots were generated.
            
        Example
        -------
        >>> vs = VisualizeSkeleton(...)
        >>> vs.plot_neurons()
        >>> vs.plot_individuals(output_format=['png', 'html'], views=['front', 'top'])
        >>> vs.plot_individuals(summary_format=['pdf', 'pptx'])  # Generate both PDF and PPTX
        >>> vs.plot_individuals(export_method='webdriver')  # Force webdriver for large figures
        """
        import copy
        
        if not hasattr(self, 'fig_3d') or self.fig_3d is None:
            self._vprint('⚠️  No figure found. Please run plot_neurons() first.')
            return None
            
        if self.backend != 'plotly':
            self._vprint('⚠️  plot_individuals() only supports plotly backend.')
            return None
        
        # Use specified export_method or fall back to self.export_method
        actual_export_method = export_method if export_method is not None else self.export_method
        
        # Use self.export_scale if scale not specified
        # For webdriver modes, cap inherited scale at 5 for stability (user can override explicitly)
        if scale is None:
            scale = self.export_scale
            if actual_export_method in ('webdriver', 'webdriver-fast') and scale > 5:
                scale = 5
                self._vprint(f'   ℹ️  Capped inherited scale to 5 for webdriver mode (explicit scale= overrides)')
        
        # Normalize inputs
        if isinstance(output_format, str):
            output_format = [output_format]
        if isinstance(views, str):
            views = [views]
            
        # Validate inputs
        valid_formats = {'png', 'html'}
        output_format = [f.lower() for f in output_format]
        for fmt in output_format:
            if fmt not in valid_formats:
                self._vprint(f'⚠️  Invalid output format: {fmt}. Use "png" or "html".')
                return None
                
        valid_views = {'front', 'back', 'top', 'bottom', 'left', 'right'}
        views = [v.lower() for v in views]
        for view in views:
            if view not in valid_views:
                self._vprint(f'⚠️  Invalid view: {view}. Use one of {valid_views}.')
                return None
        
        # Create output directory
        output_dir = os.path.join(self.save_folder, 'individual_profiles')
        os.makedirs(output_dir, exist_ok=True)
        
        self._vprint(f'\n📊 Generating individual plots...')
        self._vprint(f'   Output: {output_dir}')
        
        # View cameras for PNG export
        view_cameras = {
            'front': dict(eye=dict(x=0, y=0, z=-2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
            'back': dict(eye=dict(x=0, y=0, z=2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
            'top': dict(eye=dict(x=0, y=-2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=1)),
            'bottom': dict(eye=dict(x=0, y=2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
            'left': dict(eye=dict(x=-2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
            'right': dict(eye=dict(x=2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
        }
        
        # Adjust for MANC (Male Adult Nerve Cord)
        # Fix Front view (-Z shows Tail) -> Should look from +Z (Neck)
        if 'manc' in self.dataset.lower():
             view_cameras = {
                'front': dict(eye=dict(x=0, y=0, z=2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                'back': dict(eye=dict(x=0, y=0, z=-2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                'top': dict(eye=dict(x=0, y=-2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=1)),
                'bottom': dict(eye=dict(x=0, y=2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=1)),
                'left': dict(eye=dict(x=-2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
                'right': dict(eye=dict(x=2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=-1, z=0)),
            }

        # Adjust for hemibrain template (JRCFIB2018F) - front/back swapped due to Y-axis orientation
        # JRCFIB2018F coordinate system:
        #   X-axis: Left-Right (as normal fly brain)
        #   Y-axis: Posterior→Anterior (front is +Y direction)
        #   Z-axis: Ventral→Dorsal (up is +Z direction, but we use -Z as "up" for standard brain viewing)
        if 'hemibrain' in self.dataset.lower() and self.brain_mesh == 'template':
            view_cameras = {
                'front': dict(eye=dict(x=0, y=2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
                'back': dict(eye=dict(x=0, y=-2.5, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
                'top': dict(eye=dict(x=0, y=0, z=-2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=1, z=0)),
                'bottom': dict(eye=dict(x=0, y=0, z=2.5), center=dict(x=0, y=0, z=0), up=dict(x=0, y=1, z=0)),
                'left': dict(eye=dict(x=-2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
                'right': dict(eye=dict(x=2.5, y=0, z=0), center=dict(x=0, y=0, z=0), up=dict(x=0, y=0, z=-1)),
            }
        
        # Get all traces from the main figure
        all_traces = list(self.fig_3d.data)
        n_traces = len(all_traces)
        
        # Default neuron_alpha for individual plots (higher than main plot for better visibility)
        individual_alpha = neuron_alpha if neuron_alpha is not None else 0.8
        
        # Helper function to modify alpha in RGBA color string
        def _modify_color_alpha(color_str, new_alpha):
            """Modify the alpha value in an RGBA color string.
            
            navis encodes alpha in the color as RGBA (e.g., 'rgba(255,0,0,0.2)'),
            not in the opacity attribute. To change effective alpha, we must
            modify the color string directly.
            """
            import re
            if color_str is None:
                return None
            color_str = str(color_str)
            
            # Match rgba(r,g,b,a) format
            rgba_match = re.match(r'rgba?\(([^,]+),\s*([^,]+),\s*([^,]+)(?:,\s*([^)]+))?\)', color_str)
            if rgba_match:
                r, g, b = rgba_match.group(1), rgba_match.group(2), rgba_match.group(3)
                return f'rgba({r},{g},{b},{new_alpha})'
            
            # Match rgb(r,g,b) format - add alpha
            rgb_match = re.match(r'rgb\(([^,]+),\s*([^,]+),\s*([^)]+)\)', color_str)
            if rgb_match:
                r, g, b = rgb_match.group(1), rgb_match.group(2), rgb_match.group(3)
                return f'rgba({r},{g},{b},{new_alpha})'
            
            # For other formats (hex, named colors), try matplotlib
            try:
                import matplotlib.colors as mcolors
                rgba = mcolors.to_rgba(color_str)
                r, g, b = int(rgba[0]*255), int(rgba[1]*255), int(rgba[2]*255)
                return f'rgba({r},{g},{b},{new_alpha})'
            except:
                return color_str  # Return unchanged if can't parse
        
        # Store original visibility, opacity, and color states to restore later
        original_visibility = []
        original_opacity = []
        original_colors = []
        for trace in all_traces:
            original_visibility.append(getattr(trace, 'visible', True))
            original_opacity.append(getattr(trace, 'opacity', None))
            original_colors.append(getattr(trace, 'color', None))
        
        # Store original camera and layout settings
        original_layout = copy.deepcopy(self.fig_3d.layout)
        
        # Identify unique legend entries (excluding hidden legends and mesh/synapse traces)
        legend_entries = {}  # {legend_name: [trace_indices]}
        background_indices = []  # mesh/synapse traces to always show
        
        # Get mesh_roi names for matching
        mesh_roi_names = [r.lower() for r in self.mesh_roi] if self.mesh_roi else []
        
        # Template/mesh names to always include as background
        mesh_keywords = ['mesh', 'brain region', 'template', 'vnc']
        template_names = ['JRCFIB', 'MANC', 'JRC2018', 'FLYWIRE', 'FAFB', 'jrcfib', 'flywire', 'fafb']
        
        for idx, trace in enumerate(all_traces):
            trace_name = getattr(trace, 'name', '')
            show_legend = getattr(trace, 'showlegend', True)
            legend_group = getattr(trace, 'legendgroup', None)
            trace_name_lower = trace_name.lower() if trace_name else ''
            
            # Identify mesh/roi traces (keep visible as background)
            # Include brain-region traces, standard templates, and user-specified mesh_roi
            is_background = False
            
            if trace_name:
                # Check for mesh-related keywords
                if any(kw in trace_name_lower for kw in mesh_keywords):
                    is_background = True
                # Check for template names (case-insensitive for most, case-sensitive for acronyms)
                elif any(tn in trace_name for tn in template_names):
                    is_background = True
                # Check for user-specified ROI meshes
                elif any(roi_name in trace_name_lower for roi_name in mesh_roi_names):
                    is_background = True
            
            if is_background:
                background_indices.append(idx)
                continue
                
            # Identify synapse traces (keep visible as background)
            if trace_name and ('synapse' in trace_name_lower or 'pre-syn' in trace_name_lower or 'post-syn' in trace_name_lower):
                background_indices.append(idx)
                continue
            
            # Use legendgroup as key if available (for merged traces), else use name
            key = legend_group if legend_group else trace_name
            if key and show_legend:
                if key not in legend_entries:
                    legend_entries[key] = []
                legend_entries[key].append(idx)
            elif key and not show_legend and legend_group:
                # Traces with same legendgroup but showlegend=False
                if key not in legend_entries:
                    legend_entries[key] = []
                legend_entries[key].append(idx)
        
        if not legend_entries:
            self._vprint('⚠️  No legend entries found to plot individually.')
            return None
        
        self._vprint(f'   Found {len(legend_entries)} individual legend entries')
        
        # Generate individual plots by hiding/showing traces
        generated_files = {'png': {}, 'html': []}
        
        # Track if PNG export fails to skip remaining individuals
        png_export_failed = False
        
        # Helper to sanitize filenames
        def _sanitize_filename(name):
            safe_name = "".join(c if c.isalnum() or c in '.+_- ' else '_' for c in str(name))
            safe_name = safe_name.strip().replace(' ', '_')
            while '__' in safe_name:
                safe_name = safe_name.replace('__', '_')
            return safe_name.rstrip('_')
        
        from tqdm import tqdm
        legend_names = list(legend_entries.keys())
        
        # ==============================================================================
        # OPTIMIZED WEBDRIVER EXPORT: Open browser once, toggle visibility via JavaScript
        # ==============================================================================
        if actual_export_method in ('webdriver', 'webdriver-fast') and 'png' in output_format:
            self._vprint(f'   Using optimized WebDriver export (single browser session)')
            
            # Use simplified HTML if available from previous export (e.g., export_views)
            simplified_html = getattr(self, '_simplified_html_path', None)
            if simplified_html and os.path.exists(simplified_html):
                self._vprint(f'   ✓ Reusing simplified HTML from previous export')
                temp_html = simplified_html
            else:
                # Use simplified figure if available, otherwise use original
                export_fig = getattr(self, '_simplified_export_fig', None) or self.fig_3d
                if export_fig is not self.fig_3d:
                    self._vprint(f'   ✓ Reusing simplified figure from previous export')
                
                # Save main HTML with all traces visible for WebDriver
                temp_html = os.path.join(output_dir, '_temp_main_figure.html')
                export_fig.write_html(temp_html, auto_open=False, 
                                       include_plotlyjs='cdn',
                                       config={'displayModeBar': False})
            
            # Use the optimized helper function
            result = export_individuals_webdriver(
                html_path=temp_html,
                output_dir=output_dir,
                legend_entries=legend_entries,
                background_indices=background_indices,
                total_traces=n_traces,
                views=views,
                view_cameras=view_cameras,
                scale=scale,
                width=900,
                height=900,
                timeout=300,
                verbose=True,
                auto_crop=auto_crop,
                crop_margin=crop_margin,
                background_color=self.background_color
            )
            
            # Clean up temp HTML (but not the permanent simplified HTML)
            if temp_html and not temp_html.endswith('_simplified.html'):
                try:
                    os.remove(temp_html)
                except:
                    pass
            
            if result['success']:
                generated_files['png'] = result['files']
            else:
                png_export_failed = True
                self._vprint(f'   ⚠️  WebDriver export failed: {result["error"]}')
            
            # Export HTML files separately if requested (in traditional loop)
            if 'html' in output_format:
                for legend_name in tqdm(legend_names, desc='Exporting HTML files'):
                    trace_indices = legend_entries[legend_name]
                    safe_name = _sanitize_filename(legend_name)
                    
                    # Set visibility for this individual
                    for idx in range(n_traces):
                        if idx in trace_indices or idx in background_indices:
                            self.fig_3d.data[idx].visible = True
                            if idx in trace_indices:
                                trace = self.fig_3d.data[idx]
                                if hasattr(trace, 'color') and trace.color is not None:
                                    trace.color = _modify_color_alpha(trace.color, individual_alpha)
                                if hasattr(trace, 'opacity'):
                                    trace.opacity = individual_alpha
                        else:
                            self.fig_3d.data[idx].visible = False
                    
                    self.fig_3d.update_layout(
                        title=dict(text='', x=0.5),
                        margin=dict(l=0, r=0, b=0, t=0),
                        showlegend=False,
                        updatemenus=[],  # Remove view selection dropdown
                        annotations=[],  # Remove controls hint
                        paper_bgcolor=self.background_color,
                        plot_bgcolor=self.background_color,
                        scene=dict(bgcolor=self.background_color),
                    )
                    
                    html_filename = f'{safe_name}.html'
                    html_path = os.path.join(output_dir, html_filename)
                    self.fig_3d.write_html(html_path, include_plotlyjs='cdn', full_html=True)
                    generated_files['html'].append(html_path)
        
        # ==============================================================================
        # FALLBACK: Traditional loop for kaleido or HTML-only export
        # ==============================================================================
        else:
            for legend_name in tqdm(legend_names, desc='Plotting individuals'):
                # Skip if PNG export already failed
                if png_export_failed and 'png' in output_format and 'html' not in output_format:
                    continue
                
                trace_indices = legend_entries[legend_name]
                safe_name = _sanitize_filename(legend_name)
                
                # Hide all neuron traces, show only this legend's traces + background
                for idx in range(n_traces):
                    if idx in trace_indices or idx in background_indices:
                        self.fig_3d.data[idx].visible = True
                        if idx in trace_indices:
                            trace = self.fig_3d.data[idx]
                            if hasattr(trace, 'color') and trace.color is not None:
                                trace.color = _modify_color_alpha(trace.color, individual_alpha)
                            if hasattr(trace, 'opacity'):
                                trace.opacity = individual_alpha
                    else:
                        self.fig_3d.data[idx].visible = False
                
                self.fig_3d.update_layout(
                    title=dict(text='', x=0.5),
                    margin=dict(l=0, r=0, b=0, t=0),
                    sliders=[],
                    updatemenus=[],
                    annotations=[],  # Remove controls hint
                    showlegend=False,
                    paper_bgcolor=self.background_color,
                    plot_bgcolor=self.background_color,
                    scene=dict(
                        domain=dict(x=[0, 1], y=[0, 1]),  # Full viewport, no margin
                        bgcolor=self.background_color,
                    ),
                )
                
                # Export HTML if requested
                if 'html' in output_format:
                    html_filename = f'{safe_name}.html'
                    html_path = os.path.join(output_dir, html_filename)
                    self.fig_3d.write_html(html_path, include_plotlyjs='cdn', full_html=True)
                    generated_files['html'].append(html_path)
                
                # Export PNG(s) if requested (kaleido path)
                if 'png' in output_format and not png_export_failed:
                    if safe_name not in generated_files['png']:
                        generated_files['png'][safe_name] = []
                    
                    for view_name in views:
                        if png_export_failed:
                            break
                            
                        camera = view_cameras[view_name]
                        self.fig_3d.update_layout(scene_camera=camera)
                        
                        png_filename = f'{view_name}_{safe_name}.png'
                        png_path = os.path.join(output_dir, png_filename)
                        
                        success, msg, final_scale = self._export_png_with_timeout(
                            self.fig_3d, png_path, 
                            width=900, height=900, 
                            scale=scale, 
                            timeout=300,
                            auto_crop=auto_crop,
                            crop_margin=crop_margin
                        )
                        
                        if success:
                            generated_files['png'][safe_name].append((png_path, view_name))
                        else:
                            self._vprint(f'   ⚠️  PNG export failed for {legend_name} ({view_name}): {msg}')
                            png_export_failed = True
                            self._vprint(f'   💡 Skipping remaining exports. Try reducing scale or skeleton_mesh_simplification.')
                            break
        
        # Restore original figure state (visibility, opacity, and colors)
        for idx in range(n_traces):
            self.fig_3d.data[idx].visible = original_visibility[idx]
            # Restore opacity (even if None, to reset any changes)
            if hasattr(self.fig_3d.data[idx], 'opacity'):
                self.fig_3d.data[idx].opacity = original_opacity[idx]
            # Restore color (which contains alpha in RGBA format for navis traces)
            if hasattr(self.fig_3d.data[idx], 'color') and original_colors[idx] is not None:
                self.fig_3d.data[idx].color = original_colors[idx]
        
        # Restore original layout (includes resetting scene domain)
        self.fig_3d.update_layout(original_layout)
        
        # Normalize summary_format
        if summary_format is True:
            # Default to PDF when True
            summary_format = ['pdf']
        elif isinstance(summary_format, str):
            summary_format = [summary_format.lower()]
        elif summary_format:
            summary_format = [f.lower() for f in summary_format]
        else:
            # False or None - no summary
            summary_format = []
        
        # Generate PDF/PPTX summaries if PNG images were created
        if 'png' in output_format and generated_files['png']:
            # Save summaries in parent folder (parallel to individual_profiles/)
            parent_dir = os.path.dirname(output_dir)
            base_title = pdf_title or self.saveas
            
            # Generate PDF if requested
            if 'pdf' in summary_format:
                # For single view, generate one PDF without suffix (organized by view)
                # For multiple views, generate both _by_view and _by_name PDFs
                if len(views) == 1:
                    self._vprint(f'\n📄 Generating PDF summary...')
                    pdf_path = self._create_individual_pdf(
                        output_dir=parent_dir,
                        images_dict=generated_files['png'],
                        images_per_page=pdf_images_per_page,
                        title=base_title,
                        organize_by='view',  # Organize by view for single-view PDF
                        views=views,
                        pdf_suffix='',
                        background_color=self.background_color,
                    )
                    if pdf_path:
                        self._vprint(f'   ✅ PDF saved: {pdf_path}')
                else:
                    self._vprint(f'\n📄 Generating PDF summaries...')
                    # Generate PDF organized by view
                    pdf_path_view = self._create_individual_pdf(
                        output_dir=parent_dir,
                        images_dict=generated_files['png'],
                        images_per_page=pdf_images_per_page,
                        title=base_title,
                        organize_by='view',
                        views=views,
                        pdf_suffix='_by_view',
                        background_color=self.background_color,
                    )
                    if pdf_path_view:
                        self._vprint(f'   ✅ PDF saved: {pdf_path_view}')
                    
                    # Generate PDF organized by name
                    pdf_path_name = self._create_individual_pdf(
                        output_dir=parent_dir,
                        images_dict=generated_files['png'],
                        images_per_page=pdf_images_per_page,
                        title=base_title,
                        organize_by='name',
                        views=views,
                        pdf_suffix='_by_name',
                        background_color=self.background_color,
                    )
                    if pdf_path_name:
                        self._vprint(f'   ✅ PDF saved: {pdf_path_name}')
            
            # Generate PPTX if requested
            if 'pptx' in summary_format:
                if len(views) == 1:
                    self._vprint(f'\n📊 Generating PPTX summary...')
                    pptx_path = self._create_individual_pptx(
                        output_dir=parent_dir,
                        images_dict=generated_files['png'],
                        images_per_page=pdf_images_per_page,
                        title=base_title,
                        organize_by='view',
                        views=views,
                        pptx_suffix='',
                        background_color=self.background_color,
                    )
                    if pptx_path:
                        self._vprint(f'   ✅ PPTX saved: {pptx_path}')
                else:
                    self._vprint(f'\n📊 Generating PPTX summaries...')
                    # Generate PPTX organized by view
                    pptx_path_view = self._create_individual_pptx(
                        output_dir=parent_dir,
                        images_dict=generated_files['png'],
                        images_per_page=pdf_images_per_page,
                        title=base_title,
                        organize_by='view',
                        views=views,
                        pptx_suffix='_by_view',
                        background_color=self.background_color,
                    )
                    if pptx_path_view:
                        self._vprint(f'   ✅ PPTX saved: {pptx_path_view}')
                    
                    # Generate PPTX organized by name
                    pptx_path_name = self._create_individual_pptx(
                        output_dir=parent_dir,
                        images_dict=generated_files['png'],
                        images_per_page=pdf_images_per_page,
                        title=base_title,
                        organize_by='name',
                        views=views,
                        pptx_suffix='_by_name',
                        background_color=self.background_color,
                    )
                    if pptx_path_name:
                        self._vprint(f'   ✅ PPTX saved: {pptx_path_name}')
        
        # Summary
        n_png = sum(len(v) for v in generated_files['png'].values())
        n_html = len(generated_files['html'])
        self._vprint(f'\n✅ Individual plots complete!')
        self._vprint(f'   PNG files: {n_png}')
        self._vprint(f'   HTML files: {n_html}')
        self._vprint(f'   Output folder: {output_dir}')
        
        return output_dir

    def _create_individual_pdf(
        self,
        output_dir: str,
        images_dict: dict,
        images_per_page: tuple = (4, 3),
        title: str = None,
        organize_by: str = 'name',
        views: list = None,
        pdf_suffix: str = '',
        background_color: str = 'white',
    ) -> str | None:
        """
        Create a PDF summary from individual profile PNG images.
        
        Parameters
        ----------
        output_dir : str
            Directory where PDF will be saved.
        images_dict : dict
            Dictionary mapping legend names to list of (image_path, view_name) tuples.
            e.g., {'neuron1': [('/path/front.png', 'front'), ('/path/top.png', 'top')], ...}
        images_per_page : tuple
            (columns, rows) - number of images per page.
        title : str, optional
            Title for the PDF document.
        organize_by : str
            How images are organized: 'name' or 'view'
        views : list, optional
            List of view names for organizing by view
        pdf_suffix : str, optional
            Suffix to add to PDF filename (e.g., '_by_view', '_by_name')
            
        Returns
        -------
        str or None
            Path to created PDF, or None if creation failed.
        """
        try:
            from reportlab.lib.pagesizes import A4, landscape as rl_landscape
            from reportlab.lib.units import inch
            from reportlab.pdfgen import canvas
            from PIL import Image
        except ImportError:
            self._vprint('⚠️  PDF generation requires reportlab and Pillow.')
            self._vprint('   Install with: pip install reportlab Pillow')
            return None
        
        from pathlib import Path
        
        # Natural sort function for rank-based names like r1, r2, ..., r10, r11
        def natural_sort_key(s):
            """Sort strings containing numbers in natural order (r1, r2, ..., r10, r11)."""
            import re
            return [int(c) if c.isdigit() else c.lower() for c in re.split(r'(\d+)', s)]
        
        # Organize images based on organize_by option
        # Group by view when organize_by='view', otherwise by name
        if organize_by == 'view' and views:
            # Group images by view, each view gets its own section
            images_by_category = {}
            for view_name in views:
                images_by_category[view_name] = []
            
            for legend_name, img_info_list in sorted(images_dict.items(), key=lambda x: natural_sort_key(x[0])):
                for img_info in img_info_list:
                    if isinstance(img_info, tuple):
                        img_path, view_name = img_info
                    else:
                        img_path = img_info
                        view_name = views[0] if views else 'front'
                    if os.path.exists(img_path) and view_name in images_by_category:
                        images_by_category[view_name].append((legend_name, img_path, view_name))
        else:
            # Group images by name (default) with natural sorting
            images_by_category = {}
            for legend_name, img_info_list in sorted(images_dict.items(), key=lambda x: natural_sort_key(x[0])):
                if legend_name not in images_by_category:
                    images_by_category[legend_name] = []
                for img_info in img_info_list:
                    if isinstance(img_info, tuple):
                        img_path, view_name = img_info
                    else:
                        img_path = img_info
                        view_name = ''
                    if os.path.exists(img_path):
                        images_by_category[legend_name].append((legend_name, img_path, view_name))
        
        # Flatten but keep category boundaries for page breaks
        all_categories = list(images_by_category.keys())
        
        if not any(images_by_category.values()):
            self._vprint('⚠️  No images found for PDF generation.')
            return None
        
        # Output path
        pdf_path = os.path.join(output_dir, f'individual_profiles_summary{pdf_suffix}.pdf')
        
        # Page setup (landscape A4)
        page_width, page_height = rl_landscape(A4)
        cols, rows = images_per_page
        margin = 0.3 * inch  # Reduced from 0.5 to minimize blank space
        title_height = 20  # Reduced from 25 to minimize blank space
        
        # Determine if dark background for text color
        is_dark_bg = self._is_dark_background(background_color)
        text_color = (1, 1, 1) if is_dark_bg else (0, 0, 0)  # RGB 0-1 for reportlab
        
        # Parse background color for reportlab
        def parse_color_for_reportlab(color_str):
            """Convert CSS color to reportlab RGB tuple (0-1 range)."""
            color_str = color_str.lower().strip()
            if color_str == 'black':
                return (0, 0, 0)
            elif color_str == 'white':
                return (1, 1, 1)
            elif color_str.startswith('#'):
                hex_color = color_str.lstrip('#')
                if len(hex_color) == 3:
                    hex_color = ''.join([c*2 for c in hex_color])
                r, g, b = tuple(int(hex_color[i:i+2], 16) / 255 for i in (0, 2, 4))
                return (r, g, b)
            elif color_str.startswith('rgb'):
                import re
                nums = re.findall(r'\d+', color_str)
                if len(nums) >= 3:
                    return tuple(int(n) / 255 for n in nums[:3])
            # Default to white for unknown colors
            return (1, 1, 1)
        
        bg_color = parse_color_for_reportlab(background_color)
        
        # Calculate cell dimensions
        usable_width = page_width - 2 * margin
        usable_height = page_height - 2 * margin - title_height
        cell_width = usable_width / cols
        cell_height = usable_height / rows
        
        # Create PDF
        c = canvas.Canvas(pdf_path, pagesize=(page_width, page_height))
        
        images_per_full_page = cols * rows
        
        # Process each category separately (don't mix categories on same page)
        for category_name in all_categories:
            category_images = images_by_category[category_name]
            if not category_images:
                continue
            
            # Calculate pages needed for this category
            total_pages_for_category = (len(category_images) + images_per_full_page - 1) // images_per_full_page
            
            for page_idx in range(total_pages_for_category):
                # Draw background rectangle
                c.setFillColorRGB(*bg_color)
                c.rect(0, 0, page_width, page_height, fill=True, stroke=False)
                
                # Page title with category info
                c.setFillColorRGB(*text_color)
                c.setFont("Helvetica-Bold", 14)
                if organize_by == 'view':
                    # Use '{view} view' as title when organized by view
                    page_title = f"{category_name} view"
                else:
                    # Use layer_name as title when organized by name
                    page_title = str(category_name)
                if total_pages_for_category > 1:
                    page_title += f" ({page_idx + 1}/{total_pages_for_category})"
                c.drawCentredString(page_width / 2, page_height - margin - 5, page_title)
                
                # Get images for this page from the category
                start_idx = page_idx * images_per_full_page
                end_idx = min(start_idx + images_per_full_page, len(category_images))
                page_images = category_images[start_idx:end_idx]
                
                # Draw images
                for i, (legend_name, img_path, view_name) in enumerate(page_images):
                    row = i // cols
                    col = i % cols
                    
                    # Calculate position
                    x = margin + col * cell_width
                    y = page_height - margin - title_height - (row + 1) * cell_height
                    
                    try:
                        with Image.open(img_path) as img:
                            img_width, img_height = img.size
                            
                            # Calculate scaling - minimize padding between images
                            padding = 0  # Reduced from 5 to minimize blank space
                            label_height = 12  # Reduced from 15 to minimize blank space
                            max_width = cell_width - 2 * padding
                            max_height = cell_height - 2 * padding - label_height
                            
                            scale_w = max_width / img_width
                            scale_h = max_height / img_height
                            scale_factor = min(scale_w, scale_h)
                            
                            draw_width = img_width * scale_factor
                            draw_height = img_height * scale_factor
                            
                            # Center horizontally in cell, leave space at top for label
                            draw_x = x + (cell_width - draw_width) / 2
                            draw_y = y + (cell_height - label_height - draw_height) / 2
                            
                            # Draw image
                            c.drawImage(
                                img_path,
                                draw_x, draw_y,
                                width=draw_width,
                                height=draw_height,
                                preserveAspectRatio=True
                            )
                            
                            # Draw label on TOP of image
                            c.setFillColorRGB(*text_color)
                            c.setFont("Helvetica", 12)
                            label = str(legend_name)
                            if view_name and organize_by != 'view':
                                # Only add view suffix if not organizing by view
                                label += f" ({view_name})"
                            if len(label) > 30:
                                label = label[:27] + '...'
                            label_y = y + cell_height - label_height + 2  # Reduced from 3
                            c.drawCentredString(x + cell_width / 2, label_y, label)
                            
                    except Exception as e:
                        self._vprint(f'   ⚠️  Could not process: {img_path} - {e}', level='full')
                
                c.showPage()
        
        c.save()
        return pdf_path

    def _create_individual_pptx(
        self,
        output_dir: str,
        images_dict: dict,
        images_per_page: tuple = (4, 3),
        title: str = None,
        organize_by: str = 'name',
        views: list = None,
        pptx_suffix: str = '',
        label_fontsize: int = 20,
        title_fontsize: int = 24,
        background_color: str = 'white',
    ) -> str | None:
        """
        Create a PPTX summary from individual profile PNG images.
        
        Parameters
        ----------
        output_dir : str
            Directory where PPTX will be saved.
        images_dict : dict
            Dictionary mapping legend names to list of (image_path, view_name) tuples.
            e.g., {'neuron1': [('/path/front.png', 'front'), ('/path/top.png', 'top')], ...}
        images_per_page : tuple
            (columns, rows) - number of images per slide.
        title : str, optional
            Title for the PPTX document.
        organize_by : str
            How images are organized: 'name' or 'view'
        views : list, optional
            List of view names for organizing by view
        pptx_suffix : str, optional
            Suffix to add to PPTX filename (e.g., '_by_view', '_by_name')
        label_fontsize : int, default 20
            Font size for image labels in points.
        title_fontsize : int, default 24
            Font size for slide titles in points.
        background_color : str, default 'white'
            Background color for slides. Supports 'black', 'white', hex codes, etc.
            
        Returns
        -------
        str or None
            Path to created PPTX, or None if creation failed.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
        except ImportError:
            self._vprint('⚠️  PPTX generation requires python-pptx.')
            self._vprint('   Install with: pip install python-pptx')
            return None
        
        from PIL import Image
        
        # Natural sort function for rank-based names like r1, r2, ..., r10, r11
        def natural_sort_key(s):
            """Sort strings containing numbers in natural order (r1, r2, ..., r10, r11)."""
            import re
            return [int(c) if c.isdigit() else c.lower() for c in re.split(r'(\d+)', s)]
        
        # Organize images based on organize_by option
        if organize_by == 'view' and views:
            images_by_category = {}
            for view_name in views:
                images_by_category[view_name] = []
            
            for legend_name, img_info_list in sorted(images_dict.items(), key=lambda x: natural_sort_key(x[0])):
                for img_info in img_info_list:
                    if isinstance(img_info, tuple):
                        img_path, view_name = img_info
                    else:
                        img_path = img_info
                        view_name = views[0] if views else 'front'
                    if os.path.exists(img_path) and view_name in images_by_category:
                        images_by_category[view_name].append((legend_name, img_path, view_name))
        else:
            images_by_category = {}
            for legend_name, img_info_list in sorted(images_dict.items(), key=lambda x: natural_sort_key(x[0])):
                if legend_name not in images_by_category:
                    images_by_category[legend_name] = []
                for img_info in img_info_list:
                    if isinstance(img_info, tuple):
                        img_path, view_name = img_info
                    else:
                        img_path = img_info
                        view_name = ''
                    if os.path.exists(img_path):
                        images_by_category[legend_name].append((legend_name, img_path, view_name))
        
        all_categories = list(images_by_category.keys())
        
        if not any(images_by_category.values()):
            self._vprint('⚠️  No images found for PPTX generation.')
            return None
        
        # Output path
        pptx_path = os.path.join(output_dir, f'individual_profiles_summary{pptx_suffix}.pptx')
        
        # Determine if dark background for text color
        is_dark_bg = self._is_dark_background(background_color)
        font_color = RGBColor(255, 255, 255) if is_dark_bg else RGBColor(0, 0, 0)
        
        # Parse background color for pptx
        def parse_color_for_pptx(color_str):
            """Convert CSS color to pptx RGBColor."""
            color_str = color_str.lower().strip()
            if color_str == 'black':
                return RGBColor(0, 0, 0)
            elif color_str == 'white':
                return RGBColor(255, 255, 255)
            elif color_str.startswith('#'):
                hex_color = color_str.lstrip('#')
                if len(hex_color) == 3:
                    hex_color = ''.join([c*2 for c in hex_color])
                r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                return RGBColor(r, g, b)
            elif color_str.startswith('rgb'):
                import re
                nums = re.findall(r'\d+', color_str)
                if len(nums) >= 3:
                    return RGBColor(int(nums[0]), int(nums[1]), int(nums[2]))
            # Default to white for unknown colors
            return RGBColor(255, 255, 255)
        
        bg_color = parse_color_for_pptx(background_color)
        
        # Slide setup (widescreen 16:9)
        slide_width, slide_height = 13.333, 7.5
        cols, rows = images_per_page
        margin = 0.3  # inches
        title_height_inches = 0.5
        label_height_inches = (label_fontsize / 72) * 1.5
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)
        blank_layout = prs.slide_layouts[6]  # Blank slide
        
        # Calculate cell dimensions
        usable_width = slide_width - 2 * margin
        usable_height = slide_height - margin - title_height_inches - margin
        cell_width = usable_width / cols
        cell_height = usable_height / rows
        
        images_per_full_page = cols * rows
        
        # Process each category separately
        for category_name in all_categories:
            category_images = images_by_category[category_name]
            if not category_images:
                continue
            
            total_pages_for_category = (len(category_images) + images_per_full_page - 1) // images_per_full_page
            
            for page_idx in range(total_pages_for_category):
                slide = prs.slides.add_slide(blank_layout)
                
                # Set slide background color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = bg_color
                
                # Build title
                if organize_by == 'view':
                    slide_title = f"{category_name} view"
                else:
                    slide_title = str(category_name)
                if total_pages_for_category > 1:
                    slide_title += f" ({page_idx + 1}/{total_pages_for_category})"
                
                # Add title
                txBox = slide.shapes.add_textbox(
                    Inches(margin),
                    Inches(margin / 2),
                    Inches(slide_width - 2 * margin),
                    Inches(title_height_inches)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(title_fontsize)
                p.font.bold = True
                p.font.color.rgb = font_color
                p.alignment = PP_ALIGN.CENTER
                
                # Get images for this page
                start_idx = page_idx * images_per_full_page
                end_idx = min(start_idx + images_per_full_page, len(category_images))
                page_images = category_images[start_idx:end_idx]
                
                content_top = margin + title_height_inches
                
                for i, (legend_name, img_path, view_name) in enumerate(page_images):
                    row = i // cols
                    col = i % cols
                    
                    cell_left = margin + col * cell_width
                    cell_top = content_top + row * cell_height
                    
                    try:
                        with Image.open(img_path) as img:
                            img_width, img_height = img.size
                            
                            # Calculate scaling
                            max_width = cell_width - 0.1
                            max_height = cell_height - label_height_inches - 0.1
                            
                            scale_w = max_width / (img_width / 96)
                            scale_h = max_height / (img_height / 96)
                            scale_factor = min(scale_w, scale_h)
                            
                            final_width = (img_width / 96) * scale_factor
                            final_height = (img_height / 96) * scale_factor
                            
                            # Center image in cell
                            img_left = cell_left + (cell_width - final_width) / 2
                            img_top = cell_top + (cell_height - label_height_inches - final_height) / 2
                            
                            # Add image
                            slide.shapes.add_picture(
                                img_path,
                                Inches(img_left),
                                Inches(img_top),
                                Inches(final_width),
                                Inches(final_height)
                            )
                            
                            # Add label
                            label = str(legend_name)
                            if view_name and organize_by != 'view':
                                label += f" ({view_name})"
                            max_chars = int(cell_width * 8)
                            if len(label) > max_chars:
                                label = label[:max_chars-3] + '...'
                            
                            txBox = slide.shapes.add_textbox(
                                Inches(cell_left),
                                Inches(cell_top + cell_height - label_height_inches),
                                Inches(cell_width),
                                Inches(label_height_inches)
                            )
                            tf = txBox.text_frame
                            p = tf.paragraphs[0]
                            p.text = label
                            p.font.size = Pt(label_fontsize)
                            p.font.color.rgb = font_color
                            p.alignment = PP_ALIGN.CENTER
                            
                    except Exception as e:
                        self._vprint(f'   ⚠️  Could not process: {img_path} - {e}', level='full')
        
        prs.save(pptx_path)
        return pptx_path

    def _to_rgba(self, color, alpha=None):
        # Convert color to uint8 RGBA for trimesh.
        import numpy as np

        try:
            r, g, b, a = extract_rgba_tuple(color, default_alpha=1.0)
            if alpha is not None:
                a = float(alpha)
            c = np.array([r, g, b, max(0.0, min(1.0, a))], dtype=float)
        except:
            c = np.array([128.0, 128.0, 128.0, 255.0], dtype=float)
            return c.astype(np.uint8)

        return np.array([
            int(np.clip(round(c[0]), 0, 255)),
            int(np.clip(round(c[1]), 0, 255)),
            int(np.clip(round(c[2]), 0, 255)),
            int(np.clip(round(c[3] * 255), 0, 255)),
        ], dtype=np.uint8)

    def _extract_trimesh(self, mesh):
        """Extract a trimesh.Trimesh from navis/trimesh-like objects."""
        import trimesh

        if mesh is None:
            return None
        if isinstance(mesh, trimesh.Trimesh):
            return mesh
        if isinstance(mesh, go.Mesh3d):
            vertices = np.column_stack([np.asarray(mesh.x), np.asarray(mesh.y), np.asarray(mesh.z)])
            faces = np.column_stack([np.asarray(mesh.i), np.asarray(mesh.j), np.asarray(mesh.k)]).astype(np.int64)
            if len(vertices) == 0 or len(faces) == 0:
                return None
            return trimesh.Trimesh(vertices=vertices, faces=faces, process=False)
        if hasattr(mesh, 'trimesh'):
            return mesh.trimesh
        if hasattr(mesh, 'mesh') and isinstance(mesh.mesh, trimesh.Trimesh):
            return mesh.mesh
        if hasattr(mesh, 'vertices') and hasattr(mesh, 'faces'):
            return trimesh.Trimesh(vertices=mesh.vertices, faces=mesh.faces, process=False)
        return None

    def _append_exportable_mesh(self, mesh, color=None, alpha=None, name=None, role='mesh'):
        """Append a colored mesh copy to the cached export scene."""
        try:
            tm = self._extract_trimesh(mesh)
            if tm is None:
                return False

            tm = tm.copy()
            rgba = self._to_rgba(color if color is not None else 'gray', alpha=alpha)
            tm.visual.face_colors = rgba
            tm.metadata['export_rgba'] = rgba.tolist()
            if name:
                tm.metadata['name'] = name
            tm.metadata['export_role'] = role
            self.exportable_meshes.append(tm)
            return True
        except Exception:
            return False

    def _get_glb_export_rgba(self, mesh):
        """Return the exact cached RGBA color for GLB export."""
        metadata_rgba = getattr(mesh, 'metadata', {}).get('export_rgba')
        if metadata_rgba is not None:
            return np.asarray(metadata_rgba, dtype=np.uint8)

        face_colors = np.asarray(mesh.visual.face_colors)
        if face_colors.ndim == 2 and len(face_colors) > 0:
            return face_colors[0].astype(np.uint8).copy()

        return np.array([128, 128, 128, 255], dtype=np.uint8)

    def _prepare_glb_mesh(self, mesh, geometry_name):
        """Convert a cached mesh into a GLB-friendly mesh with explicit PBR material."""
        import trimesh

        tm = mesh.copy()
        rgba = self._get_glb_export_rgba(tm)
        material = trimesh.visual.material.PBRMaterial(
            name=geometry_name,
            baseColorFactor=rgba,
            metallicFactor=0.0,
            roughnessFactor=1.0,
            doubleSided=True,
            alphaMode='BLEND' if rgba[3] < 255 else 'OPAQUE',
        )
        tm.visual = trimesh.visual.texture.TextureVisuals(material=material)
        tm.metadata.update(getattr(mesh, 'metadata', {}))
        return tm

    def export_3d_model(self, filename=None, format='glb'):
        """Export the 3D scene to a model file (GLB, OBJ, STL).
        
        Exports the current scene (neurons, synapses, ROIs) to a 3D model file.
        Useful for importing into Blender, Unity, or other 3D software.
        
        Parameters
        ----------
        filename : str, optional
            Output filename. If None, uses self.saveas.
        format : str, default 'glb'
            Output format: 'glb' (binary glTF, recommended), 'obj', 'stl', 'ply'.
        
        Returns
        -------
        str
            Path to the saved file.
        """
        if not self.exportable_meshes and not hasattr(self, 'fig_3d'):
            self._vprint('⚠️  No meshes to export. Run plot_neurons() first.')
            return None
            
        if filename is None:
            filename = self.saveas
            
        # Ensure extension
        if not filename.lower().endswith(f'.{format}'):
            filename += f'.{format}'
            
        filepath = os.path.join(self.save_folder, filename)
        
        self._vprint(f'Exporting 3D model to {filepath}...', level='full')
        
        try:
            import trimesh
            
            # Collect all meshes
            meshes = []
            if hasattr(self, 'exportable_meshes'):
                for index, mesh in enumerate(self.exportable_meshes):
                    geometry_name = getattr(mesh, 'metadata', {}).get('name', f'mesh_{index}')
                    meshes.append(self._prepare_glb_mesh(mesh, geometry_name))
                
            if not meshes:
                self._vprint('⚠️  No meshes found to export.')
                return None
                
            # Combine scene
            scene = trimesh.Scene()
            for index, mesh in enumerate(meshes):
                geometry_name = getattr(mesh, 'metadata', {}).get('name', f'mesh_{index}')
                scene.add_geometry(mesh, geom_name=geometry_name, node_name=geometry_name)
            
            # Export
            scene.export(filepath)
            self._vprint(f'✓ Saved {filepath}')
            return filepath
            
        except ImportError:
            self._vprint('⚠️  trimesh not installed. Cannot export 3D model.')
            self._vprint('   pip install trimesh')
            return None
        except Exception as e:
            self._vprint(f'⚠️  Export failed: {e}')
            return None

    def export_video(self, fps=30, degree_per_frame=1.0, rotate='horizontal', rotate_plane=None, 
                    view_direction=None, view_distance=None, synapse_size=1, 
                    html_file=None, output_dir=None, use_existing_images=True, 
                    export_gif=True, gif_scale=0.2, gif_optimize=True,
                    auto_crop=True, crop_margin=30, export_method: str = None, **kwargs):
        '''
        Export a rotating 3D visualization to MP4 video.
        
        Can be used in two modes:
        1. After plot_neurons(): Uses the current figure in memory
        2. Standalone with html_file: Loads figure from existing HTML file
        
        For standalone usage without VisualizeSkeleton initialization, use the
        module-level function `export_video_from_html()` instead.
        
        Parameters
        ----------
        fps : int, default 30
            Frames per second for the output video.
        degree_per_frame : float, default 1.0
            Rotation angle in degrees per frame. Controls rotation speed.
            - 1.0 → 360 frames for full rotation (12 sec video at 30 fps)
            - 2.0 → 180 frames for full rotation (6 sec video at 30 fps)
            - 0.5 → 720 frames for full rotation (24 sec video at 30 fps)
        rotate : str, default 'horizontal'
            Rotation direction:
            - 'horizontal': Rotate around Y-axis (turntable motion)
            - 'vertical': Rotate around X-axis (tumbling motion)
        rotate_plane : str, optional (deprecated)
            Legacy parameter. Use 'rotate' instead.
            Plane to rotate: 'xy', 'xz', or 'yz'.
        view_direction : tuple, optional, default (1, -1)
            Camera direction multipliers for sin/cos components.
            Options: (1, 1), (1, -1), (-1, 1), or (-1, -1).
        view_distance : float, optional, default 2.2
            Relative camera distance from center (1.0 = close, 3.0 = far).
        synapse_size : int, default 1
            Size of synapse markers in the video (1-10 recommended).
        html_file : str, optional
            Path to existing HTML file to load figure data from.
            Enables standalone usage without calling plot_neurons() first.
            Example: '/path/to/my_neurons.html'
        output_dir : str, optional
            Directory to save video output. If None and html_file is provided,
            uses the directory containing the html_file.
            If None and using plot_neurons(), uses self.save_folder.
        use_existing_images : bool, default True
            If True, skip rendering and reuse cached images from previous export.
            Useful for regenerating video with different fps without re-rendering.
        export_gif : bool, default True
            If True, automatically convert videos to GIF format after export.
        gif_scale : float, default 0.2
            Scale factor for GIF resolution (0.1-1.0). Lower values = smaller file size.
            Example: 0.2 = 20% of original video resolution.
        gif_optimize : bool, default True
            Enable GIF compression optimization for smaller file sizes.
        auto_crop : bool, default True
            If True, automatically crop whitespace/background from each frame
            with consistent sizing across all frames (computes max projection bounds).
            This ensures uniform frame dimensions for smooth video playback.
        crop_margin : int, default 30
            Margin (in pixels) to preserve around content when auto_crop=True.
            Has no effect unless auto_crop=True.
        export_method : str, optional
            Export method to use for frame rendering: 'webdriver' or 'kaleido'.
            If None, uses the class-level export_method attribute.
            - 'webdriver': Uses Chrome WebDriver (recommended for large figures)
            - 'kaleido': Uses Kaleido static image export (default)
        **kwargs : dict
            Additional arguments for plotly write_image():
            - scale : int - Resolution multiplier. Defaults to min(self.export_scale, 3).
              For video, scale is capped at 3 for reasonable performance.
              To use scale=4 (very slow), pass it explicitly: scale=4
            - width : int - Video width in pixels (default 1200)
            - height : int - Video height in pixels (default 900)
        
        Returns
        -------
        int
            0 on success, 1 on failure
        
        Output Files
        ------------
        - {output_dir}/pics_{fps}fps_{rotate_plane}/ : Cached frame images
        - {output_dir}/{name}_video_h_forward.mp4 : Forward rotation video (horizontal)
        - {output_dir}/{name}_video_h_backward.mp4 : Reverse rotation video (horizontal)
        - {output_dir}/{name}_video_v_forward.mp4 : Forward rotation video (vertical)
        - {output_dir}/{name}_video_v_backward.mp4 : Reverse rotation video (vertical)
        - {output_dir}/{name}_video_h_forward.gif : Forward rotation GIF (horizontal, if export_gif=True)
        - {output_dir}/{name}_video_h_backward.gif : Reverse rotation GIF (horizontal, if export_gif=True)
        - {output_dir}/{name}_video_v_forward.gif : Forward rotation GIF (vertical, if export_gif=True)
        - {output_dir}/{name}_video_v_backward.gif : Reverse rotation GIF (vertical, if export_gif=True)
        
        Examples
        --------
        # Mode 1: After plot_neurons()
        vs = VisualizeSkeleton(dataset='hemibrain:v1.2.1', neuron_layers=['EB'])
        vs.plot_neurons()
        vs.export_video(fps=30, degree_per_frame=1.0)
        
        # Faster rotation (shorter video)
        vs.export_video(fps=30, degree_per_frame=2.0)
        
        # Vertical rotation
        vs.export_video(fps=30, rotate='vertical')
        
        # High quality export
        vs.export_video(fps=30, scale=4, width=1920, height=1080)
        
        # Mode 2: From existing HTML file (output to same directory)
        vs.export_video(html_file='/path/to/existing_plot.html')
        
        # Mode 3: Standalone function (no VisualizeSkeleton needed)
        from visualize_skeleton import export_video_from_html
        export_video_from_html('/path/to/plot.html', fps=30, degree_per_frame=1.0)
        
        # Reuse cached images (fast video regeneration)
        vs.export_video(fps=60, use_existing_images=True)
        '''
        # Handle rotate parameter - overrides rotate_plane
        if rotate == 'horizontal':
            rotate_plane = 'xz'  # Rotate around vertical (Y) axis
        elif rotate == 'vertical':
            rotate_plane = 'yz'  # Rotate around horizontal (X) axis
        elif rotate_plane is None:
            # Default to horizontal rotation if neither specified
            rotate_plane = 'xz'
        # else: use the explicitly provided rotate_plane
        
        # Determine which export method to use
        actual_export_method = export_method if export_method else self.export_method
        
        if view_direction is None:
            if 'manc' in self.dataset.lower():
                view_direction = (1, 1)  # MANC: +Z is front
            else:
                view_direction = (1, -1)
        if view_distance is None:
            view_distance = 2.2
        
        # Warn if crop_margin is set but auto_crop is disabled
        if crop_margin != 30 and not auto_crop:
            self._vprint(f'⚠️  crop_margin={crop_margin} has no effect without auto_crop=True')
            self._vprint(f'   Add auto_crop=True to enable automatic frame cropping')
        
        # Warn about kaleido limitations for video export
        if actual_export_method == 'kaleido':
            self._vprint('💡 Tip: For video export, consider using export_method="webdriver" for better stability.')
            self._vprint('   Kaleido can timeout on complex figures. WebDriver is more reliable for animations.')
        
        # Set default scale from self.export_scale if not specified in kwargs
        # For kaleido: cap at scale=3 for reasonable performance (scale=4 is very slow)
        # For webdriver/webdriver-fast: cap at scale=5 for stability
        # Users can explicitly set higher scale to override the cap
        if kwargs.get('scale') is None and kwargs.get('width') is None and kwargs.get('height') is None:
            if actual_export_method in ('webdriver', 'webdriver-fast'):
                # WebDriver can handle higher scales - cap at 5
                default_scale = min(getattr(self, 'export_scale', 5), 5)
            else:
                # Kaleido - cap at 3
                default_scale = min(getattr(self, 'export_scale', 3), 3)
            kwargs['scale'] = default_scale
        elif kwargs.get('scale') is not None:
            # User explicitly provided scale - check if > cap and warn
            user_scale = kwargs['scale']
            cap = 5 if actual_export_method in ('webdriver', 'webdriver-fast') else 3
            if user_scale > cap:
                self._vprint(f'⚠️  Using scale={user_scale} for video export (uncapped). This may be slow.')
            elif user_scale > getattr(self, 'export_scale', 3):
                # User requested higher than default, allow it
                pass
        
        # Use explicit degree_per_frame instead of calculating from fps
        step = degree_per_frame
        
        # Determine output directory and filename
        if output_dir is not None:
            save_folder = output_dir
            # Extract filename from html_file or use default
            if html_file is not None:
                saveas = os.path.splitext(os.path.basename(html_file))[0]
            else:
                saveas = 'video_export'
            os.makedirs(save_folder, exist_ok=True)
        elif html_file is not None:
            # Use the directory containing the html_file
            save_folder = os.path.dirname(os.path.abspath(html_file))
            saveas = os.path.splitext(os.path.basename(html_file))[0]
        elif hasattr(self, 'save_folder') and self.save_folder:
            save_folder = self.save_folder
            saveas = self.saveas if hasattr(self, 'saveas') and self.saveas else 'video_export'
        else:
            raise ValueError(
                'No output directory specified. Either:\n'
                '  1. Run plot_neurons() first, or\n'
                '  2. Provide html_file parameter (output goes to same directory), or\n'
                '  3. Provide output_dir parameter explicitly'
            )
        
        # Load figure from existing HTML file if provided (OPTIMIZATION)
        if html_file is not None:
            self._vprint(f'📂 Loading figure from existing HTML: {html_file}')
            if not os.path.exists(html_file):
                raise FileNotFoundError(f'HTML file not found: {html_file}')
            
            # Read and parse the HTML file to extract figure data
            import plotly.io as pio
            try:
                fig_loaded = pio.read_html(html_file)
                fig_traces = fig_loaded.data
                self._vprint(f'✓ Loaded {len(fig_traces)} traces from HTML file')
            except Exception as e:
                raise RuntimeError(f'Failed to load figure from HTML: {e}')
        else:
            # Use current figure
            if not hasattr(self, 'fig_path') or not os.path.exists(self.fig_path+'.html'):
                raise RuntimeError(
                    'No figure found. Either run plot_neurons() first or provide html_file parameter.'
                )
            html_size = os.path.getsize(self.fig_path+'.html') / 1024 / 1024 # in MB
            html_size_cap = self._get_html_size_cap()
            if html_size > html_size_cap:
                self._vprint(f'⚠️  Figure is large ({html_size:.1f} MB). Rendering may be slow.')
                if actual_export_method == 'kaleido':
                    self._vprint(f'   ⚠️  Kaleido may timeout on figures > {html_size_cap}MB.')
                    self._vprint(f'   💡 Recommended: Use export_method="webdriver" for large figures.')
                    self._vprint(f'   💡 Or increase skeleton_mesh_simplification (e.g., 0.97-0.99).')
                else:
                    self._vprint(f'   Consider using higher skeleton_mesh_simplification if export is slow.')
            
            # Use simplified HTML file if available from previous export (e.g., export_views)
            simplified_fig = getattr(self, '_simplified_export_fig', None)
            if simplified_fig is not None:
                self._vprint(f'   ✓ Using simplified figure from previous export')
                fig_traces = simplified_fig.data
            else:
                # Use original figure
                fig_traces = self.fig_3d.data
        # Configure figure for video export
        for trace in fig_traces:
            trace.showlegend = False
            if hasattr(trace,'marker'):
                trace.marker.size = synapse_size
        
        fig_layout = go.Layout(
            margin=dict(l=0, r=0, b=0, t=0, pad=0),
        )
        fig_new = go.Figure(data=fig_traces, layout=fig_layout)
        
        # Set camera parameters - always use frontal view for consistency
        scene_camera_parameters = dict(
            up=dict(x=0, y=-1, z=0),
            eye=dict(x=0, y=0, z=-view_distance),
        )
        
        if 'hemibrain' in self.dataset.lower() and self.brain_mesh == 'template':
             scene_camera_parameters = dict(
                up=dict(x=0, y=0, z=-1),
                eye=dict(x=0, y=view_distance, z=0),  # +Y is front (anterior)
            )
        
        # Adjust for MANC (Male Adult Nerve Cord)
        if 'manc' in self.dataset.lower():
             # MANC: Anterior is +Z, Posterior is -Z. Dorsal is -Y.
             # "Front" view should be looking from Anterior -> Posterior, so eye at +Z.
             scene_camera_parameters = dict(
                up=dict(x=0, y=-1, z=0),
                eye=dict(x=0, y=0, z=view_distance),
            )

        
        fig_new.update_layout(
            sliders=[],  # Remove sliders for cleaner video
            updatemenus=[],  # Remove view selection dropdown
            annotations=[],  # Remove controls hint
            paper_bgcolor=self.background_color,
            plot_bgcolor=self.background_color,
            scene=dict(
                dragmode='orbit',
                xaxis={'visible':False}, 
                yaxis={'visible':False},
                zaxis={'visible':False},
                bgcolor=self.background_color,
            ),
            scene_camera=scene_camera_parameters,
            legend=dict(
                bgcolor=self.background_color,
                font=dict(color='white' if self._is_dark_background() else 'black'),
            ),
        )
        
        # Set up image folder
        pic_folder = os.path.join(save_folder, f'pics_{fps}fps_{rotate_plane}')
        
        # Calculate rotation steps
        if step > 0:
            steps_to_write = np.linspace(0, 360, int(360/step), endpoint=False)
        elif step < 0:
            steps_to_write = np.linspace(360, 0, int(360/step), endpoint=False)
        
        # OPTIMIZATION: Skip image rendering if use_existing_images=True
        # For WebDriver modes, also check for partial completion (resume support)
        resume_mode = False
        if use_existing_images and os.path.exists(pic_folder):
            existing_images = [f for f in os.listdir(pic_folder) 
                             if f.startswith('deg_') and f.endswith('.jpeg') 
                             and os.path.getsize(os.path.join(pic_folder, f)) >= 1024]
            if len(existing_images) == len(steps_to_write):
                self._vprint(f'✓ Using {len(existing_images)} existing images from {pic_folder}')
                self._vprint(f'  Skipping image rendering (use_existing_images=True)')
            elif len(existing_images) > 0 and actual_export_method in ('webdriver', 'webdriver-fast'):
                # Partial completion - enable resume mode for WebDriver
                self._vprint(f'🔄 Found {len(existing_images)}/{len(steps_to_write)} completed frames')
                self._vprint(f'   Resuming from frame {len(existing_images) + 1}...')
                resume_mode = True
                use_existing_images = False  # Continue rendering
            else:
                self._vprint(f'⚠️  Found {len(existing_images)} images but need {len(steps_to_write)}')
                self._vprint(f'  Re-rendering images...')
                use_existing_images = False
        else:
            use_existing_images = False
        
        # Render images if needed
        if not use_existing_images:
            if os.path.exists(pic_folder) and not resume_mode:
                # Only delete folder if NOT in resume mode
                shutil.rmtree(pic_folder)
            if not os.path.exists(pic_folder):
                os.makedirs(pic_folder)
            
            self._vprint(f'🎬 Rendering {len(steps_to_write)} frames at {fps} fps...')
            self._vprint(f'   Resolution: scale={kwargs.get("scale", "auto")}', end='')
            if 'width' in kwargs and 'height' in kwargs:
                self._vprint(f', size={kwargs["width"]}x{kwargs["height"]}')
            else:
                self._vprint('')
            
            # Ensure dimensions are set to avoid blank images if not provided
            if 'width' not in kwargs: kwargs['width'] = 1200
            if 'height' not in kwargs: kwargs['height'] = 900
            
            t0 = time.time()
            current_scale = kwargs.get('scale', 2)
            frame_export_failed = False
            use_kaleido_fallback = False
            
            # Choose export method
            if actual_export_method in ('webdriver', 'webdriver-fast'):
                # WebDriver method - efficient: open browser once for all frames
                # Uses canvas.toDataURL() for fast, high-quality capture
                self._vprint(f'   Using WebDriver for frame export (canvas.toDataURL)...')
                
                # Get render_wait from attribute (None = auto-calibrate, 0 = fastest)
                render_wait = getattr(self, 'webdriver_render_wait', None)
                if render_wait is not None:
                    self._vprint(f'   Render wait: {render_wait}s')
                
                # Save figure to temp HTML
                temp_html = os.path.join(pic_folder, '_temp_video.html')
                fig_new.write_html(temp_html, auto_open=False, 
                                   include_plotlyjs='cdn',
                                   config={'displayModeBar': False})
                
                # Only save simplified HTML copy if we actually used a simplified figure
                # Check if _simplified_export_fig exists AND was used for this export
                simplified_fig = getattr(self, '_simplified_export_fig', None)
                if simplified_fig is not None:
                    html_size_mb = os.path.getsize(temp_html) / 1024 / 1024
                    simplified_html_path = os.path.join(save_folder, f"{saveas}_simplified.html")
                    if not os.path.exists(simplified_html_path):
                        import shutil
                        shutil.copy(temp_html, simplified_html_path)
                        self._vprint(f'   ✓ Saved simplified HTML: {os.path.basename(simplified_html_path)} ({html_size_mb:.1f}MB)')
                
                # Retry logic for Chrome crashes - with RESUME capability
                # Retries reset when progress is made past previous crash point
                max_retries = 3
                webdriver_success = False
                last_error = None
                resume_from_frame = 0  # Track where to resume from
                last_crash_frame = -1  # Track where last crash occurred
                consecutive_crashes_at_same_point = 0  # Only count crashes at same position
                
                while consecutive_crashes_at_same_point < max_retries:
                    # Check how many frames were completed before this attempt
                    existing_frames = [f for f in os.listdir(pic_folder) 
                                      if f.startswith('deg_') and f.endswith('.jpeg')
                                      and os.path.getsize(os.path.join(pic_folder, f)) >= 1024]
                    frames_completed = len(existing_frames)
                    
                    if frames_completed > resume_from_frame:
                        # Progress was made - this is a resume, not a retry at same point
                        resume_from_frame = frames_completed
                        if consecutive_crashes_at_same_point > 0:
                            self._vprint(f'\n   🔄 RESUMING from frame {frames_completed + 1} ({frames_completed} frames saved)...')
                        
                        # Brief pause before resume to let system resources settle
                        time.sleep(2)
                        time.sleep(2)
                    
                    try:
                        with WebDriverExportSession(
                            width=kwargs['width'], height=kwargs['height'],
                            scale=current_scale,
                            timeout=300,
                            render_wait=render_wait
                        ) as session:
                            session.load_html(temp_html, wait_for_render=True, render_wait=3, background_color=self.background_color)
                            self._vprint(f'   ✓ HTML loaded in browser (render_wait={session._render_wait:.2f}s)')
                            
                            # Get initial camera from HTML - this defines "front view"
                            initial_camera = session.get_current_camera()
                            if initial_camera:
                                initial_eye = initial_camera.get('eye', {'x': 0, 'y': 0, 'z': -view_distance})
                                initial_up = initial_camera.get('up', {'x': 0, 'y': -1, 'z': 0})
                                self._vprint(f'   Initial camera: eye=({initial_eye.get("x", 0):.2f}, {initial_eye.get("y", 0):.2f}, {initial_eye.get("z", 0):.2f})')
                            else:
                                initial_eye = {'x': 0, 'y': 0, 'z': -view_distance}
                                initial_up = {'x': 0, 'y': -1, 'z': 0}
                            
                            # Compute view distance from initial camera
                            cam_distance = np.sqrt(
                                initial_eye.get('x', 0)**2 + 
                                initial_eye.get('y', 0)**2 + 
                                initial_eye.get('z', 0)**2
                            )
                            if cam_distance < 0.1:
                                cam_distance = view_distance
                            
                            for i, deg in enumerate(steps_to_write):
                                if frame_export_failed:
                                    break
                                
                                # Skip frames that are already completed (resume support)
                                fig_path = os.path.join(pic_folder, f'deg_{deg:.1f}.jpeg')
                                if os.path.exists(fig_path) and os.path.getsize(fig_path) >= 1024:
                                    # Frame already exists and is valid, skip
                                    continue
                                
                                rad_i = np.deg2rad(deg)
                                sin_val = np.sin(rad_i)
                                cos_val = np.cos(rad_i)
                                
                                # Small offset to avoid gimbal lock
                                offset = cam_distance * 0.01
                                
                                # Rotation logic - independent from kaleido
                                # Rotate camera around the object based on rotation plane
                                # Note: Initial camera is eye={0,0,-z}, up={0,-1,0}
                                if rotate_plane == 'xy':
                                    # Horizontal rotation around Z axis
                                    eye = {
                                        'x': cam_distance * sin_val,
                                        'y': cam_distance * cos_val,
                                        'z': offset
                                    }
                                    up = {'x': 0, 'y': 0, 'z': 1}
                                    
                                elif rotate_plane == 'yz':
                                    # Vertical rotation around X axis (front→bottom→back→top)
                                    eye = {
                                        'x': offset,
                                        'y': cam_distance * sin_val,
                                        'z': -cam_distance * cos_val
                                    }
                                    up = {
                                        'x': 0,
                                        'y': -cos_val,
                                        'z': -sin_val
                                    }
                                    
                                elif rotate_plane == 'xz':
                                    # Horizontal rotation around Y axis
                                    z_sign = 1 if 'manc' in self.dataset.lower() else -1
                                    eye = {
                                        'x': cam_distance * sin_val,
                                        'y': offset,
                                        'z': z_sign * cam_distance * cos_val
                                    }
                                    up = {'x': 0, 'y': -1, 'z': 0}
                                
                                # Rotate camera via JavaScript
                                session.set_camera(eye=eye, up=up)
                                
                                # Convert background string to RGB tuple for screenshot
                                bg_rgb = (255, 255, 255)
                                if self.background_color.lower().strip() == 'black':
                                    bg_rgb = (0, 0, 0)
                                
                                # Take screenshot
                                session.screenshot(fig_path, convert_to_jpeg=True, jpeg_quality=95,
                                                 auto_crop=False, background_color=bg_rgb)
                                
                                # Verify
                                if not os.path.exists(fig_path) or os.path.getsize(fig_path) < 1024:
                                    self._vprint(f'\n⚠️  Frame {i+1} export failed')
                                    frame_export_failed = True
                                    break
                                
                                # Count completed frames for progress display
                                completed_count = len([f for f in os.listdir(pic_folder) 
                                                      if f.startswith('deg_') and f.endswith('.jpeg') 
                                                      and os.path.getsize(os.path.join(pic_folder, f)) >= 1024])
                                
                                elapsed = time.time() - t0
                                avg_time = elapsed / max(1, completed_count - resume_from_frame)
                                remaining = avg_time * (len(steps_to_write) - completed_count)
                                print(f'\r  Frame {completed_count}/{len(steps_to_write)} | '
                                      f'Elapsed: {elapsed:.1f}s | '
                                      f'ETA: {remaining:.1f}s | '
                                      f'{avg_time:.2f}s/frame', end='    ', flush=True)
                        
                        # Check if all frames completed
                        final_frame_count = len([f for f in os.listdir(pic_folder) 
                                                if f.startswith('deg_') and f.endswith('.jpeg')
                                                and os.path.getsize(os.path.join(pic_folder, f)) >= 1024])
                        
                        if final_frame_count >= len(steps_to_write) and not frame_export_failed:
                            webdriver_success = True
                            break  # Exit retry loop on success
                        elif frame_export_failed:
                            # Frame export failed but not due to crash - don't retry
                            break
                            
                    except Exception as e:
                        last_error = e
                        error_msg = str(e).lower()
                        # Check if this is a Chrome crash (empty message or specific patterns)
                        is_chrome_crash = (
                            'Message: \n' in str(e) or 
                            str(e) == '' or
                            'chrome not reachable' in error_msg or
                            'session deleted' in error_msg or
                            'no such window' in error_msg or
                            'tab crashed' in error_msg or
                            'target window already closed' in error_msg or
                            'disconnected' in error_msg
                        )
                        
                        if is_chrome_crash:
                            # Count how many frames were saved before crash
                            frames_saved = len([f for f in os.listdir(pic_folder) 
                                              if f.startswith('deg_') and f.endswith('.jpeg')
                                              and os.path.getsize(os.path.join(pic_folder, f)) >= 1024])
                            
                            # Check if this crash is at a new position (progress was made)
                            if frames_saved > last_crash_frame:
                                # Progress was made! Reset retry counter
                                consecutive_crashes_at_same_point = 1
                                last_crash_frame = frames_saved
                                self._vprint(f'\n   ⚠️  Chrome crashed unexpectedly at frame {frames_saved}. Will resume from there...')
                            else:
                                # Crashed at same position again
                                consecutive_crashes_at_same_point += 1
                                self._vprint(f'\n   ⚠️  Chrome crashed again at frame {frames_saved} (attempt {consecutive_crashes_at_same_point}/{max_retries})')
                            
                            # Do NOT delete frames - keep them for resume
                            # Reset timer for next attempt's ETA calculation
                            t0 = time.time()
                            continue
                        else:
                            # Not a crash - some other error, don't retry
                            break
                
                # Handle final result
                if not webdriver_success:
                    if last_error:
                        self._vprint(f'\n⚠️  WebDriver export failed after {consecutive_crashes_at_same_point} consecutive crashes at same point: {last_error}')
                    self._vprint(f'      Falling back to kaleido...')
                    self._vprint(f'      💡 Tip: Set export_method="kaleido" for future use if WebDriver continues to fail.')
                    use_kaleido_fallback = True
                    frame_export_failed = False  # Reset to try kaleido
                    # Reset timer for kaleido attempt
                    t0 = time.time()
                
                # Clean up temp HTML
                try:
                    os.remove(temp_html)
                except:
                    pass
                
            if actual_export_method == 'kaleido' or use_kaleido_fallback:
                # Default: kaleido export with timeout
                # Check HTML size and auto-simplify if > size cap
                html_size_mb = 0
                if hasattr(self, 'fig_path') and os.path.exists(self.fig_path + '.html'):
                    html_size_mb = os.path.getsize(self.fig_path + '.html') / 1024 / 1024
                
                html_size_cap = self._get_html_size_cap()
                
                # Auto-simplify for kaleido if HTML > size cap
                # But first check if we already have a simplified figure from previous export
                if html_size_mb > html_size_cap:
                    simplified_fig = getattr(self, '_simplified_export_fig', None)
                    simplified_html = getattr(self, '_simplified_html_path', None)
                    
                    if simplified_fig is not None:
                        self._vprint(f'⚠️  HTML size ({html_size_mb:.0f}MB) exceeds {html_size_cap}MB - using previously simplified figure')
                        # Rebuild fig_new from simplified figure traces
                        fig_new = go.Figure(data=simplified_fig.data, layout=fig_layout)
                        fig_new.update_layout(
                            sliders=[],
                            scene=dict(
                                dragmode='orbit',
                                xaxis={'visible':False}, 
                                yaxis={'visible':False},
                                zaxis={'visible':False},
                            ),
                            scene_camera=scene_camera_parameters,
                        )
                    else:
                        # Target 50% file size reduction with overhead estimation
                        html_overhead_mb = 25  # Estimated fixed overhead
                        target_size_mb = html_size_mb * 0.5  # 50% target
                        data_size = html_size_mb - html_overhead_mb
                        target_data_size = target_size_mb - html_overhead_mb
                        simplification_factor = max(0.1, target_data_size / data_size) if data_size > 0 else 0.5
                        
                        self._vprint(f'⚠️  HTML size ({html_size_mb:.0f}MB) exceeds {html_size_cap}MB - kaleido may timeout')
                        self._vprint(f'   Target: {target_size_mb:.0f}MB (50% reduction), simplification factor: {simplification_factor:.2f}...')
                        
                        # Create simplified copy with dynamic reduction
                        fig_new = self._simplify_figure_for_kaleido(fig_new, simplification_factor)
                        
                        # Save simplified HTML for future reuse
                        simplified_html_path = os.path.join(save_folder, f"{saveas}_simplified.html")
                        fig_new.write_html(simplified_html_path, auto_open=False,
                                          include_plotlyjs='cdn',
                                          config={'displayModeBar': False})
                        simplified_size_mb = os.path.getsize(simplified_html_path) / 1024 / 1024
                        self._vprint(f'   ✓ Saved simplified HTML: {os.path.basename(simplified_html_path)} ({simplified_size_mb:.0f}MB)')
                        
                        # Store for reuse
                        self._simplified_export_fig = fig_new
                        self._simplified_html_path = simplified_html_path
                        
                        self._vprint(f'   💡 For better quality, consider using export_method="webdriver"')
                
                # Set up timeout handler (Unix only)
                old_handler = None
                has_alarm_support = False
                try:
                    old_handler = signal.signal(signal.SIGALRM, _timeout_handler)
                    has_alarm_support = True
                except (AttributeError, ValueError):
                    has_alarm_support = False
                
                # Get base timeout from attribute
                base_timeout = getattr(self, 'export_timeout', 60)
                frame1_success = False  # Track if frame 1 exported successfully
                
                # Sequential rendering with timeout control
                self._vprint(f'   Rendering frames... (First frame may take longer to initialize Kaleido)')
                
                try:
                    for i, deg in enumerate(steps_to_write):
                        if frame_export_failed:
                            break
                            
                        rad_i = np.deg2rad(deg)
                        x = view_distance * np.sin(rad_i) * view_direction[0]
                        z = view_distance * np.cos(rad_i) * view_direction[1]
                        
                        # Add small offset to avoid gimbal lock at axis-aligned positions
                        # This prevents camera flipping at 90, 180, 270 degrees
                        y_offset = view_distance * 0.01  # 1% offset
                        
                        if rotate_plane == 'xy':
                            fig_new.update_layout(scene_camera=dict(eye=dict(x=x, y=z, z=y_offset)))
                        elif rotate_plane == 'yz':
                            fig_new.update_layout(scene_camera=dict(eye=dict(x=y_offset, y=x, z=z)))
                        elif rotate_plane == 'xz':
                            fig_new.update_layout(scene_camera=dict(eye=dict(x=x, y=y_offset, z=z)))
                        
                        fig_path = os.path.join(pic_folder, f'deg_{deg:.1f}.jpeg')
                        
                        try:
                            # Timeout: first frame gets base_timeout, subsequent frames get 2x if frame1 succeeded
                            # This is because frame 1 success verifies export is working (just slower)
                            if i == 0:
                                timeout_sec = base_timeout
                            else:
                                timeout_sec = base_timeout * 2 if frame1_success else base_timeout
                            
                            if has_alarm_support:
                                signal.alarm(timeout_sec)
                            
                            kwargs_copy = kwargs.copy()
                            kwargs_copy['scale'] = current_scale
                            fig_new.write_image(fig_path, **kwargs_copy)
                            
                            if has_alarm_support:
                                signal.alarm(0)  # Cancel timeout
                            
                            # Mark frame 1 as successful - allows 2x timeout tolerance for subsequent frames
                            if i == 0:
                                frame1_success = True
                                
                        except PNGExportTimeout:
                            if has_alarm_support:
                                signal.alarm(0)
                            
                            if current_scale > 1:
                                self._vprint(f'\n⚠️  Frame {i+1} timed out at scale={current_scale}, retrying with scale=1...')
                                current_scale = 1
                                try:
                                    if has_alarm_support:
                                        signal.alarm(base_timeout * 2)
                                    kwargs_copy = kwargs.copy()
                                    kwargs_copy['scale'] = 1
                                    fig_new.write_image(fig_path, **kwargs_copy)
                                    if has_alarm_support:
                                        signal.alarm(0)
                                    self._vprint(f'   ✓ Frame {i+1} exported with scale=1')
                                    if i == 0:
                                        frame1_success = True
                                except (PNGExportTimeout, Exception) as retry_e:
                                    if has_alarm_support:
                                        signal.alarm(0)
                                    self._vprint(f'\n⚠️  Frame {i+1} failed even at scale=1: {retry_e}')
                                    self._vprint('   💡 Figure is too complex. Recommendations:')
                                    self._vprint('      1. Use export_method=\"webdriver\" (more reliable for complex figures)')
                                    self._vprint('      2. Increase skeleton_mesh_simplification (e.g., 0.97-0.99)')
                                    self._vprint('      3. Increase export_timeout (current: {base_timeout}s)')
                                    frame_export_failed = True
                                    break
                            else:
                                self._vprint(f'\n⚠️  Frame {i+1} timed out at scale=1. Figure is too complex.')
                                self._vprint('   💡 Recommendations:')
                                self._vprint('      1. Use export_method=\"webdriver\" (more reliable for complex figures)')
                                self._vprint('      2. Increase skeleton_mesh_simplification (e.g., 0.97-0.99)')
                                self._vprint(f'      3. Increase export_timeout (current: {base_timeout}s)')
                                frame_export_failed = True
                                break
                                
                        except Exception as e:
                            if has_alarm_support:
                                signal.alarm(0)
                            self._vprint(f'\n⚠️  Frame {i+1} failed: {e}')
                            if i == 0 and current_scale > 1:
                                self._vprint(f'   Retrying with scale=1...')
                                current_scale = 1
                                try:
                                    kwargs_copy = kwargs.copy()
                                    kwargs_copy['scale'] = 1
                                    fig_new.write_image(fig_path, **kwargs_copy)
                                    self._vprint(f'   ✓ Frame {i+1} exported with scale=1')
                                    frame1_success = True
                                except Exception as retry_e:
                                    self._vprint(f'   ⚠️  Retry failed: {retry_e}')
                                    self._vprint('   💡 Use export_method=\"webdriver\" or increase skeleton_mesh_simplification.')
                                    frame_export_failed = True
                                    break
                            else:
                                frame_export_failed = True
                                break
                        
                        elapsed = time.time() - t0
                        avg_time = elapsed / (i + 1)
                        remaining = avg_time * (len(steps_to_write) - i - 1)
                        print(f'\r  Frame {i+1}/{len(steps_to_write)} | '
                              f'Elapsed: {elapsed:.1f}s | '
                              f'ETA: {remaining:.1f}s | '
                              f'{avg_time:.2f}s/frame', end='    ', flush=True)
                              
                finally:
                    # Always cancel any pending alarm and restore original handler
                    if has_alarm_support:
                        try:
                            signal.alarm(0)  # Cancel any pending alarm
                            if old_handler is not None:
                                signal.signal(signal.SIGALRM, old_handler)
                        except:
                            pass
            
            if frame_export_failed:
                self._vprint(f'\n⚠️  Frame rendering aborted. Skipping video generation.')
                return 1
            
            print('\n✓ Image rendering complete')
            
            # Apply consistent cropping if auto_crop is enabled
            # This ensures all frames have the same dimensions during rotation
            if auto_crop:
                self._vprint(f'   Applying consistent auto-crop across all frames...')
                crop_result = self._apply_consistent_crop(pic_folder, margin=crop_margin)
                if crop_result:
                    crop_w, crop_h = crop_result
                    self._vprint(f'   ✓ All frames cropped consistently to {crop_w}x{crop_h}')
                else:
                    self._vprint(f'   ⚠️  Auto-crop failed, using original frame sizes')
        
        # Generate videos from images
        self._vprint(f'\nGenerating videos...')
        imglist = os.listdir(pic_folder)
        img_eg = cv2.imread(os.path.join(pic_folder, imglist[0]))
        height, width, layers = img_eg.shape
        
        self._vprint(f'   Video resolution: {width}x{height}')

        # Add rotation direction indicator to filename
        rotation_suffix = 'h' if rotate_plane == 'xz' else 'v' if rotate_plane == 'yz' else rotate_plane
        
        # Forward video - OPTIMIZED with faster codec
        video_path_forward = os.path.join(save_folder, f'{saveas}_video_{rotation_suffix}_forward.mp4')
        # Use H.264 codec for better compression and compatibility
        fourcc = cv2.VideoWriter_fourcc(*'avc1')  # H.264 codec (faster than mp4v)
        out = cv2.VideoWriter(video_path_forward, fourcc, fps, frameSize=(width, height))
        
        t0 = time.time()
        for i, deg in enumerate(steps_to_write):
            img = cv2.imread(os.path.join(pic_folder, f'deg_{deg:.1f}.jpeg'))
            out.write(img)
            if (i + 1) % 10 == 0 or i == len(steps_to_write) - 1:
                print(f'\r  Forward video: {i+1}/{len(steps_to_write)} frames', end='  ')
        out.release()
        t1 = time.time()
        print(f'\n✓ Forward video: {video_path_forward} ({t1-t0:.1f}s)')
        
        # Backward video
        video_path_backward = os.path.join(save_folder, f'{saveas}_video_{rotation_suffix}_backward.mp4')
        out = cv2.VideoWriter(video_path_backward, fourcc, fps, frameSize=(width, height))
        
        t0 = time.time()
        for i, deg in enumerate(steps_to_write[::-1]):
            img = cv2.imread(os.path.join(pic_folder, f'deg_{deg:.1f}.jpeg'))
            out.write(img)
            if (i + 1) % 10 == 0 or i == len(steps_to_write) - 1:
                print(f'\r  Backward video: {i+1}/{len(steps_to_write)} frames', end='  ')
        out.release()
        t1 = time.time()
        print(f'\n✓ Backward video: {video_path_backward} ({t1-t0:.1f}s)')
        
        print(f'\n✅ Video export complete!')
        self._vprint(f'   Image cache: {pic_folder}')
        self._vprint(f'   Tip: Use use_existing_images=True to skip re-rendering next time')
        
        # Convert to GIF if requested
        if export_gif:
            self._vprint(f'\n🎞️  Converting videos to GIF format...')
            self._vprint(f'   Scale: {gif_scale} | Optimize: {gif_optimize}')
            
            # Convert forward video to GIF
            gif_path_forward = video_path_forward.replace('.mp4', '.gif')
            try:
                video2gif(
                    video_path_forward,
                    gif_path_forward,
                    fps=fps,
                    scale=gif_scale,
                    optimize=gif_optimize
                )
                self._vprint(f'   ✓ Forward GIF: {gif_path_forward}')
            except Exception as e:
                self._vprint(f'   ⚠️  Forward GIF conversion failed: {e}')
            
            # Convert backward video to GIF
            gif_path_backward = video_path_backward.replace('.mp4', '.gif')
            try:
                video2gif(
                    video_path_backward,
                    gif_path_backward,
                    fps=fps,
                    scale=gif_scale,
                    optimize=gif_optimize
                )
                self._vprint(f'   ✓ Backward GIF: {gif_path_backward}')
            except Exception as e:
                self._vprint(f'   ⚠️  Backward GIF conversion failed: {e}')
        
        return 0


def _detect_content_bounds_standalone(img, background_color=(255, 255, 255)):
    """Detect the bounding box of content in an image (standalone version)."""
    from PIL import Image
    import numpy as np
    
    if img.mode == 'RGBA':
        bg = Image.new('RGB', img.size, background_color)
        bg.paste(img, mask=img.split()[3])
        img_rgb = bg
    else:
        img_rgb = img.convert('RGB')
    
    arr = np.array(img_rgb)
    tolerance = 10
    bg_r, bg_g, bg_b = background_color
    
    non_bg_mask = (
        (np.abs(arr[:, :, 0].astype(int) - bg_r) > tolerance) |
        (np.abs(arr[:, :, 1].astype(int) - bg_g) > tolerance) |
        (np.abs(arr[:, :, 2].astype(int) - bg_b) > tolerance)
    )
    
    rows = np.any(non_bg_mask, axis=1)
    cols = np.any(non_bg_mask, axis=0)
    
    if not rows.any() or not cols.any():
        return None
    
    row_min, row_max = np.where(rows)[0][[0, -1]]
    col_min, col_max = np.where(cols)[0][[0, -1]]
    
    return (row_min, row_max, col_min, col_max)


def _apply_consistent_crop_standalone(pic_folder, margin=20, background_color=(255, 255, 255)):
    """
    Apply consistent cropping to all images in a folder (standalone version).
    
    This finds the maximum extent of content across all frames and crops
    all frames to the same unified bounds + margin.
    
    Parameters
    ----------
    pic_folder : str
        Folder containing the frame images (deg_*.jpeg)
    margin : int
        Margin to preserve around content
    background_color : tuple
        RGB tuple for background color
        
    Returns
    -------
    tuple or None
        (width, height) of the cropped images, or None if failed
    """
    from PIL import Image
    import glob
    
    # Find all JPEG images in folder
    image_paths = sorted(glob.glob(os.path.join(pic_folder, 'deg_*.jpeg')))
    
    if not image_paths:
        return None
    
    # Sample frames for efficiency (every 10 degrees for 360 rotation)
    sample_count = min(36, len(image_paths))
    if sample_count < len(image_paths):
        indices = [int(i * len(image_paths) / sample_count) for i in range(sample_count)]
        sampled_paths = [image_paths[i] for i in indices]
    else:
        sampled_paths = image_paths
    
    # Compute unified bounds across sampled frames
    unified_row_min = None
    unified_row_max = None
    unified_col_min = None
    unified_col_max = None
    
    for path in sampled_paths:
        try:
            with Image.open(path) as img:
                bounds = _detect_content_bounds_standalone(img, background_color)
            
            if bounds is None:
                continue
                
            row_min, row_max, col_min, col_max = bounds
            
            if unified_row_min is None:
                unified_row_min = row_min
                unified_row_max = row_max
                unified_col_min = col_min
                unified_col_max = col_max
            else:
                unified_row_min = min(unified_row_min, row_min)
                unified_row_max = max(unified_row_max, row_max)
                unified_col_min = min(unified_col_min, col_min)
                unified_col_max = max(unified_col_max, col_max)
        except Exception:
            continue
    
    if unified_row_min is None:
        return None
    
    # Get image dimensions from first image
    with Image.open(image_paths[0]) as img:
        img_height = img.height
        img_width = img.width
    
    # Add margin to unified bounds
    row_min = max(0, unified_row_min - margin)
    row_max = min(img_height - 1, unified_row_max + margin)
    col_min = max(0, unified_col_min - margin)
    col_max = min(img_width - 1, unified_col_max + margin)
    
    final_width = col_max - col_min + 1
    final_height = row_max - row_min + 1
    
    # Apply consistent crop to all images
    for path in image_paths:
        try:
            with Image.open(path) as img:
                cropped = img.crop((col_min, row_min, col_max + 1, row_max + 1))
                if cropped.mode == 'RGBA':
                    rgb_img = Image.new('RGB', cropped.size, background_color)
                    rgb_img.paste(cropped, mask=cropped.split()[3])
                    cropped = rgb_img
                elif cropped.mode != 'RGB':
                    cropped = cropped.convert('RGB')
                cropped.save(path, 'JPEG', quality=95)
        except Exception:
            pass
    
    return (final_width, final_height)


def export_video_from_html(html_file, fps=30, degree_per_frame=1.0, rotate='horizontal',
                           output_dir=None, use_existing_images=True, 
                           export_gif=True, gif_scale=0.2, gif_optimize=True,
                           auto_crop=False, crop_margin=30, **kwargs):
    '''
    Standalone function to export a rotating video from an existing Plotly HTML file.
    
    This function does NOT require VisualizeSkeleton initialization or NeuPrint client.
    It directly loads the HTML figure and renders the video.
    
    Parameters
    ----------
    html_file : str
        Path to existing Plotly HTML file to load figure data from.
    fps : int, default 30
        Frames per second for the output video.
    degree_per_frame : float, default 1.0
        Rotation angle in degrees per frame.
        - 1.0 → 360 frames for full rotation (12 sec video at 30 fps)
        - 2.0 → 180 frames for full rotation (6 sec video at 30 fps)
    rotate : str, default 'horizontal'
        Rotation direction: 'horizontal' or 'vertical'.
    output_dir : str, optional
        Directory to save video output. If None, uses the directory containing html_file.
    use_existing_images : bool, default True
        If True, reuse cached images from previous export if available.
    export_gif : bool, default True
        If True, automatically convert videos to GIF format after export.
    gif_scale : float, default 0.2
        Scale factor for GIF resolution (0.1-1.0). Lower values = smaller file size.
    gif_optimize : bool, default True
        Enable GIF compression optimization for smaller file sizes.
    auto_crop : bool, default False
        If True, auto-crop frames to content bounds with consistent sizing across all frames.
        This ensures uniform frame dimensions during rotation for smooth video playback.
    crop_margin : int, default 30
        Margin in pixels around content when auto_crop is enabled.
    **kwargs : dict
        Additional arguments for plotly write_image():
        - scale : int, default 2
        - width : int, default 1200
        - height : int, default 900
    
    Returns
    -------
    int
        0 on success, 1 on failure
    
    Examples
    --------
    # Basic usage - output to same directory as HTML file
    from visualize_skeleton import export_video_from_html
    export_video_from_html('/path/to/my_neurons.html')
    
    # Custom settings
    export_video_from_html(
        '/path/to/my_neurons.html',
        fps=60,
        degree_per_frame=0.5,  # Slower rotation
        rotate='vertical',
        scale=4  # Higher quality
    )
    
    # Specify output directory
    export_video_from_html(
        '/path/to/my_neurons.html',
        output_dir='/path/to/output/'
    )
    '''
    import plotly.io as pio
    import plotly.graph_objects as go
    import cv2
    import shutil
    import time
    
    # Validate input
    if not os.path.exists(html_file):
        raise FileNotFoundError(f'HTML file not found: {html_file}')
    
    # Determine output directory
    if output_dir is None:
        save_folder = os.path.dirname(os.path.abspath(html_file))
    else:
        save_folder = output_dir
        os.makedirs(save_folder, exist_ok=True)
    
    saveas = os.path.splitext(os.path.basename(html_file))[0]
    
    # Handle rotate parameter
    if rotate == 'horizontal':
        rotate_plane = 'xz'
    elif rotate == 'vertical':
        rotate_plane = 'yz'
    else:
        rotate_plane = 'xz'
    
    # Set defaults
    view_direction = kwargs.pop('view_direction', (1, -1))
    view_distance = kwargs.pop('view_distance', 2.2)
    synapse_size = kwargs.pop('synapse_size', 1)
    
    if kwargs.get('scale') is None and kwargs.get('width') is None and kwargs.get('height') is None:
        kwargs['scale'] = 2
    
    # Load figure from HTML
    print(f'📂 Loading figure from: {html_file}')
    try:
        fig_loaded = pio.read_html(html_file)
        fig_traces = fig_loaded.data
        print(f'✓ Loaded {len(fig_traces)} traces from HTML file')
    except Exception as e:
        raise RuntimeError(f'Failed to load figure from HTML: {e}')
    
    # Configure figure for video
    for trace in fig_traces:
        trace.showlegend = False
        if hasattr(trace, 'marker'):
            trace.marker.size = synapse_size
    
    fig_layout = go.Layout(margin=dict(l=1, r=1, b=1, t=1, pad=0))
    fig_new = go.Figure(data=fig_traces, layout=fig_layout)
    
    fig_new.update_layout(
        sliders=[],
        scene=dict(
            dragmode='orbit',
            xaxis={'visible': False},
            yaxis={'visible': False},
            zaxis={'visible': False},
        ),
        scene_camera=dict(
            up=dict(x=0, y=-1, z=0),
            eye=dict(x=0, y=0, z=-view_distance),
        ),
    )
    
    # Set up image folder
    pic_folder = os.path.join(save_folder, f'pics_{fps}fps_{rotate_plane}')
    
    # Calculate rotation steps
    step = degree_per_frame
    steps_to_write = np.linspace(0, 360, int(360/step), endpoint=False)
    
    # Check for existing images
    if use_existing_images and os.path.exists(pic_folder):
        existing_images = [f for f in os.listdir(pic_folder) if f.endswith('.jpeg')]
        if len(existing_images) == len(steps_to_write):
            print(f'✓ Using {len(existing_images)} existing images from {pic_folder}')
        else:
            print(f'⚠️  Found {len(existing_images)} images but need {len(steps_to_write)}, re-rendering...')
            use_existing_images = False
    else:
        use_existing_images = False
    
    # Render images if needed
    if not use_existing_images:
        if os.path.exists(pic_folder):
            shutil.rmtree(pic_folder)
        os.makedirs(pic_folder)
        
        if 'width' not in kwargs:
            kwargs['width'] = 1200
        if 'height' not in kwargs:
            kwargs['height'] = 900
        
        print(f'🎬 Rendering {len(steps_to_write)} frames at {fps} fps...')
        t0 = time.time()
        
        for i, deg in enumerate(steps_to_write):
            rad_i = np.deg2rad(deg)
            x = view_distance * np.sin(rad_i) * view_direction[0]
            y = view_distance * np.cos(rad_i) * view_direction[1]
            
            if rotate_plane == 'xy':
                fig_new.update_layout(scene_camera=dict(eye=dict(x=x, y=y, z=0)))
            elif rotate_plane == 'yz':
                fig_new.update_layout(scene_camera=dict(eye=dict(x=0, y=x, z=y)))
            elif rotate_plane == 'xz':
                fig_new.update_layout(scene_camera=dict(eye=dict(x=x, y=0, z=y)))
            
            fig_path = os.path.join(pic_folder, f'deg_{deg:.1f}.jpeg')
            
            try:
                fig_new.write_image(fig_path, **kwargs)
            except Exception as e:
                print(f'\n⚠️  Frame {i+1} failed: {e}')
                if i == 0:
                    print('   Try reducing "scale" (e.g. scale=1)')
                    return 1
            
            elapsed = time.time() - t0
            avg_time = elapsed / (i + 1)
            remaining = avg_time * (len(steps_to_write) - i - 1)
            print(f'\r  Frame {i+1}/{len(steps_to_write)} | Elapsed: {elapsed:.1f}s | ETA: {remaining:.1f}s', end='  ', flush=True)
        
        print('\n✓ Image rendering complete')
        
        # Apply consistent cropping if auto_crop is enabled
        if auto_crop:
            print(f'   Applying consistent auto-crop across all frames...')
            crop_result = _apply_consistent_crop_standalone(pic_folder, margin=crop_margin)
            if crop_result:
                crop_w, crop_h = crop_result
                print(f'   ✓ All frames cropped consistently to {crop_w}x{crop_h}')
            else:
                print(f'   ⚠️  Auto-crop failed, using original frame sizes')
    
    # Generate videos
    print(f'\nGenerating videos...')
    imglist = os.listdir(pic_folder)
    img_eg = cv2.imread(os.path.join(pic_folder, imglist[0]))
    height, width, layers = img_eg.shape
    
    fourcc = cv2.VideoWriter_fourcc(*'avc1')
    
    # Add rotation direction indicator to filename
    rotation_suffix = 'h' if rotate_plane == 'xz' else 'v' if rotate_plane == 'yz' else rotate_plane
    
    # Forward video
    video_path_forward = os.path.join(save_folder, f'{saveas}_video_{rotation_suffix}_forward.mp4')
    out = cv2.VideoWriter(video_path_forward, fourcc, fps, frameSize=(width, height))
    for deg in steps_to_write:
        img = cv2.imread(os.path.join(pic_folder, f'deg_{deg:.1f}.jpeg'))
        out.write(img)
    out.release()
    print(f'✓ Forward video: {video_path_forward}')
    
    # Backward video
    video_path_backward = os.path.join(save_folder, f'{saveas}_video_{rotation_suffix}_backward.mp4')
    out = cv2.VideoWriter(video_path_backward, fourcc, fps, frameSize=(width, height))
    for deg in steps_to_write[::-1]:
        img = cv2.imread(os.path.join(pic_folder, f'deg_{deg:.1f}.jpeg'))
        out.write(img)
    out.release()
    print(f'✓ Backward video: {video_path_backward}')
    
    print(f'\n✅ Video export complete!')
    
    # Convert to GIF if requested
    if export_gif:
        print(f'\n🎞️  Converting videos to GIF format...')
        print(f'   Scale: {gif_scale} | Optimize: {gif_optimize}')
        
        # Convert forward video to GIF
        gif_path_forward = video_path_forward.replace('.mp4', '.gif')
        try:
            video2gif(
                video_path_forward,
                gif_path_forward,
                fps=fps,
                scale=gif_scale,
                optimize=gif_optimize
            )
            print(f'   ✓ Forward GIF: {gif_path_forward}')
        except Exception as e:
            print(f'   ⚠️  Forward GIF conversion failed: {e}')
        
        # Convert backward video to GIF
        gif_path_backward = video_path_backward.replace('.mp4', '.gif')
        try:
            video2gif(
                video_path_backward,
                gif_path_backward,
                fps=fps,
                scale=gif_scale,
                optimize=gif_optimize
            )
            print(f'   ✓ Backward GIF: {gif_path_backward}')
        except Exception as e:
            print(f'   ⚠️  Backward GIF conversion failed: {e}')
    
    return 0


def export_video_webdriver(
    html_file: str,
    fps: int = 30,
    degree_per_frame: float = 1.0,
    rotate: str = 'horizontal',
    output_dir: str = None,
    width: int = 1200,
    height: int = 900,
    scale: int = 2,
    view_distance: float = 2.2,
    export_gif: bool = True,
    gif_scale: float = 0.2,
    gif_optimize: bool = True,
    timeout: int = 120,
    auto_crop: bool = False,
    crop_margin: int = 30,
) -> int:
    """
    Export a rotating video from an existing Plotly HTML file using WebDriver.
    
    This is an EFFICIENT alternative to export_video_from_html() that:
    - Opens the browser ONCE and keeps it open
    - Rotates the camera using JavaScript (no figure regeneration)
    - Takes screenshots in series without reopening the browser
    
    This method is significantly faster for large/complex 3D figures because:
    1. The HTML only needs to load once (not per-frame)
    2. Camera rotation uses JavaScript instead of Python figure update + kaleido export
    3. Works with WebGL for smooth rendering
    
    Parameters
    ----------
    html_file : str
        Path to existing Plotly HTML file to load.
    fps : int, default 30
        Frames per second for the output video.
    degree_per_frame : float, default 1.0
        Rotation angle in degrees per frame.
        - 1.0 → 360 frames for full rotation (12 sec video at 30 fps)
        - 2.0 → 180 frames for full rotation (6 sec video at 30 fps)
    rotate : str, default 'horizontal'
        Rotation direction: 'horizontal' or 'vertical'.
    output_dir : str, optional
        Directory to save video output. If None, uses the directory containing html_file.
    width : int, default 1200
        Browser viewport width.
    height : int, default 900
        Browser viewport height.
    scale : int, default 2
        Scale factor for screenshot resolution (actual size = width*scale x height*scale).
    view_distance : float, default 2.2
        Camera distance from the center (affects zoom level).
    export_gif : bool, default True
        If True, automatically convert videos to GIF format after export.
    gif_scale : float, default 0.2
        Scale factor for GIF resolution (0.1-1.0).
    gif_optimize : bool, default True
        Enable GIF compression optimization.
    timeout : int, default 120
        Maximum time in seconds to wait for page load.
    auto_crop : bool, default False
        If True, auto-crop frames to content bounds with consistent sizing across all frames.
        This ensures uniform frame dimensions during rotation for smooth video playback.
    crop_margin : int, default 30
        Margin in pixels around content when auto_crop is enabled.
    
    Returns
    -------
    int
        0 on success, 1 on failure
    
    Notes
    -----
    Requires: selenium, webdriver-manager
    
    On macOS, headless Chrome doesn't support WebGL, so we use an offscreen
    window (positioned at -10000,-10000) instead of true headless mode.
    
    Examples
    --------
    # Basic usage
    from visualize_skeleton import export_video_webdriver
    export_video_webdriver('/path/to/my_neurons.html')
    
    # Faster rotation, higher quality
    export_video_webdriver(
        '/path/to/my_neurons.html',
        fps=60,
        degree_per_frame=2.0,  # Faster rotation
        scale=3  # Higher quality
    )
    
    # Vertical rotation
    export_video_webdriver(
        '/path/to/my_neurons.html',
        rotate='vertical',
        view_distance=2.5
    )
    """
    try:
        from selenium import webdriver
        from selenium.webdriver.chrome.service import Service as ChromeService
        from selenium.webdriver.chrome.options import Options as ChromeOptions
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC
        from selenium.webdriver.common.by import By
    except ImportError:
        print("❌ Error: selenium is required for export_video_webdriver()")
        print("   Install with: pip install selenium webdriver-manager")
        return 1
    
    import shutil
    import time
    
    # Validate input
    if not os.path.exists(html_file):
        raise FileNotFoundError(f'HTML file not found: {html_file}')
    
    # Determine output directory
    if output_dir is None:
        save_folder = os.path.dirname(os.path.abspath(html_file))
    else:
        save_folder = output_dir
        os.makedirs(save_folder, exist_ok=True)
    
    saveas = os.path.splitext(os.path.basename(html_file))[0]
    
    # Handle rotate parameter
    if rotate == 'horizontal':
        rotate_plane = 'xz'
    elif rotate == 'vertical':
        rotate_plane = 'yz'
    else:
        rotate_plane = 'xz'
    
    # Calculate rotation steps
    step = degree_per_frame
    steps_to_write = np.linspace(0, 360, int(360/step), endpoint=False)
    
    # Set up image folder
    pic_folder = os.path.join(save_folder, f'pics_{fps}fps_{rotate_plane}_webdriver')
    if os.path.exists(pic_folder):
        shutil.rmtree(pic_folder)
    os.makedirs(pic_folder)
    
    # Calculate actual browser dimensions
    actual_width = width * scale
    actual_height = height * scale
    
    # Set up Chrome options
    # Use --headless=new (Chrome 109+) for WebGL support in headless mode
    chrome_options = ChromeOptions()
    chrome_options.add_argument('--no-sandbox')
    chrome_options.add_argument('--disable-dev-shm-usage')
    chrome_options.add_argument(f'--window-size={actual_width},{actual_height}')
    chrome_options.add_argument('--headless=new')  # Modern headless with WebGL support
    
    # Initialize ChromeDriver using webdriver-manager (cross-platform)
    driver = None
    try:
        from webdriver_manager.chrome import ChromeDriverManager
        service = ChromeService(ChromeDriverManager().install())
        driver = webdriver.Chrome(service=service, options=chrome_options)
    except ImportError:
        try:
            driver = webdriver.Chrome(options=chrome_options)
        except Exception as e:
            print(f"❌ Error: Could not initialize Chrome WebDriver: {e}")
            print(f"   Install dependencies: pip install selenium webdriver-manager")
            return 1
    except Exception as e:
        try:
            driver = webdriver.Chrome(options=chrome_options)
        except Exception as e2:
            print(f"❌ Error: Could not initialize Chrome WebDriver: {e2}")
            print(f"   Ensure Chrome 109+ is installed and webdriver-manager is up to date")
            return 1
    
    try:
        print(f'📂 Loading HTML file: {html_file}')
        file_url = f'file://{os.path.abspath(html_file)}'
        driver.get(file_url)
        
        # Wait for Plotly to render
        print(f'   Waiting for Plotly to render...')
        wait = WebDriverWait(driver, timeout)
        wait.until(EC.presence_of_element_located((By.CLASS_NAME, "plotly")))
        
        # Additional wait for WebGL rendering
        time.sleep(3)
        print(f'✓ Page loaded and rendered')
        
        # JavaScript to update camera position
        # Plotly stores camera in layout.scene.camera.eye
        js_set_camera = """
        var gd = document.querySelector('.js-plotly-plot');
        if (gd && gd.layout && gd.layout.scene) {
            Plotly.relayout(gd, {
                'scene.camera.eye': {x: %f, y: %f, z: %f},
                'scene.camera.up': {x: 0, y: -1, z: 0}
            });
        }
        """
        
        print(f'🎬 Rendering {len(steps_to_write)} frames at {fps} fps...')
        t0 = time.time()
        
        for i, deg in enumerate(steps_to_write):
            rad_i = np.deg2rad(deg)
            
            # Add small offset to avoid gimbal lock at axis-aligned positions
            # This prevents camera flipping at 90, 180, 270 degrees
            offset = view_distance * 0.01  # 1% offset
            
            # Use view_direction (1, -1) to match original export_video_from_html behavior
            # The -1 for cos component ensures consistent rotation direction
            sin_component = view_distance * np.sin(rad_i)  # * 1
            cos_component = view_distance * np.cos(rad_i) * (-1)  # * -1
            
            if rotate_plane == 'xz':  # Horizontal rotation
                eye_x = sin_component
                eye_y = offset  # Small offset instead of 0
                eye_z = cos_component
            elif rotate_plane == 'yz':  # Vertical rotation
                eye_x = offset  # Small offset instead of 0
                eye_y = sin_component
                eye_z = cos_component
            else:  # xy plane
                eye_x = sin_component
                eye_y = cos_component
                eye_z = offset  # Small offset instead of 0
            
            # Update camera via JavaScript
            driver.execute_script(js_set_camera % (eye_x, eye_y, eye_z))
            
            # Brief wait for rendering
            time.sleep(0.1)
            
            # Take screenshot
            frame_path = os.path.join(pic_folder, f'deg_{deg:.1f}.jpeg')
            
            # Get screenshot as PNG, then convert to JPEG
            screenshot_png = os.path.join(pic_folder, f'temp_{deg:.1f}.png')
            driver.save_screenshot(screenshot_png)
            
            # Convert PNG to JPEG for consistency with existing code
            from PIL import Image
            img = Image.open(screenshot_png)
            img = img.convert('RGB')
            img.save(frame_path, 'JPEG', quality=95)
            os.remove(screenshot_png)
            
            # Progress update
            elapsed = time.time() - t0
            avg_time = elapsed / (i + 1)
            remaining = avg_time * (len(steps_to_write) - i - 1)
            print(f'\r   Frame {i+1}/{len(steps_to_write)} | Elapsed: {elapsed:.1f}s | ETA: {remaining:.1f}s', end='  ', flush=True)
        
        print('\n✓ Image rendering complete')
        
        # Apply consistent cropping if auto_crop is enabled
        if auto_crop:
            print(f'   Applying consistent auto-crop across all frames...')
            crop_result = _apply_consistent_crop_standalone(pic_folder, margin=crop_margin)
            if crop_result:
                crop_w, crop_h = crop_result
                print(f'   ✓ All frames cropped consistently to {crop_w}x{crop_h}')
            else:
                print(f'   ⚠️  Auto-crop failed, using original frame sizes')
        
    finally:
        driver.quit()
    
    # Generate videos
    print(f'\nGenerating videos...')
    imglist = os.listdir(pic_folder)
    imglist = [f for f in imglist if f.endswith('.jpeg')]
    
    img_eg = cv2.imread(os.path.join(pic_folder, imglist[0]))
    frame_height, frame_width, layers = img_eg.shape
    
    fourcc = cv2.VideoWriter_fourcc(*'avc1')
    
    # Forward video
    video_path_forward = os.path.join(save_folder, f'{saveas}_video_forward_webdriver.mp4')
    out = cv2.VideoWriter(video_path_forward, fourcc, fps, frameSize=(frame_width, frame_height))
    for deg in steps_to_write:
        img = cv2.imread(os.path.join(pic_folder, f'deg_{deg:.1f}.jpeg'))
        out.write(img)
    out.release()
    print(f'✓ Forward video: {video_path_forward}')
    
    # Backward video
    video_path_backward = os.path.join(save_folder, f'{saveas}_video_backward_webdriver.mp4')
    out = cv2.VideoWriter(video_path_backward, fourcc, fps, frameSize=(frame_width, frame_height))
    for deg in steps_to_write[::-1]:
        img = cv2.imread(os.path.join(pic_folder, f'deg_{deg:.1f}.jpeg'))
        out.write(img)
    out.release()
    print(f'✓ Backward video: {video_path_backward}')
    
    print(f'\n✅ Video export complete!')
    
    # Convert to GIF if requested
    if export_gif:
        print(f'\n🎞️  Converting videos to GIF format...')
        print(f'   Scale: {gif_scale} | Optimize: {gif_optimize}')
        
        # Convert forward video to GIF
        gif_path_forward = video_path_forward.replace('.mp4', '.gif')
        try:
            video2gif(
                video_path_forward,
                gif_path_forward,
                fps=fps,
                scale=gif_scale,
                optimize=gif_optimize
            )
            print(f'   ✓ Forward GIF: {gif_path_forward}')
        except Exception as e:
            print(f'   ⚠️  Forward GIF conversion failed: {e}')
        
        # Convert backward video to GIF
        gif_path_backward = video_path_backward.replace('.mp4', '.gif')
        try:
            video2gif(
                video_path_backward,
                gif_path_backward,
                fps=fps,
                scale=gif_scale,
                optimize=gif_optimize
            )
            print(f'   ✓ Backward GIF: {gif_path_backward}')
        except Exception as e:
            print(f'   ⚠️  Backward GIF conversion failed: {e}')
    
    return 0


def export_png_webdriver(
    html_file: str,
    output_path: str = None,
    width: int = 1200,
    height: int = 900,
    scale: int = 3,
    timeout: int = 60,
) -> str:
    """
    Export a single PNG from an existing Plotly HTML file using WebDriver.
    
    This is an alternative to kaleido-based PNG export that works reliably
    with WebGL for complex 3D scenes.
    
    Parameters
    ----------
    html_file : str
        Path to existing Plotly HTML file to load.
    output_path : str, optional
        Path for the output PNG file. If None, uses the same path as html_file
        with .png extension.
    width : int, default 1200
        Browser viewport width.
    height : int, default 900
        Browser viewport height.
    scale : int, default 3
        Scale factor for screenshot resolution (actual size = width*scale x height*scale).
    timeout : int, default 60
        Maximum time in seconds to wait for page load.
    
    Returns
    -------
    str
        Path to the created PNG file on success, None on failure.
    
    Notes
    -----
    Requires: selenium, webdriver-manager
    
    On macOS, headless Chrome doesn't support WebGL, so we use an offscreen
    window (positioned at -10000,-10000) instead of true headless mode.
    
    Examples
    --------
    # Basic usage
    from visualize_skeleton import export_png_webdriver
    export_png_webdriver('/path/to/my_neurons.html')
    
    # Higher resolution
    export_png_webdriver('/path/to/my_neurons.html', scale=4)
    
    # Custom output path
    export_png_webdriver('/path/to/my_neurons.html', output_path='/path/to/output.png')
    """
    try:
        from selenium import webdriver
        from selenium.webdriver.chrome.service import Service as ChromeService
        from selenium.webdriver.chrome.options import Options as ChromeOptions
        from selenium.webdriver.support.ui import WebDriverWait
        from selenium.webdriver.support import expected_conditions as EC
        from selenium.webdriver.common.by import By
    except ImportError:
        print("❌ Error: selenium is required for export_png_webdriver()")
        print("   Install with: pip install selenium webdriver-manager")
        return None
    
    import time
    
    # Validate input
    if not os.path.exists(html_file):
        raise FileNotFoundError(f'HTML file not found: {html_file}')
    
    # Determine output path
    if output_path is None:
        output_path = os.path.splitext(html_file)[0] + '_webdriver.png'
    
    # Calculate actual browser dimensions
    actual_width = width * scale
    actual_height = height * scale
    
    # Set up Chrome options
    # Use --headless=new (Chrome 109+) for WebGL support in headless mode
    chrome_options = ChromeOptions()
    chrome_options.add_argument('--no-sandbox')
    chrome_options.add_argument('--disable-dev-shm-usage')
    chrome_options.add_argument(f'--window-size={actual_width},{actual_height}')
    chrome_options.add_argument('--headless=new')  # Modern headless with WebGL support
    
    # Initialize ChromeDriver using webdriver-manager (cross-platform)
    driver = None
    try:
        from webdriver_manager.chrome import ChromeDriverManager
        service = ChromeService(ChromeDriverManager().install())
        driver = webdriver.Chrome(service=service, options=chrome_options)
    except ImportError:
        try:
            driver = webdriver.Chrome(options=chrome_options)
        except Exception as e:
            print(f"❌ Error: Could not initialize Chrome WebDriver: {e}")
            print(f"   Install dependencies: pip install selenium webdriver-manager")
            return None
    except Exception as e:
        try:
            driver = webdriver.Chrome(options=chrome_options)
        except Exception as e2:
            print(f"❌ Error: Could not initialize Chrome WebDriver: {e2}")
            print(f"   Ensure Chrome 109+ is installed and webdriver-manager is up to date")
            return None
    
    try:
        print(f'📂 Loading HTML file: {html_file}')
        file_url = f'file://{os.path.abspath(html_file)}'
        driver.get(file_url)
        
        # Wait for Plotly to render
        print(f'   Waiting for Plotly to render...')
        wait = WebDriverWait(driver, timeout)
        wait.until(EC.presence_of_element_located((By.CLASS_NAME, "plotly")))
        
        # Additional wait for WebGL rendering
        time.sleep(3)
        print(f'✓ Page loaded and rendered')
        
        # Take screenshot
        driver.save_screenshot(output_path)
        
        print(f'✅ PNG exported: {output_path}')
        
        # Report file size
        output_size = os.path.getsize(output_path) / (1024 * 1024)
        print(f'   Resolution: {actual_width}x{actual_height}')
        print(f'   File size: {output_size:.2f} MB')
        
        return output_path
        
    except Exception as e:
        print(f'❌ Error during export: {e}')
        return None
        
    finally:
        driver.quit()


def video2gif(
    input_video: str,
    output_gif: str = None,
    fps: int = None,
    scale: float = 1.0,
    optimize: bool = True,
    loop: int = 0,
) -> str:
    """
    Convert a video file (MP4) to an animated GIF with adjustable compression and fps.
    
    This is a static helper function that can be called independently.
    
    Parameters
    ----------
    input_video : str
        Path to the input video file (MP4 or other formats supported by cv2).
    output_gif : str, optional
        Path for the output GIF file. If None, uses the same path as input with .gif extension.
    fps : int, optional
        Target frames per second for the GIF. If None, uses the original video fps.
        Lower fps = smaller file size, choppier animation.
    scale : float, default 1.0
        Scale factor for the output dimensions (0.0-1.0 for compression).
        - 1.0: Original resolution
        - 0.5: Half resolution (75% file size reduction)
        - 0.25: Quarter resolution
    optimize : bool, default True
        Whether to optimize the GIF palette for smaller file size.
        When True and ffmpeg is on PATH, uses ffmpeg palettegen/paletteuse
        (scene-wide palette, ~10-15% smaller than PIL on rendered scenes).
        Otherwise falls back to PIL's optimize and disposal settings.
    loop : int, default 0
        Number of times the GIF should loop.
        - 0: Loop forever
        - 1: Play once
        - n: Loop n times
    
    Returns
    -------
    str
        Path to the created GIF file.
    
    Examples
    --------
    # Basic conversion
    from visualize_skeleton import video2gif
    video2gif('/path/to/video.mp4')
    
    # With compression (half size, 15 fps)
    video2gif('/path/to/video.mp4', fps=15, scale=0.5)
    
    # Custom output path
    video2gif('/path/to/video.mp4', output_gif='/path/to/output.gif', scale=0.75)
    """
    from PIL import Image
    import shutil
    
    if not os.path.exists(input_video):
        raise FileNotFoundError(f"Input video not found: {input_video}")
    
    # Set output path
    if output_gif is None:
        output_gif = os.path.splitext(input_video)[0] + '.gif'
    
    # Open video with cv2
    cap = cv2.VideoCapture(input_video)
    if not cap.isOpened():
        raise ValueError(f"Could not open video file: {input_video}")
    
    # Get video properties
    original_fps = cap.get(cv2.CAP_PROP_FPS)
    frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
    height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
    
    # Use original fps if not specified
    target_fps = fps if fps is not None else int(original_fps)
    
    # Calculate frame skip for target fps
    if target_fps >= original_fps:
        frame_skip = 1
    else:
        frame_skip = int(original_fps / target_fps)
    
    # Calculate new dimensions
    new_width = int(width * scale)
    new_height = int(height * scale)
    
    print(f'🎬 Converting video to GIF...')
    print(f'   Input: {input_video}')
    print(f'   Original: {width}x{height} @ {original_fps:.1f} fps, {frame_count} frames')
    print(f'   Output: {new_width}x{new_height} @ {target_fps} fps')
    
    # Preferred path: ffmpeg palettegen/paletteuse builds one scene-wide
    # optimized palette, producing smaller GIFs than PIL's per-frame
    # adaptive palettes. Fall back to PIL if ffmpeg is missing or fails.
    if optimize and shutil.which('ffmpeg') is not None:
        try:
            print(f'   Using ffmpeg palettegen/paletteuse...')
            _video2gif_ffmpeg(
                input_video,
                output_gif,
                fps=target_fps if target_fps < original_fps else None,
                scale=scale,
                loop=loop,
            )
            input_size = os.path.getsize(input_video) / (1024 * 1024)
            output_size = os.path.getsize(output_gif) / (1024 * 1024)
            print(f'✅ GIF created: {output_gif}')
            print(f'   Input size: {input_size:.2f} MB')
            print(f'   Output size: {output_size:.2f} MB')
            print(f'   Compression ratio: {output_size/input_size:.2%}')
            return output_gif
        except Exception as e:
            print(f'   ⚠️  ffmpeg conversion failed ({e}); falling back to PIL')
    
    # Read frames
    frames = []
    frame_idx = 0
    
    while True:
        ret, frame = cap.read()
        if not ret:
            break
        
        if frame_idx % frame_skip == 0:
            # Convert BGR to RGB
            frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            
            # Resize if needed
            if scale != 1.0:
                frame_rgb = cv2.resize(frame_rgb, (new_width, new_height), 
                                       interpolation=cv2.INTER_AREA)
            
            # Convert to PIL Image
            pil_frame = Image.fromarray(frame_rgb)
            frames.append(pil_frame)
        
        frame_idx += 1
    
    cap.release()
    
    if not frames:
        raise ValueError("No frames extracted from video")
    
    print(f'   Extracted {len(frames)} frames')
    
    # Calculate frame duration in milliseconds
    duration = int(1000 / target_fps)
    
    # Save as GIF with progress bar
    # PIL's save doesn't have a progress callback, so we save frame by frame
    print(f'   Saving GIF ({len(frames)} frames)...')
    
    import io
    
    # For large GIFs, save incrementally to show progress
    total_frames = len(frames)
    
    # Use a temporary buffer approach with progress reporting
    print(f'   [', end='', flush=True)
    bar_width = 40
    
    # We'll save all at once but show a simple progress indicator during optimization
    # Since PIL doesn't support progress callbacks, we simulate with frame processing info
    for i, frame in enumerate(frames):
        # Show progress bar
        progress = (i + 1) / total_frames
        filled = int(bar_width * progress)
        print(f'\r   [{"="*filled}{">" if filled < bar_width else ""}{" "*(bar_width-filled-1 if filled < bar_width else 0)}] {i+1}/{total_frames}', end='', flush=True)
    
    print(f'\r   [{"="*bar_width}] Optimizing...', end='', flush=True)
    
    frames[0].save(
        output_gif,
        save_all=True,
        append_images=frames[1:],
        duration=duration,
        loop=loop,
        optimize=optimize,
        disposal=2,  # Clear frame before drawing next (better for animations)
    )
    print(f'\r   [{"="*bar_width}] Done!          ')
    
    # Report file sizes
    input_size = os.path.getsize(input_video) / (1024 * 1024)
    output_size = os.path.getsize(output_gif) / (1024 * 1024)
    
    print(f'✅ GIF created: {output_gif}')
    print(f'   Input size: {input_size:.2f} MB')
    print(f'   Output size: {output_size:.2f} MB')
    print(f'   Compression ratio: {output_size/input_size:.2%}')
    
    return output_gif


def _video2gif_ffmpeg(input_video, output_gif, fps=None, scale=1.0, loop=0):
    """
    Convert a video to GIF using ffmpeg palettegen/paletteuse.
    
    Two-pass conversion: pass 1 builds a scene-wide 256-color palette
    weighted toward pixels that change between frames (stats_mode=diff),
    pass 2 renders all frames against that palette. This yields smaller
    files than PIL's per-frame adaptive palettes.
    
    Raises on any ffmpeg failure so callers can fall back to PIL.
    """
    import subprocess
    import tempfile
    
    filters = []
    if scale != 1.0:
        # Keep dimensions even for encoder safety
        filters.append(
            f'scale=trunc(iw*{scale}/2)*2:trunc(ih*{scale}/2)*2:flags=lanczos'
        )
    if fps is not None:
        filters.append(f'fps={fps}')
    
    palette_path = os.path.join(
        tempfile.gettempdir(), f'drocat_gif_palette_{os.getpid()}.png'
    )
    try:
        # Pass 1: build the palette
        r = subprocess.run(
            ['ffmpeg', '-y', '-v', 'error', '-i', input_video,
             '-vf', ','.join(filters + ['palettegen=stats_mode=diff']),
             palette_path],
            capture_output=True, text=True,
        )
        if r.returncode != 0:
            raise RuntimeError(f'palettegen failed: {r.stderr.strip().splitlines()[-1] if r.stderr else r.returncode}')
        
        # Pass 2: render frames against the palette
        if filters:
            lavfi = f"{','.join(filters)}[x];[x][1:v]paletteuse=dither=sierra2_4a"
        else:
            lavfi = 'paletteuse=dither=sierra2_4a'
        r = subprocess.run(
            ['ffmpeg', '-y', '-v', 'error', '-i', input_video, '-i', palette_path,
             '-lavfi', lavfi, '-loop', str(loop), output_gif],
            capture_output=True, text=True,
        )
        if r.returncode != 0:
            raise RuntimeError(f'paletteuse failed: {r.stderr.strip().splitlines()[-1] if r.stderr else r.returncode}')
    finally:
        if os.path.exists(palette_path):
            os.remove(palette_path)


def img2pptx(
    input_path: str | list,
    output_pptx: str = None,
    images_per_slide: tuple = (4, 2),
    slide_title: str = None,
    slide_size: str = 'widescreen',
    margin: float = 0.3,
    title_height: int = 60,
    label_fontsize: int = 20,
    title_fontsize: int = 24,
    label_position: str = 'below',
    label_overlay_alpha: float = 0.7,
    cell_padding: float = 0.05,
    include_subfolders: bool = False,
    group_by_subfolder: bool = True,
    font_color: tuple = (0, 0, 0),
    font: str = 'Arial',
    background_color: tuple | str = None,
) -> str:
    """
    Aggregate images to PowerPoint (PPTX) with proper layout, or convert PDF pages to PPTX.
    
    This function has been moved to src/utils/report_utils.py.
    This wrapper is provided for backward compatibility.
    """
    try:
        from utils.report_utils import img2pptx as _img2pptx
    except ImportError:
        try:
            from src.utils.report_utils import img2pptx as _img2pptx
        except ImportError:
             raise ImportError("Could not import img2pptx from utils.report_utils. Please ensure the file exists.")
    
    return _img2pptx(
        input_path=input_path,
        output_pptx=output_pptx,
        images_per_slide=images_per_slide,
        slide_title=slide_title,
        slide_size=slide_size,
        margin=margin,
        title_height=title_height,
        label_fontsize=label_fontsize,
        title_fontsize=title_fontsize,
        label_position=label_position,
        label_overlay_alpha=label_overlay_alpha,
        cell_padding=cell_padding,
        include_subfolders=include_subfolders,
        group_by_subfolder=group_by_subfolder,
        font_color=font_color,
        font=font,
        background_color=background_color
    )