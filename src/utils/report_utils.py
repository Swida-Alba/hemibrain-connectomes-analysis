"""
Report Utilities Module
=======================

This module provides utilities for generating reports, such as PPTX presentations
from images or PDFs. It is designed as a lightweight module with minimal dependencies
to avoid import issues when used from other modules.

Key Functions
-------------
img2pptx : function
    Aggregate images into a PowerPoint presentation with configurable layout,
    or convert PDF pages to PPTX slides.

Dependencies
------------
Required:
    - python-pptx: For PPTX generation
    - Pillow (PIL): For image processing and color parsing

Optional:
    - PyMuPDF (fitz): For PDF to PPTX conversion

Usage
-----
>>> from utils.report_utils import img2pptx
>>> 
>>> # Aggregate images from a folder with black background
>>> img2pptx('/path/to/images/', background_color='black')
>>>
>>> # Convert PDF to PPTX
>>> img2pptx('/path/to/document.pdf')

Notes
-----
This module was extracted from visualize_skeleton.py to provide a lightweight
import path for PPTX generation without requiring heavy dependencies like
navis, neuprint, or selenium.
"""

import os
import io

def img2pptx(
    input_path: str | list,
    output_pptx: str = None,
    images_per_slide: tuple = (4, 2),
    slide_title: str = None,
    slide_size: str = 'widescreen',
    margin: float = 0.3,
    title_height: int = 60,
    label_fontsize: int = 20,
    title_fontsize: int = 24,
    label_position: str = 'below',
    label_overlay_alpha: float = 0.7,
    cell_padding: float = 0.05,
    include_subfolders: bool = False,
    group_by_subfolder: bool = True,
    font_color: tuple = (0, 0, 0),
    font: str = 'Arial',
    background_color: tuple | str = None,
) -> str:
    """
    Aggregate images to PowerPoint (PPTX) with proper layout, or convert PDF pages to PPTX.
    
    This is a static helper function that can be called independently.
    Supports:
    - List of image files → PPTX with grid layout
    - Single PDF file → PPTX with one slide per page
    - Directory of images → PPTX with grid layout
    - Directory with subfolders → PPTX with images from all subfolders
    
    Parameters
    ----------
    input_path : str or list
        Path(s) to input files. Can be:
        - A single PDF file path (converts pages to slides)
        - A single directory path (aggregates all images in the folder)
        - A list of image file paths (aggregates into PPTX)
    output_pptx : str, optional
        Path for the output PPTX file. If None, auto-generated based on input.
    images_per_slide : tuple, default (4, 3)
        (columns, rows) - number of images per slide when aggregating images.
        Not used for PDF conversion.
    slide_title : str, optional
        Title to add to each slide. For image aggregation, can use {page} placeholder
        for page number, {subfolder} for subfolder name. For PDF, defaults to showing page numbers.
    slide_size : str, default 'widescreen'
        Slide dimensions:
        - 'widescreen': 13.333" x 7.5" (16:9)
        - 'standard': 10" x 7.5" (4:3)
        - 'a4': 11.69" x 8.27" (A4 landscape)
    margin : float, default 0.3
        Margin in inches from slide edges.
    title_height : int, default 0
        Height reserved for title in points (pt). Set to 0 to disable title space.
        Recommended: 20-30 for visible titles.
    label_fontsize : int, default 20
        Font size for labels.
    title_fontsize : int, default 24
        Font size for titles.
    label_position : str, default 'below'
        Label placement relative to image.
    label_overlay_alpha : float, default 0.7
        Opacity of label background.
    cell_padding : float, default 0
        Padding around images.
    include_subfolders : bool, default False
        Recursively find images.
    group_by_subfolder : bool, default True
        Group images by subfolder.
    font_color : tuple, default (0, 0, 0)
        Text color (R, G, B).
    font : str, default 'Arial'
        Font family.
    background_color : tuple or str, optional
        Background color involved.
        - Tuple: (R, G, B) e.g., (0, 0, 0) for black
        - String: Hex string '#000000' or named color 'black'
        If None, uses default (white).
    
    Returns
    -------
    str
        Path to the created PPTX file.
    
    Examples
    --------
    # Convert PDF to PPTX
    from utils.report_utils import img2pptx
    img2pptx('/path/to/document.pdf')
    
    # Aggregate images from a folder
    img2pptx('/path/to/image_folder/', images_per_slide=(3, 2))
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        raise ImportError(
            "python-pptx is required for PPTX generation.\n"
            "Install with: pip install python-pptx"
        )
    
    from PIL import Image, ImageColor
    import re
    
    def parse_color(color):
        """Parse color from tuple, hex string, or named color. Returns (r, g, b)."""
        if isinstance(color, (tuple, list)) and len(color) >= 3:
            return tuple(int(c) for c in color[:3])
        if isinstance(color, str):
            if color.startswith('#'):
                hex_col = color.lstrip('#')
                return tuple(int(hex_col[i:i+2], 16) for i in (0, 2, 4))
            return ImageColor.getrgb(color)[:3]
        return (255, 255, 255)  # Default white
    
    # Natural sort function for proper ordering
    def natural_sort_key(s):
        return [int(c) if c.isdigit() else c.lower() for c in re.split(r'(\d+)', str(s))]
    
    # Slide size presets (width, height in inches)
    size_presets = {
        'widescreen': (13.333, 7.5),
        'standard': (10, 7.5),
        'a4': (11.69, 8.27),
    }
    
    if slide_size in size_presets:
        slide_width, slide_height = size_presets[slide_size]
    else:
        slide_width, slide_height = size_presets['widescreen']
    
    # Calculate label height based on fontsize
    label_height_inches = (label_fontsize / 72) * 1.5  # 1.5x line height
    
    # Convert title_height from points to inches
    title_height_inches = title_height / 72 if title_height > 0 else 0
    
    # Auto-detect font color from background if using default black font on dark background
    def get_luminance(color_tuple):
        """Calculate luminance from RGB tuple (0-255)."""
        r, g, b = color_tuple[:3]
        return (r * 0.299 + g * 0.587 + b * 0.114) / 255
    
    effective_font_color = font_color
    if background_color is not None and font_color == (0, 0, 0):
        # Auto-adjust font color for dark backgrounds
        bg_rgb = parse_color(background_color)
        if get_luminance(bg_rgb) < 0.5:
            effective_font_color = (255, 255, 255)  # White text on dark background
    
    # Handle font color (convert 0-1 float to 0-255 int if needed)
    r, g, b = effective_font_color
    if all(isinstance(c, (int, float)) and c <= 1.0 for c in effective_font_color) and not all(c == 0 for c in effective_font_color):
        # Heuristic: if all values are <= 1.0 (and not all 0), assume float 0-1 and convert to 0-255
        print(f"ℹ️  Converting font_color {effective_font_color} from 0-1 range to 0-255 range.")
        r, g, b = [int(c * 255) for c in effective_font_color]
    else:
        r, g, b = [int(c) for c in effective_font_color]
    
    font_color_rgb = (r, g, b)
    
    # Determine input type and gather files
    is_pdf = False
    image_files = []  # List of (path, subfolder_name) tuples
    pdf_path = None
    valid_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif'}
    
    def collect_images_from_dir(dir_path, subfolder_name=''):
        """Collect images from a directory, optionally recursively."""
        collected = []
        for f in sorted(os.listdir(dir_path), key=natural_sort_key):
            full_path = os.path.join(dir_path, f)
            if os.path.isfile(full_path) and os.path.splitext(f)[1].lower() in valid_extensions:
                collected.append((full_path, subfolder_name))
            elif os.path.isdir(full_path) and include_subfolders:
                # Recursively collect from subfolder
                sub_name = f if group_by_subfolder else subfolder_name
                collected.extend(collect_images_from_dir(full_path, sub_name))
        return collected
    
    if isinstance(input_path, str):
        if input_path.lower().endswith('.pdf'):
            is_pdf = True
            pdf_path = input_path
            if not os.path.exists(pdf_path):
                raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        elif os.path.isdir(input_path):
            # Directory of images
            image_files = collect_images_from_dir(input_path, '')
            if not image_files:
                raise ValueError(f"No image files found in directory: {input_path}")
        else:
            # Single image file
            if os.path.exists(input_path):
                image_files = [(input_path, '')]
            else:
                raise FileNotFoundError(f"File not found: {input_path}")
    elif isinstance(input_path, list):
        # List of image paths
        for p in input_path:
            if os.path.exists(p):
                image_files.append((p, ''))
            else:
                print(f"⚠️  Skipping missing file: {p}")
        if not image_files:
            raise ValueError("No valid image files provided")
        image_files = sorted(image_files, key=lambda x: natural_sort_key(x[0]))
    
    # Set output path
    if output_pptx is None:
        if is_pdf:
            output_pptx = os.path.splitext(pdf_path)[0] + '.pptx'
        elif isinstance(input_path, str) and os.path.isdir(input_path):
            output_pptx = os.path.join(input_path, 'aggregated_images.pptx')
        else:
            base_dir = os.path.dirname(image_files[0][0]) if image_files else '.'
            output_pptx = os.path.join(base_dir, 'aggregated_images.pptx')
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    
    # Get blank layout
    blank_layout = prs.slide_layouts[6]  # Blank slide
    
    if is_pdf:
        # Convert PDF pages to PPTX slides
        print(f'📄 Converting PDF to PPTX...')
        print(f'   Input: {pdf_path}')
        
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError(
                "PyMuPDF is required for PDF conversion.\n"
                "Install with: pip install pymupdf"
            )
        
        pdf_doc = fitz.open(pdf_path)
        num_pages = len(pdf_doc)
        print(f'   Pages: {num_pages}')
        
        for page_num in range(num_pages):
            page = pdf_doc[page_num]
            
            # Render page to image with good quality
            zoom = 2.0  # 2x zoom for better quality
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            
            # Convert to PIL Image
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            
            # Create slide
            slide = prs.slides.add_slide(blank_layout)

            # Set background color
            if background_color:
                try:
                    bg = slide.background
                    fill = bg.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(*parse_color(background_color))
                except Exception:
                    pass
            
            # Add title if specified
            content_top = margin
            if slide_title:
                title_text = slide_title.format(page=page_num + 1, subfolder='')
                txBox = slide.shapes.add_textbox(
                    Inches(margin), 
                    Inches(margin / 2),
                    Inches(slide_width - 2 * margin),
                    Inches(title_height_inches)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = title_text
                p.font.size = Pt(title_fontsize)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*font_color_rgb)
                p.alignment = PP_ALIGN.CENTER
                content_top = margin + title_height_inches
            
            # Calculate image placement (fit to slide)
            usable_width = slide_width - 2 * margin
            usable_height = slide_height - content_top - margin
            
            img_width, img_height = img.size
            scale_w = usable_width / (img_width / 72)  # Convert pixels to inches
            scale_h = usable_height / (img_height / 72)
            scale_factor = min(scale_w, scale_h, 1.0)
            
            final_width = (img_width / 72) * scale_factor
            final_height = (img_height / 72) * scale_factor
            
            # Center on slide
            left = (slide_width - final_width) / 2
            top = content_top + (usable_height - final_height) / 2
            
            # Save image temporarily and add to slide
            with io.BytesIO() as img_buffer:
                img.save(img_buffer, format='PNG')
                img_buffer.seek(0)
                slide.shapes.add_picture(
                    img_buffer,
                    Inches(left),
                    Inches(top),
                    Inches(final_width),
                    Inches(final_height)
                )
            
            print(f'\r   Processing page {page_num + 1}/{num_pages}...', end='', flush=True)
        
        pdf_doc.close()
        print(f'\n✅ PPTX created: {output_pptx}')
        print(f'   Slides: {num_pages}')
    
    else:
        # Aggregate images to PPTX with grid layout
        print(f'📊 Aggregating images to PPTX...')
        print(f'   Images: {len(image_files)}')
        print(f'   Layout: {images_per_slide[0]} columns × {images_per_slide[1]} rows')
        if include_subfolders:
            subfolders = set(sf for _, sf in image_files if sf)
            if subfolders:
                print(f'   Subfolders: {len(subfolders)}')
        
        cols, rows = images_per_slide
        images_per_page = cols * rows
        
        # Group images by subfolder if needed
        if group_by_subfolder and include_subfolders:
            # Group by subfolder
            from collections import OrderedDict
            grouped_images = OrderedDict()
            for img_path, subfolder in image_files:
                key = subfolder if subfolder else '_root_'
                if key not in grouped_images:
                    grouped_images[key] = []
                grouped_images[key].append(img_path)
        else:
            # All images in one group
            grouped_images = {'': [img_path for img_path, _ in image_files]}
        
        # Calculate cell dimensions (account for label position)
        # Reserve space for title if title_height is set (> 0)
        has_title_space = title_height > 0
        content_top = margin if not has_title_space else margin + title_height_inches
        usable_width = slide_width - 2 * margin
        usable_height = slide_height - content_top - margin
        cell_width = usable_width / cols
        cell_height = usable_height / rows
        
        total_slides = 0
        total_images_added = 0
        
        for group_name, group_images in grouped_images.items():
            num_slides_for_group = (len(group_images) + images_per_page - 1) // images_per_page
            
            for slide_idx in range(num_slides_for_group):
                slide = prs.slides.add_slide(blank_layout)

                # Set background color
                if background_color:
                    try:
                        bg = slide.background
                        fill = bg.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(*parse_color(background_color))
                    except Exception:
                        pass
                
                # Build title text
                if slide_title:
                    subfolder_display = group_name if group_name != '_root_' else ''
                    title_text = slide_title.format(page=slide_idx + 1, subfolder=subfolder_display)
                elif group_name and group_name != '_root_':
                    title_text = group_name
                    if num_slides_for_group > 1:
                        title_text += f" ({slide_idx + 1}/{num_slides_for_group})"
                else:
                    title_text = ""
                
                if has_title_space and title_text:
                    txBox = slide.shapes.add_textbox(
                        Inches(margin), 
                        Inches(margin / 2),
                        Inches(slide_width - 2 * margin),
                        Inches(title_height_inches)
                    )
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = title_text
                    p.font.size = Pt(title_fontsize)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(*font_color_rgb)
                    # p.alignment = PP_ALIGN.CENTER
                
                # Add images
                start_idx = slide_idx * images_per_page
                end_idx = min(start_idx + images_per_page, len(group_images))
                current_batch = group_images[start_idx:end_idx]
                
                for i, img_path in enumerate(current_batch):
                    col = i % cols
                    row = i // cols
                    
                    cell_left = margin + col * cell_width
                    cell_top = content_top + row * cell_height
                    
                    # Effective cell area with padding (all in inches)
                    eff_left = cell_left + cell_padding
                    eff_top = cell_top + cell_padding
                    eff_width = cell_width - 2 * cell_padding
                    eff_height = cell_height - 2 * cell_padding
                    
                    # Reserve space for label if needed
                    # label_height_inches is roughly 0.3 for 20pt font
                    if label_position in ['below', 'above']:
                        eff_height -= (label_height_inches + 0.05)
                        if label_position == 'above':
                            img_top = eff_top + (label_height_inches + 0.05)
                            label_top = eff_top
                        else:
                            img_top = eff_top
                            label_top = eff_top + eff_height + 0.05
                    else:
                        img_top = eff_top
                    
                    try:
                        # Load image to inspect size/aspect ratio
                        with Image.open(img_path) as img:
                            img_w, img_h = img.size
                        
                        # Calculate scaling to fit in effective area (all in inches)
                        img_width_inches = img_w / 96  # Assume 96 DPI
                        img_height_inches = img_h / 96
                        
                        scale_w = eff_width / img_width_inches
                        scale_h = eff_height / img_height_inches
                        scale = min(scale_w, scale_h)
                        
                        pic_width = img_width_inches * scale
                        pic_height = img_height_inches * scale
                        
                        # Center image in cell area (all in inches)
                        pic_left = eff_left + (eff_width - pic_width) / 2
                        pic_top = img_top + (eff_height - pic_height) / 2
                        
                        # Add picture (convert to Inches EMU)
                        slide.shapes.add_picture(
                            img_path, 
                            Inches(pic_left), 
                            Inches(pic_top), 
                            Inches(pic_width), 
                            Inches(pic_height)
                        )
                        
                        # Add label
                        if label_position != 'none':
                            filename = os.path.basename(img_path)
                            label_text = os.path.splitext(filename)[0]
                            
                            if label_position == 'overlay':
                                # Center at bottom of image
                                label_width = pic_width
                                label_left = pic_left
                                label_top_pos = pic_top + pic_height - label_height_inches
                            else:
                                # Use predefined label position
                                label_width = eff_width
                                label_left = eff_left
                                label_top_pos = label_top
                            
                            txBox = slide.shapes.add_textbox(
                                Inches(label_left), 
                                Inches(label_top_pos), 
                                Inches(label_width), 
                                Inches(label_height_inches)
                            )
                            
                            tf = txBox.text_frame
                            tf.word_wrap = True # Allow wrap
                            p = tf.paragraphs[0]
                            p.text = label_text
                            p.font.size = Pt(label_fontsize)
                            p.font.color.rgb = RGBColor(*font_color_rgb)
                            p.alignment = PP_ALIGN.CENTER
                            
                            # If overlay, maybe add background? (Not easily supported in python-pptx for textframe background, 
                            # would need a shape behind it. Keeping it simple for now)
                            
                        total_images_added += 1
                        print(f'\r   Added image {total_images_added}/{len(image_files)}: {os.path.basename(img_path)}', end='', flush=True)
                        
                    except Exception as e:
                        print(f'\n   ⚠️ Error adding {os.path.basename(img_path)}: {e}')
            
            total_slides += num_slides_for_group
        
        print(f'\n✅ PPTX created: {output_pptx}')
        print(f'   Slides: {total_slides} (Images: {total_images_added})')
    
    # Save presentation
    prs.save(output_pptx)
    
    return output_pptx
