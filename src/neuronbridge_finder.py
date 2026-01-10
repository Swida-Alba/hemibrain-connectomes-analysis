"""
NeuronBridge Finder Module

This module provides a wrapper around the NeuronBridge Python API for finding
correspondences between EM reconstructed neurons (bodyIds) and light microscopy 
(LM) GAL4 driver lines.

Key features:
- Convert neuron bodyIds to matching driver lines (id_to_lines)
- Convert neuron queries (bodyId/instance/type) to lines (neuron_to_lines)
- Reverse search: find neurons matching a driver line (line_to_neuron)
- Support for both CDS and PPPM matching algorithms
- Multi-dataset support (hemibrain, male-cns, FlyWire, etc.)
- Automatic dataset detection from NeuronBridge image metadata
- Caching of API results to reduce redundant calls
- CSV export/import for offline analysis

Dependencies:
- neuronbridge-python
- pandas

Example usage:
    from neuronbridge_finder import NeuronBridgeFinder
    
    nbf = NeuronBridgeFinder()
    
    # Find lines matching a bodyId
    lines = nbf.id_to_lines(636798093, top_n=10)
    
    # Find lines matching a neuron type
    results = nbf.neuron_to_lines('MBON01', top_n=5)
    
    # Find neurons matching a driver line (returns neurons from all datasets)
    neurons = nbf.line_to_neuron('LH173', top_n=10)

Author: Generated for hemibrain-connectomes-analysis project
"""

import json
import os
import re
import time
import warnings
from dataclasses import dataclass, field
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple, Union, TYPE_CHECKING

import numpy as np
import pandas as pd

if TYPE_CHECKING:
    from comparison.label_mapper import LabelMapper

# Try to import tqdm for progress bars
try:
    from tqdm import tqdm
    HAS_TQDM = True
except ImportError:
    HAS_TQDM = False
    tqdm = None

# Try to import LabelMapper for cross-dataset queries
try:
    from comparison.label_mapper import LabelMapper
    HAS_LABELMAPPER = True
except ImportError:
    try:
        from src.comparison.label_mapper import LabelMapper
        HAS_LABELMAPPER = True
    except ImportError:
        HAS_LABELMAPPER = False
        LabelMapper = None  # type: ignore

# Try to import neuronbridge
# Apply patch to fix API compatibility issue with pydantic validation
# (NeuronBridge API added new fields 'defaultSearchLibrary' that the Python client doesn't recognize)
NEURONBRIDGE_AVAILABLE = False

def _patch_neuronbridge_models():
    """
    Patch neuronbridge pydantic models to allow extra fields.
    
    The NeuronBridge API has added new fields that the Python client doesn't
    recognize, causing pydantic validation errors. This patches the models
    to ignore unknown fields.
    """
    try:
        import neuronbridge.client as nb_client
        
        # Get the original CustomSearchConfig class
        OriginalCustomSearchConfig = nb_client.CustomSearchConfig
        
        # Create a new class that allows extra fields
        from pydantic import BaseModel
        
        class PatchedCustomSearchConfig(BaseModel):
            """Patched version that allows extra fields."""
            model_config = {'extra': 'ignore'}
            searchFolder: str = ""
            lmLibraries: list = []
            emLibraries: list = []
        
        # Replace in the module
        nb_client.CustomSearchConfig = PatchedCustomSearchConfig
        
        # Also need to patch DataStore which contains customSearch
        OriginalDataStore = nb_client.DataStore
        
        class PatchedDataStore(BaseModel):
            """Patched version that allows extra fields."""
            model_config = {'extra': 'ignore'}
            label: str = ""
            anatomicalArea: str = ""
            prefixes: dict = {}
            customSearch: PatchedCustomSearchConfig = None
            
            def __init__(self, **data):
                # Handle customSearch conversion
                if 'customSearch' in data and data['customSearch'] is not None:
                    if not isinstance(data['customSearch'], PatchedCustomSearchConfig):
                        data['customSearch'] = PatchedCustomSearchConfig(**data['customSearch'])
                super().__init__(**data)
        
        nb_client.DataStore = PatchedDataStore
        
        # Patch DataConfig
        OriginalDataConfig = nb_client.DataConfig
        
        class PatchedDataConfig(BaseModel):
            """Patched version that allows extra fields."""
            model_config = {'extra': 'ignore'}
            stores: dict = {}
            
            def __init__(self, **data):
                # Convert stores dict values to PatchedDataStore
                if 'stores' in data:
                    new_stores = {}
                    for k, v in data['stores'].items():
                        if isinstance(v, dict):
                            new_stores[k] = PatchedDataStore(**v)
                        else:
                            new_stores[k] = v
                    data['stores'] = new_stores
                super().__init__(**data)
        
        nb_client.DataConfig = PatchedDataConfig
        
        return True
    except Exception as e:
        import warnings
        warnings.warn(f"Failed to patch neuronbridge models: {e}")
        return False

try:
    # Apply patches before creating client
    _patch_neuronbridge_models()
    from neuronbridge.client import Client as NBClient
    NEURONBRIDGE_AVAILABLE = True
except ImportError:
    warnings.warn(
        "neuronbridge-python is not installed. "
        "Install it with: pip install neuronbridge-python"
    )
except Exception as e:
    warnings.warn(f"Failed to initialize neuronbridge: {e}")

# Mapping from NeuronBridge library names to our local dataset folder names
LIBRARY_TO_DATASET = {
    'FlyEM_Hemibrain_v1.2.1': 'hemibrain_v1_2_1',
    'FlyEM_Hemibrain_v1.2': 'hemibrain_v1_2_1',
    'FlyEM_Hemibrain': 'hemibrain_v1_2_1',
    'FlyEM_MANC_v1.2.1': 'manc_v1_2_1',
    'FlyEM_MANC_v1.0': 'manc_v1_0',
    'FlyEM_MANC': 'manc_v1_0',
    'FlyEM_Male_CNS_Brain_v0.9': 'male-cns_v0_9',
    'FlyEM_Male_CNS_VNC_v0.9': 'male-cns_v0_9',
    'FlyEM_Male_CNS_v0.9': 'male-cns_v0_9',
    'FlyEM_Male_CNS': 'male-cns_v0_9',
    'FlyWire_FAFB': 'flywire_FAFB_v783',
    'FlyWire_FAFB_v783': 'flywire_FAFB_v783',
    'FlyWire_FAFB_v783_realign': 'flywire_FAFB_v783',
    'FlyWire_BANC': 'flywire_BANC_v626',
    'FlyWire_BANC_v626': 'flywire_BANC_v626',
    'FlyEM_Optic_Lobe': 'optic-lobe_v1_1',
}

# Mapping to human-readable dataset names (for display/neuprint format)
LIBRARY_TO_DATASET_NAME = {
    'FlyEM_Hemibrain_v1.2.1': 'hemibrain:v1.2.1',
    'FlyEM_Hemibrain_v1.2': 'hemibrain:v1.2.1',
    'FlyEM_Hemibrain': 'hemibrain:v1.2.1',
    'FlyEM_MANC_v1.2.1': 'manc:v1.2.1',
    'FlyEM_MANC_v1.0': 'manc:v1.0',
    'FlyEM_MANC': 'manc:v1.0',
    'FlyEM_Male_CNS_Brain_v0.9': 'male-cns:v0.9',
    'FlyEM_Male_CNS_VNC_v0.9': 'male-cns:v0.9',
    'FlyEM_Male_CNS_v0.9': 'male-cns:v0.9',
    'FlyEM_Male_CNS': 'male-cns:v0.9',
    'FlyWire_FAFB': 'flywire_FAFB_v783',
    'FlyWire_FAFB_v783': 'flywire_FAFB_v783',
    'FlyWire_FAFB_v783_realign': 'flywire_FAFB_v783',
    'FlyWire_BANC': 'flywire_BANC_v626',
    'FlyWire_BANC_v626': 'flywire_BANC_v626',
    'FlyEM_Optic_Lobe': 'optic-lobe:v1.1',
}

# Dataset abbreviations for expression matrix type naming
# Format: {dataset_folder_name: abbreviation}
DATASET_ABBREVIATIONS = {
    'hemibrain_v1_2_1': 'HEMI',
    'hemibrain:v1.2.1': 'HEMI',
    'male-cns_v0_9': 'MCNS',
    'male-cns:v0.9': 'MCNS',
    'manc_v1_0': 'MANC',
    'manc:v1.0': 'MANC',
    'manc_v1_2_1': 'MANC',
    'manc:v1.2.1': 'MANC',
    'flywire_FAFB_v783': 'FAFB',
    'flywire_BANC_v626': 'BANC',
    'optic-lobe_v1_1': 'OLOB',
    'optic-lobe:v1.1': 'OLOB',
}


def _to_int_bodyid(val):
    """
    Convert a bodyId value to integer, handling various input formats.
    
    Handles:
    - Integer values (passthrough)
    - String values like '5813128953' 
    - Float values like 5813128953.0
    - String floats like '5813128953.0'
    
    Returns the original value if conversion fails.
    """
    import numpy as np
    
    # Already an int
    if isinstance(val, (int, np.integer)):
        return int(val)
    
    # Float (including numpy float)
    if isinstance(val, (float, np.floating)):
        return int(val)
    
    # String
    if isinstance(val, str):
        # Try direct int conversion first (handles '5813128953')
        try:
            return int(val)
        except ValueError:
            pass
        # Try float then int (handles '5813128953.0')
        try:
            return int(float(val))
        except ValueError:
            pass
    
    # Return original if conversion fails
    return val


# Line name prefixes for classification
# GAL4/LexA lines: typically VT (Vienna Tile), R (Rubin), GMR, etc.
GAL4_LEXA_PREFIXES = ('VT', 'R', 'GMR')
# Split-GAL4 lines: SS (Split Screen), LH (Lateral Horn), MB (Mushroom Body), etc.
SPLIT_GAL4_PREFIXES = ('SS', 'LH', 'MB', 'IS', 'OL', 'LC', 'LLPC', 'LPC', 'JRC_SS', 'BJD_SS')

# Valid input values for validation (case-insensitive)
VALID_MATCH_TYPES = {'cds', 'pppm', 'both'}
VALID_REGIONS = {'brain', 'vnc', 'all'}
VALID_SIMILARITY_METHODS = {'jaccard', 'weighted_jaccard', 'rank_correlation'}
VALID_SORT_BY = {'completeness', 'max'}


@dataclass
class NeuronBridgeFinder:
    """
    A class for finding correspondences between EM neurons and LM driver lines.
    
    This class wraps the NeuronBridge API to provide convenient methods for
    mapping between EM body IDs and GAL4 driver lines across multiple datasets.
    
    Parameters
    ----------
    datasets_path : str, optional
        Path to the datasets folder containing neuron_df CSV files.
        Default: auto-detect from module location.
    use_cache : bool
        Whether to cache API results locally. Default: True
    cache_folder : str, optional
        Folder for cached results. Default: auto-detect.
    verbose : bool
        Print progress messages. Default: True
    separate_splitgal4 : bool
        If True, separate results into GAL4/LexA and Split-GAL4 categories.
        When enabled, download_img_for_top_n_lines applies separately to each category.
        Default: False
    neuprint_token : str, optional
        NeuPrint API token for pulling missing datasets. If not provided, will check
        NEUPRINT_TOKEN or NEUPRINT_APPLICATION_CREDENTIALS environment variables.
        Without a token, local dataset features (type lookups, specificity) will be skipped.
        Get your token at: https://neuprint.janelia.org/account
    neuprint_server : str
        NeuPrint server URL. Default: 'https://neuprint.janelia.org'
    match_type : str
        Default match algorithm: 'cds', 'pppm', or 'both' (case-insensitive). Default: 'cds'
    region : str
        Filter images by anatomical region: 'Brain', 'VNC', or 'All' (case-insensitive). Default: 'All'
    max_api_images_per_line : int
        Maximum LM images to process per driver line for API calls. Use -1 for unlimited. Default: -1
        Images are pre-filtered by match_type availability before limiting.
    
    Attributes
    ----------
    _client : NBClient
        The NeuronBridge API client.
    _neuron_dfs : dict
        Dictionary mapping dataset names to loaded neuron DataFrames.
    
    Example
    -------
    >>> nbf = NeuronBridgeFinder()
    >>> lines = nbf.id_to_lines(636798093, top_n=5)
    >>> print(lines)
    
    # Separate GAL4/LexA from Split-GAL4 lines
    >>> nbf = NeuronBridgeFinder(separate_splitgal4=True)
    >>> results = nbf.find_lines_batch('MBON01', download_img_for_top_n_lines=5)  # 5 GAL4 + 5 Split-GAL4
    
    # With NeuPrint token for pulling missing datasets
    >>> nbf = NeuronBridgeFinder(neuprint_token='your_token_here')
    """
    
    datasets_path: Optional[str] = None
    use_cache: bool = True
    cache_folder: Optional[str] = None
    verbose: bool = True
    separate_splitgal4: bool = False
    neuprint_token: Optional[str] = None
    neuprint_server: str = 'https://neuprint.janelia.org'
    match_type: str = 'cds'
    region: str = 'All'
    max_api_images_per_line: int = -1
    
    # Private fields
    _client: Any = field(init=False, repr=False, default=None)
    _neuron_dfs: Dict[str, pd.DataFrame] = field(init=False, repr=False, default_factory=dict)
    _suppress_loading_msgs: bool = field(init=False, repr=False, default=False)
    _batch_mode: bool = field(init=False, repr=False, default=False)
    _warning_collector: List[str] = field(init=False, repr=False, default_factory=list)
    
    def __post_init__(self):
        """Initialize the finder after dataclass initialization."""
        if not NEURONBRIDGE_AVAILABLE:
            raise ImportError(
                "neuronbridge-python is required. Install with: pip install neuronbridge-python"
            )
        
        # Validate and normalize match_type (case-insensitive)
        if isinstance(self.match_type, str):
            normalized_match_type = self.match_type.lower().strip()
            if normalized_match_type not in VALID_MATCH_TYPES:
                raise ValueError(
                    f"Invalid match_type: '{self.match_type}'. "
                    f"Must be one of: {', '.join(sorted(VALID_MATCH_TYPES))}"
                )
            # Use object.__setattr__ for frozen dataclass compatibility
            object.__setattr__(self, 'match_type', normalized_match_type)
        
        # Validate and normalize region (case-insensitive)
        if isinstance(self.region, str):
            normalized_region = self.region.lower().strip()
            if normalized_region not in VALID_REGIONS:
                raise ValueError(
                    f"Invalid region: '{self.region}'. "
                    f"Must be one of: {', '.join(sorted(VALID_REGIONS))}"
                )
            # Store in title case for display consistency (Brain, VNC, All)
            region_display = {'brain': 'Brain', 'vnc': 'VNC', 'all': 'All'}
            object.__setattr__(self, 'region', region_display[normalized_region])
        
        # Set default datasets path
        if self.datasets_path is None:
            module_dir = os.path.dirname(os.path.abspath(__file__))
            project_root = os.path.dirname(module_dir)
            self.datasets_path = os.path.join(project_root, 'datasets')
        
        # Set default cache folder
        if self.cache_folder is None:
            module_dir = os.path.dirname(os.path.abspath(__file__))
            project_root = os.path.dirname(module_dir)
            self.cache_folder = os.path.join(project_root, 'cache', 'neuronbridge')
        
        # Create cache folder if needed
        if self.use_cache:
            os.makedirs(self.cache_folder, exist_ok=True)
        
        # Initialize client
        self._init_client()
    
    def _vprint(self, msg: str, end: str = '\n', force: bool = False):
        """
        Print message if verbose mode is enabled.
        
        Uses tqdm.write() if tqdm is available and progress bars might be active,
        otherwise falls back to regular print().
        
        Parameters
        ----------
        msg : str
            Message to print
        end : str
            String appended after the message (default: newline)
        force : bool
            If True, print even in batch mode (default: False)
        """
        if self.verbose and (not self._batch_mode or force):
            # Use tqdm.write if available to avoid conflicts with progress bars
            if HAS_TQDM:
                try:
                    from tqdm import tqdm
                    tqdm.write(msg, end=end)
                except:
                    print(msg, end=end)
            else:
                print(msg, end=end)
    
    def _validate_match_type(self, match_type: str) -> str:
        """
        Validate and normalize match_type parameter.
        
        Parameters
        ----------
        match_type : str
            Match type to validate ('cds', 'pppm', or 'both')
            
        Returns
        -------
        str
            Normalized (lowercase) match_type
            
        Raises
        ------
        ValueError
            If match_type is not valid
        """
        normalized = match_type.lower().strip()
        if normalized not in VALID_MATCH_TYPES:
            raise ValueError(
                f"Invalid match_type: '{match_type}'. "
                f"Must be one of: {', '.join(sorted(VALID_MATCH_TYPES))}"
            )
        return normalized
    
    def _validate_similarity_method(self, method: str) -> str:
        """
        Validate and normalize similarity_method parameter.
        
        Parameters
        ----------
        method : str
            Similarity method to validate
            
        Returns
        -------
        str
            Normalized (lowercase) similarity method
            
        Raises
        ------
        ValueError
            If method is not valid
        """
        normalized = method.lower().strip()
        if normalized not in VALID_SIMILARITY_METHODS:
            raise ValueError(
                f"Invalid similarity_method: '{method}'. "
                f"Must be one of: {', '.join(sorted(VALID_SIMILARITY_METHODS))}"
            )
        return normalized
    
    def _validate_sort_by(self, sort_by: str) -> str:
        """
        Validate and normalize sort_by parameter.
        
        Parameters
        ----------
        sort_by : str
            Sort method: 'completeness' or 'max'
            
        Returns
        -------
        str
            Normalized (lowercase) sort_by value
            
        Raises
        ------
        ValueError
            If sort_by is not valid
        """
        normalized = sort_by.lower().strip()
        if normalized not in VALID_SORT_BY:
            raise ValueError(
                f"Invalid sort_by: '{sort_by}'. "
                f"Must be one of: {', '.join(sorted(VALID_SORT_BY))}"
            )
        return normalized
    
    def _print_warning_summary(self):
        """Print collected warnings as a summary."""
        if not self._warning_collector:
            return
        
        from collections import Counter
        
        # Group warnings by type
        server_unreachable = []
        no_gal4_files = []
        no_files_at_all = []
        
        for warning in self._warning_collector:
            if "server not accessible" in warning:
                server_unreachable.append(warning)
            elif "No files from" in warning:
                no_gal4_files.append(warning)
            elif "No images found" in warning:
                no_files_at_all.append(warning)
        
        # Print summary
        if server_unreachable:
            # Extract unique servers and line counts
            server_issues = {}
            for w in server_unreachable:
                if "flimg.janelia.org" in w:
                    line = w.split("for ")[-1].split(":")[0].strip()
                    server_issues.setdefault("flimg.janelia.org (VT GAL4)", []).append(line)
            
            if server_issues:
                self._vprint("\n⚠️  Server Access Issues:", force=True)
                for server, lines_affected in server_issues.items():
                    self._vprint(f"   • {server}: {len(lines_affected)} line(s) unreachable", force=True)
                    self._vprint(f"     (Attempted MCFO fallback from S3)", force=True)
        
        if no_gal4_files:
            lines = [w.split("for ")[-1].split(",")[0].strip() for w in no_gal4_files]
            self._vprint(f"\n📋 {len(lines)} line(s) used MCFO fallback (no GAL4/SplitGAL4 images)", force=True)
        
        if no_files_at_all:
            lines = [w.split("for ")[-1].split(" in")[0].strip() for w in no_files_at_all]
            if lines:
                self._vprint(f"\n❌ {len(lines)} line(s) had no images in any collection:", force=True)
                self._vprint(f"   {', '.join(lines[:10])}{'...' if len(lines) > 10 else ''}", force=True)
        
        # Clear warnings after printing
        self._warning_collector.clear()
    
    def _retry_with_backoff(
        self,
        func,
        *args,
        max_retries: int = 3,
        initial_delay: float = 1.0,
        **kwargs
    ):
        """
        Retry a function with exponential backoff for transient network errors.
        
        Parameters
        ----------
        func : callable
            Function to retry
        max_retries : int
            Maximum number of retry attempts (default: 3)
        initial_delay : float
            Initial delay in seconds (default: 1.0)
        *args, **kwargs
            Arguments to pass to func
            
        Returns
        -------
        Any
            Result from func, or None if all retries failed
        """
        last_error = None
        
        for attempt in range(max_retries):
            try:
                return func(*args, **kwargs)
            except Exception as e:
                last_error = e
                error_str = str(e).lower()
                
                # Check if it's a retryable network error
                is_retryable = any([
                    'incompleteread' in error_str,
                    'ssl' in error_str and 'eof' in error_str,
                    'connection' in error_str and 'broken' in error_str,
                    'max retries exceeded' in error_str,
                    'connection reset' in error_str,
                    'timeout' in error_str
                ])
                
                if not is_retryable or attempt >= max_retries - 1:
                    # Not retryable or last attempt - give up
                    raise
                
                # Wait with exponential backoff
                delay = initial_delay * (2 ** attempt)
                time.sleep(delay)
        
        # Should never reach here, but just in case
        if last_error:
            raise last_error
    
    def _normalize_dataset_name(self, dataset: str) -> str:
        """
        Normalize dataset name for comparison.
        
        Handles different naming conventions:
        - 'flywire_FAFB_v783' -> 'flywire-fafb'
        - 'flywire_fafb:v783' -> 'flywire-fafb'
        - 'male-cns:v0.9' -> 'male-cns'
        - 'hemibrain:v1.2.1' -> 'hemibrain'
        - 'hemibrain_v1_2_1' -> 'hemibrain'
        
        Parameters
        ----------
        dataset : str
            Dataset name in any format
            
        Returns
        -------
        str
            Normalized dataset name (lowercase, no version, hyphens for separators)
        """
        import re
        
        # Convert to lowercase
        ds = dataset.lower()
        
        # Remove version info in various formats:
        # - ':v1.2.1' or ':v783' (colon format)
        # - '_v1_2_1' or '_v783' (underscore format)
        # First handle colon-separated version
        ds = ds.split(':')[0]
        
        # Remove underscore-separated version suffix (e.g., '_v783', '_v1_2_1')
        # Match _v followed by digits and optional underscored sub-versions
        ds = re.sub(r'_v\d+(_\d+)*$', '', ds)
        
        # Normalize separators: convert underscores to hyphens
        ds = ds.replace('_', '-')
        
        return ds
    
    def _save_parameters(
        self,
        output_path: str,
        function_name: str,
        function_params: Dict[str, Any],
        filename: str = 'parameters.json'
    ) -> str:
        """
        Save all module-level and function-level parameters to a JSON file.
        
        This ensures reproducibility by recording all configuration used
        to generate the analysis results.
        
        Parameters
        ----------
        output_path : str
            Directory to save the parameters file.
        function_name : str
            Name of the function being called (e.g., 'analyze_colabeling').
        function_params : dict
            Dictionary of function-level parameters.
        filename : str
            Output filename. Default: 'parameters.json'
            
        Returns
        -------
        str
            Path to the saved parameters file.
        """
        import json
        from datetime import datetime
        
        # Collect module-level (instance) parameters
        module_params = {
            'datasets_path': self.datasets_path,
            'use_cache': self.use_cache,
            'cache_folder': self.cache_folder,
            'verbose': self.verbose if isinstance(self.verbose, bool) else str(self.verbose),
            'separate_splitgal4': self.separate_splitgal4,
            'neuprint_server': self.neuprint_server,
            'match_type': self.match_type,
            'region': self.region,
            'max_api_images_per_line': self.max_api_images_per_line,
            # Note: neuprint_token is sensitive, store presence only
            'has_neuprint_token': self.neuprint_token is not None
        }
        
        # Process function params to make them JSON serializable
        serializable_params = {}
        for key, value in function_params.items():
            if isinstance(value, (str, int, float, bool, type(None))):
                serializable_params[key] = value
            elif isinstance(value, (list, tuple)):
                # Handle lists/tuples - convert to list
                serializable_params[key] = list(value)
            elif isinstance(value, dict):
                serializable_params[key] = value
            elif hasattr(value, 'tolist'):
                # Handle numpy arrays
                serializable_params[key] = value.tolist()
            elif hasattr(value, '__dict__'):
                # Handle objects with __dict__
                serializable_params[key] = str(value)
            else:
                serializable_params[key] = str(value)
        
        # Build complete parameters dict
        params = {
            'metadata': {
                'timestamp': datetime.now().isoformat(),
                'function': function_name,
                'neuronbridge_finder_version': '3.1'
            },
            'module_params': module_params,
            'function_params': serializable_params
        }
        
        # Save to JSON file
        params_path = os.path.join(output_path, filename)
        with open(params_path, 'w', encoding='utf-8') as f:
            json.dump(params, f, indent=2, ensure_ascii=False)
        
        return params_path
    
    def _filter_images_by_region(self, images: List[Any]) -> List[Any]:
        """
        Filter images based on the anatomical region setting.
        
        Filters LM or EM images based on their anatomicalArea attribute.
        Only processes filtering when self.region is 'Brain' or 'VNC'.
        
        Parameters
        ----------
        images : list
            List of LM or EM image objects from NeuronBridge API.
            
        Returns
        -------
        list
            Filtered list of images matching the region criteria.
        """
        if self.region == 'All':
            return images
        
        filtered = []
        for img in images:
            # Get anatomicalArea attribute if it exists
            area = getattr(img, 'anatomicalArea', None)
            if area is None:
                # If no area specified, include by default
                filtered.append(img)
                continue
            
            # Normalize area name (case-insensitive)
            area_lower = area.lower()
            
            if self.region == 'Brain':
                # Include if area contains 'brain' but not 'vnc'
                if 'brain' in area_lower and 'vnc' not in area_lower:
                    filtered.append(img)
            elif self.region == 'VNC':
                # Include if area contains 'vnc'
                if 'vnc' in area_lower:
                    filtered.append(img)
        
        return filtered
    
    def _filter_images_by_match_availability(
        self, 
        images: List[Any], 
        match_type: str
    ) -> List[Any]:
        """
        Filter images to only include those with available match results for the specified type.
        
        Pre-scans image metadata to check for CDSResults/PPPMResults URLs,
        skipping images that don't have the required match type data.
        
        Parameters
        ----------
        images : list
            List of LM image objects from NeuronBridge API.
        match_type : str
            Match algorithm: 'cds', 'pppm', or 'both'.
            
        Returns
        -------
        list
            Filtered list of images that have the required match results available.
        """
        if match_type == 'both':
            # For 'both', include images that have either CDS or PPPM
            filtered = []
            for img in images:
                files = getattr(img, 'files', None)
                if files is None:
                    continue
                has_cds = getattr(files, 'CDSResults', None) is not None
                has_pppm = getattr(files, 'PPPMResults', None) is not None
                if has_cds or has_pppm:
                    filtered.append(img)
            return filtered
        
        elif match_type == 'cds':
            # Only include images with CDS results
            filtered = []
            for img in images:
                files = getattr(img, 'files', None)
                if files is None:
                    continue
                if getattr(files, 'CDSResults', None) is not None:
                    filtered.append(img)
            return filtered
        
        elif match_type == 'pppm':
            # Only include images with PPPM results
            filtered = []
            for img in images:
                files = getattr(img, 'files', None)
                if files is None:
                    continue
                if getattr(files, 'PPPMResults', None) is not None:
                    filtered.append(img)
            return filtered
        
        # Unknown match_type, return all
        return images
    
    def _classify_line_type(self, line_name: str) -> str:
        """
        Classify a driver line as 'gal4_lexa' or 'split_gal4'.
        
        Classification rules:
        - GAL4/LexA: Lines starting with VT, R, GMR (e.g., VT037867, R10A06)
        - Split-GAL4: Lines starting with SS, LH, MB, IS, OL, LC, etc.
        
        Parameters
        ----------
        line_name : str
            The driver line name to classify.
            
        Returns
        -------
        str
            'gal4_lexa' or 'split_gal4'
        """
        if not line_name:
            return 'gal4_lexa'  # default
        
        line_upper = line_name.upper()
        
        # Check Split-GAL4 prefixes first (more specific)
        for prefix in SPLIT_GAL4_PREFIXES:
            if line_upper.startswith(prefix.upper()):
                return 'split_gal4'
        
        # Check GAL4/LexA prefixes
        for prefix in GAL4_LEXA_PREFIXES:
            if line_upper.startswith(prefix.upper()):
                return 'gal4_lexa'
        
        # Default to gal4_lexa for unknown prefixes
        return 'gal4_lexa'
    
    def _separate_lines_by_type(
        self, 
        lines_df: pd.DataFrame,
        line_column: str = 'line'
    ) -> Dict[str, pd.DataFrame]:
        """
        Separate a DataFrame of lines into GAL4/LexA and Split-GAL4 categories.
        
        Parameters
        ----------
        lines_df : pd.DataFrame
            DataFrame containing line results.
        line_column : str
            Column name containing line names. Default: 'line'
            
        Returns
        -------
        dict
            Dictionary with 'gal4_lexa' and 'split_gal4' keys, each containing
            a DataFrame of lines for that category.
        """
        if lines_df.empty or line_column not in lines_df.columns:
            return {'gal4_lexa': pd.DataFrame(), 'split_gal4': pd.DataFrame()}
        
        # Add line_type column
        lines_df = lines_df.copy()
        lines_df['line_type'] = lines_df[line_column].apply(self._classify_line_type)
        
        return {
            'gal4_lexa': lines_df[lines_df['line_type'] == 'gal4_lexa'].copy(),
            'split_gal4': lines_df[lines_df['line_type'] == 'split_gal4'].copy()
        }
    
    def _calculate_expression_entropy(self, type_counts: Dict[str, int]) -> float:
        """
        Calculate Shannon entropy of expression pattern.
        
        H = -Σ p_i × log2(p_i)
        
        Lower entropy = more specific (labels fewer types more strongly)
        Higher entropy = more promiscuous (labels many types equally)
        
        Parameters
        ----------
        type_counts : dict
            Dictionary mapping neuron types to their occurrence counts.
            
        Returns
        -------
        float
            Shannon entropy in bits. Range: 0 (perfectly specific) to log2(N) (uniform).
        """
        if not type_counts:
            return 0.0
        
        total = sum(type_counts.values())
        if total == 0:
            return 0.0
        
        entropy = 0.0
        for count in type_counts.values():
            if count > 0:
                p = count / total
                entropy -= p * np.log2(p)
        
        return entropy
    
    def _calculate_weighted_specificity(
        self, 
        neurons_df: pd.DataFrame,
        queried_types_lower: set,
        score_column: str = 'score'
    ) -> Dict[str, float]:
        """
        Calculate expression-strength-weighted specificity metrics.
        
        Weights each neuron's contribution by its NeuronBridge match score,
        so high-confidence matches count more than low-confidence ones.
        
        Parameters
        ----------
        neurons_df : pd.DataFrame
            DataFrame from line_to_neuron with 'type' and score columns.
        queried_types_lower : set
            Set of queried neuron types (lowercase).
        score_column : str
            Name of score column to use for weighting.
            
        Returns
        -------
        dict
            Dictionary with:
            - 'weighted_type_proportion': Score-weighted proportion of queried types
            - 'weighted_queried_score': Total score for queried types
            - 'weighted_total_score': Total score for all types
            - 'mean_queried_score': Mean score for queried types
        """
        if neurons_df.empty or score_column not in neurons_df.columns:
            return {
                'weighted_type_proportion': 0.0,
                'weighted_queried_score': 0.0,
                'weighted_total_score': 0.0,
                'mean_queried_score': 0.0
            }
        
        # Calculate total weighted scores
        total_score = neurons_df[score_column].sum()
        
        # Get scores for queried types
        queried_mask = neurons_df['type'].fillna('').str.lower().isin(queried_types_lower)
        queried_scores = neurons_df.loc[queried_mask, score_column]
        
        queried_total = queried_scores.sum()
        queried_mean = queried_scores.mean() if len(queried_scores) > 0 else 0.0
        
        # Weighted proportion
        weighted_proportion = queried_total / total_score if total_score > 0 else 0.0
        
        return {
            'weighted_type_proportion': weighted_proportion,
            'weighted_queried_score': queried_total,
            'weighted_total_score': total_score,
            'mean_queried_score': queried_mean
        }
    
    def _build_colabeling_matrix(
        self,
        lines: List[str],
        match_type: str = 'cds',
        top_n: int = 100,
        similarity_method: str = 'weighted_jaccard',
        min_score: float = 0.0,
        min_type_avg_score: float = 0.0
    ) -> Tuple[pd.DataFrame, Dict[str, set]]:
        """
        Build a co-labeling matrix showing how often pairs of lines label the same cell types.
        
        Supports multiple similarity measures:
        - 'jaccard': Binary Jaccard similarity (presence/absence of types)
        - 'weighted_jaccard': Jaccard weighted by match scores
        - 'rank_correlation': Spearman correlation of type rankings based on scores
        
        Parameters
        ----------
        lines : list of str
            List of driver line names.
        match_type : str
            Match algorithm for line_to_neuron.
        top_n : int
            Number of top matches to consider per line.
        similarity_method : str
            Similarity method: 'jaccard', 'weighted_jaccard', or 'rank_correlation'.
            Default: 'weighted_jaccard'
        min_score : float
            Minimum score threshold for individual neurons. Default: 0.0 (no filter).
        min_type_avg_score : float
            Minimum average score threshold for types (across all lines). Default: 0.0 (no filter).
            
        Returns
        -------
        tuple
            (co_labeling_matrix, line_type_sets)
            - co_labeling_matrix: DataFrame with similarities based on cell types
            - line_type_sets: Dict mapping line names to sets of cell type names
        """
        # Validate inputs using helper methods
        match_type = self._validate_match_type(match_type)
        similarity_method = self._validate_similarity_method(similarity_method)
        
        self._vprint(f"\n🔗 Building co-labeling matrix for {len(lines)} lines...")
        self._vprint(f"   📊 Similarity method: {similarity_method}")
        if min_score > 0:
            self._vprint(f"   📊 Min neuron score: {min_score:,.0f}")
        if min_type_avg_score > 0:
            self._vprint(f"   📊 Min type avg score: {min_type_avg_score:,.0f}")
        self._vprint(f"   ⏱️  Note: Fetching neurons for each line to build type-based similarity matrix")
        
        # Collect type sets AND scores for each line
        line_neuron_sets = {}
        line_type_scores = {}
        
        # Enable batch mode to suppress individual cache messages
        self._batch_mode = True
        
        if HAS_TQDM and self.verbose:
            from tqdm import tqdm as tqdm_progress
            iterator = tqdm_progress(
                lines, 
                desc="   🔍 Collecting neurons per line", 
                unit="line",
                bar_format='{desc}: {percentage:3.0f}%|{bar}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]',
                ncols=100
            )
        else:
            iterator = lines
        
        for idx, line_name in enumerate(iterator):
            if HAS_TQDM and self.verbose:
                iterator.set_description(f"   🔍 [{idx+1}/{len(lines)}] Fetching neurons for {line_name}")
            try:
                neurons_df = self.line_to_neuron(line_name, match_type=match_type)
                if not neurons_df.empty:
                    neurons_top = neurons_df.head(top_n)
                    # Apply minimum score filter if specified
                    if min_score > 0:
                        neurons_top = neurons_top[neurons_top['score'] >= min_score]
                    # Use cell types for similarity (not bodyIds) - this gives meaningful co-labeling
                    # Two lines labeling the same cell types are considered similar
                    type_set = set()
                    type_scores = {}  # type -> score mapping
                    for _, row in neurons_top.iterrows():
                        n_type = row.get('type', '')
                        score = row.get('score', 0.0)
                        if n_type and str(n_type).lower() not in ['', 'nan', 'none', 'unknown']:
                            type_key = str(n_type).lower()
                            type_set.add(type_key)
                            # Keep the maximum score if type appears multiple times
                            type_scores[type_key] = max(type_scores.get(type_key, 0.0), float(score))
                    line_neuron_sets[line_name] = type_set
                    line_type_scores[line_name] = type_scores
                else:
                    line_neuron_sets[line_name] = set()
                    line_type_scores[line_name] = {}
            except Exception as e:
                line_neuron_sets[line_name] = set()
                line_type_scores[line_name] = {}
                if self.verbose:
                    self._vprint(f"   ⚠️ Error getting neurons for {line_name}: {e}")
        
        # Disable batch mode after fetching
        self._batch_mode = False
        
        # Apply min_type_avg_score filter: remove types with low average scores
        if min_type_avg_score > 0:
            # Collect all types and their average scores across lines
            all_type_scores = {}  # type -> list of scores
            for line_scores in line_type_scores.values():
                for type_key, score in line_scores.items():
                    if type_key not in all_type_scores:
                        all_type_scores[type_key] = []
                    all_type_scores[type_key].append(score)
            
            # Filter types by average score
            types_to_keep = set()
            for type_key, scores in all_type_scores.items():
                avg_score = sum(scores) / len(scores)
                if avg_score >= min_type_avg_score:
                    types_to_keep.add(type_key)
            
            # Remove filtered types from line sets and scores
            removed_types = set(all_type_scores.keys()) - types_to_keep
            if removed_types:
                self._vprint(f"   📊 Filtered {len(removed_types)} types with avg score < {min_type_avg_score:,.0f}")
                self._vprint(f"   📊 Keeping {len(types_to_keep)} high-confidence types")
                for line_name in line_neuron_sets:
                    line_neuron_sets[line_name] = line_neuron_sets[line_name] & types_to_keep
                    line_type_scores[line_name] = {
                        k: v for k, v in line_type_scores[line_name].items() 
                        if k in types_to_keep
                    }
        
        # Build similarity matrix based on selected method
        self._vprint(f"   🔢 Computing {similarity_method} similarities between {len(lines)} lines...")
        n = len(lines)
        matrix = np.zeros((n, n))
        
        if similarity_method == 'jaccard':
            # Binary Jaccard (original implementation)
            for i, line_i in enumerate(lines):
                set_i = line_neuron_sets.get(line_i, set())
                for j, line_j in enumerate(lines):
                    if i == j:
                        matrix[i, j] = 1.0
                    elif j > i:
                        set_j = line_neuron_sets.get(line_j, set())
                        if set_i and set_j:
                            intersection = len(set_i & set_j)
                            union = len(set_i | set_j)
                            jaccard = intersection / union if union > 0 else 0.0
                        else:
                            jaccard = 0.0
                        matrix[i, j] = jaccard
                        matrix[j, i] = jaccard
                        
        elif similarity_method == 'weighted_jaccard':
            # Weighted Jaccard using match scores
            for i, line_i in enumerate(lines):
                scores_i = line_type_scores.get(line_i, {})
                for j, line_j in enumerate(lines):
                    if i == j:
                        matrix[i, j] = 1.0
                    elif j > i:
                        scores_j = line_type_scores.get(line_j, {})
                        if scores_i and scores_j:
                            all_types = set(scores_i.keys()) | set(scores_j.keys())
                            intersection_sum = sum(
                                min(scores_i.get(t, 0.0), scores_j.get(t, 0.0)) 
                                for t in all_types
                            )
                            union_sum = sum(
                                max(scores_i.get(t, 0.0), scores_j.get(t, 0.0)) 
                                for t in all_types
                            )
                            w_jaccard = intersection_sum / union_sum if union_sum > 0 else 0.0
                        else:
                            w_jaccard = 0.0
                        matrix[i, j] = w_jaccard
                        matrix[j, i] = w_jaccard
                        
        elif similarity_method == 'rank_correlation':
            # Spearman correlation on overlapping types
            from scipy.stats import spearmanr
            for i, line_i in enumerate(lines):
                scores_i = line_type_scores.get(line_i, {})
                for j, line_j in enumerate(lines):
                    if i == j:
                        matrix[i, j] = 1.0
                    elif j > i:
                        scores_j = line_type_scores.get(line_j, {})
                        if scores_i and scores_j:
                            # Find types in common
                            common_types = set(scores_i.keys()) & set(scores_j.keys())
                            if len(common_types) > 1:
                                # Get scores for common types
                                values_i = [scores_i[t] for t in common_types]
                                values_j = [scores_j[t] for t in common_types]
                                # Compute Spearman correlation
                                corr, _ = spearmanr(values_i, values_j)
                                # Handle NaN (can occur with constant values)
                                if np.isnan(corr):
                                    corr = 0.0
                            else:
                                corr = 0.0
                        else:
                            corr = 0.0
                        matrix[i, j] = corr
                        matrix[j, i] = corr
        else:
            raise ValueError(f"Unknown similarity method: {similarity_method}")
        
        co_labeling_df = pd.DataFrame(matrix, index=lines, columns=lines)
        
        return co_labeling_df, line_neuron_sets
    
    def _calculate_colabeling_sparsity(self, co_labeling_matrix: pd.DataFrame, threshold: float = 0.1) -> Dict[str, float]:
        """
        Calculate sparsity metrics from co-labeling matrix.
        
        Sparsity measures how unique each line's labeling pattern is.
        High sparsity = line labels neurons that few other lines label.
        
        Parameters
        ----------
        co_labeling_matrix : pd.DataFrame
            Jaccard similarity matrix from _build_colabeling_matrix.
        threshold : float
            Similarity threshold to consider lines as "co-labeling".
            
        Returns
        -------
        dict
            Dictionary mapping line names to sparsity scores (0-1, higher = more unique).
        """
        sparsity_scores = {}
        
        for line in co_labeling_matrix.index:
            # Get similarities with other lines (exclude self)
            similarities = co_labeling_matrix.loc[line].drop(line)
            
            # Count how many lines have significant overlap
            n_colabeling = (similarities > threshold).sum()
            total_other_lines = len(similarities)
            
            # Sparsity = 1 - (proportion of co-labeling lines)
            sparsity = 1.0 - (n_colabeling / total_other_lines) if total_other_lines > 0 else 1.0
            
            # Also compute mean non-self similarity (lower = more unique)
            mean_similarity = similarities.mean() if len(similarities) > 0 else 0.0
            
            sparsity_scores[line] = {
                'colabel_sparsity': sparsity,
                'n_colabeling_lines': int(n_colabeling),
                'mean_colabel_similarity': mean_similarity
            }
        
        return sparsity_scores
    
    def _sort_expression_matrix(
        self,
        expression_df: pd.DataFrame,
        as_types_rows: bool = True
    ) -> pd.DataFrame:
        """
        Sort expression matrix by co-labeling quality.
        
        Sorting criteria (when as_types_rows=True, types as rows):
        1. Types labeled in ALL lines (no zeros) come first
        2. Within complete types: sorted by min_score (higher = more consistent)
        3. Types with partial labeling: sorted by num_nonzero (desc), total_score (desc)
        
        Parameters
        ----------
        expression_df : pd.DataFrame
            Expression matrix. If as_types_rows=False, expects Lines × Types format.
        as_types_rows : bool
            If True (default), output has types as rows and lines as columns.
            If False, output has lines as rows and types as columns.
            
        Returns
        -------
        pd.DataFrame
            Sorted expression matrix with types as rows (if as_types_rows=True).
        """
        # Determine if input needs transposing to get Types × Lines format
        # When as_types_rows=True (default), we want types as rows (many rows, few columns)
        # Input from _calculate_mutual_information is Lines × Types (few rows, many columns)
        # Heuristic: if significantly more columns than rows, it's likely Lines × Types format
        needs_transpose = (
            expression_df.index.name == 'line' or 
            (as_types_rows and len(expression_df.columns) > len(expression_df) * 2)
        )
        
        if needs_transpose:
            expression_transposed = expression_df.T
        else:
            expression_transposed = expression_df.copy()
        
        if len(expression_transposed) == 0 or len(expression_transposed.columns) == 0:
            return expression_transposed if as_types_rows else expression_transposed.T
        
        n_lines = len(expression_transposed.columns)
        
        # Calculate metrics for each type (row)
        nonzero_count = (expression_transposed > 0).sum(axis=1)
        is_complete = (nonzero_count == n_lines).astype(int)
        
        # Create sorting key:
        # - Primary: whether labeled in all lines (1 if all, 0 if not) - descending
        # - Secondary: min_score except 0s iteratedly (tuple of sorted non-zero scores) - descending
        #   This sorts by the smallest non-zero score, then the next smallest, etc.
        
        # Helper to get sorted non-zeros tuple
        def get_nonzero_tuple(row):
            return tuple(sorted(row[row > 0]))
            
        nonzero_tuples = expression_transposed.apply(get_nonzero_tuple, axis=1)
        
        # Create temp dataframe for sorting
        sort_df = pd.DataFrame({
            'is_complete': is_complete,
            'nonzero_count': nonzero_count,
            'nonzero_tuples': nonzero_tuples
        }, index=expression_transposed.index)
        
        # Sort by the composite key
        sorted_index = sort_df.sort_values(
            by=['is_complete', 'nonzero_count', 'nonzero_tuples'], 
            ascending=[False, False, False]
        ).index
        expression_transposed = expression_transposed.loc[sorted_index]
        
        if as_types_rows:
            return expression_transposed
        else:
            return expression_transposed.T
    
    def visualize_colabeling_matrix(
        self,
        co_labeling_matrix: pd.DataFrame,
        output_path: str,
        title: str = "Co-Labeling Matrix",
        color_scale: str = 'purple',
        filename: str = 'colabeling_matrix.html',
        zmin: Optional[float] = 0.0,
        zmax: Optional[float] = 1.0
    ) -> str:
        """
        Visualize co-labeling matrix as a heatmap using VisConnMatInteractive.
        
        Parameters
        ----------
        co_labeling_matrix : pd.DataFrame
            Similarity matrix from _build_colabeling_matrix.
        output_path : str
            Directory to save the heatmap HTML file.
        title : str
            Title for the heatmap.
        color_scale : str
            Color scale preset: 'purple', 'green', 'blue', 'orange', 'red'.
        filename : str
            Filename for the HTML file. Default: 'colabeling_matrix.html'
        zmin : float, optional
            Minimum value for color scale. Default: 0.0
        zmax : float, optional
            Maximum value for color scale. Default: 1.0 (for similarity matrices)
            
        Returns
        -------
        str
            Path to the created heatmap file.
        """
        try:
            from vispath_pkg import VisConnMatInteractive
        except ImportError:
            # Try to add vispath-subproject/src to path
            import sys
            from pathlib import Path
            
            # Find repo root (assuming we are in src/)
            current_file = Path(__file__).resolve()
            repo_root = current_file.parent.parent
            vispath_src = repo_root / 'vispath-subproject' / 'src'
            
            if vispath_src.exists() and str(vispath_src) not in sys.path:
                sys.path.append(str(vispath_src))
            
            try:
                from vispath_pkg import VisConnMatInteractive
            except ImportError:
                self._vprint("   ⚠️ Could not import VisConnMatInteractive from vispath_pkg")
                return ""
        
        # Create output directory if needed
        os.makedirs(output_path, exist_ok=True)
        
        # Resolve color scale preset to Plotly format
        color_scales = {
            'green': [[0, 'rgb(255,255,255)'], [1, 'rgb(14,83,13)']],
            'purple': [[0, 'rgb(255,255,255)'], [1, 'rgb(104,55,164)']],
            'orange': [[0, 'rgb(255,255,255)'], [1, 'rgb(204,102,0)']],
            'blue': [[0, 'rgb(255,255,255)'], [1, 'rgb(31,119,180)']],
            'red': [[0, 'rgb(255,255,255)'], [1, 'rgb(214,39,40)']],
        }
        resolved_color_scale = color_scales.get(color_scale, color_scales['purple'])
        
        # Create heatmap file path
        full_path = os.path.join(output_path, filename)
        
        # Create heatmap using VisConnMatInteractive with fixed color scale range
        VisConnMatInteractive(
            co_labeling_matrix,
            filename=full_path,
            title=title,
            color_scale=resolved_color_scale,
            zmin=zmin,
            zmax=zmax,
            showfig=False,
            verbose=self.verbose
        )
        
        self._vprint(f"   📊 Created heatmap: {full_path}")
        return full_path

    def _calculate_mutual_information(
        self,
        lines: List[str],
        queried_types: List[str],
        match_type: str = 'cds',
        top_n: int = 100,
        output_path: Optional[str] = None,
        min_score: float = 0.0,
        min_type_avg_score: float = 0.0
    ) -> Tuple[pd.DataFrame, pd.DataFrame, Dict[str, pd.DataFrame]]:
        """
        Calculate mutual information between driver lines and neuron types.
        
        MI(L; T) = Σ p(l,t) × log2(p(l,t) / (p(l) × p(t)))
        
        This measures how much knowing a line's expression tells you about neuron type.
        High MI = line expression is highly informative about neuron type.
        
        Parameters
        ----------
        lines : list of str
            List of driver line names to analyze.
        queried_types : list of str
            List of queried neuron types (for highlighting in visualization).
        match_type : str
            Match algorithm for line_to_neuron.
        top_n : int
            Number of top matches to consider per line.
        output_path : str, optional
            If provided, save per-line neuron details CSVs to this directory.
        min_score : float
            Score threshold for visualization marker only. Does NOT filter data.
            Expression matrix includes ALL neurons. Default: 0.0.
        min_type_avg_score : float
            NOT USED in this function. Kept for API compatibility.
            Type filtering is done in _build_colabeling_matrix for similarity calculations.
            Expression matrix includes ALL types regardless of this parameter.
            
        Returns
        -------
        tuple
            (mi_df, expression_matrix, line_neurons_dict, labeling_info)
            - mi_df: DataFrame with MI values per line
            - expression_matrix: Score-based expression matrix (lines × types)
              Types are prefixed with dataset abbreviation: {ABBREV}_{type}
            - line_neurons_dict: Dict mapping line name to neurons DataFrame
            - labeling_info: DataFrame with type, dataset, and per-line score columns
        """
        self._vprint(f"\n📊 Calculating mutual information for {len(lines)} lines...")
        if min_score > 0:
            self._vprint(f"   📊 Min neuron score (visualization threshold): {min_score:,.0f}")
        # Note: min_type_avg_score is NOT used here - filtering happens in _build_colabeling_matrix
        self._vprint(f"   ⏱️  Note: Fetching neuron types for each line (may take time)")
        
        # Normalize queried types for comparison (but keep original case)
        queried_types_lower = set(t.lower() for t in queried_types if t)
        
        # Enable batch mode to suppress cache messages
        self._batch_mode = True
        
        # Build expression matrix: rows = lines, cols = neuron types with dataset prefix
        # Value = max score if line labels that type, 0 otherwise
        # Format: {dataset_abbrev}_{type} e.g., HEMI_DM4, MCNS_dm4
        all_prefixed_types = set()  # Case-sensitive types with dataset prefix
        line_type_sets = {}  # Lowercase types per line (for MI calculation)
        line_prefixed_type_scores = {}  # Prefixed type -> score per line
        line_neurons_dict = {}  # Store neurons DataFrame for each line
        
        # For labeling_info: store (type, dataset) -> {line: score}
        labeling_info_data = {}  # (type, dataset) -> {line: score}
        
        if HAS_TQDM and self.verbose:
            from tqdm import tqdm as tqdm_progress
            iterator = tqdm_progress(
                lines, 
                desc="   🧬 Building expression matrix", 
                unit="line",
                bar_format='{desc}: {percentage:3.0f}%|{bar}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]',
                ncols=100
            )
        else:
            iterator = lines
        
        for idx, line_name in enumerate(iterator):
            if HAS_TQDM and self.verbose:
                iterator.set_description(f"   🧬 [{idx+1}/{len(lines)}] Processing {line_name}")
            try:
                neurons_df = self.line_to_neuron(line_name, match_type=match_type, top_n=top_n)
                if not neurons_df.empty:
                    # NO min_score filtering for expression matrix - include ALL data
                    # min_score is only used for visualization (labeling distribution plots)
                    
                    # Get lowercase types for MI calculation (from ALL neurons)
                    types = neurons_df['type'].fillna('Unknown').unique()
                    types_lower = set(t.lower() for t in types)
                    line_type_sets[line_name] = types_lower
                    
                    # Build prefixed type scores: {ABBREV}_{type} with original case
                    # Use ALL data for expression matrix (no min_score filter)
                    line_prefixed_scores = {}
                    for dataset in neurons_df['dataset'].dropna().unique():
                        ds_df = neurons_df[neurons_df['dataset'] == dataset]
                        # Get dataset abbreviation
                        ds_abbrev = DATASET_ABBREVIATIONS.get(dataset, dataset[:4].upper())
                        
                        # Group by original case type and get max score
                        for type_name, type_df in ds_df.groupby(ds_df['type'].fillna('Unknown')):
                            max_score = type_df['score'].max()
                            prefixed_type = f"{ds_abbrev}_{type_name}"
                            
                            # Store for expression matrix
                            all_prefixed_types.add(prefixed_type)
                            line_prefixed_scores[prefixed_type] = max(
                                line_prefixed_scores.get(prefixed_type, 0.0), 
                                max_score
                            )
                    
                    # Build labeling_info from ALL data (same as expression matrix)
                    for dataset in neurons_df['dataset'].dropna().unique():
                        ds_df = neurons_df[neurons_df['dataset'] == dataset]
                        ds_abbrev = DATASET_ABBREVIATIONS.get(dataset, dataset[:4].upper())
                        
                        for type_name, type_df in ds_df.groupby(ds_df['type'].fillna('Unknown')):
                            max_score = type_df['score'].max()
                            
                            # Store for labeling_info
                            key = (type_name, dataset)
                            if key not in labeling_info_data:
                                labeling_info_data[key] = {}
                            labeling_info_data[key][line_name] = max(
                                labeling_info_data[key].get(line_name, 0.0),
                                max_score
                            )
                    
                    line_prefixed_type_scores[line_name] = line_prefixed_scores
                    
                    # Store neurons DataFrame with visualization marker column
                    # _passes_min_score is used by visualization functions only
                    neurons_to_store = neurons_df.copy()
                    neurons_to_store['_passes_min_score'] = neurons_to_store['score'] >= (min_score if min_score > 0 else 0)
                    line_neurons_dict[line_name] = neurons_to_store
                else:
                    line_type_sets[line_name] = set()
                    line_prefixed_type_scores[line_name] = {}
                    line_neurons_dict[line_name] = pd.DataFrame()
            except Exception as e:
                line_type_sets[line_name] = set()
                line_prefixed_type_scores[line_name] = {}
                line_neurons_dict[line_name] = pd.DataFrame()
        
        # Disable batch mode after fetching
        self._batch_mode = False
        
        if not all_prefixed_types:
            self._vprint("   ⚠️ No types found for any line")
            return pd.DataFrame(), pd.DataFrame(), line_neurons_dict, pd.DataFrame()
        
        # NOTE: min_type_avg_score filtering is NOT applied here
        # Expression matrix includes ALL types regardless of score
        # Type filtering for similarity/clustering is done in _build_colabeling_matrix
        
        # Create score-based expression matrix with prefixed types (case-sensitive)
        # Value = max match score if line labels that type, 0 otherwise
        all_prefixed_types_list = sorted(all_prefixed_types)
        expression_matrix = np.zeros((len(lines), len(all_prefixed_types_list)))
        
        for i, line in enumerate(lines):
            for j, prefixed_type in enumerate(all_prefixed_types_list):
                score = line_prefixed_type_scores.get(line, {}).get(prefixed_type, 0.0)
                expression_matrix[i, j] = score
        
        expression_df = pd.DataFrame(
            expression_matrix,
            index=lines,
            columns=all_prefixed_types_list
        )
        
        # Create labeling_info DataFrame
        # Columns: type, dataset, {line1_score}, {line2_score}, ...
        labeling_rows = []
        for (type_name, dataset), line_scores in labeling_info_data.items():
            row = {'type': type_name, 'dataset': dataset}
            for line in lines:
                row[line] = line_scores.get(line, 0.0)
            labeling_rows.append(row)
        
        labeling_info = pd.DataFrame(labeling_rows)
        if not labeling_info.empty:
            # Sort by max score across lines (descending)
            score_cols = [col for col in labeling_info.columns if col not in ['type', 'dataset']]
            labeling_info['_max_score'] = labeling_info[score_cols].max(axis=1)
            labeling_info['_min_score'] = labeling_info[score_cols].min(axis=1)
            labeling_info['_nonzero'] = (labeling_info[score_cols] > 0).sum(axis=1)
            # Sort: all-lines-labeled first (by min_score), then partial (by nonzero, max_score)
            labeling_info['_is_complete'] = labeling_info['_nonzero'] == len(score_cols)
            labeling_info = labeling_info.sort_values(
                ['_is_complete', '_min_score', '_nonzero', '_max_score'],
                ascending=[False, False, False, False]
            )
            labeling_info = labeling_info.drop(columns=['_max_score', '_min_score', '_nonzero', '_is_complete'])
        
        # Calculate MI for each line using lowercase types (original logic)
        # MI(L; T) = H(T) - H(T|L)
        all_types_lower = set()
        for types_set in line_type_sets.values():
            all_types_lower.update(types_set)
        all_types_lower_list = sorted(all_types_lower)
        
        n_lines = len(lines)
        n_types = len(all_types_lower_list)
        
        # Build binary matrix for MI calculation (based on lowercase types)
        mi_binary_matrix = np.zeros((len(lines), n_types))
        for i, line in enumerate(lines):
            types_in_line = line_type_sets.get(line, set())
            for j, ntype in enumerate(all_types_lower_list):
                if ntype in types_in_line:
                    mi_binary_matrix[i, j] = 1.0
        
        # Marginal probability of each type (across all lines)
        type_counts = mi_binary_matrix.sum(axis=0)  # How many lines label each type
        p_type = type_counts / n_lines  # P(type)
        p_type = np.clip(p_type, 1e-10, 1)  # Avoid log(0)
        
        # Entropy of type distribution H(T)
        H_T = -np.sum(p_type * np.log2(p_type))
        
        mi_results = []
        
        for i, line in enumerate(lines):
            types_labeled = line_type_sets.get(line, set())
            n_labeled = len(types_labeled)
            
            if n_labeled == 0:
                mi_results.append({
                    'line': line,
                    'mutual_information': 0.0,
                    'normalized_mi': 0.0,
                    'n_types_labeled': 0,
                    'queried_type_coverage': 0.0
                })
                continue
            
            # Conditional entropy H(T|L=l)
            # P(t|L=l) = 1/n_labeled if line labels type t, 0 otherwise
            p_t_given_l = np.zeros(n_types)
            for j, ntype in enumerate(all_types_lower_list):
                if ntype in types_labeled:
                    p_t_given_l[j] = 1.0 / n_labeled
            
            # H(T|L=l) = -Σ p(t|l) log2(p(t|l))
            p_t_given_l_nonzero = p_t_given_l[p_t_given_l > 0]
            H_T_given_L = -np.sum(p_t_given_l_nonzero * np.log2(p_t_given_l_nonzero))
            
            # Point-wise MI for this line
            # This is simplified: we compute reduction in entropy
            mi_line = H_T - H_T_given_L
            
            # Normalized MI (0-1 scale)
            max_mi = H_T  # Maximum possible MI is H(T)
            normalized_mi = mi_line / max_mi if max_mi > 0 else 0.0
            
            # Queried type coverage
            covered = len(types_labeled & queried_types_lower)
            coverage = covered / len(queried_types_lower) if queried_types_lower else 0.0
            
            mi_results.append({
                'line': line,
                'mutual_information': mi_line,
                'normalized_mi': normalized_mi,
                'n_types_labeled': n_labeled,
                'queried_type_coverage': coverage
            })
        
        mi_df = pd.DataFrame(mi_results)
        
        self._vprint(f"   ✓ MI calculated for {len(mi_df)} lines")
        self._vprint(f"   Total prefixed types: {len(all_prefixed_types_list)}")
        self._vprint(f"   Total unique types (lowercase): {n_types}")
        self._vprint(f"   Type entropy H(T): {H_T:.3f} bits")
        
        return mi_df, expression_df, line_neurons_dict, labeling_info
    
    def visualize_expression_matrix(
        self,
        expression_df: pd.DataFrame,
        output_path: str,
        queried_types: Optional[List[str]] = None,
        title: str = "Line × Type Expression Matrix",
        color_scale: str = 'green',
        top_n_types: int = 100
    ) -> str:
        """
        Visualize the expression matrix (lines × types) as a heatmap.
        
        Parameters
        ----------
        expression_df : pd.DataFrame
            Binary expression matrix from _calculate_mutual_information.
        output_path : str
            Directory to save the heatmap HTML file.
        queried_types : list of str, optional
            Queried types to highlight in the matrix.
        title : str
            Title for the heatmap.
        color_scale : str
            Color scale preset.
        top_n_types : int
            Maximum number of types to display (default: 100).
            Types are sorted by total expression score across all lines.
            
        Returns
        -------
        str
            Path to the created heatmap file.
        """
        try:
            from vispath_pkg import VisConnMatInteractive
        except ImportError:
            # Try to add vispath-subproject/src to path
            import sys
            from pathlib import Path
            
            # Find repo root (assuming we are in src/)
            current_file = Path(__file__).resolve()
            repo_root = current_file.parent.parent
            vispath_src = repo_root / 'vispath-subproject' / 'src'
            
            if vispath_src.exists() and str(vispath_src) not in sys.path:
                sys.path.append(str(vispath_src))
            
            try:
                from vispath_pkg import VisConnMatInteractive
            except ImportError:
                self._vprint("   ⚠️ Could not import VisConnMatInteractive from vispath_pkg")
                return ""
        
        os.makedirs(output_path, exist_ok=True)
        
        # Transpose first to get types as rows
        # Preserve input order (already sorted by co-labeling quality)
        expression_transposed = expression_df.T
        
        # Limit to top N types by total expression score
        total_types = len(expression_transposed)
        if total_types > top_n_types:
            # Filter to top N by total score, but preserve quality-based sorting
            type_totals = expression_transposed.sum(axis=1)
            top_types = set(type_totals.nlargest(top_n_types).index)
            # Keep only top types but maintain original order (sorted by quality)
            expression_transposed = expression_transposed.loc[
                [idx for idx in expression_transposed.index if idx in top_types]
            ]
            # Re-sort by quality (min_score descending within complete/incomplete groups)
            expression_transposed = self._sort_expression_matrix(
                expression_transposed.T, as_types_rows=True
            )
            self._vprint(f"   📊 Showing top {top_n_types} types (of {total_types} total, filtered by expression score)")
            # Add filter info to title
            title = f"{title} [Top {top_n_types} types]"
        else:
            self._vprint(f"   📊 Showing all {total_types} types")
        
        # Resolve color scale preset to Plotly format
        color_scales = {
            'green': [[0, 'rgb(255,255,255)'], [1, 'rgb(14,83,13)']],
            'purple': [[0, 'rgb(255,255,255)'], [1, 'rgb(104,55,164)']],
            'orange': [[0, 'rgb(255,255,255)'], [1, 'rgb(204,102,0)']],
            'blue': [[0, 'rgb(255,255,255)'], [1, 'rgb(31,119,180)']],
            'red': [[0, 'rgb(255,255,255)'], [1, 'rgb(214,39,40)']],
        }
        resolved_color_scale = color_scales.get(color_scale, color_scales['green'])
        
        # Create heatmap file path
        filename = os.path.join(output_path, 'expression_matrix.html')
        
        # Calculate initial dimensions based on matrix size
        # For expression matrices: width ~8px per column, height ~20px per row
        # Minimum: 800x800, Maximum: 2400x4000
        n_rows = len(expression_transposed)
        n_cols = len(expression_transposed.columns)
        init_width = max(800, min(2400, n_cols * 80 + 200))  # 80px per column + margin
        init_height = max(800, min(4000, n_rows * 20 + 200))  # 20px per row + margin
        
        # Create heatmap using VisConnMatInteractive (transposed: types × lines)
        VisConnMatInteractive(
            expression_transposed,
            filename=filename,
            title=title.replace("Lines × Types", "Types × Lines"),
            color_scale=resolved_color_scale,
            showfig=False,
            verbose=self.verbose,
            init_width=init_width,
            init_height=init_height
        )
        
        self._vprint(f"   📊 Created expression matrix heatmap: {filename}")
        return filename

    def visualize_expression_matrix_merged(
        self,
        expression_df: pd.DataFrame,
        output_path: str,
        queried_types: Optional[List[str]] = None,
        title: str = "Line × Type Expression Matrix (Merged)",
        color_scale: str = 'green',
        top_n_types: int = 100,
        aggregation: str = 'max'
    ) -> str:
        """
        Visualize the expression matrix with same neuron types merged across datasets.
        
        Types from different datasets (e.g., 'MCNS_aMe12', 'FAFB_aMe12', 'HEMI_aMe12')
        are merged into a single row (e.g., 'aMe12') using the specified aggregation.
        
        Parameters
        ----------
        expression_df : pd.DataFrame
            Expression matrix from _calculate_mutual_information (Lines × Types format).
            Column names should be prefixed with dataset abbreviation (e.g., 'MCNS_aMe12').
        output_path : str
            Directory to save the heatmap HTML file.
        queried_types : list of str, optional
            Queried types to highlight in the matrix.
        title : str
            Title for the heatmap.
        color_scale : str
            Color scale preset.
        top_n_types : int
            Maximum number of types to display (default: 100).
        aggregation : str
            How to aggregate scores across datasets: 'max' (default), 'mean', 'sum'.
            
        Returns
        -------
        str
            Path to the created heatmap file.
        """
        try:
            from vispath_pkg import VisConnMatInteractive
        except ImportError:
            # Try to add vispath-subproject/src to path
            import sys
            from pathlib import Path
            
            # Find repo root (assuming we are in src/)
            current_file = Path(__file__).resolve()
            repo_root = current_file.parent.parent
            vispath_src = repo_root / 'vispath-subproject' / 'src'
            
            if vispath_src.exists() and str(vispath_src) not in sys.path:
                sys.path.append(str(vispath_src))
            
            try:
                from vispath_pkg import VisConnMatInteractive
            except ImportError:
                self._vprint("   ⚠️ Could not import VisConnMatInteractive from vispath_pkg")
                return ""
        
        os.makedirs(output_path, exist_ok=True)
        
        # Step 1: Extract base type names from prefixed columns
        # e.g., 'MCNS_aMe12' -> 'aMe12', 'FAFB_Dm4' -> 'Dm4'
        def extract_base_type(prefixed_type: str) -> str:
            """Extract base type name from prefixed type (e.g., 'MCNS_aMe12' -> 'aMe12')."""
            parts = prefixed_type.split('_', 1)
            if len(parts) > 1:
                return parts[1]  # Return everything after first underscore
            return prefixed_type  # Return as-is if no underscore
        
        # Build mapping: base_type -> list of prefixed columns
        base_type_to_columns = {}
        for col in expression_df.columns:
            base_type = extract_base_type(col)
            if base_type not in base_type_to_columns:
                base_type_to_columns[base_type] = []
            base_type_to_columns[base_type].append(col)
        
        # Step 2: Merge columns by base type using specified aggregation
        merged_data = {}
        for base_type, columns in base_type_to_columns.items():
            subset = expression_df[columns]
            if aggregation == 'max':
                merged_data[base_type] = subset.max(axis=1)
            elif aggregation == 'mean':
                # Mean of non-zero values only
                merged_data[base_type] = subset.apply(
                    lambda row: row[row > 0].mean() if (row > 0).any() else 0.0, axis=1
                )
            elif aggregation == 'sum':
                merged_data[base_type] = subset.sum(axis=1)
            else:
                merged_data[base_type] = subset.max(axis=1)  # Default to max
        
        merged_df = pd.DataFrame(merged_data)
        
        # Log merge statistics
        original_types = len(expression_df.columns)
        merged_types = len(merged_df.columns)
        self._vprint(f"   📊 Merged {original_types} prefixed types → {merged_types} base types (aggregation: {aggregation})")
        
        # Step 3: Transpose to get types as rows
        expression_transposed = merged_df.T
        
        # Step 4: Sort by co-labeling quality (same logic as original)
        n_lines = len(expression_transposed.columns)
        nonzero_count = (expression_transposed > 0).sum(axis=1)
        total_score = expression_transposed.sum(axis=1)
        min_score_col = expression_transposed.min(axis=1)
        
        is_complete = (nonzero_count == n_lines).astype(int)
        sort_df = pd.DataFrame({
            'is_complete': -is_complete,
            'min_score': -min_score_col,
            'nonzero_count': -nonzero_count,
            'total_score': -total_score
        }, index=expression_transposed.index)
        sorted_index = sort_df.sort_values(['is_complete', 'min_score', 'nonzero_count', 'total_score']).index
        expression_transposed = expression_transposed.loc[sorted_index]
        
        # Step 5: Limit to top N types
        total_types = len(expression_transposed)
        if total_types > top_n_types:
            type_totals = expression_transposed.sum(axis=1)
            top_types = set(type_totals.nlargest(top_n_types).index)
            expression_transposed = expression_transposed.loc[
                [idx for idx in expression_transposed.index if idx in top_types]
            ]
            # Re-sort after filtering
            n_lines = len(expression_transposed.columns)
            nonzero_count = (expression_transposed > 0).sum(axis=1)
            total_score = expression_transposed.sum(axis=1)
            min_score_col = expression_transposed.min(axis=1)
            is_complete = (nonzero_count == n_lines).astype(int)
            sort_df = pd.DataFrame({
                'is_complete': -is_complete,
                'min_score': -min_score_col,
                'nonzero_count': -nonzero_count,
                'total_score': -total_score
            }, index=expression_transposed.index)
            sorted_index = sort_df.sort_values(['is_complete', 'min_score', 'nonzero_count', 'total_score']).index
            expression_transposed = expression_transposed.loc[sorted_index]
            
            self._vprint(f"   📊 Showing top {top_n_types} merged types (of {total_types} total)")
            title = f"{title} [Top {top_n_types} types]"
        else:
            self._vprint(f"   📊 Showing all {total_types} merged types")
        
        # Step 6: Resolve color scale
        color_scales = {
            'green': [[0, 'rgb(255,255,255)'], [1, 'rgb(14,83,13)']],
            'purple': [[0, 'rgb(255,255,255)'], [1, 'rgb(104,55,164)']],
            'orange': [[0, 'rgb(255,255,255)'], [1, 'rgb(204,102,0)']],
            'blue': [[0, 'rgb(255,255,255)'], [1, 'rgb(31,119,180)']],
            'red': [[0, 'rgb(255,255,255)'], [1, 'rgb(214,39,40)']],
        }
        resolved_color_scale = color_scales.get(color_scale, color_scales['green'])
        
        # Step 7: Create heatmap
        filename = os.path.join(output_path, 'expression_matrix_merged.html')
        
        n_rows = len(expression_transposed)
        n_cols = len(expression_transposed.columns)
        init_width = max(800, min(2400, n_cols * 80 + 200))
        init_height = max(800, min(4000, n_rows * 20 + 200))
        
        VisConnMatInteractive(
            expression_transposed,
            filename=filename,
            title=title.replace("Lines × Types", "Types × Lines"),
            color_scale=resolved_color_scale,
            showfig=False,
            verbose=self.verbose,
            init_width=init_width,
            init_height=init_height
        )
        
        # Also save the merged CSV
        csv_filename = os.path.join(output_path, 'expression_matrix_merged.csv')
        expression_transposed.to_csv(csv_filename)
        self._vprint(f"   💾 Saved merged expression matrix CSV: {csv_filename}")
        
        self._vprint(f"   📊 Created merged expression matrix heatmap: {filename}")
        return filename

    def visualize_labeling_distribution(
        self,
        data: pd.DataFrame,
        output_path: str,
        score_column: str = 'score',
        label_column: str = 'type',
        title: str = "Labeling Distribution",
        color: str = '#1f77b4',
        filename: str = 'labeling_distribution.html',
        show_threshold: Optional[float] = None,
        group_by: Optional[str] = None
    ) -> str:
        """
        Visualize labeling distribution as a mountain-shaped histogram.
        
        The plot shows scores with highest values in the center, descending
        symmetrically to both sides (like a mountain/pyramid shape).
        
        Parameters
        ----------
        data : pd.DataFrame
            DataFrame with score and label columns.
        output_path : str
            Directory to save the HTML file.
        score_column : str
            Column name for scores. Default: 'score'
        label_column : str
            Column name for labels (types/neurons). Default: 'type'
        title : str
            Plot title. Default: "Labeling Distribution"
        color : str
            Bar color. Default: '#1f77b4' (Category10 blue)
        filename : str
            Output filename. Default: 'labeling_distribution.html'
        show_threshold : float, optional
            If provided, draw a horizontal threshold line.
        group_by : str, optional
            If provided, group data by this column and create subplots.
            
        Returns
        -------
        str
            Path to the created HTML file.
        """
        import plotly.graph_objects as go
        from plotly.subplots import make_subplots
        
        os.makedirs(output_path, exist_ok=True)
        
        if data.empty or score_column not in data.columns:
            self._vprint(f"   ⚠️ No data for labeling distribution")
            return ""
        
        # Category10 palette (same as bokeh.palettes.Category10)
        CATEGORY10 = ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd',
                      '#8c564b', '#e377c2', '#7f7f7f', '#bcbd22', '#17becf']
        
        def create_mountain_order(scores, labels):
            """Rearrange so highest score is in middle, descending to both sides."""
            sorted_pairs = sorted(zip(scores, labels), key=lambda x: x[0], reverse=True)
            n = len(sorted_pairs)
            if n == 0:
                return [], []
            
            result_scores = [0] * n
            result_labels = [''] * n
            
            left = n // 2 - 1 if n % 2 == 0 else n // 2
            right = n // 2 if n % 2 == 0 else n // 2
            
            for i, (score, label) in enumerate(sorted_pairs):
                if i == 0:
                    center = n // 2
                    result_scores[center] = score
                    result_labels[center] = label
                elif i % 2 == 1:
                    result_scores[left] = score
                    result_labels[left] = label
                    left -= 1
                else:
                    result_scores[right] = score
                    result_labels[right] = label
                    right += 1
            
            return result_scores, result_labels
        
        full_path = os.path.join(output_path, filename)
        
        if group_by and group_by in data.columns:
            groups = data[group_by].dropna().unique()
            n_groups = len(groups)
            
            if n_groups == 0:
                return ""
            
            # Use 2 columns layout
            n_cols = 2
            n_rows = (n_groups + n_cols - 1) // n_cols
            
            fig = make_subplots(
                rows=n_rows, cols=n_cols,
                subplot_titles=[f"{g}" for g in groups],
                vertical_spacing=0.1,
                horizontal_spacing=0.08
            )
            
            for idx, group in enumerate(groups):
                row = idx // n_cols + 1
                col = idx % n_cols + 1
                group_data = data[data[group_by] == group].copy()
                
                if label_column in group_data.columns:
                    agg_data = group_data.groupby(label_column)[score_column].max().reset_index()
                else:
                    agg_data = group_data[[score_column]].copy()
                    agg_data[label_column] = range(len(agg_data))
                
                scores = agg_data[score_column].tolist()
                labels = agg_data[label_column].tolist()
                
                if len(scores) > 0:
                    mountain_scores, mountain_labels = create_mountain_order(scores, labels)
                    
                    fig.add_trace(
                        go.Bar(
                            x=list(range(len(mountain_scores))),
                            y=mountain_scores,
                            marker_color=CATEGORY10[idx % len(CATEGORY10)],
                            name=str(group),
                            hovertemplate=f'<b>%{{customdata}}</b><br>Score: %{{y:,.0f}}<extra>{group}</extra>',
                            customdata=mountain_labels,
                            width=1.0  # No gap between bars
                        ),
                        row=row, col=col
                    )
                    
                    if show_threshold:
                        fig.add_hline(
                            y=show_threshold, 
                            line_dash="dash", 
                            line_color="red",
                            row=row, col=col
                        )
                    
                    # Clean axis style
                    fig.update_xaxes(
                        showticklabels=False, 
                        showgrid=False, 
                        zeroline=False,
                        showline=False,
                        row=row, col=col
                    )
                    fig.update_yaxes(
                        showgrid=False, 
                        zeroline=False,
                        showline=False,
                        row=row, col=col
                    )
            
            fig.update_layout(
                title=dict(text=title, x=0.5),
                height=250 * n_rows,
                showlegend=False,
                plot_bgcolor='rgba(0,0,0,0)',
                paper_bgcolor='rgba(0,0,0,0)',
                bargap=0
            )
        else:
            if label_column in data.columns:
                agg_data = data.groupby(label_column)[score_column].max().reset_index()
            else:
                agg_data = data[[score_column]].copy()
                agg_data[label_column] = range(len(agg_data))
            
            scores = agg_data[score_column].tolist()
            labels = agg_data[label_column].tolist()
            
            if len(scores) == 0:
                return ""
            
            mountain_scores, mountain_labels = create_mountain_order(scores, labels)
            
            fig = go.Figure()
            
            fig.add_trace(go.Bar(
                x=list(range(len(mountain_scores))),
                y=mountain_scores,
                marker_color=color,
                hovertemplate='<b>%{customdata}</b><br>Score: %{y:,.0f}<extra></extra>',
                customdata=mountain_labels,
                width=1.0
            ))
            
            if show_threshold:
                fig.add_hline(
                    y=show_threshold, 
                    line_dash="dash", 
                    line_color="red",
                    annotation_text=f"Threshold: {show_threshold:,.0f}"
                )
            
            fig.update_layout(
                title=dict(text=title, x=0.5),
                xaxis_title="Neurons/Types (sorted by score, highest in center)",
                yaxis_title="Score",
                xaxis=dict(showticklabels=False, showgrid=False, zeroline=False, showline=False),
                yaxis=dict(showgrid=False, zeroline=False, showline=False),
                height=400,
                plot_bgcolor='rgba(0,0,0,0)',
                paper_bgcolor='rgba(0,0,0,0)',
                bargap=0
            )
        
        fig.write_html(full_path)
        self._vprint(f"   📊 Created labeling distribution: {full_path}")
        return full_path

    def visualize_colabeling_distribution(
        self,
        line_neurons_dict: Dict[str, pd.DataFrame],
        output_path: str,
        min_score: float = 0.0,
        title: str = "Co-Labeling Score Distribution"
    ) -> Tuple[str, str]:
        """
        Visualize labeling distribution for multiple lines (co-labeling analysis).
        
        Creates two multi-panel plots:
        1. By neuron: individual neuron scores for each line
        2. By type: aggregated type scores for each line
        
        Parameters
        ----------
        line_neurons_dict : dict
            Dictionary mapping line names to neurons DataFrames.
        output_path : str
            Directory to save the HTML files.
        min_score : float
            Minimum score threshold to highlight. Default: 0.0
        title : str
            Plot title.
            
        Returns
        -------
        tuple of str
            Paths to the created HTML files (by_neuron, by_type).
        """
        import plotly.graph_objects as go
        from plotly.subplots import make_subplots
        
        os.makedirs(output_path, exist_ok=True)
        
        if not line_neurons_dict:
            return "", ""
        
        # Filter out empty DataFrames - keep all data for saving
        valid_lines_all = {k: v for k, v in line_neurons_dict.items() if not v.empty}
        if not valid_lines_all:
            return "", ""
        
        # Create filtered version for visualization (apply min_score filter if score column exists)
        valid_lines = {}
        for k, df in valid_lines_all.items():
            if min_score > 0 and 'score' in df.columns:
                # Use _passes_min_score column if available (from _calculate_mutual_information)
                if '_passes_min_score' in df.columns:
                    filtered_df = df[df['_passes_min_score']].copy()
                else:
                    filtered_df = df[df['score'] >= min_score].copy()
                if not filtered_df.empty:
                    valid_lines[k] = filtered_df
            else:
                valid_lines[k] = df
        
        # If all data was filtered out, use all data for visualization
        if not valid_lines:
            valid_lines = valid_lines_all
        
        lines = list(valid_lines.keys())
        n_lines = len(lines)
        
        # Category10 palette (same as bokeh.palettes.Category10)
        CATEGORY10 = ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd',
                      '#8c564b', '#e377c2', '#7f7f7f', '#bcbd22', '#17becf']
        
        def create_mountain_order(scores, labels):
            """Rearrange so highest score is in middle, descending to both sides."""
            sorted_pairs = sorted(zip(scores, labels), key=lambda x: x[0], reverse=True)
            n = len(sorted_pairs)
            if n == 0:
                return [], []
            
            result_scores = [0] * n
            result_labels = [''] * n
            
            left = n // 2 - 1 if n % 2 == 0 else n // 2
            right = n // 2 if n % 2 == 0 else n // 2
            
            for i, (score, label) in enumerate(sorted_pairs):
                if i == 0:
                    center = n // 2
                    result_scores[center] = score
                    result_labels[center] = label
                elif i % 2 == 1:
                    result_scores[left] = score
                    result_labels[left] = label
                    left -= 1
                else:
                    result_scores[right] = score
                    result_labels[right] = label
                    right += 1
            
            return result_scores, result_labels
        
        def create_distribution_plot(by_type: bool, filename: str) -> str:
            """Create distribution plot (by neuron or by type).
            
            Each line gets its own line-specific mountain shape distribution.
            
            Key features:
            - c_line: count of items above min_score threshold for each line
            - If c_line < 100, expand to top-100 items (get t_line = score at expansion boundary)
            - X-axis limited by max(c_line) across all INDIVIDUAL lines (not combined)
            - Y-axis: 0.9 * t_line to max_score for each subplot
            - Peak is ALWAYS centered by padding x-axis
            """
            # Use 2 columns layout, with combined as last panel
            n_panels = n_lines + 1  # lines + combined
            n_cols = 2
            n_rows = (n_panels + n_cols - 1) // n_cols
            
            subplot_titles = [f"{line}" for line in lines]
            subplot_titles.append("Combined (avg)" if by_type else "All Neurons")
            
            fig = make_subplots(
                rows=n_rows, cols=n_cols,
                subplot_titles=subplot_titles,
                vertical_spacing=0.08,
                horizontal_spacing=0.06
            )
            
            # Expansion target when items above threshold < 100
            EXPANSION_TARGET = 100
            
            # Pre-calculate all data and per-line statistics
            line_all_data = {}  # Store all data per line (sorted by score desc)
            line_max_scores = {}  # Store max score per line
            line_c_values = {}  # c_line: count of items above threshold
            line_t_values = {}  # t_line: threshold for y-axis (score at expansion boundary)
            line_display_counts = {}  # Number of items to display per line
            
            for line, df in valid_lines.items():
                # Get corresponding unfiltered data from valid_lines_all
                df_all = valid_lines_all.get(line, df)
                
                if by_type and 'type' in df_all.columns:
                    type_scores = df_all.groupby('type')['score'].max()
                    sorted_scores = type_scores.sort_values(ascending=False)
                    line_all_data[line] = sorted_scores
                    line_max_scores[line] = sorted_scores.max() if len(sorted_scores) > 0 else 50000
                    
                    # Calculate c_line (count above threshold)
                    c_line = (sorted_scores >= min_score).sum() if min_score > 0 else len(sorted_scores)
                    line_c_values[line] = c_line
                    
                    # Determine display count and t_line
                    if c_line < EXPANSION_TARGET:
                        # Expand to top-100 (or all if < 100 items available)
                        display_count = min(EXPANSION_TARGET, len(sorted_scores))
                        # t_line is the score at the expansion boundary
                        if display_count > 0:
                            t_line = sorted_scores.iloc[display_count - 1] if display_count <= len(sorted_scores) else sorted_scores.iloc[-1]
                        else:
                            t_line = min_score if min_score > 0 else 0
                    else:
                        # Use only items above threshold
                        display_count = c_line
                        t_line = min_score if min_score > 0 else 0
                    
                    line_display_counts[line] = display_count
                    line_t_values[line] = t_line
                else:
                    sorted_df = df_all.sort_values('score', ascending=False)
                    line_all_data[line] = sorted_df
                    line_max_scores[line] = sorted_df['score'].max() if len(sorted_df) > 0 else 50000
                    
                    # Calculate c_line (count above threshold)
                    c_line = (sorted_df['score'] >= min_score).sum() if min_score > 0 else len(sorted_df)
                    line_c_values[line] = c_line
                    
                    # Determine display count and t_line
                    if c_line < EXPANSION_TARGET:
                        display_count = min(EXPANSION_TARGET, len(sorted_df))
                        if display_count > 0:
                            t_line = sorted_df['score'].iloc[display_count - 1] if display_count <= len(sorted_df) else sorted_df['score'].iloc[-1]
                        else:
                            t_line = min_score if min_score > 0 else 0
                    else:
                        display_count = c_line
                        t_line = min_score if min_score > 0 else 0
                    
                    line_display_counts[line] = display_count
                    line_t_values[line] = t_line
            
            # X-axis range is limited by max(c_line) across individual lines
            # This ensures the combined plot doesn't inflate the x-axis
            max_c_line = max(line_c_values.values()) if line_c_values else EXPANSION_TARGET
            # But we need at least the display count to show expanded items
            max_display_count = max(line_display_counts.values()) if line_display_counts else EXPANSION_TARGET
            global_x_range = max(max_c_line, max_display_count)
            
            # Plot each line with its OWN mountain shape
            for idx, line in enumerate(lines):
                row = idx // n_cols + 1
                col = idx % n_cols + 1
                
                display_count = line_display_counts.get(line, EXPANSION_TARGET)
                t_line = line_t_values.get(line, min_score)
                
                if by_type:
                    # Get type scores for this line (sorted by score desc)
                    type_scores = line_all_data.get(line, pd.Series(dtype=float))
                    
                    # Take display_count items
                    top_types = type_scores.head(display_count)
                    
                    scores = top_types.values.tolist()
                    type_names = top_types.index.tolist()
                    
                    hover_texts = []
                    for type_name, score in zip(type_names, scores):
                        if score >= min_score:
                            hover_texts.append(f"{type_name}<br>Score: {score:,.0f}")
                        else:
                            hover_texts.append(f"{type_name}<br>Score: {score:,.0f}<br>(below threshold)")
                    
                    mountain_scores, _ = create_mountain_order(scores, type_names)
                    _, mountain_hovers = create_mountain_order(scores, hover_texts)
                else:
                    # Get all neurons for this line (sorted by score desc)
                    df_all = line_all_data.get(line, pd.DataFrame())
                    
                    # Take display_count items
                    top_neurons = df_all.head(display_count)
                    
                    scores = top_neurons['score'].tolist()
                    if 'bodyId' in top_neurons.columns:
                        labels = top_neurons['bodyId'].astype(str).tolist()
                    elif 'type' in top_neurons.columns:
                        labels = top_neurons['type'].tolist()
                    else:
                        labels = [f"neuron_{i}" for i in range(len(scores))]
                    
                    hover_texts = []
                    for _, row_data in top_neurons.iterrows():
                        body_id = row_data.get('bodyId', 'Unknown')
                        neuron_type = row_data.get('type', 'Unknown')
                        dataset = row_data.get('dataset', 'Unknown')
                        score = row_data.get('score', 0)
                        suffix = "<br>(below threshold)" if score < min_score else ""
                        hover_texts.append(f"bodyId: {body_id}<br>Type: {neuron_type}<br>Dataset: {dataset}{suffix}")
                    
                    mountain_scores, _ = create_mountain_order(scores, labels)
                    _, mountain_hovers = create_mountain_order(scores, hover_texts)
                
                if len(mountain_scores) > 0:
                    bar_color = CATEGORY10[idx % len(CATEGORY10)]
                    
                    # Center the peak by calculating offset
                    actual_n_items = len(mountain_scores)
                    offset = (global_x_range - actual_n_items) // 2
                    x_positions = [i + offset for i in range(actual_n_items)]
                    
                    fig.add_trace(
                        go.Bar(
                            x=x_positions,
                            y=mountain_scores,
                            marker=dict(
                                color=bar_color,
                                line=dict(width=0)
                            ),
                            name=line,
                            hovertemplate=f'<b>%{{customdata}}</b><extra>{line}</extra>',
                            customdata=mountain_hovers,
                            width=1.0
                        ),
                        row=row, col=col
                    )
                    
                    if min_score > 0:
                        fig.add_hline(
                            y=min_score, 
                            line_dash="dash", 
                            line_color="red",
                            row=row, col=col
                        )
                    
                    # Y-axis: 0.9 * t_line to max_score * 1.1
                    y_min_subplot = t_line * 0.9 if t_line > 0 else 0
                    y_max_subplot = line_max_scores.get(line, 50000) * 1.1
                    
                    fig.update_xaxes(
                        showticklabels=False, 
                        showgrid=False, 
                        zeroline=False,
                        showline=False,
                        range=[-0.5, global_x_range - 0.5],
                        row=row, col=col
                    )
                    fig.update_yaxes(
                        showgrid=False, 
                        zeroline=False,
                        showline=False,
                        range=[y_min_subplot, y_max_subplot],
                        row=row, col=col
                    )
            
            # Combined plot
            combined_idx = n_lines
            combined_row = combined_idx // n_cols + 1
            combined_col = combined_idx % n_cols + 1
            
            if by_type:
                # Collect all types across all lines with their average scores
                all_types_scores = {}  # type -> list of scores
                for line, type_scores in line_all_data.items():
                    if isinstance(type_scores, pd.Series):
                        for type_name, score in type_scores.items():
                            if type_name not in all_types_scores:
                                all_types_scores[type_name] = []
                            all_types_scores[type_name].append(score)
                
                # Calculate average score per type
                type_avg_scores = {t: sum(s) / len(s) for t, s in all_types_scores.items()}
                sorted_types = sorted(type_avg_scores.items(), key=lambda x: x[1], reverse=True)
                
                # Combined plot uses global_x_range (based on max(c_line))
                combined_labels = [t[0] for t in sorted_types[:global_x_range]]
                combined_scores = [t[1] for t in sorted_types[:global_x_range]]
                combined_hovers = [f"{t}<br>Avg: {s:,.0f}" for t, s in sorted_types[:global_x_range]]
                
                mountain_scores, _ = create_mountain_order(combined_scores, combined_labels)
                _, mountain_hovers = create_mountain_order(combined_scores, combined_hovers)
            else:
                # All neurons from all lines
                combined_scores = []
                combined_labels = []
                combined_hovers = []
                for line, df_all in line_all_data.items():
                    if isinstance(df_all, pd.DataFrame):
                        # Only take items up to display_count for this line
                        display_count = line_display_counts.get(line, EXPANSION_TARGET)
                        df_subset = df_all.head(display_count)
                        combined_scores.extend(df_subset['score'].tolist())
                        if 'bodyId' in df_subset.columns:
                            combined_labels.extend(df_subset['bodyId'].astype(str).tolist())
                        else:
                            combined_labels.extend([f"{line}_{i}" for i in range(len(df_subset))])
                        
                        for _, row_data in df_subset.iterrows():
                            body_id = row_data.get('bodyId', 'Unknown')
                            neuron_type = row_data.get('type', 'Unknown')
                            dataset = row_data.get('dataset', 'Unknown')
                            combined_hovers.append(f"bodyId: {body_id}<br>Type: {neuron_type}<br>Dataset: {dataset}<br>Line: {line}")
                
                # Sort and take top items limited by global_x_range
                if combined_scores:
                    sorted_indices = sorted(range(len(combined_scores)), key=lambda i: combined_scores[i], reverse=True)
                    combined_target = min(global_x_range * n_lines, len(combined_scores))
                    combined_scores = [combined_scores[i] for i in sorted_indices[:combined_target]]
                    combined_labels = [combined_labels[i] for i in sorted_indices[:combined_target]]
                    combined_hovers = [combined_hovers[i] for i in sorted_indices[:combined_target]]
                
                mountain_scores, _ = create_mountain_order(combined_scores, combined_labels)
                _, mountain_hovers = create_mountain_order(combined_scores, combined_hovers)
            
            if mountain_scores:
                # Combined plot y-axis: use min_score * 0.9 as base
                combined_max = max(mountain_scores) if mountain_scores else 50000
                y_max_combined = combined_max * 1.1
                y_min_combined = min_score * 0.9 if min_score > 0 else 0
                
                # Center the combined plot
                actual_combined_items = len(mountain_scores)
                # Combined plot x-range should match individual plots (global_x_range)
                offset_combined = (global_x_range - actual_combined_items) // 2 if actual_combined_items < global_x_range else 0
                x_positions_combined = [i + offset_combined for i in range(actual_combined_items)]
                
                fig.add_trace(
                    go.Bar(
                        x=x_positions_combined,
                        y=mountain_scores,
                        marker=dict(
                            color='#17becf',
                            line=dict(width=0)
                        ),
                        name='Combined',
                        hovertemplate='<b>%{customdata}</b><br>Score: %{y:,.0f}<extra>Combined</extra>',
                        customdata=mountain_hovers,
                        width=1.0
                    ),
                    row=combined_row, col=combined_col
                )
                
                if min_score > 0:
                    fig.add_hline(
                        y=min_score, 
                        line_dash="dash", 
                        line_color="red",
                        annotation_text=f"Threshold: {min_score:,.0f}",
                        row=combined_row, col=combined_col
                    )
                
                # Combined plot x-range matches individual line plots
                combined_x_max = max(global_x_range, actual_combined_items)
                
                fig.update_xaxes(
                    showticklabels=False, 
                    showgrid=False, 
                    zeroline=False,
                    showline=False,
                    range=[-0.5, combined_x_max - 0.5],
                    row=combined_row, col=combined_col
                )
                fig.update_yaxes(
                    showgrid=False, 
                    zeroline=False,
                    showline=False,
                    range=[y_min_combined, y_max_combined],
                    row=combined_row, col=combined_col
                )
            
            suffix = "by Type" if by_type else "by Neuron"
            fig.update_layout(
                title=dict(text=f"{title} ({suffix})", x=0.5),
                height=250 * n_rows,
                showlegend=False,
                plot_bgcolor='rgba(0,0,0,0)',
                paper_bgcolor='rgba(0,0,0,0)',
                bargap=0
            )
            
            full_path = os.path.join(output_path, filename)
            fig.write_html(full_path)
            return full_path
        
        # Create both visualizations
        path_by_type = create_distribution_plot(by_type=True, filename='labeling_distribution_by_type.html')
        path_by_neuron = create_distribution_plot(by_type=False, filename='labeling_distribution_by_neuron.html')
        
        # Create stacked mountain-shaped histogram
        path_stacked = self._create_stacked_mountain_plot(
            line_neurons_dict=valid_lines,
            output_path=output_path,
            min_score=min_score,
            title=title,
            colors=CATEGORY10
        )
        
        # Save ALL visualization data as CSV (unfiltered, for reproducibility)
        # Use valid_lines_all which contains all data regardless of min_score filter
        self._save_distribution_data(
            line_neurons_dict=valid_lines_all,
            output_path=output_path
        )
        
        self._vprint(f"   📊 Created labeling distribution (by type): {path_by_type}")
        self._vprint(f"   📊 Created labeling distribution (by neuron): {path_by_neuron}")
        self._vprint(f"   📊 Created labeling distribution (stacked): {path_stacked}")
        
        return path_by_type, path_by_neuron
    
    def _save_distribution_data(
        self,
        line_neurons_dict: Dict[str, pd.DataFrame],
        output_path: str
    ) -> None:
        """
        Save the data used for distribution visualizations as CSV files.
        
        Saves two CSV files:
        1. distribution_data_by_neuron.csv - All neurons with scores per line
        2. distribution_data_by_type.csv - Aggregated type scores per line
        
        Parameters
        ----------
        line_neurons_dict : dict
            Dictionary mapping line names to neurons DataFrames.
        output_path : str
            Directory to save the CSV files.
        """
        # Save by neuron (all data)
        all_neurons = []
        for line, df in line_neurons_dict.items():
            if not df.empty:
                df_copy = df.copy()
                df_copy['source_line'] = line
                all_neurons.append(df_copy)
        
        if all_neurons:
            combined_neurons = pd.concat(all_neurons, ignore_index=True)
            neurons_path = os.path.join(output_path, 'distribution_data_by_neuron.csv')
            combined_neurons.to_csv(neurons_path, index=False)
        
        # Save by type (aggregated max score per type per line)
        type_data = []
        for line, df in line_neurons_dict.items():
            if 'type' in df.columns and not df.empty:
                type_scores = df.groupby('type')['score'].max().reset_index()
                type_scores['source_line'] = line
                # Also get dataset info if available
                if 'dataset' in df.columns:
                    type_datasets = df.groupby('type')['dataset'].first().reset_index()
                    type_scores = type_scores.merge(type_datasets, on='type', how='left')
                type_data.append(type_scores)
        
        if type_data:
            combined_types = pd.concat(type_data, ignore_index=True)
            types_path = os.path.join(output_path, 'distribution_data_by_type.csv')
            combined_types.to_csv(types_path, index=False)
    
    def _create_stacked_mountain_plot(
        self,
        line_neurons_dict: Dict[str, pd.DataFrame],
        output_path: str,
        min_score: float = 0.0,
        title: str = "Co-Labeling Score Distribution",
        colors: List[str] = None
    ) -> str:
        """
        Create a stacked mountain-shaped histogram showing all lines overlaid.
        
        Each line's scores are shown as bars stacked on top of each other,
        with transparency to see overlapping regions.
        
        Parameters
        ----------
        line_neurons_dict : dict
            Dictionary mapping line names to neurons DataFrames.
        output_path : str
            Directory to save the HTML file.
        min_score : float
            Minimum score threshold to highlight.
        title : str
            Plot title.
        colors : list
            List of colors to use for each line.
            
        Returns
        -------
        str
            Path to the created HTML file.
        """
        import plotly.graph_objects as go
        
        if colors is None:
            colors = ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd',
                      '#8c564b', '#e377c2', '#7f7f7f', '#bcbd22', '#17becf']
        
        lines = list(line_neurons_dict.keys())
        n_lines = len(lines)
        
        # Collect all types across all lines with their scores per line
        all_types_data = {}  # type_name -> {line_name: max_score}
        
        for line, df in line_neurons_dict.items():
            if 'type' in df.columns:
                type_scores = df.groupby('type')['score'].max()
                for type_name, score in type_scores.items():
                    if type_name not in all_types_data:
                        all_types_data[type_name] = {}
                    all_types_data[type_name][line] = score
        
        if not all_types_data:
            return ""
        
        # Calculate total score for each type (sum across all lines)
        type_totals = {t: sum(scores.values()) for t, scores in all_types_data.items()}
        
        # Sort types by total score (descending) for mountain shape
        sorted_types = sorted(type_totals.keys(), key=lambda t: type_totals[t], reverse=True)
        
        # Create mountain shape (highest in middle)
        n_types = len(sorted_types)
        mountain_types = [''] * n_types
        left = n_types // 2 - 1 if n_types % 2 == 0 else n_types // 2
        right = n_types // 2 if n_types % 2 == 0 else n_types // 2
        
        for i, type_name in enumerate(sorted_types):
            if i == 0:
                mountain_types[n_types // 2] = type_name
            elif i % 2 == 1:
                mountain_types[left] = type_name
                left -= 1
            else:
                mountain_types[right] = type_name
                right += 1
        
        # Create figure with stacked bars
        fig = go.Figure()
        
        # Add traces for each line (in reverse order so first line is on top visually)
        for idx, line in enumerate(reversed(lines)):
            line_idx = len(lines) - 1 - idx
            scores = []
            hover_texts = []
            
            for type_name in mountain_types:
                score = all_types_data.get(type_name, {}).get(line, 0)
                scores.append(score)
                hover_texts.append(f"{type_name}<br>{line}: {score:,.0f}")
            
            # Convert hex color to rgba with alpha
            hex_color = colors[line_idx % len(colors)]
            # Parse hex to RGB
            r = int(hex_color[1:3], 16)
            g = int(hex_color[3:5], 16)
            b = int(hex_color[5:7], 16)
            rgba_color = f'rgba({r}, {g}, {b}, 0.7)'
            
            fig.add_trace(
                go.Bar(
                    x=list(range(len(mountain_types))),
                    y=scores,
                    name=line,
                    marker=dict(
                        color=rgba_color,
                        line=dict(width=0)  # No edge
                    ),
                    hovertemplate='<b>%{customdata}</b><extra></extra>',
                    customdata=hover_texts,
                    width=1.0
                )
            )
        
        # Add threshold line if specified
        if min_score > 0:
            fig.add_hline(
                y=min_score,
                line_dash="dash",
                line_color="red",
                annotation_text=f"Threshold: {min_score:,.0f}",
                annotation_position="top right"
            )
        
        fig.update_layout(
            title=dict(text=f"{title} (Stacked by Type)", x=0.5),
            barmode='stack',
            height=500,
            showlegend=True,
            legend=dict(
                orientation="h",
                yanchor="bottom",
                y=1.02,
                xanchor="center",
                x=0.5,
                bgcolor='rgba(255,255,255,0.8)',
                bordercolor='rgba(0,0,0,0.3)',
                borderwidth=1
            ),
            plot_bgcolor='rgba(0,0,0,0)',
            paper_bgcolor='rgba(0,0,0,0)',
            bargap=0,
            xaxis=dict(
                showticklabels=False,
                showgrid=False,
                zeroline=False,
                showline=False,
                title="Types (sorted by total score, mountain-shaped)"
            ),
            yaxis=dict(
                showgrid=False,
                zeroline=False,
                showline=False,
                title="Score"
            )
        )
        
        full_path = os.path.join(output_path, 'labeling_distribution_stacked.html')
        fig.write_html(full_path)
        return full_path

    def _calculate_line_specificity(
        self,
        line_stats: pd.DataFrame,
        queried_types: List[str],
        match_type: str = 'cds',
        top_n: int = 100,
        max_lines: int = 100
    ) -> pd.DataFrame:
        """
        Calculate comprehensive specificity metrics for driver lines.
        
        For each line, computes:
        - rank_sum: Sum of ranks for queried neuron types in line_to_neuron results
        - type_proportion: Proportion of queried types among all labeled types (top N)
        - n_queried_types: Number of queried types labeled by this line
        - n_total_types: Total number of distinct types labeled by this line
        - selectivity: 1 / n_total_types (higher = more selective)
        - expression_entropy: Shannon entropy of type distribution (lower = more specific)
        - normalized_entropy: Entropy normalized by max possible (0-1, lower = more specific)
        - weighted_type_proportion: Type proportion weighted by match scores
        - mean_queried_score: Mean match score for queried types
        - specificity_score: Composite score combining all metrics
        
        Parameters
        ----------
        line_stats : pd.DataFrame
            DataFrame with 'line' column from find_lines_batch aggregation.
        queried_types : list of str
            List of queried neuron types (from the original query).
        match_type : str
            Match algorithm for line_to_neuron. Default: 'cds'
        top_n : int
            Number of top matches to consider for type proportion. Default: 100
        max_lines : int
            Maximum number of lines to process (top N by ranking). Default: 100
            
        Returns
        -------
        pd.DataFrame
            Updated line_stats with specificity columns added.
        """
        if line_stats.empty or 'line' not in line_stats.columns:
            return line_stats
        
        if not queried_types:
            self._vprint("  ℹ️  No neuron types to calculate specificity (query was by bodyId)")
            return line_stats
        
        # Normalize queried types for comparison
        queried_types_lower = set(t.lower() for t in queried_types if t)
        
        # Limit to top N lines to reduce API calls
        total_lines = len(line_stats)
        if max_lines and total_lines > max_lines:
            self._vprint(f"\n📊 Calculating specificity for top {max_lines} of {total_lines} lines...")
            lines_to_process = line_stats['line'].head(max_lines).tolist()
        else:
            self._vprint(f"\n📊 Calculating specificity metrics for {total_lines} lines...")
            lines_to_process = line_stats['line'].tolist()
        
        self._vprint(f"   Queried types: {list(queried_types)[:5]}{'...' if len(queried_types) > 5 else ''}")
        self._vprint(f"   ⏱️  Note: Each line requires an API call to fetch neuron matches (may take time)")
        
        # Initialize new columns (for ALL lines, not just processed ones)
        line_stats = line_stats.copy()
        line_stats['rank_sum'] = float('inf')
        line_stats['type_proportion'] = 0.0
        line_stats['n_queried_types'] = 0
        line_stats['n_total_types'] = 0
        line_stats['selectivity'] = 0.0
        # New entropy columns
        line_stats['expression_entropy'] = 0.0
        line_stats['normalized_entropy'] = 0.0
        # New weighted columns
        line_stats['weighted_type_proportion'] = 0.0
        line_stats['mean_queried_score'] = 0.0
        
        # Process each line (with progress bar if available)
        if HAS_TQDM and self.verbose:
            from tqdm import tqdm as tqdm_progress
            iterator = tqdm_progress(
                lines_to_process, 
                desc="   🔬 Analyzing specificity", 
                unit="line",
                bar_format='{desc}: {percentage:3.0f}%|{bar:30}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}, {rate_fmt}]',
                ncols=110,
                position=0
            )
        else:
            iterator = lines_to_process
        
        import time
        start_time = time.time()
        
        # Enable batch mode to suppress individual line_to_neuron messages
        self._batch_mode = True
        
        for idx, line_name in enumerate(iterator):
            # Update progress description with current line
            if HAS_TQDM and self.verbose:
                iterator.set_description(f"   🔬 [{idx+1}/{len(lines_to_process)}] {line_name}")
            try:
                # Check if cached before calling (use same cache key format as line_to_neuron)
                region_key = self.region if self.region else 'all'
                max_imgs_key = self.max_api_images_per_line if self.max_api_images_per_line > 0 else 'all'
                cache_key = f"{line_name}_{match_type}_{region_key}_{max_imgs_key}"
                is_cached = self._load_from_cache('line_to_neuron', cache_key) is not None
                
                if HAS_TQDM and self.verbose:
                    status = "💾" if is_cached else "🌐"
                    iterator.set_postfix_str(status)
                
                # Get neurons labeled by this line (use cached if available)
                neurons_df = self.line_to_neuron(line_name, match_type=match_type)
                
                if neurons_df.empty:
                    continue
                
                # Get top N matches
                neurons_top = neurons_df.head(top_n)
                
                # Get all unique neuron types (across all datasets)
                all_types = neurons_top['type'].dropna().unique().tolist()
                all_types_lower = set(t.lower() for t in all_types if t)
                
                # Count how many queried types are labeled
                matched_types = queried_types_lower & all_types_lower
                n_queried = len(matched_types)
                n_total = len(all_types_lower)
                
                # Calculate rank sum for queried types
                rank_sum = 0
                for idx, row in neurons_top.iterrows():
                    neuron_type = str(row.get('type', '')).lower()
                    if neuron_type in queried_types_lower:
                        # Rank is 1-indexed position
                        rank = neurons_top.index.get_loc(idx) + 1
                        rank_sum += rank
                
                # If no matches, set rank_sum to a high value
                if rank_sum == 0 and n_queried > 0:
                    rank_sum = top_n * n_queried  # Penalty for not appearing in top N
                
                # Calculate type proportion
                type_proportion = n_queried / n_total if n_total > 0 else 0.0
                
                # Calculate selectivity (inverse of diversity)
                selectivity = 1.0 / n_total if n_total > 0 else 0.0
                
                # === NEW: Calculate Expression Entropy ===
                # Count occurrences of each type
                type_counts = neurons_top['type'].fillna('Unknown').value_counts().to_dict()
                expression_entropy = self._calculate_expression_entropy(type_counts)
                
                # Normalize entropy by max possible (log2(n_types))
                max_entropy = np.log2(n_total) if n_total > 1 else 1.0
                normalized_entropy = expression_entropy / max_entropy if max_entropy > 0 else 0.0
                
                # === NEW: Calculate Weighted Specificity ===
                weighted_metrics = self._calculate_weighted_specificity(
                    neurons_top, queried_types_lower, score_column='score'
                )
                
                # Update the row
                mask = line_stats['line'] == line_name
                line_stats.loc[mask, 'rank_sum'] = rank_sum
                line_stats.loc[mask, 'type_proportion'] = type_proportion
                line_stats.loc[mask, 'n_queried_types'] = n_queried
                line_stats.loc[mask, 'n_total_types'] = n_total
                line_stats.loc[mask, 'selectivity'] = selectivity
                line_stats.loc[mask, 'expression_entropy'] = expression_entropy
                line_stats.loc[mask, 'normalized_entropy'] = normalized_entropy
                line_stats.loc[mask, 'weighted_type_proportion'] = weighted_metrics['weighted_type_proportion']
                line_stats.loc[mask, 'mean_queried_score'] = weighted_metrics['mean_queried_score']
                
            except Exception as e:
                if self.verbose:
                    self._vprint(f"   ⚠️ Error calculating specificity for {line_name}: {e}")
        
        # Disable batch mode after processing
        self._batch_mode = False
        
        # Calculate composite specificity score
        # Combines: type_proportion, entropy (inverted), rank, and weighted proportion
        max_rank = line_stats['rank_sum'].replace(float('inf'), 0).max()
        
        if max_rank > 0:
            normalized_rank = line_stats['rank_sum'].replace(float('inf'), max_rank) / max_rank
        else:
            normalized_rank = 0
        
        # Specificity = high type_proportion + low entropy + low rank + high weighted_proportion
        # Formula: 0.3*type_prop + 0.2*(1-norm_entropy) + 0.2*(1-norm_rank) + 0.3*weighted_prop
        line_stats['specificity_score'] = (
            0.30 * line_stats['type_proportion'] +
            0.20 * (1.0 - line_stats['normalized_entropy']) +
            0.20 * (1.0 - normalized_rank) +
            0.30 * line_stats['weighted_type_proportion']
        )
        
        # Replace inf with max + 1 for display
        max_finite = line_stats.loc[line_stats['rank_sum'] != float('inf'), 'rank_sum'].max()
        line_stats['rank_sum'] = line_stats['rank_sum'].replace(float('inf'), max_finite + 1 if pd.notna(max_finite) else 9999)
        
        # Calculate and report timing statistics
        elapsed_time = time.time() - start_time
        lines_processed = min(len(lines_to_process), max_lines) if max_lines else len(lines_to_process)
        avg_time_per_line = elapsed_time / lines_processed if lines_processed > 0 else 0
        
        self._vprint(f"   ✓ Specificity calculated for {len(line_stats)} lines")
        self._vprint(f"   ⏱️  Processing time: {elapsed_time:.1f}s total, ~{avg_time_per_line:.2f}s per line")
        
        return line_stats

    def _init_client(self):
        """Initialize the NeuronBridge client with retry logic."""
        import time
        
        max_retries = 3
        retry_delays = [2, 5, 10]  # Seconds to wait before each retry
        
        self._vprint("🔌 Initializing NeuronBridge client...")
        
        for attempt in range(max_retries):
            try:
                # Clear any previous client state before retry
                # Use getattr to safely check if _client exists and is not None
                if getattr(self, '_client', None) is not None:
                    try:
                        del self._client
                    except AttributeError:
                        pass
                    import gc
                    gc.collect()
                
                self._client = NBClient()
                # Fix S3 URL version mismatch if needed
                self._fix_store_prefixes()
                self._vprint("  ✓ Client initialized successfully")
                return  # Success
                
            except Exception as e:
                error_msg = str(e)
                if attempt < max_retries - 1:
                    wait_time = retry_delays[attempt]
                    self._vprint(f"  ⚠️  Initialization failed (attempt {attempt + 1}/{max_retries}): {error_msg}")
                    self._vprint(f"  ⏳ Waiting {wait_time}s before retry...")
                    time.sleep(wait_time)
                else:
                    raise RuntimeError(
                        f"Failed to initialize NeuronBridge client after {max_retries} attempts: {error_msg}\n"
                        f"Please check your network connection and try again later.\n"
                        f"See TROUBLESHOOTING.md for common solutions."
                    )
    
    def _fix_store_prefixes(self):
        """
        Fix S3 URL version mismatch in the NeuronBridge client.
        
        The config may reference v3_8_0 but data might exist at v3_8_1.
        Try production bucket first, as it's more reliable.
        """
        if hasattr(self._client, 'config') and hasattr(self._client.config, 'stores'):
            stores = self._client.config.stores
            # stores is a dictionary of DataStore objects
            if isinstance(stores, dict):
                for store_key, store in stores.items():
                    if hasattr(store, 'prefixes') and store.prefixes:
                        # prefixes is a dict like {'CDSResults': 'url', ...}
                        if isinstance(store.prefixes, dict):
                            for key, url in store.prefixes.items():
                                if isinstance(url, str):
                                    new_url = url
                                    # Keep production bucket but fix version
                                    if 'v3_8_0' in new_url:
                                        new_url = new_url.replace('v3_8_0', 'v3_8_1')
                                    # Switch dev to prod (prod is more reliable)
                                    if 'data-dev' in new_url:
                                        new_url = new_url.replace('data-dev', 'data-prod')
                                    store.prefixes[key] = new_url
    
    def _parse_library_name(self, library_name: str) -> Tuple[str, str]:
        """
        Parse a NeuronBridge library name into base name and version.
        
        Parameters
        ----------
        library_name : str
            Library name like 'FlyEM_Hemibrain_v1.2.1' or 'FlyEM_MANC_v1.2.1'
            
        Returns
        -------
        Tuple[str, str]
            (base_name, version) e.g. ('FlyEM_Hemibrain', '1.2.1') or ('FlyEM_MANC', '1.2.1')
        """
        if not library_name:
            return ('', '')
        
        # Split on '_v' to get base and version
        if '_v' in library_name:
            parts = library_name.rsplit('_v', 1)
            return (parts[0], parts[1] if len(parts) > 1 else '')
        return (library_name, '')
    
    def _get_dataset_from_library(self, library_name: str) -> Optional[str]:
        """
        Get local dataset folder name from NeuronBridge library name.
        
        Parameters
        ----------
        library_name : str
            NeuronBridge library name (e.g., 'FlyEM_Hemibrain_v1.2.1')
            
        Returns
        -------
        str or None
            Local dataset folder name, or None if not found.
        """
        if not library_name:
            return None
        
        # Direct match
        if library_name in LIBRARY_TO_DATASET:
            return LIBRARY_TO_DATASET[library_name]
        
        # Try to parse and match by base name with different versions
        base_name, version = self._parse_library_name(library_name)
        if base_name and version:
            # Look for any existing mapping with the same base name
            for lib_pattern, dataset in LIBRARY_TO_DATASET.items():
                pattern_base, _ = self._parse_library_name(lib_pattern)
                if pattern_base == base_name:
                    # Found a matching base - try to construct the dataset folder
                    # Convert version to folder format (1.2.1 -> v1_2_1)
                    version_folder = 'v' + version.replace('.', '_')
                    # Extract dataset base name from the mapped folder
                    dataset_base = dataset.rsplit('_v', 1)[0] if '_v' in dataset else dataset.replace('_', '-')
                    new_dataset_folder = f"{dataset_base}_{version_folder}"
                    
                    # Print warning about unmapped library
                    self._vprint(f"  ⚠️  UNMAPPED LIBRARY: '{library_name}'")
                    self._vprint(f"      Attempting auto-mapping to: {new_dataset_folder}")
                    self._vprint(f"      (Based on similar library: {lib_pattern} -> {dataset})")
                    self._vprint(f"      Consider adding this to LIBRARY_TO_DATASET in neuronbridge_finder.py")
                    
                    return new_dataset_folder
        
        # Partial match (try prefixes) - legacy fallback
        for lib_pattern, dataset in LIBRARY_TO_DATASET.items():
            if library_name.startswith(lib_pattern.split('_v')[0]):
                self._vprint(f"  ⚠️  PARTIAL MATCH: '{library_name}' -> {dataset}")
                self._vprint(f"      Consider adding explicit mapping to LIBRARY_TO_DATASET")
                return dataset
        
        # No match found - print explicit warning
        self._vprint(f"  ⚠️  UNKNOWN LIBRARY: '{library_name}'")
        self._vprint(f"      This library is not mapped to any local dataset.")
        self._vprint(f"      To add support, update LIBRARY_TO_DATASET and LIBRARY_TO_DATASET_NAME")
        self._vprint(f"      in src/neuronbridge_finder.py with the appropriate mappings.")
        
        return None
    
    def _get_dataset_name_from_library(self, library_name: str) -> Optional[str]:
        """
        Get human-readable dataset name from NeuronBridge library name.
        
        Parameters
        ----------
        library_name : str
            NeuronBridge library name (e.g., 'FlyEM_Male_CNS_Brain_v0.9')
            
        Returns
        -------
        str or None
            Human-readable dataset name (e.g., 'male-cns:v0.9'), or None if not found.
        """
        if not library_name:
            return None
        
        # Direct match
        if library_name in LIBRARY_TO_DATASET_NAME:
            return LIBRARY_TO_DATASET_NAME[library_name]
        
        # Try to parse and match by base name with different versions
        base_name, version = self._parse_library_name(library_name)
        if base_name and version:
            # Look for any existing mapping with the same base name
            for lib_pattern, dataset_name in LIBRARY_TO_DATASET_NAME.items():
                pattern_base, _ = self._parse_library_name(lib_pattern)
                if pattern_base == base_name:
                    # Found a matching base - construct the dataset name
                    # Extract dataset base from the mapped name (e.g., 'hemibrain' from 'hemibrain:v1.2.1')
                    if ':v' in dataset_name:
                        dataset_base = dataset_name.split(':v')[0]
                        new_dataset_name = f"{dataset_base}:v{version}"
                        return new_dataset_name
                    elif '_v' in dataset_name:
                        # Handle folder-style names
                        dataset_base = dataset_name.rsplit('_v', 1)[0]
                        return f"{dataset_base}:v{version}"
        
        # Partial match (try prefixes) - legacy fallback
        for lib_pattern, dataset_name in LIBRARY_TO_DATASET_NAME.items():
            if library_name.startswith(lib_pattern.split('_v')[0]):
                return dataset_name
        
        return None
    
    def _ensure_datasets_loaded(self, datasets: List[str]) -> List[str]:
        """
        Ensure all required datasets are loaded, pulling from NeuPrint if needed.
        
        This method identifies unique datasets from NeuronBridge results and
        ensures each one is available locally (loading from cache or pulling
        from NeuPrint).
        
        Parameters
        ----------
        datasets : list of str
            List of dataset names (e.g., ['hemibrain:v1.2.1', 'male-cns:v0.9', 'manc:v1.2.1'])
            
        Returns
        -------
        list of str
            List of successfully loaded dataset folder names.
        """
        loaded = []
        
        for dataset in datasets:
            # Convert dataset name to folder format
            dataset_folder = dataset.replace(':', '_').replace('.', '_')
            
            # Check if already loaded in cache
            if dataset_folder in self._neuron_dfs:
                loaded.append(dataset_folder)
                continue
            
            # Try to load from local file
            if self.datasets_path:
                neuron_df_path = os.path.join(
                    self.datasets_path, dataset_folder, f"{dataset_folder}_allneurons_neuron_df.csv"
                )
                
                if os.path.exists(neuron_df_path):
                    try:
                        df = pd.read_csv(neuron_df_path, dtype={'bodyId': str}, low_memory=False)
                        self._neuron_dfs[dataset_folder] = df
                        self._vprint(f"   ✓ Loaded {len(df):,} neurons from {dataset_folder}")
                        loaded.append(dataset_folder)
                        continue
                    except Exception as e:
                        self._vprint(f"   ⚠️ Could not load {dataset_folder}: {e}")
            
            # Try to pull from NeuPrint
            self._vprint(f"   ⏳ {dataset}: not found locally, attempting to pull from NeuPrint...")
            pulled_df = self._pull_and_load_dataset(dataset_folder)
            
            if pulled_df is not None and not pulled_df.empty:
                self._vprint(f"   ✓ Pulled {len(pulled_df):,} neurons for {dataset}")
                loaded.append(dataset_folder)
            else:
                # Try fallback via FindNeuronConnection
                self._vprint(f"   ⏳ Trying FindNeuronConnection fallback for {dataset}...")
                fnc_df = self._fetch_neuron_df_via_fnc(dataset_folder)
                if fnc_df is not None and not fnc_df.empty:
                    self._vprint(f"   ✓ Loaded {len(fnc_df):,} neurons via FNC for {dataset}")
                    loaded.append(dataset_folder)
                else:
                    self._vprint(f"   ⚠️ Could not load dataset: {dataset}")
                    self._vprint(f"      This dataset may not be available in NeuPrint.")
                    self._vprint(f"      Neurons from this dataset will be skipped in visualization.")
        
        return loaded
    
    def _load_neuron_df_for_dataset(self, dataset: str) -> Optional[pd.DataFrame]:
        """
        Load neuron DataFrame for a specific dataset.
        
        Parameters
        ----------
        dataset : str
            Dataset folder name (e.g., 'hemibrain_v1_2_1')
            
        Returns
        -------
        pd.DataFrame or None
            Loaded DataFrame, or None if not found.
        """
        # Check cache
        if dataset in self._neuron_dfs:
            return self._neuron_dfs[dataset]
        
        # Try to load from datasets folder
        if self.datasets_path:
            neuron_df_path = os.path.join(
                self.datasets_path, dataset, f"{dataset}_allneurons_neuron_df.csv"
            )
            
            if os.path.exists(neuron_df_path):
                try:
                    df = pd.read_csv(neuron_df_path, dtype={'bodyId': str}, low_memory=False)
                    self._neuron_dfs[dataset] = df
                    # Only print loading message if not suppressing
                    if not self._suppress_loading_msgs:
                        self._vprint(f"  ✓ Loaded {len(df):,} neurons from {dataset}")
                    return df
                except Exception as e:
                    self._vprint(f"  ⚠️ Could not load neuron data for {dataset}: {e}")
            else:
                # Try to pull dataset using statvis.pull_dataset
                return self._pull_and_load_dataset(dataset)
        
        # Try fallback to FindNeuronConnection
        return self._fetch_neuron_df_via_fnc(dataset)
    
    def _pull_and_load_dataset(self, dataset: str) -> Optional[pd.DataFrame]:
        """
        Pull dataset from NeuPrint and load neuron DataFrame.
        
        Uses statvis.pull_dataset() to fetch data when local files don't exist.
        
        Parameters
        ----------
        dataset : str
            Dataset folder name (e.g., 'hemibrain_v1_2_1')
            
        Returns
        -------
        pd.DataFrame or None
            Loaded DataFrame, or None if failed.
        """
        try:
            # Try relative import first (when imported as package)
            try:
                from .statvis import pull_dataset
            except ImportError:
                # Fall back to absolute import (when running from scripts)
                from statvis import pull_dataset
            
            # Convert folder name to dataset format (hemibrain_v1_2_1 -> hemibrain:v1.2.1)
            dataset_name = self._folder_to_dataset_name(dataset)
            
            self._vprint(f"  ⏳ Dataset not found locally, attempting to pull {dataset_name}...")
            
            # Initialize NeuPrint client for this dataset
            try:
                from neuprint import Client, set_default_client, default_client
                
                # Check if there's an existing default client for the same dataset
                try:
                    existing = default_client()
                    if existing.dataset != dataset_name:
                        # Different dataset, need new client
                        existing = None
                except (RuntimeError, AttributeError):
                    existing = None
                
                if existing is None:
                    # Get token using TokenManager
                    token = self.neuprint_token
                    try:
                        from .utils.token_manager import token_manager
                        token = token_manager.get_token('NEUPRINT_TOKEN', token)
                    except ImportError:
                        try:
                            from src.utils.token_manager import token_manager
                            token = token_manager.get_token('NEUPRINT_TOKEN', token)
                        except ImportError:
                            # Fallback to env vars if TokenManager missing
                            if not token:
                                token = os.environ.get('NEUPRINT_TOKEN', os.environ.get('NEUPRINT_APPLICATION_CREDENTIALS', ''))
                    
                    if not token:
                        self._vprint(f"")
                        self._vprint(f"  ⚠️  No NeuPrint token found - cannot pull dataset.")
                        self._vprint(f"  ╔══════════════════════════════════════════════════════════════╗")
                        self._vprint(f"  ║  To enable dataset pulling, provide your NeuPrint token:    ║")
                        self._vprint(f"  ║                                                              ║")
                        self._vprint(f"  ║  Option 1: Pass directly to NeuronBridgeFinder:             ║")
                        self._vprint(f"  ║    nbf = NeuronBridgeFinder(neuprint_token='YOUR_TOKEN')    ║")
                        self._vprint(f"  ║                                                              ║")
                        self._vprint(f"  ║  Option 2: Set environment variable:                        ║")
                        self._vprint(f"  ║    export NEUPRINT_TOKEN='YOUR_TOKEN'                       ║")
                        self._vprint(f"  ║                                                              ║")
                        self._vprint(f"  ║  Get your token at:                                         ║")
                        self._vprint(f"  ║  👉 https://neuprint.janelia.org/account                    ║")
                        self._vprint(f"  ╚══════════════════════════════════════════════════════════════╝")
                        self._vprint(f"")
                        self._vprint(f"  ℹ️  Skipping local dataset features (type lookups, specificity metrics).")
                        return None
                    
                    # Parse server URL (remove https:// if present for neuprint client)
                    server = self.neuprint_server
                    if server.startswith('https://'):
                        server = server[8:]
                    if server.startswith('http://'):
                        server = server[7:]
                    
                    client = Client(server, dataset_name, token)
                    set_default_client(client)
                    self._vprint(f"  ✓ Connected to NeuPrint ({dataset_name})")
                    
            except ImportError:
                self._vprint(f"  ⚠️ neuprint package not available. Install with: pip install neuprint-python")
                return None
            except Exception as e:
                self._vprint(f"  ⚠️ Could not connect to NeuPrint: {e}")
                return None
            
            # Create dataset directory
            dataset_dir = os.path.join(self.datasets_path, dataset)
            os.makedirs(dataset_dir, exist_ok=True)
            
            save_path = os.path.join(dataset_dir, f"{dataset}_allneurons")
            
            try:
                pull_dataset(dataset_name, save_path=save_path, omitNoneType=False)
                
                # Load the pulled data
                neuron_df_path = save_path + '_neuron_df.csv'
                if os.path.exists(neuron_df_path):
                    df = pd.read_csv(neuron_df_path, dtype={'bodyId': str}, low_memory=False)
                    self._neuron_dfs[dataset] = df
                    self._vprint(f"  ✓ Pulled and loaded {len(df):,} neurons from {dataset}")
                    return df
            except Exception as e:
                self._vprint(f"  ⚠️ Could not pull dataset {dataset_name}: {e}")
                
            return None
            
        except ImportError:
            self._vprint(f"  ⚠️ statvis module not available for pulling dataset")
            return None
        except Exception as e:
            self._vprint(f"  ⚠️ Error pulling dataset: {e}")
            return None
    
    def _fetch_neuron_df_via_fnc(self, dataset: str) -> Optional[pd.DataFrame]:
        """
        Try to fetch neuron DataFrame via FindNeuronConnection.
        
        This is a fallback when local dataset files don't exist.
        
        Parameters
        ----------
        dataset : str
            Dataset name
            
        Returns
        -------
        pd.DataFrame or None
            Fetched DataFrame, or None if failed.
        """
        try:
            from scripts.FindNeuronConnection import FindNeuronConnection
            
            self._vprint(f"  ⏳ Fetching neuron data for {dataset} via FindNeuronConnection...")
            
            fnc = FindNeuronConnection(dataset=dataset)
            
            # Check if we can get neuron index
            cache_dir = os.path.join(self.datasets_path, '..', 'cache', dataset)
            csv_path = os.path.join(cache_dir, f'{dataset}_allneurons_neuron_df.csv')
            
            if os.path.exists(csv_path):
                df = pd.read_csv(csv_path, dtype={'bodyId': str}, low_memory=False)
                self._neuron_dfs[dataset] = df
                self._vprint(f"  ✓ Loaded {len(df):,} neurons via FindNeuronConnection cache")
                return df
            
            return None
            
        except ImportError:
            self._vprint(f"  ⚠️ FindNeuronConnection not available")
            return None
        except Exception as e:
            self._vprint(f"  ⚠️ Could not fetch via FindNeuronConnection: {e}")
            return None
    
    def _enrich_match_with_dataset_info(self, match: Dict[str, Any], em_image) -> Dict[str, Any]:
        """
        Enrich a match dictionary with dataset-specific information.
        
        Parameters
        ----------
        match : dict
            Match dictionary to enrich
        em_image : EMImage
            NeuronBridge EM image object
            
        Returns
        -------
        dict
            Enriched match dictionary
        """
        # Get library name and dataset
        library_name = getattr(em_image, 'libraryName', '')
        dataset_folder = self._get_dataset_from_library(library_name)
        dataset_name = self._get_dataset_name_from_library(library_name)
        
        # Add dataset info to match (use human-readable name, folder for internal lookup)
        match['dataset'] = dataset_name or 'unknown'
        match['dataset_folder'] = dataset_folder or 'unknown'
        match['library'] = library_name
        
        # Get neuron type and instance from EM image if available
        match['neuronType'] = getattr(em_image, 'neuronType', '') or ''
        match['neuronInstance'] = getattr(em_image, 'neuronInstance', '') or ''
        
        # Try to enrich from local dataset
        if dataset_folder:
            neuron_df = self._load_neuron_df_for_dataset(dataset_folder)
            if neuron_df is not None and not neuron_df.empty:
                body_id = str(match.get('bodyId', ''))
                neuron_row = neuron_df[neuron_df['bodyId'] == body_id]
                if not neuron_row.empty:
                    row = neuron_row.iloc[0]
                    # Use local data, fallback to NB data
                    match['type'] = row.get('type', '') or match['neuronType']
                    match['instance'] = row.get('instance', '') or match['neuronInstance']
                    match['status'] = row.get('status', '')
                else:
                    # Use NeuronBridge data
                    match['type'] = match['neuronType']
                    match['instance'] = match['neuronInstance']
                    match['status'] = ''
            else:
                match['type'] = match['neuronType']
                match['instance'] = match['neuronInstance']
                match['status'] = ''
        else:
            match['type'] = match['neuronType']
            match['instance'] = match['neuronInstance']
            match['status'] = ''
        
        return match
    
    def _get_cache_path(self, cache_type: str, identifier: str) -> str:
        """Get the cache file path for a given type and identifier."""
        safe_id = str(identifier).replace('/', '_').replace(':', '_')
        return os.path.join(self.cache_folder, f"{cache_type}_{safe_id}.csv")
    
    def _load_from_cache(self, cache_type: str, identifier: str) -> Optional[pd.DataFrame]:
        """Load cached results if available."""
        if not self.use_cache:
            return None
        
        cache_path = self._get_cache_path(cache_type, identifier)
        if os.path.exists(cache_path):
            try:
                df = pd.read_csv(cache_path)
                # Only print cache loads when not in batch mode (verbose individual loads suppressed)
                if not self._batch_mode:
                    self._vprint(f"  ⏩ Loaded from cache: {cache_path}")
                return df
            except Exception:
                return None
        return None
    
    def _save_to_cache(self, cache_type: str, identifier: str, df: pd.DataFrame):
        """Save results to cache."""
        if not self.use_cache or df.empty:
            return
        
        cache_path = self._get_cache_path(cache_type, identifier)
        try:
            df.to_csv(cache_path, index=False)
            # Only print cache saves when not in batch mode
            if not self._batch_mode:
                self._vprint(f"  💾 Saved to cache: {cache_path}")
        except Exception as e:
            warnings.warn(f"Failed to save cache: {e}")
    
    # =========================================================================
    # Image-based Cache System (indexed by image_id)
    # =========================================================================
    # This cache system stores API results per LM image (lm_sample) and match_type
    # allowing reuse regardless of region, top_n, or max_api_images_per_line settings
    #
    # Line-to-Image Mapping:
    # - Tracks which images belong to which line (for each region)
    # - Records cached match types per image
    # - Does NOT mark lines as "complete" since server images can update
    # - Always checks online for new images, uses cache for existing ones
    
    def _get_image_cache_dir(self) -> str:
        """Get the directory for image-based cache."""
        cache_dir = os.path.join(self.cache_folder, 'image_cache')
        os.makedirs(cache_dir, exist_ok=True)
        return cache_dir
    
    def _get_image_cache_path(self, image_id: str, match_type: str) -> str:
        """Get cache file path for a specific LM image and match type."""
        safe_id = str(image_id).replace('/', '_').replace(':', '_')
        cache_dir = self._get_image_cache_dir()
        return os.path.join(cache_dir, f"{match_type}_{safe_id}.csv")
    
    def _get_line_mapping_path(self) -> str:
        """Get the path to the line-image mapping file."""
        cache_dir = self._get_image_cache_dir()
        return os.path.join(cache_dir, 'line_image_mapping.json')
    
    def _load_line_mapping(self) -> Dict[str, Any]:
        """
        Load the line-to-image mapping file.
        
        Structure:
        {
            "lines": {
                "VT037867": {
                    "Brain": {
                        "image_ids": ["123", "456", ...],
                        "last_checked": "2024-12-29T10:30:00"
                    },
                    "VNC": {...}
                }
            },
            "images": {
                "123": {
                    "line": "VT037867",
                    "cached_types": ["cds", "pppm"]
                }
            }
        }
        """
        mapping_path = self._get_line_mapping_path()
        if os.path.exists(mapping_path):
            try:
                with open(mapping_path, 'r') as f:
                    return json.load(f)
            except Exception:
                pass
        return {"lines": {}, "images": {}}
    
    def _save_line_mapping(self, mapping: Dict[str, Any]):
        """Save the line-to-image mapping file."""
        mapping_path = self._get_line_mapping_path()
        try:
            with open(mapping_path, 'w') as f:
                json.dump(mapping, f, indent=2)
        except Exception as e:
            warnings.warn(f"Failed to save line mapping: {e}")
    
    def _update_line_mapping(
        self, 
        line_name: str, 
        region: str, 
        image_ids: List[str],
        image_id: Optional[str] = None,
        match_type: Optional[str] = None
    ):
        """
        Update the line-image mapping with new information.
        
        Parameters
        ----------
        line_name : str
            The driver line name
        region : str
            The anatomical region (Brain, VNC, etc.)
        image_ids : list
            List of all image IDs for this line/region (from current API check)
        image_id : str, optional
            Specific image ID to update cached_types for
        match_type : str, optional
            Match type that was just cached for image_id
        """
        mapping = self._load_line_mapping()
        
        # Update lines section
        if line_name not in mapping["lines"]:
            mapping["lines"][line_name] = {}
        
        from datetime import datetime
        mapping["lines"][line_name][region] = {
            "image_ids": image_ids,
            "last_checked": datetime.now().isoformat(),
            "image_count": len(image_ids)
        }
        
        # Update images section
        for img_id in image_ids:
            if img_id not in mapping["images"]:
                mapping["images"][img_id] = {
                    "line": line_name,
                    "cached_types": []
                }
            else:
                # Update line info (in case image is shared)
                mapping["images"][img_id]["line"] = line_name
        
        # Update cached_types for specific image
        if image_id and match_type:
            if image_id not in mapping["images"]:
                mapping["images"][image_id] = {
                    "line": line_name,
                    "cached_types": []
                }
            if match_type not in mapping["images"][image_id]["cached_types"]:
                mapping["images"][image_id]["cached_types"].append(match_type)
        
        self._save_line_mapping(mapping)
    
    def _get_cached_types_for_image(self, image_id: str) -> List[str]:
        """Get the list of cached match types for an image."""
        mapping = self._load_line_mapping()
        if image_id in mapping.get("images", {}):
            cached = mapping["images"][image_id].get("cached_types", [])
            if cached:
                return cached
        
        # Fallback: check actual cache files
        cached_types = []
        for mt in ['cds', 'pppm', 'both']:
            cache_path = self._get_image_cache_path(image_id, mt)
            if os.path.exists(cache_path):
                cached_types.append(mt)
        return cached_types
    
    def sync_mapping_from_cache_files(self) -> Dict[str, int]:
        """
        Sync the line-image mapping with existing cache files.
        
        Scans all cache files in image_cache directory and updates the mapping
        with the cached types for each image.
        
        Returns
        -------
        dict
            Statistics: {'images_scanned', 'types_updated'}
        """
        stats = {'images_scanned': 0, 'types_updated': 0}
        
        cache_dir = self._get_image_cache_dir()
        mapping = self._load_line_mapping()
        
        # Scan all cache files
        for filename in os.listdir(cache_dir):
            if not filename.endswith('.csv'):
                continue
            
            # Parse filename: {match_type}_{image_id}.csv
            parts = filename[:-4].split('_', 1)  # Remove .csv and split
            if len(parts) != 2:
                continue
            
            match_type, image_id = parts
            if match_type not in ['cds', 'pppm', 'both']:
                continue
            
            stats['images_scanned'] += 1
            
            # Update mapping
            if image_id not in mapping.get("images", {}):
                mapping["images"][image_id] = {
                    "line": "",  # Unknown line
                    "cached_types": []
                }
            
            if match_type not in mapping["images"][image_id].get("cached_types", []):
                if "cached_types" not in mapping["images"][image_id]:
                    mapping["images"][image_id]["cached_types"] = []
                mapping["images"][image_id]["cached_types"].append(match_type)
                stats['types_updated'] += 1
        
        self._save_line_mapping(mapping)
        self._vprint(f"✓ Synced mapping: {stats['images_scanned']} files, {stats['types_updated']} types updated")
        return stats
    
    def _load_image_cache(self, image_id: str, match_type: str) -> Optional[pd.DataFrame]:
        """
        Load cached matches for a specific LM image.
        
        Parameters
        ----------
        image_id : str
            The LM image ID (lm_sample)
        match_type : str
            Match algorithm: 'cds' or 'pppm'
            
        Returns
        -------
        pd.DataFrame or None
            Cached matches DataFrame, or None if not found
        """
        if not self.use_cache:
            return None
        
        cache_path = self._get_image_cache_path(image_id, match_type)
        if os.path.exists(cache_path):
            try:
                return pd.read_csv(cache_path)
            except Exception:
                return None
        return None
    
    def _save_image_cache(
        self, 
        image_id: str, 
        match_type: str, 
        matches: List[Dict[str, Any]],
        line_name: Optional[str] = None
    ):
        """
        Save matches for a specific LM image to cache.
        
        Parameters
        ----------
        image_id : str
            The LM image ID (lm_sample)
        match_type : str
            Match algorithm: 'cds', 'pppm', or 'both'
        matches : list
            List of match dictionaries
        line_name : str, optional
            The line name for updating the mapping
        """
        if not self.use_cache or not matches:
            return
        
        cache_path = self._get_image_cache_path(image_id, match_type)
        try:
            df = pd.DataFrame(matches)
            df.to_csv(cache_path, index=False)
            
            # Update mapping with cached type info
            if line_name:
                mapping = self._load_line_mapping()
                if image_id not in mapping.get("images", {}):
                    mapping["images"][image_id] = {
                        "line": line_name,
                        "cached_types": []
                    }
                if match_type not in mapping["images"][image_id].get("cached_types", []):
                    mapping["images"][image_id]["cached_types"].append(match_type)
                self._save_line_mapping(mapping)
        except Exception as e:
            warnings.warn(f"Failed to save image cache: {e}")
    
    def _fetch_matches_from_api(
        self,
        lm_image,
        match_type: str
    ) -> List[Dict[str, Any]]:
        """
        Fetch matches from API for a single match type (cds or pppm).
        
        Parameters
        ----------
        lm_image : LMImage
            The LM image object
        match_type : str
            Match algorithm: 'cds' or 'pppm' (not 'both')
            
        Returns
        -------
        list
            List of match dictionaries
        """
        image_id = getattr(lm_image, 'id', '')
        matches = []
        
        try:
            if match_type == 'cds':
                api_matches = self._retry_with_backoff(
                    self._client.get_cds_matches,
                    lm_image,
                    max_retries=3,
                    initial_delay=1.0
                )
            else:  # pppm
                api_matches = self._retry_with_backoff(
                    self._client.get_ppp_matches,
                    lm_image,
                    max_retries=3,
                    initial_delay=1.0
                )
            
            for match in api_matches:
                if hasattr(match, 'image') and hasattr(match.image, 'type'):
                    if match.image.type == 'EMImage':
                        body_id = self._extract_body_id(match.image)
                        match_dict = {
                            'bodyId': body_id,
                            'score': getattr(match, 'normalizedScore', 0),
                            'image_id': getattr(match.image, 'id', ''),
                            'lm_sample': str(image_id),
                            'match_type': match_type
                        }
                        # Enrich with dataset info
                        match_dict = self._enrich_match_with_dataset_info(match_dict, match.image)
                        matches.append(match_dict)
                        
        except Exception:
            pass  # Return empty matches on error
        
        return matches
    
    def _get_image_matches_cached(
        self, 
        lm_image, 
        match_type: str,
        line_name: Optional[str] = None
    ) -> Tuple[List[Dict[str, Any]], bool, bool]:
        """
        Get matches for an LM image, using cache if available.
        
        Handles 'both' match_type specially:
        - If 'both' cache exists, use it
        - If only 'cds' cached, use it and fetch 'pppm', then combine and save as 'both'
        - If only 'pppm' cached, use it and fetch 'cds', then combine and save as 'both'
        - If neither cached, fetch both and save as 'both'
        
        Parameters
        ----------
        lm_image : LMImage
            The LM image object
        match_type : str
            Match algorithm: 'cds', 'pppm', or 'both'
        line_name : str, optional
            Line name for mapping updates
            
        Returns
        -------
        tuple
            (matches_list, from_cache, partial_cache) 
            - matches_list: list of match dicts
            - from_cache: True if ALL data came from cache
            - partial_cache: True if SOME data came from cache (for 'both')
        """
        image_id = getattr(lm_image, 'id', '')
        if not image_id:
            return [], False, False
        
        image_id_str = str(image_id)
        
        # Handle simple cases (cds or pppm)
        if match_type in ['cds', 'pppm']:
            cached_df = self._load_image_cache(image_id_str, match_type)
            if cached_df is not None and not cached_df.empty:
                return cached_df.to_dict('records'), True, False
            
            # Fetch from API
            matches = self._fetch_matches_from_api(lm_image, match_type)
            if matches:
                self._save_image_cache(image_id_str, match_type, matches, line_name)
            return matches, False, False
        
        # Handle 'both' match_type with upgrade logic
        # Check if 'both' cache already exists
        both_cache = self._load_image_cache(image_id_str, 'both')
        if both_cache is not None and not both_cache.empty:
            return both_cache.to_dict('records'), True, False
        
        # Check existing cds and pppm caches
        cds_cache = self._load_image_cache(image_id_str, 'cds')
        pppm_cache = self._load_image_cache(image_id_str, 'pppm')
        
        cds_matches = []
        pppm_matches = []
        cds_from_cache = False
        pppm_from_cache = False
        
        # Get CDS matches (from cache or API)
        if cds_cache is not None and not cds_cache.empty:
            cds_matches = cds_cache.to_dict('records')
            cds_from_cache = True
        else:
            cds_matches = self._fetch_matches_from_api(lm_image, 'cds')
            if cds_matches:
                self._save_image_cache(image_id_str, 'cds', cds_matches, line_name)
        
        # Get PPPM matches (from cache or API)
        if pppm_cache is not None and not pppm_cache.empty:
            pppm_matches = pppm_cache.to_dict('records')
            pppm_from_cache = True
        else:
            pppm_matches = self._fetch_matches_from_api(lm_image, 'pppm')
            if pppm_matches:
                self._save_image_cache(image_id_str, 'pppm', pppm_matches, line_name)
        
        # Combine CDS and PPPM matches
        all_matches = cds_matches + pppm_matches
        
        # Save as 'both' cache for future use
        if all_matches:
            self._save_image_cache(image_id_str, 'both', all_matches, line_name)
        
        # Determine cache status
        all_from_cache = cds_from_cache and pppm_from_cache
        partial_cache = cds_from_cache or pppm_from_cache
        
        return all_matches, all_from_cache, partial_cache
    
    def migrate_cache_to_image_format(self, dry_run: bool = False) -> Dict[str, int]:
        """
        Migrate existing line_to_neuron cache files to new image-based format.
        
        Reads existing cache files, extracts per-image results, and saves them
        to the new image_cache directory.
        
        Parameters
        ----------
        dry_run : bool
            If True, only count files without actually migrating
            
        Returns
        -------
        dict
            Statistics: {'files_processed', 'images_extracted', 'errors'}
        """
        stats = {'files_processed': 0, 'images_extracted': 0, 'errors': 0}
        
        # Find all line_to_neuron cache files
        cache_files = [f for f in os.listdir(self.cache_folder) 
                       if f.startswith('line_to_neuron_') and f.endswith('.csv')]
        
        self._vprint(f"📦 Found {len(cache_files)} line_to_neuron cache files to migrate")
        
        for cache_file in cache_files:
            try:
                cache_path = os.path.join(self.cache_folder, cache_file)
                df = pd.read_csv(cache_path)
                
                if df.empty or 'lm_sample' not in df.columns:
                    continue
                
                # Extract match_type from the data or filename
                if 'match_type' in df.columns:
                    file_match_types = df['match_type'].unique()
                else:
                    # Try to extract from filename (e.g., line_to_neuron_VT037867_cds_Brain_5.csv)
                    parts = cache_file.replace('.csv', '').split('_')
                    if 'cds' in parts:
                        file_match_types = ['cds']
                    elif 'pppm' in parts:
                        file_match_types = ['pppm']
                    else:
                        file_match_types = ['cds']  # Default
                    df['match_type'] = file_match_types[0]
                
                # Group by lm_sample and match_type
                for (lm_sample, match_type), group_df in df.groupby(['lm_sample', 'match_type']):
                    if pd.isna(lm_sample):
                        continue
                    
                    image_id = str(int(lm_sample) if isinstance(lm_sample, float) else lm_sample)
                    
                    if not dry_run:
                        # Check if already migrated
                        existing = self._load_image_cache(image_id, match_type)
                        if existing is not None:
                            stats['images_extracted'] += 1
                            continue
                        
                        # Save to image cache
                        matches = group_df.to_dict('records')
                        self._save_image_cache(image_id, match_type, matches)
                    
                    stats['images_extracted'] += 1
                
                stats['files_processed'] += 1
                
            except Exception as e:
                stats['errors'] += 1
                self._vprint(f"  ⚠️ Error processing {cache_file}: {e}")
        
        action = "Would migrate" if dry_run else "Migrated"
        self._vprint(f"✓ {action} {stats['images_extracted']} images from {stats['files_processed']} files")
        if stats['errors'] > 0:
            self._vprint(f"  ⚠️ {stats['errors']} files had errors")
        
        return stats
    
    def _get_em_image_for_dataset(
        self,
        body_id: int,
        expected_dataset: Optional[str] = None
    ):
        """
        Get EM image for a body ID, optionally filtering by expected dataset.
        
        NeuronBridge can return multiple EM images for the same body ID from
        different datasets (e.g., vnc:v0.5, manc:v1.2.1, male-cns:v0.9).
        This method finds the one matching the expected dataset.
        
        Parameters
        ----------
        body_id : int
            The body ID to look up.
        expected_dataset : str, optional
            Expected dataset name (e.g., 'male-cns:v0.9').
            If None, returns the first result.
            
        Returns
        -------
        EMImage or None
            The matching EM image, or None if not found.
        """
        if not self._client:
            return None
        
        try:
            # Use get_em_images (plural) to get ALL results for this body ID
            em_images = self._client.get_em_images(body_id)
            
            if not em_images:
                return None
            
            # Convert to list if needed
            if not isinstance(em_images, list):
                em_images = list(em_images)
            
            if not em_images:
                return None
            
            # If no expected dataset, return first result
            if not expected_dataset:
                return em_images[0]
            
            # Normalize expected dataset for comparison
            expected_base = expected_dataset.lower().split(':')[0].replace('_', '-')
            
            # Find the EM image matching the expected dataset
            for em_image in em_images:
                pub_name = getattr(em_image, 'publishedName', '')
                if ':' in pub_name:
                    actual_base = pub_name.lower().split(':')[0].replace('_', '-')
                    if actual_base == expected_base:
                        return em_image
            
            # If no exact match, return first result
            return em_images[0]
            
        except Exception:
            # Fallback to singular get_em_image
            try:
                return self._client.get_em_image(body_id)
            except Exception:
                return None
    
    def _validate_body_id_dataset(
        self, 
        body_id: int, 
        expected_dataset: str
    ) -> Optional[str]:
        """
        Validate that a body ID exists in the expected dataset in NeuronBridge.
        
        Parameters
        ----------
        body_id : int
            The body ID to validate.
        expected_dataset : str
            Expected dataset name (e.g., 'male-cns:v0.9').
            
        Returns
        -------
        str or None
            The actual dataset name if found matching expected, or None if not found.
        """
        if not self._client:
            return expected_dataset  # Can't validate without client
        
        try:
            # Use get_em_images (plural) to get ALL results for this body ID
            em_images = self._client.get_em_images(body_id)
            
            if not em_images:
                return None
            
            # Convert to list if needed
            if not isinstance(em_images, list):
                em_images = list(em_images)
            
            if not em_images:
                return None
            
            # Normalize expected dataset for comparison
            expected_base = expected_dataset.lower().split(':')[0].replace('_', '-')
            
            # Check if any EM image matches the expected dataset
            for em_image in em_images:
                pub_name = getattr(em_image, 'publishedName', '')
                if ':' in pub_name:
                    parts = pub_name.split(':')
                    actual_base = parts[0].lower().replace('_', '-')
                    if actual_base == expected_base:
                        # Found matching dataset
                        return f"{parts[0]}:{parts[1]}" if len(parts) >= 2 else parts[0]
            
            # No matching dataset found - return what was found
            pub_name = getattr(em_images[0], 'publishedName', '')
            if ':' in pub_name:
                parts = pub_name.split(':')
                return f"{parts[0]}:{parts[1]}" if len(parts) >= 2 else parts[0]
            return None
            
        except Exception:
            return expected_dataset  # On error, assume expected dataset
    
    def _sort_matches_by_rank(
        self, 
        matches: List[Dict[str, Any]], 
        key_field: str = 'line'
    ) -> List[Dict[str, Any]]:
        """
        Sort matches by combined rank when both CDS and PPPM results exist.
        
        For each unique line/bodyId, compute:
        - CDS rank (1-based, by descending score)
        - PPPM rank (1-based, by descending score)
        - Combined rank = CDS_rank + PPPM_rank
        
        Sort by combined rank ascending (lower is better).
        
        Parameters
        ----------
        matches : list
            List of match dictionaries with 'match_type', 'score', and key_field.
        key_field : str
            Field to use as unique identifier ('line' for EM->LM, 'bodyId' for LM->EM).
            
        Returns
        -------
        list
            Sorted list of match dictionaries with added 'cds_rank', 'pppm_rank', 
            'combined_rank' columns.
        """
        if not matches:
            return matches
        
        # Separate CDS and PPPM matches
        cds_matches = [m for m in matches if m.get('match_type') == 'cds']
        pppm_matches = [m for m in matches if m.get('match_type') == 'pppm']
        
        # Sort each by score descending and assign ranks
        cds_matches.sort(key=lambda x: x.get('score', 0), reverse=True)
        pppm_matches.sort(key=lambda x: x.get('score', 0), reverse=True)
        
        # Build rank dictionaries (key -> rank)
        cds_ranks = {}
        for rank, m in enumerate(cds_matches, 1):
            key = m.get(key_field)
            if key and key not in cds_ranks:
                cds_ranks[key] = rank
        
        pppm_ranks = {}
        for rank, m in enumerate(pppm_matches, 1):
            key = m.get(key_field)
            if key and key not in pppm_ranks:
                pppm_ranks[key] = rank
        
        # Default rank for missing matches (use max+1 or a high number)
        max_cds_rank = len(cds_matches) + 1 if cds_matches else 1
        max_pppm_rank = len(pppm_matches) + 1 if pppm_matches else 1
        
        # Combine matches by key, keeping best score per match_type
        combined = {}
        for m in matches:
            key = m.get(key_field)
            if not key:
                continue
            
            if key not in combined:
                combined[key] = {
                    key_field: key,
                    'library': m.get('library', ''),
                    'image_id': m.get('image_id', ''),
                    'cds_score': 0,
                    'pppm_score': 0,
                    'cds_rank': max_cds_rank,
                    'pppm_rank': max_pppm_rank,
                }
                # Copy other fields for LM->EM matches
                for field in ['bodyId', 'dataset', 'instance', 'type', 'status', 'lm_sample']:
                    if field in m:
                        combined[key][field] = m[field]
            
            # Update scores and ranks
            if m.get('match_type') == 'cds':
                if m.get('score', 0) > combined[key]['cds_score']:
                    combined[key]['cds_score'] = m.get('score', 0)
                combined[key]['cds_rank'] = cds_ranks.get(key, max_cds_rank)
            elif m.get('match_type') == 'pppm':
                if m.get('score', 0) > combined[key]['pppm_score']:
                    combined[key]['pppm_score'] = m.get('score', 0)
                combined[key]['pppm_rank'] = pppm_ranks.get(key, max_pppm_rank)
        
        # Calculate combined rank and sort
        result = []
        for key, data in combined.items():
            data['combined_rank'] = data['cds_rank'] + data['pppm_rank']
            # Use higher score as primary score for display
            data['score'] = max(data['cds_score'], data['pppm_score'])
            data['match_type'] = 'both'
            result.append(data)
        
        # Sort by combined rank ascending
        result.sort(key=lambda x: x['combined_rank'])
        
        return result
    
    def _get_em_matches(
        self, 
        body_id: int, 
        match_type: Optional[str] = None,
        expected_dataset: Optional[str] = None
    ) -> List[Dict[str, Any]]:
        """
        Get LM matches for an EM body ID.
        
        Parameters
        ----------
        body_id : int
            The EM body ID to search for.
        match_type : str, optional
            Match algorithm: 'cds', 'pppm', or 'both'. 
            If None, uses self.match_type. Default: None
        expected_dataset : str, optional
            Expected dataset name (e.g., 'male-cns:v0.9').
            If provided, will filter EM images to match this dataset.
            
        Returns
        -------
        list
            List of match dictionaries, sorted by:
            - 'cds' or 'pppm': score descending
            - 'both': combined rank ascending (sum of CDS rank + PPPM rank)
        """
        # Use class-level match_type if not specified
        if match_type is None:
            match_type = self.match_type
            
        try:
            # Use helper to get EM image for the expected dataset
            em_image = self._get_em_image_for_dataset(body_id, expected_dataset)
            if em_image is None:
                self._vprint(f"  ⚠️ No EM image found for body ID {body_id}")
                return []
            
            # Check if em_image has required URL for fetching matches
            # Some body IDs exist but don't have valid match URLs in NeuronBridge
            em_files = getattr(em_image, 'files', None)
            if em_files is None:
                self._vprint(f"  ⚠️ No files metadata for body ID {body_id}")
                return []
            
            all_matches = []
            
            cds_failed = False
            pppm_failed = False
            
            # Get CDS matches with retry logic
            if match_type in ['cds', 'both']:
                try:
                    # Check if CDSResults URL exists
                    cds_url = getattr(em_files, 'CDSResults', None)
                    if cds_url is None:
                        cds_failed = True
                    else:
                        cds_matches = self._retry_with_backoff(
                            self._client.get_cds_matches,
                            em_image,
                            max_retries=3,
                            initial_delay=1.0
                        )
                        for match in cds_matches:
                            if hasattr(match, 'image') and hasattr(match.image, 'type'):
                                if match.image.type == 'LMImage':
                                    all_matches.append({
                                        'line': getattr(match.image, 'publishedName', ''),
                                        'library': getattr(match.image, 'libraryName', ''),
                                        'score': getattr(match, 'normalizedScore', 0),
                                        'image_id': getattr(match.image, 'id', ''),
                                        'match_type': 'cds'
                                    })
                except Exception as e:
                    cds_failed = True
                    # Suppress individual errors - will be caught by calling method
            
            # Get PPPM matches with retry logic
            if match_type in ['pppm', 'both']:
                try:
                    # Check if PPPMResults URL exists
                    pppm_url = getattr(em_files, 'PPPMResults', None)
                    if pppm_url is None:
                        pppm_failed = True
                    else:
                        pppm_matches = self._retry_with_backoff(
                            self._client.get_ppp_matches,
                            em_image,
                            max_retries=3,
                            initial_delay=1.0
                        )
                        for match in pppm_matches:
                            if hasattr(match, 'image') and hasattr(match.image, 'type'):
                                if match.image.type == 'LMImage':
                                    all_matches.append({
                                        'line': getattr(match.image, 'publishedName', ''),
                                        'library': getattr(match.image, 'libraryName', ''),
                                        'score': getattr(match, 'normalizedScore', 0),
                                        'image_id': getattr(match.image, 'id', ''),
                                        'match_type': 'pppm'
                                    })
                except Exception as e:
                    pppm_failed = True
                    # Suppress individual errors - will be caught by calling method
            
            # Only warn if both failed when match_type='both'
            if match_type == 'both' and cds_failed and pppm_failed:
                self._vprint(f"  ⚠️ Both CDS and PPPM matches failed for body ID {body_id}")
            
            # Sort results based on match_type
            if match_type == 'both' and all_matches:
                # Use rank-based sorting for combined results
                all_matches = self._sort_matches_by_rank(all_matches, key_field='line')
            else:
                # Sort by score descending for single match type
                all_matches.sort(key=lambda x: x.get('score', 0), reverse=True)
            
            return all_matches
            
        except Exception as e:
            self._vprint(f"  ⚠️ Error getting matches for body ID {body_id}: {e}")
            return []
    
    def _get_lm_matches(
        self, 
        line_name: str, 
        match_type: Optional[str] = None
    ) -> List[Dict[str, Any]]:
        """
        Get EM matches for an LM line name.
        
        Uses image-based caching to store results per LM image, allowing reuse
        regardless of region, top_n, or max_api_images_per_line settings.
        
        Always checks online for current images (since server images can update),
        but uses cached match results when available.
        
        Parameters
        ----------
        line_name : str
            The driver line name to search for.
        match_type : str, optional
            Match algorithm: 'cds', 'pppm', or 'both'. 
            If None, uses self.match_type. Default: None
            
        Returns
        -------
        list
            List of match dictionaries with dataset info, sorted by:
            - 'cds' or 'pppm': score descending
            - 'both': combined rank ascending (sum of CDS rank + PPPM rank)
        """
        # Use class-level match_type if not specified
        if match_type is None:
            match_type = self.match_type
            
        try:
            # Always check online for current images (server images can update)
            lm_images = self._client.get_lm_images(line_name)
            if not lm_images:
                self._vprint(f"  ⚠️ No LM images found for line '{line_name}'")
                return []
            
            # Convert generator to list if needed
            if not isinstance(lm_images, list):
                lm_images = list(lm_images)
            
            # Filter images by region
            lm_images = self._filter_images_by_region(lm_images)
            if not lm_images:
                self._vprint(f"  ⚠️ No LM images for line '{line_name}' in region '{self.region}'")
                return []
            
            # Filter images by match_type availability (pre-scan to skip unavailable)
            lm_images = self._filter_images_by_match_availability(lm_images, match_type)
            if not lm_images:
                self._vprint(f"  ⚠️ No LM images with {match_type.upper()} results for line '{line_name}'")
                return []
            
            # Update line-image mapping with current images from server
            all_image_ids = [str(getattr(img, 'id', '')) for img in lm_images if getattr(img, 'id', '')]
            if all_image_ids:
                self._update_line_mapping(line_name, self.region, all_image_ids)
            
            # Apply max_api_images_per_line limit
            original_count = len(lm_images)
            if self.max_api_images_per_line > 0 and len(lm_images) > self.max_api_images_per_line:
                lm_images = lm_images[:self.max_api_images_per_line]
                self._vprint(f"  ℹ️  Using {len(lm_images)}/{original_count} images (max_api_images_per_line={self.max_api_images_per_line})")
            
            n_images = len(lm_images)
            # Only print verbose message if not in a progress bar context
            if not HAS_TQDM or not self.verbose:
                self._vprint(f"  Found {n_images} LM images for '{line_name}'")
            
            all_matches = []
            cache_hits = 0
            partial_cache_hits = 0
            api_fetches = 0
            errors = 0
            
            # Create progress bar for image processing if we have multiple images and tqdm is available
            show_image_progress = HAS_TQDM and self.verbose and n_images > 1
            
            if show_image_progress:
                from tqdm import tqdm as tqdm_progress
                # Suppress loading messages while progress bar is active
                self._suppress_loading_msgs = True
                image_iterator = tqdm_progress(
                    lm_images,
                    desc=f"  🖼️  Processing images",
                    unit="img",
                    leave=False,
                    bar_format='  {desc}: {percentage:3.0f}%|{bar:20}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]',
                    ncols=90
                )
            else:
                image_iterator = lm_images
            
            for img_idx, lm_image in enumerate(image_iterator, 1):
                # Use the unified _get_image_matches_cached which handles 'both' internally
                matches, from_cache, partial_cache = self._get_image_matches_cached(
                    lm_image, match_type, line_name
                )
                
                if from_cache:
                    cache_hits += 1
                elif partial_cache:
                    partial_cache_hits += 1
                else:
                    api_fetches += 1
                
                if matches:
                    all_matches.extend(matches)
                elif not from_cache and not partial_cache:
                    errors += 1
            
            # Restore loading messages flag after processing images
            if show_image_progress:
                self._suppress_loading_msgs = False
            
            # Report cache/API statistics explicitly
            total_lookups = cache_hits + partial_cache_hits + api_fetches
            if total_lookups > 0:
                cache_count = cache_hits + partial_cache_hits
                self._vprint(f"  📊 Image data: {cache_count}/{total_lookups} from cache, {api_fetches}/{total_lookups} API fetches")
            
            # Report error summary if there were failures
            if errors > 0:
                self._vprint(f"  ℹ️  {errors}/{n_images} images had no matches")
            
            # Sort results based on match_type
            if match_type == 'both' and all_matches:
                # Use rank-based sorting for combined results
                unique_matches = self._sort_matches_by_rank(all_matches, key_field='bodyId')
            else:
                # Sort by score descending for single match type
                all_matches.sort(key=lambda x: x.get('score', 0), reverse=True)
                
                # Remove duplicates (same bodyId + dataset), keeping highest score
                seen = set()
                unique_matches = []
                for match in all_matches:
                    key = (match.get('bodyId'), match.get('dataset'))
                    if key not in seen:
                        seen.add(key)
                        unique_matches.append(match)
            
            return unique_matches
            
        except Exception as e:
            self._vprint(f"  ⚠️ Error getting matches for line '{line_name}': {e}")
            return []
    
    def _extract_body_id(self, em_image) -> str:
        """Extract body ID from an EM image object."""
        # Try publishedName first (format: 'hemibrain:v1.2.1:BODY_ID')
        published_name = getattr(em_image, 'publishedName', '')
        if published_name and ':' in published_name:
            parts = published_name.split(':')
            if len(parts) >= 3:
                return parts[-1]
        
        # Try id field
        image_id = getattr(em_image, 'id', '')
        if image_id:
            # Extract digits
            digits = ''.join(c for c in str(image_id) if c.isdigit())
            if digits:
                return digits
        
        return ''
    
    def _find_bodyIds_by_query(
        self, 
        query: Union[str, int, List],
        dataset: Optional[Union[str, List[str]]] = None
    ) -> List[Dict[str, Any]]:
        """
        Find body IDs matching a query in the neuron DataFrame(s).
        
        Parameters
        ----------
        query : str, int, or list
            - If int: treated as bodyId
            - If str: searched in instance and type columns (supports regex)
            - If list: each item processed individually
        dataset : str, list of str, or None
            Specific dataset(s) to search (e.g., 'hemibrain:v1.2.1', 'male-cns:v0.9').
            Can be a single string or a list of datasets.
            If None, searches all available datasets.
            
        Returns
        -------
        list
            List of dicts with 'bodyId' and 'dataset' keys.
        """
        # Determine which datasets to search
        datasets_to_search = []
        
        if dataset:
            # Handle list of datasets
            dataset_list = [dataset] if isinstance(dataset, str) else list(dataset)
            
            for ds in dataset_list:
                # Map display name to folder name if needed
                dataset_folder = self._dataset_name_to_folder(ds)
                if dataset_folder:
                    datasets_to_search.append(dataset_folder)
                else:
                    datasets_to_search.append(ds.replace(':', '_').replace('.', '_'))
            # Remove duplicates while preserving order
            datasets_to_search = list(dict.fromkeys(datasets_to_search))
        else:
            # Search all available datasets
            # Get unique dataset folder names
            datasets_to_search = list(set(LIBRARY_TO_DATASET.values()))
        
        if isinstance(query, list):
            all_results = []
            for q in query:
                all_results.extend(self._find_bodyIds_by_query(q, dataset))
            # Remove duplicates (same bodyId + dataset)
            seen = set()
            unique_results = []
            for r in all_results:
                key = (r['bodyId'], r['dataset'])
                if key not in seen:
                    seen.add(key)
                    unique_results.append(r)
            return unique_results
        
        if isinstance(query, int):
            # For direct body ID, we don't know the dataset - return for all with matching ID
            results = []
            for ds_folder in datasets_to_search:
                neuron_df = self._load_neuron_df_for_dataset(ds_folder)
                if neuron_df is not None and not neuron_df.empty:
                    if str(query) in neuron_df['bodyId'].astype(str).values:
                        ds_name = self._folder_to_dataset_name(ds_folder)
                        results.append({'bodyId': str(query), 'dataset': ds_name, 'dataset_folder': ds_folder})
            # If not found in any, still return it
            if not results:
                results = [{'bodyId': str(query), 'dataset': 'unknown', 'dataset_folder': ''}]
            return results
        
        query_str = str(query)
        matched_results = []
        
        # Try exact bodyId match first
        if query_str.isdigit():
            for ds_folder in datasets_to_search:
                neuron_df = self._load_neuron_df_for_dataset(ds_folder)
                if neuron_df is not None and not neuron_df.empty:
                    if query_str in neuron_df['bodyId'].astype(str).values:
                        ds_name = self._folder_to_dataset_name(ds_folder)
                        matched_results.append({
                            'bodyId': query_str, 
                            'dataset': ds_name,
                            'dataset_folder': ds_folder
                        })
            if matched_results:
                return matched_results
        
        # Search in type/instance columns across all datasets
        for ds_folder in datasets_to_search:
            neuron_df = self._load_neuron_df_for_dataset(ds_folder)
            if neuron_df is None or neuron_df.empty:
                continue
            
            ds_name = self._folder_to_dataset_name(ds_folder)
            
            # Search in 'type' column
            if 'type' in neuron_df.columns:
                try:
                    type_matches = neuron_df[
                        neuron_df['type'].astype(str).str.match(f'^{query_str}$', case=False, na=False)
                    ]
                    for body_id in type_matches['bodyId'].astype(str).tolist():
                        matched_results.append({
                            'bodyId': body_id, 
                            'dataset': ds_name,
                            'dataset_folder': ds_folder
                        })
                except re.error:
                    type_matches = neuron_df[
                        neuron_df['type'].astype(str).str.lower() == query_str.lower()
                    ]
                    for body_id in type_matches['bodyId'].astype(str).tolist():
                        matched_results.append({
                            'bodyId': body_id, 
                            'dataset': ds_name,
                            'dataset_folder': ds_folder
                        })
            
            # Search in 'instance' column
            if 'instance' in neuron_df.columns:
                try:
                    instance_matches = neuron_df[
                        neuron_df['instance'].astype(str).str.match(f'^{query_str}$', case=False, na=False)
                    ]
                    for body_id in instance_matches['bodyId'].astype(str).tolist():
                        matched_results.append({
                            'bodyId': body_id, 
                            'dataset': ds_name,
                            'dataset_folder': ds_folder
                        })
                except re.error:
                    instance_matches = neuron_df[
                        neuron_df['instance'].astype(str).str.lower() == query_str.lower()
                    ]
                    for body_id in instance_matches['bodyId'].astype(str).tolist():
                        matched_results.append({
                            'bodyId': body_id, 
                            'dataset': ds_name,
                            'dataset_folder': ds_folder
                        })
        
        # Remove duplicates (same bodyId + dataset)
        seen = set()
        unique_results = []
        for r in matched_results:
            key = (r['bodyId'], r['dataset'])
            if key not in seen:
                seen.add(key)
                unique_results.append(r)
        
        return unique_results
    
    def _dataset_name_to_folder(self, dataset_name: str) -> Optional[str]:
        """Convert display dataset name to folder name."""
        # Try direct lookup via reverse mapping
        for lib, name in LIBRARY_TO_DATASET_NAME.items():
            if name.lower() == dataset_name.lower():
                return LIBRARY_TO_DATASET.get(lib)
        # Try folder names directly
        for lib, folder in LIBRARY_TO_DATASET.items():
            if folder.lower() == dataset_name.lower().replace(':', '_').replace('.', '_'):
                return folder
        return None
    
    def _folder_to_dataset_name(self, folder: str) -> str:
        """Convert folder name to display dataset name."""
        # Reverse lookup in LIBRARY_TO_DATASET
        for lib, lib_folder in LIBRARY_TO_DATASET.items():
            if lib_folder == folder:
                return LIBRARY_TO_DATASET_NAME.get(lib, folder)
        return folder
    
    # =========================================================================
    # Public Methods
    # =========================================================================
    
    def id_to_lines(
        self, 
        body_id: int, 
        match_type: str = 'cds',
        expected_dataset: Optional[str] = None
    ) -> pd.DataFrame:
        """
        Find driver lines matching a given EM body ID.
        
        Parameters
        ----------
        body_id : int
            The EM body ID to search for.
        match_type : str
            Match algorithm: 'cds', 'pppm', or 'both'. Default: 'cds'
        expected_dataset : str, optional
            Expected dataset name (e.g., 'male-cns:v0.9').
            If provided, will filter EM images to match this dataset.
            
        Returns
        -------
        pd.DataFrame
            DataFrame with columns: line, library, score, image_id, match_type
            Sorted by score descending (cds/pppm) or combined_rank ascending (both)
            
        Example
        -------
        >>> nbf = NeuronBridgeFinder()
        >>> lines = nbf.id_to_lines(636798093)
        >>> print(lines)
        """
        # Validate and normalize match_type
        match_type = self._validate_match_type(match_type)
        
        self._vprint(f"🔍 Searching for lines matching body ID: {body_id}")
        
        # Check cache (include all parameters in key)
        ds_key = expected_dataset.replace(':', '_') if expected_dataset else 'any'
        region_key = self.region if self.region else 'all'
        max_imgs_key = self.max_api_images_per_line if self.max_api_images_per_line > 0 else 'all'
        cache_key = f"{body_id}_{match_type}_{ds_key}_{region_key}_{max_imgs_key}"
        cached = self._load_from_cache('id_to_lines', cache_key)
        if cached is not None:
            return cached
        
        # Fetch from API with expected dataset
        matches = self._get_em_matches(
            body_id, 
            match_type=match_type, 
            expected_dataset=expected_dataset
        )
        
        if not matches:
            self._vprint(f"  ℹ️ No matches found for body ID {body_id}")
            cols = ['line', 'library', 'score', 'image_id', 'match_type']
            if match_type == 'both':
                cols.extend(['cds_score', 'pppm_score', 'cds_rank', 'pppm_rank', 'combined_rank'])
            return pd.DataFrame(columns=cols)
        
        df = pd.DataFrame(matches)
        
        # Cache results
        self._save_to_cache('id_to_lines', cache_key, df)
        
        self._vprint(f"  ✓ Found {len(df)} matches")
        
        return df
    
    def neuron_to_lines(
        self, 
        query: Union[str, int, List],
        dataset: Optional[Union[str, List[str]]] = None,
        match_type: str = 'cds'
    ) -> Dict[str, pd.DataFrame]:
        """
        Find driver lines matching neurons specified by bodyId, instance, or type.
        
        Parameters
        ----------
        query : str, int, or list
            - If int: treated as bodyId
            - If str: searched in instance and type columns (supports regex)
            - If list: each item processed individually
        dataset : str, list of str, or None
            Specific dataset(s) to search (e.g., 'hemibrain:v1.2.1', 'male-cns:v0.9').
            Can be a single string or a list of datasets to search multiple.
            If None, searches all available datasets.
        match_type : str
            Match algorithm: 'cds', 'pppm', or 'both'. Default: 'cds'
            
        Returns
        -------
        dict
            Dictionary mapping bodyId -> DataFrame of matching lines
            All results sorted by score descending (cds/pppm) or combined_rank ascending (both)
            
        Example
        -------
        >>> nbf = NeuronBridgeFinder()
        >>> results = nbf.neuron_to_lines('MBON01')
        >>> for body_id, lines_df in results.items():
        ...     print(f"Body ID {body_id}: {len(lines_df)} matches")
        """
        # Validate and normalize match_type
        match_type = self._validate_match_type(match_type)
        
        self._vprint(f"🔍 Searching for lines matching query: {query}")
        
        # Find matching bodyIds (now returns list of dicts with bodyId and dataset)
        body_info_list = self._find_bodyIds_by_query(query, dataset)
        
        if not body_info_list:
            self._vprint(f"  ⚠️ No neurons found matching query: {query}")
            return {}
        
        self._vprint(f"  Found {len(body_info_list)} neurons to search")
        
        results = {}
        skipped_count = 0
        
        # Use progress bar for multiple body IDs
        if HAS_TQDM and self.verbose and len(body_info_list) > 1:
            from tqdm import tqdm as tqdm_progress
            
            # Enable batch mode to suppress individual messages
            self._batch_mode = True
            
            pbar = tqdm_progress(
                body_info_list,
                desc=f"  🔄 Processing {len(body_info_list)} neurons",
                unit="neuron",
                bar_format='{desc}: {percentage:3.0f}%|{bar:30}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]',
                ncols=110,
                position=0,
                leave=False
            )
            
            for body_info in pbar:
                body_id = body_info['bodyId']
                expected_ds = body_info.get('dataset', 'unknown')
                pbar.set_postfix_str(f"{body_id}")
                
                # Validate body ID against NeuronBridge to ensure dataset match
                actual_ds = self._validate_body_id_dataset(int(body_id), expected_ds)
                
                # Check for dataset mismatch
                if actual_ds and expected_ds != 'unknown':
                    expected_base = self._normalize_dataset_name(expected_ds)
                    actual_base = self._normalize_dataset_name(actual_ds)
                    
                    if expected_base != actual_base:
                        skipped_count += 1
                        continue
                
                try:
                    lines_df = self.id_to_lines(
                        int(body_id), 
                        match_type=match_type,
                        expected_dataset=expected_ds
                    )
                    if not lines_df.empty:
                        lines_df = lines_df.copy()
                        lines_df['source_dataset'] = actual_ds or expected_ds
                        lines_df['source_bodyId'] = body_id
                    results[body_id] = lines_df
                except Exception as e:
                    results[body_id] = pd.DataFrame()
            
            pbar.close()
            self._batch_mode = False
            
            if skipped_count > 0:
                self._vprint(f"  ℹ️  Skipped {skipped_count} body IDs due to dataset mismatch")
        else:
            # Single neuron or no tqdm - original behavior
            for i, body_info in enumerate(body_info_list):
                body_id = body_info['bodyId']
                expected_ds = body_info.get('dataset', 'unknown')
                
                # Validate body ID against NeuronBridge to ensure dataset match
                actual_ds = self._validate_body_id_dataset(int(body_id), expected_ds)
                
                # Check for dataset mismatch
                if actual_ds and expected_ds != 'unknown':
                    # Normalize for comparison - extract base dataset name
                    # Handle different naming conventions:
                    # - 'flywire_FAFB_v783' (underscore + version suffix)
                    # - 'flywire_fafb:v783' (colon-separated version)
                    # - 'male-cns:v0.9' -> 'male-cns'
                    # - 'hemibrain:v1.2.1' -> 'hemibrain'
                    expected_base = self._normalize_dataset_name(expected_ds)
                    actual_base = self._normalize_dataset_name(actual_ds)
                    
                    if expected_base != actual_base:
                        self._vprint(f"  ⏭️  Skipping {body_id}: belongs to {actual_ds}, not {expected_ds} (dataset mismatch)")
                        skipped_count += 1
                        continue
                
                self._vprint(f"  Processing {i+1}/{len(body_info_list)}: {body_id} ({actual_ds or expected_ds})")
                try:
                    # Pass expected dataset to get correct EM image
                    lines_df = self.id_to_lines(
                        int(body_id), 
                        match_type=match_type,
                        expected_dataset=expected_ds
                    )
                    # Add source dataset info (use actual dataset from NeuronBridge)
                    if not lines_df.empty:
                        lines_df = lines_df.copy()
                        lines_df['source_dataset'] = actual_ds or expected_ds
                        lines_df['source_bodyId'] = body_id
                    results[body_id] = lines_df
                except Exception as e:
                    self._vprint(f"    ⚠️ Error processing body ID {body_id}: {e}")
                    results[body_id] = pd.DataFrame()
            
            if skipped_count > 0:
                self._vprint(f"  ℹ️  Skipped {skipped_count} body IDs due to dataset mismatch")
        
        return results
    
    def line_to_neuron(
        self, 
        line_name: str,
        match_type: str = 'cds',
        top_n: int = -1
    ) -> pd.DataFrame:
        """
        Find EM neurons matching a driver line name.
        
        Parameters
        ----------
        line_name : str
            The driver line name to search for (e.g., 'LH173').
        match_type : str
            Match algorithm: 'cds', 'pppm', or 'both'. Default: 'cds'
        top_n : int
            Maximum number of matches to return. Default: -1 (all matches)
            
        Returns
        -------
        pd.DataFrame
            DataFrame with columns: bodyId, dataset, instance, type, status, 
            score, image_id, lm_sample, match_type, library
            Sorted by score descending (cds/pppm) or combined_rank ascending (both)
            
        Example
        -------
        >>> nbf = NeuronBridgeFinder()
        >>> neurons = nbf.line_to_neuron('LH173')
        >>> print(neurons)
        """
        # Validate and normalize match_type
        match_type = self._validate_match_type(match_type)
        
        # Only print search message if not in batch/progress mode
        # (check if we're being called from within a progress bar loop)
        import inspect
        frame = inspect.currentframe()
        in_progress_context = False
        if frame and frame.f_back:
            # Check if caller's local variables include tqdm iterator
            caller_locals = frame.f_back.f_locals
            in_progress_context = any('tqdm' in str(type(v)).lower() for v in caller_locals.values())
        
        if not in_progress_context:
            self._vprint(f"🔍 Searching for neurons matching line: {line_name}")
        
        # Check cache - include all relevant parameters
        region_key = self.region if self.region else 'all'
        max_imgs_key = self.max_api_images_per_line if self.max_api_images_per_line > 0 else 'all'
        cache_key = f"{line_name}_{match_type}_{region_key}_{max_imgs_key}"
        cached = self._load_from_cache('line_to_neuron', cache_key)
        if cached is not None:
            # Indicate cache hit in verbose mode
            if self.verbose and not any(c in str(self._vprint.__code__) for c in ['silent', 'quiet']):
                # Only print if not suppressed by progress bar
                pass  # Progress bar will show cache status
            return cached
        
        # Fetch from API (matches are already enriched with dataset info)
        matches = self._get_lm_matches(line_name, match_type=match_type)
        
        if not matches:
            self._vprint(f"  ℹ️ No matches found for line '{line_name}'")
            cols = ['bodyId', 'dataset', 'instance', 'type', 'status', 'score', 
                    'image_id', 'lm_sample', 'match_type', 'library']
            if match_type == 'both':
                cols.extend(['cds_score', 'pppm_score', 'cds_rank', 'pppm_rank', 'combined_rank'])
            return pd.DataFrame(columns=cols)
        
        df = pd.DataFrame(matches)
        
        # Reorder columns - include dataset and rank columns if present
        base_cols = ['bodyId', 'dataset', 'instance', 'type', 'status', 'score', 
                     'image_id', 'lm_sample', 'match_type', 'library']
        rank_cols = ['cds_score', 'pppm_score', 'cds_rank', 'pppm_rank', 'combined_rank']
        cols = [c for c in base_cols if c in df.columns]
        cols.extend([c for c in rank_cols if c in df.columns])
        df = df[cols]
        
        # Apply top_n limit if specified
        if top_n > 0 and len(df) > top_n:
            df = df.head(top_n)
        
        # Cache results
        self._save_to_cache('line_to_neuron', cache_key, df)
        
        self._vprint(f"  ✓ Found {len(df)} matches")
        
        return df
    
    def save_results(
        self, 
        results: Union[pd.DataFrame, Dict[str, pd.DataFrame]], 
        output_path: str,
        include_timestamp: bool = True
    ) -> str:
        """
        Save results to a CSV file.
        
        Parameters
        ----------
        results : DataFrame or dict
            Results from id_to_lines, neuron_to_lines, or line_to_neuron.
        output_path : str
            Output file path.
        include_timestamp : bool
            Whether to include timestamp in filename. Default: True
            
        Returns
        -------
        str
            Path to the saved file.
        """
        if include_timestamp:
            base, ext = os.path.splitext(output_path)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_path = f"{base}_{timestamp}{ext}"
        
        if isinstance(results, dict):
            # Combine dict of DataFrames into one
            combined = []
            for body_id, df in results.items():
                if not df.empty:
                    df_copy = df.copy()
                    df_copy.insert(0, 'query_bodyId', body_id)
                    combined.append(df_copy)
            
            if combined:
                final_df = pd.concat(combined, ignore_index=True)
            else:
                final_df = pd.DataFrame()
        else:
            final_df = results
        
        final_df.to_csv(output_path, index=False)
        self._vprint(f"💾 Saved results to: {output_path}")
        
        return output_path
    
    def _save_dataset_categorized_files(
        self,
        neurons_df: pd.DataFrame,
        line_name: str,
        output_path: str,
        verbose: bool = True
    ) -> None:
        """
        Save dataset-categorized neuron files and type summary files.
        
        Creates:
        - {line}_{dataset}_neurons.csv: Neurons for each dataset
        - {line}_{dataset}_types.csv: Type summary with labeled_N and typed_N_in_dataset
        
        Parameters
        ----------
        neurons_df : pd.DataFrame
            DataFrame with matched neurons (must have 'dataset' column)
        line_name : str
            Driver line name
        output_path : str
            Output directory path
        verbose : bool
            Whether to print progress messages
        """
        if 'dataset' not in neurons_df.columns:
            return
        
        # Group by dataset
        for dataset, ds_df in neurons_df.groupby('dataset'):
            # Normalize dataset name for filename (replace : with _)
            ds_filename = dataset.replace(':', '_').replace('.', '_')
            
            # Save dataset-specific neurons file
            ds_neurons_file = os.path.join(output_path, f'{line_name}_{ds_filename}_neurons.csv')
            ds_df.to_csv(ds_neurons_file, index=False)
            if verbose:
                self._vprint(f"   💾 Saved: {ds_neurons_file}")
            
            # Create type summary
            type_summary = self._create_type_summary(ds_df, dataset)
            if not type_summary.empty:
                ds_types_file = os.path.join(output_path, f'{line_name}_{ds_filename}_types.csv')
                type_summary.to_csv(ds_types_file, index=False)
                if verbose:
                    self._vprint(f"   💾 Saved: {ds_types_file}")
    
    def _create_type_summary(
        self,
        neurons_df: pd.DataFrame,
        dataset: str
    ) -> pd.DataFrame:
        """
        Create a type summary DataFrame with labeled_N and typed_N_in_dataset.
        
        Parameters
        ----------
        neurons_df : pd.DataFrame
            DataFrame with matched neurons
        dataset : str
            Dataset name for looking up total type counts
            
        Returns
        -------
        pd.DataFrame
            Type summary with columns: type, labeled_N, avg_score, typed_N_in_dataset
        """
        # Determine type column - use 'type' if available, else use bodyId
        if 'type' in neurons_df.columns:
            # For untyped neurons, use 'unknown_{bodyId}' format
            neurons_df = neurons_df.copy()
            neurons_df['type_label'] = neurons_df.apply(
                lambda row: f"unknown_{row['bodyId']}" if pd.isna(row['type']) or row['type'] == '' else row['type'],
                axis=1
            )
        else:
            neurons_df = neurons_df.copy()
            neurons_df['type_label'] = neurons_df['bodyId'].apply(lambda x: f"unknown_{x}")
        
        # Group by type and calculate statistics
        score_col = 'score' if 'score' in neurons_df.columns else None
        
        if score_col:
            type_stats = neurons_df.groupby('type_label').agg(
                labeled_N=('bodyId', 'count'),
                avg_score=(score_col, 'mean')
            ).reset_index()
        else:
            type_stats = neurons_df.groupby('type_label').agg(
                labeled_N=('bodyId', 'count')
            ).reset_index()
            type_stats['avg_score'] = None
        
        type_stats = type_stats.rename(columns={'type_label': 'type'})
        
        # Get total count of each type in the dataset
        type_stats['typed_N_in_dataset'] = type_stats['type'].apply(
            lambda t: self._get_type_count_in_dataset(t, dataset)
        )
        
        # Sort by avg_score descending (primary), then by labeled_N descending (secondary)
        # This ensures top-N types are the ones with highest match scores
        if score_col:
            type_stats = type_stats.sort_values(['avg_score', 'labeled_N'], ascending=[False, False])
        else:
            type_stats = type_stats.sort_values('labeled_N', ascending=False)
        
        # Reorder columns
        cols = ['type', 'labeled_N', 'typed_N_in_dataset']
        if score_col:
            cols.insert(2, 'avg_score')
        type_stats = type_stats[cols]
        
        return type_stats
    
    def _get_type_count_in_dataset(self, type_name: str, dataset: str) -> int:
        """
        Get the total count of a neuron type in the dataset.
        
        Uses _load_neuron_df_for_dataset() which will pull the dataset from
        NeuPrint if not available locally.
        
        Parameters
        ----------
        type_name : str
            Neuron type name (or bodyId if untyped)
        dataset : str
            Dataset name (e.g., 'hemibrain:v1.2.1' or 'male-cns:v0.9')
            
        Returns
        -------
        int
            Total count of the type in the dataset, or 1 if it's a bodyId or lookup fails
        """
        # If type_name is unknown_{bodyId} or looks like a bodyId (numeric), return 1
        if type_name.startswith('unknown_') or type_name.isdigit():
            return 1
        
        # Try to load the dataset neuron index
        try:
            # Convert dataset name to folder format (e.g., 'hemibrain:v1.2.1' -> 'hemibrain_v1_2_1')
            dataset_folder = self._dataset_name_to_folder(dataset)
            
            if dataset_folder is None:
                # Try direct conversion if mapping fails
                dataset_folder = dataset.replace(':', '_').replace('.', '_').replace('-', '_')
            
            # Use _load_neuron_df_for_dataset which handles loading and pulling if needed
            neuron_df = self._load_neuron_df_for_dataset(dataset_folder)
            
            if neuron_df is not None and 'type' in neuron_df.columns:
                # Count neurons of this type (case-sensitive match)
                count = len(neuron_df[neuron_df['type'] == type_name])
                return count if count > 0 else 1
        except Exception:
            pass
        
        return 1  # Default if lookup fails
    
    def _apply_type_filter(
        self,
        type_list: List[Tuple[int, str]],
        type_filter: Optional[Dict[str, Union[str, List[str]]]]
    ) -> List[Tuple[int, str]]:
        """
        Filter type names based on filter criteria while preserving original ranks.
        
        Filter logic:
        - Multiple values within a key (list): OR logic (match any)
        - Across different keys: AND logic (must match all filter types)
        
        Parameters
        ----------
        type_list : list of (rank, type_name) tuples
            List of types with their original ranks (before filtering)
        type_filter : dict, optional
            Filter criteria with keys: 'contains', 'startswith', 'endswith', 'regex'
            Each value can be a string or list of strings.
            Example: {'contains': ['DN', 'AN'], 'startswith': 'M'}
            
        Returns
        -------
        list of (rank, type_name) tuples
            Filtered list preserving original ranks
        """
        if type_filter is None or not type_filter:
            return type_list
        
        import re
        
        def matches_filter(type_name: str) -> bool:
            """Check if type_name matches all filter criteria."""
            # Each key must be satisfied (AND logic across keys)
            for filter_type, filter_values in type_filter.items():
                # Convert single value to list for uniform processing
                if isinstance(filter_values, str):
                    filter_values = [filter_values]
                
                # OR logic within a key: match any value
                matched = False
                
                if filter_type == 'contains':
                    for val in filter_values:
                        if val in type_name:
                            matched = True
                            break
                elif filter_type == 'startswith':
                    for val in filter_values:
                        if type_name.startswith(val):
                            matched = True
                            break
                elif filter_type == 'endswith':
                    for val in filter_values:
                        if type_name.endswith(val):
                            matched = True
                            break
                elif filter_type == 'regex':
                    for val in filter_values:
                        try:
                            if re.search(val, type_name):
                                matched = True
                                break
                        except re.error:
                            self._vprint(f"   ⚠️  Invalid regex pattern: {val}")
                            continue
                else:
                    # Unknown filter type, skip
                    self._vprint(f"   ⚠️  Unknown filter type: {filter_type}")
                    matched = True  # Don't fail on unknown types
                
                if not matched:
                    return False  # AND logic: all keys must match
            
            return True
        
        # Filter while preserving ranks
        filtered = [(rank, name) for rank, name in type_list if matches_filter(name)]
        
        return filtered
    
    def _get_top_types_fallback(
        self,
        ds_df: pd.DataFrame,
        top_n: Optional[int] = None
    ) -> List[str]:
        """
        Fallback method to get top N types by score when labeling_info is not available.
        
        Parameters
        ----------
        ds_df : pd.DataFrame
            Dataset-filtered DataFrame with 'type_label' and 'score' columns
        top_n : int, optional
            Number of top types to return. If None, returns all types sorted by score.
            
        Returns
        -------
        list
            List of top type names (case-sensitive)
        """
        score_col = 'score' if 'score' in ds_df.columns else None
        
        if score_col:
            type_stats = ds_df.groupby('type_label').agg(
                avg_score=(score_col, 'mean'),
                count=('bodyId', 'count')
            ).reset_index()
            type_stats = type_stats.sort_values('avg_score', ascending=False)
            if top_n is not None:
                type_stats = type_stats.head(top_n)
            return type_stats['type_label'].tolist()
        else:
            # Fallback to count-based if no score column
            type_counts = ds_df['type_label'].value_counts()
            if top_n is not None:
                type_counts = type_counts.head(top_n)
            return type_counts.index.tolist()
    
    def _visualize_top_types(
        self,
        combined_df: pd.DataFrame,
        top_n: int,
        output_path: str,
        per_dataset: bool = True,
        source_line: str = '',
        visualize_by: str = 'type',
        generate_individual_profiles: Union[bool, List[str]] = None,
        pdf_images_per_page: Tuple[int, int] = (4, 3),
        labeling_info: Optional[pd.DataFrame] = None,
        type_filter: Optional[Dict[str, Union[str, List[str]]]] = None,
        datasets_to_visualize: Union[str, List[str]] = 'all',
    ) -> None:
        """
        Visualize top N types/bodyIds per dataset using VisualizeSkeleton.
        
        Parameters
        ----------
        combined_df : pd.DataFrame
            Combined DataFrame with all matched neurons
        top_n : int
            Number of top types/bodyIds to visualize per dataset
        output_path : str
            Output directory for visualizations
        per_dataset : bool
            If True, create separate visualization per dataset
        source_line : str
            Source line name for folder naming
        visualize_by : str
            How to organize: 'type' (merge) or 'bodyId' (individual)
        generate_individual_profiles : list of str, bool, or None
            Formats to generate: ['pdf'], ['pptx'], or ['pdf', 'pptx']
            Set to False or None to disable generation
        pdf_images_per_page : tuple
            (columns, rows) for PDF layout
        labeling_info : pd.DataFrame, optional
            DataFrame with case-sensitive types and dataset column for filtering.
            Columns: type, dataset, {line1_score}, {line2_score}, ...
        type_filter : dict, optional
            Filter neuron types by name pattern. When specified, gets ALL types
            first, applies filter, then takes top N from filtered results.
            Keys: 'contains', 'startswith', 'endswith', 'regex'
            Values: str or list of str
            Logic: OR within same key, AND across keys
            Original ranks are preserved in output labels (r{N} reflects rank
            before filtering, allowing tracking of original ranking position).
        datasets_to_visualize : str or list, default 'all'
            Constrain which datasets to visualize. 'all' for all datasets.
        """
        try:
            from visualize_skeleton import VisualizeSkeleton
        except ImportError:
            try:
                from .visualize_skeleton import VisualizeSkeleton
            except ImportError:
                self._vprint("⚠️  VisualizeSkeleton not available for visualization")
                return
        
        if 'dataset' not in combined_df.columns:
            self._vprint("⚠️  No dataset column for visualization")
            return
        
        mode_label = 'types' if visualize_by == 'type' else 'bodyIds'
        self._vprint(f"\n🎨 Visualizing top {top_n} {mode_label}...")
        
        # Get all unique datasets from results
        all_datasets = combined_df['dataset'].unique().tolist()
        
        # Apply datasets_to_visualize filter
        if datasets_to_visualize is None or datasets_to_visualize == 'all':
            datasets = all_datasets
        elif isinstance(datasets_to_visualize, str):
            # Single dataset specified
            datasets = [d for d in all_datasets if d == datasets_to_visualize]
            if not datasets:
                self._vprint(f"   ⚠️  Dataset '{datasets_to_visualize}' not found in results")
                self._vprint(f"   Available datasets: {all_datasets}")
                return
        else:
            # List of datasets specified
            datasets = [d for d in all_datasets if d in datasets_to_visualize]
            if not datasets:
                self._vprint(f"   ⚠️  None of the specified datasets found in results")
                self._vprint(f"   Requested: {datasets_to_visualize}")
                self._vprint(f"   Available: {all_datasets}")
                return
            skipped = [d for d in datasets_to_visualize if d not in all_datasets]
            if skipped:
                self._vprint(f"   ⚠️  Skipping unavailable datasets: {skipped}")
        
        if type_filter:
            self._vprint(f"   🔍 Type filter: {type_filter}")
        if datasets_to_visualize not in (None, 'all'):
            self._vprint(f"   📂 Visualizing datasets: {datasets}")
        
        # Pre-load all required datasets (pull from NeuPrint if not available locally)
        self._vprint(f"\n📦 Loading datasets for visualization...")
        loaded_datasets = self._ensure_datasets_loaded(datasets)
        
        if not loaded_datasets:
            self._vprint("   ⚠️  No datasets could be loaded for visualization")
            return
        
        self._vprint(f"   ✓ {len(loaded_datasets)} dataset(s) ready for visualization")
        
        for dataset in datasets:
            ds_df = combined_df[combined_df['dataset'] == dataset].copy()
            
            if visualize_by == 'type':
                # Group by type
                # Create type label (use 'unknown_{bodyId}' for untyped)
                if 'type' in ds_df.columns:
                    ds_df['type_label'] = ds_df.apply(
                        lambda row: f"unknown_{row['bodyId']}" if pd.isna(row['type']) or row['type'] == '' else row['type'],
                        axis=1
                    )
                else:
                    ds_df['type_label'] = ds_df['bodyId'].apply(lambda x: f"unknown_{x}")
                
                # Filter out unknown types for visualization (untyped neurons)
                ds_df_typed = ds_df[~ds_df['type_label'].str.startswith('unknown_')]
                if ds_df_typed.empty:
                    self._vprint(f"   ⚠️  {dataset}: No typed neurons for visualization")
                    continue
                ds_df = ds_df_typed
                
                # Get top N types using labeling_info if available (case-sensitive, properly sorted)
                # Otherwise fallback to avg_score
                # If type_filter is set, get ALL types first, then filter, then take top N
                get_all_types = type_filter is not None and type_filter
                
                if labeling_info is not None and not labeling_info.empty and 'dataset' in labeling_info.columns:
                    # Filter labeling_info for this dataset
                    ds_labeling = labeling_info[labeling_info['dataset'] == dataset].copy()
                    if not ds_labeling.empty:
                        # labeling_info is already sorted by quality (complete types first, then by min_score)
                        if get_all_types:
                            # Get all types for filtering
                            all_types = ds_labeling['type'].tolist()
                        else:
                            # Get only top N types
                            all_types = ds_labeling['type'].head(top_n).tolist()
                        self._vprint(f"   📋 Using labeling_info for {dataset}: {len(all_types)} candidate types")
                    else:
                        # Fallback if no labeling info for this dataset
                        all_types = self._get_top_types_fallback(ds_df, None if get_all_types else top_n)
                else:
                    # Fallback to score-based ranking
                    all_types = self._get_top_types_fallback(ds_df, None if get_all_types else top_n)
                
                if not all_types:
                    continue
                
                # Create list of (original_rank, type_name) tuples
                ranked_types = [(i, name) for i, name in enumerate(all_types, start=1)]
                
                # Apply type filter if specified
                if type_filter:
                    filtered_types = self._apply_type_filter(ranked_types, type_filter)
                    if not filtered_types:
                        self._vprint(f"   ⚠️  {dataset}: No types match filter criteria")
                        continue
                    self._vprint(f"   🔍 Filter applied: {len(filtered_types)}/{len(ranked_types)} types match")
                    # Take top N from filtered results (preserving original ranks)
                    top_ranked_types = filtered_types[:top_n]
                else:
                    top_ranked_types = ranked_types[:top_n]
                
                # Build neuron_layers as nested list (one sublist per type)
                # Use r{rank}_{type}_x{N} format for legend names (rank is ORIGINAL rank before filtering)
                neuron_layers = []
                layer_names = []
                
                for original_rank, type_name in top_ranked_types:
                    # Case-sensitive type matching
                    type_neurons = ds_df[ds_df['type_label'] == type_name]['bodyId'].tolist()
                    # Convert to int (handles strings, floats, and string floats like '123.0')
                    type_neurons = [_to_int_bodyid(n) for n in type_neurons]
                    
                    if len(type_neurons) > 0:
                        neuron_layers.append(type_neurons)
                        # Create legend name: r{rank}_{type}_x{N} - using ORIGINAL rank
                        n_neurons = len(type_neurons)
                        layer_names.append(f"r{original_rank}_{type_name}_x{n_neurons}")
                
                if not neuron_layers:
                    continue
                
                self._vprint(f"   📊 {dataset}: {len(neuron_layers)} types, {sum(len(l) for l in neuron_layers)} neurons")
            
            else:  # visualize_by == 'bodyId'
                # Get top N bodyIds by score, but group by type for visualization
                # Each type becomes a layer, but legend_mode='single' shows individual neurons
                score_col = 'score' if 'score' in ds_df.columns else None
                
                # Create type label (use 'unknown_{bodyId}' for untyped)
                if 'type' in ds_df.columns:
                    ds_df['type_label'] = ds_df.apply(
                        lambda row: f"unknown_{row['bodyId']}" if pd.isna(row['type']) or row['type'] == '' else row['type'],
                        axis=1
                    )
                else:
                    ds_df['type_label'] = ds_df['bodyId'].apply(lambda x: f"unknown_{x}")
                
                # Filter out unknown types for visualization (untyped neurons)
                ds_df_typed = ds_df[~ds_df['type_label'].str.startswith('unknown_')]
                if ds_df_typed.empty:
                    self._vprint(f"   ⚠️  {dataset}: No typed neurons for visualization")
                    continue
                ds_df = ds_df_typed
                
                # Sort by score - get ALL bodyIds first if filter is set, then filter, then limit
                if score_col:
                    ds_df_sorted = ds_df.sort_values(score_col, ascending=False)
                else:
                    ds_df_sorted = ds_df
                
                # Group ALL bodyIds by type for layer organization
                # Track the minimum rank (best score position) for each type for sorting
                type_to_bodyids = {}
                type_min_rank = {}  # Track minimum rank (best) for each type
                for rank_idx, (_, row) in enumerate(ds_df_sorted.iterrows(), start=1):
                    bodyid = row['bodyId']
                    # Convert to int (handles strings, floats, and string floats like '123.0')
                    bodyid_val = _to_int_bodyid(bodyid)
                    type_label = row['type_label']
                    
                    if type_label not in type_to_bodyids:
                        type_to_bodyids[type_label] = []
                        type_min_rank[type_label] = rank_idx  # First occurrence is minimum rank
                    type_to_bodyids[type_label].append(bodyid_val)
                
                # Sort types by their minimum rank (best ranked type first)
                sorted_types = sorted(type_to_bodyids.keys(), key=lambda t: type_min_rank[t])
                
                # Apply type filter if specified
                if type_filter:
                    # Create list of (rank, type_name) tuples for filtering
                    ranked_types = [(type_min_rank[t], t) for t in sorted_types]
                    filtered_types = self._apply_type_filter(ranked_types, type_filter)
                    if not filtered_types:
                        self._vprint(f"   ⚠️  {dataset}: No types match filter criteria")
                        continue
                    self._vprint(f"   🔍 Filter applied: {len(filtered_types)}/{len(ranked_types)} types match")
                    # Take top N from filtered results, extract type names
                    sorted_types = [t for _, t in filtered_types[:top_n]]
                else:
                    # No filter - just take top N types
                    sorted_types = sorted_types[:top_n]
                
                # Build neuron_layers: one layer per type, containing all bodyIds of that type
                # Use r{rank}_{type}_x{N} format for legend names (rank is ORIGINAL rank before filtering)
                neuron_layers = []
                layer_names = []
                
                for type_label in sorted_types:
                    bodyids = type_to_bodyids[type_label]
                    neuron_layers.append(bodyids)
                    # Layer name: r{rank}_{type}_x{N} where rank is the best (min) rank for this type
                    # This is the ORIGINAL rank from score ranking, preserved after filtering
                    rank = type_min_rank[type_label]
                    n_neurons = len(bodyids)
                    layer_names.append(f"r{rank}_{type_label}_x{n_neurons}")
                
                if not neuron_layers:
                    continue
                
                self._vprint(f"   📊 {dataset}: {len(neuron_layers)} types ({sum(len(l) for l in neuron_layers)} bodyIds)")
            
            # Verify bodyIds exist in local dataset before attempting visualization
            # This prevents the "No neurons matching" error from NeuPrint
            dataset_folder = dataset.replace(':', '_').replace('.', '_')
            local_neuron_df = self._load_neuron_df_for_dataset(dataset_folder)
            
            if local_neuron_df is None or local_neuron_df.empty:
                self._vprint(f"   ⚠️  Skipping visualization for {dataset}: local dataset not available")
                continue
            
            # Check if at least some bodyIds can be found
            # Local dataset is loaded with dtype={'bodyId': str}, so convert to strings for comparison
            sample_ids = [bid for layer in neuron_layers for bid in layer[:3]]  # Sample bodyIds
            
            # Convert sample_ids to strings for comparison (local df uses strings)
            # Use _to_int_bodyid first to normalize, then convert to string
            test_ids = [str(_to_int_bodyid(bid)) for bid in sample_ids]
            
            found_count = local_neuron_df[local_neuron_df['bodyId'].isin(test_ids)].shape[0]
            
            if found_count == 0:
                # Debug: show what we're looking for vs what exists
                self._vprint(f"   ⚠️  Skipping visualization for {dataset}: bodyIds not found in local dataset")
                self._vprint(f"      Looking for: {test_ids[:3]}")
                sample_local = local_neuron_df['bodyId'].head(3).tolist()
                self._vprint(f"      Sample local bodyIds: {sample_local}")
                self._vprint(f"      (This can happen if NeuronBridge matches come from a different dataset version)")
                continue
            
            try:
                # Determine brain_mesh based on dataset
                brain_mesh = 'template'  # Use template for better performance
                if 'vnc' in dataset.lower() or 'manc' in dataset.lower():
                    brain_mesh = 'template'
                elif 'cns' in dataset.lower():
                    brain_mesh = 'template'
                
                # Set skeleton_mesh_simplification based on dataset
                # FAFB/FlyWire needs more simplification (0.95) due to larger meshes
                # Hemibrain and male-cns use 0.9
                if 'fafb' in dataset.lower() or 'flywire' in dataset.lower():
                    skeleton_simplification = 0.95
                else:
                    skeleton_simplification = 0.9
                
                # Set legend_mode based on visualize_by mode
                legend_mode = 'layer' if visualize_by == 'type' else 'single'
                
                # Determine whether to show VNC mesh based on region and dataset
                # Show VNC when dataset is manc, male-cns, or region is VNC/All
                is_vnc_dataset = ('manc' in dataset.lower() or 'cns' in dataset.lower() 
                                  or 'vnc' in dataset.lower())
                show_vnc_mesh = is_vnc_dataset or self.region in ('VNC', 'All')
                
                # Determine export_views based on region setting
                # - region='VNC': only bottom view (shows VNC from below)
                # - region='All': both front and bottom views (shows brain + VNC)
                # - region='Brain': only front view (standard brain view)
                if self.region == 'VNC':
                    export_views = ['bottom']
                elif self.region == 'All':
                    export_views = ['front', 'bottom']
                else:  # 'Brain' or default
                    export_views = ['front']
                
                # Custom folder name: plot3d_{dataset_folder} (VisualizeSkeleton prepends 'plot3d_')
                custom_saveas = dataset_folder
                
                vs = VisualizeSkeleton(
                    dataset=dataset,
                    output_dir=output_path,
                    neuron_layers=neuron_layers,
                    custom_layer_names=layer_names,
                    saveas=custom_saveas,
                    include_timestamp=False,  # No timestamp for cleaner folder names
                    skip_synapse=True,
                    neuron_alpha=0.3,
                    skeleton_mode='tube',
                    legend_mode=legend_mode,  # 'layer' for type, 'single' for bodyId
                    brain_mesh=brain_mesh,
                    vnc_mesh=show_vnc_mesh,  # Show VNC for manc, male-cns, or VNC region
                    export_views=export_views,  # Region-based view selection
                    skeleton_mesh_simplification=skeleton_simplification,
                    roi_mesh_simplification=0.95,
                    cache_neurons=True,
                    show_fig=False,
                    verbose='full',  # Full verbose to see simplification logs
                )
                vs.plot_neurons()
                
                # Generate individual profiles if requested
                # Use region-appropriate views for individual profile generation
                if generate_individual_profiles:
                    # Use same views as export_views for consistency
                    profile_views = export_views
                    vs.plot_individuals(
                        output_format='png',
                        views=profile_views,
                        scale=3,
                        pdf_images_per_page=pdf_images_per_page,
                        pdf_title=f"{source_line} - {dataset}",
                        summary_format=generate_individual_profiles
                    )
                
                self._vprint(f"   ✅ Visualization saved to: {vs.save_folder}")
            except Exception as e:
                self._vprint(f"   ⚠️  Visualization failed for {dataset}: {e}")
    
    # =========================================================================
    # Batch Processing Methods (for simplified script usage)
    # =========================================================================
    
    def find_neurons_batch(
        self,
        line_names: Union[str, List[str]],
        top_n: int = -1,
        match_type: Optional[str] = None,
        output_dir: Optional[str] = None,
        visualize_top_n: int = 0,
        visualize_by: str = 'type',
        visualize_per_dataset: bool = True,
        generate_individual_profiles: Union[bool, List[str]] = None,
        pdf_images_per_page: Tuple[int, int] = (4, 3),
        type_filter: Optional[Dict[str, Union[str, List[str]]]] = None,
        datasets_to_visualize: Union[str, List[str]] = 'all',
    ) -> pd.DataFrame:
        """
        Find EM neurons for multiple driver lines with automatic saving.
        
        This is a convenience method that processes multiple lines, adds source
        information, and optionally saves results to files.
        
        Parameters
        ----------
        line_names : str or list
            Driver line name(s). Can be comma-separated string or list.
        top_n : int
            Maximum matches per line. Default: -1 (all matches)
        match_type : str, optional
            Match algorithm: 'cds', 'pppm', or 'both'. 
            If None, uses self.match_type. Default: None
        output_dir : str, optional
            Directory to save results. If provided, saves individual and combined CSVs.
        visualize_top_n : int
            Visualize top N types/bodyIds per dataset using 3D skeleton visualization.
            Set to 0 to disable (default). Requires VisualizeSkeleton module.
        visualize_by : str
            How to organize visualization: 'type' or 'bodyId'.
            - 'type': Group neurons by type (legend_mode='layer')
            - 'bodyId': Show individual neurons (legend_mode='single')
            Default: 'type'
        visualize_per_dataset : bool
            If True (default), create separate visualizations per dataset.
            If False, combine all datasets in one visualization.
        generate_individual_profiles : list of str, bool, or None
            Formats to generate: ['pdf'], ['pptx'], or ['pdf', 'pptx']
            Set to False or None to disable generation. Default: None.
        pdf_images_per_page : tuple
            (columns, rows) for PDF layout. Default: (4, 3).
        type_filter : dict, optional
            Filter neuron types for visualization by name pattern.
            Keys: 'contains', 'startswith', 'endswith', 'regex'
            Values: str or list of str patterns
            Multiple patterns within same key use OR logic.
            Multiple keys use AND logic.
            Ranks are preserved from original ranking before filtering.
            Example: {'contains': 'DN', 'startswith': ['IN', 'DN']}
        datasets_to_visualize : str or list, default 'all'
            Constrain which datasets to visualize.
            - 'all' or None: Visualize all datasets found in results
            - List of dataset names: Only visualize specified datasets
            - Single dataset name: Only visualize that dataset
            
        Returns
        -------
        pd.DataFrame
            Combined DataFrame with all results, including 'source_line' column.
            
        Output Files
        ------------
        When output_dir is specified:
        - {line}_neurons.csv: All matched neurons for the line
        - {line}_{dataset}_neurons.csv: Neurons categorized by dataset
        - {line}_{dataset}_types.csv: Type summary with labeled_N and typed_N_in_dataset
        - all_neurons.csv: Combined results from all searches
            
        Example
        -------
        >>> nbf = NeuronBridgeFinder()
        >>> results = nbf.find_neurons_batch('LH173,VT037867', output_dir='./output')
        >>> # With type filtering
        >>> results = nbf.find_neurons_batch(
        ...     'SS29633',
        ...     visualize_top_n=20,
        ...     type_filter={'contains': 'DN'},
        ...     datasets_to_visualize=['manc:v1.0']
        ... )
        """
        # Use class-level match_type if not specified
        if match_type is None:
            match_type = self.match_type
        else:
            # Validate and normalize match_type using helper
            match_type = self._validate_match_type(match_type)
            
        # Parse line names
        if isinstance(line_names, str):
            lines = [ln.strip() for ln in line_names.split(',') if ln.strip()]
        else:
            lines = list(line_names)
        
        if not lines:
            self._vprint("❌ No line names provided")
            return pd.DataFrame()
        
        self._vprint(f"🔍 Finding neurons for {len(lines)} line(s)")
        
        # Create output directory if needed
        output_path = None
        if output_dir:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            # Add line names to folder name
            line_info = '_'.join(lines[:3])  # First 3 lines
            if len(lines) > 3:
                line_info += '_etc'
            # Sanitize folder name (remove special characters)
            line_info = ''.join(c if c.isalnum() or c in '-_' else '_' for c in line_info)
            output_path = os.path.join(output_dir, f'findneuron_{line_info}_{timestamp}')
            os.makedirs(output_path, exist_ok=True)
            self._vprint(f"   Output: {output_path}")
            
            # Save parameters for reproducibility
            self._save_parameters(
                output_path=output_path,
                function_name='find_neurons_batch',
                function_params={
                    'line_names': lines,
                    'top_n': top_n,
                    'match_type': match_type,
                    'visualize_top_n': visualize_top_n,
                    'visualize_by': visualize_by,
                    'visualize_per_dataset': visualize_per_dataset,
                    'generate_individual_profiles': generate_individual_profiles,
                    'pdf_images_per_page': pdf_images_per_page,
                    'type_filter': type_filter,
                    'datasets_to_visualize': datasets_to_visualize
                }
            )
            self._vprint(f"   💾 Parameters: parameters.json")
        
        # Process each line
        all_results = []
        
        # Phase 1: Check cache and identify what needs to be fetched
        if HAS_TQDM and self.verbose and len(lines) > 1:
            from tqdm import tqdm as tqdm_progress
            
            self._vprint("\n📦 Checking cache...")
            cached_lines = []
            uncached_lines = []
            
            cache_pbar = tqdm_progress(
                lines,
                desc="   💾 Loading cache",
                unit="line",
                bar_format='{desc}: {percentage:3.0f}%|{bar:30}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]',
                ncols=110,
                position=0,
                leave=False
            )
            
            for line_name in cache_pbar:
                # Check cache status
                region_key = self.region if self.region else 'all'
                max_imgs_key = self.max_api_images_per_line if self.max_api_images_per_line > 0 else 'all'
                cache_key = f"{line_name}_{match_type}_{region_key}_{max_imgs_key}"
                cached_data = self._load_from_cache('line_to_neuron', cache_key)
                
                if cached_data is not None:
                    cached_lines.append((line_name, cached_data))
                    cache_pbar.set_postfix_str(f"✓ {line_name}")
                else:
                    uncached_lines.append(line_name)
            
            cache_pbar.close()
            
            if cached_lines:
                self._vprint(f"   ✓ Loaded {len(cached_lines)} line result(s) from cache")
            if uncached_lines:
                self._vprint(f"   🌐 Need to fetch {len(uncached_lines)} from API")
            
            # Suppress dataset loading messages during cache processing
            # (datasets will be loaded explicitly before visualization)
            old_suppress = getattr(self, '_suppress_loading_msgs', False)
            self._suppress_loading_msgs = True
            
            # Process cached results first
            for line_name, cached_data in cached_lines:
                cached_data = cached_data.copy()
                cached_data['source_line'] = line_name
                all_results.append(cached_data)
                
                # Save individual cached results if output_path specified
                if output_path:
                    output_file = os.path.join(output_path, f'{line_name}_neurons.csv')
                    cached_data.to_csv(output_file, index=False)
                    
                    # Save dataset-categorized files
                    if 'dataset' in cached_data.columns:
                        self._save_dataset_categorized_files(
                            cached_data, line_name, output_path, verbose=False
                        )
            
            # Restore loading message setting
            self._suppress_loading_msgs = old_suppress
            
            # Phase 2: Fetch uncached data with progress bar
            if uncached_lines:
                self._vprint("\n🌐 Fetching new data...")
                fetch_pbar = tqdm_progress(
                    uncached_lines,
                    desc="   🔄 Fetching",
                    unit="line",
                    bar_format='{desc}: {percentage:3.0f}%|{bar:30}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]',
                    ncols=110,
                    position=0,
                    leave=False
                )
                
                for line_name in fetch_pbar:
                    fetch_pbar.set_postfix_str(line_name[:20])
                    
                    try:
                        neurons_df = self.line_to_neuron(
                            line_name,
                            top_n=top_n,
                            match_type=match_type
                        )
                        
                        if neurons_df.empty:
                            continue
                        
                        # Add source line info
                        neurons_df = neurons_df.copy()
                        neurons_df['source_line'] = line_name
                        all_results.append(neurons_df)
                        
                        # Save individual results
                        if output_path:
                            output_file = os.path.join(output_path, f'{line_name}_neurons.csv')
                            neurons_df.to_csv(output_file, index=False)
                            
                            # Save dataset-categorized files
                            if 'dataset' in neurons_df.columns:
                                self._save_dataset_categorized_files(
                                    neurons_df, line_name, output_path, verbose=False
                                )
                        
                    except Exception as e:
                        pass  # Silent in progress mode
                
                fetch_pbar.close()
                self._vprint(f"   ✓ Fetched {len(uncached_lines)} from API")
        
        else:
            # Single line or no progress bar - original behavior
            for idx, line_name in enumerate(lines):
                self._vprint(f"\n📋 Processing: {line_name}")
                
                try:
                    neurons_df = self.line_to_neuron(
                        line_name,
                        top_n=top_n,
                        match_type=match_type
                    )
                    
                    if neurons_df.empty:
                        self._vprint(f"   ⚠️ No matching neurons found")
                        continue
                    
                    # Add source line info
                    neurons_df = neurons_df.copy()
                    neurons_df['source_line'] = line_name
                    all_results.append(neurons_df)
                    
                    self._vprint(f"   ✅ Found {len(neurons_df)} matching neurons")
                    
                    # Show dataset distribution
                    if 'dataset' in neurons_df.columns:
                        datasets = neurons_df['dataset'].value_counts()
                        for ds, count in datasets.items():
                            self._vprint(f"      {ds}: {count}")
                    
                    # Save individual results
                    if output_path:
                        output_file = os.path.join(output_path, f'{line_name}_neurons.csv')
                        neurons_df.to_csv(output_file, index=False)
                        self._vprint(f"   💾 Saved: {output_file}")
                        
                        # Save dataset-categorized files
                        if 'dataset' in neurons_df.columns:
                            self._save_dataset_categorized_files(
                                neurons_df, line_name, output_path, verbose=True
                            )
                        
                except Exception as e:
                    self._vprint(f"   ❌ Error: {e}")
        
        # Combine results
        if all_results:
            combined_df = pd.concat(all_results, ignore_index=True)
            
            # Show summary statistics
            self._vprint(f"\n📊 Summary:")
            self._vprint(f"   Total neurons found: {len(combined_df)}")
            self._vprint(f"   From {len(lines)} line(s), {len(all_results)} had matches")
            
            # Dataset distribution
            if 'dataset' in combined_df.columns:
                self._vprint(f"   Dataset distribution:")
                datasets = combined_df['dataset'].value_counts()
                for ds, count in datasets.items():
                    self._vprint(f"      {ds}: {count}")
            
            # Save combined results
            if output_path:
                combined_file = os.path.join(output_path, 'all_neurons.csv')
                combined_df.to_csv(combined_file, index=False)
                self._vprint(f"\n💾 Combined results: {combined_file}")
                
                # Create labeling distribution visualization
                self.visualize_labeling_distribution(
                    data=combined_df,
                    output_path=output_path,
                    score_column='score',
                    label_column='type',
                    title=f"Labeling Score Distribution ({len(lines)} line(s))",
                    group_by='source_line' if 'source_line' in combined_df.columns and len(lines) > 1 else None
                )
            
            # Visualize top N types per dataset if requested
            if visualize_top_n > 0 and output_path:
                # For multiple lines, visualize each line separately
                if len(lines) > 1:
                    self._vprint(f"\n🎨 Visualizing each line separately (multiple lines detected)...")
                    for line_name in lines:
                        line_df = combined_df[combined_df['source_line'] == line_name]
                        if line_df.empty:
                            continue
                        
                        # Create line-specific output folder
                        line_output_path = os.path.join(output_path, f'viz_{line_name}')
                        os.makedirs(line_output_path, exist_ok=True)
                        
                        self._vprint(f"\n   📍 {line_name}: {len(line_df)} neurons")
                        
                        self._visualize_top_types(
                            combined_df=line_df,
                            top_n=visualize_top_n,
                            output_path=line_output_path,
                            per_dataset=visualize_per_dataset,
                            source_line=line_name,
                            visualize_by=visualize_by,
                            generate_individual_profiles=generate_individual_profiles,
                            pdf_images_per_page=pdf_images_per_page,
                            type_filter=type_filter,
                            datasets_to_visualize=datasets_to_visualize,
                        )
                    
                    # Recommend colabeling analysis for multiple lines
                    self._vprint(f"\n" + "="*70)
                    self._vprint(f"💡 TIP: For {len(lines)} driver lines, consider running co-labeling analysis")
                    self._vprint(f"   to find shared neuron types and understand overlap patterns:")
                    self._vprint(f"")
                    self._vprint(f"   >>> nbf.analyze_colabeling(")
                    self._vprint(f"   ...     lines={lines},")
                    self._vprint(f"   ...     output_dir='{output_dir}',")
                    self._vprint(f"   ...     visualize_top_n={visualize_top_n},")
                    self._vprint(f"   ... )")
                    self._vprint(f"="*70)
                else:
                    # Single line - original behavior
                    self._visualize_top_types(
                        combined_df=combined_df,
                        top_n=visualize_top_n,
                        output_path=output_path,
                        per_dataset=visualize_per_dataset,
                        source_line=lines[0],
                        visualize_by=visualize_by,
                        generate_individual_profiles=generate_individual_profiles,
                        pdf_images_per_page=pdf_images_per_page,
                        type_filter=type_filter,
                        datasets_to_visualize=datasets_to_visualize,
                    )
            
            return combined_df
        
        self._vprint(f"\n⚠️ No neurons found for any of the {len(lines)} line(s)")
        return pd.DataFrame()
    
    def analyze_colabeling(
        self,
        lines: Union[str, List[str]],
        match_type: Optional[str] = None,
        top_n_neurons: int = -1,
        similarity_methods: Union[str, List[str]] = ['jaccard', 'weighted_jaccard'],
        output_dir: Optional[str] = None,
        generate_report: bool = True,
        visualize: bool = True,
        visualize_top_n: int = 0,
        generate_individual_profiles: Union[bool, List[str]] = None,
        pdf_images_per_page: Tuple[int, int] = (3, 2),
        min_score: float = 20000.0,
        min_type_avg_score: float = 10000.0
    ) -> Dict[str, Any]:
        """
        Analyze co-labeling patterns among given driver lines.
        
        This method performs a comprehensive co-labeling analysis to understand
        how different driver lines overlap in their neuron labeling patterns.
        
        Parameters
        ----------
        lines : str or list of str
            Driver line names to analyze. Can be:
            - Single line: 'LH173'
            - Multiple as string: 'LH173,VT037867,SS00731'
            - Multiple as list: ['LH173', 'VT037867', 'SS00731']
        match_type : str, optional
            Match algorithm: 'cds', 'pppm', or 'both'.
            If None, uses self.match_type. Default: None
        top_n_neurons : int
            Number of top neuron matches to consider per line. Default: -1 (all)
        similarity_methods : str or list of str
            Similarity method(s) for co-labeling matrix:
            - 'jaccard': Binary Jaccard similarity (presence/absence of types)
            - 'weighted_jaccard': Jaccard weighted by match scores
            - 'rank_correlation': Spearman correlation of type rankings
            Default: ['jaccard', 'weighted_jaccard']
        output_dir : str, optional
            Directory to save results. If None, returns results without saving.
        generate_report : bool
            Generate a comprehensive HTML/text report. Default: True
        visualize : bool
            Generate heatmap visualizations. Default: True
        visualize_top_n : int
            Visualize top N types per dataset using 3D skeleton. Default: 0 (disabled)
        generate_individual_profiles : list of str, bool, or None
            Formats to generate: ['pdf'], ['pptx'], or ['pdf', 'pptx']
            Set to False or None to disable generation. Default: None
        pdf_images_per_page : tuple
            (columns, rows) for PDF layout. Default: (3, 2)
        min_score : float
            Score threshold for visualization marker only (labeling distribution plots).
            Does NOT filter data in expression matrix - all neurons are included.
            Default: 30000.0
        min_type_avg_score : float
            Minimum average score threshold for types in similarity matrix. Default: 20000.0
            Types with average score < threshold may be excluded from clustering.
            Note: Expression matrix includes ALL types regardless of this threshold.
            
        Returns
        -------
        dict
            Dictionary containing:
            - 'expression_matrix': pd.DataFrame - Type × Line expression matrix (scores)
              Types are prefixed with dataset abbreviation: {ABBREV}_{type}
              (HEMI=hemibrain, MCNS=male-cns, FAFB=FlyWire FAFB)
            - 'labeling_info': pd.DataFrame - Case-sensitive types with dataset column
            - 'colabeling_matrices': Dict[str, pd.DataFrame] - Similarity matrices per method
            - 'line_neurons': Dict[str, pd.DataFrame] - Neurons per line with scores
            - 'line_summary': pd.DataFrame - Summary stats per line
            - 'report_path': str - Path to generated report (if generate_report=True)
            
        Output Files (when output_dir is provided)
        ------------------------------------------
        - expression_matrix.csv: Type × Line matrix with match scores
          (types prefixed with dataset abbreviation: HEMI_, MCNS_, FAFB_)
        - labeling_info.csv: Case-sensitive types with dataset column for per-dataset filtering
        - expression_matrix.html: Interactive heatmap visualization
        - colabeling_matrix_{method}.csv: Line × Line similarity matrix
        - colabeling_matrix_{method}.html: Interactive heatmap
        - line_labeled_neurons/{line}_neurons.csv: Per-line neuron details
        - line_summary.csv: Summary statistics per line
        - colabeling_report.html: Comprehensive analysis report
        
        Example
        -------
        >>> nbf = NeuronBridgeFinder()
        >>> results = nbf.analyze_colabeling(
        ...     lines=['LH173', 'VT037867', 'SS00731'],
        ...     output_dir='./colabel_analysis'
        ... )
        >>> print(results['line_summary'])
        """
        from datetime import datetime
        
        # Parse lines input
        if isinstance(lines, str):
            line_list = [l.strip() for l in lines.split(',') if l.strip()]
        else:
            line_list = list(lines)
        
        if len(line_list) < 2:
            self._vprint("❌ At least 2 lines are required for co-labeling analysis")
            return {}
        
        # Use class-level match_type if not specified
        if match_type is None:
            match_type = self.match_type
        else:
            # Validate and normalize match_type using helper
            match_type = self._validate_match_type(match_type)
            
        # Normalize similarity_methods
        if isinstance(similarity_methods, str):
            similarity_methods = [similarity_methods]
        
        # Validate similarity_methods using helper
        similarity_methods = [self._validate_similarity_method(m) for m in similarity_methods]
        
        self._vprint(f"\n{'='*60}")
        self._vprint(f"🔬 Co-Labeling Analysis")
        self._vprint('='*60)
        self._vprint(f"   Lines: {len(line_list)}")
        self._vprint(f"   Match type: {match_type}")
        self._vprint(f"   Top neurons per line: {top_n_neurons}")
        self._vprint(f"   Similarity methods: {', '.join(similarity_methods)}")
        self._vprint(f"   Min neuron score: {min_score:,.0f}")
        self._vprint(f"   Min type avg score: {min_type_avg_score:,.0f}")
        
        # Create output directory if needed
        output_path = None
        if output_dir:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            # Create subfolder
            lines_info = '_'.join(line_list[:3])
            if len(line_list) > 3:
                lines_info += f'_etc{len(line_list)}'
            # Sanitize folder name
            lines_info = ''.join(c if c.isalnum() or c in '-_' else '_' for c in lines_info)
            output_path = os.path.join(output_dir, f'colabel_{lines_info}_{timestamp}')
            os.makedirs(output_path, exist_ok=True)
            self._vprint(f"   Output: {output_path}")
            
            # Save parameters for reproducibility
            self._save_parameters(
                output_path=output_path,
                function_name='analyze_colabeling',
                function_params={
                    'lines': line_list,
                    'match_type': match_type,
                    'top_n_neurons': top_n_neurons,
                    'similarity_methods': similarity_methods,
                    'generate_report': generate_report,
                    'visualize': visualize,
                    'visualize_top_n': visualize_top_n,
                    'generate_individual_profiles': generate_individual_profiles,
                    'pdf_images_per_page': pdf_images_per_page,
                    'min_score': min_score,
                    'min_type_avg_score': min_type_avg_score
                }
            )
            self._vprint(f"   💾 Parameters: parameters.json")
        
        results = {
            'expression_matrix': None,
            'labeling_info': None,
            'colabeling_matrices': {},
            'line_neurons': {},
            'line_summary': None,
            'report_path': None
        }
        
        # Step 1: Fetch neurons for each line and build expression matrix
        self._vprint(f"\n📊 Step 1: Fetching neurons and building expression matrix...")
        
        # Use _calculate_mutual_information which already does what we need
        mi_df, expression_df, line_neurons_dict, labeling_info = self._calculate_mutual_information(
            lines=line_list,
            queried_types=[],  # No specific queried types for pure co-labeling
            match_type=match_type,
            top_n=top_n_neurons,
            output_path=output_path,
            min_score=min_score,
            min_type_avg_score=min_type_avg_score
        )
        
        results['line_neurons'] = line_neurons_dict
        results['labeling_info'] = labeling_info
        
        if expression_df.empty:
            self._vprint("   ⚠️ No expression data found for any lines")
            return results
        
        # Sort expression matrix by co-labeling quality using shared helper
        expression_transposed = self._sort_expression_matrix(expression_df, as_types_rows=True)
        
        n_lines = len(expression_transposed.columns)
        nonzero_count = (expression_transposed > 0).sum(axis=1)
        n_complete = (nonzero_count == n_lines).sum()
        self._vprint(f"   Sorted {len(expression_transposed)} types: {n_complete} in all lines, {len(expression_transposed) - n_complete} partial")
        
        results['expression_matrix'] = expression_transposed
        
        # Combine all line neurons into a single DataFrame for dataset splitting and visualization
        combined_neurons_df = pd.DataFrame()
        if line_neurons_dict:
            dfs = []
            for line_name, neurons_df in line_neurons_dict.items():
                if not neurons_df.empty:
                    neurons_copy = neurons_df.copy()
                    neurons_copy['source_line'] = line_name
                    dfs.append(neurons_copy)
            if dfs:
                combined_neurons_df = pd.concat(dfs, ignore_index=True)
        
        # Save expression matrix, labeling_info, and per-line neuron details
        if output_path:
            expr_csv = os.path.join(output_path, 'expression_matrix.csv')
            expression_transposed.to_csv(expr_csv)
            self._vprint(f"   💾 Expression matrix: {expr_csv}")
            
            # Save labeling_info.csv (case-sensitive types with dataset column)
            if not labeling_info.empty:
                labeling_csv = os.path.join(output_path, 'labeling_info.csv')
                labeling_info.to_csv(labeling_csv, index=False)
                self._vprint(f"   💾 Labeling info: {labeling_csv}")
            
            # Save per-line neuron details (simple version, split by dataset below)
            if line_neurons_dict:
                neurons_dir = os.path.join(output_path, 'line_labeled_neurons')
                os.makedirs(neurons_dir, exist_ok=True)
                
                for line_name, neurons_df in line_neurons_dict.items():
                    if not neurons_df.empty:
                        safe_name = line_name.replace('/', '_').replace('\\', '_')
                        # Save the combined line neurons file
                        neurons_csv = os.path.join(neurons_dir, f'{safe_name}_neurons.csv')
                        neurons_df.to_csv(neurons_csv, index=False)
                        
                        # Split by dataset and save dataset-specific files with type summaries
                        self._save_dataset_categorized_files(neurons_df, safe_name, neurons_dir, verbose=False)
                
                self._vprint(f"   💾 Line neurons: {neurons_dir}/ ({len(line_neurons_dict)} lines, split by dataset)")
        
        # Step 2: Build co-labeling matrices
        self._vprint(f"\n📊 Step 2: Building co-labeling matrices...")
        
        for method in similarity_methods:
            self._vprint(f"   Computing {method} similarity...")
            
            matrix, line_type_sets = self._build_colabeling_matrix(
                lines=line_list,
                match_type=match_type,
                top_n=top_n_neurons,
                similarity_method=method,
                min_score=min_score,
                min_type_avg_score=min_type_avg_score
            )
            
            results['colabeling_matrices'][method] = matrix
            
            # Save matrix CSV
            if output_path:
                csv_filename = f'colabeling_matrix_{method}.csv'
                csv_path = os.path.join(output_path, csv_filename)
                matrix.to_csv(csv_path)
                self._vprint(f"   💾 {csv_filename}")
                
                # Create visualization
                if visualize:
                    method_titles = {
                        'jaccard': 'Co-Labeling Matrix (Binary Jaccard)',
                        'weighted_jaccard': 'Co-Labeling Matrix (Weighted Jaccard)',
                        'rank_correlation': 'Co-Labeling Matrix (Rank Correlation)'
                    }
                    title = method_titles.get(method, f'Co-Labeling Matrix ({method})')
                    html_filename = f'colabeling_matrix_{method}.html'
                    self.visualize_colabeling_matrix(
                        co_labeling_matrix=matrix,
                        output_path=output_path,
                        title=title,
                        color_scale='purple',
                        filename=html_filename
                    )
        
        # Visualize expression matrix (use already sorted expression_transposed.T to get Lines × Types)
        if output_path and visualize:
            # Pass the sorted matrix (transpose back to Lines × Types format expected by visualize_expression_matrix)
            self.visualize_expression_matrix(
                expression_df=expression_transposed.T,
                output_path=output_path,
                queried_types=[],
                title=f"Expression Matrix ({len(line_list)} Lines × Types)"
            )
            
            # Create merged dataset version (same types across datasets combined)
            self.visualize_expression_matrix_merged(
                expression_df=expression_transposed.T,
                output_path=output_path,
                queried_types=[],
                title=f"Expression Matrix ({len(line_list)} Lines × Types) - Merged Datasets",
                aggregation='max'  # Use max score across datasets for same type
            )
            
            # Visualize labeling distribution (mountain-shaped histogram)
            self.visualize_colabeling_distribution(
                line_neurons_dict=line_neurons_dict,
                output_path=output_path,
                min_score=min_score,
                title=f"Labeling Score Distribution ({len(line_list)} Lines)"
            )
        
        # Step 3: Calculate line summary statistics
        self._vprint(f"\n📊 Step 3: Computing line statistics...")
        
        line_summary_data = []
        
        # Get the primary similarity matrix (weighted_jaccard if available)
        primary_matrix = results['colabeling_matrices'].get(
            'weighted_jaccard', 
            results['colabeling_matrices'].get('jaccard', pd.DataFrame())
        )
        
        # Calculate sparsity metrics
        sparsity_scores = {}
        if not primary_matrix.empty:
            sparsity_scores = self._calculate_colabeling_sparsity(primary_matrix)
        
        for line_name in line_list:
            neurons_df = line_neurons_dict.get(line_name, pd.DataFrame())
            
            # Basic stats
            n_neurons = len(neurons_df)
            n_types = neurons_df['type'].nunique() if not neurons_df.empty else 0
            
            # Score statistics
            mean_score = neurons_df['score'].mean() if not neurons_df.empty else 0.0
            max_score = neurons_df['score'].max() if not neurons_df.empty else 0.0
            
            # Half-max score threshold
            half_max_score = max_score / 2.0 if max_score > 0 else 0.0
            
            # n_neurons_HMS: neurons with score >= half_max_score
            n_neurons_HMS = len(neurons_df[neurons_df['score'] >= half_max_score]) if not neurons_df.empty else 0
            
            # n_types_HMS: unique types with score >= half_max_score
            if not neurons_df.empty:
                hms_neurons = neurons_df[neurons_df['score'] >= half_max_score]
                n_types_HMS = hms_neurons['type'].nunique() if not hms_neurons.empty else 0
            else:
                n_types_HMS = 0
            
            # n_neurons_MS: neurons at max score (within 0.1% tolerance)
            if not neurons_df.empty and max_score > 0:
                ms_threshold = max_score * 0.999  # Allow 0.1% tolerance
                n_neurons_MS = len(neurons_df[neurons_df['score'] >= ms_threshold])
            else:
                n_neurons_MS = 0
            
            # n_types_MS: unique types at max score
            if not neurons_df.empty and max_score > 0:
                ms_threshold = max_score * 0.999
                ms_neurons = neurons_df[neurons_df['score'] >= ms_threshold]
                n_types_MS = ms_neurons['type'].nunique() if not ms_neurons.empty else 0
            else:
                n_types_MS = 0
            
            # Qf: Quality factor = max_score / n_types_HMS (higher = more specific/selective)
            Qf = max_score / n_types_HMS if n_types_HMS > 0 else 0.0
            
            # Sparsity metrics
            sparsity_data = sparsity_scores.get(line_name, {})
            
            line_summary_data.append({
                'line': line_name,
                'n_neurons': n_neurons,
                'n_types': n_types,
                'mean_score': round(mean_score, 4),
                'max_score': round(max_score, 4),
                'n_neurons_HMS': n_neurons_HMS,
                'n_types_HMS': n_types_HMS,
                'n_neurons_MS': n_neurons_MS,
                'n_types_MS': n_types_MS,
                'Qf': round(Qf, 2),
                'colabel_sparsity': round(sparsity_data.get('colabel_sparsity', 0), 4),
            })
        
        line_summary = pd.DataFrame(line_summary_data)
        results['line_summary'] = line_summary
        
        # Save line summary
        if output_path:
            summary_csv = os.path.join(output_path, 'line_summary.csv')
            line_summary.to_csv(summary_csv, index=False)
            self._vprint(f"   💾 Line summary: {summary_csv}")
        
        # Step 4: Generate comprehensive report
        if generate_report and output_path:
            self._vprint(f"\n📝 Step 4: Generating analysis report...")
            report_path = self._generate_colabeling_report(
                results=results,
                line_list=line_list,
                match_type=match_type,
                top_n_neurons=top_n_neurons,
                output_path=output_path
            )
            results['report_path'] = report_path
        
        # Step 5: Visualize top N types per dataset (3D skeleton)
        if visualize_top_n > 0 and output_path and not combined_neurons_df.empty:
            self._vprint(f"\n🎨 Step 5: Visualizing top {visualize_top_n} types per dataset...")
            
            # Create a source_line label for folder naming
            lines_label = '_'.join(line_list[:3])
            if len(line_list) > 3:
                lines_label += f'_etc{len(line_list)}'
            
            self._visualize_top_types(
                combined_df=combined_neurons_df,
                top_n=visualize_top_n,
                output_path=output_path,
                per_dataset=True,
                source_line=lines_label,
                visualize_by='type',  # Default to type-based visualization
                generate_individual_profiles=generate_individual_profiles,
                pdf_images_per_page=pdf_images_per_page,
                labeling_info=labeling_info,  # Pass case-sensitive type info for per-dataset filtering
            )
        
        # Summary
        self._vprint(f"\n{'='*60}")
        self._vprint(f"✅ Co-Labeling Analysis Complete!")
        self._vprint('='*60)
        self._vprint(f"   Lines analyzed: {len(line_list)}")
        self._vprint(f"   Total unique types: {len(expression_transposed)}")
        self._vprint(f"   Output: {output_path}")
        
        return results
    
    def _generate_colabeling_report(
        self,
        results: Dict[str, Any],
        line_list: List[str],
        match_type: str,
        top_n_neurons: int,
        output_path: str
    ) -> str:
        """Generate a comprehensive HTML report for co-labeling analysis."""
        from datetime import datetime
        
        report_path = os.path.join(output_path, 'colabeling_report.html')
        
        expression_matrix = results.get('expression_matrix', pd.DataFrame())
        line_summary = results.get('line_summary', pd.DataFrame())
        colabeling_matrices = results.get('colabeling_matrices', {})
        
        # Get top co-labeling pairs
        top_pairs = []
        if 'weighted_jaccard' in colabeling_matrices:
            matrix = colabeling_matrices['weighted_jaccard']
            for i, line_i in enumerate(matrix.index):
                for j, line_j in enumerate(matrix.columns):
                    if i < j:  # Only upper triangle
                        similarity = matrix.iloc[i, j]
                        if similarity > 0.1:  # Only significant pairs
                            top_pairs.append({
                                'line1': line_i,
                                'line2': line_j,
                                'similarity': similarity
                            })
            top_pairs = sorted(top_pairs, key=lambda x: x['similarity'], reverse=True)[:20]
        
        # Find most specific lines (high sparsity)
        most_specific = []
        if not line_summary.empty:
            sorted_by_sparsity = line_summary.sort_values('colabel_sparsity', ascending=False)
            most_specific = sorted_by_sparsity.head(10).to_dict('records')
        
        # Build HTML report
        html_content = f'''<!DOCTYPE html>
<html>
<head>
    <title>Co-Labeling Analysis Report</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; background-color: #f5f5f5; }}
        .container {{ max-width: 1200px; margin: 0 auto; background: white; padding: 30px; border-radius: 10px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); }}
        h1 {{ color: #2c3e50; border-bottom: 3px solid #3498db; padding-bottom: 10px; }}
        h2 {{ color: #34495e; margin-top: 30px; }}
        h3 {{ color: #7f8c8d; }}
        .summary-box {{ background: #ecf0f1; padding: 20px; border-radius: 8px; margin: 20px 0; }}
        .metric {{ display: inline-block; margin: 10px 20px 10px 0; }}
        .metric-value {{ font-size: 24px; font-weight: bold; color: #2980b9; }}
        .metric-label {{ font-size: 14px; color: #7f8c8d; }}
        table {{ border-collapse: collapse; width: 100%; margin: 15px 0; }}
        th, td {{ border: 1px solid #ddd; padding: 10px; text-align: left; }}
        th {{ background-color: #3498db; color: white; }}
        tr:nth-child(even) {{ background-color: #f9f9f9; }}
        tr:hover {{ background-color: #e8f4fc; }}
        .highlight {{ background-color: #fff3cd; }}
        .file-link {{ color: #3498db; text-decoration: none; }}
        .file-link:hover {{ text-decoration: underline; }}
        .pair-card {{ background: #e8f6f3; padding: 15px; margin: 10px 0; border-radius: 8px; border-left: 4px solid #1abc9c; }}
        .similarity-bar {{ height: 20px; background: #3498db; border-radius: 4px; }}
        .timestamp {{ color: #95a5a6; font-size: 12px; }}
    </style>
</head>
<body>
    <div class="container">
        <h1>🔬 Co-Labeling Analysis Report</h1>
        <p class="timestamp">Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
        
        <div class="summary-box">
            <div class="metric">
                <div class="metric-value">{len(line_list)}</div>
                <div class="metric-label">Lines Analyzed</div>
            </div>
            <div class="metric">
                <div class="metric-value">{len(expression_matrix)}</div>
                <div class="metric-label">Unique Types</div>
            </div>
            <div class="metric">
                <div class="metric-value">{match_type.upper()}</div>
                <div class="metric-label">Match Type</div>
            </div>
            <div class="metric">
                <div class="metric-value">{top_n_neurons}</div>
                <div class="metric-label">Top N Neurons/Line</div>
            </div>
        </div>
        
        <h2>📊 Line Summary</h2>
        <p><strong>Metrics Explanation:</strong></p>
        <ul>
            <li><strong>Neurons</strong>: Total number of neurons matched by this line</li>
            <li><strong>Types</strong>: Total number of unique neuron types labeled</li>
            <li><strong>Max Score</strong>: Highest NeuronBridge match score (50000 = perfect match)</li>
            <li><strong>n_HMS</strong>: Neurons/Types above Half Max Score (score ≥ max_score/2)</li>
            <li><strong>n_MS</strong>: Neurons/Types at Max Score (within 0.1% of max_score)</li>
            <li><strong>Qf</strong>: Quality Factor = max_score / n_types_HMS. Higher = more selective line. A high Qf means high match quality with few high-scoring types.</li>
            <li><strong>Sparsity</strong>: 1 - (fraction of lines with similar pattern). Higher = more unique labeling.</li>
        </ul>
        <table>
            <tr>
                <th>Line</th>
                <th>Neurons</th>
                <th>Types</th>
                <th>Max Score</th>
                <th>n_neurons_HMS</th>
                <th>n_types_HMS</th>
                <th>n_neurons_MS</th>
                <th>n_types_MS</th>
                <th>Qf</th>
                <th>Sparsity</th>
            </tr>
'''
        
        for _, row in line_summary.iterrows():
            html_content += f'''
            <tr>
                <td><strong>{row['line']}</strong></td>
                <td>{row['n_neurons']}</td>
                <td>{row['n_types']}</td>
                <td>{row['max_score']:.0f}</td>
                <td>{row['n_neurons_HMS']}</td>
                <td>{row['n_types_HMS']}</td>
                <td>{row['n_neurons_MS']}</td>
                <td>{row['n_types_MS']}</td>
                <td>{row['Qf']:.0f}</td>
                <td>{row['colabel_sparsity']:.3f}</td>
            </tr>
'''
        
        html_content += '''
        </table>
        
        <h2>🔗 Top Co-Labeling Pairs</h2>
        <p>Lines with highest overlap in labeled neuron types (weighted Jaccard similarity &gt; 0.1):</p>
'''
        
        if top_pairs:
            for pair in top_pairs[:10]:
                bar_width = int(pair['similarity'] * 100)
                html_content += f'''
        <div class="pair-card">
            <strong>{pair['line1']}</strong> ↔ <strong>{pair['line2']}</strong>
            <div style="margin-top: 8px;">
                <div class="similarity-bar" style="width: {bar_width}%;"></div>
            </div>
            <small>Similarity: {pair['similarity']:.3f}</small>
        </div>
'''
        else:
            html_content += '<p><em>No significant co-labeling pairs found (all similarities ≤ 0.1)</em></p>'
        
        html_content += '''
        <h2>🎯 Most Specific Lines</h2>
        <p>Lines with highest sparsity (most unique labeling patterns):</p>
        <table>
            <tr>
                <th>Rank</th>
                <th>Line</th>
                <th>Sparsity</th>
                <th>Quality Factor (Qf)</th>
                <th>Types (HMS)</th>
            </tr>
'''
        
        for i, line_data in enumerate(most_specific[:10], 1):
            html_content += f'''
            <tr>
                <td>{i}</td>
                <td><strong>{line_data['line']}</strong></td>
                <td>{line_data['colabel_sparsity']:.3f}</td>
                <td>{line_data.get('Qf', 0):.2f}</td>
                <td>{line_data.get('n_types_HMS', 0)}</td>
            </tr>
'''
        
        html_content += '''
        </table>
        
        <h2>� Visualizations</h2>
        
        <h3>Expression Matrix Heatmap</h3>
        <p>Interactive heatmap showing which neuron types each line labels. Types prefixed with dataset abbreviations (HEMI_, MCNS_, FAFB_).</p>
        <iframe src="expression_matrix.html" width="100%" height="700" frameborder="0" style="border: 1px solid #ddd; border-radius: 8px;"></iframe>
        
        <h3>Expression Matrix (Merged Datasets)</h3>
        <p>Same expression matrix but with types merged across datasets. Same neuron types from different datasets (e.g., MCNS_aMe12, FAFB_aMe12) are combined into a single row (aMe12) using max score.</p>
        <iframe src="expression_matrix_merged.html" width="100%" height="700" frameborder="0" style="border: 1px solid #ddd; border-radius: 8px;"></iframe>
        
        <h3>Co-Labeling Similarity Matrix</h3>
        <p>Pairwise similarity between lines based on their labeled neuron types (weighted Jaccard).</p>
        <iframe src="colabeling_matrix_weighted_jaccard.html" width="100%" height="700" frameborder="0" style="border: 1px solid #ddd; border-radius: 8px;"></iframe>
        
        <h3>Labeling Distribution (by Type)</h3>
        <p>Score distribution for each line aggregated by neuron type. Mountain-shaped with highest scores in center.</p>
        <iframe src="labeling_distribution_by_type.html" width="100%" height="800" frameborder="0" style="border: 1px solid #ddd; border-radius: 8px;"></iframe>
        
        <h3>Labeling Distribution (by Neuron)</h3>
        <p>Score distribution for each line showing individual neuron scores.</p>
        <iframe src="labeling_distribution_by_neuron.html" width="100%" height="800" frameborder="0" style="border: 1px solid #ddd; border-radius: 8px;"></iframe>
        
        <h3>Stacked Labeling Distribution</h3>
        <p>All lines overlaid to show comparative labeling patterns.</p>
        <iframe src="labeling_distribution_stacked.html" width="100%" height="600" frameborder="0" style="border: 1px solid #ddd; border-radius: 8px;"></iframe>
        
        <h2>📁 Output Files</h2>
        <ul>
            <li><a class="file-link" href="expression_matrix.csv">expression_matrix.csv</a> - Type × Line score matrix (types prefixed with dataset: HEMI_, MCNS_, FAFB_)</li>
            <li><a class="file-link" href="expression_matrix_merged.csv">expression_matrix_merged.csv</a> - Type × Line score matrix with types merged across datasets</li>
            <li><a class="file-link" href="labeling_info.csv">labeling_info.csv</a> - Case-sensitive type × Line matrix with dataset column</li>
            <li><a class="file-link" href="expression_matrix.html">expression_matrix.html</a> - Interactive heatmap (per-dataset types)</li>
            <li><a class="file-link" href="expression_matrix_merged.html">expression_matrix_merged.html</a> - Interactive heatmap (merged types)</li>
            <li><a class="file-link" href="colabeling_matrix_weighted_jaccard.csv">colabeling_matrix_weighted_jaccard.csv</a> - Line similarity matrix</li>
            <li><a class="file-link" href="colabeling_matrix_weighted_jaccard.html">colabeling_matrix_weighted_jaccard.html</a> - Interactive heatmap</li>
            <li><a class="file-link" href="labeling_distribution_by_type.html">labeling_distribution_by_type.html</a> - Score distribution by type</li>
            <li><a class="file-link" href="labeling_distribution_by_neuron.html">labeling_distribution_by_neuron.html</a> - Score distribution by neuron</li>
            <li><a class="file-link" href="labeling_distribution_stacked.html">labeling_distribution_stacked.html</a> - Stacked distribution</li>
            <li><a class="file-link" href="distribution_data_by_neuron.csv">distribution_data_by_neuron.csv</a> - Raw neuron data (unfiltered)</li>
            <li><a class="file-link" href="distribution_data_by_type.csv">distribution_data_by_type.csv</a> - Aggregated type data (unfiltered)</li>
            <li><a class="file-link" href="line_labeled_neurons/">line_labeled_neurons/</a> - Per-line neuron details</li>
            <li><a class="file-link" href="line_summary.csv">line_summary.csv</a> - Summary statistics</li>
            <li><a class="file-link" href="parameters.json">parameters.json</a> - Analysis parameters for reproducibility</li>
        </ul>
        
        <h2>📖 Interpretation Guide</h2>
        <h3>Expression Matrix</h3>
        <p>Shows which neuron types each line labels, with values representing NeuronBridge match scores. 
        Types are prefixed with dataset abbreviations (HEMI=hemibrain, MCNS=male-cns, FAFB=FlyWire FAFB).
        Higher scores indicate stronger morphological matches.</p>
        
        <h3>Expression Matrix (Merged)</h3>
        <p>Same data as the expression matrix but with types merged across datasets. For example, 
        MCNS_aMe12 and FAFB_aMe12 are combined into a single "aMe12" row. The merged score uses the 
        maximum value across all datasets. This view is useful when you care about neuron type identity 
        regardless of which EM dataset it comes from.</p>
        
        <h3>Labeling Info</h3>
        <p>Provides case-sensitive neuron types with their source dataset. Use this file for per-dataset
        filtering and to identify which neurons from each dataset match a given driver line.</p>
        
        <h3>Co-Labeling Matrix</h3>
        <p>Shows pairwise similarity between lines based on their labeled neuron types:</p>
        <ul>
            <li><strong>Weighted Jaccard:</strong> Accounts for match scores - lines labeling the same types with similar scores have higher similarity</li>
            <li><strong>Binary Jaccard:</strong> Simple overlap - what proportion of types are labeled by both lines</li>
        </ul>
        
        <h3>Sparsity</h3>
        <p>Measures how unique a line's labeling pattern is. High sparsity (close to 1.0) means the line labels types 
        that few other lines label - useful for identifying highly specific driver lines.</p>
        
    </div>
</body>
</html>
'''
        
        with open(report_path, 'w') as f:
            f.write(html_content)
        
        self._vprint(f"   📝 Report: {report_path}")
        
        return report_path

    def find_lines_batch(
        self,
        queries: Union[str, int, List[Union[str, int]], Any],  # str, int, List, or LabelMapper
        dataset: Optional[Union[str, List[str]]] = None,
        match_type: Optional[str] = None,
        sort_by: str = 'max',
        output_dir: Optional[str] = None,
        download_images: Optional[str] = 'flylight',
        download_img_for_top_n_lines: Optional[int] = 10,
        image_formats: Union[str, List[str]] = ['png','jpg'],
        image_types: Union[str, List[str]] = 'all',
        max_download_images_per_line: Optional[int] = 12,
        flylight_category: Optional[Union[str, List[str]]] = ['GAL4/LEXA', 'SplitGAL4'],
        organize_by_region: bool = False,
        simple_mode: bool = False,
        pdf_images_per_page: Tuple[int, int] = (3, 2),
        pdf_landscape: bool = True,
        summary_format: Union[str, List[str]] = 'pdf',
    ) -> pd.DataFrame:
        """
        Find driver lines for multiple EM neurons with automatic saving.
        
        This is a convenience method that processes multiple queries (bodyIds,
        types, or instances), adds source information, optionally saves results,
        and optionally downloads images.
        
        Parameters
        ----------
        queries : str, int, list, or LabelMapper
            Neuron query(s). Can be:
            - bodyId (int)
            - type/instance name (str)
            - comma-separated string of names/ids
            - list of queries
            - LabelMapper object for cross-dataset unified naming (extracts queries
              for each dataset from source_mapping or target_mapping)
        dataset : str, list of str, or None
            Dataset(s) for type/instance lookups (e.g., 'hemibrain:v1.2.1').
            Can be a single string or a list of datasets to search multiple.
            Set to None to search ALL available datasets.
            When using LabelMapper, this is automatically derived from the mapper.
        match_type : str, optional
            Match algorithm: 'cds', 'pppm', or 'both'. 
            If None, uses self.match_type. Default: None
        sort_by : str
            Sorting method for line summary results (case-insensitive):
            - 'completeness': Sort by weighted_score (prioritizes lines labeling ALL queried neurons)
            - 'max': Sort by agg_max_score (prioritizes lines with highest individual match scores)
            Default: 'max'
        output_dir : str, optional
            Directory to save results. If provided, saves individual and combined CSVs.
        download_images : str, optional
            Image download source (case-insensitive):
            - 'neuronbridge': Download CDM images from NeuronBridge
            - 'flylight': Download images from FlyLight (S3/HTTP CDN)
            - 'both': Download from both sources
            - None/False: No image download (default)
        download_img_for_top_n_lines : int, optional
            Download images only for top N lines (by aggregate score/rank).
            Default: None (download for all lines)
        image_formats : str or list
            File formats to download. For neuronbridge: 'png', 'jpg'.
            For flylight: 'png', 'jpg', 'h5j', 'lsm', 'mp4', 'json', 'all'.
            Default: 'png'
        image_types : str or list
            Image types to download. For neuronbridge: 'cdm', 'mip'.
            For flylight: 'mip', 'cdm', 'aligned', 'translation', 'metadata', 'all'.
            Default: 'all' (download all available image types)
        max_download_images_per_line : int, optional
            Maximum images to download per line. Default: 20
        flylight_category : str or list, optional
            FlyLight collection category for flylight downloads.
            Options: 'GAL4/LEXA', 'SplitGAL4', 'MCFO', 'RawImages', 'All'.
            Default: ['GAL4/LEXA', 'SplitGAL4'] (search these collections)
        organize_by_region : bool
            If True, organize downloaded images into Brain/VNC subfolders
            based on the anatomical region in the filename. Default: False
        simple_mode : bool
            If True, apply filename filtering to reduce download volume:
            - Split-GAL4 collections: only files with '20x' AND 'multichannel' in filename
            - GAL4/LexA collections: only files with 'total' in filename
            Default: False
        pdf_images_per_page : tuple of (int, int)
            (columns, rows) - number of images per page/slide in the summary.
            Default: (3, 2) = 6 images per page/slide
        pdf_landscape : bool
            Use landscape orientation for PDF. Default: True (horizontal A4)
        summary_format : str or list of str
            Format(s) for summary file generation.
            Options: 'pdf', 'pptx', or list like ['pdf', 'pptx']
            Default: 'pdf'
        
        Notes
        -----
        When `separate_splitgal4=True` is set on the NeuronBridgeFinder instance:
        - Results will include a 'line_type' column ('gal4_lexa' or 'split_gal4')
        - download_img_for_top_n_lines applies separately to each category
        - GAL4/LexA lines: start with VT, R, GMR (e.g., VT037867, R10A06)
        - Split-GAL4 lines: start with SS, LH, MB, IS, OL, LC, etc.
        - Separate CSV files saved: gal4_lexa_lines.csv, split_gal4_lines.csv
        
        Weighted Score and Multi-Type Query:
            The weighted_score column prioritizes lines labeling ALL queried neurons:
            
            weighted_score = agg_mean_score × (match_count / total_query_neurons)
            
            - agg_mean_score: Average NeuronBridge match score
            - match_count: Number of unique query neurons this line labels
            - total_query_neurons: Total unique neurons in the query
            - coverage_ratio: match_count / total_query_neurons
            
            All *_summary.csv files are SORTED BY weighted_score (descending).
            This ensures lines labeling ALL queried types rank highest.
            
            ** IMPORTANT: Multi-type query behavior **
            When querying 'aMe12,MBON01', the program finds lines labeling BOTH types.
            If you want lines for DIFFERENT neuron groups, query separately:
            - Query 1: 'aMe12' → best lines for aMe12
            - Query 2: 'MBON01' → best lines for MBON01
        
        Cross-Dataset Scoring (when multiple datasets):
            When searching across multiple datasets, the results include:
            - 'cross_dataset_score': Mean of max scores across datasets
            - 'datasets_labeled': Number of datasets where the line labels the query
            - 'min_score_per_dataset': Minimum score across all datasets
        
        For Specificity/Selectivity Analysis:
            Use analyze_colabeling() method separately to study how specific each
            driver line is to your target neuron types. This provides:
            - Co-labeling matrices showing overlap between lines
            - Expression matrices showing which types each line labels
            - Specificity score distributions
            
        Returns
        -------
        pd.DataFrame
            Combined DataFrame with all results, including:
            - 'source_query': original query string
            - 'source_type': neuron type (for type queries)
            - 'source_bodyId': matching bodyId for type/instance queries
            - 'source_dataset': source dataset
            - For match_type='both': cds_score, pppm_score, cds_rank, pppm_rank, combined_rank
            - For multi-dataset: cross_dataset_score, datasets_labeled, min_score_per_dataset
            
        Example
        -------
        >>> nbf = NeuronBridgeFinder()
        >>> # Basic search
        >>> results = nbf.find_lines_batch('aMe12,MBON01', dataset='hemibrain:v1.2.1')
        >>> # With NeuronBridge images (top 50 lines only)
        >>> results = nbf.find_lines_batch('aMe12', download_images='neuronbridge', 
        ...                                 download_img_for_top_n_lines=50, output_dir='./output')
        >>> # With FlyLight images
        >>> results = nbf.find_lines_batch('aMe12', download_images='flylight', output_dir='./output')
        >>> # With simple mode (reduced download volume)
        >>> results = nbf.find_lines_batch('aMe12', download_images='flylight',  
        ...                                 simple_mode=True, output_dir='./output')
        >>> # Separate GAL4/LexA from Split-GAL4 (set on instance)
        >>> nbf = NeuronBridgeFinder(separate_splitgal4=True)
        >>> results = nbf.find_lines_batch('MBON01', download_img_for_top_n_lines=5)  # 5 GAL4 + 5 Split-GAL4
        >>> # With LabelMapper for cross-dataset unified naming
        >>> from comparison.label_mapper import LabelMapper
        >>> mapper = LabelMapper(source_mapping_file='my_mappings.json')
        >>> results = nbf.find_lines_batch(mapper, output_dir='./output')
        """
        # Handle LabelMapper input - expand to dataset-specific queries
        label_mapper_info = None  # Track LabelMapper for type info
        if HAS_LABELMAPPER and isinstance(queries, LabelMapper):
            label_mapper_info = queries  # Keep reference
            self._vprint("📋 Expanding LabelMapper to cross-dataset queries...")
            
            # Build query list from LabelMapper's source or target mappings
            expanded_queries = []
            mapper_datasets = set()
            
            # Get datasets and standard labels from the mapper
            for role in ['source', 'target', 'intermediate']:
                all_labels = queries.get_all_std_labels(role)
                for std_label in all_labels:
                    ds_list = queries.get_datasets(role)
                    for ds in ds_list:
                        neurons = queries.get_neurons_for_label(std_label, ds, role)
                        if neurons:
                            mapper_datasets.add(ds)
                            # Add as tuples: (neuron_id, dataset, std_label)
                            for neuron in neurons:
                                expanded_queries.append({
                                    'query': neuron,
                                    'dataset': ds,
                                    'std_label': std_label
                                })
            
            # Override dataset parameter with mapper datasets if not specified
            if dataset is None and mapper_datasets:
                dataset = list(mapper_datasets)
                self._vprint(f"   Using datasets from LabelMapper: {dataset}")
            
            # Build query_list from expanded queries - store metadata for later
            query_list = []
            query_metadata = {}  # {query_str: {std_label, dataset}}
            for eq in expanded_queries:
                q = eq['query']
                query_list.append(q)
                # Store metadata for type info enrichment
                key = f"{q}_{eq['dataset']}"
                query_metadata[key] = {
                    'std_label': eq['std_label'],
                    'dataset': eq['dataset']
                }
            
            # Remove duplicates while preserving order
            seen = set()
            unique_queries = []
            for q in query_list:
                if q not in seen:
                    seen.add(q)
                    unique_queries.append(q)
            query_list = unique_queries
            
            self._vprint(f"   Expanded to {len(query_list)} unique queries from LabelMapper")
        elif isinstance(queries, str):
            query_list = [q.strip() for q in queries.split(',') if q.strip()]
            query_metadata = {}
        elif isinstance(queries, int):
            query_list = [queries]
            query_metadata = {}
        else:
            query_list = list(queries)
            query_metadata = {}
        
        # Use class-level match_type if not specified
        if match_type is None:
            match_type = self.match_type
        else:
            # Validate and normalize match_type using helper
            match_type = self._validate_match_type(match_type)
        
        # Validate and normalize sort_by
        sort_by = self._validate_sort_by(sort_by)
        
        if not query_list:
            self._vprint("❌ No queries provided")
            return pd.DataFrame()
        
        # Determine if multi-dataset search
        is_multi_dataset = isinstance(dataset, list) and len(dataset) > 1
        
        self._vprint(f"🔍 Finding lines for {len(query_list)} query(s)")
        if is_multi_dataset:
            self._vprint(f"   📊 Multi-dataset mode: {len(dataset)} datasets")
        
        # Create output directory if needed
        output_path = None
        if output_dir:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            # Add query info to folder name
            query_info = '_'.join(str(q) for q in query_list[:3])  # First 3 queries
            if len(query_list) > 3:
                query_info += '_etc'
            # Sanitize folder name (remove special characters)
            query_info = ''.join(c if c.isalnum() or c in '-_' else '_' for c in query_info)
            output_path = os.path.join(output_dir, f'findlines_{query_info}_{timestamp}')
            os.makedirs(output_path, exist_ok=True)
            self._vprint(f"   Output: {output_path}")
            
            # Save parameters for reproducibility
            self._save_parameters(
                output_path=output_path,
                function_name='find_lines_batch',
                function_params={
                    'queries': query_list,
                    'dataset': dataset,
                    'match_type': match_type,
                    'sort_by': sort_by,
                    'download_images': download_images,
                    'download_img_for_top_n_lines': download_img_for_top_n_lines,
                    'image_formats': image_formats,
                    'image_types': image_types,
                    'max_download_images_per_line': max_download_images_per_line,
                    'flylight_category': flylight_category,
                    'organize_by_region': organize_by_region,
                    'simple_mode': simple_mode,
                    'pdf_images_per_page': pdf_images_per_page,
                    'pdf_landscape': pdf_landscape
                }
            )
            self._vprint(f"   💾 Parameters: parameters.json")
        
        # Process each query
        all_results = []
        
        # Phase 1: Check cache and identify what needs to be fetched
        if HAS_TQDM and self.verbose and len(query_list) > 1:
            from tqdm import tqdm as tqdm_progress
            
            self._vprint("\n📦 Checking cache...")
            cached_queries = []
            uncached_queries = []
            
            cache_pbar = tqdm_progress(
                query_list,
                desc="   💾 Loading cache",
                unit="query",
                bar_format='{desc}: {percentage:3.0f}%|{bar:30}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]',
                ncols=110,
                position=0,
                leave=False
            )
            
            for q in cache_pbar:
                # Determine cache key based on query type
                try:
                    body_id = int(q)
                    is_body_id = True
                    cache_key = f"{body_id}_{match_type}"
                    cached_data = self._load_from_cache('id_to_lines', cache_key)
                    query_name = str(body_id)
                except (ValueError, TypeError):
                    is_body_id = False
                    # For type/instance queries, cache is per bodyId in neuron_to_lines
                    # We can't easily check cache here, so mark as uncached
                    cached_data = None
                    query_name = str(q)
                
                if cached_data is not None and is_body_id:
                    if not cached_data.empty:
                        cached_data = cached_data.copy()
                        cached_data['source_bodyId'] = body_id
                        cached_data['source_query'] = query_name
                        cached_queries.append((q, cached_data))
                        cache_pbar.set_postfix_str(f"✓ {query_name}")
                    else:
                        uncached_queries.append(q)
                else:
                    uncached_queries.append(q)
            
            cache_pbar.close()
            
            if cached_queries:
                self._vprint(f"   ✓ Loaded {len(cached_queries)} from cache")
            if uncached_queries:
                self._vprint(f"   🌐 Need to fetch {len(uncached_queries)} from API")
            
            # Process cached results
            for q, cached_data in cached_queries:
                all_results.append(cached_data)
                
                # Save individual cached results (sorted by score descending)
                if output_path:
                    query_name = str(q)
                    safe_name = ''.join(c if c.isalnum() or c in '_-' else '_' for c in query_name)
                    output_file = os.path.join(output_path, f'{safe_name}_lines.csv')
                    # Sort by score descending before saving
                    if 'score' in cached_data.columns:
                        cached_data_sorted = cached_data.sort_values('score', ascending=False)
                    else:
                        cached_data_sorted = cached_data
                    cached_data_sorted.to_csv(output_file, index=False)
            
            # Phase 2: Fetch uncached data
            if uncached_queries:
                self._vprint("\n🌐 Fetching new data...")
                fetch_pbar = tqdm_progress(
                    uncached_queries,
                    desc="   🔄 Fetching",
                    unit="query",
                    bar_format='{desc}: {percentage:3.0f}%|{bar:30}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]',
                    ncols=110,
                    position=0,
                    leave=False
                )
                
                for q in fetch_pbar:
                    query_name = str(q)
                    fetch_pbar.set_postfix_str(query_name[:20])
                    
                    try:
                        # Check if query is a body ID (integer)
                        try:
                            body_id = int(q)
                            is_body_id = True
                        except (ValueError, TypeError):
                            body_id = None
                            is_body_id = False
                        
                        if is_body_id:
                            # Direct body ID search
                            lines_df = self.id_to_lines(body_id, match_type=match_type)
                            if not lines_df.empty:
                                lines_df = lines_df.copy()
                                lines_df['source_bodyId'] = body_id
                            query_name = str(body_id)
                        else:
                            # Type/instance search
                            results_dict = self.neuron_to_lines(
                                q, dataset=dataset, match_type=match_type
                            )
                            if results_dict:
                                dfs = [df for df in results_dict.values() if not df.empty]
                                lines_df = pd.concat(dfs, ignore_index=True) if dfs else pd.DataFrame()
                            else:
                                lines_df = pd.DataFrame()
                            query_name = str(q)
                        
                        if lines_df.empty:
                            continue
                        
                        # Add source query info
                        lines_df = lines_df.copy()
                        lines_df['source_query'] = query_name
                        all_results.append(lines_df)
                        
                        # Save individual results (sorted by score descending)
                        if output_path:
                            safe_name = ''.join(c if c.isalnum() or c in '_-' else '_' for c in query_name)
                            output_file = os.path.join(output_path, f'{safe_name}_lines.csv')
                            # Sort by score descending before saving
                            if 'score' in lines_df.columns:
                                lines_df_sorted = lines_df.sort_values('score', ascending=False)
                            else:
                                lines_df_sorted = lines_df
                            lines_df_sorted.to_csv(output_file, index=False)
                        
                    except Exception as e:
                        pass  # Silent in progress mode
                
                fetch_pbar.close()
                self._vprint(f"   ✓ Fetched {len(uncached_queries)} from API")
        
        else:
            # Single query or no progress bar - original behavior
            for q in query_list:
                self._vprint(f"\n📋 Processing: {q}")
                
                try:
                    # Check if query is a body ID (integer)
                    try:
                        body_id = int(q)
                        is_body_id = True
                    except (ValueError, TypeError):
                        body_id = None
                        is_body_id = False
                    
                    if is_body_id:
                        # Direct body ID search
                        lines_df = self.id_to_lines(body_id, match_type=match_type)
                        if not lines_df.empty:
                            lines_df = lines_df.copy()
                            lines_df['source_bodyId'] = body_id
                        query_name = str(body_id)
                    else:
                        # Type/instance search
                        results_dict = self.neuron_to_lines(
                            q, dataset=dataset, match_type=match_type
                        )
                        if results_dict:
                            dfs = [df for df in results_dict.values() if not df.empty]
                            lines_df = pd.concat(dfs, ignore_index=True) if dfs else pd.DataFrame()
                        else:
                            lines_df = pd.DataFrame()
                        query_name = str(q)
                    
                    if lines_df.empty:
                        self._vprint(f"   ⚠️ No matching lines found")
                        continue
                    
                    # Add source query info
                    lines_df = lines_df.copy()
                    lines_df['source_query'] = query_name
                    all_results.append(lines_df)
                    
                    self._vprint(f"   ✅ Found {len(lines_df)} matching driver lines")
                    
                    # Save individual results (sorted by score descending)
                    if output_path:
                        safe_name = ''.join(c if c.isalnum() or c in '_-' else '_' for c in query_name)
                        output_file = os.path.join(output_path, f'{safe_name}_lines.csv')
                        # Sort by score descending before saving
                        if 'score' in lines_df.columns:
                            lines_df_sorted = lines_df.sort_values('score', ascending=False)
                        else:
                            lines_df_sorted = lines_df
                        lines_df_sorted.to_csv(output_file, index=False)
                        self._vprint(f"   💾 Saved: {output_file}")
                        
                except Exception as e:
                    self._vprint(f"   ❌ Error: {e}")
        
        # Combine results
        if all_results:
            combined_df = pd.concat(all_results, ignore_index=True)
            
            # Calculate total unique query neurons (bodyIds) for weighted_score normalization
            total_query_neurons = 0
            if 'source_bodyId' in combined_df.columns:
                total_query_neurons = len(set(str(v) for v in combined_df['source_bodyId'].dropna().unique()))
            
            # Add line_type classification if separate_splitgal4 is enabled
            if self.separate_splitgal4 and 'line' in combined_df.columns:
                combined_df['line_type'] = combined_df['line'].apply(self._classify_line_type)
            
            # Enrich with type information from neuron_df if available
            if 'source_bodyId' in combined_df.columns and 'source_dataset' in combined_df.columns:
                # Lookup type info for each bodyId
                type_cache = {}
                for _, row in combined_df.iterrows():
                    body_id = row.get('source_bodyId')
                    ds = row.get('source_dataset')
                    if pd.isna(body_id) or pd.isna(ds):
                        continue
                    key = (str(body_id), ds)
                    if key not in type_cache:
                        ds_folder = self._dataset_name_to_folder(ds)
                        neuron_df = self._load_neuron_df_for_dataset(ds_folder)
                        if neuron_df is not None and 'bodyId' in neuron_df.columns and 'type' in neuron_df.columns:
                            match = neuron_df[neuron_df['bodyId'].astype(str) == str(body_id)]
                            if not match.empty:
                                type_cache[key] = match['type'].iloc[0]
                            else:
                                type_cache[key] = None
                        else:
                            type_cache[key] = None
                
                # Add source_type column
                combined_df['source_type'] = combined_df.apply(
                    lambda row: type_cache.get((str(row.get('source_bodyId', '')), row.get('source_dataset', '')), ''),
                    axis=1
                )
            
            # Aggregate line-level results across bodyIds for ranking
            # For each unique line, compute aggregate score/rank
            if 'line' in combined_df.columns:
                if match_type == 'both' and 'combined_rank' in combined_df.columns:
                    # For 'both', use mean combined_rank (lower is better)
                    agg_dict = {
                        'combined_rank': 'mean',
                        'score': 'max',
                        # Use nunique for unique bodyId count instead of count with repeats
                        'source_bodyId': [
                            lambda x: len(set(str(v) for v in x.dropna().unique())),  # unique count
                            lambda x: ','.join(sorted(set(str(v) for v in x.dropna().unique())))  # unique list
                        ]
                    }
                    if self.separate_splitgal4 and 'line_type' in combined_df.columns:
                        agg_dict['line_type'] = 'first'
                    if 'source_type' in combined_df.columns:
                        agg_dict['source_type'] = lambda x: ','.join(sorted(set(str(v) for v in x.dropna().unique() if v)))
                    
                    line_stats = combined_df.groupby('line').agg(agg_dict).reset_index()
                    # Flatten multi-level columns
                    line_stats.columns = ['_'.join(col).strip('_') if isinstance(col, tuple) else col 
                                         for col in line_stats.columns]
                    line_stats = line_stats.rename(columns={
                        'combined_rank_mean': 'agg_combined_rank',
                        'score_max': 'agg_max_score',
                        'source_bodyId_<lambda_0>': 'match_count',
                        'source_bodyId_<lambda_1>': 'matched_bodyIds',
                        'source_type_<lambda>': 'matched_types'
                    })
                    
                    # Calculate weighted_score for 'both' match_type
                    # weighted_score = agg_max_score * (match_count / total_query_neurons)
                    if 'match_count' in line_stats.columns and total_query_neurons > 0:
                        line_stats['coverage_ratio'] = line_stats['match_count'] / total_query_neurons
                        line_stats['weighted_score'] = line_stats['agg_max_score'] * line_stats['coverage_ratio']
                        # Sort based on sort_by parameter
                        if sort_by == 'completeness':
                            # Sort by weighted_score (higher is better), then by combined_rank (lower is better)
                            line_stats = line_stats.sort_values(
                                ['weighted_score', 'agg_combined_rank'], 
                                ascending=[False, True]
                            )
                        else:  # sort_by == 'max'
                            # Sort by agg_mean_score (higher is better), then by combined_rank (lower is better)
                            line_stats = line_stats.sort_values(
                                ['agg_mean_score', 'agg_combined_rank'], 
                                ascending=[False, True]
                            )
                    else:
                        line_stats = line_stats.sort_values('agg_combined_rank', ascending=True)
                else:
                    # For cds/pppm, use mean score (higher is better)
                    # FIXED: match_count now uses nunique for unique bodyId count
                    agg_dict = {
                        'score': ['mean', 'max'],
                        # Use nunique for unique bodyId count instead of count with repeats
                        'source_bodyId': [
                            lambda x: len(set(str(v) for v in x.dropna().unique())),  # unique count
                            lambda x: ','.join(sorted(set(str(v) for v in x.dropna().unique())))  # unique list
                        ]
                    }
                    # Add type aggregation if available
                    if 'source_type' in combined_df.columns:
                        agg_dict['source_type'] = lambda x: ','.join(sorted(set(str(v) for v in x.dropna().unique() if v)))
                    
                    # Add dataset info for cross-dataset scoring
                    if 'source_dataset' in combined_df.columns and is_multi_dataset:
                        agg_dict['source_dataset'] = [
                            lambda x: len(set(str(v) for v in x.dropna().unique())),  # datasets_labeled
                            lambda x: ','.join(sorted(set(str(v) for v in x.dropna().unique())))  # dataset list
                        ]
                    
                    line_stats = combined_df.groupby('line').agg(agg_dict).reset_index()
                    # Flatten multi-level columns
                    new_cols = ['line']
                    for col in line_stats.columns[1:]:
                        if isinstance(col, tuple):
                            new_cols.append(f"{col[0]}_{col[1] if isinstance(col[1], str) else 'agg'}")
                        else:
                            new_cols.append(col)
                    line_stats.columns = new_cols
                    
                    # Rename columns properly
                    rename_map = {
                        'score_mean': 'agg_mean_score',
                        'score_max': 'agg_max_score',
                    }
                    # Find and rename source_bodyId columns
                    for col in line_stats.columns:
                        if 'source_bodyId' in col and 'lambda' in col and '0' in col:
                            rename_map[col] = 'match_count'
                        elif 'source_bodyId' in col and 'lambda' in col:
                            rename_map[col] = 'matched_bodyIds'
                        elif 'source_type' in col and 'lambda' in col:
                            rename_map[col] = 'matched_types'
                        elif 'source_dataset' in col and 'lambda' in col and '0' in col:
                            rename_map[col] = 'datasets_labeled'
                        elif 'source_dataset' in col and 'lambda' in col:
                            rename_map[col] = 'matched_datasets'
                    
                    line_stats = line_stats.rename(columns=rename_map)
                    
                    # Add line_type if available
                    if self.separate_splitgal4 and 'line_type' in combined_df.columns:
                        line_type_map = combined_df.groupby('line')['line_type'].first()
                        line_stats['line_type'] = line_stats['line'].map(line_type_map)
                    
                    # Calculate cross-dataset score if multi-dataset
                    if is_multi_dataset and 'source_dataset' in combined_df.columns:
                        # Calculate min score per dataset for each line
                        # This rewards lines that have good scores across ALL datasets
                        def calc_min_score_across_datasets(line_name):
                            line_data = combined_df[combined_df['line'] == line_name]
                            if 'source_dataset' not in line_data.columns:
                                return line_data['score'].mean() if 'score' in line_data.columns else 0
                            
                            # Get max score per dataset for this line
                            dataset_scores = line_data.groupby('source_dataset')['score'].max()
                            if len(dataset_scores) == 0:
                                return 0
                            # Return minimum of max scores (worst dataset performance)
                            return dataset_scores.min()
                        
                        def calc_cross_dataset_score(line_name):
                            """Calculate mean of max scores across datasets."""
                            line_data = combined_df[combined_df['line'] == line_name]
                            if 'source_dataset' not in line_data.columns:
                                return line_data['score'].mean() if 'score' in line_data.columns else 0
                            
                            # Get max score per dataset for this line
                            dataset_scores = line_data.groupby('source_dataset')['score'].max()
                            if len(dataset_scores) == 0:
                                return 0
                            return dataset_scores.mean()
                        
                        line_stats['min_score_per_dataset'] = line_stats['line'].apply(calc_min_score_across_datasets)
                        line_stats['cross_dataset_score'] = line_stats['line'].apply(calc_cross_dataset_score)
                    
                    # =================================================================
                    # WEIGHTED SCORE CALCULATION
                    # =================================================================
                    # weighted_score = agg_mean_score * (match_count / total_query_neurons)
                    # 
                    # This score prioritizes lines that:
                    # 1. Have high average matching scores (agg_mean_score)
                    # 2. Label MORE of the queried neurons (match_count / total_query_neurons)
                    #
                    # For multi-type queries, this finds lines that label ALL queried neurons.
                    # A line labeling all N queried neurons with score S gets weighted_score = S
                    # A line labeling only 1 of N neurons with score S gets weighted_score = S/N
                    # =================================================================
                    if 'match_count' in line_stats.columns and total_query_neurons > 0:
                        # coverage_ratio = match_count / total_query_neurons
                        # weighted_score = agg_mean_score * coverage_ratio
                        line_stats['coverage_ratio'] = line_stats['match_count'] / total_query_neurons
                        line_stats['weighted_score'] = line_stats['agg_mean_score'] * line_stats['coverage_ratio']
                        
                        # Sort based on sort_by parameter
                        if sort_by == 'completeness':
                            # Sort by weighted_score (prioritizes lines labeling ALL queried neurons)
                            line_stats = line_stats.sort_values('weighted_score', ascending=False)
                            self._vprint(f"   📊 Sorting by weighted_score (agg_mean_score × coverage_ratio)")
                        else:  # sort_by == 'max'
                            # Sort by agg_mean_score (prioritizes lines with highest average scores)
                            line_stats = line_stats.sort_values('agg_mean_score', ascending=False)
                            self._vprint(f"   📊 Sorting by agg_mean_score (average match scores)")
                        self._vprint(f"      Total query neurons: {total_query_neurons}")
                    elif is_multi_dataset:
                        # Fallback to min_score_per_dataset for multi-dataset
                        line_stats = line_stats.sort_values('min_score_per_dataset', ascending=False)
                        self._vprint(f"   📊 Multi-dataset sorting: by min_score_per_dataset (descending)")
                    else:
                        if sort_by == 'completeness':
                            line_stats = line_stats.sort_values('agg_mean_score', ascending=False)
                        else:  # sort_by == 'max'
                            line_stats = line_stats.sort_values('agg_mean_score', ascending=False)
                
                # Add aggregated stats back to combined_df
                merge_cols = ['line', 'matched_bodyIds']
                if 'matched_types' in line_stats.columns:
                    merge_cols.append('matched_types')
                combined_df = combined_df.merge(
                    line_stats[merge_cols], 
                    on='line', 
                    how='left'
                )
            
            # Save combined results (sorted by score descending)
            if output_path:
                combined_file = os.path.join(output_path, 'all_lines.csv')
                # Sort by score descending before saving
                if 'score' in combined_df.columns:
                    combined_df_sorted = combined_df.sort_values('score', ascending=False)
                else:
                    combined_df_sorted = combined_df
                combined_df_sorted.to_csv(combined_file, index=False)
                self._vprint(f"\n💾 Combined results: {combined_file}")
                self._vprint(f"   Total: {len(combined_df)} lines from {len(query_list)} query(s)")
                
                # Save line-level aggregate summary
                if 'line' in combined_df.columns:
                    summary_file = os.path.join(output_path, 'line_summary.csv')
                    
                    # Reorder columns: put weighted_score before agg_mean_score
                    # Desired order: line, match_count, weighted_score, coverage_ratio, agg_mean_score, agg_max_score, ...
                    cols = list(line_stats.columns)
                    priority_cols = ['line', 'line_type', 'match_count', 'weighted_score', 'coverage_ratio', 
                                     'agg_mean_score', 'agg_max_score', 'min_score_per_dataset', 
                                     'cross_dataset_score', 'datasets_labeled', 'matched_bodyIds', 
                                     'matched_types', 'matched_datasets']
                    # Build ordered column list: priority columns first (if they exist), then remaining
                    ordered_cols = [col for col in priority_cols if col in cols]
                    remaining_cols = [col for col in cols if col not in ordered_cols]
                    final_cols = ordered_cols + remaining_cols
                    line_stats = line_stats[final_cols]
                    
                    line_stats.to_csv(summary_file, index=False)
                    self._vprint(f"   Summary: {summary_file} ({len(line_stats)} unique lines)")
                    
                    # Save separate files for GAL4/LexA and Split-GAL4 if enabled
                    if self.separate_splitgal4 and 'line_type' in line_stats.columns:
                        # Split combined_df by line type
                        gal4_lexa_combined = combined_df[combined_df['line_type'] == 'gal4_lexa']
                        split_gal4_combined = combined_df[combined_df['line_type'] == 'split_gal4']
                        
                        # Split line_stats by line type
                        gal4_lexa_stats = line_stats[line_stats['line_type'] == 'gal4_lexa']
                        split_gal4_stats = line_stats[line_stats['line_type'] == 'split_gal4']
                        
                        # Save GAL4/LexA files (sorted by score descending)
                        if not gal4_lexa_combined.empty:
                            gal4_file = os.path.join(output_path, 'gal4_lexa_lines.csv')
                            if 'score' in gal4_lexa_combined.columns:
                                gal4_lexa_combined = gal4_lexa_combined.sort_values('score', ascending=False)
                            gal4_lexa_combined.to_csv(gal4_file, index=False)
                            gal4_summary = os.path.join(output_path, 'gal4_lexa_summary.csv')
                            gal4_lexa_stats.to_csv(gal4_summary, index=False)
                            self._vprint(f"   GAL4/LexA: {gal4_file} ({len(gal4_lexa_combined)} matches, {len(gal4_lexa_stats)} unique)")
                        
                        # Save Split-GAL4 files (sorted by score descending)
                        if not split_gal4_combined.empty:
                            split_file = os.path.join(output_path, 'split_gal4_lines.csv')
                            if 'score' in split_gal4_combined.columns:
                                split_gal4_combined = split_gal4_combined.sort_values('score', ascending=False)
                            split_gal4_combined.to_csv(split_file, index=False)
                            split_summary = os.path.join(output_path, 'split_gal4_summary.csv')
                            split_gal4_stats.to_csv(split_summary, index=False)
                            self._vprint(f"   Split-GAL4: {split_file} ({len(split_gal4_combined)} matches, {len(split_gal4_stats)} unique)")
            
            # Download images if requested
            if download_images and output_path:
                # Normalize download_images parameter (case-insensitive)
                download_source = download_images.lower() if isinstance(download_images, str) else None
                
                if download_source in ('neuronbridge', 'flylight', 'both'):
                    # Get line names for image download (use aggregated ranking)
                    if 'line' in combined_df.columns:
                        # Use pre-computed line_stats for ordering
                        all_lines = line_stats['line'].tolist()
                        
                        # Handle separate_splitgal4 mode
                        if self.separate_splitgal4:
                            # Separate lines by type and get top N from each
                            gal4_lines = [l for l in all_lines if self._classify_line_type(l) == 'gal4_lexa']
                            split_lines = [l for l in all_lines if self._classify_line_type(l) == 'split_gal4']
                            
                            if download_img_for_top_n_lines is not None and download_img_for_top_n_lines > 0:
                                gal4_lines = gal4_lines[:download_img_for_top_n_lines]
                                split_lines = split_lines[:download_img_for_top_n_lines]
                                self._vprint(f"\n🖼️  Downloading images (separate mode):")
                                self._vprint(f"      GAL4/LexA: top {len(gal4_lines)} lines")
                                self._vprint(f"      Split-GAL4: top {len(split_lines)} lines")
                            else:
                                self._vprint(f"\n🖼️  Downloading images (separate mode):")
                                self._vprint(f"      GAL4/LexA: {len(gal4_lines)} lines")
                                self._vprint(f"      Split-GAL4: {len(split_lines)} lines")
                            
                            download_lines = gal4_lines + split_lines
                        else:
                            # Normal mode: just apply top_n limit to all lines
                            if download_img_for_top_n_lines is not None and download_img_for_top_n_lines > 0:
                                download_lines = all_lines[:download_img_for_top_n_lines]
                                self._vprint(f"\n🖼️  Downloading images for top {len(download_lines)} lines...")
                            else:
                                download_lines = all_lines
                                self._vprint(f"\n🖼️  Downloading images for {len(download_lines)} lines...")
                    else:
                        download_lines = []
                    
                    if download_lines:
                        images_dir = os.path.join(output_path, 'images')
                        
                        # Download from NeuronBridge
                        if download_source in ('neuronbridge', 'both'):
                            nb_dir = os.path.join(images_dir, 'neuronbridge') if download_source == 'both' else images_dir
                            self._download_neuronbridge_images(
                                lines=download_lines,
                                output_dir=nb_dir,
                                formats=image_formats,
                                image_types=image_types,
                                max_files=max_download_images_per_line,
                                verbose=self.verbose
                            )
                        
                        # Download from FlyLight
                        if download_source in ('flylight', 'both'):
                            fl_dir = os.path.join(images_dir, 'flylight') if download_source == 'both' else images_dir
                            flylight_files, lines_without_flylight = self._download_flylight_images_with_category(
                                lines=download_lines,
                                output_dir=fl_dir,
                                formats=image_formats,
                                image_types=image_types,
                                max_files=max_download_images_per_line,
                                category=flylight_category,
                                organize_by_region=organize_by_region,
                                simple_mode=simple_mode,
                                verbose=self.verbose
                            )
                            
                            # Report lines without any FlyLight images
                            if lines_without_flylight:
                                self._vprint(f"\n⚠️  Note: No FlyLight images found for {len(lines_without_flylight)} line(s):")
                                self._vprint(f"   {', '.join(lines_without_flylight)}")
                                self._vprint("   (tried all categories including MCFO fallback)")
                            
                            # Generate PDF/PPTX summary if images were downloaded
                            images_dir = os.path.join(output_path, 'images')
                            if os.path.exists(images_dir):
                                # Normalize summary_format
                                formats = summary_format if isinstance(summary_format, list) else [summary_format]
                                formats = [f.lower() for f in formats]
                                
                                if 'pdf' in formats:
                                    self._vprint(f"\n📄 Generating PDF summary...")
                                    pdf_path = create_image_pdf(
                                        images_dir=images_dir,
                                        output_pdf=os.path.join(output_path, 'images_summary.pdf'),
                                        images_per_page=pdf_images_per_page,
                                        landscape=pdf_landscape,
                                        line_order=download_lines,  # Preserve ranking order
                                        verbose=self.verbose
                                    )
                                    if pdf_path:
                                        self._vprint(f"   ✅ PDF saved: {pdf_path}")
                                
                                if 'pptx' in formats:
                                    self._vprint(f"\n📊 Generating PPTX summary...")
                                    pptx_path = create_image_pptx(
                                        images_dir=images_dir,
                                        output_pptx=os.path.join(output_path, 'images_summary.pptx'),
                                        images_per_slide=pdf_images_per_page,
                                        line_order=download_lines,  # Preserve ranking order
                                        verbose=self.verbose
                                    )
                                    if pptx_path:
                                        self._vprint(f"   ✅ PPTX saved: {pptx_path}")
            
            return combined_df
        
        return pd.DataFrame()
    
    def clear_cache(self, cache_type: Optional[str] = None):
        """
        Clear cached results.
        
        Parameters
        ----------
        cache_type : str, optional
            Type of cache to clear: 'id_to_lines', 'line_to_neuron', or None (all).
        """
        if not os.path.exists(self.cache_folder):
            return
        
        import glob
        
        if cache_type:
            pattern = os.path.join(self.cache_folder, f"{cache_type}_*.csv")
        else:
            pattern = os.path.join(self.cache_folder, "*.csv")
        
        files = glob.glob(pattern)
        for f in files:
            try:
                os.remove(f)
            except Exception:
                pass
        
        self._vprint(f"🗑️ Cleared {len(files)} cached files")
    
    # =========================================================================
    # Image Download Methods
    # =========================================================================
    
    def download_line_images(
        self,
        line_names: Union[str, List[str]],
        output_dir: str,
        source: str = 'neuronbridge',
        formats: Union[str, List[str]] = 'png',
        image_types: Union[str, List[str]] = 'cdm',
        max_files: Optional[int] = None,
        verbose: Optional[bool] = None
    ) -> List[str]:
        """
        Download images for driver lines from NeuronBridge or FlyLight.
        
        Parameters
        ----------
        line_names : str or list
            Driver line name(s) to download images for.
            Can be comma-separated string or list.
        output_dir : str
            Directory to save downloaded images.
        source : str
            Image source: 'neuronbridge' for CDM images, 'flylight' for raw images.
            Default: 'neuronbridge'
        formats : str or list
            File formats to download. For neuronbridge: 'png', 'jpg'. 
            For flylight: 'png', 'jpg', 'h5j', 'lsm', 'mp4', 'json', 'all'.
            Default: 'png'
        image_types : str or list
            Image types to download. For neuronbridge: 'cdm', 'mip'.
            For flylight: 'mip', 'cdm', 'aligned', 'translation', 'metadata', 'all'.
            Default: 'cdm'
        max_files : int, optional
            Maximum number of files to download per line. Default: None (no limit)
        verbose : bool, optional
            Override class verbose setting. Default: use class setting.
            
        Returns
        -------
        list
            List of downloaded file paths.
            
        Example
        -------
        >>> nbf = NeuronBridgeFinder()
        >>> files = nbf.download_line_images('LH173', './images', source='neuronbridge')
        >>> files = nbf.download_line_images('VT037867', './images', source='flylight')
        """
        if verbose is None:
            verbose = self.verbose
        
        # Parse line names
        if isinstance(line_names, str):
            lines = [l.strip() for l in line_names.split(',') if l.strip()]
        else:
            lines = list(line_names)
        
        if not lines:
            if verbose:
                print("❌ No line names provided")
            return []
        
        # Create output directory
        os.makedirs(output_dir, exist_ok=True)
        
        all_downloaded = []
        
        if source.lower() == 'neuronbridge':
            all_downloaded = self._download_neuronbridge_images(
                lines, output_dir, formats, image_types, max_files, verbose
            )
        elif source.lower() == 'flylight':
            all_downloaded = self._download_flylight_images(
                lines, output_dir, formats, image_types, max_files, verbose
            )
        else:
            if verbose:
                print(f"❌ Unknown source: {source}. Use 'neuronbridge' or 'flylight'.")
        
        return all_downloaded
    
    def _download_neuronbridge_images(
        self,
        lines: List[str],
        output_dir: str,
        formats: Union[str, List[str]],
        image_types: Union[str, List[str]],
        max_files: Optional[int],
        verbose: bool
    ) -> List[str]:
        """Download CDM images from NeuronBridge with progress bar."""
        import urllib.request
        
        downloaded = []
        
        # NeuronBridge CDM URL patterns
        # CDM files are at: https://s3.amazonaws.com/janelia-flylight-color-depth/{path}
        cdm_base_url = "https://s3.amazonaws.com/janelia-flylight-color-depth/"
        
        # Phase 1: Scan all lines to count total files
        if verbose:
            print(f"📊 Scanning {len(lines)} lines for NeuronBridge images...")
        
        line_files_map = {}  # line_name -> list of (file_type, file_path)
        total_file_count = 0
        
        for line_name in lines:
            try:
                lm_images = self._client.get_lm_images(line_name)
                if not lm_images:
                    line_files_map[line_name] = []
                    continue
                
                if not isinstance(lm_images, list):
                    lm_images = list(lm_images)
                
                file_paths = []
                for lm_image in lm_images:
                    if max_files and len(file_paths) >= max_files:
                        break
                    
                    files_obj = getattr(lm_image, 'files', None)
                    if not files_obj:
                        continue
                    
                    img_types = image_types if isinstance(image_types, list) else [image_types]
                    
                    for img_type in img_types:
                        if img_type.lower() in ['cdm', 'all']:
                            cdm_path = getattr(files_obj, 'CDM', None)
                            if cdm_path:
                                file_paths.append(('cdm', cdm_path))
                        if img_type.lower() in ['mip', 'all']:
                            mip_path = getattr(files_obj, 'SignalMip', None)
                            if mip_path:
                                file_paths.append(('mip', mip_path))
                            mip_masked = getattr(files_obj, 'SignalMipMasked', None)
                            if mip_masked:
                                file_paths.append(('mip_masked', mip_masked))
                
                # Filter by format
                fmt_list = formats if isinstance(formats, list) else [formats]
                filtered_paths = []
                for file_type, file_path in file_paths:
                    ext = os.path.splitext(file_path)[1].lower().lstrip('.')
                    if 'all' in fmt_list or ext in [f.lstrip('.').lower() for f in fmt_list]:
                        filtered_paths.append((file_type, file_path))
                
                if max_files:
                    filtered_paths = filtered_paths[:max_files]
                
                line_files_map[line_name] = filtered_paths
                total_file_count += len(filtered_paths)
                
            except Exception:
                line_files_map[line_name] = []
        
        if total_file_count == 0:
            if verbose:
                print("  ⚠️ No files found for any lines")
            return []
        
        lines_with_files = len([l for l, f in line_files_map.items() if f])
        if verbose:
            print(f"  📦 Found {total_file_count} files across {lines_with_files} lines")
        
        # Phase 2: Download with progress bar
        pbar = None
        if HAS_TQDM and verbose:
            pbar = tqdm(total=total_file_count, desc="  Downloading", unit="file")
        
        files_downloaded = 0
        for line_name in lines:
            if line_name not in line_files_map or not line_files_map[line_name]:
                continue
            
            try:
                line_dir = os.path.join(output_dir, line_name)
                os.makedirs(line_dir, exist_ok=True)
                
                for file_type, file_path in line_files_map[line_name]:
                    url = cdm_base_url + file_path
                    filename = os.path.basename(file_path)
                    local_path = os.path.join(line_dir, filename)
                    
                    try:
                        if pbar:
                            pbar.set_postfix(line=line_name[:15], refresh=False)
                        urllib.request.urlretrieve(url, local_path)
                        downloaded.append(local_path)
                        files_downloaded += 1
                        if pbar:
                            pbar.update(1)
                        elif verbose:
                            print(f"  [{files_downloaded}/{total_file_count}] {line_name}: {filename}")
                    except Exception as e:
                        if verbose and not pbar:
                            print(f"  ⚠️ Failed to download {filename}: {e}")
                        
            except Exception as e:
                if verbose and not pbar:
                    print(f"  ❌ Error processing '{line_name}': {e}")
        
        if pbar:
            pbar.close()
        
        if verbose:
            print(f"  ✅ Downloaded {len(downloaded)}/{total_file_count} files")
        
        return downloaded
    
    def _download_flylight_images(
        self,
        lines: List[str],
        output_dir: str,
        formats: Union[str, List[str]],
        image_types: Union[str, List[str]],
        max_files: Optional[int],
        verbose: bool
    ) -> List[str]:
        """Download images from FlyLight using flylight_downloader module."""
        return self._download_flylight_images_with_category(
            lines=lines,
            output_dir=output_dir,
            formats=formats,
            image_types=image_types,
            max_files=max_files,
            category=None,
            organize_by_region=False,
            verbose=verbose
        )
    
    def _parse_region_from_filename(self, filename: str, full_key: str = None) -> str:
        """
        Parse anatomical region from FlyLight filename or file key path.
        
        FlyLight filename format: {line}-{date}_{sample}-{sex}-{mag}-{region}-{driver}-...
        Example: SS01015-20131220_31_C3-f-20x-brain-Split_GAL4-JRC2018_Unisex_20x_HR-CDM_1.png
        
        VT GAL4 files have region in the key path:
        Example: VT GAL4/VT037867/brain/filename.jpg
        
        Brain regions include: brain, central, dorsal, optic, protocerebrum, etc.
        VNC regions include: vnc, ventral_nerve_cord, metathoracic, prothoracic, mesothoracic
        
        Returns 'Brain', 'VNC', or 'Other' based on the region field.
        """
        # First check the full key path for region (VT GAL4 files)
        if full_key:
            key_lower = full_key.lower()
            # Check for /brain/ or /vnc/ in the path
            if '/brain/' in key_lower:
                return 'Brain'
            if '/vnc/' in key_lower:
                return 'VNC'
        
        # Split by '-' and try to find region field (MCFO and other formats)
        parts = filename.split('-')
        
        # VNC-specific keywords (check first to avoid confusion with 'ventral')
        vnc_keywords = ['vnc', 'ventral_nerve_cord', 'metathoracic', 'prothoracic', 'mesothoracic']
        
        # Brain-specific keywords
        brain_keywords = ['brain', 'central', 'dorsal', 'optic', 'protocerebrum', 
                         'left_optic_lobe', 'right_optic_lobe', 'left_dorsal', 'right_dorsal',
                         'ventral']  # 'ventral' alone (without nerve_cord) means brain ventral
        
        # Search for VNC keywords first (more specific)
        for part in parts:
            part_lower = part.lower()
            for keyword in vnc_keywords:
                if keyword in part_lower:
                    return 'VNC'
        
        # Then search for brain keywords
        for part in parts:
            part_lower = part.lower()
            for keyword in brain_keywords:
                if keyword in part_lower:
                    return 'Brain'
        
        return 'Other'
    
    def _filter_flylight_files_by_region(self, files: List[Any]) -> List[Any]:
        """
        Filter FlyLight files based on anatomical region setting.
        
        Uses filename and key path parsing to determine region (Brain/VNC).
        Only filters when self.region is 'Brain' or 'VNC'.
        
        Parameters
        ----------
        files : list
            List of FlyLightFile objects.
            
        Returns
        -------
        list
            Filtered list of files matching the region criteria.
        """
        if self.region == 'All' or not self.region:
            return files
        
        filtered = []
        for file in files:
            # Get full key and filename
            full_key = file.key if hasattr(file, 'key') else ''
            filename = full_key.split('/')[-1] if full_key else ''
            if not filename and hasattr(file, 'url'):
                filename = file.url.split('/')[-1]
            
            if not filename:
                # If we can't determine filename, include by default
                filtered.append(file)
                continue
            
            # Parse region from filename and full key path
            file_region = self._parse_region_from_filename(filename, full_key)
            
            # Include if matches the desired region
            if file_region == self.region:
                filtered.append(file)
        
        return filtered
    
    def _reorganize_files_by_region(
        self,
        downloaded_files: List[str],
        output_dir: str,
        verbose: bool = False
    ) -> List[str]:
        """
        Reorganize downloaded files into Brain/VNC/Other subfolders.
        
        Original structure: output_dir/Collection/LineName/filename.png
        New structure: output_dir/Brain/LineName/filename.png (or VNC or Other)
        
        Returns list of new file paths.
        """
        import shutil
        from pathlib import Path
        
        reorganized = []
        region_counts = {'Brain': 0, 'VNC': 0, 'Other': 0}
        
        for file_path in downloaded_files:
            file_path = Path(file_path)
            if not file_path.exists():
                continue
            
            filename = file_path.name
            
            # Parse region and line name from filename
            region = self._parse_region_from_filename(filename)
            
            # Extract line name (first field before '-')
            line_name = filename.split('-')[0] if '-' in filename else filename.split('.')[0]
            
            # Create new path: output_dir/Region/LineName/filename
            new_dir = Path(output_dir) / region / line_name
            new_dir.mkdir(parents=True, exist_ok=True)
            new_path = new_dir / filename
            
            try:
                # Move file to new location
                shutil.move(str(file_path), str(new_path))
                reorganized.append(str(new_path))
                region_counts[region] += 1
            except Exception as e:
                # If move fails, try copy
                try:
                    shutil.copy2(str(file_path), str(new_path))
                    file_path.unlink()  # Delete original
                    reorganized.append(str(new_path))
                    region_counts[region] += 1
                except Exception as e2:
                    if verbose:
                        print(f"  ⚠️  Could not reorganize {filename}: {e2}")
                    reorganized.append(str(file_path))  # Keep original
        
        # Clean up empty collection directories
        try:
            for item in Path(output_dir).iterdir():
                if item.is_dir() and item.name not in ('Brain', 'VNC', 'Other'):
                    # Check if empty (or has only empty subdirs)
                    remaining_files = list(item.rglob('*'))
                    remaining_files = [f for f in remaining_files if f.is_file()]
                    if not remaining_files:
                        shutil.rmtree(str(item), ignore_errors=True)
        except Exception:
            pass
        
        if verbose:
            print(f"  📁 Reorganized by region: Brain={region_counts['Brain']}, VNC={region_counts['VNC']}, Other={region_counts['Other']}")
        
        return reorganized
    
    def _download_flylight_images_with_category(
        self,
        lines: List[str],
        output_dir: str,
        formats: Union[str, List[str]],
        image_types: Union[str, List[str]],
        max_files: Optional[int],
        category: Optional[Union[str, List[str]]],
        organize_by_region: bool = False,
        simple_mode: bool = False,
        verbose: bool = False
    ) -> Tuple[List[str], List[str]]:
        """
        Download images from FlyLight with sequential category searching.
        
        This method searches FlyLight collections in priority order, collecting
        images from each category sequentially until ``max_files`` is reached
        for each line. If a line has no images in the specified categories,
        it automatically falls back to searching 'MCFO' collection.
        
        **Category Search Order:**
        
        Categories are searched in the order specified in the ``category`` parameter.
        For example, if ``category=['GAL4/LEXA', 'SplitGAL4']``:
        
        1. First, search 'GAL4/LEXA' collection
        2. If not enough files, search 'SplitGAL4' collection
        3. If still no files found, fallback to 'MCFO' collection
        
        This sequential approach ensures:
        - Preferred collections are prioritized
        - ``max_files`` limit is respected across all categories
        - Lines without images in primary categories get MCFO fallback
        
        Parameters
        ----------
        lines : list of str
            List of driver line names to download images for.
        output_dir : str
            Directory to save downloaded images.
        formats : str or list of str
            File formats to download: 'png', 'jpg', 'h5j', 'mp4', 'all'.
        image_types : str or list of str
            Image types: 'mip', 'cdm', 'aligned', 'translation', 'all'.
        max_files : int, optional
            Maximum images to download per line. Images are collected from
            categories in order until this limit is reached.
        category : str, list of str, or None
            FlyLight collection category(s) to search, in priority order.
            Options: 'GAL4/LEXA', 'SplitGAL4', 'MCFO', 'RawImages', 'All'.
            If a list, categories are searched sequentially in the given order.
            If None, searches all collections.
            
            **Fallback behavior:** If no images are found for a line in the
            specified categories, 'MCFO' is automatically added as a fallback
            (unless already included or category is 'All'/None).
        organize_by_region : bool
            If True, organize images into Brain/VNC subfolders.
        simple_mode : bool
            If True, apply filename filtering to reduce download volume:
            - Split-GAL4 collections: only files with '20x' AND 'multichannel'
            - GAL4/LexA collections: only files with 'total' in filename
            - MCFO: all files kept (no filtering)
            
        Returns
        -------
        tuple of (list, list)
            - downloaded_files: List of downloaded file paths
            - lines_without_files: List of line names that had no FlyLight
              files available even after MCFO fallback
              
        Examples
        --------
        >>> # Search GAL4/LEXA first, then SplitGAL4, with MCFO fallback
        >>> files, missing = finder._download_flylight_images_with_category(
        ...     lines=['SS00731', 'VT000770'],
        ...     output_dir='./images',
        ...     formats=['png', 'jpg'],
        ...     image_types='mip',
        ...     max_files=10,
        ...     category=['GAL4/LEXA', 'SplitGAL4'],  # Priority order
        ...     simple_mode=True
        ... )
        >>> # SS00731 might get images from SplitGAL4
        >>> # VT000770 might get images from MCFO fallback
        """
        try:
            from flylight_downloader import FlyLightDownloader
        except ImportError:
            try:
                # Try relative import
                from .flylight_downloader import FlyLightDownloader
            except ImportError:
                if verbose:
                    print("❌ flylight_downloader module not available")
                return [], lines
        
        def apply_simple_mode_filter(files, collection_name: str):
            """Filter files based on simple_mode rules."""
            if not simple_mode or not files:
                return files
            
            collection_lower = collection_name.lower() if collection_name else ''
            filtered = []
            
            for f in files:
                filename_lower = f.filename.lower()
                
                # Split-GAL4: only include '20x' AND 'multichannel', exclude 'image1'/'image2'
                if 'splitgal4' in collection_lower or 'split-gal4' in collection_lower or 'split_gal4' in collection_lower:
                    if '20x' in filename_lower and 'multichannel' in filename_lower:
                        # Exclude duplicate images (image1, image2, etc.)
                        if 'image1' not in filename_lower and 'image2' not in filename_lower:
                            filtered.append(f)
                # GAL4/LexA: only include 'total'
                elif 'gal4' in collection_lower or 'lexa' in collection_lower:
                    if 'total' in filename_lower:
                        filtered.append(f)
                else:
                    # For other collections (MCFO, etc.), keep all files
                    filtered.append(f)
            
            return filtered
        
        def verify_vt_file_accessible(file) -> tuple:
            """Check if a VT file URL is accessible (HEAD request).
            
            Returns:
                tuple: (is_accessible: bool, error_msg: str or None)
            """
            if file.source != 'http':
                return True, None  # S3 files don't need verification
            try:
                import urllib.request
                req = urllib.request.Request(file.url, method='HEAD')
                with urllib.request.urlopen(req, timeout=5) as resp:
                    return resp.status == 200, None
            except Exception as e:
                # Extract server from URL
                import urllib.parse
                parsed = urllib.parse.urlparse(file.url)
                server = parsed.netloc
                error_msg = f"{server} unreachable: {type(e).__name__}"
                return False, error_msg
        
        def get_files_for_line_sequential(line_name: str, categories: List[str], 
                                          downloader_kwargs: dict, max_files_limit: int,
                                          verify_vt: bool = True) -> List:
            """
            Get files for a line by searching categories sequentially until max_files reached.
            
            Parameters
            ----------
            line_name : str
                The driver line name.
            categories : list of str
                Categories to search in priority order.
            downloader_kwargs : dict
                Base kwargs for FlyLightDownloader.
            max_files_limit : int
                Maximum files to collect.
                
            Returns
            -------
            list
                List of FlyLightFile objects, up to max_files_limit.
            """
            collected_files = []
            
            for cat in categories:
                if len(collected_files) >= max_files_limit:
                    break
                    
                try:
                    # Create downloader for this category
                    dl = FlyLightDownloader(
                        output_dir=downloader_kwargs.get('output_dir', output_dir),
                        collection_category=cat,
                        formats=downloader_kwargs.get('formats', formats),
                        image_types=downloader_kwargs.get('image_types', image_types),
                        verbose=False
                    )
                    
                    files = dl.get_filtered_files(line_name)
                    if files:
                        collection = files[0].collection if files else cat
                        files = apply_simple_mode_filter(files, collection)
                        
                        # Apply region filter (important for MCFO fallback)
                        files = self._filter_flylight_files_by_region(files)
                        
                        # For VT GAL4 files (HTTP source), verify the server is accessible
                        # This catches cases where flimg.janelia.org is down
                        if verify_vt and files and files[0].source == 'http':
                            is_accessible, error_msg = verify_vt_file_accessible(files[0])
                            if not is_accessible:
                                # Server down - collect warning and skip files for MCFO fallback
                                collection_name = files[0].collection or cat
                                warning_msg = f"⚠️  {collection_name} server not accessible for {line_name}: {error_msg}"
                                self._warning_collector.append(warning_msg)
                                continue
                        
                        # Add files up to the limit
                        remaining = max_files_limit - len(collected_files)
                        collected_files.extend(files[:remaining])
                except Exception:
                    continue
            
            return collected_files
        
        downloaded = []
        
        # Normalize category to list
        if category is None:
            categories = ['All']
        elif isinstance(category, str):
            categories = [category]
        else:
            categories = list(category)
        
        # Check if MCFO fallback is needed (not already included and not searching all)
        categories_lower = [c.lower() for c in categories]
        needs_mcfo_fallback = ('mcfo' not in categories_lower and 
                               'all' not in categories_lower)
        
        # Separate VT lines from other lines (they need different format settings)
        vt_lines = [l for l in lines if l.upper().startswith('VT')]
        other_lines = [l for l in lines if not l.upper().startswith('VT')]
        
        # Phase 1: Scan all lines to count total files (with sequential category search)
        if verbose:
            cat_str = ' → '.join(categories)
            if needs_mcfo_fallback:
                cat_str += ' → MCFO (fallback)'
            if simple_mode:
                print(f"📊 Scanning {len(lines)} lines (simple mode, categories: {cat_str})...")
            else:
                print(f"📊 Scanning {len(lines)} lines (categories: {cat_str})...")
        
        line_files_map = {}  # line_name -> list of FlyLightFile
        total_file_count = 0
        max_files_limit = max_files if max_files else 999999
        
        # Create a scanning progress bar
        scan_pbar = None
        if HAS_TQDM and verbose:
            scan_pbar = tqdm(total=len(lines), desc="  Scanning", unit="line", leave=False)
        
        # Scan other lines (non-VT)
        for line_name in other_lines:
            try:
                # Sequential category search
                files = get_files_for_line_sequential(
                    line_name, 
                    categories,
                    {'formats': formats, 'image_types': image_types},
                    max_files_limit
                )
                
                # MCFO fallback if no files found and fallback is enabled
                if not files and needs_mcfo_fallback:
                    warning_msg = f"ℹ️  No files from {', '.join(categories)} for {line_name}, trying MCFO fallback..."
                    self._warning_collector.append(warning_msg)
                    
                    files = get_files_for_line_sequential(
                        line_name,
                        ['MCFO'],
                        {'formats': formats, 'image_types': image_types},
                        max_files_limit
                    )
                    # Only print success message if no progress bar (otherwise it interrupts the bar)
                    if files and verbose and not scan_pbar:
                        print(f"  ✓ Found {len(files)} MCFO images for {line_name}")
                    elif not files:
                        warning_msg = f"⚠️  No images found for {line_name} in any collection (including MCFO)"
                        self._warning_collector.append(warning_msg)
                
                line_files_map[line_name] = files
                total_file_count += len(files)
            except Exception:
                line_files_map[line_name] = []
            
            if scan_pbar:
                scan_pbar.set_postfix(files=total_file_count, refresh=False)
                scan_pbar.update(1)
        
        # Scan VT lines (need different format settings)
        if vt_lines:
            vt_formats = formats
            if isinstance(vt_formats, str):
                vt_formats = [vt_formats]
            if 'png' in vt_formats and 'jpg' not in vt_formats:
                vt_formats = list(vt_formats) + ['jpg']
            
            for line_name in vt_lines:
                try:
                    # Sequential category search for VT lines
                    files = get_files_for_line_sequential(
                        line_name,
                        categories,
                        {'formats': vt_formats, 'image_types': 'all'},
                        max_files_limit
                    )
                    
                    # MCFO fallback if no files found
                    if not files and needs_mcfo_fallback:
                        warning_msg = f"ℹ️  No files from {', '.join(categories)} for {line_name}, trying MCFO fallback..."
                        self._warning_collector.append(warning_msg)
                        
                        files = get_files_for_line_sequential(
                            line_name,
                            ['MCFO'],
                            {'formats': vt_formats, 'image_types': 'all'},
                            max_files_limit
                        )
                        # Only print success message if no progress bar (otherwise it interrupts the bar)
                        if files and verbose and not scan_pbar:
                            print(f"  ✓ Found {len(files)} MCFO images for {line_name}")
                        elif not files:
                            warning_msg = f"⚠️  No images found for {line_name} in any collection (including MCFO)"
                            self._warning_collector.append(warning_msg)
                    
                    line_files_map[line_name] = files
                    total_file_count += len(files)
                except Exception:
                    line_files_map[line_name] = []
                
                if scan_pbar:
                    scan_pbar.set_postfix(files=total_file_count, refresh=False)
                    scan_pbar.update(1)
        
        if scan_pbar:
            scan_pbar.close()
        
        # Track lines without any FlyLight files
        lines_without_files = [l for l, f in line_files_map.items() if not f]
        
        if total_file_count == 0:
            if verbose:
                print("  ⚠️ No files found for any lines (tried all categories including MCFO)")
            return [], lines
        
        lines_with_files = len([l for l, f in line_files_map.items() if f])
        if verbose:
            print(f"  📦 Found {total_file_count} files across {lines_with_files} lines")
            if lines_without_files:
                print(f"  ⚠️ No files for: {', '.join(lines_without_files[:5])}{'...' if len(lines_without_files) > 5 else ''}")
        
        # Phase 2: Download with combined progress bar
        downloaded = []
        files_downloaded = [0]  # Use list to allow modification in callback
        current_line = ['']  # Track current line for progress display
        
        # Create a single progress bar for all files
        pbar = None
        if HAS_TQDM and verbose:
            pbar = tqdm(total=total_file_count, desc="  Downloading", unit="file")
        
        def on_file_downloaded(file_path, line_name):
            """Callback to update progress after each file download."""
            files_downloaded[0] += 1
            if pbar:
                # Update postfix to show current line
                if line_name != current_line[0]:
                    current_line[0] = line_name
                pbar.set_postfix(line=line_name[:15], refresh=False)
                pbar.update(1)
            elif verbose:
                # Fallback without tqdm
                print(f"  [{files_downloaded[0]}/{total_file_count}] {line_name}: {file_path.name if hasattr(file_path, 'name') else file_path}")
        
        # Download all lines using pre-scanned files
        for line_name in lines:
            if line_name in line_files_map and line_files_map[line_name]:
                try:
                    # Determine the correct formats for this line
                    line_formats = formats
                    line_image_types = image_types
                    if line_name.upper().startswith('VT'):
                        if isinstance(line_formats, str):
                            line_formats = [line_formats]
                        if 'png' in line_formats and 'jpg' not in line_formats:
                            line_formats = list(line_formats) + ['jpg']
                        line_image_types = 'all'
                    
                    # Create downloader (category doesn't matter since we're using pre-scanned files)
                    downloader = FlyLightDownloader(
                        output_dir=output_dir,
                        formats=line_formats,
                        image_types=line_image_types,
                        verbose=False
                    )
                    
                    dl_files = downloader.download(
                        line_name=line_name,
                        max_files=max_files,
                        on_file_downloaded=on_file_downloaded,
                        flat_structure=True,
                        add_timestamp=False,  # Don't add timestamp - organize by line name only for summary generation
                        files=line_files_map[line_name]  # Use pre-filtered files
                    )
                    downloaded.extend([str(f) for f in dl_files])
                except Exception as e:
                    if verbose and not pbar:
                        print(f"  ❌ {line_name}: {e}")
        
        if pbar:
            pbar.close()
        
        if verbose:
            print(f"  ✅ Downloaded {len(downloaded)}/{total_file_count} files")
        
        # Print warning summary
        if verbose:
            self._print_warning_summary()
        
        # Reorganize files by anatomical region if requested
        if organize_by_region and downloaded:
            if verbose:
                print("  🔄 Reorganizing files by anatomical region...")
            downloaded = self._reorganize_files_by_region(
                downloaded_files=downloaded,
                output_dir=output_dir,
                verbose=verbose
            )
        
        return downloaded, lines_without_files
    
    def find_lines_batch_with_images(
        self,
        queries: Union[str, int, List],
        dataset: Optional[str] = None,
        top_n: int = -1,
        match_type: str = 'cds',
        output_dir: Optional[str] = None,
        download_images: Union[bool, str] = False,
        image_source: str = 'neuronbridge',
        image_formats: Union[str, List[str]] = 'png',
        image_types: Union[str, List[str]] = 'cdm',
        max_download_images_per_line: Optional[int] = 5
    ) -> pd.DataFrame:
        """
        Find driver lines with optional image download.
        
        .. deprecated::
            This method is deprecated. Use find_lines_batch() with download_images parameter instead.
        
        Parameters
        ----------
        queries : str, int, or list
            Neuron query(s). Can be bodyId, type, instance, or comma-separated string.
        dataset : str, optional
            Dataset for type/instance lookups.
        top_n : int
            Maximum matches per neuron. Default: -1 (all matches)
        match_type : str
            Match algorithm: 'cds', 'pppm', or 'both'. Default: 'cds'
        output_dir : str, optional
            Directory to save results and images.
        download_images : bool or str
            Whether to download images. True uses image_source, or pass source directly.
        image_source : str
            'neuronbridge' for CDM images, 'flylight' for raw confocal. Default: 'neuronbridge'
        image_formats : str or list
            File formats for images. Default: 'png'
        image_types : str or list
            Image types to download. Default: 'cdm'
        max_download_images_per_line : int, optional
            Maximum images to download per line. Default: 5
            
        Returns
        -------
        pd.DataFrame
            Combined results DataFrame.
        """
        import warnings
        warnings.warn(
            "find_lines_batch_with_images() is deprecated. Use find_lines_batch() with download_images parameter.",
            DeprecationWarning,
            stacklevel=2
        )
        
        # Convert legacy parameters to new format
        if download_images:
            if isinstance(download_images, bool):
                dl_source = image_source
            else:
                dl_source = download_images
        else:
            dl_source = None
        
        return self.find_lines_batch(
            queries=queries,
            dataset=dataset,
            match_type=match_type,
            output_dir=output_dir,
            download_images=dl_source,
            image_formats=image_formats,
            image_types=image_types,
            max_download_images_per_line=max_download_images_per_line
        )


def create_image_pdf(
    images_dir: str,
    output_pdf: Optional[str] = None,
    images_per_page: Tuple[int, int] = (5, 3),
    page_size: str = 'A4',
    landscape: bool = True,
    title_font_size: int = 14,
    margin: float = 0.5,
    line_order: Optional[List[str]] = None,
    verbose: bool = True,
    background_color: Union[str, Tuple[float, float, float]] = None
) -> Optional[str]:
    """
    Create a PDF file from downloaded images, organized by line name.
    
    Each page shows images for one driver line with the line name as title.
    Images are arranged in a grid (default 5 columns x 3 rows = 15 images per page).
    If a line has more images than fit on one page, additional pages are created.
    
    Parameters
    ----------
    images_dir : str
        Directory containing images, organized by line name subdirectories.
        Expected structure: images_dir/LineName/*.png (or .jpg)
    output_pdf : str, optional
        Path for output PDF file. If None, saves to images_dir/images_summary.pdf
    images_per_page : tuple of (int, int)
        (columns, rows) - number of images per page. Default: (5, 3) = 15 images
    page_size : str
        Page size: 'A4', 'Letter', etc. Default: 'A4'
    landscape : bool
        Use landscape orientation. Default: True (horizontal A4)
    title_font_size : int
        Font size for line name title. Default: 14
    margin : float
        Page margin in inches. Default: 0.5
    line_order : list of str, optional
        Ordered list of line names for page ordering. Lines are sorted in this order,
        preserving the ranking used to select top-N lines (e.g., by weighted_score).
        Lines not in this list are appended at the end alphabetically.
        If None, lines are sorted alphabetically.
    verbose : bool
        Print progress messages. Default: True
    background_color : str or tuple, optional
        Background color for PDF pages. Can be named color ('black', 'white'),
        hex string ('#000000'), or RGB tuple (0-1). Default: None (white)
        
    Returns
    -------
    str or None
        Path to the created PDF file, or None if creation failed.
        
    Example
    -------
    >>> create_image_pdf(
    ...     images_dir='/path/to/images',
    ...     output_pdf='/path/to/summary.pdf',
    ...     images_per_page=(5, 3),
    ...     landscape=True,
    ...     background_color='black'
    ... )
    """
    try:
        from reportlab.lib.pagesizes import A4, letter, landscape as rl_landscape
        from reportlab.lib.units import inch
        from reportlab.pdfgen import canvas
        from reportlab.lib.utils import ImageReader
        from reportlab.lib import colors
        from PIL import Image
        HAS_REPORTLAB = True
    except ImportError:
        HAS_REPORTLAB = False
        
    if not HAS_REPORTLAB:
        if verbose:
            print("⚠️  PDF generation requires reportlab and Pillow.")
            print("   Install with: pip install reportlab Pillow")
        return None
    
    from pathlib import Path
    
    images_path = Path(images_dir)
    if not images_path.exists():
        if verbose:
            print(f"⚠️  Images directory not found: {images_dir}")
        return None
    
    # Find all line directories with images
    line_images = {}
    image_extensions = {'.png', '.jpg', '.jpeg'}
    
    for item in sorted(images_path.iterdir()):
        if item.is_dir():
            # Line name subdirectory
            line_name = item.name
            images = sorted([
                f for f in item.iterdir()
                if f.suffix.lower() in image_extensions
            ])
            if images:
                line_images[line_name] = images
        elif item.suffix.lower() in image_extensions:
            # Images directly in the images folder (no subdirectory)
            # Group by line name prefix (before first '-')
            line_name = item.stem.split('-')[0] if '-' in item.stem else 'Unknown'
            if line_name not in line_images:
                line_images[line_name] = []
            line_images[line_name].append(item)
    
    if not line_images:
        if verbose:
            print(f"⚠️  No images found in: {images_dir}")
        return None
    
    # Sort images within each line
    for line_name in line_images:
        line_images[line_name] = sorted(line_images[line_name])
    
    # Determine line ordering for PDF pages
    if line_order:
        # Use provided order (preserves ranking like weighted_score)
        # Lines in line_order come first in that order, others appended alphabetically
        ordered_lines = []
        remaining_lines = set(line_images.keys())
        for line in line_order:
            if line in remaining_lines:
                ordered_lines.append(line)
                remaining_lines.remove(line)
        # Append any remaining lines alphabetically
        ordered_lines.extend(sorted(remaining_lines))
    else:
        # Default: alphabetical order
        ordered_lines = sorted(line_images.keys())
    
    # Set output path
    if output_pdf is None:
        output_pdf = images_path / 'images_summary.pdf'
    output_pdf = Path(output_pdf)
    
    # Page setup
    if page_size.upper() == 'A4':
        base_size = A4
    else:
        base_size = letter
    
    if landscape:
        page_width, page_height = rl_landscape(base_size)
    else:
        page_width, page_height = base_size
    
    cols, rows = images_per_page
    margin_pts = margin * inch
    
    # Calculate available space for images
    usable_width = page_width - 2 * margin_pts
    usable_height = page_height - 2 * margin_pts - 30  # 30 pts for title
    
    # Calculate cell size
    cell_width = usable_width / cols
    cell_height = usable_height / rows
    
    # Create PDF
    c = canvas.Canvas(str(output_pdf), pagesize=(page_width, page_height))
    
    total_pages = 0
    total_images = 0
    
    if verbose:
        if line_order:
            print(f"📄 Creating PDF from {len(line_images)} lines (ordered by ranking)...")
        else:
            print(f"📄 Creating PDF from {len(line_images)} lines (alphabetical order)...")
    
    for line_name in ordered_lines:
        images = line_images[line_name]
        # Calculate how many pages needed for this line
        images_per_full_page = cols * rows
        num_pages = (len(images) + images_per_full_page - 1) // images_per_full_page
        
        for page_idx in range(num_pages):
            # Apply background color
            text_color = colors.black
            if background_color:
                try:
                    # Parse color
                    if isinstance(background_color, str):
                        if background_color.startswith('#'):
                            bg_col = colors.HexColor(background_color)
                        else:
                            bg_col = getattr(colors, background_color, colors.white)
                    elif isinstance(background_color, (tuple, list)) and len(background_color) >= 3:
                        bg_col = colors.Color(*background_color[:3])
                    else:
                        bg_col = colors.white
                    
                    c.setFillColor(bg_col)
                    c.rect(0, 0, page_width, page_height, fill=1, stroke=0)
                    
                    # Set text color contrast
                    # Simple luminance check
                    if hasattr(bg_col, 'red'):
                        luminance = bg_col.red*0.299 + bg_col.green*0.587 + bg_col.blue*0.114
                        if luminance < 0.5:
                             text_color = colors.white
                except Exception:
                    pass

            # Draw title
            c.setFillColor(text_color)
            c.setFont("Helvetica-Bold", title_font_size)
            title = line_name
            if num_pages > 1:
                title += f" ({page_idx + 1}/{num_pages})"
            c.drawCentredString(page_width / 2, page_height - margin_pts - 5, title)
            
            # Get images for this page
            start_idx = page_idx * images_per_full_page
            end_idx = min(start_idx + images_per_full_page, len(images))
            page_images = images[start_idx:end_idx]
            
            # Calculate optimal layout for this page
            num_images = len(page_images)
            if num_images < images_per_full_page:
                # Auto-arrange for fewer images
                # Try to fill rows as much as possible
                actual_rows = (num_images + cols - 1) // cols
                if actual_rows < rows:
                    # Recalculate cell height to use more space
                    actual_cell_height = usable_height / actual_rows
                else:
                    actual_cell_height = cell_height
            else:
                actual_rows = rows
                actual_cell_height = cell_height
            
            # Draw images
            for i, img_path in enumerate(page_images):
                row = i // cols
                col = i % cols
                
                # Calculate position (top-left corner of cell)
                x = margin_pts + col * cell_width
                y = page_height - margin_pts - 25 - (row + 1) * actual_cell_height
                
                try:
                    # Open image to get dimensions
                    with Image.open(img_path) as img:
                        img_width, img_height = img.size
                        
                        # Calculate scaling to fit in cell with padding
                        padding = 5  # pixels padding
                        max_width = cell_width - 2 * padding
                        max_height = actual_cell_height - 2 * padding
                        
                        scale_w = max_width / img_width
                        scale_h = max_height / img_height
                        scale = min(scale_w, scale_h)
                        
                        draw_width = img_width * scale
                        draw_height = img_height * scale
                        
                        # Center in cell
                        draw_x = x + (cell_width - draw_width) / 2
                        draw_y = y + (actual_cell_height - draw_height) / 2
                        
                        # Draw the image
                        c.drawImage(
                            str(img_path),
                            draw_x, draw_y,
                            width=draw_width,
                            height=draw_height,
                            preserveAspectRatio=True
                        )
                        total_images += 1
                        
                except Exception as e:
                    # Skip problematic images
                    if verbose:
                        print(f"   ⚠️  Could not process: {img_path.name} - {e}")
            
            # Add new page
            c.showPage()
            total_pages += 1
    
    # Save PDF
    c.save()
    
    if verbose:
        print(f"✅ Created PDF: {output_pdf}")
        print(f"   {total_pages} pages, {total_images} images from {len(line_images)} lines")
    
    return str(output_pdf)


def create_image_pptx(
    images_dir: str,
    output_pptx: Optional[str] = None,
    images_per_slide: Tuple[int, int] = (5, 3),
    slide_size: str = 'widescreen',
    title_font_size: int = 24,
    label_font_size: int = 20,
    margin: float = 0.3,
    line_order: Optional[List[str]] = None,
    font_color: Tuple[int, int, int] = (0, 0, 0),
    verbose: bool = True
) -> Optional[str]:
    """
    Create a PPTX file from downloaded images, organized by line name.
    
    Each slide shows images for one driver line with the line name as title.
    Images are arranged in a grid (default 5 columns x 3 rows = 15 images per slide).
    If a line has more images than fit on one slide, additional slides are created.
    
    Parameters
    ----------
    images_dir : str
        Directory containing images, organized by line name subdirectories.
        Expected structure: images_dir/LineName/*.png (or .jpg)
    output_pptx : str, optional
        Path for output PPTX file. If None, saves to images_dir/images_summary.pptx
    images_per_slide : tuple of (int, int)
        (columns, rows) - number of images per slide. Default: (5, 3) = 15 images
    slide_size : str
        Slide dimensions: 'widescreen' (16:9), 'standard' (4:3), 'a4'. Default: 'widescreen'
    title_font_size : int
        Font size for line name title. Default: 24
    label_font_size : int
        Font size for image labels. Default: 20
    margin : float
        Slide margin in inches. Default: 0.3
    line_order : list of str, optional
        Ordered list of line names for slide ordering. Lines are sorted in this order,
        preserving the ranking used to select top-N lines (e.g., by weighted_score).
        Lines not in this list are appended at the end alphabetically.
        If None, lines are sorted alphabetically.
    font_color : tuple of (int, int, int), default (0, 0, 0)
        RGB color tuple for label text (r, g, b), each value 0-255. Default is black.
    verbose : bool
        Print progress messages. Default: True
        
    Returns
    -------
    str or None
        Path to the created PPTX file, or None if creation failed.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        HAS_PPTX = True
    except ImportError:
        HAS_PPTX = False
    
    if not HAS_PPTX:
        if verbose:
            print("⚠️  PPTX generation requires python-pptx.")
            print("   Install with: pip install python-pptx")
        return None
    
    from pathlib import Path
    from PIL import Image
    
    images_path = Path(images_dir)
    if not images_path.exists():
        if verbose:
            print(f"⚠️  Images directory not found: {images_dir}")
        return None
    
    # Find all line directories with images
    line_images = {}
    image_extensions = {'.png', '.jpg', '.jpeg'}
    
    for item in sorted(images_path.iterdir()):
        if item.is_dir():
            line_name = item.name
            images = sorted([
                f for f in item.iterdir()
                if f.suffix.lower() in image_extensions
            ])
            if images:
                line_images[line_name] = images
        elif item.suffix.lower() in image_extensions:
            line_name = item.stem.split('-')[0] if '-' in item.stem else 'Unknown'
            if line_name not in line_images:
                line_images[line_name] = []
            line_images[line_name].append(item)
    
    if not line_images:
        if verbose:
            print(f"⚠️  No images found in: {images_dir}")
        return None
    
    # Sort images within each line
    for line_name in line_images:
        line_images[line_name] = sorted(line_images[line_name])
    
    # Determine line ordering
    if line_order:
        ordered_lines = []
        remaining_lines = set(line_images.keys())
        for line in line_order:
            if line in remaining_lines:
                ordered_lines.append(line)
                remaining_lines.remove(line)
        ordered_lines.extend(sorted(remaining_lines))
    else:
        ordered_lines = sorted(line_images.keys())
    
    # Set output path
    if output_pptx is None:
        output_pptx = images_path / 'images_summary.pptx'
    output_pptx = Path(output_pptx)
    
    # Slide setup
    size_presets = {
        'widescreen': (13.333, 7.5),
        'standard': (10, 7.5),
        'a4': (11.69, 8.27),
    }
    slide_width, slide_height = size_presets.get(slide_size, size_presets['widescreen'])
    
    cols, rows = images_per_slide
    title_height = 0.5
    label_height = (label_font_size / 72) * 1.5
    
    # Calculate cell dimensions
    usable_width = slide_width - 2 * margin
    usable_height = slide_height - margin - title_height - margin
    cell_width = usable_width / cols
    cell_height = usable_height / rows
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    blank_layout = prs.slide_layouts[6]
    
    total_slides = 0
    total_images = 0
    
    if verbose:
        if line_order:
            print(f"📊 Creating PPTX from {len(line_images)} lines (ordered by ranking)...")
        else:
            print(f"📊 Creating PPTX from {len(line_images)} lines (alphabetical order)...")
    
    for line_name in ordered_lines:
        images = line_images[line_name]
        images_per_full_slide = cols * rows
        num_slides = (len(images) + images_per_full_slide - 1) // images_per_full_slide
        
        for slide_idx in range(num_slides):
            slide = prs.slides.add_slide(blank_layout)
            
            # Build title
            slide_title = line_name
            if num_slides > 1:
                slide_title += f" ({slide_idx + 1}/{num_slides})"
            
            # Add title
            txBox = slide.shapes.add_textbox(
                Inches(margin),
                Inches(margin / 2),
                Inches(slide_width - 2 * margin),
                Inches(title_height)
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(title_font_size)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # Get images for this slide
            start_idx = slide_idx * images_per_full_slide
            end_idx = min(start_idx + images_per_full_slide, len(images))
            slide_images = images[start_idx:end_idx]
            
            content_top = margin + title_height
            
            for i, img_path in enumerate(slide_images):
                row = i // cols
                col = i % cols
                
                cell_left = margin + col * cell_width
                cell_top = content_top + row * cell_height
                
                try:
                    with Image.open(img_path) as img:
                        img_width, img_height = img.size
                        
                        max_width = cell_width - 0.1
                        max_height = cell_height - label_height - 0.1
                        
                        scale_w = max_width / (img_width / 96)
                        scale_h = max_height / (img_height / 96)
                        scale_factor = min(scale_w, scale_h)
                        
                        final_width = (img_width / 96) * scale_factor
                        final_height = (img_height / 96) * scale_factor
                        
                        img_left = cell_left + (cell_width - final_width) / 2
                        img_top = cell_top + (cell_height - label_height - final_height) / 2
                        
                        slide.shapes.add_picture(
                            str(img_path),
                            Inches(img_left),
                            Inches(img_top),
                            Inches(final_width),
                            Inches(final_height)
                        )
                        
                        # Add label
                        label = img_path.stem
                        max_chars = int(cell_width * 8)
                        if len(label) > max_chars:
                            label = label[:max_chars-3] + '...'
                        
                        label_box = slide.shapes.add_textbox(
                            Inches(cell_left),
                            Inches(cell_top + cell_height - label_height),
                            Inches(cell_width),
                            Inches(label_height)
                        )
                        tf = label_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = label
                        p.font.size = Pt(label_font_size)
                        from pptx.dml.color import RGBColor
                        p.font.color.rgb = RGBColor(*font_color)
                        p.alignment = PP_ALIGN.CENTER
                        
                        total_images += 1
                        
                except Exception as e:
                    if verbose:
                        print(f"   ⚠️  Could not process: {img_path.name} - {e}")
            
            total_slides += 1
    
    prs.save(str(output_pptx))
    
    if verbose:
        print(f"✅ Created PPTX: {output_pptx}")
        print(f"   {total_slides} slides, {total_images} images from {len(line_images)} lines")
    
    return str(output_pptx)


# Convenience function for quick usage
def find_lines_for_body(body_id: int, match_type: str = 'cds') -> pd.DataFrame:
    """
    Quick function to find driver lines matching a body ID.
    
    Parameters
    ----------
    body_id : int
        The EM body ID to search for.
    match_type : str
        Match algorithm: 'cds', 'pppm', or 'both'. Default: 'cds'
        
    Returns
    -------
    pd.DataFrame
        DataFrame of matching lines, sorted by score (cds/pppm) or combined_rank (both).
    """
    nbf = NeuronBridgeFinder(verbose=False)
    return nbf.id_to_lines(body_id, match_type=match_type)


def find_neurons_for_line(line_name: str, match_type: str = 'cds') -> pd.DataFrame:
    """
    Quick function to find neurons matching a driver line.
    
    Parameters
    ----------
    line_name : str
        The driver line name to search for.
    match_type : str
        Match algorithm: 'cds', 'pppm', or 'both'. Default: 'cds'
        
    Returns
    -------
    pd.DataFrame
        DataFrame of matching neurons, sorted by score (cds/pppm) or combined_rank (both).
    """
    nbf = NeuronBridgeFinder(verbose=False)
    return nbf.line_to_neuron(line_name, match_type=match_type)


if __name__ == '__main__':
    # Simple test
    print("Testing NeuronBridgeFinder...")
    
    # Initialize
    nbf = NeuronBridgeFinder()
    
    # Test id_to_lines
    print("\n--- Test id_to_lines ---")
    lines = nbf.id_to_lines(636798093)
    print(lines.head(10))
    
    # Test line_to_neuron
    print("\n--- Test line_to_neuron ---")
    neurons = nbf.line_to_neuron('LH173')
    print(neurons.head(10))
